import streamlit as st
import os
import json
import asyncio
import base64
import pandas as pd
import pymupdf  # PyMuPDF
from typing import List, Dict, Optional, Union
from pydantic import BaseModel, Field, ValidationError, field_validator, ConfigDict
from google import genai
from google.genai import types
import docx
from pptx import Presentation
import csv
import openpyxl
import re
from collections import Counter, defaultdict
import time
import io
import datetime

# Configure page
st.set_page_config(
    page_title="Multimodal Document Processor",
    page_icon="📊",
    layout="wide"
)

###############################
# DATA MODELS
###############################

class VisualElement(BaseModel):
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data")
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")

class NumericalDataPoint(BaseModel):
    value: str = Field(..., description="The numerical value")
    description: str = Field(default="", description="What the number represents")
    context: str = Field(default="", description="Surrounding text context")
    model_config = ConfigDict(arbitrary_types_allowed=True, extra="ignore")
    @field_validator('value', 'description', 'context', mode='before')
    def convert_to_string(cls, v): return str(v) if v is not None else ""

class TableData(BaseModel):
    table_content: str = Field(default="", description="Markdown formatted table content")
    title: Optional[str] = None
    summary: Optional[str] = None
    page_number: int = 0
    model_config = ConfigDict(arbitrary_types_allowed=True)
    @field_validator('table_content')
    def clean_table_content(cls, v): return str(v) if v is not None else ""

class FinancialMetric(BaseModel):
    name: str = Field(..., description="Name of the financial metric")
    value: str = Field(..., description="Metric value")
    period: Optional[str] = Field(None, description="Reporting period")
    context: Optional[str] = Field(None, description="Surrounding text context")
    trend: Optional[str] = Field(None, description="Trend direction if available")

class FinancialStatement(BaseModel):
    type: str = Field(..., description="Type of statement")
    period: str = Field(..., description="Reporting period")
    metrics: List[FinancialMetric] = Field(default_factory=list, description="Key metrics")

class PageContent(BaseModel):
    model_config = ConfigDict(arbitrary_types_allowed=True)
    page_number: int = Field(..., description="Page number in document")
    text: str = Field("", description="Full text content")
    title: str = Field("Untitled", description="Brief title for this page")
    topics: List[str] = Field(default_factory=list, description="Key topics discussed")
    summary: str = Field("", description="Summary of key points")
    entities: List[str] = Field(default_factory=list, description="Important entities mentioned")
    has_tables: bool = Field(False, description="True if page contains tables")
    has_visuals: bool = Field(False, description="True if page contains visuals")
    has_numbers: bool = Field(False, description="True if page contains key numerical data")
    tables: List[Union[TableData, dict]] = Field(default_factory=list, description="Extracted tables")
    visuals: List[Union[VisualElement, dict]] = Field(default_factory=list, description="Extracted visual elements")
    numbers: List[Union[NumericalDataPoint, dict]] = Field(default_factory=list, description="Extracted numerical data")
    dates: List[str] = Field(default_factory=list, description="Important dates mentioned")
    financial_statements: List[Union[FinancialStatement, dict]] = Field(default_factory=list)
    key_metrics: List[Union[FinancialMetric, dict]] = Field(default_factory=list)
    financial_terms: List[str] = Field(default_factory=list)

def convert_to_model(data, model_class):
    """Convert dictionary data to a Pydantic model with error handling."""
    if not isinstance(data, dict):
        print(f"Warning: Expected dict for {model_class.__name__}, got {type(data)}")
        return None
    
    try:
        # Extract only the fields that the model expects
        model_fields = model_class.__annotations__.keys()
        filtered_data = {k: v for k, v in data.items() if k in model_fields}
        return model_class(**filtered_data)
    except Exception as e:
        print(f"Error converting to {model_class.__name__}: {str(e)}")
        return None
            
class DocumentSummary(BaseModel):
    """Document-level summary."""
    title: str = Field(..., description="Concise title/summary")
    themes: List[str] = Field(..., description="Concept/theme tags")
    questions: List[str] = Field(..., description="Hypothetical questions")
    summary: str = Field(..., description="Comprehensive summary")
    tables_summary: Optional[str] = Field(None, description="Summary of key tables")
    visuals_summary: Optional[str] = Field(None, description="Summary of key visuals")

class KeyTopic(BaseModel):
    """Represents a key topic in the document."""
    name: str = Field(..., description="Short topic name (1-3 words)")
    description: str = Field(..., description="Brief description of the topic")
    relevance: str = Field(..., description="Relevance level (High/Medium/Low)")
    sentiment: str = Field(..., description="Sentiment analysis (Positive/Neutral/Mixed/Cautionary/Negative)")
    analysis: str = Field(..., description="Brief analysis of the topic")

class QuotedStatement(BaseModel):
    """Represents an important quoted statement in the document."""
    speaker: str = Field(..., description="Person or entity who made the statement")
    quote: str = Field(..., description="The quoted text")
    page: int = Field(..., description="Page number where the quote appears")

class DocumentReport(BaseModel):
    """UI-friendly document report structure for presentation."""
    file_name: str = Field(..., description="Original filename")
    page_count: int = Field(..., description="Number of pages in document")
    title: str = Field(..., description="Document title")
    title_summary: str = Field(..., description="Brief summary/subtitle")
    concept_theme_hashtags: List[str] = Field(default_factory=list, description="Theme tags as hashtags")
    date_published: Optional[str] = Field(None, description="Publication date if available")
    source: Optional[str] = Field(None, description="Document source")
    confidence: str = Field("Medium", description="Analysis confidence (High/Medium/Low)")
    document_summary: str = Field(..., description="Comprehensive summary")
    key_insights: List[str] = Field(default_factory=list, description="Key takeaways and insights")
    key_topics: List[KeyTopic] = Field(default_factory=list, description="Detailed topic analysis")
    quoted_statements: List[QuotedStatement] = Field(default_factory=list, description="Important quotes from document")
    content_excerpt: Optional[str] = Field(None, description="Highlighted document content")

class ProcessedDocument(BaseModel):
    """Fully processed document with all content."""
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages")
    summary: Optional[DocumentSummary] = None

###############################
# API AND UTILITY FUNCTIONS
###############################

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
        st.stop()
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)
            
            # Validate JSON if applicable
            config = kwargs.get('config')
            if (config and hasattr(config, 'response_mime_type') and 
                config.response_mime_type == 'application/json' and 
                hasattr(response, 'candidates') and response.candidates):
                try:
                    # Try to parse the JSON response
                    json_text = response.candidates[0].content.parts[0].text
                    json.loads(clean_json_response(json_text))  # Use our clean function first
                except json.JSONDecodeError as e:
                    if attempt < max_retries - 1:
                        print(f"Received malformed JSON on attempt {attempt+1}, retrying: {e}")
                        await asyncio.sleep(2 ** attempt)
                        continue
            
            return response
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            print(f"API call failed on attempt {attempt+1}, retrying: {e}")
            await asyncio.sleep(2 ** attempt)

def clean_json_response(json_text: str, extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        return "{}"
    
    # Ensure it's a string
    json_text = str(json_text)
    
    try:
        # Handle markdown code blocks
        if json_text.startswith("```"):
            blocks = json_text.split("```")
            if len(blocks) >= 3:
                json_text = blocks[1]
                if json_text.startswith("json"):
                    json_text = json_text[4:].strip()
            else:
                json_text = json_text.replace("```", "").strip()
        
        elif json_text.startswith("json"):
            json_text = json_text[4:].strip()
        
        # Remove any leading/trailing whitespace
        json_text = json_text.strip()
        
        # Attempt to parse the JSON as is
        try:
            json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # Continue with fixes
            pass
        
        # Apply all the regex fixes and other cleanup methods as before...
        # [Your existing JSON cleanup code here]
        
        # Final attempt to parse with Python's JSON parser
        try:
            data = json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # If we're extracting text on failure and can't fix the JSON
            if extract_text_on_failure:
                print("JSON parsing failed - extracting text content instead")
                
                # Extract what appears to be the main text content
                # Look for patterns that would indicate text fields in a JSON response
                text_content = ""
                
                # Try to find text between quotes after "text": or "content": patterns
                text_matches = re.findall(r'"(?:text|content)"\s*:\s*"([^"]+)"', json_text)
                if text_matches:
                    text_content = " ".join(text_matches)
                
                # If that doesn't work, take any text between quotes as potential content
                if not text_content:
                    # Get all quoted strings that look reasonably long (over 20 chars)
                    quoted_text = re.findall(r'"([^"]{20,})"', json_text)
                    if quoted_text:
                        text_content = " ".join(quoted_text)
                
                # If still nothing, just take the longest line that's not mostly symbols
                if not text_content:
                    lines = json_text.split('\n')
                    content_lines = [line for line in lines if len(line) > 50 and 
                                    sum(c.isalpha() or c.isspace() for c in line) / len(line) > 0.7]
                    if content_lines:
                        text_content = " ".join(content_lines)
                
                # Fallback - create a minimal valid JSON with the extracted text
                fallback_json = {
                    "text": text_content or "Extracted text content not found",
                    "title": "Text Extraction Fallback",
                    "topics": ["text extraction"],
                    "summary": "Structured JSON parsing failed - basic text extraction used as fallback.",
                    "entities": [],
                    "has_tables": False,
                    "has_visuals": False,
                    "has_numbers": False,
                    "dates": [],
                    "tables": [],
                    "visuals": [],
                    "numbers": [],
                    "financial_statements": [],
                    "key_metrics": [],
                    "financial_terms": []
                }
                return json.dumps(fallback_json)
            else:
                # Return minimal valid JSON as before
                return '{"text":"Error parsing JSON response","error":"Invalid JSON structure"}'
            
    except Exception as e:
        print(f"Error cleaning JSON response: {e}")
        
        if extract_text_on_failure:
            # Create a fallback with whatever raw content we have
            raw_text = json_text.strip()
            if len(raw_text) > 10000:  # Truncate very long responses
                raw_text = raw_text[:10000] + "... (truncated)"
                
            fallback_json = {
                "text": raw_text,
                "title": "Raw Text Fallback",
                "topics": ["raw text"],
                "summary": "Response processing failed - returning raw text.",
                "entities": [],
                "has_tables": False,
                "has_visuals": False,
                "has_numbers": False,
                "dates": [],
                "tables": [],
                "visuals": [],
                "numbers": [],
                "financial_statements": [],
                "key_metrics": [],
                "financial_terms": []
            }
            return json.dumps(fallback_json)
        else:
            return '{"text":"Error cleaning JSON response","error":"' + str(e).replace('"', '\\"') + '"}'

def run_async(func, *args, **kwargs):
    """Run an async function from Streamlit."""
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        # If func is already a coroutine, run it directly
        if asyncio.iscoroutine(func):
            return loop.run_until_complete(func)
        # Otherwise call it with the provided args
        return loop.run_until_complete(func(*args, **kwargs))
    finally:
        loop.close()

###############################
# CONTENT EXTRACTION FUNCTIONS
###############################

async def process_page_with_fallback(client, image_part, page_num, filename=None):
    """Process page with multi-level fallback including pure text extraction."""
    try:
        # First attempt: Full structured extraction
        return await extract_page_content_from_memory(client, image_part, page_num)
    except Exception as e:
        print(f"Complex extraction failed for page {page_num}, trying simple extraction: {str(e)}")
        
        try:
            # Second attempt: Simple extraction with basic prompt
            simple_prompt = f"""
            Extract just the basic information from page {page_num}.
            Return a JSON with these fields:
            {{
                "text": "full text content",
                "title": "brief title",
                "topics": ["topic1", "topic2"]
            }}
            """
            
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.0-flash",
                contents=[
                    types.Content(parts=[image_part, types.Part.from_text(text=simple_prompt)]),
                ],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            
            if response.candidates:
                # Use the enhanced clean_json_response with text extraction fallback
                json_text = clean_json_response(response.candidates[0].content.parts[0].text, 
                                              extract_text_on_failure=True)
                data = json.loads(json_text)
                
                return PageContent(
                    page_number=page_num,
                    text=data.get("text", ""),
                    title=data.get("title", f"Page {page_num} (Simple Extraction)"),
                    topics=data.get("topics", ["text extraction"]),
                    summary=data.get("summary", "Extracted using fallback method."),
                    entities=data.get("entities", []),
                    has_tables=data.get("has_tables", False),
                    has_visuals=data.get("has_visuals", False),
                    has_numbers=data.get("has_numbers", False),
                    tables=[],
                    visuals=[],
                    numbers=[],
                    dates=data.get("dates", []),
                    financial_statements=[],
                    key_metrics=[],
                    financial_terms=data.get("financial_terms", [])
                )
        except Exception as fallback_error:
            print(f"Simple extraction failed: {str(fallback_error)}, trying pure text extraction")
            
            # Third attempt: Pure text extraction (no JSON)
            try:
                text_only_prompt = f"Extract ONLY the plain text content from this page. Do NOT format as JSON. Just return the raw text."
                
                text_response = await client.aio.models.generate_content(
                    model="gemini-2.0-flash",
                    contents=[
                        types.Content(parts=[image_part, types.Part.from_text(text=text_only_prompt)]),
                    ]
                )
                
                if text_response.candidates:
                    raw_text = text_response.candidates[0].content.parts[0].text
                    
                    return PageContent(
                        page_number=page_num,
                        text=raw_text,
                        title=f"Page {page_num} (Text Only)",
                        topics=["text extraction"],
                        summary="Only plain text could be extracted from this page.",
                        entities=[],
                        has_tables=False,
                        has_visuals=False,
                        has_numbers=False,
                        tables=[],
                        visuals=[],
                        numbers=[],
                        dates=[],
                        financial_statements=[],
                        key_metrics=[],
                        financial_terms=[]
                    )
            except Exception as text_error:
                print(f"Pure text extraction also failed: {str(text_error)}")
        
        # If all extraction methods fail, return error page
        return create_error_page(page_num, "All extraction methods failed", None)

async def extract_page_content_from_memory(client, image_part, page_num):
    """Extract content from a single page using Gemini with direct memory upload in one API call."""
    # Initialize page data with defaults
    page_data = {
        "page_number": page_num,
        "text": "",
        "title": f"Page {page_num}",
        "topics": [],
        "summary": "",
        "entities": [],
        "has_tables": False,
        "has_visuals": False,
        "has_numbers": False,
        "dates": [],
        "tables": [],
        "visuals": [],
        "numbers": [],
        "financial_statements": [],
        "key_metrics": [],
        "financial_terms": []
    }

    # Combined prompt that handles all extraction in one pass
    combined_prompt = f"""
    Analyze page {page_num} and return JSON with this exact structure:
    {{
        "text": "full text content",
        "title": "brief title (2-5 words)",
        "topics": ["topic1", "topic2"],
        "summary": "key points summary (3-5 sentences)",
        "entities": ["entity1", "entity2"],
        "has_tables": true/false,
        "has_visuals": true/false,
        "has_numbers": true/false,
        "dates": ["date1", "date2"],
        "financial_terms": ["term1", "term2"],
        "tables": [
            {{
                "table_content": "markdown formatted table",
                "title": "optional table title",
                "summary": "optional table summary",
                "page_number": {page_num}
            }}
        ],
        "visuals": [
            {{
                "type": "chart/graph type",
                "description": "description of visual",
                "data_summary": "summary of data shown",
                "page_numbers": [{page_num}]
            }}
        ],
        "numbers": [
            {{
                "value": "string value",
                "description": "what the number represents",
                "context": "surrounding text context"
            }}
        ],
        "key_metrics": [
            {{
                "name": "metric name",
                "value": "string value",
                "period": "time period if available",
                "trend": "trend direction if available",
                "context": "context if available"
            }}
        ]
    }}

    Rules:
    1. All fields must be included even if empty
    2. Maintain exact field names and structure
    3. For empty arrays, use [] not null
    4. All numerical values must be strings
    5. Tables should be in markdown format with pipes (|)
    6. CRITICAL: Ensure all JSON is properly formatted with commas between all elements
    7. CRITICAL: Ensure all strings have closing quotes
    8. CRITICAL: Return ONLY valid JSON with no additional text

    """

    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",  # Using 2.0 for better multimodal understanding
            contents=[
                types.Content(parts=[image_part, types.Part.from_text(text=combined_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            # Handle case where API returns a list instead of a dict
            if isinstance(data, list):
                if data and isinstance(data[0], dict):
                    data = data[0]  # Take the first item if it's a dictionary
                else:
                    data = {}  # Empty dict as fallback
            
            # Update page data with extracted content
            page_data.update({
                "text": data.get("text", ""),
                "title": data.get("title", f"Page {page_num}"),
                "topics": data.get("topics", []),
                "summary": data.get("summary", ""),
                "entities": data.get("entities", []),
                "has_tables": data.get("has_tables", False),
                "has_visuals": data.get("has_visuals", False),
                "has_numbers": data.get("has_numbers", False),
                "dates": data.get("dates", []),
                "financial_terms": data.get("financial_terms", []),
                "tables": [
                    convert_to_model(table, TableData) if isinstance(table, dict) else None
                    for table in data.get("tables", []) if table is not None
                ],
                "visuals": [
                    convert_to_model(visual, VisualElement) if isinstance(visual, dict) else None
                    for visual in data.get("visuals", []) if visual is not None
                ],
                "numbers": [
                    convert_to_model(num, NumericalDataPoint) if isinstance(num, dict) else None
                    for num in data.get("numbers", []) if num is not None
                ],
                "key_metrics": [
                    convert_to_model(metric, FinancialMetric) if isinstance(metric, dict) else None
                    for metric in data.get("key_metrics", []) if metric is not None
                ]
            })
            
            # Filter out None values that might have been introduced
            page_data["tables"] = [t for t in page_data["tables"] if t is not None]
            page_data["visuals"] = [v for v in page_data["visuals"] if v is not None]
            page_data["numbers"] = [n for n in page_data["numbers"] if n is not None]
            page_data["key_metrics"] = [m for m in page_data["key_metrics"] if m is not None]

    except Exception as e:
        print(f"Error processing page {page_num}: {str(e)}")
        page_data["text"] = f"Error processing page: {str(e)}"
        page_data["title"] = f"Page {page_num} (Error)"

    # Convert to PageContent model with validation
    try:
        return PageContent(**page_data)
    except ValidationError as e:
        print(f"Validation error creating PageContent: {e}")
        return create_error_page(page_num, str(e), e.errors())


def parse_basic_response(response, page_num):
    """Parse basic page content response with robust error handling."""
    try:
        if not response.candidates or not response.candidates[0].content.parts:
            return {}
            
        json_text = response.candidates[0].content.parts[0].text
        
        # Add null check and string conversion
        if json_text is None:
            return {}
        json_text = str(json_text)
        
        try:
            data = json.loads(json_text)
        except json.JSONDecodeError:
            st.error(f"Invalid JSON response: {json_text[:200]}...")
            return {}
        
        # Handle both list and dict responses
        if isinstance(data, list):
            if data and isinstance(data[0], dict):
                data = data[0]
            else:
                data = {}
        
        return {
            "text": data.get("text", ""),
            "title": data.get("title", f"Page {page_num}"),
            "topics": data.get("topics", []),
            "summary": data.get("summary", ""),
            "entities": data.get("entities", []),
            "has_tables": data.get("has_tables", False),
            "has_visuals": data.get("has_visuals", False),
            "has_numbers": data.get("has_numbers", False),
            "dates": data.get("dates", []),
            "financial_terms": data.get("financial_terms", [])
        }
    except Exception as e:
        st.error(f"Basic parse error: {str(e)}")
        return {}

###############################
# DOCUMENT PROCESSING FUNCTIONS
###############################

async def process_single_page_with_semaphore(semaphore, client, page_info: dict, uploaded_files):
    """
    Process a single page using semaphore for rate limiting.
    Returns a dictionary: {"info": page_info, "result": PageContent object}.
    Handles potential errors during processing.
    """
    # Get doc_name and page_num early for use in logic and error reporting
    doc_name = page_info.get("doc_name", "Unknown_Document")
    page_num = page_info.get("page_num", 0) # Default to 0 if missing

    async with semaphore:
        result_page_obj = None # Initialize result
        try:
            # Call process_single_pdf_page, which returns a DICTIONARY
            page_dict = await process_single_pdf_page(client, page_info, uploaded_files)

            # Ensure essential keys for Pydantic initialization exist
            page_dict.setdefault("page_number", page_num)
            # Add defaults for required fields if not present in the dict returned by Gemini
            # This ensures PageContent(**page_dict) is less likely to fail basic validation
            # (though Gemini *should* return them based on the prompt)
            for field, default_value in PageContent.model_fields.items():
                 if default_value.is_required() and field not in page_dict:
                      # Try to provide a sensible default based on annotation
                      if default_value.annotation == str: page_dict.setdefault(field, "")
                      elif default_value.annotation == int: page_dict.setdefault(field, 0)
                      elif default_value.annotation == bool: page_dict.setdefault(field, False)
                      elif default_value.annotation == list or getattr(default_value.annotation, '__origin__', None) == list: page_dict.setdefault(field, [])
                      else: page_dict.setdefault(field, None) # Fallback

            # Initialize PageContent with the dictionary
            result_page_obj = PageContent(**page_dict)

        except ValidationError as ve:
            st.warning(f"Pydantic validation error creating PageContent for page {page_num} of {doc_name} AFTER initial processing: {ve}")
            # Create error page (pass doc_name in the message)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Pydantic Validation Error in '{doc_name}': {ve}", # Include doc_name in msg
                validation_errors=ve.errors()
            )

        except Exception as e:
            st.error(f"Unexpected error in semaphore task for page {page_num} of {doc_name}: {str(e)}")
            import traceback
            st.error(f"Traceback: {traceback.format_exc()}")
            # Create error page (pass doc_name in the message)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Unexpected processing error in semaphore task for '{doc_name}': {str(e)}" # Include doc_name in msg
            )

        # Return the result packaged with the original info
        return {"info": page_info, "result": result_page_obj}


async def process_single_pdf_page(client, page_info: dict, uploaded_files) -> dict:
    """
    Process a single PDF page image using Gemini with fallback options, returning extracted data as a dictionary.
    Handles errors by returning a dictionary representing an error page structure.
    The returned dictionary does NOT include 'doc_name' as a key.
    """
    page_num = page_info.get("page_num", 0)
    doc_name = page_info.get("doc_name", "Unknown_Document")  # Keep track of it locally
    page_dict = {}  # Initialize the dictionary to be returned

    try:
        # --- Ensure required info is present ---
        if "image_b64" not in page_info or not page_info["image_b64"]:
             # Create error object, convert to dict, and return
             st.warning(f"Missing image data for page {page_num} of {doc_name}.")
             error_page_obj = create_error_page(
                 page_num=page_num,
                 error_msg="Missing or empty image data (image_b64)",
                 validation_errors=None
             )
             return error_page_obj.model_dump()  # Return dict representation

        if "page_num" not in page_info:
             st.warning(f"Missing page number for a page in {doc_name}.")
             # Use 0 or handle as appropriate, create error object/dict
             error_page_obj = create_error_page(
                 page_num=0,  # Or decide on a better default/error indicator
                 error_msg="Missing page number information",
                 validation_errors=None
             )
             return error_page_obj.model_dump()

        # --- Prepare for Gemini Call ---
        img_data = base64.b64decode(page_info["image_b64"])
        image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")  # Assuming JPEG, adjust if needed

        # --- Call Gemini with fallback strategy instead of direct extraction ---
        page_content_obj = await process_page_with_fallback(
            client, image_part, page_num, doc_name  # Pass page_num and doc_name
        )

        # --- Convert successful PageContent object to a dictionary ---
        # This dictionary contains only the fields defined in the PageContent model
        page_dict = page_content_obj.model_dump()

        return page_dict  # Return the clean dictionary

    except Exception as e:
        st.error(f"Error during Gemini processing or data handling for page {page_num} of {doc_name}: {str(e)}")
        # Create an error PageContent object first
        error_page_obj = create_error_page(
            page_num=page_num,
            error_msg=f"Gemini API or processing error: {str(e)}",
            validation_errors=None
        )
        # Convert the error PageContent object to a dictionary to maintain return type consistency
        page_dict = error_page_obj.model_dump()

        return page_dict

    
def preprocess_page_data(page_data):
    """Preprocess page data with special handling for numerical data"""
    if not isinstance(page_data, dict):
        if hasattr(page_data, 'model_dump'):
            page_data = page_data.model_dump()
        else:
            raise ValueError("Page data must be a dictionary or Pydantic model")
    
    # Handle numbers with null contexts
    if "numbers" in page_data:
        processed_numbers = []
        for num in page_data["numbers"]:
            if isinstance(num, dict):
                processed_numbers.append({
                    "value": str(num.get("value", "")),
                    "description": str(num.get("description", "")),
                    "context": str(num.get("context", ""))
                })
            elif hasattr(num, 'model_dump'):
                num_dict = num.model_dump()
                processed_numbers.append({
                    "value": str(num_dict.get("value", "")),
                    "description": str(num_dict.get("description", "")),
                    "context": str(num_dict.get("context", ""))
                })
            else:
                processed_numbers.append({
                    "value": str(getattr(num, 'value', '')),
                    "description": str(getattr(num, 'description', '')),
                    "context": str(getattr(num, 'context', ''))
                })
        page_data["numbers"] = processed_numbers
    
    # Similar handling for other numerical fields
    for list_field in ["key_metrics", "financial_statements"]:
        if page_data.get(list_field) is None:
            page_data[list_field] = []
    
    return page_data

async def finalize_document(client, doc_name, pages):
    """Create document summary and final structure with proper model handling"""
    try:
        # Convert page data to dictionaries first
        validated_pages = []
        for page in pages:
            try:
                # Convert to dict if model
                page_dict = page.model_dump() if hasattr(page, 'model_dump') else page
                
                # Ensure required fields
                page_dict.setdefault("title", f"Page {page_dict.get('page_number', 0)}")
                page_dict.setdefault("text", "")
                
                validated_pages.append(PageContent(**page_dict))
            except Exception as e:
                st.warning(f"Error validating page: {e}")
                validated_pages.append(PageContent(
                    page_number=page.get("page_number", len(validated_pages)+1),
                    text=f"Validation error: {str(e)}",
                    title=f"Page {len(validated_pages)+1} (Invalid)"
                ))

        # Generate document summary as dict
        summary_dict = await generate_financial_summary(validated_pages, doc_name, client)
        
        # Create the final result structure
        result = {
            "raw_extracted_content": {
                "filename": doc_name,
                "pages": validated_pages,
                "summary": summary_dict
            }
        }
        
        return result
        
    except Exception as e:
        print(f"Error finalizing document {doc_name}: {str(e)}")
        # Return basic document structure on error
        return {
            "raw_extracted_content": {
                "filename": doc_name,
                "pages": pages,
                "summary": None
            }
        }

# Remove doc_name parameter and assignment
def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List] = None) -> PageContent:
    """
    Creates a PageContent object representing an error state.
    The document name association must be handled by the caller.
    """
    # Include context in the error message if available (passed in error_msg)
    error_text = f"Error processing page {page_num}: {error_msg}"

    if validation_errors:
        try:
            error_details = "\nValidation Errors:\n" + "\n".join(
                f"- Field '{'.'.join(map(str, err.get('loc', ['unknown'])))}': {err.get('msg', 'No message')}"
                for err in validation_errors
            )
            error_text += error_details
        except Exception as format_error:
            error_text += f"\n(Could not format validation errors: {format_error})"

    # Initialize the PageContent object with standard fields
    # Title can reflect the error state
    error_page = PageContent(
        page_number=page_num,
        text=error_text,
        title=f"Page {page_num} (Error)", # Keep it simple, context is in error_text
        topics=["error"],
        summary=f"Failed to process page {page_num}.",
        entities=[],
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        dates=[],
        financial_statements=[],
        key_metrics=[],
        financial_terms=[]
    )

    # *** REMOVED: error_page.doc_name = doc_name ***

    return error_page

async def process_text_page(client, text_content: str, page_num: int, filename: str):
    """Process full text content of a page without size restrictions"""
    try:
        prompt = f"""
        Analyze page {page_num} from document {filename} and extract:
        1. Complete text content
        2. Key topics
        3. Numerical data
        4. Financial terms
        5. Tables (if any)
        6. Important entities
        
        Return JSON with this exact structure:
        {{
            "text": "full content",
            "title": "page title",
            "topics": ["topic1", "topic2"],
            "entities": ["entity1", "entity2"],
            "financial_terms": ["term1", "term2"],
            "tables": [
                {{
                    "table_content": "| Header1 | Header2 |\\n|---|---| etc...",
                    "title": "Optional table title",
                    "summary": "Optional table summary",
                    "page_number": {page_num}
                }}
            ],
            "summary": "summary of page content",
            "numerical_data": [{{
                "value": "123",
                "description": "description"
            }}]
        }}
        """
        
        # Combine content and prompt
        full_content = f"DOCUMENT CONTENT:\n{text_content}\n\nPROMPT:\n{prompt}"
        
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",  # Using 1.5 for larger context
            contents=[
                types.Content(
                    role="user",
                    parts=[types.Part.from_text(text=full_content)]
                )
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            return PageContent(
                page_number=page_num,
                text=data.get("text", text_content),
                title=data.get("title", f"Page {page_num}"),  # Fallback title
                topics=data.get("topics", []),
                entities=data.get("entities", []),
                financial_terms=data.get("financial_terms", []),
                tables=[TableData(
                    table_content="\n".join([
                        f"| {' | '.join(str(cell) if cell is not None else '' for cell in table.get('headers', []))} |",
                        f"| {' | '.join(['---' for _ in table.get('headers', [])])} |",
                        *[f"| {' | '.join(str(cell) if cell is not None else '' for cell in row)} |" 
                          for row in table.get('rows', [])]
                    ]),
                    title=table.get("title"),
                    summary=table.get("summary"),
                    page_number=page_num
                ) for table in data.get("tables", [])],
                numbers=[NumericalDataPoint(
                    value=num["value"],
                    description=num.get("description", "")
                ) for num in data.get("numerical_data", [])],
                has_tables=bool(data.get("tables")),
                has_numbers=bool(data.get("numerical_data"))
            )
        
        return None
    except Exception as e:
        return PageContent(
            page_number=page_num,
            text=f"Error processing page: {str(e)}",
            title=f"Page {page_num} (Error)"
        )

async def process_word_document(client, file_content: bytes, filename: str):
    """Process Word document from bytes in memory"""
    try:
        # Create in-memory document
        doc = docx.Document(io.BytesIO(file_content))
        
        # Extract all paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Detect page breaks
        pages = []
        current_page = []
        
        for para in paragraphs:
            if "PAGE_BREAK" in para or "### PAGE" in para:
                if current_page:
                    pages.append("\n".join(current_page))
                    current_page = []
                continue
            current_page.append(para)
        
        if current_page:
            pages.append("\n".join(current_page))
        
        # If no page breaks, treat as single page
        if not pages:
            pages = ["\n".join(paragraphs)]
        
        # Process each page
        processed_pages = []
        for i, page_text in enumerate(pages, 1):
            page = await process_text_page(client, page_text, i, filename)
            if page:
                processed_pages.append(page)
                
        return processed_pages
        
    except Exception as e:
        st.error(f"Error processing Word document {filename}: {str(e)}")
        return []

async def process_pptx_document(client, file_content: bytes, filename: str):
    """Process PowerPoint files from bytes in memory"""
    try:
        from pptx import Presentation
        import io
        
        # Create in-memory presentation
        ppt = Presentation(io.BytesIO(file_content))
        
        # Extract text from slides
        pages = []
        for i, slide in enumerate(ppt.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                pages.append({
                    "page_number": i,
                    "text": "\n".join(slide_text),
                    "title": f"Slide {i}"
                })
        
        # Process each slide as a page
        processed_pages = []
        for page in pages:
            processed_page = await process_text_page(
                client,
                page["text"],
                page["page_number"],
                filename
            )
            if processed_page:
                processed_pages.append(processed_page)
                
        return processed_pages
        
    except Exception as e:
        st.error(f"Error processing PowerPoint file {filename}: {str(e)}")
        return []

async def process_text_document(client, file_data: dict):
    """Process plain text document with page breaks"""
    file_buffer = io.BytesIO(file_data["content"])
    filename = file_data["name"]
    
    try:
        text_content = file_buffer.read().decode('utf-8')
        
        # Split by page breaks (customize based on your document format)
        page_breaks = re.compile(r"\n{3,}|=== Page \d+ ===|\f")
        pages = [p.strip() for p in page_breaks.split(text_content) if p.strip()]
        
        # If no page breaks, treat as single page
        if not pages:
            pages = [text_content]
        
        processed_pages = []
        for i, page_text in enumerate(pages, 1):
            page = await process_text_page(client, page_text, i, filename)
            if page:
                processed_pages.append(page)
                
        return processed_pages
        
    finally:
        file_buffer.close()

async def extract_financial_metrics(data: List[Dict]) -> List[FinancialMetric]:
    """Identify and extract financial metrics from data."""
    metrics = []
    financial_terms = [
        "revenue", "profit", "ebitda", "net income", "assets", 
        "liabilities", "equity", "eps", "roi", "roa", "roe"
    ]
    
    for row in data:
        for key, value in row.items():
            if isinstance(value, (int, float)) or (isinstance(value, str) and value.replace('.','',1).isdigit()):
                # Check if key contains financial terms
                if isinstance(key, str) and any(term in key.lower() for term in financial_terms):
                    metrics.append(FinancialMetric(
                        name=key,
                        value=str(value),
                        context=f"Found in {key}"
                    ))
    
    return metrics

async def process_financial_data(content: Dict, filename: str) -> List[PageContent]:
    """Process financial data files (Excel, CSV)."""
    pages = []
    
    if isinstance(content, dict) and "data" in content:
        # Process as tabular data
        page = PageContent(
            page_number=1,
            text=f"Financial data from {filename}",
            title=f"Financial Data: {filename}",
            topics=["financial data"],
            has_tables=True,
            has_numbers=True
        )
        
        # Convert data to table format
        table = TableData(
            table_content="\n".join([
                f"| {' | '.join(str(cell) if cell is not None else '' for cell in table.get('headers', []))} |",
                f"| {' | '.join(['---' for _ in table.get('headers', [])])} |",
                *[f"| {' | '.join(str(cell) if cell is not None else '' for cell in row)} |" 
                  for row in table.get('rows', [])]
            ]),
            title=f"Data from {filename}",
            summary=table.get("summary"),
            page_number=1
        )
        page.tables.append(table)
        
        # Extract financial metrics
        metrics = await extract_financial_metrics(content["data"])
        page.key_metrics = metrics
        
        pages.append(page)
    
    return pages

async def generate_financial_summary(pages: List[Union[Dict, PageContent]], filename: str, client) -> Dict:
    """Generate enhanced financial summary using Gemini and return as dict."""
    # Prepare summary input
    summary_input = {
        "page_titles": [],
        "key_metrics": [],
        "financial_terms": [],
        "has_tables": False,
        "has_visuals": False
    }

    # Track if we have any tables or visuals
    has_tables = False
    has_visuals = False

    for page in pages:
        # Handle both dict and PageContent instances
        if isinstance(page, PageContent):
            page_dict = page.model_dump()
        else:
            page_dict = page

        # Add page title
        summary_input["page_titles"].append(page_dict.get("title", f"Page {page_dict.get('page_number', '?')}"))

        # Collect metrics - handle both dict and model instances
        for metric in page_dict.get("key_metrics", []):
            if hasattr(metric, 'model_dump'):  # If it's a Pydantic model
                metric_dict = metric.model_dump()
            else:
                metric_dict = metric
            
            summary_input["key_metrics"].append({
                **metric_dict,
                "page": page_dict.get("page_number", 0)
            })

        # Collect other data
        summary_input["financial_terms"].extend(page_dict.get("financial_terms", []))
        
        # Track if we have tables or visuals
        if page_dict.get("has_tables", False):
            has_tables = True
        if page_dict.get("has_visuals", False):
            has_visuals = True

    summary_input["has_tables"] = has_tables
    summary_input["has_visuals"] = has_visuals

    # Generate summary with Gemini
    try:
        summary_response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_text(text=f"""
                            Create a financial summary for "{filename}" with this structure:
                            {{
                                "title": "concise title",
                                "themes": ["theme1", "theme2"],
                                "questions": ["question1", "question2"],
                                "summary": "detailed summary",
                                "tables_summary": "summary of key tables",
                                "visuals_summary": "summary of key visuals"
                            }}
                            Input Data:
                            {json.dumps(summary_input, indent=2)}
                        """)
                    ]
                )
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json"
            )
        )
        
        if summary_response.candidates:
            json_text = summary_response.candidates[0].content.parts[0].text
            json_text = clean_json_response(json_text)
            summary_data = json.loads(json_text)
            
            # Handle case where API returns a list
            if isinstance(summary_data, list):
                summary_data = summary_data[0] if summary_data else {}
            
            return summary_data
            
    except Exception as e:
        print(f"Summary generation error: {str(e)}")
    
    # Fallback summary as dict
    return {
        "title": f"Financial Summary: {filename}",
        "themes": ["financial", "analysis"],
        "questions": ["What are the key financial trends?"],
        "summary": "Generated summary unavailable",
        "tables_summary": "",
        "visuals_summary": ""
    }

##############
# Document Processing
##############

async def process_all_documents_async(file_data):
    """Process all documents with proper model conversion and error handling."""
    # ... (Keep initial setup: API key, client, status elements, steps, update_step_status)
    api_key = get_gemini_api_key()
    # Ensure client configuration is correct for your setup
    client = genai.Client(api_key=api_key,
                          http_options={
                              "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'), # Use proxy or direct URL
                              "headers": {
                                  "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                                  # Only include target URL if using Helicone proxy
                                  # "helicone-target-url": 'https://generativelanguage.googleapis.com'
                              }
                          })
    processed_docs = []
    uploaded_files = [] # Keep if used elsewhere

    status_container = st.session_state.get("status_container")
    progress_bar = st.session_state.get("progress_bar")
    time_info = st.session_state.get("time_info")

    processing_status = {
        "active": True, "current_step": "", "total_steps": 7, "current_step_num": 0,
        "step_progress": 0.0, "start_time": time.time(), "step_start_time": time.time(),
        "estimated_time_remaining": None, "parallel_count": 1
    }
    steps = {
        "Detecting document types": 1, "Processing PDF documents": 2, "Analyzing PDF pages": 3,
        "Processing non-PDF documents": 4, "Generating document summaries": 5,
        "Finalizing results": 6, "Cleaning up resources": 7
    }

    def update_step_status(step, current=None, total=None, parallel=1):
        #...(same implementation as before)...
        processing_status["current_step"] = step
        processing_status["current_step_num"] = steps.get(step, processing_status["current_step_num"])

        step_progress = 0.0
        if current is not None and total is not None and total > 0:
             step_progress = min(1.0, float(current) / float(total))
        processing_status["step_progress"] = step_progress

        overall_progress = (processing_status["current_step_num"] - 1 + processing_status["step_progress"]) / processing_status["total_steps"]
        overall_progress = max(0.0, min(1.0, overall_progress))

        status_msg = f"{processing_status['current_step_num']}/{processing_status['total_steps']}: {step}"
        if current is not None and total is not None:
            status_msg += f" ({current}/{total})"
        if parallel > 1:
            status_msg += f" | Parallelism: {parallel}"

        eta_msg = ""
        if overall_progress > 0.01:
            time_elapsed = time.time() - processing_status["start_time"]
            est_total_time = time_elapsed / overall_progress if overall_progress > 0 else 0 # Avoid division by zero
            remaining = max(0, est_total_time - time_elapsed)
            mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
            mins_remaining, secs_remaining = divmod(int(remaining), 60)
            eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: ~{mins_remaining}m {secs_remaining}s"

        if status_container:
             # Use try-except for robustness if elements disappear
             try:
                 status_container.update(label=status_msg, expanded=True)
             except Exception: pass # Ignore if container not found
        if progress_bar:
             try:
                 progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
             except Exception: pass
        if time_info:
             try:
                 time_info.markdown(f"`{eta_msg}`")
             except Exception: pass


    final_docs = []

    try:
        # Step 1: Detect document types
        update_step_status("Detecting document types", 0, len(file_data))
        pdf_docs = []
        other_docs = []
        for i, file in enumerate(file_data):
            update_step_status("Detecting document types", i + 1, len(file_data))
            # Robust file type detection
            file_name = file.get("name", "")
            file_ext = file.get("type", "").lower() or (file_name.split('.')[-1].lower() if '.' in file_name else "")
            file["type"] = file_ext # Ensure type is set in the dict
            if file_ext == 'pdf':
                pdf_docs.append(file)
            else:
                other_docs.append(file)
        update_step_status("Detecting document types", len(file_data), len(file_data))

        # --- PDF PROCESSING SECTION ---
        if pdf_docs:
            update_step_status("Processing PDF documents", 0, len(pdf_docs))

            MAX_PARALLEL_CONVERSIONS = 5 # Limit concurrent PDF rendering
            all_pdf_pages_info = [] # List to hold page_info dicts from extraction

            # --- Convert PDF pages to images ---
            pdf_conversion_progress = 0
            conversion_tasks = []
            for file in pdf_docs:
                if "content" in file and "name" in file:
                    # Wrap the async call to extract pages
                    conversion_tasks.append(extract_pages_from_pdf_bytes(file["content"], file["name"]))
                else:
                    st.warning(f"Skipping PDF file due to missing 'content' or 'name': {file.get('name', 'Unknown')}")

            # Use a semaphore to limit concurrent PDF rendering if extract_pages_from_pdf_bytes is IO heavy
            pdf_render_semaphore = asyncio.Semaphore(MAX_PARALLEL_CONVERSIONS)
            async def run_with_semaphore(task):
                 async with pdf_render_semaphore:
                     return await task

            tasks_with_semaphore = [run_with_semaphore(task) for task in conversion_tasks]

            # Gather results with progress update
            conversion_results = await asyncio.gather(*tasks_with_semaphore, return_exceptions=True)

            for i, result in enumerate(conversion_results):
                 pdf_conversion_progress += 1 # Increment progress regardless of result
                 if isinstance(result, list):
                     all_pdf_pages_info.extend(result)
                 elif isinstance(result, Exception):
                     st.error(f"PDF page extraction failed for '{pdf_docs[i].get('name', 'Unknown')}': {result}")
                 else:
                     st.warning(f"Unexpected result type during PDF extraction: {type(result)}")

                 update_step_status(
                    "Processing PDF documents",
                    pdf_conversion_progress,
                    len(pdf_docs) # Total PDFs to convert
                    # Parallel count isn't easily tracked here, maybe remove or estimate
                )
            # Final update for this step
            update_step_status("Processing PDF documents", len(pdf_docs), len(pdf_docs))


            # --- Analyze PDF Pages ---
            pdf_page_count = len(all_pdf_pages_info)
            if pdf_page_count > 0:
                update_step_status("Analyzing PDF pages", 0, pdf_page_count)

                MAX_PARALLEL_PAGES = 25 # Gemini API concurrency limit
                semaphore = asyncio.Semaphore(MAX_PARALLEL_PAGES)

                # --- Create tasks ---
                page_tasks = [] # List to hold the Future objects
                valid_page_info_count = 0

                for page_info in all_pdf_pages_info:
                    # Basic validation of the page_info structure
                    if isinstance(page_info, dict) and "doc_name" in page_info and "page_num" in page_info and "image_b64" in page_info:
                        coro = process_single_page_with_semaphore(
                            semaphore, client, page_info, uploaded_files
                        )
                        page_tasks.append(asyncio.create_task(coro)) # Create and store task
                        valid_page_info_count += 1
                    else:
                         st.warning(f"Skipping invalid or incomplete page_info data: {str(page_info)[:100]}...") # Log truncated data

                # Adjust total page count if some were invalid
                pdf_page_count = valid_page_info_count
                if pdf_page_count == 0:
                    st.warning("No valid PDF page data found to analyze after extraction.")
                else:
                    # --- Collect results using as_completed ---
                    processed_results_with_context = [] # Store dicts: {"doc_name": ..., "page_content": PageContent}
                    processed_count = 0

                    for future in asyncio.as_completed(page_tasks):
                        try:
                            # Get the result dictionary returned by the wrapper
                            task_output = await future # This is {"info": ..., "result": ...}

                            # Extract the original info and the actual page content/error
                            original_info = task_output.get("info", {}) # Safely get info
                            result_page_obj = task_output.get("result") # Safely get result (PageContent or Error PageContent)

                            doc_name = original_info.get("doc_name", "Unknown_Context") # Use fallback
                            page_num = original_info.get("page_num", -1) # Use fallback

                            if isinstance(result_page_obj, PageContent):
                                processed_results_with_context.append({
                                    "doc_name": doc_name,
                                    "page_content": result_page_obj # Already has page_num inside
                                })
                            elif result_page_obj is None:
                                 st.warning(f"Page processing task for page {page_num} of {doc_name} returned None result payload.")
                                 # Optionally create an error page here too if needed
                                 error_page = create_error_page(page_num, f"Task returned None payload in '{doc_name}'")
                                 processed_results_with_context.append({"doc_name": doc_name, "page_content": error_page})
                            else:
                                # This case should ideally not happen if the wrapper always returns PageContent or None
                                st.warning(f"Unexpected item in result payload for {doc_name} page {page_num}: {type(result_page_obj)}. Storing as error.")
                                error_page = create_error_page(page_num, f"Unexpected result type {type(result_page_obj)} in '{doc_name}'")
                                processed_results_with_context.append({"doc_name": doc_name, "page_content": error_page})


                        except Exception as e:
                            # Handle exceptions raised by 'await future' itself (e.g., task cancellation)
                            # We don't have reliable context (doc_name/page_num) here if the task itself failed critically
                            st.error(f"Error awaiting page processing task result: {e}")
                            # Log a generic error page if context is lost
                            processed_count += 1 # Increment even on failure to keep progress moving
                            # We can't easily associate this error with a specific doc/page here
                            # Maybe log it separately or create a generic error entry later?

                        # Increment progress AFTER processing the result or handling the exception
                        processed_count += 1
                        update_step_status(
                            "Analyzing PDF pages",
                            processed_count,
                            pdf_page_count,
                            min(MAX_PARALLEL_PAGES, pdf_page_count - processed_count + 1) # Active tasks
                        )

                    # Final update for this step
                    update_step_status("Analyzing PDF pages", processed_count, pdf_page_count)

                    # --- Organize PDF pages using collected context ---
                    doc_pages = defaultdict(list)
                    for item in processed_results_with_context:
                        # item is {"doc_name": ..., "page_content": PageContent_object}
                        doc_name = item["doc_name"]
                        page_obj = item["page_content"] # This is the PageContent object
                        if doc_name != "Unknown_Context": # Avoid grouping unknown pages
                            doc_pages[doc_name].append(page_obj)
                        else:
                            st.warning(f"Could not determine document name for processed page: {getattr(page_obj, 'page_number', '?')}")


                    # --- Generate Summaries for PDF documents ---
                    doc_count = len(doc_pages)
                    if doc_count > 0:
                        update_step_status("Generating document summaries", 0, doc_count)
                        summary_tasks = []
                        doc_name_order = [] # Keep track of the order for results

                        for doc_name, pages_list in doc_pages.items():
                            # Ensure pages are PageContent objects and sort them
                            valid_pages = sorted(
                                [p for p in pages_list if isinstance(p, PageContent)],
                                key=lambda p: getattr(p, 'page_number', 0) # Safe access to page_number
                            )
                            if not valid_pages:
                                st.warning(f"No valid pages found for document '{doc_name}' during summary phase.")
                                continue # Skip summary if no valid pages

                            # Pass the list of PageContent objects to finalize_document
                            summary_tasks.append(finalize_document(client, doc_name, valid_pages))
                            doc_name_order.append(doc_name)

                        if summary_tasks: # Only run gather if there are tasks
                            summary_results = await asyncio.gather(*summary_tasks, return_exceptions=True)

                            for i, result in enumerate(summary_results):
                                update_step_status("Generating document summaries", i + 1, len(doc_name_order)) # Use len(doc_name_order) as total
                                doc_name = doc_name_order[i] # Get doc_name based on original order
                                if isinstance(result, Exception):
                                    st.error(f"Failed to finalize summary for {doc_name}: {result}")
                                    # Convert original pages (from doc_pages) to dicts for error output
                                    pages_as_dicts = []
                                    for p in doc_pages.get(doc_name, []): # Use .get for safety
                                        try:
                                            pages_as_dicts.append(p.model_dump())
                                        except Exception as dump_err:
                                            pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})

                                    processed_docs.append({
                                        "raw_extracted_content": {
                                            "filename": doc_name,
                                            "pages": pages_as_dicts,
                                            "summary": None,
                                            "error": f"Summary generation failed: {result}"
                                        }
                                    })
                                elif isinstance(result, dict) and "raw_extracted_content" in result:
                                    # Ensure pages within the final result are dicts
                                    raw_content = result["raw_extracted_content"]
                                    if "pages" in raw_content and isinstance(raw_content["pages"], list):
                                        pages_as_dicts = []
                                        for p in raw_content["pages"]:
                                             try:
                                                 # Convert PageContent to dict if it isn't already
                                                 pages_as_dicts.append(p.model_dump() if hasattr(p, 'model_dump') else p)
                                             except Exception as dump_err:
                                                 pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})
                                        raw_content["pages"] = pages_as_dicts # Replace list with dicts
                                    processed_docs.append(result) # Add the fully processed doc structure
                                else:
                                    st.warning(f"Unexpected result type during summary finalization for {doc_name}: {type(result)}")
                                    # Convert original pages to dicts for error output
                                    pages_as_dicts = []
                                    for p in doc_pages.get(doc_name, []):
                                        try:
                                             pages_as_dicts.append(p.model_dump())
                                        except Exception as dump_err:
                                             pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})
                                    processed_docs.append({
                                        "raw_extracted_content": {
                                            "filename": doc_name,
                                            "pages": pages_as_dicts,
                                            "summary": None,
                                            "error": "Unexpected finalization result type"
                                        }
                                   })
                        else:
                            st.info("No valid documents found to summarize.")

                        # Final status update for summary generation
                        update_step_status("Generating document summaries", len(doc_name_order), len(doc_name_order))
                    else:
                         st.info("No documents found after page processing to summarize.")

            else: # pdf_page_count == 0 or no valid page info
                 st.info("No PDF pages were successfully extracted or processed, skipping PDF analysis and summary.")


        # --- Non-PDF Document Processing ---
        if other_docs:
            update_step_status("Processing non-PDF documents", 0, len(other_docs))

            MAX_PARALLEL_NON_PDF = 4
            semaphore_non_pdf = asyncio.Semaphore(MAX_PARALLEL_NON_PDF)
            non_pdf_tasks = []
            valid_non_pdf_files = [] # Keep track of files we are actually processing

            for file in other_docs:
                 # Ensure file type and content are present
                 if "type" not in file or not file["type"]:
                      file_name = file.get("name", "Unknown")
                      if '.' in file_name:
                           file["type"] = file_name.split('.')[-1].lower()
                      else:
                           st.warning(f"Cannot determine type for non-PDF file: {file_name}. Skipping.")
                           continue
                 if "content" not in file or not file["content"]:
                     st.warning(f"Missing content for non-PDF file: {file.get('name', 'Unknown')}. Skipping.")
                     continue

                 valid_non_pdf_files.append(file) # Add to list of files to process

                 # Wrapper to apply semaphore to the processing function
                 async def process_non_pdf_wrapper(f):
                      async with semaphore_non_pdf:
                           try:
                               # process_single_document_memory should return the nested dict { "raw_extracted_content": ... }
                               result_dict = await process_single_document_memory(client, f, uploaded_files)

                               # --- Ensure pages are dicts in the final result ---
                               if result_dict and "raw_extracted_content" in result_dict:
                                   raw_content = result_dict["raw_extracted_content"]
                                   if "pages" in raw_content and isinstance(raw_content["pages"], list):
                                       pages_as_dicts = []
                                       for p in raw_content["pages"]:
                                           try:
                                               pages_as_dicts.append(p.model_dump() if hasattr(p, 'model_dump') else p)
                                           except Exception as dump_err:
                                                pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})
                                       raw_content["pages"] = pages_as_dicts
                               return result_dict # Return the modified dict
                           except Exception as e:
                                st.error(f"Error processing non-PDF {f.get('name', 'Unknown')}: {e}")
                                # Return an error structure consistent with successful returns
                                return {
                                    "raw_extracted_content": {
                                        "filename": f.get('name', 'Unknown_NonPDF_Error'),
                                        "pages": [],
                                        "summary": None,
                                        "error": f"Processing failed: {e}"
                                    }
                                }
                 non_pdf_tasks.append(process_non_pdf_wrapper(file)) # Pass the valid file

            non_pdf_results = []
            non_pdf_progress = 0
            total_non_pdf_to_process = len(non_pdf_tasks) # Based on actual tasks created

            if total_non_pdf_to_process > 0:
                # Use asyncio.gather for potentially better performance if tasks are independent
                gathered_results = await asyncio.gather(*non_pdf_tasks, return_exceptions=True)

                for result in gathered_results:
                    non_pdf_progress += 1
                    if isinstance(result, Exception):
                        # This catches exceptions from the gather itself or unhandled ones in the wrapper
                        st.error(f"Non-PDF processing task failed critically: {result}")
                        # Optionally add an error placeholder if needed, but the wrapper should handle most
                    elif result and isinstance(result, dict) and "raw_extracted_content" in result:
                        non_pdf_results.append(result) # Add successful or handled error results
                    else:
                        st.warning(f"Received invalid/unexpected result structure from non-PDF processing: {result}")
                    update_step_status("Processing non-PDF documents", non_pdf_progress, total_non_pdf_to_process, min(MAX_PARALLEL_NON_PDF, total_non_pdf_to_process - non_pdf_progress + 1))

                processed_docs.extend(non_pdf_results)
            else:
                st.info("No valid non-PDF documents found to process.")

            # Update status one last time after loop/gather finishes
            update_step_status("Processing non-PDF documents", total_non_pdf_to_process, total_non_pdf_to_process)


        # --- Final Processing Steps ---
        update_step_status("Finalizing results", 1, 1)
        # Ensure all page lists within processed_docs contain only dictionaries
        # This step might be redundant if conversion happens correctly earlier, but acts as a safeguard
        final_cleaned_docs = []
        for doc_result in processed_docs:
             if doc_result and "raw_extracted_content" in doc_result: # Check if doc_result is valid
                 raw_content = doc_result["raw_extracted_content"]
                 if "pages" in raw_content and isinstance(raw_content["pages"], list):
                     pages_as_dicts = []
                     for p in raw_content["pages"]:
                          if isinstance(p, dict):
                              pages_as_dicts.append(p)
                          elif hasattr(p, 'model_dump'):
                              try:
                                  pages_as_dicts.append(p.model_dump())
                              except Exception as final_dump_err:
                                   pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Final dump failed: {final_dump_err}"})
                          else: # Fallback for unexpected types
                               pages_as_dicts.append({"page_number": -1, "error": "Unknown page format", "data": str(p)[:100]})
                     raw_content["pages"] = pages_as_dicts
                 final_cleaned_docs.append(doc_result) # Add the cleaned doc result
             else:
                 st.warning(f"Skipping invalid document result during finalization: {doc_result}")


        final_docs = final_cleaned_docs # Assign the cleaned list
        update_step_status("Cleaning up resources", 1, 1)
        processing_status["active"] = False
        if status_container:
            try:
                 status_container.update(label="✅ Document processing complete!", state="complete", expanded=False)
            except Exception: pass


    except asyncio.CancelledError:
         st.warning("Document processing was cancelled.")
         if status_container:
             try:
                 status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
             except Exception: pass
         processing_status["active"] = False
         return [] # Return empty list on cancellation
    except Exception as e:
        st.error(f"❌ An error occurred during document processing: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        if status_container:
            try:
                 status_container.update(label=f"❌ Error: {str(e)}", state="error", expanded=True)
            except Exception: pass
        processing_status["active"] = False
        # Return potentially partially processed docs
        return final_docs # Return what was processed so far

    finally:
        processing_status["active"] = False
        st.session_state.processing_active = False
        # Clean up status elements from session state if needed
        # It's often better to leave them until the next run starts or clear explicitly
        # if "status_container" in st.session_state: del st.session_state["status_container"]
        # if "progress_bar" in st.session_state: del st.session_state["progress_bar"]
        # if "time_info" in st.session_state: del st.session_state["time_info"]

    # Ensure return type consistency (list of dictionaries)
    return final_docs

        
async def extract_pages_from_pdf_bytes(pdf_bytes, file_name):
    """Extract pages from PDF bytes using PyMuPDF."""
    loop = asyncio.get_running_loop()
    pages = []
    try:
        pdf_stream = io.BytesIO(pdf_bytes)
        # Run synchronous PyMuPDF code in an executor thread
        def sync_extract():
            doc = pymupdf.open(stream=pdf_stream, filetype="pdf")
            extracted = []
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    # Moderate DPI for balance
                    pix = page.get_pixmap(dpi=150, alpha=False) # No alpha needed for JPEG
                    img_bytes = pix.tobytes("jpeg", jpg_quality=85) # Control quality
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    extracted.append({
                        "page_num": page_num + 1,
                        "image_b64": img_b64,
                        "doc_name": file_name,
                        # "doc_bytes": None # Avoid storing large bytes here
                    })
                except Exception as page_e:
                     st.warning(f"Error extracting page {page_num+1} from {file_name}: {page_e}")
            doc.close()
            return extracted

        pages = await loop.run_in_executor(None, sync_extract)
        return pages

    except Exception as e:
        st.error(f"Error opening or processing PDF {file_name}: {str(e)}")
        return [] # Return empty list on error

async def run_in_executor(func, *args):
    """Run synchronous IO-bound functions in executor"""
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, func, *args)

async def process_single_document_memory(client, file, uploaded_files):
    """Process a single non-PDF document from memory"""
    try:
        file_name = file["name"]
        file_type = file["type"]
        file_content = file["content"]
        
        # Process based on file type
        if file_type in ["xlsx", "xls"]:
            content = await run_in_executor(process_excel_memory, file_content)
            processed_pages = await process_financial_data(content, file_name)
        elif file_type == "csv":
            content = await run_in_executor(process_csv_memory, file_content)
            processed_pages = await process_financial_data(content, file_name)
        elif file_type == "docx":
            processed_pages = await process_word_document(client, file_content, file_name)
        elif file_type == "pptx":
            processed_pages = await process_pptx_document(client, file_content, file_name)
        elif file_type == "txt":
            processed_pages = await process_text_document(client, {"content": file_content, "name": file_name})
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        if not processed_pages:
            raise Exception(f"No pages were processed in {file_name}")

        # Generate document summary
        document_summary = await generate_financial_summary(processed_pages, file_name, client)
        
        return {
            "raw_extracted_content": {
                "filename": file_name,
                "pages": processed_pages,
                "summary": document_summary
            }
        }
    except Exception as e:
        st.error(f"Error processing document {file_name}: {str(e)}")
        return None
    
def process_excel_memory(file_content: bytes):
    """Process Excel files directly from bytes"""
    content = {}
    try:
        # Use openpyxl to load from bytes
        wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            data = []
            
            # Get headers from first row
            headers = []
            for row in sheet.iter_rows(min_row=1, max_row=1):
                headers = [cell.value for cell in row if cell.value]
                break
            
            # Process data rows
            for row in sheet.iter_rows(min_row=2):
                row_data = {}
                for header, cell in zip(headers, row):
                    if header and cell.value:
                        row_data[header] = cell.value
                if row_data:
                    data.append(row_data)
            
            content[sheet_name] = {
                "headers": headers,
                "data": data,
                "type": "excel_sheet"
            }
            
        return content
        
    except Exception as e:
        st.error(f"Error processing Excel file: {str(e)}")
        return {}
    finally:
        if 'wb' in locals():
            wb.close()

def process_csv_memory(file_buffer):
    """Process CSV files from memory buffer."""
    content = {"data": []}
    
    # Read as text
    text_content = file_buffer.read().decode('utf-8-sig')
    
    # Parse with CSV reader
    reader = csv.DictReader(io.StringIO(text_content))
    content["headers"] = reader.fieldnames
    content["data"] = [row for row in reader]
    
    return content

###############################
# UI DISPLAY FUNCTIONS
###############################

def display_table(table):
    """Display table that could be either TableData model or dict."""
    if hasattr(table, 'model_dump'):
        table_data = table.model_dump()
    else:
        table_data = table
    
    # Display table title if available
    if table_data.get("title"):
        st.subheader(table_data["title"])
    
    # Display the table content
    if table_data.get("table_content"):
        try:
            df = pd.read_csv(
                io.StringIO(table_data["table_content"]),
                sep="|",
                engine="python",
                skipinitialspace=True
            )
            df = df.dropna(axis=1, how="all")
            st.dataframe(df)
        except Exception:
            st.markdown(f"```markdown\n{table_data['table_content']}\n```")
        
def display_visual_element(visual: VisualElement):
    """Display a visual element."""
    st.markdown(f"**{visual.type.capitalize()}**")
    st.markdown(f"*Description:* {visual.description}")
    if visual.data_summary:
        st.markdown(f"*Data Summary:* {visual.data_summary}")

def display_page_details(document):
    """Display page details, handling both model objects and dictionaries."""
    st.header("📑 Page Level Details")
    
    if hasattr(document, 'pages') and document.pages:
        # Create tab labels that handle both dict and model objects
        tab_labels = []
        for page in document.pages:
            if isinstance(page, dict):
                page_num = page.get('page_number', '?')
            else:
                page_num = getattr(page, 'page_number', '?')
            tab_labels.append(f"Page {page_num}")
        
        tabs = st.tabs(tab_labels)
        
        for i, page in enumerate(document.pages):
            # Convert to dict if it's a model
            if hasattr(page, 'model_dump'):
                page_dict = page.model_dump()
            else:
                page_dict = page
                
            with tabs[i]:
                col1, col2 = st.columns([2, 1])
                
                with col1:
                    st.subheader(page_dict.get('title', f'Page {page_dict.get("page_number", "?")}'))
                    st.markdown("#### Topics")
                    st.write(", ".join(page_dict.get('topics', [])))
                    
                    if page_dict.get('summary'):
                        st.markdown("#### Summary")
                        st.markdown(page_dict.get('summary', ''))
                
                with col2:
                    st.markdown("#### Entities")
                    st.write(", ".join(page_dict.get('entities', [])))
                    
                    st.markdown("#### Content Flags")
                    st.write(f"Tables: {'✅' if page_dict.get('has_tables', False) else '❌'}")
                    st.write(f"Visuals: {'✅' if page_dict.get('has_visuals', False) else '❌'}")
                    st.write(f"Numbers: {'✅' if page_dict.get('has_numbers', False) else '❌'}")
                
                # Display tables
                if page_dict.get('tables'):
                    st.markdown("#### Tables")
                    for table in page_dict.get('tables', []):
                        display_table(table)
                
                # Display visuals
                if page_dict.get('visuals'):
                    st.markdown("#### Visual Elements")
                    for visual in page_dict.get('visuals', []):
                        if isinstance(visual, dict):
                            visual_type = visual.get('type', 'Unknown')
                            description = visual.get('description', '')
                            data_summary = visual.get('data_summary', '')
                        else:
                            visual_type = getattr(visual, 'type', 'Unknown')
                            description = getattr(visual, 'description', '')
                            data_summary = getattr(visual, 'data_summary', '')
                            
                        st.markdown(f"**{visual_type.capitalize()}**")
                        st.markdown(f"*Description:* {description}")
                        if data_summary:
                            st.markdown(f"*Data Summary:* {data_summary}")
                
                # Display numerical data
                if page_dict.get('numbers'):
                    st.markdown("#### Numerical Data")
                    for num in page_dict.get('numbers', []):
                        value = num.get('value', '') if isinstance(num, dict) else getattr(num, 'value', '')
                        description = num.get('description', '') if isinstance(num, dict) else getattr(num, 'description', '')
                        context = num.get('context', '') if isinstance(num, dict) else getattr(num, 'context', '')
                        
                        st.markdown(f"**{value}**: {description}")
                        if context:
                            st.markdown(f"*Context: {context}*")
                
                # Display page text
                st.markdown("#### Page Content")
                st.text_area("Page Content", page_dict.get('text', ''), height=300, 
                             key=f"page_text_{page_dict.get('page_number', i)}", 
                             label_visibility="collapsed")  # This hides the label visually but keeps it for accessibility

#############################
# DISPLAY PROCESSING STATUS
#############################

def render_unified_document_report(document_data):
    """All document data in one scrollable view with detailed page tabs"""
    # Convert to dict if it's a Pydantic model
    if hasattr(document_data, 'model_dump'):
        doc = document_data.model_dump()
    else:
        doc = document_data
    
    # Safely get raw_content
    raw_content = doc.get('raw_extracted_content', doc)
    
    # --- Header Section ---
    st.title(raw_content.get('filename', 'Document Report'))
    
    # Count tables/visuals (works with both dict and model)
    def count_attr(items, attr):
        return sum(
            1 for item in items 
            if (isinstance(item, dict) and item.get(attr)) 
            or (hasattr(item, attr) and getattr(item, attr))
        )
    
    st.caption(
        f"🔢 {len(raw_content.get('pages', []))} pages | " +
        f"📊 {count_attr(raw_content.get('pages', []), 'has_tables')} tables | " +
        f"📈 {count_attr(raw_content.get('pages', []), 'has_visuals')} visuals"
    )

    # Create tabs for document view modes
    tab1, tab2 = st.tabs(["Executive Summary", "Detailed Page View"])
    
    with tab1:
        # --- Summary Card ---
        with st.container(border=True):
            summary = raw_content.get('summary', {})
            if isinstance(summary, dict):
                summary_data = summary
            else:
                summary_data = summary.model_dump() if hasattr(summary, 'model_dump') else {}
            
            st.subheader("Executive Summary")
            
            # Display summary text
            if summary_data.get('summary'):
                st.markdown(summary_data.get('summary'))
            else:
                st.info("No summary available for this document.")
            
            # Display themes if available
            if summary_data.get('themes'):
                st.markdown("#### Key Themes")
                themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" 
                                     for tag in summary_data.get('themes', [])])
                st.markdown(themes_html, unsafe_allow_html=True)
            
            # Metrics in columns
            if summary_data.get('key_metrics'):
                st.markdown("#### Key Metrics")
                cols = st.columns(3)
                metrics = summary_data['key_metrics']
                if not isinstance(metrics, list):
                    metrics = list(metrics) if hasattr(metrics, '__iter__') else []
                
                for i, metric in enumerate(metrics[:6]):
                    metric_data = metric if isinstance(metric, dict) else metric.model_dump()
                    cols[i%3].metric(
                        label=metric_data.get('name', 'Metric'),
                        value=metric_data.get('value', 'N/A'),
                        delta=metric_data.get('trend'),
                        help=metric_data.get('context')
                    )
                    
        # --- Simple Page Navigation (keep this for the summary view) ---
        st.divider()
        pages = raw_content.get('pages', [])
        page_options = []
        for p in pages:
            if isinstance(p, dict):
                title = p.get('title', '')
                num = p.get('page_number', len(page_options)+1)
            else:
                title = getattr(p, 'title', '')
                num = getattr(p, 'page_number', len(page_options)+1)
            page_options.append(f"Page {num} - {title}")
        
        selected_page = st.selectbox(
            "Navigate to page:",
            options=page_options,
            key="page_nav_summary"
        )
        page_idx = int(selected_page.split()[1]) - 1
        page_data = pages[page_idx]
        
        # Convert page data to dict if needed
        if hasattr(page_data, 'model_dump'):
            page_data = page_data.model_dump()
        
        # --- Selected Page Content ---
        with st.container(border=True):
            # Text Content
            st.subheader(page_data.get('title', selected_page))
            st.text_area("Full Text", 
                        page_data.get('text', ''), 
                        height=200,
                        label_visibility="collapsed")

            # Tables
            tables = page_data.get('tables', [])
            if tables:
                st.subheader(f"Tables ({len(tables)})")
                for table in tables:
                    if hasattr(table, 'model_dump'):
                        table = table.model_dump()
                    display_table(table)
    
    with tab2:
        # Use your existing detailed page view but adapt it for the new structure
        pages = raw_content.get('pages', [])
        
        # Convert all pages to PageContent objects or compatible dicts
        processed_pages = []
        for page in pages:
            if hasattr(page, 'model_dump'):
                processed_pages.append(page)  # Already a model
            else:
                # It's a dict, either use as is or convert to PageContent
                processed_pages.append(page)
        
        # Create a wrapper object compatible with display_page_details
        document_wrapper = type('ProcessedDocument', (), {'pages': processed_pages})
        
        # Call your existing function
        display_page_details(document_wrapper)

def display_sidebar_chat():
    """Manages the sidebar content: file upload, chat, and triggers processing."""
    st.sidebar.title("📝 Input & Chat")
    st.sidebar.markdown("Upload documents and chat about their content.")

    # Define supported file types
    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt"]

    # Initialize session state keys if they don't exist
    if "processed_file_names" not in st.session_state:
        st.session_state.processed_file_names = set()  # Store strings only
    if "processing_active" not in st.session_state:
        st.session_state.processing_active = False
    if "last_uploaded_files" not in st.session_state:
        st.session_state.last_uploaded_files = []
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "processed_documents" not in st.session_state:
        st.session_state.processed_documents = []  # Can store both dicts and models
    if "selected_docs" not in st.session_state:
        st.session_state.selected_docs = []
    if "show_results" not in st.session_state:
        st.session_state.show_results = False
    if "selected_doc_to_view" not in st.session_state:
        st.session_state.selected_doc_to_view = None

    # --- File Upload Section ---
    st.sidebar.subheader("📁 Upload Documents")
    
    uploaded_files = st.sidebar.file_uploader(
        "Select files:",
        type=supported_types,
        accept_multiple_files=True,
        key="sidebar_file_uploader_main"
    )

    # Detect *new* files to process automatically
    new_files_to_process = []
    current_files = {f.name for f in uploaded_files} if uploaded_files else set()
    previous_files = {f.name for f in st.session_state.last_uploaded_files} if st.session_state.last_uploaded_files else set()
    
    # If there are new files and we're not already processing
    if uploaded_files and (current_files != previous_files) and not st.session_state.processing_active:
        st.session_state.last_uploaded_files = uploaded_files
        
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                new_files_to_process.append(file)
    
    # If we have new files to process, prepare and trigger processing directly
    if new_files_to_process and not st.session_state.processing_active:
        st.sidebar.info(f"Processing {len(new_files_to_process)} new files...")
        
        # Create the status elements directly in the sidebar
        status_container = st.sidebar.status("Starting document processing...", expanded=True)
        progress_bar = st.sidebar.progress(0)
        time_info = st.sidebar.empty()
        
        # Store in session state for access
        st.session_state.status_container = status_container
        st.session_state.progress_bar = progress_bar
        st.session_state.time_info = time_info
        
        # Prepare file data
        files_data = []
        for file in new_files_to_process:
            try:
                file_content = file.getvalue()
                file_ext = file.name.split('.')[-1].lower()
                files_data.append({
                    "name": file.name,
                    "content": file_content,
                    "type": file_ext
                })
            except Exception as e:
                st.sidebar.error(f"Error reading file {file.name}: {e}")

        if files_data:
            # Process documents directly
            try:
                st.session_state.processing_active = True
                
                # Call the main processing function
                processed_docs = run_async(
                    process_all_documents_async,
                    files_data
                )

                # Update session state with results
                new_processed_docs = [doc for doc in processed_docs if doc is not None]
                st.session_state.processed_documents.extend(new_processed_docs)

                # Update the set of processed file names
                for file_info in files_data:
                    st.session_state.processed_file_names.add(file_info["name"])

                # Set flag to show results
                st.session_state.show_results = True
                
                status_container.update(label="✅ Document processing finished!", state="complete")
                
            except Exception as e:
                st.sidebar.error(f"Error during document processing: {e}")
                status_container.update(label=f"❌ Processing Error: {e}", state="error")
            finally:
                # Ensure processing flag is reset
                st.session_state.processing_active = False
    
    # Display already processed files
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---")
        st.sidebar.write("Processed files:")
        for filename in sorted(st.session_state.processed_file_names):
            st.sidebar.caption(f"✓ {filename}")
        
        # Add a button to view results
        if st.sidebar.button("View Document Reports"):
            st.session_state.show_results = True
            st.rerun()

    # --- Chat Section --- 
    st.sidebar.divider()
    st.sidebar.subheader("💬 Document Chat")

    # Document selection for chat context
    if st.session_state.processed_documents:
        doc_options = []
        for doc_result in st.session_state.processed_documents:
            filename = None
            
            if isinstance(doc_result, dict):
                raw_content = doc_result.get("raw_extracted_content", {})
                filename = raw_content.get("filename")
            elif hasattr(doc_result, 'filename'):
                filename = doc_result.filename
            elif hasattr(doc_result, 'raw_extracted_content'):
                filename = doc_result.raw_extracted_content.filename
                
            if filename:
                doc_options.append(filename)

        if doc_options:
            st.session_state.selected_docs = st.sidebar.multiselect(
                "Select documents for chat context:",
                options=sorted(list(set(doc_options))),
                default=st.session_state.selected_docs,
                key="doc_context_selector_sidebar"
            )
        else:
            st.sidebar.caption("No processed documents available to reference.")
    else:
        st.sidebar.caption("No documents processed yet. You can still chat, but I won't be able to reference any document content.")

    # Create a container with border in the sidebar for chat history
    with st.sidebar.container(border=True, height=600):
        st.subheader("Chat History")
        # Display chat history
        if not st.session_state.get("messages", []):
            st.info("No messages yet. Start a conversation below!")
        else:
            for message in st.session_state.messages:
                role = message.get("role", "")
                content = message.get("content", "")
                timestamp = message.get("timestamp", "")
                
                if role == "user":
                    with st.chat_message("user"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Sent: {timestamp}")
                elif role == "assistant":
                    with st.chat_message("assistant"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Received: {timestamp}")
                elif role == "system":
                    with st.chat_message("system"):
                        st.markdown(f"*{content}*")
                        if timestamp:
                            st.caption(f"System: {timestamp}")

        # Chat input
        if prompt := st.sidebar.chat_input("Ask a question...", key="sidebar_chat_input_main"):
            current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages.append({"role": "user", "content": prompt, "timestamp": current_time})
            run_async(process_chat_message, prompt)
            st.rerun()
        
async def process_chat_message(message):
    """Process a chat message with clear user feedback."""
    # Make sure the message list exists
    if "messages" not in st.session_state:
        st.session_state.messages = []
        
    # Add a thinking indicator that will be updated with the real response
    current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.session_state.messages.append({"role": "assistant", "content": "Thinking...", "timestamp": current_time})
    
    try:
        # Get API client
        api_key = get_gemini_api_key()
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": 'https://gateway.helicone.ai',
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets["HELICONE_API_KEY"]}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                }
            })
        
        # Build context from documents (if any)
        context = ""
        selected_docs = st.session_state.get("selected_docs", [])
        processed_docs = st.session_state.get("processed_documents", [])
        
        if selected_docs and processed_docs:
            # Extract content from selected documents for context
            context = "Using document context from: " + ", ".join(selected_docs) + "\n\n"
            
            for doc_name in selected_docs:
                # Find the document in processed_documents
                selected_doc = None
                for doc in processed_docs:
                    if isinstance(doc, dict) and "raw_extracted_content" in doc:
                        filename = doc["raw_extracted_content"].get("filename", "")
                        if filename == doc_name:
                            selected_doc = doc
                            break
                
                if selected_doc and "raw_extracted_content" in selected_doc:
                    # Extract summary and key pages
                    tech_data = selected_doc["raw_extracted_content"]
                    
                    # Add summary if available
                    if "summary" in tech_data and tech_data["summary"]:
                        summary = tech_data["summary"]
                        context += f"Document: {doc_name}\n"
                        context += f"Title: {summary.get('title', '')}\n"
                        context += f"Summary: {summary.get('summary', '')}\n\n"
                    
                    # Add content from up to 3 key pages
                    if "pages" in tech_data and tech_data["pages"]:
                        pages = tech_data["pages"]
                        # Sort pages by importance (tables, numbers, length)
                        sorted_pages = sorted(pages, 
                                            key=lambda p: (p.get("has_tables", False), 
                                                          p.get("has_numbers", False), 
                                                          len(p.get("text", ""))), 
                                            reverse=True)[:3]
                        
                        for page in sorted_pages:
                            context += f"Page {page.get('page_number', '')}: {page.get('text', '')[:500]}...\n\n"
        
        # Simple prompt for demonstration
        prompt = f"""
        You are a helpful assistant that analyzes documents. 
        
        Document context:
        {context}
        
        Please respond to: {message}
        """
        
        # Generate content - using the correct async pattern
        response = await client.aio.models.generate_content(
            model="gemini-2.0-flash",
            contents=[
                types.Content(
                    role="user",
                    parts=[types.Part.from_text(text=prompt)]
                )
            ]
        )
        
        # Update the thinking message with the actual response
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            if response and response.candidates:
                content = response.candidates[0].content.parts[0].text
                st.session_state.messages[-1]["content"] = content
                st.session_state.messages[-1]["timestamp"] = timestamp
            else:
                st.session_state.messages[-1]["content"] = "Sorry, I couldn't generate a response."
                st.session_state.messages[-1]["timestamp"] = timestamp
        
    except Exception as e:
        # Update thinking message with the error
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages[-1]["content"] = f"Sorry, an error occurred: {str(e)}"
            st.session_state.messages[-1]["timestamp"] = timestamp
    
    # Trigger UI update
    st.rerun()

##############
# Main
##############

def main():
    # Initialize session state
    if "processed_documents" not in st.session_state:
        st.session_state.processed_documents = []
    
    # --- SIDEBAR ---
    display_sidebar_chat()
    
    # --- MAIN CONTENT ---
    st.title("📑 Unified Document Report")
    
    if st.session_state.get("processing_active"):
        with st.status("Processing documents..."):
            st.write("This may take a few minutes")
    elif st.session_state.processed_documents:
        # SAFER DOCUMENT SELECTION LOGIC
        if len(st.session_state.processed_documents) > 1:
            try:
                # Get list of valid document names
                doc_options = []
                for d in st.session_state.processed_documents:
                    if isinstance(d, dict):
                        name = d.get("filename") or d.get("raw_extracted_content", {}).get("filename")
                    else:
                        name = getattr(d, "filename", None)
                    if name:
                        doc_options.append(name)
                
                # Default to first document if no selection yet
                default_idx = 0
                if "selected_doc_index" in st.session_state:
                    default_idx = min(st.session_state.selected_doc_index, len(doc_options)-1)
                
                selected_name = st.selectbox(
                    "Choose document:",
                    options=doc_options,
                    index=default_idx,
                    key="doc_selector"
                )
                
                # Find matching document
                doc = None
                for d in st.session_state.processed_documents:
                    current_name = None
                    if isinstance(d, dict):
                        current_name = d.get("filename") or d.get("raw_extracted_content", {}).get("filename")
                    else:
                        current_name = getattr(d, "filename", None)
                    
                    if current_name == selected_name:
                        doc = d
                        break
                
                if not doc:
                    st.warning("Document not found, showing first available")
                    doc = st.session_state.processed_documents[0]
                    
            except Exception as e:
                st.error(f"Document selection error: {e}")
                doc = st.session_state.processed_documents[0]
        else:
            doc = st.session_state.processed_documents[0]
        
        render_unified_document_report(doc)
    else:
        st.info("Upload documents in the sidebar to begin")

if __name__ == "__main__":
    main()
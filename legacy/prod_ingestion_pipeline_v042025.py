import streamlit as st
import os
import json
import asyncio
import base64
import pandas as pd
import pymupdf  # PyMuPDF
from typing import List, Dict, Optional, Union, Any, Set, Tuple
import pydantic
from pydantic import BaseModel, Field, ValidationError, field_validator, ConfigDict
from google import genai
from google.genai import types
import docx
from pptx import Presentation
import csv
import openpyxl
import re
from collections import Counter, defaultdict
import time
import io
import datetime
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import logging
import traceback
import pathlib
import random
###############################
# DATA MODELS
###############################

class VisualElement(BaseModel):
    """Model for visual elements like charts, graphs, etc."""
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data, or N/A") # Updated description to reflect N/A possibility
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")
    source_url: Optional[str] = Field(None, description="Source URL of the visual")
    alt_text: Optional[str] = Field(None, description="Alternative text for the visual")
    visual_id: str = Field(..., description="Unique identifier for the visual") # Made required as per prompt

class NumericalDataPoint(BaseModel):
    """Model for numerical data points extracted from text."""
    value: str = Field(..., description="The numerical value as a string") # Clarified string type
    description: str = Field(..., description="What the number represents, with units/context") # Made required
    context: Optional[str] = Field(None, description="Surrounding text context, if available") # Changed to Optional[str] and default to None

class TableData(BaseModel):
    """Model for tables extracted from documents."""
    table_content: str = Field(..., description="Markdown formatted table content") # Made required
    title: Optional[str] = Field(None, description="Optional table title") # Kept optional
    summary: Optional[str] = Field(None, description="Optional table summary") # Kept optional
    page_number: int = Field(..., description="Page number where the table appears") # Made required
    table_id: str = Field(..., description="Unique identifier for the table") # Made required

class Subsection(BaseModel):
    """Model for subsections extracted from pages."""
    subsection_id: str = Field(..., description="Unique identifier for the subsection")
    order: int = Field(..., description="Order of the subsection within the page")
    title: str = Field(..., description="Title of the subsection (less than 7 words)") # Added constraint hint
    text: str = Field(..., description="Full text content of the subsection") # Renamed from content to match prompt output
    description: str = Field(..., description="One line description summarizing the main point")  # Made required as per schema
    is_cutoff: bool = Field(..., description="True if content appears to be cut off by page break") # Made required
    referenced_visuals: List[str] = Field(default_factory=list, description="IDs of referenced visuals")
    referenced_tables: List[str] = Field(default_factory=list, description="IDs of referenced tables")

# --- Simplified PageContent Model ---
class PageContent(BaseModel):
    """
    Model for structured content extracted from a single page.
    Stage 1 populates raw_text, tables, visuals, numbers.
    Stage 2 populates subsections from raw_text.
    """
    page_number: int = Field(..., description="Page number in the original document")
    has_tables: bool = Field(..., description="Indicates if any tables were extracted from this page")
    has_visuals: bool = Field(..., description="Indicates if any visuals (charts, images, etc.) were extracted")
    has_numbers: bool = Field(..., description="Indicates if any key numerical data points were extracted")
    tables: List[TableData] = Field(default_factory=list, description="List of tables extracted from the page")
    visuals: List[VisualElement] = Field(default_factory=list, description="List of visual elements extracted from the page")
    numbers: List[NumericalDataPoint] = Field(default_factory=list, description="List of numerical data points extracted from the page")
    # NEW: Store raw text from initial extraction
    raw_text: Optional[str] = Field(None, description="Full raw text extracted initially from the page/chunk")
    # Subsections are now populated in a second step
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections extracted from raw_text")

    model_config = ConfigDict(arbitrary_types_allowed=True, extra='forbid')

# --- New Models for Qdrant Structure ---
class LinkedNode(BaseModel):
    """Structure for linked nodes metadata in Qdrant."""
    target_node_id: str = Field(..., description="ID of the target node (e.g., type_name_chX)")
    relationship_description: str = Field(..., description="Nature of the relationship from source to target")
    relationship_keywords: List[str] = Field(default_factory=list, description="Keywords summarizing the relationship")
    relationship_strength: float = Field(..., description="Strength of the relationship (1-10)")

    @field_validator('relationship_strength', mode='before')
    def strength_to_float(cls, v):
        try:
            return float(v)
        except (ValueError, TypeError):
            return 5.0 # Default strength on conversion error

class QdrantNode(BaseModel):
    """Structure for nodes to be stored in Qdrant, using a properties dictionary."""
    node_id: str = Field(..., description="Unique node ID (e.g., type_sanitizedname_chX)")
    node_type: str = Field(..., description="'entity' or 'concept'")
    name: str = Field(..., description="Entity or concept name")
    # description: str = Field(..., description="Comprehensive description from text") # <-- REMOVED
    chapter_id_context: str = Field(..., description="Chapter ID where this node was primarily identified")
    # document_context: str = Field(..., description="Identifier of the source document") # <-- REMOVED (now in properties.source)

    properties: Dict[str, str] = Field(default_factory=dict, description="Intrinsic properties of the node")
    # Example keys in properties:
    #   title: str
    #   source: str (filename/URL)
    #   source_credibility: str
    #   information_type: str
    #   core_claim: str (replaces description)
    #   key_entities: str (JSON stringified list)
    #   sentiment_tone: str
    #   tags: str (JSON stringified list)
    #   stored_at: str (ISO timestamp)

    linked_nodes: List[LinkedNode] = Field(default_factory=list, description="List of outgoing relationships")

    # Optional: Backward compatibility property (as per your instructions)
    @property
    def description(self) -> str | None:
        """Alias for core_claim for backward compatibility."""
        return self.properties.get("core_claim")

    @description.setter
    def description(self, value: str):
        """Set core_claim when description is set."""
        self.properties["core_claim"] = value

    model_config = ConfigDict(extra='ignore') # Allow ignoring extra fields if needed

class Chapter(BaseModel):
    """Model for chapters composed of subsections, including Qdrant node data."""
    chapter_id: str = Field(..., description="Unique identifier for the chapter")
    title: str = Field(..., description="Title of the chapter")
    summary: str = Field(..., description="Summary of the chapter")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections in this chapter")
    # REMOVE THIS LINE:
    # entity_relationships: List[EntityRelationship] = Field(default_factory=list, description="Legacy entity relationships within the chapter")
    order: int = Field(..., description="Order of the chapter in the document")
    # New field to store structured node data for Qdrant
    qdrant_nodes: Optional[List[QdrantNode]] = Field(None, description="Structured nodes and relationships for Qdrant")
    
class DocumentSummaryDetails(BaseModel):
    """Structure for summary details generated from node analysis."""
    title: str = Field(..., description="Concise title for the document based on nodes")
    themes: List[str] = Field(default_factory=list, description="Main themes/topics derived from nodes")
    questions: List[str] = Field(default_factory=list, description="Sample questions reflecting node content")
    summary: str = Field(..., description="Comprehensive summary synthesizing node information")
    # Note: No chapters, entity_relationships, table/visual summaries here

class KeyTopic(BaseModel):
    name: str = Field(..., description="Short topic name (1-3 words)")
    description: str = Field(..., description="Brief description of the topic")
    relevance: str = Field(..., description="Relevance level (High/Medium/Low)")
    sentiment: str = Field(..., description="Sentiment analysis (Positive/Neutral/Mixed/Cautionary/Negative)")
    analysis: str = Field(..., description="Brief analysis of the topic")

class QuotedStatement(BaseModel):
    speaker: str = Field(..., description="Person or entity who made the statement")
    quote: str = Field(..., description="The quoted text")
    page: int = Field(..., description="Page number where the quote appears")

class DocumentReport(BaseModel):
    file_name: str = Field(..., description="Original filename")
    page_count: int = Field(..., description="Number of pages in document")
    title: str = Field(..., description="Document title")
    title_summary: str = Field(..., description="Brief summary/subtitle")
    concept_theme_hashtags: List[str] = Field(default_factory=list, description="Theme tags as hashtags")
    date_published: Optional[str] = Field(None, description="Publication date if available")
    source: Optional[str] = Field(None, description="Document source")
    confidence: str = Field("Medium", description="Analysis confidence (High/Medium/Low)")
    document_summary: str = Field(..., description="Comprehensive summary")
    key_insights: List[str] = Field(default_factory=list, description="Key takeaways and insights")
    key_topics: List[KeyTopic] = Field(default_factory=list, description="Detailed topic analysis")
    quoted_statements: List[QuotedStatement] = Field(default_factory=list, description="Important quotes from document")
    content_excerpt: Optional[str] = Field(None, description="Highlighted document content")
    chapters: List[Chapter] = Field(default_factory=list, description="Document chapters")

class ProcessedDocument(BaseModel):
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages") # This now uses the simplified PageContent
    summary: Optional[DocumentSummaryDetails] = None

class ProjectOntology(BaseModel):
    title: str = Field(..., description="Project title")
    overview: str = Field(..., description="Project overview based on synthesized nodes")
    document_count: int = Field(..., description="Number of documents analyzed")
    documents: List[str] = Field(..., description="List of document filenames included in analysis")
    global_themes: List[str] = Field(..., description="High-level project themes derived from nodes")
    key_concepts: List[str] = Field(..., description="Key concepts identified across all documents from nodes")
    # NEW FIELD: Stores the aggregated nodes including inter-document relationships
    project_graph_nodes: Optional[List[QdrantNode]] = Field(None,
                                description="Aggregated nodes and their relationships (including inter-document links) across the project.")

    model_config = ConfigDict(arbitrary_types_allowed=True) # If QdrantNode uses non-standard types


###############################
# API AND UTILITY FUNCTIONS
###############################

def convert_to_model(data, model_class):
    """Convert dictionary data to a Pydantic model with error handling."""
    if not isinstance(data, dict):
        print(f"Warning: Expected dict for {model_class.__name__}, got {type(data)}")
        return None
    
    try:
        # Extract only the fields that the model expects
        model_fields = model_class.__annotations__.keys()
        filtered_data = {k: v for k, v in data.items() if k in model_fields}
        return model_class(**filtered_data)
    except Exception as e:
        print(f"Error converting to {model_class.__name__}: {str(e)}")
        return None

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
        st.stop()
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)
            
            # Validate JSON if applicable
            config = kwargs.get('config')
            if (config and hasattr(config, 'response_mime_type') and 
                config.response_mime_type == 'application/json' and 
                hasattr(response, 'candidates') and response.candidates):
                try:
                    # Try to parse the JSON response
                    json_text = response.candidates[0].content.parts[0].text
                    json.loads(clean_json_response(json_text))  # Use our clean function first
                except json.JSONDecodeError as e:
                    if attempt < max_retries - 1:
                        print(f"Received malformed JSON on attempt {attempt+1}, retrying: {e}")
                        await asyncio.sleep(2 ** attempt)
                        continue
            
            return response
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            print(f"API call failed on attempt {attempt+1}, retrying: {e}")
            await asyncio.sleep(2 ** attempt)

def clean_json_response(json_text: str, extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        return "{}"
    
    # Ensure it's a string
    json_text = str(json_text)
    
    try:
        # Handle markdown code blocks
        if json_text.startswith("```"):
            blocks = json_text.split("```")
            if len(blocks) >= 3:
                json_text = blocks[1]
                if json_text.startswith("json"):
                    json_text = json_text[4:].strip()
            else:
                json_text = json_text.replace("```", "").strip()
        
        elif json_text.startswith("json"):
            json_text = json_text[4:].strip()
        
        # Remove any leading/trailing whitespace
        json_text = json_text.strip()
        
        # Attempt to parse the JSON as is
        try:
            json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # Continue with fixes
            pass
        
        # Common JSON fixes
        # Fix missing quotes around property names
        json_text = re.sub(r'([{,]\s*)(\w+)(\s*:)', r'\1"\2"\3', json_text)
        
        # Fix trailing commas in arrays/objects
        json_text = re.sub(r',(\s*[}\]])', r'\1', json_text)
        
        # Fix missing quotes around string values
        # This is trickier and might cause issues with valid JSON if not careful
        
        # Final attempt to parse with Python's JSON parser
        try:
            data = json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # If we're extracting text on failure and can't fix the JSON
            if extract_text_on_failure:
                logging.warning(f"Initial JSON parsing failed. Raw Response Text: '{json_text[:2000]}...'") # Log raw response
                
                # Extract what appears to be the main text content
                # Look for patterns that would indicate text fields in a JSON response
                text_content = ""
                
                # Try to find text between quotes after "text": or "content": patterns
                text_matches = re.findall(r'"(?:text|content)"\s*:\s*"([^"]+)"', json_text)
                if text_matches:
                    text_content = " ".join(text_matches)
                
                # If that doesn't work, take any text between quotes as potential content
                if not text_content:
                    # Get all quoted strings that look reasonably long (over 20 chars)
                    quoted_text = re.findall(r'"([^"]{20,})"', json_text)
                    if quoted_text:
                        text_content = " ".join(quoted_text)
                
                # If still nothing, just take the longest line that's not mostly symbols
                if not text_content:
                    lines = json_text.split('\n')
                    content_lines = [line for line in lines if len(line) > 50 and 
                                    sum(c.isalpha() or c.isspace() for c in line) / len(line) > 0.7]
                    if content_lines:
                        text_content = " ".join(content_lines)
                
                # Fallback - create a minimal valid JSON with the extracted text
                fallback_json = {
                    "text": text_content or "Extracted text content not found",
                    "title": "Text Extraction Fallback",
                    "topics": ["text extraction"],
                    "summary": "Structured JSON parsing failed - basic text extraction used as fallback.",
                    "entities": [],
                    "has_tables": False,
                    "has_visuals": False,
                    "has_numbers": False,
                    "dates": [],
                    "tables": [],
                    "visuals": [],
                    "numbers": [],
                    "financial_statements": [],
                    "key_metrics": [],
                    "financial_terms": [],
                    "subsections": []
                }
                return json.dumps(fallback_json)
            else:
                # Return minimal valid JSON as before
                return '{"text":"Error parsing JSON response","error":"Invalid JSON structure"}'
            
    except Exception as e:
        logging.error(f"Error cleaning JSON response: {e}")
        
        if extract_text_on_failure:
            # Create a fallback with whatever raw content we have
            raw_text = json_text.strip()
            if len(raw_text) > 10000:  # Truncate very long responses
                raw_text = raw_text[:10000] + "... (truncated)"
                
            fallback_json = {
                "text": raw_text,
                "title": "Raw Text Fallback",
                "topics": ["raw text"],
                "summary": "Response processing failed - returning raw text.",
                "entities": [],
                "has_tables": False,
                "has_visuals": False,
                "has_numbers": False,
                "dates": [],
                "tables": [],
                "visuals": [],
                "numbers": [],
                "financial_statements": [],
                "key_metrics": [],
                "financial_terms": [],
                "subsections": []
            }
            return json.dumps(fallback_json)
        else:
            return '{"text":"Error cleaning JSON response","error":"' + str(e).replace('"', '\\"') + '"}'

async def run_in_executor(func, *args):
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, func, *args) # Use default executor

def run_async(func, *args, **kwargs):
    """Run an async function from Streamlit."""
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        # If func is already a coroutine, run it directly
        if asyncio.iscoroutine(func):
            return loop.run_until_complete(func)
        # Otherwise call it with the provided args
        return loop.run_until_complete(func(*args, **kwargs))
    finally:
        loop.close()

###############################
# CONTENT EXTRACTION FUNCTIONS
###############################

async def process_page_with_fallback(client, image_part, page_num, filename=None):
    """Process page with multi-level fallback including pure text extraction."""
    source_identifier = filename or "Unknown_Document"
    
    try:
        # First attempt: Full structured extraction
        logging.debug(f"Attempting full subsection extraction for page {page_num}")
        page_content = await extract_page_content_from_memory(client, image_part, page_num, source_identifier)

        # Check if it failed severely (e.g., returned an error page type, or text is placeholder)
        is_error = any(s.title == "Processing Error" for s in page_content.subsections) if page_content.subsections else False # Check error subsection
        if is_error or page_content.raw_text == "[LLM failed to extract raw text]":
             logging.warning(f"Initial structure extraction failed or yielded no text for page {page_num} of '{source_identifier}'. Triggering fallback.")
             raise ValueError("Initial extraction failed or returned no text") # Go to fallback

        logging.info(f"Successfully extracted structure (incl. raw_text) for page {page_num} of '{source_identifier}'")
        return page_content

    except Exception as e:
        logging.warning(f"Initial extraction failed for page {page_num} of '{source_identifier}': {e}. Trying pure text fallback.")

        # --- Fallback attempt: Pure text extraction ---
        try:
            logging.debug(f"Attempting pure text extraction fallback for page {page_num} of '{source_identifier}'")
            text_only_prompt = f"Extract ONLY the plain text content from page {page_num} of this document. Preserve paragraph breaks. Do NOT format as JSON."

            text_response = await retry_api_call(
                 client.aio.models.generate_content,
                 model="gemini-2.0-flash", # Model suitable for text extraction
                 contents=[
                    types.Content(parts=[image_part, types.Part.from_text(text=text_only_prompt)]),
                 ],
            )

            raw_text = ""
            if text_response.candidates:
                raw_text = text_response.candidates[0].content.parts[0].text.strip()

            if not raw_text:
                 raw_text = "[No text extracted during fallback]"
                 logging.warning(f"Pure text fallback for page {page_num} resulted in empty text.")

            logging.info(f"Successfully extracted pure text fallback for page {page_num} of '{source_identifier}'")
            # Create a PageContent object with ONLY raw_text and page_number populated
            return PageContent(
                page_number=page_num,
                raw_text=raw_text,
                # Set defaults for other required fields
                has_tables=False, has_visuals=False, has_numbers=False,
                tables=[], visuals=[], numbers=[], subsections=[]
            )

        except Exception as text_error:
            logging.error(f"Pure text extraction fallback also failed for page {page_num} of '{source_identifier}': {text_error}", exc_info=True)
            # Final fallback: return an error page object
            return create_error_page(page_num, f"All extraction methods failed: {text_error}")

# --- Define the NEW prompt for subsection extraction from text ---
subsection_extraction_prompt = """
Analyze the provided text content from a single page ({page_num}) of a document, originating from source '{source_identifier}'. Your goal is to segment this text into logical subsections, providing concise yet **informative descriptions geared towards identifying entities and relationships later**.

Input Text:
--- START TEXT ---
{raw_text_content}
--- END TEXT ---

Instructions:
1.  Read the text and identify logical breaks based on topic shifts, headings (if any within the text), or paragraph structure.
2.  Segment the text into sequential subsections.
3.  For each subsection, determine:
    *   `subsection_id`: Generate a unique ID formatted as `page_{page_num}_section_ORDER` where ORDER is the sequential order number.
    *   `order`: The sequential order number (integer, starting from 1).
    *   `title`: A concise title (less than 7 words). Use headings from the text if present and appropriate, otherwise generate a descriptive title.
    *   `text`: The full text content belonging *only* to this specific subsection. Ensure all original text is captured across subsections.
    *   `description`: **CRITICAL FOR RELATIONSHIP ANALYSIS:** A concise, one-sentence summary (approx. 15-25 words) that **highlights the key entities, concepts, or specific actions/interactions** discussed **within this specific subsection**. This description acts as a vital preview for downstream entity and relationship identification.
        *   **Focus:** Mention the main actors (e.g., 'Parsely', 'GPT-40', 'CafeCorner LLC') and the core action or relationship involving them *in this text block* (e.g., 'developed by', 'utilizes', 'integrates', 'enhances', 'piloted for').
        *   **Example Good Descriptions:** "Explains how Parsely integrates Llamalndex and Qdrant Vector DB.", "Details CafeCorner LLC's role in developing the Parsely system.", "Introduces GPT-40 as the LLM used for internal analysis.", "Discusses the enhancement of Pitchbook data distribution via Parsely.", "Outlines the piloting of Parsely for JPMorgan's ChatIQ assistant."
        *   **Example Less Useful Description:** "This section talks about the system." (Too generic - doesn't mention specific entities or actions).
    *   `is_cutoff`: Set to `true` ONLY if this specific text chunk appears to end mid-sentence or mid-thought, suggesting it was truncated *before* being passed to you. Otherwise `false`.
    *   `referenced_visuals`: An empty list `[]`. (Visuals/Tables extracted previously).
    *   `referenced_tables`: An empty list `[]`. (Visuals/Tables extracted previously).

4.  Return ONLY a valid JSON array containing the subsection objects, structured exactly like this example:
    ```json
    [
      {{
        "subsection_id": "page_{page_num}_section_1",
        "order": 1,
        "title": "Parsely System Overview",
        "text": "The full text content of the first subsection detailing Parsely...",
        "description": "Introduces the Parsely system developed at CafeCorner LLC for data enhancement.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }},
      {{
        "subsection_id": "page_{page_num}_section_2",
        "order": 2,
        "title": "Core Technologies",
        "text": "Details on the integration of various technologies...",
        "description": "Details Parsely's use of GPT-40, Llamalndex, and Qdrant Vector DB for analysis.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }}
      // ... more subsection objects
    ]
    ```
CRITICAL: Your entire response MUST be ONLY the JSON array, starting with `[` and ending with `]`. Do NOT include ```json markdown fences, introductory text, explanations, or any other characters outside the JSON structure. Ensure all string values are enclosed in double quotes and properly escaped. Validate your JSON structure carefully before outputting.
"""

# --- NEW Prompt for JSON Repair ---
json_repair_prompt_template = """
The following text was generated by an AI attempting to create a valid JSON array of subsection objects based on previous instructions. However, the text has syntax errors and is not valid JSON.

Your task is to analyze the provided text, correct the JSON syntax errors, and return ONLY the corrected, valid JSON array. Do NOT add any introductory text, explanations, or markdown formatting like ```json. The output must start with `[` and end with `]`.

Ensure:
- All strings are enclosed in double quotes.
- Internal quotes within strings are properly escaped (e.g., \\").
- Objects within the array are correctly separated by commas.
- There are no trailing commas after the last element in the array or the last property in an object.
- All brackets (`[]`) and braces (`{{}}`) are correctly paired and closed.

Faulty Text to Repair:
--- START FAULTY TEXT ---
{faulty_json_text}
--- END FAULTY TEXT ---

Return ONLY the valid JSON array:
"""


async def extract_subsections_from_text(
    client: genai.Client,
    raw_text: str,
    page_num: int,
    source_identifier: str,
    max_initial_retries: int = 2, # Retries for the initial API call
    retry_delay: int = 2
) -> List[Subsection]:
    """
    Takes raw text from a page and uses an LLM to segment it into Subsection objects.
    Includes robustness for dict-wrapped lists and a secondary LLM call to attempt
    JSON repair if initial parsing fails.
    Incorporates enhanced logging for debugging.

    Args:
        client: The configured GenAI client.
        raw_text: The raw text content of the page.
        page_num: The page number.
        source_identifier: Identifier of the source document (e.g., filename).
        max_initial_retries: Max attempts for the first LLM call.
        retry_delay: Base delay (seconds) between retries.

    Returns:
        A list of validated Subsection objects, or a list containing a single
        error fallback subsection if processing fails after all attempts.
    """
    function_name = "extract_subsections_from_text_with_repair" # Updated name
    logging.debug(f"[{function_name}] Starting subsection extraction for page {page_num} of '{source_identifier}'")

    # --- Input Validation ---
    if not raw_text or not raw_text.strip():
        logging.warning(f"[{function_name}] No raw text provided for page {page_num} of '{source_identifier}'. Returning empty list.")
        return []

    # --- Prepare Initial Prompt ---
    MAX_TEXT_LIMIT = 100000 # Adjust as needed
    text_for_prompt = raw_text
    if len(raw_text) > MAX_TEXT_LIMIT:
        logging.warning(f"[{function_name}] Truncating raw text for page {page_num} of '{source_identifier}' from {len(raw_text)} to {MAX_TEXT_LIMIT} chars.")
        text_for_prompt = raw_text[:MAX_TEXT_LIMIT] + "... [Content Truncated]"

    try:
        initial_prompt = subsection_extraction_prompt.format(
            page_num=page_num,
            raw_text_content=text_for_prompt,
            source_identifier=source_identifier
        )
    except KeyError as fmt_err:
         logging.error(f"[{function_name}] Error formatting initial prompt for page {page_num}: Missing key {fmt_err}. Using default prompt.")
         initial_prompt = f"Extract subsections as JSON for page {page_num}. Text: {text_for_prompt}"

    # --- Initial LLM Call and Retry Logic ---
    final_list_data: Optional[List[Dict]] = None # Variable to hold the successfully parsed list data
    last_exception: Optional[Exception] = None # Store the last exception encountered
    initial_call_failed_parsing = False # Flag if the first call resulted in parse failure

    for attempt in range(max_initial_retries):
        logging.debug(f"[{function_name}] Initial LLM Call Attempt {attempt + 1}/{max_initial_retries} for page {page_num}")
        last_exception = None # Reset exception for this attempt
        try:
            response = await retry_api_call( # Handles API-level retries (timeouts, server errors)
                client.aio.models.generate_content,
                # Consider using a more robust model here if issues persist
                model="gemini-2.5-flash-preview-04-17", 
                contents=[types.Content(parts=[types.Part.from_text(text=initial_prompt)])],
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.2 # Lower temperature for consistency
                    ),
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates from initial LLM call (Attempt {attempt + 1}) for page {page_num}.")
                # If no candidates after retries, we can't proceed to repair
                if attempt == max_initial_retries - 1:
                    final_list_data = None
                    initial_call_failed_parsing = False # Call itself failed, not parsing
                    last_exception = ValueError("LLM response had no candidates after retries.")
                # Continue to next retry if not last attempt
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue

            # --- Clean and Attempt Initial Parse ---
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"[{function_name}] Initial Cleaned JSON Response (Attempt {attempt + 1}, Page {page_num}): {json_text[:1000]}...") # Log more for debug

            try:
                # Check if it looks like a list before parsing fully
                if not isinstance(json_text, str) or not json_text.strip().startswith('['):
                     logging.warning(f"[{function_name}] Initial response doesn't start with '[' (Attempt {attempt+1}, Page {page_num}). Raw: {json_text[:200]}...")
                     # Treat this as a potential parsing failure for retry/repair logic
                     raise json.JSONDecodeError("Response does not start with '['", json_text, 0)

                subsections_data = json.loads(json_text)

                # --- Initial Structure Check ---
                if isinstance(subsections_data, list):
                    final_list_data = subsections_data # Correct structure received
                    initial_call_failed_parsing = False # Success!
                    logging.debug(f"[{function_name}] Initial call successful with list structure on attempt {attempt + 1} for page {page_num}.")
                    break # Success, exit retry loop

                elif isinstance(subsections_data, dict):
                    logging.warning(f"[{function_name}] Initial call returned a dict (Attempt {attempt + 1}, Page {page_num}). Checking common keys...")
                    possible_keys = ["subsections", "data", "items", "result", "page_subsections", "sections"]
                    found = False
                    for key in possible_keys:
                        if key in subsections_data and isinstance(subsections_data.get(key), list):
                            logging.info(f"[{function_name}] Found list under key '{key}' in dict. Using this list.")
                            final_list_data = subsections_data[key]
                            initial_call_failed_parsing = False # Recovered successfully
                            found = True
                            break
                    if found:
                        break # Success (recovered), exit retry loop
                    else:
                         # Could not find list within dict, treat as structural failure for retry/repair
                         logging.warning(f"[{function_name}] Could not find list within dict structure (Attempt {attempt+1}, Page {page_num}).")
                         raise json.JSONDecodeError("Dict received, but known list keys not found", json_text, 0)
                else:
                    # Incorrect type altogether
                    logging.warning(f"[{function_name}] Initial call response was not list or dict (type: {type(subsections_data)}) on attempt {attempt+1}, Page {page_num}.")
                    raise json.JSONDecodeError(f"Expected list or dict, got {type(subsections_data)}", json_text, 0)

            except json.JSONDecodeError as json_err:
                last_exception = json_err # Store the parsing error
                logging.warning(f"[{function_name}] Initial JSONDecodeError on attempt {attempt + 1} for page {page_num}: {json_err}. Raw text starts with: {json_text[:200]}...")
                if attempt < max_initial_retries - 1:
                    logging.info(f"[{function_name}] Retrying initial call after delay due to JSON decode error.")
                    await asyncio.sleep(retry_delay * (attempt + 1))
                    continue # Go to next initial call attempt
                else:
                    logging.error(f"[{function_name}] Failed to decode JSON from initial call after {max_initial_retries} attempts for page {page_num}.")
                    initial_call_failed_parsing = True # Mark that parsing failed on the last attempt
                    final_list_data = None # Ensure we don't proceed with bad data
                    # Store the faulty text for potential repair attempt
                    faulty_json_text_to_repair = json_text
                    break # Exit initial retry loop, proceed to repair attempt

        except Exception as e:
            last_exception = e # Store any other exception
            logging.error(f"[{function_name}] Error during initial LLM call/processing (Attempt {attempt + 1}) for page {page_num}: {e}", exc_info=True)
            if attempt < max_initial_retries - 1:
                logging.info(f"[{function_name}] Retrying initial call after delay due to unexpected error.")
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue
            else:
                 logging.error(f"[{function_name}] Initial call failed after {max_initial_retries} attempts for page {page_num} due to error: {e}")
                 final_list_data = None # Signal failure
                 initial_call_failed_parsing = False # The call itself failed, not parsing
                 break # Exit retry loop

    # --- JSON Repair Attempt (if initial parsing failed on last try) ---
    if initial_call_failed_parsing and faulty_json_text_to_repair:
        logging.info(f"[{function_name}] Initial JSON parsing failed for page {page_num}. Attempting self-repair LLM call.")
        try:
            repair_prompt = json_repair_prompt_template.format(faulty_json_text=faulty_json_text_to_repair)

            repair_response = await retry_api_call( # Use retry for the repair call too (1 attempt usually enough)
                client.aio.models.generate_content,
                # Use a capable model for repair
                model="gemini-2.0-flash", 
                contents=[types.Content(parts=[types.Part.from_text(text=repair_prompt)])],
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.1 # Very low temperature for repair
                ),
                max_retries=1 # Only retry the repair call once on API errors
            )

            if repair_response and repair_response.candidates:
                repaired_json_text = clean_json_response(repair_response.candidates[0].content.parts[0].text)
                logging.debug(f"[{function_name}] Repaired JSON Response (Page {page_num}): {repaired_json_text[:1000]}...")
                try:
                    # Attempt to parse the REPAIRED text
                    repaired_data = json.loads(repaired_json_text)
                    if isinstance(repaired_data, list):
                        logging.info(f"[{function_name}] Successfully parsed JSON after repair call for page {page_num}.")
                        final_list_data = repaired_data # Use the repaired data
                    elif isinstance(repaired_data, dict) and "subsections" in repaired_data and isinstance(repaired_data["subsections"], list):
                        # Handle if repair somehow wrapped it again
                        logging.warning(f"[{function_name}] Repaired JSON was dict with 'subsections' list for page {page_num}. Using list.")
                        final_list_data = repaired_data["subsections"]
                    else:
                        logging.error(f"[{function_name}] Repaired JSON was not a list or expected dict structure for page {page_num}. Type: {type(repaired_data)}")
                        final_list_data = None # Repair failed structurally
                        last_exception = json.JSONDecodeError("Repaired JSON structure invalid", repaired_json_text, 0)

                except json.JSONDecodeError as repair_err:
                    logging.error(f"[{function_name}] Failed to decode JSON even after repair attempt for page {page_num}: {repair_err}")
                    final_list_data = None # Repair failed parsing
                    last_exception = repair_err # Store the repair parsing error
            else:
                logging.error(f"[{function_name}] Repair LLM call returned no candidates for page {page_num}.")
                final_list_data = None # Repair call failed
                last_exception = ValueError("Repair LLM call had no candidates.")

        except Exception as repair_call_err:
            logging.error(f"[{function_name}] Error during repair LLM call for page {page_num}: {repair_call_err}", exc_info=True)
            final_list_data = None # Repair call failed
            last_exception = repair_call_err


    # --- Process the final data (if successful) or fallback ---
    if final_list_data is not None and isinstance(final_list_data, list): # Check if we successfully got a list
        validated_subsections: List[Subsection] = []
        logging.debug(f"[{function_name}] Processing {len(final_list_data)} potential subsection items for page {page_num}.")
        items_skipped_validation = 0
        for i, sub_dict in enumerate(final_list_data):
            # Skip items that are not dicts or have empty text
            if not isinstance(sub_dict, dict):
                logging.warning(f"[{function_name}] Skipping item {i+1} (not a dict) for page {page_num}: {str(sub_dict)[:100]}...")
                items_skipped_validation += 1
                continue
            subsection_text = sub_dict.get("text", "")
            if not subsection_text or not subsection_text.strip():
                logging.warning(f"[{function_name}] Skipping item {i+1} with missing/empty 'text' for page {page_num}: {str(sub_dict)[:100]}...")
                items_skipped_validation += 1
                continue

            # Proceed with validation
            try:
                # Set defaults before validation
                sub_dict.setdefault("referenced_visuals", [])
                sub_dict.setdefault("referenced_tables", [])
                sub_dict.setdefault("is_cutoff", False)

                # Ensure order is int, default to index
                try:
                    sub_dict['order'] = int(sub_dict.get('order', i + 1))
                except (ValueError, TypeError):
                    logging.warning(f"[{function_name}] Converting order '{sub_dict.get('order')}' to int failed for item {i+1} page {page_num}. Using index {i+1}.")
                    sub_dict['order'] = i + 1

                subsection_obj = Subsection(**sub_dict)
                validated_subsections.append(subsection_obj)

            except ValidationError as ve:
                logging.warning(f"[{function_name}] Pydantic validation failed for subsection item {i+1} on page {page_num}: {ve}. Skipping. Data: {sub_dict}")
                items_skipped_validation += 1
            except Exception as val_err:
                logging.error(f"[{function_name}] Unexpected error during subsection validation {i+1} on page {page_num}: {val_err}. Skipping. Data: {sub_dict}", exc_info=True)
                items_skipped_validation += 1

        # Post-process validated subsections
        if validated_subsections:
            # Sort by order
            validated_subsections.sort(key=lambda s: getattr(s, 'order', float('inf')))

            # Re-assign sequential order
            for i, sub in enumerate(validated_subsections):
                if sub.order != i + 1:
                    logging.debug(f"[{function_name}] Re-assigning order for subsection '{sub.subsection_id}' from {sub.order} to {i+1} on page {page_num}.")
                    sub.order = i + 1

            logging.info(f"[{function_name}] Successfully extracted and validated {len(validated_subsections)} non-empty subsections for page {page_num} (skipped {items_skipped_validation} during validation).")
            return validated_subsections
        else:
            # Fallback if validation filtered everything out
            logging.warning(f"[{function_name}] Processing resulted in zero valid subsections after validation for page {page_num}. Using fallback.")
            return [create_error_fallback_subsection(page_num, raw_text, "LLM response processed, but no valid/non-empty subsections remained after validation.")]

    else:
        # Fallback if all initial attempts AND the repair attempt failed
        fail_reason = "Initial LLM calls failed or returned invalid structure."
        if initial_call_failed_parsing:
             fail_reason = "Initial JSON parsing failed, and subsequent repair attempt also failed."
        if last_exception:
             fail_reason += f" Last known error: {type(last_exception).__name__}: {str(last_exception)[:100]}..." # Add error context

        logging.error(f"[{function_name}] All subsection extraction attempts failed for page {page_num}. Using fallback. Reason: {fail_reason}")
        logging.debug(f"[{function_name}] Input raw_text for failed page {page_num}:\n--- START RAW TEXT ---\n{raw_text[:1000]}...\n--- END RAW TEXT ---")
        return [create_error_fallback_subsection(page_num, raw_text, fail_reason)]


def create_error_fallback_subsection(page_num: int, raw_text: str, error_msg: str) -> Subsection:
     """
     Creates a single fallback Subsection object when subsection extraction fails.
     This ensures the raw text is preserved along with the error context.
     """
     fallback_id = f"page_{page_num}_section_fallback_error"
     fallback_title = "Full Page Content (Subsection Extraction Failed)"
     # Combine error and original text for context
     fallback_text = (
         f"--- ERROR DURING SUBSECTION EXTRACTION ---\n"
         f"{error_msg}\n"
         f"--- END ERROR ---\n\n"
         f"--- ORIGINAL RAW TEXT FOR PAGE {page_num} ---\n"
         f"{raw_text}\n"
         f"--- END RAW TEXT ---"
     )
     fallback_description = f"Could not segment page {page_num} into subsections due to error: {error_msg[:100]}..." # Truncate long errors

     try:
          # Create and return a single Subsection object
          return Subsection(
              subsection_id=fallback_id,
              order=1, # Always the first (and only) subsection in this fallback case
              title=fallback_title,
              text=fallback_text,
              description=fallback_description,
              is_cutoff=False, # Assume not cutoff unless original text was truncated
              referenced_visuals=[], # No references identified
              referenced_tables=[]  # No references identified
          )
     except ValidationError as ve_fb:
         # This should be rare if the inputs (page_num, raw_text, error_msg) are valid types
         logging.error(f"[create_error_fallback_subsection] Critical: Failed to create even the fallback Subsection for page {page_num}: {ve_fb}", exc_info=True)
         # Re-raise the error as we cannot create a valid object, indicating a deeper issue
         raise ve_fb
     except Exception as e:
        # Catch any other unexpected errors during creation
        logging.error(f"[create_error_fallback_subsection] Unexpected error creating fallback Subsection for page {page_num}: {e}", exc_info=True)
        # Depending on desired behavior, could raise or return a dummy/error object
        # Re-raising is often safer to signal the failure clearly
        raise e


# Prompt for initial structure extraction (raw_text + elements)
# --- New Task-Oriented Prompt for Initial Page Structure Extraction ---

page_structure_prompt = """
--- Goal ---
Analyze the input, which represents a single page ({page_num}) of a document (could be an image or text). Your primary task is to extract the complete raw text content AND identify and structure specific elements like tables, visuals, and key numerical data points found on that page.

--- Input ---
- Content of page {page_num} (provided as image or text).

--- Instructions ---
Perform the following tasks based *only* on the provided input content for page {page_num}:

**Task 1: Extract Raw Text**
1a. Extract ALL text content visible on the page.
1b. Preserve original formatting like paragraphs and line breaks as accurately as possible.
1c. Store this complete text in the `raw_text` field of the output JSON.

**Task 2: Extract Tables**
2a. Identify all distinct tables present on the page.
2b. For EACH table found:
    *   Extract the complete cell content accurately.
    *   Format this content STRICTLY as GitHub Flavored Markdown within the `table_content` field (using '|' separators and a '---' header separator line). Do NOT include surrounding paragraphs or explanatory text in `table_content`.
    *   If a clear title or caption is associated with the table, extract it into the `title` field. Otherwise, use `null`.
    *   Write a brief (1-2 sentence) `summary` of the table's main content or purpose, if possible. Otherwise, use `null`.
    *   Generate a unique `table_id` formatted as `page_{page_num}_table_N` (where N is 1, 2, 3...).
    *   Set `page_number` to {page_num}.
    *   Add the complete table object to the `tables` list in the output JSON.
2c. If NO tables are found on the page, the `tables` list MUST be an empty list `[]`.

**Task 3: Extract Visual Elements**
3a. Identify all visual elements (e.g., charts, graphs, diagrams, images, photographs, equations, code blocks).
3b. For EACH visual element found:
    *   Determine its `type` (e.g., "bar chart", "line graph", "diagram", "image", "equation", "code block").
    *   Write a detailed `description` covering what the visual shows, its purpose, and key elements.
    *   Provide a `data_summary` summarizing key data points, trends, or findings presented in the visual. If the visual doesn't present data (e.g., a decorative image) or a summary isn't applicable, use the string "N/A".
    *   Generate a unique `visual_id` formatted as `page_{page_num}_visual_N` (where N is 1, 2, 3...).
    *   Set `page_numbers` to `[{page_num}]`.
    *   Set `source_url` and `alt_text` to `null` unless explicitly available.
    *   Add the complete visual object to the `visuals` list in the output JSON.
3c. If NO visual elements are found, the `visuals` list MUST be an empty list `[]`.

**Task 4: Extract Numerical Data**
4a. Identify significant numerical data points mentioned in the page's text (excluding those already inside tables).
4b. For EACH significant number found:
    *   Extract the `value` as a string (e.g., "123.45", "50%", "$1.2M").
    *   Write a `description` explaining what the number represents, including units or context (e.g., "increase in revenue", "percentage completion").
    *   Extract the surrounding text snippet (~10-20 words) providing immediate context into the `context` field. If no clear surrounding context is available, use `null`.
    *   Add the complete number object to the `numbers` list in the output JSON.
4c. If NO significant numerical data points are found, the `numbers` list MUST be an empty list `[]`.

**Task 5: Set Boolean Flags**
5a. Set `has_tables` to `true` if the `tables` list is NOT empty, otherwise set it to `false`.
5b. Set `has_visuals` to `true` if the `visuals` list is NOT empty, otherwise set it to `false`.
5c. Set `has_numbers` to `true` if the `numbers` list is NOT empty, otherwise set it to `false`.

**Task 6: Assemble Final JSON**
6a. Combine all extracted information (`raw_text`, `has_tables`, `tables`, `has_visuals`, `visuals`, `has_numbers`, `numbers`) into a single JSON object adhering precisely to the schema defined below.

--- Output Format ---
Return ONLY a single, valid JSON object. Do NOT include any text, comments, or markdown formatting before or after the JSON object. It MUST strictly follow this schema:

```json
{{
  "raw_text": "string", // REQUIRED. Full text from the page.
  "has_tables": "boolean", // REQUIRED. True if 'tables' list is not empty.
  "has_visuals": "boolean", // REQUIRED. True if 'visuals' list is not empty.
  "has_numbers": "boolean", // REQUIRED. True if 'numbers' list is not empty.
  "tables": [ // REQUIRED list (can be empty []).
    {{
      "table_id": "string", // REQUIRED. e.g., "page_{page_num}_table_1"
      "table_content": "string", // REQUIRED. Table content as GitHub Flavored Markdown.
      "title": "string | null", // Optional
      "summary": "string | null", // Optional
      "page_number": "integer" // REQUIRED. e.g., {page_num}
    }}
  ],
  "visuals": [ // REQUIRED list (can be empty []).
    {{
      "visual_id": "string", // REQUIRED. e.g., "page_{page_num}_visual_1"
      "type": "string", // REQUIRED. e.g., "bar chart", "image"
      "description": "string", // REQUIRED. Detailed description.
      "data_summary": "string", // REQUIRED. Summary of data or "N/A".
      "page_numbers": ["integer"], // REQUIRED. e.g., [{page_num}]
      "source_url": "string | null", // Optional
      "alt_text": "string | null" // Optional
    }}
  ],
  "numbers": [ // REQUIRED list (can be empty []).
    {{
      "value": "string", // REQUIRED. Numerical value as string.
      "description": "string", // REQUIRED. What the number represents.
      "context": "string | null" // Optional. Surrounding text.
    }}
  ]
}}

--- Critical Rules ---
VALID JSON ONLY: The entire output MUST be a single JSON object starting with {{ and ending with }}.
SCHEMA ADHERENCE: Follow the schema EXACTLY. All REQUIRED fields must be present with the correct data types. Use null for optional string fields if no value is applicable/found. Use [] for empty lists.
MARKDOWN TABLES: Table content MUST be formatted as GitHub Flavored Markdown in the table_content field.
TEXTUAL BASIS: All extracted information must be derived SOLELY from the provided page content. Do not infer or add external knowledge.

Generate the JSON output for page {page_num} now:
"""


async def extract_structure_from_text_content(client: genai.Client, text_content: str, page_num: int, filename: str) -> PageContent:
    """
    Uses Gemini to extract raw text and structural elements (tables, visuals, numbers)
    from the input TEXT content. Populates PageContent WITHOUT subsections initially.
    Includes fallback for pure text preservation.
    """
    function_name = "extract_structure_from_text_content"
    logging.debug(f"[{function_name}] Starting initial structure extraction from text for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    if not text_content or not text_content.strip():
         logging.warning(f"[{function_name}] Input text_content is empty for page {page_num} of '{filename}'. Returning basic PageContent.")
         page_data["raw_text"] = "" # Set raw text to empty string
         try:
              return PageContent(**page_data)
         except ValidationError as ve:
             logging.error(f"[{function_name}] Failed validation even for empty text page {page_num}: {ve}")
             # Return hardcoded error if basic creation fails
             return create_error_page(page_num, f"Validation error on empty text input: {ve}")


    # Limit text content size
    MAX_TEXT_LENGTH = 100000
    if len(text_content) > MAX_TEXT_LENGTH:
        logging.warning(f"[{function_name}] Truncating text input for page {page_num} of '{filename}' from {len(text_content)} to {MAX_TEXT_LENGTH} chars.")
        text_content_for_api = text_content[:MAX_TEXT_LENGTH] + "\n... [Content Truncated] ..."
    else:
        text_content_for_api = text_content

    try:
        prompt_for_page = page_structure_prompt.format(page_num=page_num)

        full_prompt_with_data = f"""
        {prompt_for_page}
        
        {text_content_for_api}
        """

        # --- Attempt structured extraction (raw_text + elements) from text ---
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=full_prompt_with_data)]),
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {}
            if not isinstance(data, dict): data = {}

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text") # Get raw text from LLM's perspective
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False) # Less likely but check
            page_data["has_numbers"] = data.get("has_numbers", False)
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])
            page_data["subsections"] = [] # Explicitly ensure subsections are empty here

            # If LLM failed to echo raw_text, use the original input text as fallback
            if not page_data["raw_text"]:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num}. Using original input text.")
                 page_data["raw_text"] = text_content # Use original full text

            # Create PageContent object (validates tables, visuals, numbers)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted initial structure from text for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for text structure extraction on page {page_num}.")
            raise ValueError("LLM response had no candidates") # Trigger fallback

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.warning(f"[{function_name}] Initial structure extraction from text failed for page {page_num}: {e}. Falling back to text preservation.")

        # --- Fallback: Preserve original text ---
        # Create a PageContent object with ONLY raw_text and page_number populated
        page_data["raw_text"] = text_content # Ensure original text is saved
        page_data["has_tables"] = False # Reset flags
        page_data["has_visuals"] = False
        page_data["has_numbers"] = False
        page_data["tables"] = []
        page_data["visuals"] = []
        page_data["numbers"] = []
        page_data["subsections"] = []

        try:
             # Return a valid PageContent object even in fallback
             return PageContent(**page_data)
        except ValidationError as ve_fallback:
            logging.error(f"[{function_name}] Failed to create fallback PageContent for text page {page_num}: {ve_fallback}")
            return create_error_page(page_num, f"Text structure extraction failed AND fallback creation failed: {ve_fallback}")


async def extract_page_content_from_memory(client, image_part, page_num, filename=None):
    """
    Extracts raw text and structural elements (tables, visuals, numbers) from a single page image.
    Populates PageContent WITHOUT subsections initially.
    """
    function_name = "extract_page_structure" # For logging
    logging.debug(f"[{function_name}] Starting extraction for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    page_structure_prompt = f"""
    Analyze page {page_num} from the input document/image. Extract the full raw text content and identify structural elements like tables, visuals, and key numbers. Return ONLY a valid JSON object with the exact structure specified below.

    JSON Structure:
    {{
    "raw_text": "Complete text content extracted from the page. Preserve formatting like paragraphs and line breaks.",
    "has_tables": true/false,
    "has_visuals": true/false,
    "has_numbers": true/false,
    "tables": [ {{ "table_id": "page_{page_num}_table_N", "table_content": "markdown...", "title": "optional", "summary": "optional", "page_number": {page_num} }} ],
    "visuals": [ {{ "visual_id": "page_{page_num}_visual_N", "type": "...", "description": "...", "data_summary": "...", "page_numbers": [{page_num}], "source_url": null, "alt_text": null }} ],
    "numbers": [ {{ "value": "...", "description": "...", "context": "..." }} ]
    }}

    Rules & Guidelines:
    1.  **Raw Text:** Extract ALL text content into the `raw_text` field. Maintain original paragraph structure.
    2.  **Elements:** Accurately identify and extract all tables, visuals (charts, images, diagrams, equations), and significant numbers.
    3.  **IDs:** Generate unique IDs for tables and visuals starting N from 1 for the page.
    4.  **Structure:** Adhere STRICTLY to the specified JSON structure. Use `[]` for empty lists.
    5.  **JSON Validity:** Ensure the output is ONLY a single, valid JSON object.

    EXAMPLE JSON Schema (Conceptual):
    {{
    "type": "object",
    "properties": {{
        "raw_text": {{ "type": ["string", "null"] }},
        "has_tables": {{ "type": "boolean" }},
        "has_visuals": {{ "type": "boolean" }},
        "has_numbers": {{ "type": "boolean" }},
        "tables": {{ "type": "array", "items": {{ "$ref": "#/definitions/TableData" }} }},
        "visuals": {{ "type": "array", "items": {{ "$ref": "#/definitions/VisualElement" }} }},
        "numbers": {{ "type": "array", "items": {{ "$ref": "#/definitions/NumericalDataPoint" }} }}
    }},
    "required": ["raw_text", "has_tables", "has_visuals", "has_numbers", "tables", "visuals", "numbers"]
    // Define TableData, VisualElement, NumericalDataPoint schemas conceptually if needed
    }}
    """

    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[image_part, types.Part.from_text(text=page_structure_prompt)]),
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )



        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {} # Handle list response
            if not isinstance(data, dict): data = {} # Ensure data is a dict

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text")
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False)
            page_data["has_numbers"] = data.get("has_numbers", False)
            # Pydantic will validate items during PageContent creation
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])

            # Check if raw_text was extracted
            if page_data["raw_text"] is None:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num} of '{filename}'.")
                 page_data["raw_text"] = "[LLM failed to extract raw text]" # Provide placeholder

            # Create PageContent object (will validate tables, visuals, numbers if models defined)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted structure for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for page {page_num} of '{filename}'.")
            # Return error page content (subsections will be empty)
            return create_error_page(page_num, "LLM response had no candidates.")

    except (json.JSONDecodeError, ValidationError, Exception) as e:
        logging.error(f"[{function_name}] Failed for page {page_num} of '{filename}': {e}", exc_info=True)
        # Return error page content (subsections will be empty)
        error_page = create_error_page(page_num, f"Initial extraction failed: {e}")
        # Ensure raw_text is None or set appropriately in error page if needed downstream
        error_page.raw_text = page_data.get("raw_text") # Preserve if extracted before error
        return error_page


def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List[Dict[str, Any]]] = None) -> PageContent:
    """
    Creates a simplified PageContent object representing an error state,
    placing the error details within a single subsection.
    """
    logging.error(f"Creating error page for page {page_num} due to: {error_msg}")
    # --- Format the detailed error text ---
    error_text = f"Error processing page {page_num}: {error_msg}"

    if validation_errors:
        # Ensure validation_errors is a list of dicts before processing
        if isinstance(validation_errors, list) and all(isinstance(item, dict) for item in validation_errors):
            try:
                error_details = "\n\nValidation Errors:\n" + "\n".join(
                    # Safely access keys with .get()
                    f"- Field '{'.'.join(map(str, err.get('loc', ['unknown']))) if err.get('loc') else 'unknown'}': {err.get('msg', 'No message')}"
                    for err in validation_errors
                )
                error_text += error_details
                logging.debug(f"Formatted validation errors for page {page_num}: {error_details}")
            except Exception as format_error:
                 logging.error(f"Could not format validation errors for page {page_num}: {format_error}")
                 error_text += f"\n(Could not format validation errors: {format_error})"
        else:
            logging.warning(f"Received validation_errors in unexpected format for page {page_num}: {type(validation_errors)}. Skipping details.")
            error_text += "\n(Validation error details format incorrect)"

    # --- Create the error subsection ---
    error_subsection = Subsection(
        subsection_id=f"page_{page_num}_section_error",
        order=1,
        title="Processing Error",
        text=error_text, # Place the detailed error message here
        description=f"An error occurred processing page {page_num}.", # Brief summary
        is_cutoff=False,
        referenced_visuals=[],
        referenced_tables=[]
    )

    # --- Initialize the simplified PageContent object ---
    error_page = PageContent(
        page_number=page_num,
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        subsections=[error_subsection] # Put the error subsection in the list
        # Removed old fields: text, title, topics, summary, entities, dates, etc.
    )

    logging.debug(f"Generated error PageContent object for page {page_num}")
    return error_page

async def merge_cutoff_subsections(pages: List[PageContent]) -> List[PageContent]:
    """
    Merge subsections marked as cut off with the first subsection of the next page.
    Operates on the simplified PageContent model.
    """
    if not pages or len(pages) < 2:
        return pages # Nothing to merge if less than 2 pages

    # Ensure pages are sorted by page number
    # Use getattr for safe access in case of malformed objects, though Pydantic should ensure existence
    try:
        sorted_pages = sorted(pages, key=lambda p: getattr(p, 'page_number', float('inf')))
    except Exception as sort_err:
        logging.error(f"Failed to sort pages for merging: {sort_err}. Returning original order.")
        sorted_pages = pages # Proceed with original order if sorting fails


    logging.debug(f"Attempting to merge subsections across {len(sorted_pages)} pages.")
    merged_something = False # Flag to track if any merge happened

    # Iterate through pages up to the second-to-last one
    for i in range(len(sorted_pages) - 1):
        current_page = sorted_pages[i]
        next_page = sorted_pages[i + 1]

        # Basic validation of pages and subsections lists
        if not isinstance(current_page, PageContent) or not isinstance(next_page, PageContent):
             logging.warning(f"Skipping merge check between page {i+1} and {i+2} due to invalid PageContent object types.")
             continue

        if not current_page.subsections or not next_page.subsections:
            # logging.debug(f"Skipping merge check between page {current_page.page_number} and {next_page.page_number}: one has no subsections.")
            continue # Nothing to merge if either page lacks subsections

        # Sort subsections by order for safety, though they should be ordered
        current_subsections = sorted(current_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))
        next_subsections = sorted(next_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))

        # Get the last subsection of the current page and first of the next
        last_subsection = current_subsections[-1]
        first_next_subsection = next_subsections[0]

        # Check if the last subsection is actually marked as cut off
        if getattr(last_subsection, 'is_cutoff', False):
            logging.info(f"Merging cutoff subsection '{last_subsection.subsection_id}' (Page {current_page.page_number}) with '{first_next_subsection.subsection_id}' (Page {next_page.page_number})")

            # --- Perform the merge ---
            # Append text (Use space or newline as appropriate)
            last_subsection.text += "\n" + getattr(first_next_subsection, 'text', '') # Use getattr for safety

            # Mark the merged subsection as not cut off
            last_subsection.is_cutoff = False

            # Merge references (handle potential missing attributes)
            last_subsection.referenced_visuals.extend(getattr(first_next_subsection, 'referenced_visuals', []))
            last_subsection.referenced_tables.extend(getattr(first_next_subsection, 'referenced_tables', []))
            # Ensure uniqueness if needed:
            # last_subsection.referenced_visuals = list(set(last_subsection.referenced_visuals))
            # last_subsection.referenced_tables = list(set(last_subsection.referenced_tables))


            # --- Update the next page ---
            # Remove the first subsection from the next page
            next_page.subsections = next_subsections[1:]

            # Renumber the order of remaining subsections in the next page
            for idx, subsection in enumerate(next_page.subsections):
                 subsection.order = idx + 1

            merged_something = True # Mark that a merge occurred

    if merged_something:
         logging.info("Finished merging cut-off subsections.")
    else:
         logging.debug("No cut-off subsections found to merge.")

    return sorted_pages # Return the list of pages (potentially modified)

async def extract_chapters_from_subsections(client: genai.Client, pages: List[PageContent]) -> List[Chapter]:
    """
    Groups subsections from multiple pages into logical chapters using an LLM.
    Input `pages` should already have cut-off subsections merged.
    """
    all_subsections_with_context = []
    subsection_map = {} # For quick lookup later

    logging.debug(f"Extracting subsections from {len(pages)} pages for chapter generation.")
    for page in pages:
        if not isinstance(page, PageContent) or not page.subsections:
            continue # Skip pages without valid subsections
        for subsection in page.subsections:
             if isinstance(subsection, Subsection):
                # Store the subsection object itself and context for the prompt
                context = {
                    "subsection_id": subsection.subsection_id,
                    "title": subsection.title,
                    "description": subsection.description, # Use the description field
                    "page_number": page.page_number, # Get page number from parent
                    # Optional: Add a very short text preview if description isn't enough context
                    # "text_preview": subsection.text[:50] + "..." if len(subsection.text) > 50 else subsection.text
                }
                all_subsections_with_context.append(context)
                subsection_map[subsection.subsection_id] = subsection # Store the actual object
             else:
                  logging.warning(f"Skipping invalid subsection object on page {page.page_number}")


    if not all_subsections_with_context:
        logging.warning("No valid subsections found across all pages to form chapters.")
        # Return an empty list or a single default chapter if required by downstream logic
        # For now, returning empty list. Adjust if a default chapter is mandatory.
        return []

    # Sort subsections by page number and order (using the context dict)
    # This ensures the LLM sees them in the correct document order.
    sorted_subsections_context = sorted(all_subsections_with_context, key=lambda s: (s["page_number"], s.get("order", 0))) # Use get for order safety

    logging.info(f"Sending {len(sorted_subsections_context)} subsections to LLM for chapter structuring.")

    # Create prompt for chapter extraction using NEW structure
    # Pass subsection_id, title, and description
    chapter_prompt = f"""
    Analyze the following list of subsections extracted sequentially from a document. Each subsection includes its ID, title, and a brief description. Your goal is to group these subsections into logically coherent chapters.

    Subsections List:
    ```json
    {json.dumps(sorted_subsections_context, indent=2)}
    ```

    Instructions:
    1. Group the provided subsections into 3 to 10 logical chapters.
    2. Base the grouping on topic cohesion and narrative flow. Subsections within a chapter should relate to a common theme or part of the document's structure.
    3. Each chapter MUST contain at least one subsection.
    4. Ensure every subsection ID from the input list is assigned to exactly ONE chapter.
    5. Maintain the original relative order of subsections within and across chapters as much as possible.
    6. For each chapter created, provide:
        - A unique `chapter_id` (e.g., "chapter_1", "chapter_2").
        - A concise and descriptive `title` (max 7 words) reflecting the chapter's main theme.
        - A brief `summary` (1-2 sentences) describing the chapter's content or purpose.
        - A list of `subsection_ids` belonging to that chapter, in their original relative order.
    7. Number chapters sequentially starting from 1 using an `order` field.

    Return ONLY a valid JSON array containing the chapter objects, structured exactly like this example:
    ```json
    [
      {{
        "chapter_id": "chapter_1",
        "order": 1,
        "title": "Introduction and Setup",
        "summary": "Provides background information and initial setup instructions.",
        "subsection_ids": ["page_1_section_1", "page_1_section_2", "page_2_section_1"]
      }},
      {{
        "chapter_id": "chapter_2",
        "order": 2,
        "title": "Core Functionality Explained",
        "summary": "Details the main features and how they operate.",
        "subsection_ids": ["page_2_section_2", "page_3_section_1", "page_3_section_2"]
      }}
      // ... more chapters
    ]
    ```
    Ensure the final output is only the JSON array, with no extra text before or after.
    """

    chapters = []
    try:
        logging.debug("Sending chapter extraction request to LLM.")
        # Use the actual client and API call structure
        response = await retry_api_call(
            client.aio.models.generate_content, # Adjusted path might be needed based on client setup
            model="gemini-2.0-flash", # Or pro, flash might struggle with complex structuring
            contents=[
                types.Content(parts=[types.Part.from_text(text=chapter_prompt)]),
            ], # Prompt first, then text content
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"LLM response for chapters: {json_text[:500]}...") # Log beginning of response
            chapters_data = json.loads(json_text) # Expecting a List[Dict]

            if not isinstance(chapters_data, list):
                 logging.warning(f"LLM returned non-list data for chapters: {type(chapters_data)}. Attempting recovery if dict.")
                 # Handle cases where it might incorrectly return a dict containing the list
                 if isinstance(chapters_data, dict) and "chapters" in chapters_data and isinstance(chapters_data["chapters"], list):
                     chapters_data = chapters_data["chapters"]
                 else:
                     raise ValueError("LLM response for chapters was not a list.")


            processed_subsection_ids = set()
            for i, chapter_data in enumerate(chapters_data):
                # Validate chapter_data structure
                if not isinstance(chapter_data, dict) or not all(k in chapter_data for k in ["title", "summary", "subsection_ids", "order"]):
                     logging.warning(f"Skipping malformed chapter data received from LLM: {chapter_data}")
                     continue

                chapter_subsections = []
                subsection_ids_in_chapter = chapter_data.get("subsection_ids", [])
                if not isinstance(subsection_ids_in_chapter, list):
                     logging.warning(f"Subsection IDs for chapter '{chapter_data.get('title')}' is not a list. Skipping chapter.")
                     continue

                for subsection_id in subsection_ids_in_chapter:
                    if subsection_id in subsection_map:
                        chapter_subsections.append(subsection_map[subsection_id])
                        processed_subsection_ids.add(subsection_id)
                    else:
                        logging.warning(f"Subsection ID '{subsection_id}' from LLM not found in original map for chapter '{chapter_data.get('title')}'.")

                if not chapter_subsections:
                     logging.warning(f"Chapter '{chapter_data.get('title')}' created by LLM has no valid subsections. Skipping.")
                     continue

                # Create the Chapter object (ensure your Chapter model matches)
                chapter = Chapter(
                    chapter_id=chapter_data.get("chapter_id", f"chapter_{chapter_data.get('order', i+1)}"), # Use order for default ID
                    title=chapter_data.get("title", f"Chapter {chapter_data.get('order', i+1)}"),
                    summary=chapter_data.get("summary", "No summary provided."),
                    subsections=chapter_subsections, # List of actual Subsection objects
                    # entity_relationships=[], # Defer relationship extraction
                    order=chapter_data.get("order", i+1)
                )
                chapters.append(chapter)

            # Check if all original subsections were processed
            original_ids = set(subsection_map.keys())
            missing_ids = original_ids - processed_subsection_ids
            if missing_ids:
                 logging.warning(f"LLM failed to assign {len(missing_ids)} subsections to chapters: {missing_ids}. They will be excluded.")
                 # Optionally, create an "Orphaned" chapter for these

            # Sort final chapters list by order
            chapters = sorted(chapters, key=lambda c: c.order)

            if not chapters:
                 logging.warning("LLM processing resulted in zero valid chapters. Creating default chapter.")
                 raise ValueError("LLM returned no valid chapters") # Trigger default chapter creation in except block

            logging.info(f"Successfully extracted {len(chapters)} chapters from LLM.")
            return chapters

        else:
             logging.error("Chapter extraction response from LLM had no candidates.")
             raise ValueError("LLM response had no candidates") # Trigger default

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.error(f"Error processing LLM response or extracting chapters: {e}", exc_info=True)

        # --- Fallback: Create a single default chapter ---
        logging.info("Creating a single default chapter due to chapter extraction error.")
        # Get all original subsections back from the map, sorted
        all_original_subsections = sorted(subsection_map.values(), key=lambda s: (
            # Need page number back for sorting - requires passing it or storing it differently
            # Assuming subsection_map values have access or we modify the storage
            # For now, sort by ID as a fallback sort order
            getattr(s, 'subsection_id', '')
        ))

        default_chapter = Chapter(
            chapter_id="chapter_1_default",
            title="Document Content",
            summary="Content grouped into a single chapter due to an error during automatic chapter structuring.",
            subsections=all_original_subsections, # Assign all subsections here
            # entity_relationships=[],
            order=1
        )
        return [default_chapter]

async def analyze_concept_entity_relationships(
    client: genai.Client, # Using a more specific type hint if possible
    chapter: Chapter,
    source_identifier: str,
    max_structure_retries: int = 2, # Total attempts for the *initial* prompt
    retry_delay: int = 3
) -> None:
    """
    Analyzes text within a single chapter to extract key concepts/entities
    and their relationships using an LLM structured prompt. Populates
    chapter.qdrant_nodes in-place. Includes retries on structural errors,
    AND a fallback mechanism using extracted text from failed responses.

    Args:
        client: The configured GenAI client/model instance.
        chapter: The Chapter object containing subsections to analyze.
        source_identifier: Identifier for the source document.
        max_structure_retries: How many times to retry the INITIAL prompt if the LLM returns the wrong JSON structure.
        retry_delay: Base delay (in seconds) between structure retries.
    """
    function_name = "analyze_concept_entity_relationships_fallback" # Renamed
    chapter_id = getattr(chapter, 'chapter_id', 'unknown_chapter')
    chapter_title = getattr(chapter, 'title', 'Untitled Chapter')
    logging.info(f"[{function_name}] Analyzing relationships for Chapter '{chapter_id}': '{chapter_title}' from source '{source_identifier}'")

    # --- 1. Prepare Input Text ---
    if not chapter.subsections:
        logging.warning(f"[{function_name}] Chapter '{chapter_id}' has no subsections. Skipping analysis.")
        chapter.qdrant_nodes = []
        return

    combined_text = ""
    try:
        # --- Start Expansion: combined_text generation ---
        subsection_texts = []
        for subsection in chapter.subsections:
             # Use getattr for safe access to attributes
             sub_id = getattr(subsection, 'subsection_id', 'unknown_id')
             sub_title = getattr(subsection, 'title', 'Untitled Subsection')
             sub_desc = getattr(subsection, 'description', '') # Use the description for context
             sub_text = getattr(subsection, 'text', '')

             # Only include non-empty text
             if sub_text and sub_text.strip():
                 # Format clearly for the LLM
                 subsection_texts.append(
                     f"--- Subsection Start (ID: {sub_id}) ---\n"
                     f"Title: {sub_title}\n"
                     f"Description: {sub_desc}\n\n" # Description provides LLM context
                     f"{sub_text.strip()}\n" # Add the actual text
                     f"--- Subsection End ---"
                 )
             else:
                 logging.debug(f"[{function_name}] Skipping empty subsection {sub_id} in chapter {chapter_id}.")

        combined_text = "\n\n".join(subsection_texts)
        # --- End Expansion: combined_text generation ---

        MAX_TEXT_LIMIT = 150000 # Adjust based on model context window limits
        if len(combined_text) > MAX_TEXT_LIMIT:
             logging.warning(f"[{function_name}] Truncating combined text for chapter '{chapter_id}' from {len(combined_text)} to {MAX_TEXT_LIMIT} chars.")
             combined_text = combined_text[:MAX_TEXT_LIMIT] + "\n... [Content Truncated] ..."
        elif not combined_text.strip():
            logging.warning(f"[{function_name}] Combined text for chapter '{chapter_id}' is empty after processing subsections. Skipping analysis.")
            chapter.qdrant_nodes = []
            return
    except Exception as text_prep_error:
        logging.error(f"[{function_name}] Error preparing text for chapter '{chapter_id}': {text_prep_error}", exc_info=True)
        chapter.qdrant_nodes = []
        return

    # --- 2. Extract Chapter Number ---
    chapter_num_str = "1" # Default
    if chapter_id and isinstance(chapter_id, str):
        match = re.search(r'[_-](\d+)$', chapter_id)
        if match:
            chapter_num_str = match.group(1)

    # --- 3. Define the *ORIGINAL* LLM Prompt ---

    # --- Start Expansion: new_json_schema_description ---
    new_json_schema_description = """
```json
{
  "qdrant_nodes": [
    {
      "node_id": "string", // REQUIRED. Format: "{type}_{sanitized_name}_ch{chapter_num_str}". Example: "concept_machine_learning_ch1". Sanitized name: lowercase, spaces/symbols to underscores, max 50 chars.
      "node_type": "string", // REQUIRED. "entity" or "concept".
      "name": "string", // REQUIRED. Original human-readable name of the entity/concept.
      "chapter_id_context": "string", // REQUIRED. Use the provided Chapter ID. Example: "{chapter_id}"
      "properties": { // REQUIRED Dictionary.
        "title": "string", // REQUIRED. A concise title for this specific node's content within the chapter. Use headings from text if applicable, otherwise generate a descriptive title (5-10 words).
        // "source": "string", // NOTE: This will be added by the Python code later based on 'source_identifier'. Do NOT generate it here.
        "source_credibility": "string", // REQUIRED. Infer credibility based on source type/context. Examples: "Peer-Reviewed", "Company Report", "News Article", "Blog Post", "Technical Documentation", "User Manual", "Internal Memo", "Unknown".
        "information_type": "string", // REQUIRED. Classify the main type of information represented by this node's content. Examples: "Definition", "Explanation", "Data Point", "Example", "Process Description", "Case Study", "Methodology", "API Specification", "Configuration", "Finding", "Claim", "Opinion", "Requirement".
        "core_claim": "string", // REQUIRED. A 1-3 sentence summary capturing the main point, assertion, or function of this node within the chapter's text. This replaces the old 'description'. Be specific.
        "key_entities": "string", // REQUIRED. A JSON string representing a list of key entity names (people, orgs, products, locations) explicitly mentioned in the text segment MOST relevant to THIS node. Example: "[\\"Acme Corp\\", \\"Project X\\"]". Generate an empty list "[]" if none apply.
        "sentiment_tone": "string", // REQUIRED. Overall sentiment/tone of the text related to this node. Examples: "Positive", "Neutral", "Negative", "Mixed", "Cautionary", "Objective", "Promotional".
        "tags": "string", // REQUIRED. A JSON string representing a list of relevant keywords or tags (1-5 tags). Example: "[\\"AI\\", \\"optimization\\", \\"gpu\\"]". Generate an empty list "[]" if none apply.
        // "stored_at": "string" // NOTE: This will be added by the Python code later. Do NOT generate it here.
      },
      "linked_nodes": [ // List, REQUIRED (can be empty [] if no links found). Max 7 links per node.
        {
          "target_node_id": "string", // REQUIRED. node_id of another node defined in THIS response's "qdrant_nodes" list.
          "relationship_description": "string", // REQUIRED. Specific description of the link. (Use verbs, avoid generics like 'related to'). Example: 'calculates cost based on', 'developed by', 'depends on', 'provides data for'.
          "relationship_keywords": ["string"], // List, OPTIONAL. 1-3 keywords summarizing the relationship. Defaults to empty list []. Example: ["calculation", "dependency"].
          "relationship_strength": "float" // REQUIRED. Link strength/importance (1.0-10.0). Estimate based on how central the connection is.
        }
      ]
    }
  ]
}
```"""
    # --- End Expansion: new_json_schema_description ---
    initial_entity_prompt = f"""
        --- Goal ---
        Analyze the text from the single document chapter provided below. Your task is to identify key specific entities (e.g., organizations, people, products) AND important abstract concepts (e.g., processes, topics, ideas). Extract these as nodes. For each node, determine its intrinsic properties based *only* on the chapter text. Also, identify direct relationships *between these nodes within this chapter*.

        --- Input Data ---
        Chapter Title: {chapter_title}
        Chapter ID: {chapter_id}
        Source Identifier (Filename/URL): {source_identifier}
        Chapter Content:
        {combined_text}
        --- End Input Data ---

        --- Instructions ---
        Perform the following tasks based *only* on the provided "Chapter Content":

        {'''
        **Task 1: Extract Nodes, Properties, and Intra-Chapter Relationships**
        1a. **Identify Nodes:** Scan the text and identify all distinct key entities and abstract concepts. An entity is typically a concrete noun (person, place, organization, product). A concept is an abstract idea, process, topic, or technique. Aim for ~5-15 significant nodes per chapter, focusing on items central to the chapter's narrative.
        1b. **Define Node Structure:** For each node identified, structure its information as follows:
            *   `node_id`: Generate using the format `{{type}}_{{sanitized_name}}_ch{{chapter_num_str}}`. '{{type}}' is 'entity' or 'concept'. '{{sanitized_name}}' is the lowercase node name with spaces/symbols replaced by underscores `_`, truncated to ~50 characters. Ensure uniqueness within this chapter's output.
            *   `node_type`: Must be exactly "entity" or "concept".
            *   `name`: The original, human-readable name as found in the text or a standard representation (e.g., "Machine Learning" not "machine learning").
            *   `chapter_id_context`: Set to exactly "{chapter_id}".
            *   `properties`: A nested JSON object containing the following fields derived from the text segment most relevant to *this specific node*:
                *   `title`: Generate a concise title (5-10 words) capturing the essence of this node's discussion within the chapter. If a clear heading exists for the relevant section, adapt it. Otherwise, create a descriptive title. Example: "Function of the RAG System Retriever".
                *   `source_credibility`: Infer the credibility level based on the source identifier '{source_identifier}' and the text content/tone. Choose from: "Peer-Reviewed", "Company Report", "News Article", "Blog Post", "Technical Documentation", "User Manual", "Internal Memo", "Marketing Material", "Personal Communication", "Unknown".
                *   `information_type`: Classify the primary type of information this node represents within its context. Choose from: "Definition", "Explanation", "Description", "Example", "Comparison", "Contrast", "Data Point", "Metric", "Process Step", "Workflow", "Case Study", "Methodology", "Algorithm", "API Specification", "Configuration Detail", "Requirement", "Constraint", "Benefit", "Limitation", "Finding", "Claim", "Hypothesis", "Opinion", "Prediction", "Guideline", "Best Practice".
                *   `core_claim`: Write a 1-3 sentence summary capturing the central point, assertion, or function of this node *in this chapter's context*. Be specific and action-oriented where possible. Avoid generic statements. Example: "The retrieval module (R) searches external knowledge sources to find relevant context for the generation module." NOT "This node describes the retrieval module."
                *   `key_entities`: Provide a JSON stringified list of specific entity names (people, organizations, product names, locations, specific systems mentioned by name) explicitly mentioned in the immediate context of this node's primary discussion. Example: "["NVIDIA GPUs", "TensorFlow Framework"]". If none are mentioned in the direct context, use "[]".
                *   `sentiment_tone`: Assess the overall sentiment or tone of the text discussing this node. Choose from: "Positive", "Neutral", "Negative", "Mixed", "Objective", "Subjective", "Promotional", "Cautionary", "Speculative", "Instructional", "Analytical".
                *   `tags`: Provide a JSON stringified list of 1-5 relevant keywords or tags that categorize this node. Use lowercase and underscores. Example: "["rag_framework", "llm_enhancement", "vector_database"]". If no specific tags apply, use "[]".
            *   `linked_nodes`: An initially empty list `[]` where relationship objects will be added (see step 1c).

        1c. **Identify Relationships (Linked Nodes):** For each node created (the *source* node), scan the text again to find its most important (max 5-7) direct relationships *to other nodes also created in this Task* (the *target* nodes). Relationships should represent actions, dependencies, comparisons, compositions, etc.
            *   For each relationship found, create a `linked_nodes` object containing:
                *   `target_node_id`: The `node_id` (generated in step 1b) of the target node. This target node MUST be one of the other nodes defined in this response's `qdrant_nodes` list. Do NOT link to nodes outside this list.
                *   `relationship_description`: **CRITICAL: Describe the SPECIFIC ACTION or NATURE of the relationship using precise verbs.** Avoid generic terms like 'related to', 'associated with', 'discusses'. Examples: "implements algorithm", "evaluates performance of", "optimizes parameters for", "is a component of", "contrasts with", "cites findings from".
                *   `relationship_keywords`: (Optional) 1-3 lowercase keywords summarizing the relationship type (e.g., ["dependency", "causation", "comparison", "optimization"]). Default to `[]` if none clearly apply.
                *   `relationship_strength`: Estimate the importance or centrality of this link within the chapter's narrative (float 1.0-10.0). A core dependency might be 9.0, a brief mention or comparison might be 3.0.
            *   Populate the `linked_nodes` list of the source node with these relationship objects. If a node has no direct links to *other nodes in this list*, its `linked_nodes` list remains empty `[]`.

        1d. **Assemble Output:** Format all extracted node objects (including their `properties` and `linked_nodes`) into a single JSON object according to the schema specified below. Ensure the entire output is one valid JSON structure.
        '''}
        --- End Instructions ---

        --- Output Format ---
        Return ONLY a single, valid JSON object. It MUST contain exactly one top-level key: "qdrant_nodes". The value associated with "qdrant_nodes" MUST be a JSON array of node objects. Adhere STRICTLY to this schema:

        {new_json_schema_description}
        --- End Output Format ---

        --- Critical Rules ---
        *   **VALID JSON ONLY:** Entire output must be a single JSON object starting with `{{` and ending with `}}`. No introductory text, explanations, apologies, or markdown formatting outside the JSON.
        *   **SCHEMA ADHERENCE:** Follow the provided schema precisely. All REQUIRED fields/objects must be present (e.g., `properties` dict must exist, `linked_nodes` list must exist even if empty `[]`). Check data types carefully (strings, floats, lists).
        *   **PROPERTIES:** All fields within the `properties` object are REQUIRED, except `source` and `stored_at` which MUST NOT be generated here (they are added later by code). Ensure default/placeholder values are provided if info isn't in text (e.g., "Unknown" for credibility).
        *   **JSON STRINGS:** Ensure `key_entities` and `tags` values within `properties` are valid JSON strings representing lists (e.g., `"[]"`, `"["item1", "item2"]"`). Double-check escaping.
        *   **NODE ID CONSISTENCY:** Every `target_node_id` in any `linked_nodes` list MUST correspond to a `node_id` defined earlier in the main `qdrant_nodes` list within THIS SAME RESPONSE. Do not hallucinate target IDs or link outside the generated list.
        *   **RELATIONSHIP SPECIFICITY:** Relationship descriptions MUST be specific and verb-driven.
        *   **RELATIONSHIP LIMIT:** Max 5-7 `linked_nodes` objects per source node. Focus on the most direct and important links mentioned in the text.
        *   **TEXTUAL BASIS:** All information (`properties`, `relationship_description`) must be derived SOLELY from the provided "Chapter Content". Do not infer external knowledge or details not present in the text.
        *   **UNICODE:** Ensure correct handling of Unicode characters within the JSON output.

        --- End Critical Rules ---

        Generate the JSON output now:
        """

    # --- 4. Call LLM and Process Response (Modified with Fallback) ---

    chapter.qdrant_nodes = [] # Initialize
    final_data = None
    last_raw_response_text = "N/A" # Store last raw response for logging/fallback
    last_cleaned_response_text = "N/A" # Store last cleaned for logging/fallback
    last_parsed_data = None # Store last parsed data for fallback check

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Initial Prompt Attempt {attempt + 1}/{max_structure_retries} for chapter '{chapter_id}'")
        response = None # Reset response for this attempt
        json_text = ""
        data = None

        try:
            response = await retry_api_call(
                client.aio.models.generate_content, # Correct method assuming 'client' is the model instance
                model="gemini-2.5-flash-preview-04-17", # Use model name
                contents=[types.Content(parts=[types.Part.from_text(text=initial_entity_prompt)])], # Pass prompt directly as content
                config=types.GenerateContentConfig(response_mime_type="application/json") # Use GenerateContentConfig
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates (Initial Attempt {attempt + 1}) for chapter '{chapter_id}'.")
                last_raw_response_text = "No candidates received"
                last_cleaned_response_text = "{}"
                last_parsed_data = {}
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: break

            last_raw_response_text = response.candidates[0].content.parts[0].text
            json_text = clean_json_response(last_raw_response_text)
            last_cleaned_response_text = json_text
            logging.debug(f"[{function_name}] Cleaned LLM response (Initial Attempt {attempt + 1}, Chapter {chapter_id}): {json_text[:500]}...")

            try:
                data = json.loads(json_text)
                last_parsed_data = data
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError initial attempt {attempt + 1} for chapter '{chapter_id}': {json_err}.")
                # Attempt to extract text even from decode error for fallback
                if isinstance(json_text, str) and len(json_text.strip()) > 50:
                     last_parsed_data = {"error": "JSONDecodeError", "text": json_text}
                else:
                     last_parsed_data = {"error": "JSONDecodeError", "text": ""}
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: break

            # --- Structure Check ---
            if isinstance(data, dict) and "qdrant_nodes" in data and isinstance(data.get("qdrant_nodes"), list):
                logging.info(f"[{function_name}] Received correct structure initial attempt {attempt + 1} for chapter '{chapter_id}'.")
                final_data = data
                break
            else:
                logging.warning(f"[{function_name}] Incorrect JSON structure initial attempt {attempt+1} for chapter '{chapter_id}'. Got: {type(data)}. Structure: {str(data)[:300]}...")
                # Log the failing structure for debugging
                logging.debug(f"[{function_name}] Failing structure details (Attempt {attempt+1}, Chapter {chapter_id}): Type={type(data)}, Content='{str(data)[:1000]}...'")
                if attempt < max_structure_retries - 1:
                    await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else:
                    logging.info(f"[{function_name}] Initial prompt failed structure check after {max_structure_retries} attempts for chapter '{chapter_id}'. Proceeding to fallback attempt.")
                    break

        except Exception as e:
            logging.error(f"[{function_name}] Error during initial attempt {attempt + 1} for chapter '{chapter_id}': {e}", exc_info=True)
            last_raw_response_text = f"Exception: {e}"
            last_cleaned_response_text = "{}"
            last_parsed_data = {"error": str(e)}
            if attempt < max_structure_retries - 1:
                 await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else:
                 logging.error(f"[{function_name}] Initial attempts failed completely after {max_structure_retries} tries for chapter '{chapter_id}'. Proceeding to fallback attempt.")
                 break

    # --- 5. Fallback Attempt (if final_data is still None after initial retries) ---
    if final_data is None:
        logging.info(f"[{function_name}] Entering fallback: Attempting extraction using text from last failed response for chapter '{chapter_id}'.")

        extracted_text_for_fallback = None
        # Check if the last parsed data (from failed attempt) contains usable text
        if isinstance(last_parsed_data, dict) and "text" in last_parsed_data and isinstance(last_parsed_data["text"], str) and len(last_parsed_data["text"].strip()) > 50: # Basic check for usable text
            extracted_text_for_fallback = last_parsed_data["text"].strip()
            logging.debug(f"[{function_name}] Found extracted text for fallback (length {len(extracted_text_for_fallback)}) for chapter '{chapter_id}'.")
        else:
            logging.warning(f"[{function_name}] Could not find usable extracted text in the last failed response for chapter '{chapter_id}'. Cannot perform text-based fallback.")
            logging.debug(f"[{function_name}] Last parsed data type for fallback check: {type(last_parsed_data)}, keys: {list(last_parsed_data.keys()) if isinstance(last_parsed_data, dict) else 'N/A'}")


        if extracted_text_for_fallback:
            # --- Construct Fallback Prompt ---
            # Uses the same instructions, schema, and rules, but different input text
            fallback_entity_prompt = f"""
--- Goal ---
Analyze the provided text, which was extracted from a previous failed attempt to structure chapter content. Your task is to identify key specific entities AND important abstract concepts from THIS EXTRACTED TEXT ONLY. Extract these as nodes with properties and relationships.

--- Input Data ---
Chapter Title: {chapter_title} (Context Only)
Chapter ID: {chapter_id} (Context Only)
Source Identifier (Filename/URL): {source_identifier} (Context Only)
Extracted Chapter Content (Potentially Incomplete):
{extracted_text_for_fallback}
--- End Input Data ---

--- Instructions ---
Perform the following tasks based *only* on the provided "Extracted Chapter Content":
{'''
**Task 1: Extract Nodes, Properties, and Intra-Chapter Relationships**
1a. **Identify Nodes:** Scan the text and identify all distinct key entities and abstract concepts. An entity is typically a concrete noun (person, place, organization, product). A concept is an abstract idea, process, topic, or technique. Aim for ~5-15 significant nodes per chapter, focusing on items central to the chapter's narrative.
1b. **Define Node Structure:** For each node identified, structure its information as follows:
    *   `node_id`: Generate using the format `{type}_{sanitized_name}_ch{chapter_num_str}`. '{type}' is 'entity' or 'concept'. '{sanitized_name}' is the lowercase node name with spaces/symbols replaced by underscores `_`, truncated to ~50 characters. Ensure uniqueness within this chapter's output.
    *   `node_type`: Must be exactly "entity" or "concept".
    *   `name`: The original, human-readable name as found in the text or a standard representation (e.g., "Machine Learning" not "machine learning").
    *   `chapter_id_context`: Set to exactly "{chapter_id}".
    *   `properties`: A nested JSON object containing the following fields derived from the text segment most relevant to *this specific node*:
        *   `title`: Generate a concise title (5-10 words) capturing the essence of this node's discussion within the chapter. If a clear heading exists for the relevant section, adapt it. Otherwise, create a descriptive title. Example: "Function of the RAG System Retriever".
        *   `source_credibility`: Infer the credibility level based on the source identifier '{source_identifier}' and the text content/tone. Choose from: "Peer-Reviewed", "Company Report", "News Article", "Blog Post", "Technical Documentation", "User Manual", "Internal Memo", "Marketing Material", "Personal Communication", "Unknown".
        *   `information_type`: Classify the primary type of information this node represents within its context. Choose from: "Definition", "Explanation", "Description", "Example", "Comparison", "Contrast", "Data Point", "Metric", "Process Step", "Workflow", "Case Study", "Methodology", "Algorithm", "API Specification", "Configuration Detail", "Requirement", "Constraint", "Benefit", "Limitation", "Finding", "Claim", "Hypothesis", "Opinion", "Prediction", "Guideline", "Best Practice".
        *   `core_claim`: Write a 1-3 sentence summary capturing the central point, assertion, or function of this node *in this chapter's context*. Be specific and action-oriented where possible. Avoid generic statements. Example: "The retrieval module (R) searches external knowledge sources to find relevant context for the generation module." NOT "This node describes the retrieval module."
        *   `key_entities`: Provide a JSON stringified list of specific entity names (people, organizations, product names, locations, specific systems mentioned by name) explicitly mentioned in the immediate context of this node's primary discussion. Example: "["NVIDIA GPUs", "TensorFlow Framework"]". If none are mentioned in the direct context, use "[]".
        *   `sentiment_tone`: Assess the overall sentiment or tone of the text discussing this node. Choose from: "Positive", "Neutral", "Negative", "Mixed", "Objective", "Subjective", "Promotional", "Cautionary", "Speculative", "Instructional", "Analytical".
        *   `tags`: Provide a JSON stringified list of 1-5 relevant keywords or tags that categorize this node. Use lowercase and underscores. Example: "["rag_framework", "llm_enhancement", "vector_database"]". If no specific tags apply, use "[]".
    *   `linked_nodes`: An initially empty list `[]` where relationship objects will be added (see step 1c).

1c. **Identify Relationships (Linked Nodes):** For each node created (the *source* node), scan the text again to find its most important (max 5-7) direct relationships *to other nodes also created in this Task* (the *target* nodes). Relationships should represent actions, dependencies, comparisons, compositions, etc.
    *   For each relationship found, create a `linked_nodes` object containing:
        *   `target_node_id`: The `node_id` (generated in step 1b) of the target node. This target node MUST be one of the other nodes defined in this response's `qdrant_nodes` list. Do NOT link to nodes outside this list.
        *   `relationship_description`: **CRITICAL: Describe the SPECIFIC ACTION or NATURE of the relationship using precise verbs.** Avoid generic terms like 'related to', 'associated with', 'discusses'. Examples: "implements algorithm", "evaluates performance of", "optimizes parameters for", "is a component of", "contrasts with", "cites findings from".
        *   `relationship_keywords`: (Optional) 1-3 lowercase keywords summarizing the relationship type (e.g., ["dependency", "causation", "comparison", "optimization"]). Default to `[]` if none clearly apply.
        *   `relationship_strength`: Estimate the importance or centrality of this link within the chapter's narrative (float 1.0-10.0). A core dependency might be 9.0, a brief mention or comparison might be 3.0.
    *   Populate the `linked_nodes` list of the source node with these relationship objects. If a node has no direct links to *other nodes in this list*, its `linked_nodes` list remains empty `[]`.

1d. **Assemble Output:** Format all extracted node objects (including their `properties` and `linked_nodes`) into a single JSON object according to the schema specified below. Ensure the entire output is one valid JSON structure.
'''}
--- End Instructions ---

--- Output Format ---
Return ONLY a single, valid JSON object. It MUST contain exactly one top-level key: "qdrant_nodes". Adhere STRICTLY to this schema:

{new_json_schema_description}
--- End Output Format ---

--- Critical Rules ---
{'''
*   **VALID JSON ONLY:** Entire output must be a single JSON object starting with `{` and ending with `}`. No introductory text, explanations, apologies, or markdown formatting outside the JSON.
*   **SCHEMA ADHERENCE:** Follow the provided schema precisely. All REQUIRED fields/objects must be present (e.g., `properties` dict must exist, `linked_nodes` list must exist even if empty `[]`). Check data types carefully (strings, floats, lists).
*   **PROPERTIES:** All fields within the `properties` object are REQUIRED, except `source` and `stored_at` which MUST NOT be generated here (they are added later by code). Ensure default/placeholder values are provided if info isn't in text (e.g., "Unknown" for credibility).
*   **JSON STRINGS:** Ensure `key_entities` and `tags` values within `properties` are valid JSON strings representing lists (e.g., `"[]"`, `"["item1", "item2"]"`). Double-check escaping.
*   **NODE ID CONSISTENCY:** Every `target_node_id` in any `linked_nodes` list MUST correspond to a `node_id` defined earlier in the main `qdrant_nodes` list within THIS SAME RESPONSE. Do not hallucinate target IDs or link outside the generated list.
*   **RELATIONSHIP SPECIFICITY:** Relationship descriptions MUST be specific and verb-driven.
*   **RELATIONSHIP LIMIT:** Max 5-7 `linked_nodes` objects per source node. Focus on the most direct and important links mentioned in the text.
*   **TEXTUAL BASIS:** All information (`properties`, `relationship_description`) must be derived SOLELY from the provided "Extracted Chapter Content". Do not infer external knowledge or details not present in the text.
*   **UNICODE:** Ensure correct handling of Unicode characters within the JSON output.
'''}
--- End Critical Rules ---

Generate the JSON output now based *only* on the 'Extracted Chapter Content':
"""
            # --- Execute Fallback LLM Call ---
            logging.debug(f"[{function_name}] Making fallback LLM call for chapter '{chapter_id}'.")
            fallback_response = None
            fallback_json_text = ""
            fallback_data = None
            try:
                fallback_response = await retry_api_call(
                    client.aio.models.generate_content, # Use retry_api_call for robustness
                    model="gemini-2.0-flash", # Use Gemini 2.0 Flash for speed
                    contents=[types.Content(parts=[types.Part.from_text(text=fallback_entity_prompt)])],
                    config=types.GenerateContentConfig(response_mime_type="application/json"),
                )

                if fallback_response and fallback_response.candidates:
                    fallback_json_text = clean_json_response(fallback_response.candidates[0].content.parts[0].text)
                    logging.debug(f"[{function_name}] Cleaned Fallback Response (Chapter {chapter_id}): {fallback_json_text[:500]}...")
                    try:
                        fallback_data = json.loads(fallback_json_text)
                        # --- Final Structure Check (Fallback) ---
                        if isinstance(fallback_data, dict) and "qdrant_nodes" in fallback_data and isinstance(fallback_data.get("qdrant_nodes"), list):
                            logging.info(f"[{function_name}] Fallback attempt SUCCEEDED with correct structure for chapter '{chapter_id}'.")
                            final_data = fallback_data # Assign data from successful fallback
                        else:
                            logging.error(f"[{function_name}] Fallback attempt FAILED structure check for chapter '{chapter_id}'. Got type: {type(fallback_data)}")
                            # Optionally log the failing structure: logging.debug(f"Fallback failing structure: {str(fallback_data)[:1000]}")
                    except json.JSONDecodeError as fallback_json_err:
                        logging.error(f"[{function_name}] Fallback JSONDecodeError for chapter '{chapter_id}': {fallback_json_err}")
                else:
                    logging.error(f"[{function_name}] Fallback LLM call returned no candidates for chapter '{chapter_id}'.")

            except Exception as fallback_err:
                logging.error(f"[{function_name}] Error during fallback LLM call for chapter '{chapter_id}': {fallback_err}", exc_info=True)
        # End of fallback execution block (if extracted_text_for_fallback was found)

    # --- 6. Process Final Data (if successful from initial or fallback) ---
    if final_data is not None:
        qdrant_nodes_data = final_data.get("qdrant_nodes")
        if isinstance(qdrant_nodes_data, list):
            # --- 6a. Validate Individual Nodes ---
            # --- Start Expansion: Node Validation Logic ---
            validated_nodes: List[QdrantNode] = []
            node_ids_in_chapter = set() # Keep track of IDs defined in this chapter's response
            current_iso_ts = datetime.datetime.now(datetime.timezone.utc).isoformat() # Timestamp for this batch

            for node_dict in qdrant_nodes_data:
                if not isinstance(node_dict, dict):
                    logging.warning(f"[{function_name}] Skipping non-dict item found in qdrant_nodes list for chapter '{chapter_id}'.")
                    continue

                try:
                    # --- A. Extract and Enrich Properties ---
                    properties_dict = node_dict.get("properties")
                    if not isinstance(properties_dict, dict):
                        # Log the problematic node structure for debugging
                        logging.warning(f"[{function_name}] Node '{node_dict.get('node_id')}' missing or invalid 'properties' dict. Skipping node. Node data: {str(node_dict)[:500]}")
                        continue

                    # Add fields managed by Python code
                    properties_dict['source'] = source_identifier
                    properties_dict['stored_at'] = current_iso_ts

                    # Ensure required properties keys exist (provide defaults if LLM missed them)
                    required_props = ["title", "source_credibility", "information_type", "core_claim", "key_entities", "sentiment_tone", "tags"]
                    for key in required_props:
                         if key not in properties_dict or properties_dict[key] is None:
                              logging.debug(f"[{function_name}] Node '{node_dict.get('node_id')}' missing property '{key}'. Setting default.")
                              properties_dict[key] = f"Missing: {key}" # Or specific defaults like "Unknown", "[]"

                    # Ensure key_entities and tags are valid JSON strings
                    for key in ["key_entities", "tags"]:
                        val = properties_dict.get(key)
                        if isinstance(val, list): # If LLM gave a list instead of string
                            try:
                                properties_dict[key] = json.dumps(val)
                            except TypeError:
                                logging.warning(f"[{function_name}] Could not JSON stringify {key} for node '{node_dict.get('node_id')}': {val}. Setting to '[]'.")
                                properties_dict[key] = "[]"
                        elif not isinstance(val, str): # Handle other invalid types (e.g., None, int)
                            logging.warning(f"[{function_name}] Property '{key}' for node '{node_dict.get('node_id')}' is not a string or list ({type(val)}). Setting to '[]'.")
                            properties_dict[key] = "[]"
                        else:
                            # Validate if it's supposed to be a JSON list string
                            try:
                                json.loads(val) # Check if it's valid JSON
                            except json.JSONDecodeError:
                                logging.warning(f"[{function_name}] Property '{key}' for node '{node_dict.get('node_id')}' is not a valid JSON string: '{val}'. Attempting to wrap.")
                                # Try to wrap it as a list if it's a simple string
                                try:
                                     properties_dict[key] = json.dumps([val.strip()])
                                except TypeError:
                                     properties_dict[key] = "[]" # Fallback

                    # --- B. Prepare Node Dictionary for Pydantic ---
                    # The main node_dict already contains 'node_id', 'node_type', 'name', 'chapter_id_context' from LLM
                    node_dict['properties'] = properties_dict # Assign enriched properties

                    # Remove fields that should not be at the top level (handled by properties)
                    node_dict.pop('description', None)
                    node_dict.pop('document_context', None)

                    # Ensure linked_nodes is present and is a list
                    linked_nodes_list = node_dict.setdefault('linked_nodes', [])
                    if not isinstance(linked_nodes_list, list):
                        logging.warning(f"[{function_name}] Node '{node_dict.get('node_id')}' had non-list 'linked_nodes'. Resetting to [].")
                        node_dict['linked_nodes'] = []

                    # Ensure keywords within links are lists (important for Pydantic LinkedNode model)
                    for link_item in node_dict['linked_nodes']:
                        if isinstance(link_item, dict):
                            kw = link_item.setdefault('relationship_keywords', [])
                            if not isinstance(kw, list):
                                logging.debug(f"[{function_name}] Correcting non-list keywords for link in node '{node_dict.get('node_id')}'.")
                                link_item['relationship_keywords'] = []
                            # Ensure relationship_strength is float/convertible
                            strength = link_item.get('relationship_strength', 5.0)
                            try:
                                link_item['relationship_strength'] = float(strength)
                            except (ValueError, TypeError):
                                logging.warning(f"[{function_name}] Invalid relationship_strength '{strength}' for link in node '{node_dict.get('node_id')}'. Setting to 5.0.")
                                link_item['relationship_strength'] = 5.0


                    # --- C. Validate with Pydantic (QdrantNode) ---
                    node_obj = QdrantNode(**node_dict)
                    validated_nodes.append(node_obj)
                    node_ids_in_chapter.add(node_obj.node_id) # Add validated ID to set

                except ValidationError as node_error:
                    logging.warning(f"[{function_name}] Pydantic validation failed for a node in chapter '{chapter_id}': {node_error}. Skipping this node. Data: {str(node_dict)[:500]}...")
                except Exception as unexpected_err:
                    logging.error(f"[{function_name}] Unexpected error processing node data in chapter '{chapter_id}': {unexpected_err}. Skipping node. Data: {str(node_dict)[:500]}...", exc_info=True)
            # --- End Expansion: Node Validation Logic ---


            # --- 6b. Post-Validation Link Check ---
            # --- Start Expansion: Link Validation Logic ---
            final_validated_nodes = []
            MAX_LINKS_PER_NODE = 10 # Set a sensible limit

            for node in validated_nodes: # Iterate through Pydantic QdrantNode objects
                 valid_links_for_node = []
                 link_count = 0
                 # Ensure linked_nodes attribute exists and is a list
                 node_links = getattr(node, 'linked_nodes', [])
                 if node_links is None: node_links = [] # Handle None case

                 for link in node_links:
                      # Ensure link is a valid LinkedNode object with target_node_id
                      if not isinstance(link, LinkedNode) or not hasattr(link, 'target_node_id') or not link.target_node_id:
                           logging.warning(f"[{function_name}] Skipping invalid link object found for source node '{node.node_id}'. Link: {link}")
                           continue

                      if link_count >= MAX_LINKS_PER_NODE:
                           logging.debug(f"[{function_name}] Max link limit ({MAX_LINKS_PER_NODE}) reached for node '{node.node_id}'. Dropping further links.")
                           break

                      # Check if the target ID exists in the set of nodes defined *in this chapter's response*
                      if link.target_node_id not in node_ids_in_chapter:
                           logging.warning(f"[{function_name}] Node '{node.node_id}' links to '{link.target_node_id}' which was NOT defined or validated in this chapter's response. Dropping link.")
                           continue

                      # If target ID is valid and limit not reached, keep the link
                      valid_links_for_node.append(link)
                      link_count += 1

                 # Update the node's linked_nodes list with only the valid ones
                 node.linked_nodes = valid_links_for_node
                 final_validated_nodes.append(node)
            # --- End Expansion: Link Validation Logic ---

            chapter.qdrant_nodes = final_validated_nodes
            logging.info(f"[{function_name}] Assigned {len(chapter.qdrant_nodes)} validated nodes (with properties) for chapter '{chapter_id}'.")
        else:
             logging.error(f"[{function_name}] Final data had 'qdrant_nodes' key but value was not a list (type: {type(qdrant_nodes_data)}) for chapter '{chapter_id}'. Assigning empty list.")
             chapter.qdrant_nodes = []
    else:
        # Failure after initial retries AND fallback attempt
        logging.error(f"[{function_name}] All node extraction attempts (including text fallback) failed for chapter '{chapter_id}'. Assigning empty list.")
        chapter.qdrant_nodes = [] # Assign empty list on complete failure

    # Function modifies chapter in-place, no return value needed.



async def analyze_inter_chapter_relationships(client: genai.Client, chapters: List[Chapter], source_identifier: Optional[str] = "Document") -> Optional[Dict[str, Any]]:
    """
    Analyzes inter-chapter relationships AND generates a document summary based on aggregated nodes.
    Updates the `linked_nodes` attribute of the QdrantNode objects within the chapters list *in-place*.
    Returns a dictionary containing the document summary details, or None on failure.
    Uses gemini-2.0-flash.
    """
    function_name = "analyze_inter_chapter_relationships"
    logging.info(f"[{function_name}] Starting analysis for {len(chapters)} chapters of '{source_identifier}'.")

    # --- 1. Aggregate Node Information & Create Lookup Map ---
    all_nodes_for_prompt = []
    node_map: Dict[str, QdrantNode] = {} # node_id -> QdrantNode object

    for chapter in chapters:
        if not chapter.qdrant_nodes: continue
        for node in chapter.qdrant_nodes:
            if not isinstance(node, QdrantNode): continue
            all_nodes_for_prompt.append({
                "node_id": node.node_id, "node_type": node.node_type,
                "name": node.name, "description": node.description,
                "chapter_id_context": node.chapter_id_context
            })
            if node.node_id in node_map:
                 logging.warning(f"[{function_name}] Duplicate node_id '{node.node_id}' found. Overwriting map entry.")
            node_map[node.node_id] = node

    if len(all_nodes_for_prompt) == 0: # Check for zero nodes specifically
        logging.warning(f"[{function_name}] No nodes found across all chapters for '{source_identifier}'. Cannot perform analysis or generate summary.")
        return None # No nodes, no analysis, no summary

    logging.info(f"[{function_name}] Aggregated {len(all_nodes_for_prompt)} nodes from {len(chapters)} chapters for '{source_identifier}'.")

    # --- 2. Define the Combined LLM Prompt ---
    # Asks for both inter-chapter links AND document summary details
    inter_chapter_and_summary_prompt = f"""
    ---Goal---
    Analyze the provided list of concepts and entities (nodes) aggregated from all chapters of the document '{source_identifier}'. Perform two tasks:
    1. Identify significant relationships *between* nodes originating from *different* chapters.
    2. Generate a concise document-level summary based on the overall collection of nodes.


    ---Input Data---
    List of node definitions:
    ```json
    {json.dumps(all_nodes_for_prompt, indent=2)}
    ```

    ---Instructions---
    **Task 1: Identify Inter-Chapter Relationships**
    1a. Compare nodes across the entire list, focusing on potential connections between nodes from different `chapter_id_context`.
    1b. Identify pairs (`source_node`, `target_node`) with **different** `chapter_id_context` values that have a direct, meaningful relationship based on their names and descriptions.
    1c. For each significant inter-chapter relationship identified, determine: `source_node_id`, `target_node_id`, `relationship_description`, `relationship_keywords` (list), `relationship_strength` (float 1.0-10.0).
    1d. **Prioritize Quality & Limit:** Focus on the 3-10 *strongest* inter-chapter links *per source node*. Do not exceed 10 outgoing links per source.

    **Task 2: Generate Document Summary**
    2a. Synthesize the information from the *entire list* of nodes (names, descriptions, types, chapter context).
    2b. Generate the following summary components:
        *   `title`: A concise title for the document '{source_identifier}' based on the aggregated nodes (max 10 words).
        *   `themes`: A list of 3-7 main themes or topics reflected across all nodes.
        *   `questions`: A list of 2-4 insightful questions a reader might have after understanding the key concepts/entities.
        *   `summary`: A comprehensive summary paragraph (4-8 sentences) synthesizing the key nodes and their overall significance or narrative represented by the aggregation.

    **Output Format:**
    Return ONLY a single valid JSON object containing TWO top-level keys: "inter_chapter_links" and "document_summary_details".

    *   `inter_chapter_links`: A JSON array containing objects, where each object represents one identified *inter-chapter* relationship from Task 1 (with keys: `source_node_id`, `target_node_id`, `relationship_description`, `relationship_keywords`, `relationship_strength`).
    *   `document_summary_details`: A JSON object containing the results from Task 2 (with keys: `title`, `themes`, `questions`, `summary`).

    ---Explicit Exclusions---
    *   DO NOT include intra-chapter relationships (source/target in same chapter) in the `inter_chapter_links` list.
    *   DO NOT include the original node definitions in the output.

    Example Output Format:
    ```json
    {{
      "inter_chapter_links": [
        {{
          "source_node_id": "concept_machine_learning_ch2",
          "target_node_id": "entity_gpu_optimization_ch5",
          "relationship_description": "Utilizes concepts from Chapter 2...",
          "relationship_keywords": ["hardware acceleration", "optimization"],
          "relationship_strength": 8.5
        }}
      ],
      "document_summary_details": {{
        "title": "AI Development and Hardware Optimization",
        "themes": ["Machine Learning", "GPU Computing", "Optimization Techniques", "Acme Corp Strategy"],
        "questions": [
          "How does Acme Corp leverage machine learning?",
          "What are the key challenges in GPU optimization mentioned?"
        ],
        "summary": "The document outlines core machine learning concepts introduced in early chapters, detailing Acme Corp's involvement. Later chapters focus on the technical challenges and solutions related to GPU optimization, showing practical application of the earlier concepts."
      }}
    }}
    ```
    Ensure the final output is only the single JSON object described.
    """

    # --- 3. Call LLM and Process Response ---
    document_summary_output = None # Initialize summary return value

    try:
        logging.debug(f"[{function_name}] Sending combined analysis request for '{source_identifier}' to LLM ({'gemini-2.5-flash-preview-04-17'}).")
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.5-flash-preview-04-17", # Use the specified powerful model - **NOTE: Ensure correct model name per Gemini API docs**
            contents=[
                types.Content(parts=[types.Part.from_text(text=inter_chapter_and_summary_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text) # Expecting the dict with two keys

            if not isinstance(data, dict):
                raise ValueError("LLM response is not a JSON object.")

            # --- Process Inter-Chapter Links (Updates node_map in-place) ---
            newly_found_links_data = data.get("inter_chapter_links", [])
            if isinstance(newly_found_links_data, list):
                links_added_count, links_skipped_count = _merge_links_into_nodes(newly_found_links_data, node_map)
                logging.info(f"[{function_name}] Merged inter-chapter links for '{source_identifier}'. Added: {links_added_count}, Skipped: {links_skipped_count}")
            else:
                logging.warning(f"[{function_name}] Key 'inter_chapter_links' missing or not a list in LLM response for '{source_identifier}'.")

            # --- Process Document Summary Details ---
            summary_details_data = data.get("document_summary_details")
            if isinstance(summary_details_data, dict):
                 try:
                     # Validate the extracted summary details (optional but recommended)
                     validated_summary = DocumentSummaryDetails(**summary_details_data)
                     document_summary_output = validated_summary.model_dump() # Store as dict
                     logging.info(f"[{function_name}] Successfully extracted document summary details for '{source_identifier}'.")
                 except ValidationError as summary_error:
                     logging.error(f"[{function_name}] Validation failed for document_summary_details from LLM for '{source_identifier}': {summary_error}")
                     document_summary_output = None # Indicate summary failure
            else:
                logging.warning(f"[{function_name}] Key 'document_summary_details' missing or not a dict in LLM response for '{source_identifier}'.")

        else:
            logging.error(f"[{function_name}] Combined analysis response from LLM had no candidates for '{source_identifier}'.")

    except (json.JSONDecodeError, ValidationError, Exception) as e:
        logging.error(f"[{function_name}] Error during combined analysis or parsing LLM response for '{source_identifier}': {e}", exc_info=True)
        # Summary generation failed

    # The chapters list has been modified in-place by _merge_links_into_nodes via node_map
    logging.info(f"[{function_name}] Finished analysis for '{source_identifier}'. Returning summary.")
    return document_summary_output

# --- Helper function to merge links (extracted from previous version) ---
def _merge_links_into_nodes(link_data_list: List[Dict], node_map: Dict[str, QdrantNode]) -> (int, int):
    """Merges link data into the linked_nodes of QdrantNode objects in node_map."""
    links_added_count = 0
    links_skipped_count = 0
    MAX_OUTGOING_LINKS_PER_NODE = 10

    if not isinstance(link_data_list, list):
        logging.warning("[_merge_links_into_nodes] Input link_data_list is not a list.")
        return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict):
            logging.warning(f"[_merge_links_into_nodes] Skipping invalid item in link data: {link_data}")
            links_skipped_count += 1
            continue

        source_id = link_data.get("source_node_id")
        target_id = link_data.get("target_node_id")
        desc = link_data.get("relationship_description")
        keywords = link_data.get("relationship_keywords", [])
        strength_str = link_data.get("relationship_strength", "5.0")

        if not all([source_id, target_id, desc]) or not isinstance(source_id, str) or not isinstance(target_id, str):
            logging.warning(f"[_merge_links_into_nodes] Skipping incomplete/invalid link data: {link_data}")
            links_skipped_count += 1
            continue

        try: strength = float(strength_str)
        except (ValueError, TypeError): strength = 5.0

        source_node = node_map.get(source_id)
        target_node = node_map.get(target_id)

        if not source_node or not target_node:
            logging.warning(f"[_merge_links_into_nodes] Source ('{source_id}') or Target ('{target_id}') not found. Skipping.")
            links_skipped_count += 1
            continue

        if source_node.chapter_id_context == target_node.chapter_id_context:
            logging.debug(f"[_merge_links_into_nodes] Skipping INTRA-chapter link: {source_id} -> {target_id}")
            links_skipped_count += 1
            continue

        # Add Forward Link
        if source_node.linked_nodes is None: source_node.linked_nodes = []
        if len(source_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == target_id for link in source_node.linked_nodes):
            try:
                forward_link = LinkedNode(**link_data) # Directly use link_data if keys match LinkedNode
                source_node.linked_nodes.append(forward_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added FORWARD link: {source_id} -> {target_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating forward LinkedNode {source_id} -> {target_id}: {ve}. Skipping.")
                 links_skipped_count += 1
        else:
             links_skipped_count += 1 # Skipped due to limit or duplicate

        # Add Reverse Link
        if target_node.linked_nodes is None: target_node.linked_nodes = []
        if len(target_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == source_id for link in target_node.linked_nodes):
            try:
                reverse_desc = f"Related to '{source_node.name}' from {source_node.chapter_id_context} ({desc})"
                reverse_link = LinkedNode(target_node_id=source_id, relationship_description=reverse_desc, relationship_keywords=keywords, relationship_strength=strength)
                target_node.linked_nodes.append(reverse_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added REVERSE link: {target_id} -> {source_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating reverse LinkedNode {target_id} -> {source_id}: {ve}. Skipping.")
                 # links_skipped_count += 1 # Optional separate count
        else:
             # links_skipped_count += 1 # Optional separate count
             pass # Skip reverse if limit reached or duplicate


    return links_added_count, links_skipped_count

# async def generate_project_ontology(client: genai.Client,
#                                   documents: List[Dict], # List of finalized document dicts
#                                   max_structure_retries: int = 1, # Less retries needed for summary usually
#                                   retry_delay: int = 5
#                                   ) -> Optional[ProjectOntology]:
#     """
#     Generates a project-wide ontology by analyzing aggregated nodes from all
#     provided documents. Identifies inter-document relationships and synthesizes
#     project-level metadata (overview, themes, concepts).

#     Args:
#         client: The configured GenAI client.
#         documents: A list of finalized document dictionaries, each expected to
#                    contain 'raw_extracted_content' -> 'summary' -> 'chapters' -> 'qdrant_nodes'.
#         max_structure_retries: Max retries if LLM returns incorrect JSON structure.
#         retry_delay: Base delay between retries.

#     Returns:
#         A ProjectOntology object containing the synthesized information and
#         aggregated graph nodes, or None if processing fails significantly.
#     """
#     function_name = "generate_project_ontology_with_properties" # Renamed for clarity
#     logging.info(f"[{function_name}] Starting ontology generation for {len(documents)} documents.")

#     # --- 1. Aggregate Node Information & Document Details (Adjusted) ---
#     all_project_nodes_data: List[Dict] = [] # Store node dicts (as dumped from Pydantic)
#     node_map: Dict[str, Dict] = {} # node_id -> node dict (for link merging)
#     doc_filenames: List[str] = []
#     doc_count = 0

#     for doc_dict in documents:
#         if not isinstance(doc_dict, dict) or "raw_extracted_content" not in doc_dict:
#             logging.warning(f"[{function_name}] Skipping invalid document structure: {type(doc_dict)}")
#             continue

#         raw_content = doc_dict["raw_extracted_content"]
#         if not isinstance(raw_content, dict): continue
#         doc_filename = raw_content.get("filename", f"Unknown_Doc_{doc_count}")
#         doc_filenames.append(doc_filename)
#         doc_count += 1
#         summary = raw_content.get("summary")
#         if not isinstance(summary, dict): continue
#         chapters = summary.get("chapters")
#         if not isinstance(chapters, list): continue

#         nodes_found_in_doc = 0
#         for chapter in chapters:
#             if not isinstance(chapter, dict): continue
#             # qdrant_nodes should be a list of DICTIONARIES already dumped from Pydantic
#             qdrant_nodes_list = chapter.get("qdrant_nodes")
#             if isinstance(qdrant_nodes_list, list):
#                 for node_dict in qdrant_nodes_list:
#                     # Ensure it's a dict and has the core fields + properties
#                     if isinstance(node_dict, dict) and "node_id" in node_dict and "properties" in node_dict:
#                         node_id = node_dict["node_id"]
#                         # Ensure 'source' is present in properties for uniqueness check later
#                         node_dict.setdefault('properties', {}).setdefault('source', doc_filename)

#                         if node_id in node_map:
#                             logging.warning(f"[{function_name}] Duplicate node_id '{node_id}' encountered across documents ('{node_map[node_id].get('properties',{}).get('source')}' vs '{doc_filename}'). Overwriting map entry.")
#                         node_map[node_id] = node_dict # Store the dict itself
#                         all_project_nodes_data.append(node_dict) # Add dict to list for prompt
#                         nodes_found_in_doc += 1
#         logging.debug(f"[{function_name}] Aggregated {nodes_found_in_doc} nodes from '{doc_filename}'.")

#     if not all_project_nodes_data:
#         logging.warning(f"[{function_name}] No nodes found across any documents. Cannot generate graph or detailed ontology.")
#         # Return a basic ontology object indicating failure
#         return ProjectOntology(
#             title=f"Project Ontology (No Node Data)",
#             overview="No node data could be extracted from the documents to perform analysis.",
#             document_count=doc_count,
#             documents=doc_filenames,
#             global_themes=[],
#             key_concepts=[],
#             project_graph_nodes=[] # Empty list
#         )

#     logging.info(f"[{function_name}] Aggregated {len(all_project_nodes_data)} total nodes from {doc_count} documents for analysis.")

#     # --- 2. Define the LLM Prompt (Adjusted) ---
#     # Prepare node info for prompt (using properties)
#     nodes_for_prompt = []
#     for n_dict in all_project_nodes_data:
#          props = n_dict.get("properties", {})
#          nodes_for_prompt.append({
#              "node_id": n_dict.get("node_id"),
#              "node_type": n_dict.get("node_type"),
#              "name": n_dict.get("name"),
#              # Send relevant properties for context
#              "title": props.get("title"),
#              "core_claim": props.get("core_claim"),
#              "information_type": props.get("information_type"),
#              "tags": props.get("tags"), # Send the JSON string tags
#              # Get document context from 'source' property
#              "document_context": props.get("source", "unknown_document")
#          })

#     ontology_prompt = f"""
# --- Goal ---
# Analyze the provided list of concepts and entities (nodes), aggregated from ALL documents in a project. Each node includes core details and context like title, core claim, tags, and its source document. Perform two main tasks:
# 1. Identify significant relationships connecting nodes that originate from *different* source documents ('document_context').
# 2. Synthesize the information from *all* nodes to generate high-level project metadata (title, overview, themes, key concepts).

# --- Input Data ---
# List of node definitions from across the project:
# ```json
# {json.dumps(nodes_for_prompt, indent=2)}
# Total Documents Analyzed: {doc_count}
# --- Instructions ---
# Perform the following tasks based on the entire "Input Data":
# Task 1: Identify Inter-Document Relationships
# 1a. Compare nodes across the entire list. Focus specifically on finding meaningful connections between nodes that have different document_context values.
# 1b. Identify pairs (source_node, target_node) where source_node.document_context is DIFFERENT from target_node.document_context, and where a direct relationship (e.g., usage, comparison, causality, reference) can be inferred from their names, titles, core claims, or tags.
# 1c. For each significant inter-document relationship identified, create a link object with: source_node_id, target_node_id, relationship_description, relationship_keywords (optional, default []), and relationship_strength (float 1.0-10.0).
# 1d. Prioritize & Limit: Focus on the 5-10 most significant inter-document links per source node. Limit total outgoing links per source node added in this step.
# 1e. Compile these link objects into the inter_document_links list. If none found, use empty [].
# Task 2: Generate Project Summary Details
# 2a. Synthesize information from the entire list of input nodes (consider names, titles, core claims, information types, tags, document context).
# 2b. Generate the following components within the project_summary_details object:
# * title: A concise title for the overall project (max 10 words).
# * overview: A comprehensive summary paragraph (3-6 sentences) describing the project's main subjects, scope, purpose, and key findings derived from the aggregated nodes.
# * global_themes: A list of 3-7 high-level themes reflected across multiple documents/nodes (consider tags and core claims).
# * key_concepts: A list of the 5-10 most important or frequently occurring concepts identified across the project nodes (consider node names, titles, tags).
# --- Output Format ---
# Return ONLY a single, valid JSON object with TWO top-level keys: "project_summary_details" and "inter_document_links". Adhere STRICTLY to the schema:
# {{
#   "project_summary_details": {{ // REQUIRED object
#     "title": "string",
#     "overview": "string",
#     "global_themes": ["string"], // List
#     "key_concepts": ["string"] // List
#   }},
#   "inter_document_links": [ // REQUIRED list (can be empty [])
#     {{
#       "source_node_id": "string",
#       "target_node_id": "string", // Must be from different document_context
#       "relationship_description": "string",
#       "relationship_keywords": ["string"], // Optional list
#       "relationship_strength": "float"
#     }}
#   ]
# }}
# --- Critical Rules ---
# VALID JSON ONLY.
# SCHEMA ADHERENCE: Follow the output format precisely.
# INTER-DOCUMENT LINKS ONLY: inter_document_links must ONLY contain relationships where source and target nodes have different document_context values based on the input data.
# NODE ID CONSISTENCY: source_node_id and target_node_id MUST correspond to node_id values present in the input node list.
# TEXTUAL BASIS: Base summary details and relationships on the provided node information (names, titles, claims, tags, etc.).
# Generate the JSON output now:
# """

#     # --- 3. Call LLM and Process Response with Retries ---
#     llm_output_data = None
#     final_ontology = None

#     for attempt in range(max_structure_retries):
#         logging.debug(f"[{function_name}] Attempt {attempt + 1}/{max_structure_retries} for project ontology.")
#         try:
#             response = await retry_api_call(
#                 client.aio.models.generate_content,
#                 model="gemini-2.5-flash-preview-04-17", # Capable model needed
#                 contents=[types.Content(parts=[types.Part.from_text(text=ontology_prompt)])],
#                 config=types.GenerateContentConfig(response_mime_type="application/json")
#             )

#             if not response or not response.candidates:
#                 logging.error(f"[{function_name}] No candidates (Attempt {attempt + 1}) for project ontology.")
#                 llm_output_data = None; break

#             json_text = clean_json_response(response.candidates[0].content.parts[0].text)
#             logging.debug(f"[{function_name}] Cleaned LLM response (Attempt {attempt + 1}): {json_text[:500]}...")

#             try: data = json.loads(json_text)
#             except json.JSONDecodeError as json_err:
#                 logging.warning(f"[{function_name}] JSONDecodeError attempt {attempt + 1}: {json_err}.")
#                 if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
#                 else: logging.error(f"[{function_name}] Failed JSON decode after retries."); llm_output_data = None; break

#             # Structure Check
#             if isinstance(data, dict) and "project_summary_details" in data and "inter_document_links" in data:
#                 logging.info(f"[{function_name}] Received correct structure attempt {attempt + 1}.")
#                 llm_output_data = data; break
#             else:
#                 logging.warning(f"[{function_name}] Incorrect JSON structure attempt {attempt+1}. Got: {type(data)}. Keys: {list(data.keys()) if isinstance(data,dict) else 'N/A'}")
#                 if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
#                 else: logging.error(f"[{function_name}] Failed structure check after retries."); llm_output_data = None; break

#         except Exception as e:
#             logging.error(f"[{function_name}] Error attempt {attempt + 1}: {e}", exc_info=True)
#             if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
#             else: logging.error(f"[{function_name}] Failed after retries: {e}"); llm_output_data = None; break

#     # --- 4. Process LLM Output ---
#     if llm_output_data is not None:
#         summary_details = llm_output_data.get("project_summary_details", {})
#         inter_doc_links = llm_output_data.get("inter_document_links", [])

#         # Merge the new inter-document links into the nodes stored in node_map
#         if isinstance(inter_doc_links, list) and inter_doc_links:
#             logging.info(f"[{function_name}] Merging {len(inter_doc_links)} potential inter-document links.")
#             # Use the adjusted helper function
#             added_count, skipped_count = _merge_project_links_into_node_dicts_properties(inter_doc_links, node_map)
#             logging.info(f"[{function_name}] Link merge results - Added: {added_count}, Skipped: {skipped_count}")

#         # --- 5. Validate Final Nodes (Same logic, but relies on correct node_map) ---
#         validated_project_nodes: List[QdrantNode] = []
#         final_node_list_for_ontology = list(node_map.values()) # Get dicts (potentially updated with links)
#         logging.info(f"[{function_name}] Validating final {len(final_node_list_for_ontology)} project nodes with Pydantic.")
#         for node_dict in final_node_list_for_ontology:
#             try:
#                 # Ensure linked_nodes is list (might have been added/modified)
#                 if 'linked_nodes' not in node_dict or not isinstance(node_dict['linked_nodes'], list):
#                     node_dict['linked_nodes'] = []
#                 # Re-validate the potentially modified node dict using the NEW QdrantNode model
#                 node_obj = QdrantNode(**node_dict)
#                 validated_project_nodes.append(node_obj)
#             except ValidationError as node_error:
#                 logging.warning(f"[{function_name}] Pydantic validation failed for final project node '{node_dict.get('node_id')}': {node_error}. Skipping node.")
#             except Exception as final_val_err:
#                 logging.error(f"[{function_name}] Unexpected error during final node validation '{node_dict.get('node_id')}': {final_val_err}", exc_info=True)

#         # --- 6. Create Final ProjectOntology Object ---
#         try:
#             final_ontology = ProjectOntology(
#                 title=summary_details.get("title", f"Project Ontology ({doc_count} Docs)"),
#                 overview=summary_details.get("overview", "Overview not generated."),
#                 document_count=doc_count,
#                 documents=doc_filenames, # List of filenames
#                 global_themes=summary_details.get("global_themes", []),
#                 key_concepts=summary_details.get("key_concepts", []),
#                 project_graph_nodes=validated_project_nodes # Assign the validated list
#             )
#             logging.info(f"[{function_name}] Project ontology created successfully with {len(validated_project_nodes)} nodes.")

#         except ValidationError as ve:
#             logging.error(f"[{function_name}] Failed to validate final ProjectOntology object: {ve}")
#             # Return basic ontology on final validation error
#             final_ontology = ProjectOntology(
#                 title="Project Ontology (Validation Error)",
#                 overview=f"Error creating final ontology object: {ve}",
#                 document_count=doc_count, documents=doc_filenames,
#                 global_themes=[], key_concepts=[], project_graph_nodes=[]
#             )

#     else:
#         # LLM processing failed after retries
#         logging.error(f"[{function_name}] LLM processing failed for project ontology after all retries.")
#         final_ontology = ProjectOntology(
#             title="Project Ontology (LLM Error)",
#             overview="Failed to generate project overview and relationships via LLM.",
#             document_count=doc_count, documents=doc_filenames,
#             global_themes=[], key_concepts=[], project_graph_nodes=[]
#         )

#     return final_ontology

async def incrementally_generate_doc_level_ontology(
    client: genai.Client, # Replace Any with your actual genai.Client type hint
    documents: List[Dict],
    *,
    # --- Constants moved inside or passed as parameters ---
    ontology_path: str | pathlib.Path = "doc_level_ontology.json", # Default path
    char_budget: int = 100_000,
    max_new_links: int = 50,
    branch_size: int = 20,
    llm_retries: int = 2,
    llm_delay: int = 5
) -> Optional[ProjectOntology]:
    """
    Builds or incrementally updates ontology using DIRECT Gemini API calls.
    Includes nested tree summarizer and internal constants.

    Args:
        client: The configured Generative AI client instance.
        documents: List of *all* document dictionaries.
        ontology_path: Path to load/save the ontology JSON file.
        char_budget: Max characters for prompts (affects summarization).
        max_new_links: Max new links to generate in incremental step.
        branch_size: Fan-out for tree summarization.
        llm_retries: Number of retries for Gemini API calls.
        llm_delay: Base delay (seconds) between retries.

    Returns:
        The generated/updated ProjectOntology object, or None on critical failure.
    """
    # --- Setup Logging and Constants ---
    log = logging.getLogger("incremental_doc_ontology")
    path = pathlib.Path(ontology_path)
    # Constants are now local to this function scope
    CHAR_BUDGET = char_budget
    MAX_NEW_LINKS = max_new_links
    BRANCH_SIZE = branch_size
    LLM_RETRIES = llm_retries
    LLM_DELAY = llm_delay

    # ----------------------------------------------------------------------- #
    # ❹ Tree Summarizer (Nested Function)                                     #
    # ----------------------------------------------------------------------- #
    async def _tree_summarise_nested(
        input_docs: List[Dict[str, Any]],
        *,
        target_chars: int # Target chars passed explicitly
    ) -> List[Dict[str, Any]]:
        """
        NESTED: Recursively compress docs until JSON size <= target_chars.
        Uses client, BRANCH_SIZE, CHAR_BUDGET, LLM_RETRIES, LLM_DELAY from outer scope.
        """
        tree_log = logging.getLogger("tree_summarise_nested") # Use specific logger

        def blob_size(nodes_list: List[Dict]) -> int:
            try: return len(json.dumps(nodes_list, ensure_ascii=False))
            except TypeError: return float('inf')

        current_size = blob_size(input_docs)
        tree_log.debug(f"Input docs: {len(input_docs)}, Current size: {current_size}, Target size: {target_chars}")

        if current_size <= target_chars: return input_docs
        if not input_docs: return []
        if len(input_docs) == 1 and current_size > target_chars:
            tree_log.warning(f"Single document exceeds target size ({current_size} > {target_chars}). Cannot compress further.")
            return input_docs

        groups = [input_docs[i:i + BRANCH_SIZE] for i in range(0, len(input_docs), BRANCH_SIZE)]
        tree_log.info(f"Compressing {len(input_docs)} docs into {len(groups)} groups (branch size {BRANCH_SIZE}).")

        compressed_results: List[Dict] = []
        tasks = []

        # --- Direct API Call Task Creation within Nested Tree Summarizer ---
        async def call_gemini_for_summary(prompt: str, group_idx: int) -> Optional[Dict]:
            """Inner function for single API call with retry & JSON parsing"""
            # Uses LLM_RETRIES, LLM_DELAY, client from outer scope
            for attempt in range(LLM_RETRIES):
                try:
                    tree_log.debug(f"Calling Gemini for group {group_idx} (Attempt {attempt+1}/{LLM_RETRIES})")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text
                    raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    parsed_json = json.loads(raw_text)
                    return parsed_json
                except Exception as e:
                    tree_log.warning(f"Gemini call/parse failed group {group_idx} (Attempt {attempt+1}/{LLM_RETRIES}): {e}")
                    if attempt == LLM_RETRIES - 1:
                        tree_log.error(f"Gemini call failed group {group_idx} after {LLM_RETRIES} attempts.")
                        return None
                    await asyncio.sleep(LLM_DELAY * (attempt + 1))
            return None

        for idx, chunk in enumerate(groups, 1):
            if not chunk: continue
            chunk_for_prompt = []
            estimated_prompt_size = 1000
            # Use CHAR_BUDGET from outer scope
            chars_per_doc_limit = (CHAR_BUDGET - estimated_prompt_size) // len(chunk) if len(chunk) > 0 else CHAR_BUDGET
            chars_per_doc_limit = max(500, chars_per_doc_limit)
            for doc_item in chunk:
                 summary = doc_item.get("executive_summary", "")
                 truncated_summary = summary[:chars_per_doc_limit] + "..." if len(summary) > chars_per_doc_limit else summary
                 chunk_for_prompt.append({"node_id": doc_item.get("node_id"), "document_name": doc_item.get("document_name"), "executive_summary_preview": truncated_summary})

            chunk_json = json.dumps(chunk_for_prompt, indent=2, ensure_ascii=False)
            # Use CHAR_BUDGET from outer scope
            if len(chunk_json) > (CHAR_BUDGET * 0.8): chunk_json = chunk_json[:int(CHAR_BUDGET * 0.8)] + "\n... ]\n```"

            prompt = f"""
                You are an expert summarizer consolidating information from multiple document summaries.
                Input: {len(chunk)} summaries/previews:
                ```json
                {chunk_json}
                ```
                Goal: Create **one single, concise meta-summary node** representing this group.
                Instructions: Synthesize key ideas into an executive summary (<800 words). Assign node_id "meta::{idx}", document_name "meta_group_{idx}".
                Output Format: **Valid JSON only**: `{{"meta_node": {{"node_id": "meta::{idx}", "document_name": "meta_group_{idx}", "executive_summary": "<summary>"}}}}`
                Group Index: {idx}
                Generate JSON:"""
            tasks.append(asyncio.create_task(call_gemini_for_summary(prompt, idx), name=f"SummarizeGroup_{idx}"))

        task_results = await asyncio.gather(*tasks, return_exceptions=False)

        for i, result in enumerate(task_results, 1):
            if result is None: tree_log.error(f"Skipping group {i} due to API/JSON failure."); continue
            if isinstance(result, dict) and "meta_node" in result and isinstance(result["meta_node"], dict):
                meta_node = result["meta_node"]
                if ("node_id" in meta_node and "document_name" in meta_node and "executive_summary" in meta_node): compressed_results.append(meta_node)
                else: tree_log.warning(f"Invalid meta_node structure group {i}: {meta_node}")
            else: tree_log.warning(f"Unexpected result type group {i}: {type(result)}")

        if not compressed_results:
             tree_log.error("All summarization tasks failed. Returning original docs.")
             return input_docs # Return original to avoid losing data

        tree_log.debug(f"Recursion step. Input docs: {len(compressed_results)}, Target chars: {target_chars}")
        # Recursive call to the *nested* function
        return await _tree_summarise_nested(compressed_results, target_chars=target_chars)
    # ----------------------------------------------------------------------- #
    # End of Nested Tree Summarizer Definition                                #
    # ----------------------------------------------------------------------- #


    # --- 1. Load Existing Ontology or Initialize ---
    ontology = load_ontology_json(path) # Uses external helper
    first_run = ontology is None
    if first_run:
        log.info(f"No ontology file found at '{path}' – creating new.")
        ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
    else:
        if not isinstance(ontology, ProjectOntology):
             log.error(f"Loaded data from '{path}' invalid. Starting fresh.")
             first_run = True
             ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
        else:
            log.info(f"Loaded existing ontology with {ontology.document_count} documents from '{path}'.")


    # --- 2. Identify New Documents ---
    existing_files = set(ontology.documents)
    target_docs_data = []
    if first_run: target_docs_data = documents
    else:
        new_doc_count = 0
        for d in documents:
            try:
                filename = d.get("raw_extracted_content", {}).get("filename")
                if filename and isinstance(filename, str) and filename not in existing_files: target_docs_data.append(d); new_doc_count += 1
            except Exception as e: log.warning(f"Error accessing filename: {e} - Skipping.")
        log.info(f"Incremental run: Found {new_doc_count} new documents.")

    if not target_docs_data: log.info("No new documents found."); return ontology


    # --- 3. Prepare Nodes from New Document Summaries ---
    doc_nodes_for_llm: List[Dict] = []
    processed_filenames = []
    for d in target_docs_data:
        try:
            raw = d.get("raw_extracted_content", {}); fname = raw.get("filename", f"unknown_doc_{random.randint(1000,9999)}")
            summ_block = raw.get("summary", {}); summary_text = ""
            if isinstance(summ_block, dict): summary_text = (summ_block.get("executive_summary") or summ_block.get("summary") or summ_block.get("overview") or "").strip()
            if not summary_text: log.warning(f"Document '{fname}' lacks summary – skipping."); continue
            doc_nodes_for_llm.append({"node_id": f"doc::{fname}", "document_name": fname, "executive_summary": summary_text})
            processed_filenames.append(fname)
        except Exception as e: log.warning(f"Error processing doc data '{raw.get('filename', 'unknown')}': {e} - Skipping.")

    if not doc_nodes_for_llm: log.error("No usable summaries in new docs."); return ontology if not first_run else None
    log.info(f"Prepared {len(doc_nodes_for_llm)} nodes from new summaries.")


    # --- 4. Apply Tree Summarization (using nested function) if Needed ---
    summaries_json_blob = json.dumps(doc_nodes_for_llm, ensure_ascii=False)
    current_char_size = len(summaries_json_blob)
    final_doc_nodes = doc_nodes_for_llm # Default

    if current_char_size > CHAR_BUDGET: # Use local constant
        log.info(f"Summary size ({current_char_size:,} chars) > budget ({CHAR_BUDGET:,}). Summarizing...")
        try:
            # Call the NESTED summarizer
            final_doc_nodes = await _tree_summarise_nested(
                doc_nodes_for_llm,
                target_chars=CHAR_BUDGET # Pass target explicitly
            )
            final_size = len(json.dumps(final_doc_nodes, ensure_ascii=False))
            log.info(f"Summarization complete. Final size: {final_size:,} chars ({len(final_doc_nodes)} nodes).")
            if final_size > CHAR_BUDGET: log.warning(f"Summarization didn't meet budget.")
        except Exception as e:
            log.error(f"Error during tree summarization: {e}. Proceeding with original nodes.", exc_info=True)
            final_doc_nodes = doc_nodes_for_llm # Fallback
    else:
        log.info(f"Summary size ({current_char_size:,} chars) within budget.")


    # --- 5. Call LLM for Ontology Generation / Update (Direct Calls) ---
    llm_links = []
    llm_summary_details = {}
    llm_call_failed = False

    max_json_chars_in_prompt = int(CHAR_BUDGET * 0.85) # Use local constant
    final_doc_nodes_json_str = json.dumps(final_doc_nodes, indent=2, ensure_ascii=False)
    if len(final_doc_nodes_json_str) > max_json_chars_in_prompt:
        log.warning(f"Final doc nodes JSON ({len(final_doc_nodes_json_str)} chars) truncated for prompt.")
        final_doc_nodes_json_str = final_doc_nodes_json_str[:max_json_chars_in_prompt] + "\n... ]\n```"

    try:
        if first_run:
            log.info("Performing initial ontology generation LLM call (direct)...")
            init_prompt = f"""Act as KG Architect. Create initial ontology from summaries.
                Input:
                ```json
                {final_doc_nodes_json_str}
                ```
                Tasks: 1. Find inter-doc relationships. 2. Generate project metadata (title, overview, themes, key_docs).
                Output Schema (JSON ONLY): `{{"project_summary_details": {{"title": "...", "overview": "...", "global_themes": [...], "key_documents": [...]}}, "inter_document_links": [ {{ "source_node_id": "...", ...}} ]}}`
                Rules: Valid JSON, schema adherence, use input node_ids, base on summaries.
                Generate JSON:"""
            llm_out = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Initial Ontology: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=init_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    llm_out = json.loads(raw_text); log.info(f"Initial ontology LLM call successful."); break
                except Exception as e:
                    log.warning(f"Initial Ontology: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Initial Ontology: Call failed after retries."); llm_call_failed = True; break
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if llm_out: llm_links = llm_out.get("inter_document_links", []); llm_summary_details = llm_out.get("project_summary_details", {})
            else: llm_links = []; llm_summary_details = {}

        else: # Incremental update
            log.info("Performing incremental ontology update LLM calls (direct)...")
            # --- 5a. Link Discovery ---
            existing_nodes_for_context = []
            if ontology.project_graph_nodes:
                 # (Sampling logic remains the same)
                 key_doc_node_ids={f"doc::{kd}" for kd in ontology.key_concepts}; existing_node_map={n.node_id: n for n in ontology.project_graph_nodes}
                 priority_nodes=[existing_node_map[nid] for nid in key_doc_node_ids if nid in existing_node_map]; other_nodes=[n for n in ontology.project_graph_nodes if n.node_id not in key_doc_node_ids]
                 sample_size=min(len(ontology.project_graph_nodes), max(len(final_doc_nodes)*2, 30)); sample_count=min(len(ontology.project_graph_nodes), sample_size)
                 sample_nodes_pydantic=priority_nodes; remaining_sample_count=sample_count - len(sample_nodes_pydantic)
                 if remaining_sample_count > 0 and other_nodes: sample_nodes_pydantic.extend(random.sample(other_nodes, min(remaining_sample_count, len(other_nodes))))
                 max_summary_preview=800
                 for n in sample_nodes_pydantic:
                      summary_preview=n.properties.get("executive_summary", "")[:max_summary_preview] + ("..." if len(n.properties.get("executive_summary", "")) > max_summary_preview else "")
                      existing_nodes_for_context.append({"node_id": n.node_id, "document_name": n.name, "executive_summary_preview": summary_preview})

            existing_context_json_str = json.dumps(existing_nodes_for_context, indent=2, ensure_ascii=False)
            max_context_json_chars = int(CHAR_BUDGET * 0.4) # Use local constant
            if len(existing_context_json_str) > max_context_json_chars: existing_context_json_str = existing_context_json_str[:max_context_json_chars] + "\n... ]\n```"

            link_prompt = f"""Act as KG Analyst. Update ontology. Find links involving NEW docs.
                Input 1 (New Docs):
                ```json
                {final_doc_nodes_json_str}
                ```
                Input 2 (Existing Doc Sample):
                ```json
                {existing_context_json_str}
                ```
                Task: Find links (New<->New or New<->Existing). Limit: {MAX_NEW_LINKS} links.
                Output (JSON ONLY): `{{"new_inter_document_links": [ {{ "source_node_id": "...", ... }} ]}}`
                Rules: Valid JSON, use input node_ids, >=1 NEW node per link, empty list [] if none.
                Generate JSON:"""
            link_json = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Incr Links: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=link_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    link_json = json.loads(raw_text); log.info(f"Incr link discovery call successful."); break
                except Exception as e:
                    log.warning(f"Incr Links: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Incr Links: Call failed after retries."); break # Continue even if links fail
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if link_json: llm_links = link_json.get("new_inter_document_links", [])
            else: llm_links = []
            log.info(f"Found {len(llm_links)} potential new links.")

            # --- 5b. Summary Revision ---
            current_summary = {"title": ontology.title, "overview": ontology.overview, "global_themes": ontology.global_themes, "key_documents": ontology.key_concepts}
            current_summary_json_str = json.dumps(current_summary, indent=2)
            rev_prompt = f"""Act as Content Strategist. Update project summary with new docs.
                Input 1 (Current Summary):
                ```json
                {current_summary_json_str}
                ```
                Input 2 (New Docs):
                ```json
                {final_doc_nodes_json_str}
                ```
                Task: Revise title, overview, global_themes, key_documents.
                Output (JSON ONLY): `{{"revised_project_summary_details": {{ "title": "...", "overview": "...", "global_themes": [...], "key_documents": [...] }}}}`
                Rules: Valid JSON, use node_ids for key_docs, integrate new info.
                Generate JSON:"""
            summary_json = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Summary Rev: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=rev_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    summary_json = json.loads(raw_text); log.info("Summary revision call successful."); break
                except Exception as e:
                    log.warning(f"Summary Rev: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Summary Rev: Call failed after retries."); llm_call_failed = True; break
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if summary_json: llm_summary_details = summary_json.get("revised_project_summary_details", {})
            else: llm_summary_details = current_summary # Fallback to current

    except Exception as e: log.error(f"Critical error during LLM processing: {e}", exc_info=True); llm_call_failed = True

    if llm_call_failed and first_run: log.error("Initial LLM calls failed critically. Aborting."); return None
    elif llm_call_failed: log.warning("One+ incremental LLM calls failed. Proceeding with partial update.")

    # --- 6. Merge LLM Results into Ontology ---
    node_map: Dict[str, Dict] = {}
    if not first_run and ontology and ontology.project_graph_nodes:
         for node in ontology.project_graph_nodes:
              try: node_map[node.node_id] = node.model_dump(mode='json')
              except Exception as dump_err: log.warning(f"Could not dump node {node.node_id}: {dump_err}")

    for new_node_data in final_doc_nodes:
         node_id = new_node_data.get("node_id"); doc_name = new_node_data.get("document_name"); exec_summary = new_node_data.get("executive_summary")
         if not node_id or not doc_name or exec_summary is None: continue
         node_type = "document" if node_id.startswith("doc::") else "meta_summary" if node_id.startswith("meta::") else "unknown"

         # --- MODIFICATION START ---
         # Add the 'chapter_id_context' field with a default value
         # Using a specific prefix combined with doc_name seems appropriate
         chapter_context_value = f"doc_level::{doc_name}" # Assign a placeholder

         node_map[node_id] = {
             "node_id": node_id,
             "node_type": node_type,
             "name": doc_name,
             "chapter_id_context": chapter_context_value, # <--- ENSURE THIS LINE IS PRESENT AND UNCOMMENTED
             "properties": {
                 # Add other relevant properties for doc-level nodes if needed
                 "executive_summary": exec_summary,
                 "source": doc_name, # Use doc_name as the source identifier here
                 "title": f"Document Summary: {doc_name}", # Add a default title
                 # Provide defaults for other required properties of QdrantNode.properties
                 "source_credibility": "Document Summary",
                 "information_type": "Summary",
                 "core_claim": exec_summary[:500] + "..." if exec_summary and len(exec_summary) > 500 else exec_summary or "N/A",
                 "key_entities": "[]", # Default for doc level
                 "sentiment_tone": "Neutral", # Default for doc level
                 "tags": "[]", # Default for doc level
                 "stored_at": datetime.datetime.now(datetime.timezone.utc).isoformat() # Add timestamp
             },
             "linked_nodes": node_map.get(node_id, {}).get("linked_nodes", []) # Preserve existing links if updating
         }
         # --- MODIFICATION END ---



    added_links, skipped_links = _merge_project_links_into_node_dicts_properties(llm_links, node_map) # Uses external helper
    log.info(f"Merged LLM links: {added_links} added, {skipped_links} skipped.")

    # --- Validate Nodes before Final Assignment ---
    validated_nodes_list: List[QdrantNode] = []
    final_node_ids = set()
    for node_id, node_dict in node_map.items():
        try:
            if 'linked_nodes' not in node_dict or not isinstance(node_dict['linked_nodes'], list): node_dict['linked_nodes'] = []
            valid_links = []
            for link in node_dict['linked_nodes']:
                if isinstance(link, dict) and 'relationship_strength' in link:
                    try: link['relationship_strength'] = float(link['relationship_strength']); valid_links.append(link)
                    except (ValueError, TypeError): log.warning(f"Invalid strength node {node_id}, removing link: {link}")
                else: valid_links.append(link)
            node_dict['linked_nodes'] = valid_links
            validated_node = QdrantNode.model_validate(node_dict) # Uses global model
            validated_nodes_list.append(validated_node); final_node_ids.add(node_id)
        except ValidationError as e: log.error(f"Pydantic validation failed node '{node_id}'. Skipping. Error: {e}")
        except Exception as final_val_err: log.error(f"Unexpected validation error node '{node_id}': {final_val_err}", exc_info=True)

    for node in validated_nodes_list: node.linked_nodes = [link for link in node.linked_nodes if isinstance(link, dict) and link.get("target_node_id") in final_node_ids]
    log.info(f"Validated {len(validated_nodes_list)} final nodes.")

    # --- Update Ontology Metadata ---
    ontology.title = llm_summary_details.get("title", ontology.title or "Project Document Ontology")
    ontology.overview = llm_summary_details.get("overview", ontology.overview or "Analysis of project documents.")
    raw_themes = llm_summary_details.get("global_themes", ontology.global_themes)
    ontology.global_themes = [str(t) for t in raw_themes if isinstance(t, str)] if isinstance(raw_themes, list) else ontology.global_themes
    raw_concepts = llm_summary_details.get("key_documents", ontology.key_concepts)
    ontology.key_concepts = [str(c) for c in raw_concepts if isinstance(c, str)] if isinstance(raw_concepts, list) else ontology.key_concepts
    ontology.project_graph_nodes = validated_nodes_list
    final_doc_filenames_in_graph = {node.name for node in validated_nodes_list if node.node_id.startswith("doc::")}
    updated_doc_list = sorted(list(existing_files.union(final_doc_filenames_in_graph)))
    ontology.documents = updated_doc_list; ontology.document_count = len(ontology.documents)
    log.info(f"Ontology update complete. State: {ontology.document_count} docs, {len(ontology.project_graph_nodes)} nodes.")

    # --- 7. Persist Ontology & Return ---
    try:
        save_ontology_json(ontology, path) # Uses external helper
        log.info(f"Successfully saved updated ontology to '{path}'.")
    except Exception as e: log.error(f"Failed to save ontology to '{path}': {e}", exc_info=True)

    return ontology # Return the final object (or None if aborted)

def load_ontology_json(path: pathlib.Path) -> Optional[ProjectOntology]:
    """Loads ontology from JSON file, validates with Pydantic."""
    if not path.exists():
        return None
    try:
        with path.open("r", encoding="utf-8") as f:
            data = json.load(f)
        return ProjectOntology.model_validate(data)
    except (json.JSONDecodeError, ValidationError, IOError) as e:
        logging.error(f"Error loading or validating ontology from {path}: {e}", exc_info=True)
        return None

def save_ontology_json(ont: ProjectOntology, path: pathlib.Path):
    """Saves Pydantic ontology model to JSON file."""
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("w", encoding="utf-8") as f:
            json.dump(ont.model_dump(mode="json"), f, ensure_ascii=False, indent=2)
    except (IOError, TypeError) as e:
        logging.error(f"Error saving ontology to {path}: {e}", exc_info=True)


def _merge_project_links_into_node_dicts_properties(link_data_list: List[Dict], node_map: Dict[str, Dict], max_total_links: int = 15) -> Tuple[int, int]:
    """
    Merges INTER-DOCUMENT link data into the 'linked_nodes' lists of node
dictionaries stored in node_map. Uses properties['source'] for document context.
"""
    links_added_count = 0
    links_skipped_count = 0
    if not isinstance(link_data_list, list):
        logging.warning("Link merging received non-list input for link data.")
        return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict): continue

        source_id = link_data.get("source_node_id")
        target_id = link_data.get("target_node_id")
        desc = link_data.get("relationship_description")

        if not all([source_id, target_id, desc]):
            links_skipped_count += 1; continue

        source_node_dict = node_map.get(source_id)
        target_node_dict = node_map.get(target_id)

        if not source_node_dict or not target_node_dict:
            links_skipped_count += 1; continue

        # *** Get document context from PROPERTIES ***
        source_doc = source_node_dict.get('properties', {}).get('source', 'source_unknown')
        target_doc = target_node_dict.get('properties', {}).get('source', 'target_unknown')

        if source_doc == target_doc:
            # This link is INTRA-document, skip it as we only want INTER here
            links_skipped_count += 1; continue

        try:
            link_to_add = {
                "target_node_id": target_id,
                "relationship_description": desc,
                "relationship_keywords": link_data.get("relationship_keywords", []),
                "relationship_strength": float(link_data.get("relationship_strength", 5.0))
            }
            if not isinstance(link_to_add["relationship_keywords"], list): link_to_add["relationship_keywords"] = []
        except (ValueError, TypeError):
            links_skipped_count += 1; continue

        # Add Forward Link (directly modifying the dictionary in node_map)
        source_node_dict.setdefault('linked_nodes', [])
        if isinstance(source_node_dict['linked_nodes'], list):
            if len(source_node_dict['linked_nodes']) < max_total_links and not any(lnk.get('target_node_id') == target_id for lnk in source_node_dict['linked_nodes'] if isinstance(lnk, dict)):
                source_node_dict['linked_nodes'].append(link_to_add)
                links_added_count += 1
            else: links_skipped_count += 1
        else: links_skipped_count += 1 # linked_nodes is not a list

        # Optional: Add Reverse Link (directly modifying the dictionary in node_map)
        target_node_dict.setdefault('linked_nodes', [])
        if isinstance(target_node_dict['linked_nodes'], list):
            if len(target_node_dict['linked_nodes']) < max_total_links and not any(lnk.get('target_node_id') == source_id for lnk in target_node_dict['linked_nodes'] if isinstance(lnk, dict)):
                reverse_desc = f"Referenced by '{source_node_dict.get('name', source_id)}' from document '{source_doc}' ({desc})"
                reverse_link_to_add = {
                    "target_node_id": source_id,
                    "relationship_description": reverse_desc,
                    "relationship_keywords": link_to_add["relationship_keywords"],
                    "relationship_strength": link_to_add["relationship_strength"]
                }
                target_node_dict['linked_nodes'].append(reverse_link_to_add)
                # links_added_count += 1 # Count reverse separately?
            # else: links_skipped_count += 1 # Skip reverse if limit/duplicate
        # else: links_skipped_count += 1 # linked_nodes is not a list

    return links_added_count, links_skipped_count

###############################
# DOCUMENT PROCESSING FUNCTIONS
###############################

# async def process_single_page_with_semaphore(semaphore, client, page_info: dict, uploaded_files):
#     """
#     Process a single page using semaphore for rate limiting.
#     Returns a dictionary: {"info": page_info, "result": PageContent object}.
#     Handles potential errors during processing. Uses 'source_identifier'.
#     """
#     # Get source_identifier and page_num early for use in logic and error reporting
#     source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
#     page_num = page_info.get("page_num", 0)

#     async with semaphore:
#         result_page_obj = None
#         try:
#             # Call process_single_pdf_page, which returns a DICTIONARY
#             # Pass the whole page_info which now contains source_identifier
#             page_dict = await process_single_pdf_page(client, page_info, uploaded_files)

#             page_dict.setdefault("page_number", page_num)
#             logging.debug(f"Attempting Pydantic validation for page {page_num} dict: {list(page_dict.keys())}")
#             result_page_obj = PageContent(**page_dict)
#             # Use source_identifier in log
#             logging.info(f"Successfully validated PageContent for page {page_num} of {source_identifier}")

#         except pydantic.ValidationError as ve:
#             # Use source_identifier in log/error
#             logging.error(f"Pydantic validation error creating PageContent for page {page_num} of {source_identifier} from returned dictionary: {ve}")
#             error_page_obj = create_error_page(
#                 page_num=page_num,
#                 error_msg=f"Pydantic Validation Error in '{source_identifier}': {ve}", # Include identifier in msg
#                 validation_errors=ve.errors()
#             )
#             result_page_obj = error_page_obj # Already a PageContent object

#         except Exception as e:
#             # Use source_identifier in log/error
#             logging.error(f"Unexpected error in semaphore task for page {page_num} of {source_identifier}: {str(e)}")
#             error_page_obj = create_error_page(
#                 page_num=page_num,
#                 error_msg=f"Unexpected processing error in semaphore task for '{source_identifier}': {str(e)}" # Include identifier in msg
#             )
#             result_page_obj = error_page_obj

#         # Return the result packaged with the original info
#         logging.debug(f"Finished processing for page {page_num} of {source_identifier}")
#         return {"info": page_info, "result": result_page_obj}


# async def process_single_pdf_page(client, page_info: dict, uploaded_files) -> dict:
    # """
    # Process a single PDF page image using Gemini with fallback options, returning extracted data as a dictionary.
    # Uses 'source_identifier' from page_info.
    # """
    # page_num = page_info.get("page_num", 0)
    # source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
    # page_dict = {}

    # try:
    #     if "image_b64" not in page_info or not page_info["image_b64"]:
    #          st.warning(f"Missing image data for page {page_num} of {source_identifier}.")
    #          error_page_obj = create_error_page(page_num=page_num, error_msg="Missing or empty image data (image_b64)")
    #          return error_page_obj.model_dump()

    #     if "page_num" not in page_info:
    #          st.warning(f"Missing page number for a page in {source_identifier}.")
    #          error_page_obj = create_error_page(page_num=0, error_msg="Missing page number information")
    #          return error_page_obj.model_dump()

    #     img_data = base64.b64decode(page_info["image_b64"])
    #     image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")

    #     # Call Gemini with fallback strategy, passing source_identifier
    #     # Assumes process_page_with_fallback accepts source_identifier now
    #     page_content_obj = await process_page_with_fallback(
    #         client, image_part, page_num, source_identifier # <-- Pass identifier
    #     )

    #     page_dict = page_content_obj.model_dump()
    #     return page_dict

    # except Exception as e:
    #     st.error(f"Error during Gemini processing or data handling for page {page_num} of {source_identifier}: {str(e)}")
    #     error_page_obj = create_error_page(
    #         page_num=page_num,
    #         error_msg=f"Gemini API or processing error in '{source_identifier}': {str(e)}",
    #     )
    #     page_dict = error_page_obj.model_dump()
    #     return page_dict
    
# --- New/Renamed Function for PDF Page Image Processing (Stage 1) ---
async def process_single_pdf_page_image(semaphore: asyncio.Semaphore, client: genai.Client, page_info: dict) -> PageContent:
    """
    Processes a single PDF page image using Gemini (Stage 1 extraction only).
    Uses semaphore for rate limiting. Returns a PageContent object.
    Handles potential errors during processing. Uses 'source_identifier'.
    """
    source_identifier = page_info.get("source_identifier", "Unknown_Source")
    page_num = page_info.get("page_num", 0)
    result_page_obj = None

    async with semaphore:
        try:
            if "image_b64" not in page_info or not page_info["image_b64"]:
                 logging.warning(f"Missing image data for page {page_num} of {source_identifier}.")
                 result_page_obj = create_error_page(page_num=page_num, error_msg="Missing or empty image data (image_b64)")
            else:
                img_data = base64.b64decode(page_info["image_b64"])
                image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")

                # Call Gemini (Stage 1 extraction - NO subsections yet)
                # Uses process_page_with_fallback internally for robustness
                # Assumes extract_page_content_from_memory performs Stage 1 extraction
                result_page_obj = await extract_page_content_from_memory(
                    client, image_part, page_num, source_identifier
                )

                # Basic check if extraction failed severely (e.g. returned error type)
                # This check depends on how create_error_page marks errors
                if isinstance(result_page_obj, PageContent) and \
                   hasattr(result_page_obj, 'subsections') and \
                   any(getattr(s,'title', '') == "Processing Error" for s in getattr(result_page_obj, 'subsections', [])):
                     logging.warning(f"Initial structure extraction seems to have failed for page {page_num} of '{source_identifier}'.")
                     # Keep the error page object

                elif not isinstance(result_page_obj, PageContent):
                     logging.error(f"Unexpected type {type(result_page_obj)} returned from page processing for page {page_num} of {source_identifier}")
                     result_page_obj = create_error_page(page_num=page_num, error_msg=f"Unexpected type from processing: {type(result_page_obj)}")

        except pydantic.ValidationError as ve:
            logging.error(f"Pydantic validation error creating PageContent for page {page_num} of {source_identifier}: {ve}")
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Pydantic Validation Error in '{source_identifier}': {ve}",
                validation_errors=ve.errors()
            )
        except Exception as e:
            logging.error(f"Unexpected error in page processing task for page {page_num} of {source_identifier}: {str(e)}", exc_info=True)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Unexpected processing error in page task for '{source_identifier}': {str(e)}"
            )

    # Ensure we always return a PageContent object, even if it's an error one
    if result_page_obj is None:
        logging.error(f"Processing yielded None for page {page_num} of {source_identifier}")
        result_page_obj = create_error_page(page_num=page_num, error_msg=f"Processing yielded None for page {page_num} of {source_identifier}")

    # logging.debug(f"Finished initial processing for page {page_num} of {source_identifier}")
    return result_page_obj


async def finalize_document(client: genai.Client, source_identifier: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Merges cut-off subsections.
    2. Extracts chapters based on subsections.
    3. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    4. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    5. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{source_identifier}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": source_identifier, # Store the identifier here
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }

    # Attempt to dump incoming pages immediately for error reporting, handle failures
    try:
        initial_pages_dump = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in pages]
        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed to initially dump pages for '{source_identifier}' for error reporting: {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{source_identifier}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{source_identifier}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result

        # --- Step 1: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Merging subsections for '{source_identifier}'.")
        merged_pages = await merge_cutoff_subsections(pages) # Operates on PageContent objects
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{source_identifier}'. Using original pages.")
             merged_pages = pages

        # --- Step 2: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Extracting chapters for '{source_identifier}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{source_identifier}'. Cannot proceed with relationship/summary analysis.")
             # Return result with merged pages but no summary/chapter info
             final_result["raw_extracted_content"]["pages"] = [p.model_dump(mode='json') for p in merged_pages]
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = { # Basic placeholder summary
                  "title": f"Processing Incomplete: {source_identifier}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result

        # --- Step 3: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Analyzing relationships within {len(chapters)} chapters for '{source_identifier}'.")
        # This loop modifies chapter.qdrant_nodes in place
        analysis_tasks = [analyze_concept_entity_relationships(client, chapter) for chapter in chapters]
        await asyncio.gather(*analysis_tasks, return_exceptions=True) # Run in parallel, log errors if any task fails
        # Check results? For now, we assume modification happened or warnings were logged internally.

        # --- Step 4: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
        # This function modifies chapter.qdrant_nodes again (adds inter-chapter links) AND returns the summary dict
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)

        # --- Step 5: Assemble Final Result ---
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated before dumping
        summary_dict['chapters'] = [ch.model_dump(mode='json') for ch in chapters if isinstance(ch, Chapter)]

        # Add document-level relationships placeholder (could be populated later from qdrant_nodes)
        summary_dict['entity_relationships'] = []

        # Dump the final state of merged pages
        final_pages_as_dicts = [p.model_dump(mode='json') for p in merged_pages if isinstance(p, PageContent)]

        final_result = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result

    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{source_identifier}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {source_identifier}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": [], "entity_relationships": []
            }
        # Keep the initial pages dump in the error case
        return final_result


async def extract_pages_from_pdf_bytes(pdf_bytes, source_identifier: str) -> List[Dict]:
    """Extract pages from PDF bytes as image info dicts."""
    loop = asyncio.get_running_loop()
    pages_info = []
    try:
        pdf_stream = io.BytesIO(pdf_bytes)
        # Run synchronous PyMuPDF code in an executor thread
        def sync_extract():
            doc = pymupdf.open(stream=pdf_stream, filetype="pdf")
            extracted = []
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    # Reduce DPI slightly to potentially save memory/time if quality allows
                    pix = page.get_pixmap(dpi=150, alpha=False)
                    img_bytes = pix.tobytes("jpeg", jpg_quality=80) # Slightly lower quality
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    extracted.append({
                        "page_num": page_num + 1,
                        "image_b64": img_b64,
                        "source_identifier": source_identifier # Use the passed identifier
                    })
                except Exception as page_e:
                     logging.warning(f"Error extracting page {page_num+1} from {source_identifier}: {page_e}")
            doc.close()
            return extracted

        pages_info = await loop.run_in_executor(None, sync_extract)
        return pages_info

    except Exception as e:
        logging.error(f"Error opening or processing PDF {source_identifier}: {str(e)}")
        # Return empty list, caller should handle this
        return []


# Non-PDF Document Processing Functions

async def process_single_document_memory(client: genai.Client, file_data: dict) -> List[PageContent]:
    """
    Processes a single non-PDF document from memory (Stage 1 extraction only).
    Calls specific handlers (Word, PPTX, Tabular, Text) which perform the
    initial structure extraction.
    Returns a list of PageContent objects (one per page/chunk).
    """
    source_identifier = file_data.get("identifier", file_data.get("name", "Unknown_Source")) # Get name robustly
    file_type = file_data.get("type", "").lower()
    file_content = file_data.get("content")
    processed_pages: List[PageContent] = []

    logging.info(f"Starting initial page-level processing for non-PDF: '{source_identifier}' (type: {file_type})")

    if not file_content:
        logging.error(f"No content found for file '{source_identifier}'.")
        return [create_error_page(1, f"Source '{source_identifier}' has no content.")]
    if not file_type:
        logging.error(f"Cannot determine file type for '{source_identifier}'.")
        return [create_error_page(1, f"Unknown file type for '{source_identifier}'.")]

    try:
        # Dispatch based on file type
        if file_type in ["xlsx", "xls"]:
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "excel")
        elif file_type == "csv":
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "csv")
        elif file_type == "docx":
            processed_pages = await process_word_document(client, file_content, source_identifier)
        elif file_type == "pptx":
            processed_pages = await process_pptx_document(client, file_content, source_identifier)
        # Ensure all expected text types are handled
        elif file_type in ["txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]:
            # Pass the dict containing name, content, type
            processed_pages = await process_text_document(client, file_data)
        else:
            logging.warning(f"Unsupported file type '{file_type}' for '{source_identifier}'.")
            processed_pages = [create_error_page(1, f"Unsupported file type: {file_type}")]

        if not processed_pages:
            logging.warning(f"Initial processing for '{source_identifier}' resulted in no pages.")
            # Return an error page to indicate failure for this doc
            processed_pages = [create_error_page(1, f"No content pages generated during initial processing for '{source_identifier}'.")]

        logging.info(f"Finished initial page-level processing for non-PDF: '{source_identifier}'. Generated {len(processed_pages)} PageContent objects.")
        return processed_pages

    except Exception as e:
        logging.error(f"Core initial processing failed for document '{source_identifier}': {str(e)}", exc_info=True)
        return [create_error_page(1, f"Core initial processing error for '{source_identifier}': {e}")]

# --- NEW: Workflow function for a single document ---
async def process_and_finalize_single_document(
    client: genai.Client,
    file_data: dict,
    page_api_semaphore: asyncio.Semaphore, # Pass semaphore for page-level API calls
    # Optional: Add another semaphore if finalize_document also needs rate limiting
    # finalize_semaphore: asyncio.Semaphore
) -> Dict:
    """
    Handles the complete processing pipeline for a single document asynchronously.
    1. Initial page/chunk extraction (Stage 1)
    2. Finalization (Stage 2: subsections, chapters, relationships, summary)
    Returns the final dictionary structure produced by finalize_document.
    """
    function_name = "process_and_finalize_single_document"
    source_identifier = file_data.get("name", "Unknown_Document")
    file_type = file_data.get("type", "").lower()
    logging.info(f"[{function_name}] Starting workflow for '{source_identifier}' (type: {file_type})")

    initial_pages: List[PageContent] = []
    final_result_dict: Optional[Dict] = None

    try:
        # --- Stage 1: Initial Page/Chunk Extraction ---
        stage1_start = time.time()
        if file_type == 'pdf':
            pdf_content = file_data.get("content")
            if not pdf_content:
                 logging.error(f"[{function_name}] Missing PDF content for '{source_identifier}'.")
                 raise ValueError("Missing PDF content")

            # 1a. Extract page images
            logging.debug(f"[{function_name}] Extracting PDF page images for '{source_identifier}'")
            page_infos = await extract_pages_from_pdf_bytes(pdf_content, source_identifier)

            if not page_infos:
                 logging.warning(f"[{function_name}] No pages extracted from PDF '{source_identifier}'.")
                 # Create a single error page for this document
                 initial_pages = [create_error_page(1, f"Failed to extract any pages from PDF: {source_identifier}")]
            else:
                # 1b. Process page images concurrently (Stage 1 LLM)
                logging.debug(f"[{function_name}] Processing {len(page_infos)} PDF pages for '{source_identifier}' (Stage 1)")
                page_tasks = []
                valid_page_info_indices = [] # Keep track of indices corresponding to tasks
                for idx, pi in enumerate(page_infos):
                    if isinstance(pi, dict) and pi.get("image_b64"):
                        page_tasks.append(process_single_pdf_page_image(page_api_semaphore, client, pi))
                        valid_page_info_indices.append(idx)
                    else:
                        logging.warning(f"[{function_name}] Skipping invalid page info at index {idx} for PDF '{source_identifier}'")


                if page_tasks:
                     # Gather results - these are PageContent objects
                     gathered_results = await asyncio.gather(*page_tasks, return_exceptions=True)

                     # Filter out potential exceptions, replace with error pages
                     processed_pages = []
                     for i, res in enumerate(gathered_results):
                         # Get original page info index for context
                         original_info_index = valid_page_info_indices[i]
                         page_num_ctx = page_infos[original_info_index].get('page_num', original_info_index + 1)

                         if isinstance(res, PageContent):
                             processed_pages.append(res)
                         elif isinstance(res, Exception):
                             logging.error(f"[{function_name}] Error processing PDF page {page_num_ctx} of '{source_identifier}': {res}")
                             processed_pages.append(create_error_page(page_num_ctx, f"PDF page processing error: {res}"))
                         else:
                             logging.error(f"[{function_name}] Unexpected result type {type(res)} processing PDF page {page_num_ctx} of '{source_identifier}'")
                             processed_pages.append(create_error_page(page_num_ctx, f"Unexpected PDF page result type: {type(res)}"))
                     initial_pages = processed_pages
                     # Sort pages by page number
                     initial_pages.sort(key=lambda p: getattr(p, 'page_number', float('inf')))
                else:
                     logging.warning(f"[{function_name}] No valid page processing tasks created for PDF '{source_identifier}'.")
                     initial_pages = [create_error_page(1, f"No valid pages found or processed in PDF: {source_identifier}")]

        else: # Non-PDF types
            logging.debug(f"[{function_name}] Processing non-PDF '{source_identifier}' (Stage 1)")
            # This function handles different non-PDF types internally and calls Stage 1 LLM
            initial_pages = await process_single_document_memory(client, file_data)

        stage1_duration = time.time() - stage1_start
        logging.info(f"[{function_name}] Stage 1 (Initial Extraction) for '{source_identifier}' completed in {stage1_duration:.2f}s. Pages generated: {len(initial_pages)}")

        # Check if Stage 1 produced *any* valid PageContent objects (even if they contain errors internally)
        if not isinstance(initial_pages, list) or not initial_pages:
            logging.error(f"[{function_name}] Stage 1 failed catastrophically for '{source_identifier}'. No PageContent objects returned. Skipping finalization.")
            final_result_dict = {
                "raw_extracted_content": { "filename": source_identifier, "pages": [], "summary": None, "error": "Initial page extraction (Stage 1) failed completely." }
            }
            return final_result_dict

        # Further check: Did any page actually get usable text? (More robust check)
        has_content = any(isinstance(p, PageContent) and (p.raw_text or p.subsections) for p in initial_pages)
        if not has_content:
            logging.error(f"[{function_name}] Stage 1 produced PageContent objects but none contained usable text/subsections for '{source_identifier}'. Skipping finalization.")
            # Dump the potentially error-filled pages
            dumped_pages = []
            for p in initial_pages:
                 try: dumped_pages.append(p.model_dump(mode='json'))
                 except Exception: dumped_pages.append({"error": "Failed to dump page object"})

            final_result_dict = {
                "raw_extracted_content": { "filename": source_identifier, "pages": dumped_pages, "summary": None, "error": "Initial page extraction (Stage 1) failed to produce usable content." }
            }
            return final_result_dict # Return early


        # --- Stage 2: Finalization ---
        stage2_start = time.time()
        logging.debug(f"[{function_name}] Starting Stage 2 (Finalization) for '{source_identifier}'")

        # Use the existing finalize_document function
        # Optional: Add semaphore here if needed: async with finalize_semaphore: ...
        final_result_dict = await finalize_document(client, source_identifier, initial_pages)

        stage2_duration = time.time() - stage2_start
        logging.info(f"[{function_name}] Stage 2 (Finalization) for '{source_identifier}' completed in {stage2_duration:.2f}s.")

        # Log if finalization itself reported an error
        if isinstance(final_result_dict, dict) and final_result_dict.get("raw_extracted_content", {}).get("error"):
            logging.warning(f"[{function_name}] Finalization step reported an error for '{source_identifier}': {final_result_dict['raw_extracted_content']['error']}")

        return final_result_dict

    except Exception as e:
        logging.error(f"[{function_name}] Unhandled error during workflow for '{source_identifier}': {e}", exc_info=True)
        # Return an error dictionary
        # Attempt to dump initial pages if available
        dumped_initial_pages = []
        if initial_pages:
             for p in initial_pages:
                  try: dumped_initial_pages.append(p.model_dump(mode='json'))
                  except Exception: dumped_initial_pages.append({"error": "Failed to dump initial page object"})

        return {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": dumped_initial_pages,
                "summary": None,
                "error": f"Workflow failed unexpectedly: {str(e)}"
            }
        }

    
async def finalize_document(client: genai.Client, source_identifier: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Extracts subsections from raw_text for each page. <--- NEW STEP
    2. Merges cut-off subsections.
    3. Extracts chapters based on subsections.
    4. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    5. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    6. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{source_identifier}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": source_identifier, # Store the identifier here
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }
    # Attempt to dump incoming pages immediately for error reporting, handle failures
    initial_pages_dump = []
    try:
        # Dump safely, handle potential model errors during dump
        for i, p in enumerate(pages):
             if isinstance(p, PageContent):
                 try: initial_pages_dump.append(p.model_dump(mode='json'))
                 except Exception as dump_err: initial_pages_dump.append({"page_number": getattr(p, 'page_number', i+1), "error": f"Initial model dump failed: {dump_err}"})
             else: initial_pages_dump.append(p) # Keep as is if not PageContent obj

        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed during initial page dump for '{source_identifier}': {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{source_identifier}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result # Return early with initial page dump

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{source_identifier}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result # Return early

        # --- Step 1: Extract Subsections from Raw Text (NEW STEP) ---
        logging.debug(f"[{function_name}] Step 1: Extracting subsections from raw text for {len(pages)} pages of '{source_identifier}'.")
        subsection_extraction_tasks = []
        pages_needing_subsections = [] # Keep track of pages we are processing

        # Create tasks FIRST to run them concurrently
        for page in pages:
             # Only process if raw_text exists and subsections are currently empty
             if isinstance(page, PageContent) and page.raw_text and not page.subsections:
                 task = extract_subsections_from_text(client, page.raw_text, page.page_number, source_identifier)
                 subsection_extraction_tasks.append(task)
                 pages_needing_subsections.append(page) # Store reference to the page object
             elif not isinstance(page, PageContent):
                  logging.warning(f"[{function_name}] Item in pages list is not PageContent object for '{source_identifier}', page {getattr(page, 'page_number', '?')}. Skipping subsection extraction for it.")

             # else: subsections already exist or no raw text, skip

        # Run tasks concurrently
        if subsection_extraction_tasks:
             logging.info(f"[{function_name}] Running {len(subsection_extraction_tasks)} subsection extraction tasks for '{source_identifier}'.")
             results = await asyncio.gather(*subsection_extraction_tasks, return_exceptions=True)
             logging.info(f"[{function_name}] Finished subsection extraction tasks for '{source_identifier}'.")

             # Assign results back to the corresponding page objects
             if len(results) == len(pages_needing_subsections):
                 for i, result in enumerate(results):
                     page_to_update = pages_needing_subsections[i] # Get the correct page object
                     if isinstance(result, Exception):
                         logging.error(f"[{function_name}] Subsection extraction failed for page {page_to_update.page_number} of '{source_identifier}': {result}")
                         page_to_update.subsections = [] # Ensure empty list on error
                     elif isinstance(result, list): # Expected result is List[Subsection]
                         page_to_update.subsections = result # Assign generated subsections
                     else:
                         logging.error(f"[{function_name}] Unexpected result type {type(result)} for subsection extraction on page {page_to_update.page_number} of '{source_identifier}'.")
                         page_to_update.subsections = []
             else:
                  # This indicates a logic error in tracking tasks/pages
                  logging.error(f"[{function_name}] Mismatch between pages needing processing ({len(pages_needing_subsections)}) and task results ({len(results)}) for '{source_identifier}'. Subsections may be missing.")
                  # Attempt to assign based on index anyway, but it might be wrong
                  for idx, page in enumerate(pages_needing_subsections):
                     if idx < len(results):
                         result = results[idx]
                         if isinstance(result, list): page.subsections = result
                         else: page.subsections = []
                     else: page.subsections = [] # No result for this page

        else:
             logging.info(f"[{function_name}] No pages required subsection extraction for '{source_identifier}'.")


        # --- Step 2: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Step 2: Merging subsections for '{source_identifier}'.")
        # Pass the 'pages' list which now has potentially populated subsections
        merged_pages = await merge_cutoff_subsections(pages)
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{source_identifier}'. Using original pages.")
             merged_pages = pages # Fallback to original pages (already updated with subsections)


        # --- Step 3: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Step 3: Extracting chapters for '{source_identifier}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{source_identifier}'. Cannot proceed further.")
             final_pages_as_dicts = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in merged_pages] # Dump state after merge/subsection attempts
             final_result["raw_extracted_content"]["pages"] = final_pages_as_dicts
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = {
                  "title": f"Processing Incomplete: {source_identifier}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result


        # --- Step 4: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Step 4: Analyzing intra-chapter relationships for {len(chapters)} chapters in '{source_identifier}'.")
        # This loop modifies chapter.qdrant_nodes in place
        intra_chapter_tasks = [
            analyze_concept_entity_relationships(client, chapter, source_identifier) # <-- ADD source_identifier HERE
            for chapter in chapters
        ]
        intra_results = await asyncio.gather(*intra_chapter_tasks, return_exceptions=True)

        # Log any exceptions from intra-chapter analysis
        for i, res in enumerate(intra_results):
            if isinstance(res, Exception):
                 logging.error(f"[{function_name}] Intra-chapter analysis failed for chapter {chapters[i].chapter_id} of '{source_identifier}': {res}")


        # --- Step 5: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Step 5: Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)


        # --- Step 6: Assemble Final Result ---
        logging.debug(f"[{function_name}] Step 6: Assembling final result for '{source_identifier}'.")
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated (even if empty list) before dumping
        dumped_chapters = []
        for ch in chapters:
            if isinstance(ch, Chapter):
                try: dumped_chapters.append(ch.model_dump(mode='json'))
                except Exception as ch_dump_err:
                     logging.error(f"[{function_name}] Failed to dump chapter {ch.chapter_id}: {ch_dump_err}")
                     dumped_chapters.append({"chapter_id": ch.chapter_id, "title": ch.title, "error": "Failed to serialize chapter"})
            else: logging.warning(f"[{function_name}] Invalid chapter object found during final assembly.")
        summary_dict['chapters'] = dumped_chapters


        # Dump the final state of pages (use merged_pages)
        final_pages_as_dicts = []
        for p in merged_pages:
             if isinstance(p, PageContent):
                 try: final_pages_as_dicts.append(p.model_dump(mode='json'))
                 except Exception as pg_dump_err:
                      logging.error(f"[{function_name}] Failed to dump page {p.page_number}: {pg_dump_err}")
                      final_pages_as_dicts.append({"page_number": p.page_number, "error": "Failed to serialize page"})
             else: final_pages_as_dicts.append(p)


        final_result = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result

    # --- Error Handling ---
    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{source_identifier}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {source_identifier}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": []
                # Removed entity_relationships placeholder here too
            }
        # Keep the initial pages dump in the error case if available
        return final_result



#-------------------------------------
# HELPER: Create Markdown Table
#-------------------------------------
def create_markdown_table(headers, data):
    """Create a markdown table string from headers and list of lists/dicts."""
    if not data or not headers: return "<!-- Empty Table -->"
    header_row = "| " + " | ".join(map(str, headers)) + " |"
    separator_row = "| " + " | ".join("---" * len(headers)) + " |"
    data_rows = []
    for row in data:
        if isinstance(row, dict):
            values = [str(row.get(h, '')) for h in headers]
        elif isinstance(row, (list, tuple)):
            # Ensure row has same length as headers, padding if necessary
            values = [str(row[i]) if i < len(row) else '' for i in range(len(headers))]
        else: # Handle unexpected row types
             values = [str(row)] + [''] * (len(headers) - 1)
        data_rows.append("| " + " | ".join(values) + " |")
    return "\n".join([header_row, separator_row] + data_rows)

#-------------------------------------
# REWRITTEN NON-PDF PROCESSORS
#-------------------------------------

async def process_tabular_data(client: genai.Client, file_content: bytes, source_identifier: str, filetype: str) -> List[PageContent]:
    """Processes Excel/CSV by converting sheets/file to Markdown text and extracting structure (Stage 1)."""
    processed_pages = []
    page_counter = 1 # For sheets or single CSV page

    try:
        if filetype == "excel":
            try:
                 excel_data = pd.read_excel(io.BytesIO(file_content), sheet_name=None, engine='openpyxl')
                 if not excel_data:
                     raise ValueError("Excel file is empty or unreadable.")

                 for sheet_name, df in excel_data.items():
                     logging.info(f"Processing Excel sheet '{sheet_name}' for '{source_identifier}' as page {page_counter}")
                     if df.empty:
                          logging.warning(f"Sheet '{sheet_name}' in '{source_identifier}' is empty. Skipping.")
                          page_counter += 1
                          continue
                     markdown_text = f"# Sheet: {sheet_name}\n\n" + df.to_markdown(index=False)
                     # --- Calls Stage 1 Extraction ---
                     page_content = await extract_structure_from_text_content(client, markdown_text, page_counter, f"{source_identifier} (Sheet: {sheet_name})")
                     processed_pages.append(page_content)
                     page_counter += 1

            except Exception as e:
                logging.error(f"Error reading Excel file '{source_identifier}' with pandas: {e}", exc_info=True)
                processed_pages.append(create_error_page(1, f"Error reading Excel file '{source_identifier}': {e}"))

        elif filetype == "csv":
            logging.info(f"Processing CSV file '{source_identifier}' as page 1")
            try:
                df = pd.read_csv(io.BytesIO(file_content))
                if df.empty:
                     raise ValueError("CSV file is empty or unreadable.")
                markdown_text = "# CSV Data\n\n" + df.to_markdown(index=False)
                # --- Calls Stage 1 Extraction ---
                page_content = await extract_structure_from_text_content(client, markdown_text, 1, source_identifier)
                processed_pages.append(page_content)
            except Exception as e:
                 logging.error(f"Error reading CSV file '{source_identifier}' with pandas: {e}", exc_info=True)
                 processed_pages.append(create_error_page(1, f"Error reading CSV file '{source_identifier}': {e}"))
        else:
            raise ValueError(f"Unsupported tabular file type: {filetype}")

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing tabular data for '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process tabular file '{source_identifier}': {e}")]

async def process_word_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes Word docx file by extracting text and applying structure extraction (Stage 1)."""
    processed_pages = []
    try:
        logging.info(f"Processing Word document '{source_identifier}'")
        doc = docx.Document(io.BytesIO(file_content))
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        if not full_text.strip():
             logging.warning(f"Word document '{source_identifier}' contains no text.")
             return [create_error_page(1, "Word document is empty or contains no text.")]

        logging.debug(f"Extracted {len(full_text)} chars from '{source_identifier}'. Processing as single page.")
        # --- Calls Stage 1 Extraction ---
        page_content = await extract_structure_from_text_content(client, full_text, 1, source_identifier)
        processed_pages.append(page_content)

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing Word document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process Word document '{source_identifier}': {e}")]

async def process_pptx_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes PowerPoint pptx file by extracting text per slide and applying structure extraction (Stage 1)."""
    processed_pages = []
    try:
        logging.info(f"Processing PowerPoint document '{source_identifier}'")
        ppt = Presentation(io.BytesIO(file_content))

        if not ppt.slides:
             logging.warning(f"PowerPoint document '{source_identifier}' contains no slides.")
             return [create_error_page(1, "PowerPoint document has no slides.")]

        for i, slide in enumerate(ppt.slides, 1):
            slide_texts = []
            title = f"Slide {i}"
            has_title_shape = False
            try:
                for shape in slide.shapes:
                     if hasattr(shape, "is_title") and shape.is_title and shape.has_text_frame and shape.text.strip():
                         title = shape.text.strip()
                         slide_texts.append(f"# {title}")
                         has_title_shape = True
                         break
            except Exception as shape_error:
                 logging.warning(f"Error accessing title shape on slide {i} of '{source_identifier}': {shape_error}")

            try:
                for shape in slide.shapes:
                     if shape.has_text_frame and shape.text and not (has_title_shape and hasattr(shape, "is_title") and shape.is_title):
                         for para in shape.text_frame.paragraphs:
                              if para.text.strip():
                                   prefix = "- " * (para.level + 1)
                                   slide_texts.append(f"{prefix}{para.text.strip()}")
            except Exception as shape_error:
                 logging.warning(f"Error accessing text shapes on slide {i} of '{source_identifier}': {shape_error}")

            try:
                 if slide.has_notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip():
                     slide_texts.append("\n## Speaker Notes")
                     slide_texts.append(slide.notes_slide.notes_text_frame.text.strip())
            except Exception as notes_error:
                logging.warning(f"Error accessing notes on slide {i} of '{source_identifier}': {notes_error}")

            slide_full_text = "\n".join(slide_texts)

            if not slide_full_text.strip():
                 logging.warning(f"Slide {i} of '{source_identifier}' contains no text content. Skipping.")
                 continue

            logging.debug(f"Processing slide {i} of '{source_identifier}'")
            # --- Calls Stage 1 Extraction ---
            page_content = await extract_structure_from_text_content(client, slide_full_text, i, f"{source_identifier} (Slide {i})")
            processed_pages.append(page_content)

        # Ensure at least one page object is returned if slides existed but were all empty/failed
        if ppt.slides and not processed_pages:
            logging.warning(f"PPTX '{source_identifier}' had slides but none yielded processable text.")
            return [create_error_page(1, "PPTX had slides but no text could be processed.")]

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing PowerPoint document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process PowerPoint document '{source_identifier}': {e}")]

async def process_text_document(client: genai.Client, file_data: dict) -> List[PageContent]:
    """Processes plain text, markdown, json, html, xml files by chunking and extracting structure (Stage 1)."""
    source_identifier = file_data.get("name", "Unknown_Text_File")
    file_type = file_data.get("type", "txt")
    processed_pages = []
    page_counter = 1

    try:
        logging.info(f"Processing {file_type} document '{source_identifier}'")
        try:
            text_content = file_data["content"].decode('utf-8')
        except UnicodeDecodeError:
             try:
                 text_content = file_data["content"].decode('latin-1')
             except Exception as decode_err:
                 logging.error(f"Failed to decode file '{source_identifier}': {decode_err}")
                 return [create_error_page(1, f"Failed to decode file '{source_identifier}'")]

        if not text_content.strip():
            logging.warning(f"Text document '{source_identifier}' is empty.")
            return [create_error_page(1, f"Text document '{source_identifier}' is empty.")]

        # Simple chunking by character count, aiming for rough page equivalents
        MAX_CHUNK_CHARS = 10000 # Adjust based on typical LLM context window and desired granularity
        text_chunks = [text_content[i:i + MAX_CHUNK_CHARS] for i in range(0, len(text_content), MAX_CHUNK_CHARS)]

        if not text_chunks:
             text_chunks = [text_content] # Ensure at least one chunk if splitting failed

        logging.info(f"Split '{source_identifier}' into {len(text_chunks)} chunks for processing.")

        for text_chunk in text_chunks:
            if not text_chunk.strip(): continue
            # --- Calls Stage 1 Extraction ---
            page_content = await extract_structure_from_text_content(client, text_chunk, page_counter, f"{source_identifier} (Chunk {page_counter})")
            processed_pages.append(page_content)
            page_counter += 1

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing text document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process text document '{source_identifier}': {e}")]


async def process_all_documents_async(file_data: List[Dict]) -> List[Dict]:
    """
    Process all documents concurrently, including asynchronous finalization per document.
    """
    # --- Setup ---
    try:
        api_key = get_gemini_api_key()
        # Ensure Helicone proxy details are correct or remove if not used
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                } if st.secrets.get("HELICONE_PROXY_URL") and st.secrets.get("HELICONE_API_KEY") else {} # Conditionally add headers
            })
        logging.info("Gemini client initialized for document processing.")
    except Exception as client_err:
        logging.error(f"Failed to initialize Gemini client: {client_err}")
        st.error(f"Failed to initialize Gemini client: {client_err}")
        return []

    final_doc_results = [] # Stores the FINAL dictionary results from each document's workflow
    original_doc_count = len(file_data)
    docs_to_process_count = 0 # Count valid files

    # --- Status Tracking Setup ---
    status_container = st.session_state.get("status_container")
    progress_bar = st.session_state.get("progress_bar")
    time_info = st.session_state.get("time_info")
    processing_status = {
        "active": True, "current_step": "Initializing...",
        "total_steps": original_doc_count, # Track documents as steps initially
        "current_step_num": 0, "step_progress": 0.0, # Not used per-step
        "start_time": time.time(), "step_start_time": time.time(),
        "estimated_time_remaining": None, "parallel_count": 1
    }
    # We no longer have distinct pipeline phases in the status update logic

    def update_doc_progress(completed_count, total_count, current_source_identifier=None):
        """Updates the processing status based on completed documents."""
        try:
            # Update total steps if some files were skipped early
            processing_status["total_steps"] = total_count
            processing_status["current_step_num"] = completed_count
            overall_progress = 0.0
            if total_count > 0:
                overall_progress = min(1.0, float(completed_count) / float(total_count))
            overall_progress = max(0.0, min(1.0, overall_progress))

            status_msg = f"Processing Documents ({completed_count}/{total_count})"
            # Display last completed document name might be more reliable than current
            if current_source_identifier:
                status_msg = f"Completed: {current_source_identifier} ({completed_count}/{total_count})"


            eta_msg = ""
            time_elapsed = time.time() - processing_status["start_time"]
            if overall_progress > 0.01:
                est_total_time = time_elapsed / overall_progress
                remaining = max(0, est_total_time - time_elapsed)
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                mins_remaining, secs_remaining = divmod(int(remaining), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: ~{mins_remaining}m {secs_remaining}s"
            else:
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: Calculating..."

            if status_container: status_container.update(label=status_msg, expanded=True)
            if progress_bar: progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
            if time_info: time_info.markdown(f"`{eta_msg}`")

        except Exception as ui_err:
             logging.warning(f"Failed to update UI status: {ui_err}")

    # --- Main Concurrent Processing Logic ---
    try:
        update_doc_progress(0, original_doc_count) # Initial status

        # --- Concurrency Limits ---
        # Limit for page-level API calls (within each doc's Stage 1)
        MAX_PARALLEL_PAGES = 100
        page_api_semaphore = asyncio.Semaphore(MAX_PARALLEL_PAGES)

        # Limit for how many documents process concurrently overall
        MAX_PARALLEL_DOCS = 10
        doc_workflow_semaphore = asyncio.Semaphore(MAX_PARALLEL_DOCS)

        # --- Create Tasks for each Document Workflow ---
        doc_tasks = []
        valid_file_infos = []
        for file_info in file_data:
            if not isinstance(file_info, dict) or "name" not in file_info or "content" not in file_info or not file_info["content"]:
                 st.warning(f"Skipping invalid file data item (missing name, content, or empty): {str(file_info.get('name', 'N/A'))}...")
                 continue # Skip this file

            valid_file_infos.append(file_info)
            docs_to_process_count += 1 # Increment count of actual files to process

            async def run_workflow_with_semaphore(file_info_arg):
                async with doc_workflow_semaphore:
                    # This allows MAX_PARALLEL_DOCS workflows to run concurrently
                    # Each workflow uses page_api_semaphore internally for its page processing
                    logging.debug(f"Acquired doc semaphore for: {file_info_arg.get('name')}")
                    result = await process_and_finalize_single_document(
                        client,
                        file_info_arg,
                        page_api_semaphore
                        # Add finalize_semaphore if needed
                    )
                    logging.debug(f"Releasing doc semaphore for: {file_info_arg.get('name')}")
                    return result

            doc_tasks.append(run_workflow_with_semaphore(file_info))

        # Update total steps based on valid files
        processing_status["total_steps"] = docs_to_process_count
        update_doc_progress(0, docs_to_process_count) # Update progress bar total

        # --- Execute and Gather Results ---
        completed_docs = 0
        if not doc_tasks:
             logging.warning("No valid documents found to process.")
             st.warning("No valid documents found to process.")
             processing_status["active"] = False
             if status_container:
                 try: status_container.update(label="⏹️ No valid documents found.", state="error", expanded=False)
                 except Exception: pass
             return []


        logging.info(f"Launching {len(doc_tasks)} document processing workflows concurrently (max {MAX_PARALLEL_DOCS})...")
        processing_status["parallel_count"] = MAX_PARALLEL_DOCS

        # Use asyncio.as_completed to process results as they finish and update progress
        for future in asyncio.as_completed(doc_tasks):
            last_completed_name = "Unknown"
            try:
                result_dict = await future # result_dict is the final dict from finalize_document or an error dict
                if isinstance(result_dict, dict) and "raw_extracted_content" in result_dict:
                     final_doc_results.append(result_dict)
                     last_completed_name = result_dict.get("raw_extracted_content", {}).get("filename", "Unknown")
                     # Log success or failure for this doc
                     if result_dict.get("raw_extracted_content", {}).get("error"):
                          logging.warning(f"Workflow for '{last_completed_name}' completed with error: {result_dict['raw_extracted_content']['error']}")
                     else:
                          logging.info(f"Workflow for '{last_completed_name}' completed successfully.")
                else:
                     # This case indicates a severe failure within the workflow task itself
                     logging.error(f"Workflow task returned unexpected data type: {type(result_dict)}")
                     last_completed_name = "Workflow Error"
                     final_doc_results.append({
                         "raw_extracted_content": {"filename": "Unknown_Workflow_Error", "pages": [], "summary": None, "error": f"Workflow task failed unexpectedly, returned type {type(result_dict)}"}
                     })

                completed_docs += 1
                update_doc_progress(completed_docs, docs_to_process_count, last_completed_name)

            except asyncio.CancelledError:
                 logging.warning("A document processing task was cancelled.")
                 completed_docs += 1 # Count it as "done" for progress calculation
                 update_doc_progress(completed_docs, docs_to_process_count, "Cancelled Task")
                 # Optionally add a specific marker for cancelled tasks if needed
                 final_doc_results.append({
                     "raw_extracted_content": {"filename": "Cancelled_Task", "pages": [], "summary": None, "error": "Processing was cancelled."}
                 })
            except Exception as e:
                 logging.error(f"Error awaiting document task result: {e}", exc_info=True)
                 completed_docs += 1 # Count error task as "done"
                 last_completed_name = "Errored Task"
                 update_doc_progress(completed_docs, docs_to_process_count, last_completed_name)
                 # Add a generic error dict to results
                 final_doc_results.append({
                     "raw_extracted_content": {"filename": "Unknown_Workflow_Error", "pages": [], "summary": None, "error": f"Workflow task failed: {str(e)}"}
                 })

        # Final status update
        update_doc_progress(completed_docs, docs_to_process_count)
        processing_status["active"] = False
        if status_container:
             final_state = "complete" if completed_docs == docs_to_process_count else "error"
             final_label = "✅ Document processing complete!" if final_state == "complete" else "⚠️ Processing finished with errors."
             try: status_container.update(label=final_label, state=final_state, expanded=False)
             except Exception as ui_err: logging.warning(f"Failed to update final status: {ui_err}")

    # --- Exception Handling & Finally Block ---
    except asyncio.CancelledError:
         logging.warning("Document processing orchestration was cancelled.")
         st.warning("Document processing was cancelled.")
         if status_container:
             try: status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
             except Exception: pass
         processing_status["active"] = False
         return [] # Return empty list on cancellation
    except Exception as e:
        logging.error(f"❌ An error occurred during the main document processing orchestration: {str(e)}", exc_info=True)
        st.error(f"❌ An error occurred during document processing: {str(e)}")
        if status_container:
            try: status_container.update(label=f"❌ Pipeline Error: {str(e)}", state="error", expanded=True)
            except Exception: pass
        processing_status["active"] = False
        # Return whatever was collected before the error
        logging.info(f"Returning {len(final_doc_results)} potentially partially processed documents due to pipeline error.")
        return final_doc_results
    finally:
        logging.info("Document processing orchestration finished or exited.")
        processing_status["active"] = False
        if "st" in locals() and hasattr(st, "session_state"):
            st.session_state.processing_active = False # Ensure flag is reset

    # Return the list of final document dictionaries
    logging.info(f"Successfully processed {len(final_doc_results)} documents out of {docs_to_process_count} attempted.")
    return final_doc_results


###############################
# UI DISPLAY FUNCTIONS (Updated)
###############################

def display_table(table):
    """Display table data (accepts model or dict)."""
    if not table: return
    table_data = table.model_dump() if hasattr(table, 'model_dump') else table

    title = table_data.get("title")
    content = table_data.get("table_content")

    if title:
        st.subheader(f"Table: {title}")

    if content:
        # Attempt to display as DataFrame, fallback to Markdown
        try:
            if pd: # Check if pandas is available
                # Preprocess markdown for better parsing: remove leading/trailing whitespace, split lines
                lines = [line.strip() for line in content.strip().split('\n')]
                # Remove separator line if present
                if len(lines) > 1 and all(c in '-:| ' for c in lines[1]):
                    lines.pop(1)
                # Reconstruct cleaned content for StringIO
                cleaned_content = "\n".join(lines)

                df = pd.read_csv(
                    io.StringIO(cleaned_content),
                    sep="|",
                    skipinitialspace=True,
                    index_col=False # Important: Avoid assuming first column is index
                ).iloc[1:] # Skip the header row parsed as data
                # Clean up columns: remove leading/trailing spaces in names, drop empty columns
                df.columns = [col.strip() for col in df.columns]
                df = df.drop(columns=[col for col in df.columns if col == ''], errors='ignore') # Drop unnamed columns from extra separators
                df = df.dropna(axis=1, how='all')
                st.dataframe(df.reset_index(drop=True)) # Show clean index
            else:
                st.markdown(f"```markdown\n{content}\n```")
        except Exception as e:
            logging.warning(f"Pandas failed to parse table, showing markdown. Error: {e}. Content:\n{content}")
            st.markdown(f"```markdown\n{content}\n```")
    else:
        st.caption("Table content not available.")

def display_visual_element(visual):
    """Display visual element data (accepts model or dict)."""
    if not visual: return
    visual_data = visual.model_dump() if hasattr(visual, 'model_dump') else visual

    visual_type = visual_data.get('type', 'Unknown Type')
    description = visual_data.get('description', 'No description available.')
    data_summary = visual_data.get('data_summary') # Can be None or "N/A"

    st.markdown(f"**🖼️ {visual_type.capitalize()}**")
    st.markdown(f"*Description:* {description}")
    if data_summary and data_summary != "N/A":
        st.markdown(f"*Data Summary:* {data_summary}")
    # Consider adding visual_id if useful: st.caption(f"ID: {visual_data.get('visual_id')}")

def display_numerical_data_point(num_data):
    """Displays a single numerical data point."""
    if not num_data: return
    num_dict = num_data.model_dump() if hasattr(num_data, 'model_dump') else num_data

    value = num_dict.get('value', 'N/A')
    description = num_dict.get('description', 'No description.')
    context = num_dict.get('context')

    st.markdown(f"**{value}**: {description}")
    if context:
        st.caption(f"Context: {context}")

def display_qdrant_node_network(qdrant_nodes_data: Optional[List[Dict]], title="Concept/Entity Network"): # Accept List[Dict]
    """Displays relationships from QdrantNode data (as dicts) as a network graph."""
    if not qdrant_nodes_data: # Handles None or empty list
        st.info("No relationship data available to display.")
        return

    # Ensure input is a list of dictionaries
    if not isinstance(qdrant_nodes_data, list) or not all(isinstance(n, dict) for n in qdrant_nodes_data):
         st.error("Invalid data format passed to node network display function (expected list of dicts).")
         logging.error(f"display_qdrant_node_network received invalid data type: {type(qdrant_nodes_data)}")
         return

    if not nx or not plt:
        # ... (Fallback text display as before, using .get() on dicts) ...
        st.warning("NetworkX/Matplotlib not installed. Displaying text list.")
        st.markdown("#### Relationships (Text List):")
        for node_dict in qdrant_nodes_data:
             linked_nodes = node_dict.get('linked_nodes', [])
             if linked_nodes:
                 for link_dict in linked_nodes:
                     st.markdown(f"- **{node_dict.get('name','?')}** → **{link_dict.get('target_node_id','?')}**: {link_dict.get('relationship_description','?')} (Strength: {link_dict.get('relationship_strength','?')})")
        return

    # Create a network graph
    G = nx.DiGraph()
    node_labels = {}
    edge_weights = {}
    edge_labels = {}
    node_types_map = {} # Store node type for coloring

    # Add nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        node_id = node_dict.get('node_id')
        node_name = node_dict.get('name', 'Unknown Node')
        node_type = node_dict.get('node_type', 'unknown')
        if node_id and node_id not in G: # Check if node_id exists and is not None
            G.add_node(node_id)
            node_labels[node_id] = f"{node_name}\n({node_type})"
            node_types_map[node_id] = node_type # Store type

    # Add edges from linked_nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        source_id = node_dict.get('node_id')
        linked_nodes = node_dict.get('linked_nodes', [])
        if not source_id or not G.has_node(source_id) or not isinstance(linked_nodes, list):
             continue # Skip if source_id invalid or no links

        for link_dict in linked_nodes:
             if not isinstance(link_dict, dict): continue # Skip invalid links
             target_id = link_dict.get('target_node_id')

             if not target_id: continue # Skip links without target

             # Ensure target node exists in the graph
             if target_id not in G:
                 G.add_node(target_id)
                 # Try to create a basic label from ID if target details aren't in current node list
                 target_name_match = re.match(r"entity_(.*?)_ch\d+|concept_(.*?)_ch\d+", target_id)
                 target_name_part = target_name_match.group(1) or target_name_match.group(2) if target_name_match else target_id
                 node_labels[target_id] = target_name_part.replace('_',' ').title() # Basic label
                 node_types_map[target_id] = 'unknown' # Mark as unknown type

             # Add edge if not already present
             if not G.has_edge(source_id, target_id):
                 strength = float(link_dict.get('relationship_strength', 5.0)) # Use get with default
                 description = link_dict.get('relationship_description', '')
                 G.add_edge(source_id, target_id, weight=strength)
                 edge_weights[(source_id, target_id)] = strength
                 edge_labels[(source_id, target_id)] = description

    # --- Graph drawing logic remains largely the same ---
    if not G.nodes():
         st.info("No nodes found to build the graph.")
         return
    if not G.edges():
         st.info("Nodes found, but no links between them to display.")
         st.write("Nodes identified:", list(node_labels.values()))
         return

    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(12, 10))
    try:
        pos = nx.spring_layout(G, k=0.5, iterations=50, seed=42)
    except Exception as layout_err:
        logging.error(f"NetworkX layout failed: {layout_err}. Using random layout.")
        pos = nx.random_layout(G, seed=42)

    # Draw nodes with color based on type
    node_colors = ['skyblue' if node_types_map.get(n) == 'entity' else 'lightgreen' if node_types_map.get(n) == 'concept' else 'lightgray' for n in G.nodes()]
    nx.draw_networkx_nodes(G, pos, node_size=1500, node_color=node_colors, alpha=0.9, ax=ax)

    # Draw edges
    if edge_weights: # Check if edges exist before calculating weights
        min_strength = min(edge_weights.values()) if edge_weights else 1
        max_strength = max(edge_weights.values()) if edge_weights else 10
        if max_strength == min_strength: max_strength += 1
        edge_width = [1 + 4 * (edge_weights.get(edge, 5.0) - min_strength) / (max_strength - min_strength) for edge in G.edges()]
        nx.draw_networkx_edges(G, pos, width=edge_width, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)
    else: # Draw default thin edges if no weights somehow
         nx.draw_networkx_edges(G, pos, width=1.0, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)


    # Draw labels
    nx.draw_networkx_labels(G, pos, labels=node_labels, font_size=9, font_weight='normal', ax=ax)

    ax.set_title(title, fontsize=16)
    fig.tight_layout()
    plt.axis('off')
    st.pyplot(fig)
    plt.close(fig)

    # Display edge details table
    edge_data = []
    # Use edge_labels which was populated correctly
    for (source_id, target_id), desc in edge_labels.items():
        edge_data.append({
            "Source": node_labels.get(source_id, source_id),
            "Target": node_labels.get(target_id, target_id),
            "Strength": edge_weights.get((source_id, target_id), 5.0),
            "Description": desc
        })

    if edge_data:
        st.markdown("### Relationship Details")
        try:
             if pd: # Check if pandas is available
                 df = pd.DataFrame(edge_data)
                 st.dataframe(df)
             else:
                 for item in edge_data: st.write(item) # Fallback display
        except Exception as df_err:
            logging.error(f"Failed to create DataFrame for edge data: {df_err}")
            for item in edge_data: st.write(item) # Fallback display
            
def display_chapter(chapter):
    """Display chapter details including subsections and Qdrant node relationships."""
    if not chapter: return
    chapter_data = chapter.model_dump() if hasattr(chapter, 'model_dump') else chapter

    title = chapter_data.get('title', 'Untitled Chapter')
    summary = chapter_data.get('summary', 'No summary available.')
    subsections = chapter_data.get('subsections', [])
    qdrant_nodes = chapter_data.get('qdrant_nodes') # This holds QdrantNode objects/dicts

    st.markdown(f"## {title}")
    st.markdown(f"_{summary}_")

    # Display Qdrant Node relationships if available
    if qdrant_nodes:
        with st.expander("View Concepts/Entities and Relationships within this Chapter", expanded=False):
            # Pass the actual list of nodes (could be dicts or models)
            display_qdrant_node_network(qdrant_nodes, title=f"Network for '{title}'")
    else:
        st.caption("No concept/entity relationship data available for this chapter.")

    # Display subsections
    if subsections:
        st.markdown("### Subsections Contained")
        for sub_idx, subsection_data in enumerate(subsections):
            if not subsection_data: continue
            # Handle dict or model
            sub_obj = subsection_data
            if isinstance(subsection_data, dict):
                 # Minimal validation or create Subsection obj if needed for consistency
                 # sub_obj = Subsection(**subsection_data) # Optional
                 pass

            sub_title = getattr(sub_obj, 'title', f'Subsection {sub_idx + 1}')
            sub_text = getattr(sub_obj, 'text', 'No text available.')
            sub_description = getattr(sub_obj, 'description', '')
            sub_id = getattr(sub_obj, 'subsection_id', 'N/A')
            page_num = "?" # Page number not stored on subsection anymore
            is_cutoff = getattr(sub_obj, 'is_cutoff', False)
            ref_visuals = getattr(sub_obj, 'referenced_visuals', [])
            ref_tables = getattr(sub_obj, 'referenced_tables', [])

            expander_title = f"{sub_idx + 1}. {sub_title}"
            if is_cutoff: expander_title += " (Continues...)"

            with st.expander(expander_title):
                if sub_description:
                    st.caption(f"Description: {sub_description}")
                st.markdown(sub_text) # Display the actual text
                st.caption(f"Subsection ID: {sub_id}")

                # Display referenced visuals/tables if any
                if ref_visuals: st.write(f"**Ref Visuals:** {', '.join(ref_visuals)}")
                if ref_tables: st.write(f"**Ref Tables:** {', '.join(ref_tables)}")
    else:
        st.caption("No subsections listed for this chapter.")

def display_node_links(all_nodes: List[Dict | Any], # List of node dicts or Pydantic objects
                       selected_node_id: str,
                       node_map: Optional[Dict[str, Dict | Any]] = None):
    """
    Finds a specific node by ID and displays its outgoing relationships (linked_nodes),
    including the origin context (document, chapter) of the target node.

    Args:
        all_nodes: A list containing all nodes (as dictionaries or Pydantic objects).
        selected_node_id: The node_id of the node to inspect.
        node_map: Optional pre-built dictionary mapping node_id to node object/dict
                  for faster lookups, especially for target node info.
    """
    function_name = "display_node_links"
    if not selected_node_id:
        st.caption("No node selected.")
        return

    # --- Find the selected node (Source Node) ---
    selected_node = None
    if node_map and selected_node_id in node_map:
        selected_node = node_map[selected_node_id]
    else:
        # Fallback search if map not provided
        for node in all_nodes:
            is_model = not isinstance(node, dict)
            node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
            if node_id == selected_node_id:
                selected_node = node
                break

    if selected_node is None:
        st.warning(f"Node with ID '{selected_node_id}' not found.")
        return

    is_source_model = not isinstance(selected_node, dict)
    source_name = getattr(selected_node, 'name', selected_node_id) if is_source_model else selected_node.get('name', selected_node_id)
    source_doc = getattr(selected_node, 'document_context', '?') if is_source_model else selected_node.get('document_context', '?')
    source_chapter = getattr(selected_node, 'chapter_id_context', '?') if is_source_model else selected_node.get('chapter_id_context', '?')

    st.markdown(f"#### Relationships FROM: **{source_name}**")
    st.caption(f"*(Source Origin: Document '{source_doc}', Chapter '{source_chapter}')*")
    st.markdown("---")


    linked_nodes = getattr(selected_node, 'linked_nodes', []) if is_source_model else selected_node.get('linked_nodes', [])

    if not linked_nodes or not isinstance(linked_nodes, list):
        st.caption("No outgoing links found for this node.")
        return

    # --- Build Node Map if needed (for target node lookups) ---
    if node_map is None:
        logging.debug(f"[{function_name}] Building temporary node_map for target info lookup.")
        node_map = {}
        for node in all_nodes:
             inner_is_model = not isinstance(node, dict)
             inner_node_id = getattr(node, 'node_id', None) if inner_is_model else node.get('node_id')
             if inner_node_id:
                 node_map[inner_node_id] = node

    # --- Display Links with Target Context ---
    if not node_map: # Check if map is still empty after trying to build
         st.warning("Could not build node map to look up target details.")

    for i, link in enumerate(linked_nodes):
        link_is_model = not isinstance(link, dict)
        target_id = getattr(link, 'target_node_id', 'Unknown Target ID') if link_is_model else link.get('target_node_id', 'Unknown Target ID')
        description = getattr(link, 'relationship_description', 'No description') if link_is_model else link.get('relationship_description', 'No description')
        strength = getattr(link, 'relationship_strength', 0.0) if link_is_model else link.get('relationship_strength', 0.0)
        keywords = getattr(link, 'relationship_keywords', []) if link_is_model else link.get('relationship_keywords', [])
        if not isinstance(keywords, list): keywords = []

        # --- Find Target Node Info using the map ---
        target_node = node_map.get(target_id) if node_map else None
        target_name = target_id # Default to ID if not found
        target_type = "unknown"
        target_doc = "unknown"
        target_chapter = "unknown"

        if target_node:
             target_is_model = not isinstance(target_node, dict)
             target_name = getattr(target_node, 'name', target_id) if target_is_model else target_node.get('name', target_id)
             target_type = getattr(target_node, 'node_type', 'unknown') if target_is_model else target_node.get('node_type', 'unknown')
             target_doc = getattr(target_node, 'document_context', 'unknown') if target_is_model else target_node.get('document_context', 'unknown')
             target_chapter = getattr(target_node, 'chapter_id_context', 'unknown') if target_is_model else target_node.get('chapter_id_context', 'unknown')
        elif target_id != 'Unknown Target ID':
             logging.warning(f"[{function_name}] Target node '{target_id}' linked from '{selected_node_id}' not found in node map.")


        # Display formatted link info including target context
        st.markdown(f"**{i+1}. Target Node:** {target_name} (`{target_type}`)")
        # Add target origin using italics or caption
        st.markdown(f"    *Origin: Document '{target_doc}', Chapter '{target_chapter}'*")
        st.markdown(f"    **Link Description:** {description}")
        st.markdown(f"    **Strength:** {strength:.1f}/10.0")
        if keywords:
            st.markdown(f"    **Keywords:** " + ", ".join(keywords))
        st.markdown(f"    *Target Node ID:* `{target_id}`") # Keep ID easily visible
        st.markdown("---") # Separator

    if not linked_nodes: # Redundant check, but safe
        st.caption("No valid outgoing links found to display.")



def display_project_ontology(ontology: Optional[ProjectOntology]): # Type hint remains Optional[ProjectOntology]
    """Displays project-wide ontology including the aggregated graph.
       Expects a ProjectOntology Pydantic model object as input.
    """
    function_name = "display_project_ontology"
    if not ontology:
        st.warning("No project ontology data available to display.")
        return

    # --- Treat input strictly as a ProjectOntology object ---
    # Remove the is_model check and the dictionary fallback logic
    if not hasattr(ontology, '__class__') or ontology.__class__.__name__ != 'ProjectOntology':
        st.error(f"[display_project_ontology] Error: Expected a ProjectOntology object based on name, but received type {type(ontology)}. Cannot display.")
        logging.error(f"[display_project_ontology] display_project_ontology received unexpected type: {type(ontology)}")
        return

    try:
        # Access attributes directly using getattr for safety or direct access
        title = getattr(ontology, 'title', 'Project Ontology [Default]')
        overview = getattr(ontology, 'overview', 'No overview available.')
        doc_count = getattr(ontology, 'document_count', 0)
        documents = getattr(ontology, 'documents', [])
        global_themes = getattr(ontology, 'global_themes', [])
        key_concepts = getattr(ontology, 'key_concepts', [])
        project_nodes = getattr(ontology, 'project_graph_nodes', None) # Access the graph nodes field

        st.title(f"🌐 {title}")
        st.markdown(f"**Documents Analyzed:** {doc_count}")

        # Display documents included
        if documents:
            st.markdown("## Documents Included In Analysis")
            if hasattr(st, 'chip'):
                for doc_title in documents: st.chip(doc_title, icon="📄")
            else:  # Fallback for older Streamlit versions
                docs_html = " ".join([f"<span style='background-color:#c5e1ff; padding:5px; margin:2px; border-radius:5px'>{doc}</span>" for doc in documents])
                st.markdown(docs_html, unsafe_allow_html=True)

        # Display overview
        st.markdown("## Overview")
        st.markdown(overview)

        # Display global themes
        if global_themes:
            st.markdown("## Global Themes")
            if hasattr(st, 'chip'):
                for theme in global_themes: st.chip(theme, icon="🏷️")
            else: # Fallback
                themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" for tag in global_themes])
                st.markdown(themes_html, unsafe_allow_html=True)

        # Display key concepts
        if key_concepts:
            st.markdown("## Key Concepts")
            if hasattr(st, 'chip'):
                 for concept in key_concepts: st.chip(concept, icon="💡")
            else: # Fallback
                concepts_html = " ".join([f"<span style='background-color:#f0f0f0; padding:5px; margin:2px; border-radius:5px'>{concept}</span>" for concept in key_concepts])
                st.markdown(concepts_html, unsafe_allow_html=True)

        # Display Project-Level Graph
        st.divider()
        st.markdown("## Project Knowledge Graph")
        if project_nodes and isinstance(project_nodes, list):
            logging.info(f"[{function_name}] Rendering project graph with {len(project_nodes)} nodes.")
            with st.expander("View Full Project Concept/Entity Network", expanded=False):
                 try:
                     # Convert Pydantic nodes back to dicts IF display_qdrant_node_network expects dicts
                     # If it accepts Pydantic objects, you can pass project_nodes directly
                     nodes_as_dicts = [node.model_dump(mode='json') if hasattr(node, 'model_dump') else node for node in project_nodes]
                     display_qdrant_node_network(nodes_as_dicts, title="Project-Wide Concept/Entity Network")
                 except NameError:
                      st.error("`display_qdrant_node_network` function not found.")
                      logging.error("`display_qdrant_node_network` function not found.")
                 except Exception as graph_err:
                      st.error(f"Failed to render project network graph: {graph_err}")
                      logging.error(f"Failed to render project network graph: {graph_err}", exc_info=True)
        else:
            st.caption("No cross-document relationship data was generated or processed to display a project graph.")


        # --- NEW: Node Relationship Explorer ---
        st.divider()
        st.markdown("## Node Relationship Explorer")

        if project_nodes and isinstance(project_nodes, list):
            # Create map for faster lookups and consistent data access
            node_map_for_display: Dict[str, Any] = {}
            node_options: Dict[str, str] = {} # Map display label -> node_id
            # Sort nodes alphabetically by name for the dropdown
            sorted_nodes = sorted(project_nodes, key=lambda n: getattr(n, 'name', '') if not isinstance(n, dict) else n.get('name', ''))

            for node in sorted_nodes:
                is_model = not isinstance(node, dict)
                node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
                node_name = getattr(node, 'name', node_id) if is_model else node.get('name', node_id)
                node_type = getattr(node, 'node_type', 'unknown') if is_model else node.get('node_type', 'unknown')

                if node_id:
                    display_label = f"{node_name} ({node_type})"
                    # Handle potential duplicate display labels (though unlikely with type added)
                    count = 1
                    original_label = display_label
                    while display_label in node_options:
                        count += 1
                        display_label = f"{original_label} ({count})"

                    node_options[display_label] = node_id
                    node_map_for_display[node_id] = node # Store original object/dict

            if node_options:
                selected_label = st.selectbox(
                    "Select a Node to see its connections:",
                    options=["Select a Node..."] + list(node_options.keys()),
                    index=0,
                    key="project_node_selector"
                )

                if selected_label != "Select a Node...":
                    selected_id = node_options.get(selected_label)
                    if selected_id:
                        # Call the helper function to display links
                        display_node_links(
                            project_nodes, # Pass the original list
                            selected_id,
                            node_map=node_map_for_display # Pass the map for efficiency
                         )
                    else:
                        st.warning("Could not find the ID for the selected node.")
            else:
                st.info("No valid nodes available to select.")
        else:
            st.caption("No nodes available to explore relationships.")
        # --- END: Node Relationship Explorer ---

    except AttributeError as attr_err:
        # Catch cases where expected attributes are missing from the model instance
        st.error(f"Error accessing project ontology data: Missing attribute - {attr_err}")
        logging.error(f"[{function_name}] Missing attribute in ProjectOntology object: {attr_err}", exc_info=True)
    except Exception as e:
        # Catch any other unexpected errors during display
        st.error(f"An unexpected error occurred while displaying the project ontology: {e}")
        logging.error(f"[{function_name}] Unexpected error displaying ontology: {e}", exc_info=True)


def display_page_details(page_content_data):
    """Displays details for a single page based on the simplified PageContent model."""
    if not page_content_data: return
    page_dict = page_content_data.model_dump() if hasattr(page_content_data, 'model_dump') else page_content_data

    page_num = page_dict.get("page_number", "?")
    st.header(f"📑 Page {page_num} Details")

    # Display Content Flags
    st.caption(f"Detected Content: "
               f"{'Tables ✅' if page_dict.get('has_tables') else 'Tables ❌'} | "
               f"{'Visuals ✅' if page_dict.get('has_visuals') else 'Visuals ❌'} | "
               f"{'Numbers ✅' if page_dict.get('has_numbers') else 'Numbers ❌'}")
    st.divider()

    # Display Subsections first as they contain the core text
    subsections = page_dict.get('subsections', [])
    if subsections:
        st.subheader("Content Subsections")
        for sub_idx, sub_data in enumerate(subsections):
             sub_title = sub_data.get('title', f'Subsection {sub_idx+1}')
             sub_desc = sub_data.get('description', '')
             sub_text = sub_data.get('text', 'No text.')
             sub_id = sub_data.get('subsection_id', 'N/A')
             is_cutoff = sub_data.get('is_cutoff', False)
             ref_v = sub_data.get('referenced_visuals', [])
             ref_t = sub_data.get('referenced_tables', [])

             exp_title = f"{sub_idx+1}. {sub_title}"
             if is_cutoff: exp_title += " (Continues...)"

             with st.expander(exp_title, expanded=True): # Expand by default
                 if sub_desc:
                     st.caption(f"Description: {sub_desc}")
                 st.markdown(sub_text)
                 # Optionally show refs if needed
                 if ref_v or ref_t:
                      refs = []
                      if ref_v: refs.append(f"Visuals: {', '.join(ref_v)}")
                      if ref_t: refs.append(f"Tables: {', '.join(ref_t)}")
                      st.caption(f"References: {'; '.join(refs)} | ID: {sub_id}")
                 else:
                      st.caption(f"ID: {sub_id}")

        st.divider()
    else:
         st.info("No subsections were extracted for this page.")

    # Display Tables
    tables = page_dict.get('tables', [])
    if tables:
        st.subheader(f"Tables on Page {page_num}")
        for table in tables:
            display_table(table) # Use the existing helper
        st.divider()

    # Display Visuals
    visuals = page_dict.get('visuals', [])
    if visuals:
        st.subheader(f"Visuals on Page {page_num}")
        for visual in visuals:
            display_visual_element(visual) # Use the existing helper
        st.divider()

    # Display Numerical Data
    numbers = page_dict.get('numbers', [])
    if numbers:
        st.subheader(f"Numerical Data on Page {page_num}")
        for num_data in numbers:
            display_numerical_data_point(num_data) # Use new helper
        st.divider()

    # --- REMOVED display of old fields: title, topics, summary, entities, page text area ---


#############################
# STREAMLIT UI COMPONENTS (Updated)
#############################

def render_unified_document_report(document_data: Optional[Dict]):
    """
    Displays the full report for a single processed document, including
    Executive Summary (with full graph), Chapters, and Detailed Page View tabs.

    Args:
        document_data: A dictionary containing the finalized document structure,
                       typically the output from finalize_document. Expected structure:
                       {
                           "raw_extracted_content": {
                               "filename": str,
                               "pages": List[Dict], # List of PageContent dicts
                               "summary": Dict | None, # Dict from DocumentSummaryDetails + chapters
                               "error": str | None
                           }
                       }
                       Can also accept the ProcessedDocument model object directly.
    """
    function_name = "render_unified_document_report"
    if not document_data:
        st.error(f"[{function_name}] No document data provided to display.")
        return

    # --- Data Extraction and Validation ---
    raw_content = {}
    # Handle both dict input and potential Pydantic model input
    if hasattr(document_data, 'model_dump'): # Check if it's a Pydantic model
        logging.debug(f"[{function_name}] Input is a Pydantic model, dumping to dict.")
        try:
            doc = document_data.model_dump(mode='json') # Use mode='json' for potential complex types
        except Exception as dump_err:
            st.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}")
            logging.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}", exc_info=True)
            return
    elif isinstance(document_data, dict):
        doc = document_data
    else:
        st.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        logging.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        return

    # Safely get raw_content, checking structure
    if isinstance(doc, dict) and "raw_extracted_content" in doc and isinstance(doc["raw_extracted_content"], dict):
        raw_content = doc["raw_extracted_content"]
        logging.debug(f"[{function_name}] Successfully extracted raw_extracted_content.")
    # Add fallback if top-level dict *is* raw_content (less likely with finalize_document output)
    elif isinstance(doc, dict) and "filename" in doc and "pages" in doc:
        logging.warning(f"[{function_name}] Input dict seems to be raw_extracted_content directly. Using it.")
        raw_content = doc
    else:
        st.error(f"[{function_name}] Could not find 'raw_extracted_content' dictionary in the provided data.")
        logging.error(f"[{function_name}] Invalid data structure, missing 'raw_extracted_content'. Data keys: {list(doc.keys()) if isinstance(doc, dict) else 'N/A'}")
        return

    # Extract key components, providing defaults
    filename = raw_content.get('filename', 'Unknown Document')
    processing_error = raw_content.get('error')
    pages = raw_content.get('pages', []) # Should be list of dicts
    summary_data = raw_content.get('summary', {}) # Should be a dict
    if not isinstance(summary_data, dict):
        logging.warning(f"[{function_name}] 'summary' data for '{filename}' is not a dict ({type(summary_data)}). Resetting to empty dict.")
        summary_data = {} # Ensure it's a dict for safe access later
    chapters_list = summary_data.get('chapters', []) # Get chapter list

    # --- Display Processing Error (if any) ---
    if processing_error:
        st.error(f"🚨 **Processing Error for {filename}:** {processing_error}")
        st.warning("Displayed data below might be incomplete or relate to a previous successful run.")
        # Decide whether to continue rendering or stop here
        # return # Option: Stop rendering if there was a critical error

    # --- Header Section ---
    st.header(f"📄 Report: {filename}")

    # Count actual elements present in pages for the caption
    num_pages = len(pages) if isinstance(pages, list) else 0
    num_tables = 0
    num_visuals = 0
    num_numbers = 0
    if isinstance(pages, list):
        for p in pages:
            if isinstance(p, dict):
                 num_tables += len(p.get('tables', [])) if isinstance(p.get('tables'), list) else 0
                 num_visuals += len(p.get('visuals', [])) if isinstance(p.get('visuals'), list) else 0
                 num_numbers += len(p.get('numbers', [])) if isinstance(p.get('numbers'), list) else 0

    caption_parts = [f"📄 {num_pages} pages"]
    if num_tables > 0: caption_parts.append(f"📊 {num_tables} tables")
    if num_visuals > 0: caption_parts.append(f"🖼️ {num_visuals} visuals")
    if num_numbers > 0: caption_parts.append(f"🔢 {num_numbers} numbers")
    st.caption(" | ".join(caption_parts))

    # Create tabs for document view modes
    tab_titles = ["Executive Summary", "Chapters", "Detailed Page View"]
    tab1, tab2, tab3 = st.tabs(tab_titles)

    # --- Tab 1: Executive Summary ---
    with tab1:
        # Display Summary, Themes, Questions
        with st.container(border=True):
            st.subheader("Executive Summary")
            summary_title = summary_data.get('title')
            if summary_title and summary_title != f"Summary Failed: {filename}" and summary_title != f"Error Finalizing {filename}":
                 st.markdown(f"**Title:** {summary_title}") # Display LLM generated title if valid

            summary_text = summary_data.get('summary')
            if summary_text and not summary_text.startswith("Failed to generate"):
                st.markdown(summary_text)
            else:
                st.info("No summary text available or summary generation failed.")

            themes = summary_data.get('themes')
            if themes and themes != ["error"]:
                st.markdown("#### Key Themes")
                if hasattr(st, 'chip'):
                    for theme in themes: st.chip(theme, icon="🏷️")
                else: # Fallback display
                     st.markdown("- " + "\n- ".join(themes))

            questions = summary_data.get('questions')
            if questions:
                 st.markdown("#### Key Questions Answered")
                 for q in questions: st.markdown(f"- {q}")

        # Display Document-Level Graph
        st.divider()
        st.subheader("Document Concept & Entity Network")

        # 1. Aggregate all nodes from all chapters
        all_document_nodes = []
        if isinstance(chapters_list, list):
            for chapter_dict in chapters_list:
                # Ensure chapter_dict is a dict and contains qdrant_nodes list
                if isinstance(chapter_dict, dict):
                    nodes_in_chapter = chapter_dict.get('qdrant_nodes', [])
                    if isinstance(nodes_in_chapter, list):
                        # Nodes should already be dicts after model_dump in finalize_document
                        all_document_nodes.extend(nodes_in_chapter)
                    else:
                        logging.warning(f"[{function_name}] Chapter '{chapter_dict.get('chapter_id')}' has non-list 'qdrant_nodes'.")
                else:
                    logging.warning(f"[{function_name}] Found non-dictionary item in chapters list for '{filename}'.")

        # 2. Display the aggregated network using the helper function
        if all_document_nodes:
            logging.info(f"[{function_name}] Rendering overall network graph with {len(all_document_nodes)} nodes for {filename}")
            with st.expander("View Full Document Network Graph", expanded=False):
                 try:
                      # Ensure display_qdrant_node_network exists and handles list of dicts
                      display_qdrant_node_network(
                          all_document_nodes,
                          title=f"Overall Concept/Entity Network for {filename}"
                      )
                 except NameError:
                      st.error("`display_qdrant_node_network` function not found.")
                      logging.error("`display_qdrant_node_network` function not found.")
                 except Exception as graph_err:
                      st.error(f"Failed to render document network graph: {graph_err}")
                      logging.error(f"Failed to render document network graph for {filename}: {graph_err}", exc_info=True)
        else:
            st.caption("No concept or entity data was extracted or processed correctly across chapters to build a document network.")

        # Page Content Preview Section
        st.divider()
        st.subheader("Page Content Preview")
        if pages and isinstance(pages, list):
            # Create options for the selectbox
            page_options = {} # Map label to index
            for i, p in enumerate(pages):
                page_num = p.get('page_number', i + 1) if isinstance(p, dict) else i + 1
                page_options[f"Page {page_num}"] = i

            if page_options:
                selected_page_label = st.selectbox(
                    "Select page to preview content:",
                    options=["Select a page..."] + list(page_options.keys()),
                    index=0, # Default to placeholder
                    key=f"page_nav_summary_{filename}" # Unique key per document
                )
                if selected_page_label != "Select a page...":
                    try:
                        page_idx = page_options[selected_page_label]
                        page_to_preview = pages[page_idx]

                        if isinstance(page_to_preview, dict):
                             with st.container(border=True):
                                  st.markdown(f"**Previewing Page {page_to_preview.get('page_number')}**")
                                  # Show subsection titles and descriptions for the selected page
                                  page_subsections = page_to_preview.get('subsections', [])
                                  if page_subsections and isinstance(page_subsections, list):
                                       st.markdown("**Subsections on this page:**")
                                       for sub_idx, sub in enumerate(page_subsections):
                                           if isinstance(sub, dict):
                                                sub_title = sub.get('title', f'Subsection {sub_idx+1}')
                                                sub_desc = sub.get('description', '')
                                                st.markdown(f"- **{sub_title}**: {sub_desc}")
                                           else: st.caption("- Invalid subsection format")
                                  else:
                                       st.caption("No subsections extracted or available for this page.")
                                  # Optionally show raw text preview if helpful
                                  raw_page_text = page_to_preview.get('raw_text')
                                  if raw_page_text:
                                       with st.expander("View Raw Text (preview)"):
                                             st.text(raw_page_text[:1000] + ("..." if len(raw_page_text) > 1000 else ""))

                        else:
                            st.warning(f"Invalid data format for page index {page_idx}.")

                    except (IndexError, KeyError, ValueError) as e:
                        st.warning(f"Could not select page for preview: {e}")
            else:
                st.info("No pages available in this document.")
        else:
            st.info("No page data available for preview.")


    # --- Tab 2: Chapters ---
    with tab2:
        if chapters_list and isinstance(chapters_list, list):
            chapter_titles = [ch.get('title', f"Chapter {ch.get('order', i+1)}") for i, ch in enumerate(chapters_list) if isinstance(ch, dict)]
            if chapter_titles:
                try:
                    chapter_tabs = st.tabs(chapter_titles)
                    for i, (tab, chapter_dict) in enumerate(zip(chapter_tabs, chapters_list)):
                        # Ensure we only process valid chapter dicts
                        if isinstance(chapter_dict, dict):
                            with tab:
                                try:
                                    # Ensure display_chapter exists and handles dicts
                                    display_chapter(chapter_dict)
                                except NameError:
                                    st.error("`display_chapter` function not found.")
                                    logging.error("`display_chapter` function not found.")
                                    break # Stop trying if function is missing
                                except Exception as chapter_disp_err:
                                     st.error(f"Error displaying chapter '{chapter_dict.get('title')}': {chapter_disp_err}")
                                     logging.error(f"Error displaying chapter {chapter_dict.get('chapter_id')}: {chapter_disp_err}", exc_info=True)

                        else:
                             logging.warning(f"[{function_name}] Skipping non-dictionary item found in chapters list for '{filename}'.")

                except Exception as tabs_err: # Catch errors from st.tabs if titles are invalid etc.
                     st.error(f"Could not create chapter tabs: {tabs_err}")
                     logging.error(f"Could not create chapter tabs for {filename}: {tabs_err}", exc_info=True)
            else:
                 st.info("No valid chapter titles found in the summary data.")
        else:
            st.info("No chapter information available or processed for this document.")
            if summary_data.get("error") and "chapter extraction" in summary_data.get("error","").lower():
                 st.warning(f"Note: Chapter extraction previously failed for this document. Reason: {summary_data['error']}")


    # --- Tab 3: Detailed Page View ---
    with tab3:
        if pages and isinstance(pages, list):
             # Create options mapping Page Number (int) -> Index (int)
             page_num_options: Dict[int, int] = {}
             valid_page_nums: List[int] = [] # Store the actual page numbers (integers)
             for i, p in enumerate(pages):
                 # Ensure page data is dict and page_number is an int
                 if isinstance(p, dict) and isinstance(p.get('page_number'), int):
                      page_num = p['page_number']
                      page_num_options[page_num] = i # Map the integer page number to its list index
                      valid_page_nums.append(page_num)
                 else:
                      logging.warning(f"[{function_name}] Skipping invalid page data at index {i} (page_number missing or not int) for page view selector in '{filename}'. Page data keys: {list(p.keys()) if isinstance(p,dict) else 'N/A'}")

             if valid_page_nums:
                  # Sort page numbers numerically for the selector options
                  sorted_page_nums = sorted(valid_page_nums)

                  # Use format_func to display "Page X", but the options are the integers
                  selected_page_num = st.selectbox(
                       "Select Page Number:",
                       options=sorted_page_nums,       # Pass the list of INT page numbers
                       format_func=lambda x: f"Page {x}", # Display "Page X" to the user
                       index=0,                         # Default selection to the first page
                       key=f"page_detail_selector_{filename}" # Unique key for this selectbox
                  )
                  # 'selected_page_num' now directly holds the chosen INTEGER page number (e.g., 5)

                  # Find the corresponding page dictionary using the selected integer page number
                  try:
                       # REMOVED: selected_num = int(selected_page_num.split()[-1])
                       # USE selected_page_num directly as the key
                       selected_index = page_num_options.get(selected_page_num)

                       if selected_index is not None:
                           selected_page_dict = pages[selected_index]
                           try:
                               # Ensure display_page_details exists and handles dicts
                               display_page_details(selected_page_dict)
                           except NameError:
                                st.error("`display_page_details` function not found.")
                                logging.error("`display_page_details` function not found.")
                           except Exception as page_disp_err:
                                st.error(f"Error displaying details for page {selected_page_num}: {page_disp_err}")
                                logging.error(f"Error displaying details for page {selected_page_num} of {filename}: {page_disp_err}", exc_info=True)
                       else:
                           # This should be unlikely if valid_page_nums was populated correctly
                           st.warning(f"Internal error: Could not find index for selected page number {selected_page_num}.")
                  except Exception as e: # Catch any unexpected errors during index lookup/display
                       st.error(f"Error processing page selection or display: {e}")
                       logging.error(f"Error processing page selection {selected_page_num} or display for {filename}: {e}", exc_info=True)

             else:
                  st.info("No valid page data with page numbers found to display details.")
        else:
             st.info("No page data available to display details.")


# --- Sidebar function (Mostly unchanged, ensure it calls correct ontology function if using simplified) ---

def display_sidebar_chat():
    """Manages the sidebar content: file upload, chat, and triggers processing."""
    st.sidebar.title("📝 Input & Control")
    st.sidebar.markdown("Upload documents, view reports, generate ontology, and chat.")

    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]

    # Initialize session state robustly
    if "processed_file_names" not in st.session_state: st.session_state.processed_file_names = set()
    if "processing_active" not in st.session_state: st.session_state.processing_active = False
    if "processed_documents" not in st.session_state: st.session_state.processed_documents = []
    if "show_results" not in st.session_state: st.session_state.show_results = False
    if "project_ontology" not in st.session_state: st.session_state.project_ontology = None
    if "show_ontology" not in st.session_state: st.session_state.show_ontology = False
    if "messages" not in st.session_state: st.session_state.messages = []
    if "selected_docs" not in st.session_state: st.session_state.selected_docs = []


    # --- File Upload ---
    st.sidebar.subheader("📁 Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Select files:",
        type=supported_types,
        accept_multiple_files=True,
        key="sidebar_file_uploader_main" # Ensure key consistency if needed elsewhere
    )

    # --- Processing Trigger ---
    # Logic simplified: Button triggers processing of ALL non-processed uploaded files.
    files_ready_to_process = []
    files_data_for_processing = []
    if uploaded_files:
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                files_ready_to_process.append(file)
                try:
                    file_content = file.getvalue()
                    if not file_content:
                         st.sidebar.warning(f"File '{file.name}' is empty, skipping.")
                         continue
                    file_ext = file.name.split('.')[-1].lower() if '.' in file.name else ''
                    files_data_for_processing.append({
                        "name": file.name,
                        "content": file_content,
                        "type": file_ext
                    })
                except Exception as e:
                    st.sidebar.error(f"Error reading file {file.name}: {e}")

    # Display button only if there are files to process and not already processing
    if files_ready_to_process and not st.session_state.processing_active:
        if st.sidebar.button(f"Process {len(files_ready_to_process)} New File(s)", key="process_files_btn"):
            if files_data_for_processing:
                # Create the status elements just before starting
                with st.sidebar: # Ensure elements appear in the sidebar
                    status_container = st.status("Starting document processing...", expanded=True)
                    progress_bar = st.progress(0)
                    time_info = st.empty()
                    # Store handles in session state ONLY if they might be needed by the async function
                    # In this revised structure, update_doc_progress is called within the main async func
                    # So, we pass the UI elements or let the function access session state directly
                    st.session_state.status_container = status_container
                    st.session_state.progress_bar = progress_bar
                    st.session_state.time_info = time_info

                try:
                    st.session_state.processing_active = True
                    st.sidebar.warning("Processing started. Please wait...") # Give user feedback

                    # --- Call the REVISED process_all_documents_async ---
                    processed_results_list = run_async(
                        process_all_documents_async,
                        files_data_for_processing # Pass the list of valid file dicts
                    )
                    # --- ---

                    # Update session state with potentially partial/error results
                    new_processed_docs = []
                    processed_names_this_run = set()
                    for result_dict in processed_results_list:
                        if isinstance(result_dict, dict) and "raw_extracted_content" in result_dict:
                            new_processed_docs.append(result_dict)
                            fname = result_dict["raw_extracted_content"].get("filename")
                            if fname:
                                processed_names_this_run.add(fname)
                        else:
                            logging.warning(f"Received invalid item in processed results list: {type(result_dict)}")

                    st.session_state.processed_documents.extend(new_processed_docs)
                    st.session_state.processed_file_names.update(processed_names_this_run)

                    st.session_state.show_results = True # Automatically show results
                    if 'status_container' in st.session_state and st.session_state.status_container: # Check if UI element exists
                        st.session_state.status_container.update(label="✅ Document processing finished!", state="complete", expanded=False)
                    st.sidebar.success(f"Processed {len(processed_names_this_run)} file(s).")
                    # No rerun needed here, main loop will handle showing results

                except Exception as e:
                    st.sidebar.error(f"Error during document processing orchestration: {e}")
                    logging.error(f"Error during document processing orchestration: {e}", exc_info=True)
                    if 'status_container' in st.session_state and st.session_state.status_container:
                        st.session_state.status_container.update(label=f"❌ Processing Error: {e}", state="error", expanded=False)
                finally:
                    st.session_state.processing_active = False
                    # Clean up UI element references after processing run
                    for key in ["status_container", "progress_bar", "time_info"]:
                         if key in st.session_state:
                              # Attempt to clear the elements visually if possible, then delete key
                              try:
                                   if key == "time_info" and st.session_state[key]: st.session_state[key].empty()
                                   # Status and progress bar might persist until next run without explicit clear
                              except Exception: pass # Ignore errors during cleanup
                              del st.session_state[key]
                    st.rerun() # Rerun to update sidebar state after processing finishes/fails

            else:
                st.sidebar.warning("No valid new files could be prepared for processing.")

    # --- Display Processed Files & Control Buttons ---
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---")
        st.sidebar.markdown("**Processed Files:**")
        # Sort names for consistent display
        for filename in sorted(list(st.session_state.processed_file_names)):
             # Check if the document had an error during processing
             has_error = False
             for doc in st.session_state.processed_documents:
                  if isinstance(doc, dict) and doc.get("raw_extracted_content", {}).get("filename") == filename:
                       if doc.get("raw_extracted_content", {}).get("error"):
                            has_error = True
                            break
             icon = "⚠️" if has_error else "✓"
             st.sidebar.caption(f"{icon} {filename}")

        st.sidebar.markdown("---")
        col1, col2 = st.sidebar.columns(2)

        with col1:
            # Disable button if processing is active
            if st.button("View Reports", key="view_reports_btn", disabled=st.session_state.processing_active):
                if st.session_state.processed_documents:
                    st.session_state.show_results = True
                    st.session_state.show_ontology = False
                    st.rerun()
                else:
                    st.sidebar.warning("No processed document data available.")

        with col2:
            ontology_exists = st.session_state.project_ontology is not None
            ontology_button_label = "View Ontology" if ontology_exists else "Gen. Ontology"
            # Disable button if processing is active
            if st.button(ontology_button_label, key="ontology_btn", disabled=st.session_state.processing_active):
                # Check if there are successfully processed docs suitable for ontology
                valid_docs_for_ontology = [
                    doc for doc in st.session_state.processed_documents
                    if isinstance(doc, dict) and not doc.get("raw_extracted_content", {}).get("error")
                ]

                if valid_docs_for_ontology:
                    if not ontology_exists:
                        with st.spinner("Generating project ontology... This may take a while."):
                            try:
                                api_key = get_gemini_api_key()
                                client = genai.Client(api_key=api_key,
                                    http_options={
                                        "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                                        "headers": {
                                            "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                                            "helicone-target-url": 'https://generativelanguage.googleapis.com'
                                        } if st.secrets.get("HELICONE_PROXY_URL") and st.secrets.get("HELICONE_API_KEY") else {}
                                    })
                                # Pass only the valid document dictionaries
                                # ontology_data = run_async(generate_project_ontology, client, valid_docs_for_ontology)
                                # *** USE THE NEW FUNCTION ***
                                ontology_data = run_async(
                                    incrementally_generate_doc_level_ontology, # <-- NEW FUNCTION
                                    client,
                                    valid_docs_for_ontology      # Pass the list of processed doc dicts
                                    # No need for current_ontology - this version rebuilds
                                )
                                # *** END OF MODIFIED LINE ***
                                
                                st.session_state.project_ontology = ontology_data # Store the ProjectOntology object
                                if ontology_data: # Check if generation succeeded
                                     st.sidebar.success("Project ontology generated!")
                                else:
                                     st.sidebar.error("Ontology generation failed (returned None). Check logs.")

                            except Exception as gen_err:
                                st.sidebar.error(f"Ontology generation failed: {gen_err}")
                                logging.error(f"Ontology generation failed: {gen_err}", exc_info=True)
                                st.session_state.project_ontology = None

                    # Show ontology if generated or already exists
                    if st.session_state.project_ontology:
                         st.session_state.show_ontology = True
                         st.session_state.show_results = False
                         st.rerun()
                    elif not ontology_exists: # If generation was just attempted and failed
                         pass # Error shown above
                    else: # Should not happen unless state is inconsistent
                         st.sidebar.warning("Ontology not available.")

                else:
                    st.sidebar.warning("No successfully processed documents available for ontology generation.")

    # --- Chat Section ---
    st.sidebar.divider()
    st.sidebar.subheader("💬 Document Chat")
    # Document selection for chat context
    if st.session_state.processed_documents:
        # Filter for successfully processed documents for chat selection
        valid_doc_options = []
        for doc_result in st.session_state.processed_documents:
             if isinstance(doc_result, dict) and "raw_extracted_content" in doc_result:
                  raw_content = doc_result["raw_extracted_content"]
                  filename = raw_content.get("filename")
                  # Only include docs without processing errors
                  if filename and not raw_content.get("error"):
                       valid_doc_options.append(filename)

        if valid_doc_options:
            st.session_state.selected_docs = st.sidebar.multiselect(
                "Select documents for chat context (max 3 recommended):",
                options=sorted(list(set(valid_doc_options))),
                default=st.session_state.selected_docs,
                key="doc_context_selector_sidebar",
                max_selections=3 # Limit context size for performance/cost
            )
        else:
            st.sidebar.caption("No successfully processed documents available to reference.")
    else:
        st.sidebar.caption("No documents processed yet.")

    # Chat History & Input Container
    with st.sidebar.container(border=True, height=400): # Adjusted height
        st.markdown("**Chat History**")
        # Display chat history (as before)
        if not st.session_state.get("messages", []):
            st.info("Ask a question below!")
        else:
            for message in st.session_state.messages[-10:]: # Show last 10 messages
                role = message.get("role", "")
                content = message.get("content", "")
                with st.chat_message(role):
                    st.markdown(content)

    # Chat input outside the history container
    if prompt := st.sidebar.chat_input("Ask about selected documents...", key="sidebar_chat_input_main", disabled=st.session_state.processing_active):
        current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        st.session_state.messages.append({"role": "user", "content": prompt, "timestamp": current_time})
        # Immediately show user message
        st.rerun()
        # Call async chat processing
        run_async(process_chat_message, prompt)


# --- Modified Main Function ---
def main():
    """Main function to run the Streamlit application."""
    st.set_page_config(
        page_title="Multimodal Document Processor",
        page_icon="📊",
        layout="wide"
    )

    # Initialize session state keys robustly (as before)
    default_session_state = {
        "processed_documents": [], "processed_file_names": set(), "processing_active": False,
        "last_uploaded_files": [], "messages": [], "selected_docs": [], "show_results": False,
        "selected_doc_to_view": None, "project_ontology": None, "show_ontology": False,
        # Removed UI element keys as they are now created/managed within sidebar function run
    }
    for key, default_value in default_session_state.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

    # --- SIDEBAR ---
    display_sidebar_chat() # Contains file upload, buttons, chat, and triggers processing

    # --- MAIN CONTENT AREA ---
    st.title("📊 Advanced Document Processor")
    st.markdown("Upload PDFs, Office documents, text files, and more to extract structure, relationships, and summaries.")
    st.divider()

    # Check flags in session state to determine what to display
    # Note: processing_active check might not be needed here if sidebar handles UI locking
    if st.session_state.get("show_ontology") and st.session_state.project_ontology:
        # Ensure display_project_ontology handles ProjectOntology object
        display_project_ontology(st.session_state.project_ontology)

    elif st.session_state.get("show_results") and st.session_state.processed_documents:
        st.header("Document Reports")
        # Document Selection Logic (unchanged, uses processed_documents list of dicts)
        doc_options = {}
        doc_list = st.session_state.processed_documents # List of final dicts
        for i, doc_result in enumerate(doc_list):
            filename = "Unknown Document"
            # Added check for dict structure
            if isinstance(doc_result, dict) and "raw_extracted_content" in doc_result:
                 filename = doc_result["raw_extracted_content"].get("filename", f"Doc_{i+1}")
            doc_options[f"{i+1}: {filename}"] = i

        selected_index = 0 # Default
        if len(doc_list) > 1:
             # Check if options exist before creating selectbox
             if doc_options:
                 selected_display_name = st.selectbox(
                     "Choose document report to view:",
                     options=list(doc_options.keys()), index=0, key="doc_report_selector_main"
                 )
                 selected_index = doc_options.get(selected_display_name, 0)
             else:
                 st.warning("Could not generate document options for selection.")
                 selected_index = -1 # Indicate no valid selection
        elif len(doc_list) == 1:
            # Ensure the single document is valid before setting index
            if doc_options:
                 selected_index = 0
            else:
                 st.warning("Could not generate document option for the single document.")
                 selected_index = -1
        else: # No documents
             st.info("No documents processed yet.")
             selected_index = -1

        if selected_index != -1:
            try:
                selected_doc_data = doc_list[selected_index]
                # Ensure render_unified_document_report handles the final dict structure
                render_unified_document_report(selected_doc_data)
            except IndexError:
                 st.error(f"Error: Could not retrieve document at index {selected_index}. Please refresh or reprocess.")
                 logging.error(f"IndexError accessing processed_documents at index {selected_index}, length is {len(doc_list)}")
            except Exception as display_err:
                 st.error(f"An error occurred displaying the selected document report: {display_err}")
                 logging.error(f"Error in render_unified_document_report: {display_err}", exc_info=True)

    else:
        # Welcome/Instructions Screen (unchanged)
        if not st.session_state.processing_active: # Only show welcome if not processing
             st.info("Upload documents using the sidebar to begin analysis.")
             st.markdown("""
             ## 📊 Advanced Multimodal Document Analyzer
             Leverage the power of Google's Gemini models to unlock insights from your documents. This application offers:

             - **Broad File Support**: Process PDFs (image-based), Office documents (Word, PowerPoint, Excel), plain text, Markdown, code files, and more.
             - **Intelligent Structure Extraction**: Automatically segment documents into meaningful **subsections**, identifying titles, descriptions, and content flow. Detects and extracts **tables**, **visuals** (charts, images), and key **numerical data points** within each page or chunk.
             - **Logical Chapter Generation**: Automatically group related subsections across pages into coherent **chapters** based on thematic content, providing a structured overview.
             - **Concept & Entity Graphing**: Identify key **concepts** and **entities** within each chapter and analyze their relationships, including connections **between chapters**, building a knowledge graph foundation. (Visualized per chapter).
             - **Automated Summarization**: Generate a concise **document-level summary**, including key themes and potential questions, derived from the extracted concepts and entities.
             - **Project-Wide Ontology (Optional)**: Create a high-level overview synthesizing the main themes and concepts across *all* uploaded documents, including *inter-document* relationships.
             - **Interactive Chat**: Engage in conversation with your documents; ask questions based on the extracted content and summaries (context-aware based on selected documents).

             ### Getting Started:

             1.  **Upload:** Use the sidebar to upload one or more supported documents.
             2.  **Process:** Click the "Process New File(s)" button that appears in the sidebar.
             3.  **Monitor:** Track the processing progress in the sidebar.
             4.  **Explore Reports:** Once finished, click "View Document Reports" in the sidebar. Use the tabs:
                 *   **Executive Summary:** Read the AI-generated summary and key themes. View the overall document network graph.
                 *   **Chapters:** Navigate through the structured chapters and view their concept/entity networks.
                 *   **Detailed Page View:** Examine the extracted subsections, tables, visuals, and numbers for each page.
             5.  **Chat:** Select documents in the sidebar and use the chat interface to ask specific questions.
             6.  **Synthesize:** Click "Generate Project Ontology" (optional) to get a combined overview of all processed documents and explore the project-wide knowledge graph.
             """)
             st.image("https://via.placeholder.com/800x400.png?text=Document+Structure+and+Analysis",
                      caption="From Raw Files to Structured Insights and Knowledge Graphs")

# --- Entry Point ---
if __name__ == "__main__":
    # Basic logging setup - Consider making level configurable
    log_level = os.environ.get("LOG_LEVEL", "INFO").upper()
    logging.basicConfig(
        level=getattr(logging, log_level, logging.INFO),
        format='%(asctime)s - %(levelname)s - [%(funcName)s:%(lineno)d] %(message)s'
    )
    logging.info(f"Logging level set to {log_level}")
    # Ensure all necessary Pydantic models and helper functions are defined above main()
    main()

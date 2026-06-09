import streamlit as st
import os
import json
import asyncio
import base64
import pandas as pd
import pymupdf  # PyMuPDF
from typing import List, Dict, Optional, Union, Any, Set
from pydantic import BaseModel, Field, ValidationError, field_validator, ConfigDict
from google import genai
from google.genai import types
import docx
from pptx import Presentation
import csv
import openpyxl
import re
from collections import Counter, defaultdict
import time
import io
import datetime
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np

###############################
# DATA MODELS
###############################

class VisualElement(BaseModel):
    """Model for visual elements like charts, graphs, etc."""
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data")
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")
    source_url: Optional[str] = Field(None, description="Source URL of the visual")
    alt_text: Optional[str] = Field(None, description="Alternative text for the visual")
    visual_id: Optional[str] = Field(None, description="Unique identifier for the visual")

class NumericalDataPoint(BaseModel):
    """Model for numerical data points extracted from text."""
    value: str = Field(..., description="The numerical value")
    description: str = Field(default="", description="What the number represents")
    context: str = Field(default="", description="Surrounding text context")
    model_config = ConfigDict(arbitrary_types_allowed=True, extra="ignore")
    @field_validator('value', 'description', 'context', mode='before')
    def convert_to_string(cls, v): return str(v) if v is not None else ""

class TableData(BaseModel):
    """Model for tables extracted from documents."""
    table_content: str = Field(default="", description="Markdown formatted table content")
    title: Optional[str] = None
    summary: Optional[str] = None
    page_number: int = 0
    table_id: Optional[str] = Field(None, description="Unique identifier for the table")
    model_config = ConfigDict(arbitrary_types_allowed=True)
    @field_validator('table_content')
    def clean_table_content(cls, v): return str(v) if v is not None else ""

class FinancialMetric(BaseModel):
    """Model for financial metrics extracted from documents."""
    name: str = Field(..., description="Name of the financial metric")
    value: str = Field(..., description="Metric value")
    period: Optional[str] = Field(None, description="Reporting period")
    context: Optional[str] = Field(None, description="Surrounding text context")
    trend: Optional[str] = Field(None, description="Trend direction if available")

class FinancialStatement(BaseModel):
    """Model for financial statements extracted from documents."""
    type: str = Field(..., description="Type of statement")
    period: str = Field(..., description="Reporting period")
    metrics: List[FinancialMetric] = Field(default_factory=list, description="Key metrics")

class Subsection(BaseModel):
    """Model for subsections extracted from pages."""
    subsection_id: str = Field(..., description="Unique identifier for the subsection")
    order: int = Field(..., description="Order of the subsection within the page")
    title: str = Field(..., description="Title of the subsection")
    content: str = Field(..., description="Raw content of the subsection")
    is_cutoff: bool = Field(False, description="True if content appears to be cut off")
    referenced_visuals: List[str] = Field(default_factory=list, description="IDs of referenced visuals")
    referenced_tables: List[str] = Field(default_factory=list, description="IDs of referenced tables")
    page_number: int = Field(..., description="Page number where this subsection appears")

class EntityRelationship(BaseModel):
    """Model for entity relationships within documents."""
    source_entity: str = Field(..., description="Source entity name")
    target_entity: str = Field(..., description="Target entity name")
    relationship_description: str = Field(..., description="Description of the relationship")
    relationship_keywords: List[str] = Field(default_factory=list, description="Keywords describing the relationship")
    relationship_strength: float = Field(..., description="Strength of the relationship (0-10)")

class Chapter(BaseModel):
    """Model for chapters composed of subsections."""
    chapter_id: str = Field(..., description="Unique identifier for the chapter")
    title: str = Field(..., description="Title of the chapter")
    summary: str = Field(..., description="Summary of the chapter")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections in this chapter")
    entity_relationships: List[EntityRelationship] = Field(default_factory=list, description="Entity relationships within the chapter")
    order: int = Field(..., description="Order of the chapter in the document")

class PageContent(BaseModel):
    """Model for content extracted from a single page."""
    model_config = ConfigDict(arbitrary_types_allowed=True)
    page_number: int = Field(..., description="Page number in document")
    text: str = Field("", description="Full text content")
    title: str = Field("Untitled", description="Brief title for this page")
    topics: List[str] = Field(default_factory=list, description="Key topics discussed")
    summary: str = Field("", description="Summary of key points")
    entities: List[str] = Field(default_factory=list, description="Important entities mentioned")
    has_tables: bool = Field(False, description="True if page contains tables")
    has_visuals: bool = Field(False, description="True if page contains visuals")
    has_numbers: bool = Field(False, description="True if page contains key numerical data")
    tables: List[Union[TableData, dict]] = Field(default_factory=list, description="Extracted tables")
    visuals: List[Union[VisualElement, dict]] = Field(default_factory=list, description="Extracted visual elements")
    numbers: List[Union[NumericalDataPoint, dict]] = Field(default_factory=list, description="Extracted numerical data")
    dates: List[str] = Field(default_factory=list, description="Important dates mentioned")
    financial_statements: List[Union[FinancialStatement, dict]] = Field(default_factory=list)
    key_metrics: List[Union[FinancialMetric, dict]] = Field(default_factory=list)
    financial_terms: List[str] = Field(default_factory=list)
    subsections: List[Union[Subsection, dict]] = Field(default_factory=list, description="Extracted subsections")

class DocumentSummary(BaseModel):
    """Document-level summary."""
    title: str = Field(..., description="Concise title/summary")
    themes: List[str] = Field(..., description="Concept/theme tags")
    questions: List[str] = Field(..., description="Hypothetical questions")
    summary: str = Field(..., description="Comprehensive summary")
    tables_summary: Optional[str] = Field(None, description="Summary of key tables")
    visuals_summary: Optional[str] = Field(None, description="Summary of key visuals")
    chapters: List[Chapter] = Field(default_factory=list, description="Extracted chapters")
    entity_relationships: List[EntityRelationship] = Field(default_factory=list, description="Document-level entity relationships")

class KeyTopic(BaseModel):
    """Represents a key topic in the document."""
    name: str = Field(..., description="Short topic name (1-3 words)")
    description: str = Field(..., description="Brief description of the topic")
    relevance: str = Field(..., description="Relevance level (High/Medium/Low)")
    sentiment: str = Field(..., description="Sentiment analysis (Positive/Neutral/Mixed/Cautionary/Negative)")
    analysis: str = Field(..., description="Brief analysis of the topic")

class QuotedStatement(BaseModel):
    """Represents an important quoted statement in the document."""
    speaker: str = Field(..., description="Person or entity who made the statement")
    quote: str = Field(..., description="The quoted text")
    page: int = Field(..., description="Page number where the quote appears")

class DocumentReport(BaseModel):
    """UI-friendly document report structure for presentation."""
    file_name: str = Field(..., description="Original filename")
    page_count: int = Field(..., description="Number of pages in document")
    title: str = Field(..., description="Document title")
    title_summary: str = Field(..., description="Brief summary/subtitle")
    concept_theme_hashtags: List[str] = Field(default_factory=list, description="Theme tags as hashtags")
    date_published: Optional[str] = Field(None, description="Publication date if available")
    source: Optional[str] = Field(None, description="Document source")
    confidence: str = Field("Medium", description="Analysis confidence (High/Medium/Low)")
    document_summary: str = Field(..., description="Comprehensive summary")
    key_insights: List[str] = Field(default_factory=list, description="Key takeaways and insights")
    key_topics: List[KeyTopic] = Field(default_factory=list, description="Detailed topic analysis")
    quoted_statements: List[QuotedStatement] = Field(default_factory=list, description="Important quotes from document")
    content_excerpt: Optional[str] = Field(None, description="Highlighted document content")
    chapters: List[Chapter] = Field(default_factory=list, description="Document chapters")

class ProcessedDocument(BaseModel):
    """Fully processed document with all content."""
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages")
    summary: Optional[DocumentSummary] = None

class ProjectOntology(BaseModel):
    """Project-wide ontology generated from multiple documents."""
    title: str = Field(..., description="Project title")
    overview: str = Field(..., description="Project overview")
    document_count: int = Field(..., description="Number of documents in the project")
    documents: List[str] = Field(..., description="List of document titles")
    global_themes: List[str] = Field(..., description="High-level project themes")
    entity_relationships: List[EntityRelationship] = Field(..., description="Project-wide entity relationships")
    key_concepts: List[str] = Field(..., description="Key concepts across all documents")

###############################
# API AND UTILITY FUNCTIONS
###############################

def convert_to_model(data, model_class):
    """Convert dictionary data to a Pydantic model with error handling."""
    if not isinstance(data, dict):
        print(f"Warning: Expected dict for {model_class.__name__}, got {type(data)}")
        return None
    
    try:
        # Extract only the fields that the model expects
        model_fields = model_class.__annotations__.keys()
        filtered_data = {k: v for k, v in data.items() if k in model_fields}
        return model_class(**filtered_data)
    except Exception as e:
        print(f"Error converting to {model_class.__name__}: {str(e)}")
        return None

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
        st.stop()
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)
            
            # Validate JSON if applicable
            config = kwargs.get('config')
            if (config and hasattr(config, 'response_mime_type') and 
                config.response_mime_type == 'application/json' and 
                hasattr(response, 'candidates') and response.candidates):
                try:
                    # Try to parse the JSON response
                    json_text = response.candidates[0].content.parts[0].text
                    json.loads(clean_json_response(json_text))  # Use our clean function first
                except json.JSONDecodeError as e:
                    if attempt < max_retries - 1:
                        print(f"Received malformed JSON on attempt {attempt+1}, retrying: {e}")
                        await asyncio.sleep(2 ** attempt)
                        continue
            
            return response
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            print(f"API call failed on attempt {attempt+1}, retrying: {e}")
            await asyncio.sleep(2 ** attempt)

def clean_json_response(json_text: str, extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        return "{}"
    
    # Ensure it's a string
    json_text = str(json_text)
    
    try:
        # Handle markdown code blocks
        if json_text.startswith("```"):
            blocks = json_text.split("```")
            if len(blocks) >= 3:
                json_text = blocks[1]
                if json_text.startswith("json"):
                    json_text = json_text[4:].strip()
            else:
                json_text = json_text.replace("```", "").strip()
        
        elif json_text.startswith("json"):
            json_text = json_text[4:].strip()
        
        # Remove any leading/trailing whitespace
        json_text = json_text.strip()
        
        # Attempt to parse the JSON as is
        try:
            json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # Continue with fixes
            pass
        
        # Common JSON fixes
        # Fix missing quotes around property names
        json_text = re.sub(r'([{,]\s*)(\w+)(\s*:)', r'\1"\2"\3', json_text)
        
        # Fix trailing commas in arrays/objects
        json_text = re.sub(r',(\s*[}\]])', r'\1', json_text)
        
        # Fix missing quotes around string values
        # This is trickier and might cause issues with valid JSON if not careful
        
        # Final attempt to parse with Python's JSON parser
        try:
            data = json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # If we're extracting text on failure and can't fix the JSON
            if extract_text_on_failure:
                print("JSON parsing failed - extracting text content instead")
                
                # Extract what appears to be the main text content
                # Look for patterns that would indicate text fields in a JSON response
                text_content = ""
                
                # Try to find text between quotes after "text": or "content": patterns
                text_matches = re.findall(r'"(?:text|content)"\s*:\s*"([^"]+)"', json_text)
                if text_matches:
                    text_content = " ".join(text_matches)
                
                # If that doesn't work, take any text between quotes as potential content
                if not text_content:
                    # Get all quoted strings that look reasonably long (over 20 chars)
                    quoted_text = re.findall(r'"([^"]{20,})"', json_text)
                    if quoted_text:
                        text_content = " ".join(quoted_text)
                
                # If still nothing, just take the longest line that's not mostly symbols
                if not text_content:
                    lines = json_text.split('\n')
                    content_lines = [line for line in lines if len(line) > 50 and 
                                    sum(c.isalpha() or c.isspace() for c in line) / len(line) > 0.7]
                    if content_lines:
                        text_content = " ".join(content_lines)
                
                # Fallback - create a minimal valid JSON with the extracted text
                fallback_json = {
                    "text": text_content or "Extracted text content not found",
                    "title": "Text Extraction Fallback",
                    "topics": ["text extraction"],
                    "summary": "Structured JSON parsing failed - basic text extraction used as fallback.",
                    "entities": [],
                    "has_tables": False,
                    "has_visuals": False,
                    "has_numbers": False,
                    "dates": [],
                    "tables": [],
                    "visuals": [],
                    "numbers": [],
                    "financial_statements": [],
                    "key_metrics": [],
                    "financial_terms": [],
                    "subsections": []
                }
                return json.dumps(fallback_json)
            else:
                # Return minimal valid JSON as before
                return '{"text":"Error parsing JSON response","error":"Invalid JSON structure"}'
            
    except Exception as e:
        print(f"Error cleaning JSON response: {e}")
        
        if extract_text_on_failure:
            # Create a fallback with whatever raw content we have
            raw_text = json_text.strip()
            if len(raw_text) > 10000:  # Truncate very long responses
                raw_text = raw_text[:10000] + "... (truncated)"
                
            fallback_json = {
                "text": raw_text,
                "title": "Raw Text Fallback",
                "topics": ["raw text"],
                "summary": "Response processing failed - returning raw text.",
                "entities": [],
                "has_tables": False,
                "has_visuals": False,
                "has_numbers": False,
                "dates": [],
                "tables": [],
                "visuals": [],
                "numbers": [],
                "financial_statements": [],
                "key_metrics": [],
                "financial_terms": [],
                "subsections": []
            }
            return json.dumps(fallback_json)
        else:
            return '{"text":"Error cleaning JSON response","error":"' + str(e).replace('"', '\\"') + '"}'

def run_async(func, *args, **kwargs):
    """Run an async function from Streamlit."""
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        # If func is already a coroutine, run it directly
        if asyncio.iscoroutine(func):
            return loop.run_until_complete(func)
        # Otherwise call it with the provided args
        return loop.run_until_complete(func(*args, **kwargs))
    finally:
        loop.close()

###############################
# CONTENT EXTRACTION FUNCTIONS
###############################

async def process_page_with_fallback(client, image_part, page_num, filename=None):
    """Process page with multi-level fallback including pure text extraction."""
    try:
        # First attempt: Full structured extraction
        return await extract_page_content_from_memory(client, image_part, page_num)
    except Exception as e:
        print(f"Complex extraction failed for page {page_num}, trying simple extraction: {str(e)}")
        
        try:
            # Second attempt: Simple extraction with basic prompt
            simple_prompt = f"""
            Extract just the basic information from page {page_num}.
            Return a JSON with these fields:
            {{
                "text": "full text content",
                "title": "brief title",
                "topics": ["topic1", "topic2"]
            }}
            """
            
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.0-flash",
                contents=[
                    types.Content(parts=[image_part, types.Part.from_text(text=simple_prompt)]),
                ],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            
            if response.candidates:
                # Use the enhanced clean_json_response with text extraction fallback
                json_text = clean_json_response(response.candidates[0].content.parts[0].text, 
                                              extract_text_on_failure=True)
                data = json.loads(json_text)
                
                return PageContent(
                    page_number=page_num,
                    text=data.get("text", ""),
                    title=data.get("title", f"Page {page_num} (Simple Extraction)"),
                    topics=data.get("topics", ["text extraction"]),
                    summary=data.get("summary", "Extracted using fallback method."),
                    entities=data.get("entities", []),
                    has_tables=data.get("has_tables", False),
                    has_visuals=data.get("has_visuals", False),
                    has_numbers=data.get("has_numbers", False),
                    tables=[],
                    visuals=[],
                    numbers=[],
                    dates=data.get("dates", []),
                    financial_statements=[],
                    key_metrics=[],
                    financial_terms=data.get("financial_terms", []),
                    subsections=[]
                )
        except Exception as fallback_error:
            print(f"Simple extraction failed: {str(fallback_error)}, trying pure text extraction")
            
            # Third attempt: Pure text extraction (no JSON)
            try:
                text_only_prompt = f"Extract ONLY the plain text content from this page. Do NOT format as JSON. Just return the raw text."
                
                text_response = await client.aio.models.generate_content(
                    model="gemini-2.0-flash",
                    contents=[
                        types.Content(parts=[image_part, types.Part.from_text(text=text_only_prompt)]),
                    ]
                )
                
                if text_response.candidates:
                    raw_text = text_response.candidates[0].content.parts[0].text
                    
                    return PageContent(
                        page_number=page_num,
                        text=raw_text,
                        title=f"Page {page_num} (Text Only)",
                        topics=["text extraction"],
                        summary="Only plain text could be extracted from this page.",
                        entities=[],
                        has_tables=False,
                        has_visuals=False,
                        has_numbers=False,
                        tables=[],
                        visuals=[],
                        numbers=[],
                        dates=[],
                        financial_statements=[],
                        key_metrics=[],
                        financial_terms=[],
                        subsections=[]
                    )
            except Exception as text_error:
                print(f"Pure text extraction also failed: {str(text_error)}")
        
        # If all extraction methods fail, return error page
        return create_error_page(page_num, "All extraction methods failed", None)

async def extract_page_content_from_memory(client, image_part, page_num):
    """Extract content from a single page using Gemini with direct memory upload in one API call."""
    # Initialize page data with defaults
    page_data = {
        "page_number": page_num,
        "text": "",
        "title": f"Page {page_num}",
        "topics": [],
        "summary": "",
        "entities": [],
        "has_tables": False,
        "has_visuals": False,
        "has_numbers": False,
        "dates": [],
        "tables": [],
        "visuals": [],
        "numbers": [],
        "financial_statements": [],
        "key_metrics": [],
        "financial_terms": [],
        "subsections": []
    }

    # Combined prompt that handles all extraction in one pass with subsection extraction
    combined_prompt = f"""
    Analyze page {page_num} and return JSON with this exact structure:
    {{
        "text": "full text content",
        "title": "brief title (2-5 words)",
        "topics": ["topic1", "topic2"],
        "summary": "key points summary (3-5 sentences)",
        "entities": ["entity1", "entity2"],
        "has_tables": true/false,
        "has_visuals": true/false,
        "has_numbers": true/false,
        "dates": ["date1", "date2"],
        "financial_terms": ["term1", "term2"],
        "tables": [
            {{
                "table_content": "markdown formatted table with pipes (|) for columns, including headers. Align columns consistently (left, right, or center).",
                "title": "optional table title, if available",
                "summary": "optional table summary highlighting key insights",
                "page_number": {page_num}
            }}
        ],
        "visuals": [
            {{
                "type": "chart/graph type (e.g., bar chart, scatter plot, pie chart, line graph)",
                "description": "description of visual, including purpose and elements",
                "data_summary": "summary of data shown in the visual, highlighting key trends and values",
                "page_numbers": [{page_num}],
                "source_url": "URL of the visual if available, otherwise null",
                "alt_text": "Alternative text description of the visual if available, otherwise null"
            }}
        ],
        "numbers": [
            {{
                "value": "string value (e.g., '123.45', '1,000,000')",
                "description": "what the number represents, including units of measure (e.g., dollars, percentages)",
                "context": "surrounding text context to understand the number's significance"
            }}
        ],
        "key_metrics": [
            {{
                "name": "metric name (e.g., Revenue, Profit Margin)",
                "value": "string value (e.g., '1000000', '15.5')",
                "period": "time period if available (e.g., 'Q1 2023', 'FY2022')",
                "trend": "trend direction if available (e.g., 'increasing', 'decreasing', 'stable')",
                "context": "context if available"
            }}
        ],
        "subsections": [
            {{
                "subsection_id": "page_{page_num}_section_1",
                "order": 1,
                "title": "subsection title in less than 7 words (use existing headings if available)",
                "content": "raw subsection content",
                "is_cutoff": true/false,
                "referenced_visuals": ["visual_id_1", "visual_id_2"],
                "referenced_tables": ["table_id_1", "table_id_2"]
            }}
        ]
    }}

    Rules:
    1. All fields must be included, even if empty. Use [] for empty arrays. Use {{}} for empty JSON Objects
    2. Maintain exact field names and structure as defined above.
    3. Numerical values must be strings (e.g., "123.45"). Do *NOT* enclose strings in single quotes; use double quotes (e.g., "example string").
    4. Dates should be formatted as YYYY-MM-DD. Handle date ranges appropriately.  Ambiguous dates (e.g., 'Q1 2023') are acceptable. Include time information if present.
    5. For sections, consider topic shifts, document structure (headings/subheadings), and paragraph breaks to determine boundaries. Use existing headings as subsection titles when available. Do not create a section for single sentences or phrases unless they delineate distinct concepts.
    6. For empty values or non applicable fields: Use [] for empty arrays. Use {{}} for empty JSON Objects

    7. If OCR quality is low, prioritize extracting complete chunks or sentences, even if entity extraction is imperfect.
    8. Tables should be in markdown format with pipes (|) for columns, including headers. Align columns consistently.
    9. CRITICAL: Ensure all JSON is properly formatted with commas between all elements. Avoid trailing commas.
    10. CRITICAL: Ensure all strings have closing double quotes ("). Escape special characters in strings if needed.
    11. CRITICAL: Return ONLY valid JSON with no additional text.
    12. Validate the final JSON output against the following (example) JSON Schema (you do not need to literally include a schema validator but act as if you had one, catching common data integrity errors).

    EXAMPLE JSON Schema (Conceptual):
    {{
      "type": "object",
      "properties": {{
        "text": {{ "type": "string" }},
        "title": {{ "type": "string" }},
        "topics": {{ "type": "array", "items": {{ "type": "string" }} }},
        "summary": {{ "type": "string" }},
        "entities": {{ "type": "array", "items": {{ "type": "string" }} }},
        "has_tables": {{ "type": "boolean" }},
        "has_visuals": {{ "type": "boolean" }},
        "has_numbers": {{ "type": "boolean" }},
        "dates": {{ "type": "array", "items": {{ "type": "string" }} }},
        "financial_terms": {{ "type": "array", "items": {{ "type": "string" }} }},
        "tables": {{
          "type": "array",
          "items": {{
            "type": "object",
            "properties": {{
              "table_content": {{ "type": "string" }},
              "title": {{ "type": "string" }},
              "summary": {{ "type": "string" }},
              "page_number": {{ "type": "integer" }}
            }},
            "required": ["table_content", "page_number"]
          }}
        }},
        "visuals": {{
          "type": "array",
          "items": {{
            "type": "object",
            "properties": {{
              "type": {{ "type": "string" }},
              "description": {{ "type": "string" }},
              "data_summary": {{ "type": "string" }},
              "page_numbers": {{ "type": "array", "items": {{ "type": "integer" }} }},
              "source_url": {{ "type": ["string", "null"] }},
              "alt_text": {{ "type": ["string", "null"] }}
            }},
            "required": ["type", "description", "data_summary", "page_numbers"]
          }}
        }},
        "numbers": {{
          "type": "array",
          "items": {{
            "type": "object",
            "properties": {{
              "value": {{ "type": "string" }},
              "description": {{ "type": "string" }},
              "context": {{ "type": "string" }}
            }},
            "required": ["value", "description", "context"]
          }}
        }},
         "key_metrics": {{
          "type": "array",
          "items": {{
            "type": "object",
            "properties": {{
              "name": {{ "type": "string" }},
              "value": {{ "type": "string" }},
              "period": {{ "type": "string" }},
              "trend": {{ "type": "string" }},
              "context": {{ "type": "string" }}
            }}
          }}
        }},
        "subsections": {{
          "type": "array",
          "items": {{
            "type": "object",
            "properties": {{
              "subsection_id": {{ "type": "string" }},
              "order": {{ "type": "integer" }},
              "title": {{ "type": "string" }},
              "content": {{ "type": "string" }},
              "is_cutoff": {{ "type": "boolean" }},
              "referenced_visuals": {{ "type": "array", "items": {{ "type": "string" }} }},
              "referenced_tables": {{ "type": "array", "items": {{ "type": "string" }} }}
            }},
            "required": ["subsection_id", "order", "title", "content", "is_cutoff"]
          }}
        }}
      }},
      "required": [
        "text",
        "title",
        "topics",
        "summary",
        "entities",
        "has_tables",
        "has_visuals",
        "has_numbers",
        "dates",
        "financial_terms",
        "tables",
        "visuals",
        "numbers",
        "key_metrics",
        "subsections"
      ]
    }}
    """

    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",  # Using 2.0 for better multimodal understanding
            contents=[
                types.Content(parts=[image_part, types.Part.from_text(text=combined_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            # Handle case where API returns a list instead of a dict
            if isinstance(data, list):
                if data and isinstance(data[0], dict):
                    data = data[0]  # Take the first item if it's a dictionary
                else:
                    data = {}  # Empty dict as fallback
            
            # Update page data with extracted content
            page_data.update({
                "text": data.get("text", ""),
                "title": data.get("title", f"Page {page_num}"),
                "topics": data.get("topics", []),
                "summary": data.get("summary", ""),
                "entities": data.get("entities", []),
                "has_tables": data.get("has_tables", False),
                "has_visuals": data.get("has_visuals", False),
                "has_numbers": data.get("has_numbers", False),
                "dates": data.get("dates", []),
                "financial_terms": data.get("financial_terms", []),
                "tables": [
                    convert_to_model(table, TableData) if isinstance(table, dict) else None
                    for table in data.get("tables", []) if table is not None
                ],
                "visuals": [
                    convert_to_model(visual, VisualElement) if isinstance(visual, dict) else None
                    for visual in data.get("visuals", []) if visual is not None
                ],
                "numbers": [
                    convert_to_model(num, NumericalDataPoint) if isinstance(num, dict) else None
                    for num in data.get("numbers", []) if num is not None
                ],
                "key_metrics": [
                    convert_to_model(metric, FinancialMetric) if isinstance(metric, dict) else None
                    for metric in data.get("key_metrics", []) if metric is not None
                ],
                "subsections": [
                    convert_to_model(subsection, Subsection) if isinstance(subsection, dict) else None
                    for subsection in data.get("subsections", []) if subsection is not None
                ]
            })
            
            # Filter out None values that might have been introduced
            page_data["tables"] = [t for t in page_data["tables"] if t is not None]
            page_data["visuals"] = [v for v in page_data["visuals"] if v is not None]
            page_data["numbers"] = [n for n in page_data["numbers"] if n is not None]
            page_data["key_metrics"] = [m for m in page_data["key_metrics"] if m is not None]
            page_data["subsections"] = [s for s in page_data["subsections"] if s is not None]

    except Exception as e:
        print(f"Error processing page {page_num}: {str(e)}")
        page_data["text"] = f"Error processing page: {str(e)}"
        page_data["title"] = f"Page {page_num} (Error)"

    # Convert to PageContent model with validation
    try:
        return PageContent(**page_data)
    except ValidationError as e:
        print(f"Validation error creating PageContent: {e}")
        return create_error_page(page_num, str(e), e.errors())

def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List] = None) -> PageContent:
    """
    Creates a PageContent object representing an error state.
    The document name association must be handled by the caller.
    """
    # Include context in the error message if available (passed in error_msg)
    error_text = f"Error processing page {page_num}: {error_msg}"

    if validation_errors:
        try:
            error_details = "\nValidation Errors:\n" + "\n".join(
                f"- Field '{'.'.join(map(str, err.get('loc', ['unknown'])))}': {err.get('msg', 'No message')}"
                for err in validation_errors
            )
            error_text += error_details
        except Exception as format_error:
            error_text += f"\n(Could not format validation errors: {format_error})"

    # Initialize the PageContent object with standard fields
    error_page = PageContent(
        page_number=page_num,
        text=error_text,
        title=f"Page {page_num} (Error)",
        topics=["error"],
        summary=f"Failed to process page {page_num}.",
        entities=[],
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        dates=[],
        financial_statements=[],
        key_metrics=[],
        financial_terms=[],
        subsections=[]
    )

    return error_page

async def merge_cutoff_subsections(pages):
    """Merge subsections that are cut off with the next subsection."""
    if not pages:
        return pages
    
    # Sort pages by page number to ensure correct order
    sorted_pages = sorted(pages, key=lambda p: p.page_number)
    
    # Iterate through pages to find and merge cut-off subsections
    for i in range(len(sorted_pages) - 1):
        current_page = sorted_pages[i]
        next_page = sorted_pages[i + 1]
        
        # Check if the current page has subsections and if the last one is cut off
        if not hasattr(current_page, 'subsections') or not current_page.subsections:
            continue
            
        # Sort subsections by order
        current_subsections = sorted(current_page.subsections, key=lambda s: s.order if hasattr(s, 'order') else 0)
        if not current_subsections:
            continue
        
        last_subsection = current_subsections[-1]
        
        # Check if the last subsection is marked as cut off
        if not hasattr(last_subsection, 'is_cutoff') or not last_subsection.is_cutoff:
            continue
        
        # Find the first subsection of the next page
        if not hasattr(next_page, 'subsections') or not next_page.subsections:
            continue
            
        next_subsections = sorted(next_page.subsections, key=lambda s: s.order if hasattr(s, 'order') else 0)
        if not next_subsections:
            continue
        
        first_next_subsection = next_subsections[0]
        
        # Merge the content of the cut-off subsection with the first subsection of the next page
        last_subsection.content += " " + first_next_subsection.content
        last_subsection.is_cutoff = False
        
        # Update references
        if hasattr(first_next_subsection, 'referenced_visuals') and first_next_subsection.referenced_visuals:
            if not hasattr(last_subsection, 'referenced_visuals'):
                last_subsection.referenced_visuals = []
            last_subsection.referenced_visuals.extend(first_next_subsection.referenced_visuals)
            
        if hasattr(first_next_subsection, 'referenced_tables') and first_next_subsection.referenced_tables:
            if not hasattr(last_subsection, 'referenced_tables'):
                last_subsection.referenced_tables = []
            last_subsection.referenced_tables.extend(first_next_subsection.referenced_tables)
        
        # Remove the first subsection from the next page
        next_page.subsections = next_subsections[1:] if len(next_subsections) > 1 else []
        
        # Renumber the remaining subsections in the next page
        for idx, subsection in enumerate(next_page.subsections):
            subsection.order = idx + 1
    
    return sorted_pages

async def extract_chapters_from_subsections(client, pages):
    """Extract chapters from subsections across pages."""
    # Extract all subsections from all pages
    all_subsections = []
    for page in pages:
        if hasattr(page, 'subsections') and page.subsections:
            # Add page_number field to each subsection if not already present
            for subsection in page.subsections:
                if not hasattr(subsection, 'page_number'):
                    subsection.page_number = page.page_number
            all_subsections.extend(page.subsections)
    
    if not all_subsections:
        # No subsections found, create a default chapter with all page content
        default_chapter = Chapter(
            chapter_id="chapter_1",
            title="Document Content",
            summary="Complete document content without chapter structure.",
            subsections=[],
            entity_relationships=[],
            order=1
        )
        
        # Create a subsection for each page
        for page in pages:
            default_subsection = Subsection(
                subsection_id=f"page_{page.page_number}_section_1",
                order=page.page_number,
                title=page.title,
                content=page.text,
                is_cutoff=False,
                referenced_visuals=[],
                referenced_tables=[],
                page_number=page.page_number
            )
            default_chapter.subsections.append(default_subsection)
        
        return [default_chapter]
    
    # Sort subsections by page number and order
    sorted_subsections = sorted(all_subsections, key=lambda s: (s.page_number, s.order))
    
    # Convert subsections to a format suitable for the LLM
    subsections_data = []
    for subsection in sorted_subsections:
        subsections_data.append({
            "subsection_id": subsection.subsection_id,
            "title": subsection.title,
            "page_number": subsection.page_number,
            "order": subsection.order,
            "content_preview": subsection.content[:100] + "..." if len(subsection.content) > 100 else subsection.content
        })
    
    # Create prompt for chapter extraction
    chapter_prompt = f"""
    Analyze these subsections from a document and group them into coherent chapters:
    
    {json.dumps(subsections_data, indent=2)}
    
    Return a JSON array of chapters with this structure:
    [
        {{
            "chapter_id": "chapter_1",
            "title": "Descriptive chapter title (max 7 words)",
            "order": 1,
            "summary": "Brief summary of the chapter's content and purpose",
            "subsection_ids": ["subsection_id1", "subsection_id2"]
        }}
    ]
    
    Rules:
    1. Create between 3-10 logical chapters based on content cohesion
    2. Each chapter must contain at least one subsection
    3. Maintain the original order of subsections when possible
    4. Use descriptive titles that reflect the overall theme of the grouped subsections
    5. Number chapters sequentially starting with 1
    6. Ensure every subsection is assigned to exactly one chapter
    """
    
    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=chapter_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            chapters_data = json.loads(json_text)
            
            # Convert to Chapter objects
            chapters = []
            for chapter_data in chapters_data:
                # Get subsections for this chapter
                chapter_subsections = []
                for subsection_id in chapter_data.get("subsection_ids", []):
                    for subsection in sorted_subsections:
                        if subsection.subsection_id == subsection_id:
                            chapter_subsections.append(subsection)
                            break
                
                # Create the chapter
                chapter = Chapter(
                    chapter_id=chapter_data.get("chapter_id", f"chapter_{len(chapters)+1}"),
                    title=chapter_data.get("title", f"Chapter {len(chapters)+1}"),
                    summary=chapter_data.get("summary", "No summary available."),
                    subsections=chapter_subsections,
                    entity_relationships=[],
                    order=chapter_data.get("order", len(chapters)+1)
                )
                chapters.append(chapter)
            
            # Sort chapters by order
            chapters = sorted(chapters, key=lambda c: c.order)
            
            # If no chapters were created, create a default chapter
            if not chapters:
                default_chapter = Chapter(
                    chapter_id="chapter_1",
                    title="Document Content",
                    summary="Complete document content without chapter structure.",
                    subsections=sorted_subsections,
                    entity_relationships=[],
                    order=1
                )
                chapters = [default_chapter]
            
            return chapters
    except Exception as e:
        print(f"Error extracting chapters: {str(e)}")
        
        # Create a default chapter on error
        default_chapter = Chapter(
            chapter_id="chapter_1",
            title="Document Content",
            summary="Error occurred during chapter extraction.",
            subsections=sorted_subsections,
            entity_relationships=[],
            order=1
        )
        return [default_chapter]

async def analyze_entity_relationships(client, chapter):
    """Analyze entity relationships within a chapter."""
    # Collect content from all subsections
    subsection_texts = []
    for subsection in chapter.subsections:
        subsection_texts.append(f"--- {subsection.title} ---\n{subsection.content}")
    
    # Combine texts while keeping under token limit (approx. 100k chars)
    combined_text = "\n\n".join(subsection_texts)
    if len(combined_text) > 100000:
        combined_text = combined_text[:100000] + "..."
    
    # Create prompt for entity extraction and relationship analysis
    entity_prompt = f"""
    ---Goal---
    Given this chapter text, identify all entities and relationships among the identified entities.
    Use English as output language.

    ---Steps---
    1. Identify key entities. For each identified entity, extract:
       - entity_name: Name of the entity
       - entity_type: One of the following types: [organization, person, geo, event, category, technology, concept]
    
    2. Identify all pairs of (source_entity, target_entity) that are clearly related to each other.
       For each pair of related entities, extract:
       - source_entity: name of the source entity
       - target_entity: name of the target entity
       - relationship_description: explanation of how entities are related
       - relationship_strength: a numeric score from 1-10 indicating strength
       - relationship_keywords: high-level keywords summarizing the relationship
    
    3. Return output as a JSON object with this structure:
    {{
        "entity_relationships": [
            {{
                "source_entity": "Entity A",
                "target_entity": "Entity B",
                "relationship_description": "Description of relationship",
                "relationship_keywords": ["keyword1", "keyword2"],
                "relationship_strength": 8
            }}
        ]
    }}

    Chapter: {chapter.title}
    
    Text:
    {combined_text}
    """
    
    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=entity_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            # Convert to EntityRelationship objects
            relationships = []
            for rel_data in data.get("entity_relationships", []):
                try:
                    relationship = EntityRelationship(
                        source_entity=rel_data.get("source_entity", "Unknown"),
                        target_entity=rel_data.get("target_entity", "Unknown"),
                        relationship_description=rel_data.get("relationship_description", ""),
                        relationship_keywords=rel_data.get("relationship_keywords", []),
                        relationship_strength=float(rel_data.get("relationship_strength", 5))
                    )
                    relationships.append(relationship)
                except Exception as rel_error:
                    print(f"Error creating relationship: {rel_error}")
            
            return relationships
    except Exception as e:
        print(f"Error analyzing entity relationships: {str(e)}")
        return []

async def generate_chapter_summary(client, chapter):
    """Generate a summary for a chapter."""
    # Collect content from all subsections
    subsection_texts = []
    for subsection in chapter.subsections:
        subsection_texts.append(f"--- {subsection.title} ---\n{subsection.content}")
    
    # Combine texts while keeping under token limit
    combined_text = "\n\n".join(subsection_texts)
    if len(combined_text) > 100000:
        combined_text = combined_text[:100000] + "..."
    
    # Create prompt for summary generation
    summary_prompt = f"""
    Generate a comprehensive summary for this chapter from a document:
    
    Chapter Title: {chapter.title}
    
    Content:
    {combined_text}
    
    Return a JSON object with this structure:
    {{
        "summary": "Detailed chapter summary that captures key points, themes, and insights (5-10 sentences)"
    }}
    """
    
    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=summary_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            if "summary" in data:
                chapter.summary = data["summary"]
            
            return chapter
    except Exception as e:
        print(f"Error generating chapter summary: {str(e)}")
        # No change to chapter on error
        return chapter

async def generate_project_ontology(client, documents):
    """Generate project-wide ontology from multiple documents."""
    # Collect metadata from all documents
    documents_data = []
    all_entity_relationships = []
    
    for doc in documents:
        # Extract document info based on structure
        if isinstance(doc, dict) and "raw_extracted_content" in doc:
            raw_content = doc["raw_extracted_content"]
            
            doc_info = {
                "filename": raw_content.get("filename", "Unknown"),
                "title": raw_content.get("summary", {}).get("title", raw_content.get("filename", "Unknown")),
                "themes": raw_content.get("summary", {}).get("themes", []),
                "summary": raw_content.get("summary", {}).get("summary", "No summary available.")
            }
            
            # Extract entity relationships
            entity_relationships = raw_content.get("summary", {}).get("entity_relationships", [])
            
            # Add to collected data
            documents_data.append(doc_info)
            all_entity_relationships.extend(entity_relationships)
    
    # If no documents found, return empty ontology
    if not documents_data:
        return ProjectOntology(
            title="No Documents Available",
            overview="No documents have been processed.",
            document_count=0,
            documents=[],
            global_themes=[],
            entity_relationships=[],
            key_concepts=[]
        )
    
    # Create prompt for project ontology
    ontology_prompt = f"""
    Create a project-wide ontology based on these documents:
    
    {json.dumps(documents_data, indent=2)}
    
    Return a JSON object with this structure:
    {{
        "title": "Project title (max 50 words)",
        "overview": "Comprehensive project overview that synthesizes content across all documents",
        "global_themes": ["theme1", "theme2", "theme3"],
        "key_concepts": ["concept1", "concept2", "concept3"]
    }}
    """
    
    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=ontology_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            # Deduplicate entity relationships
            unique_relationships = {}
            for rel in all_entity_relationships:
                # Convert to dict if it's a model
                if hasattr(rel, 'model_dump'):
                    rel_dict = rel.model_dump()
                else:
                    rel_dict = rel
                
                # Create a key using source and target
                key = f"{rel_dict.get('source_entity', '')}-{rel_dict.get('target_entity', '')}"
                
                # Only add if key doesn't exist or if relationship is stronger
                if key not in unique_relationships or rel_dict.get('relationship_strength', 0) > unique_relationships[key].get('relationship_strength', 0):
                    unique_relationships[key] = rel_dict
            
            # Convert back to EntityRelationship objects
            entity_relationships = []
            for rel_dict in unique_relationships.values():
                try:
                    relationship = EntityRelationship(
                        source_entity=rel_dict.get("source_entity", "Unknown"),
                        target_entity=rel_dict.get("target_entity", "Unknown"),
                        relationship_description=rel_dict.get("relationship_description", ""),
                        relationship_keywords=rel_dict.get("relationship_keywords", []),
                        relationship_strength=float(rel_dict.get("relationship_strength", 5))
                    )
                    entity_relationships.append(relationship)
                except Exception as rel_error:
                    print(f"Error creating relationship: {rel_error}")
            
            # Create ProjectOntology
            return ProjectOntology(
                title=data.get("title", "Project Ontology"),
                overview=data.get("overview", "No overview available."),
                document_count=len(documents_data),
                documents=[doc.get("title") for doc in documents_data],
                global_themes=data.get("global_themes", []),
                entity_relationships=entity_relationships,
                key_concepts=data.get("key_concepts", [])
            )
    except Exception as e:
        print(f"Error generating project ontology: {str(e)}")
        
        # Return basic ontology on error
        return ProjectOntology(
            title="Project Ontology (Error)",
            overview="Error occurred during ontology generation.",
            document_count=len(documents_data),
            documents=[doc.get("title") for doc in documents_data],
            global_themes=[],
            entity_relationships=[],
            key_concepts=[]
        )

###############################
# DOCUMENT PROCESSING FUNCTIONS
###############################

async def process_single_page_with_semaphore(semaphore, client, page_info: dict, uploaded_files):
    """
    Process a single page using semaphore for rate limiting.
    Returns a dictionary: {"info": page_info, "result": PageContent object}.
    Handles potential errors during processing.
    """
    # Get doc_name and page_num early for use in logic and error reporting
    doc_name = page_info.get("doc_name", "Unknown_Document")
    page_num = page_info.get("page_num", 0) # Default to 0 if missing

    async with semaphore:
        result_page_obj = None # Initialize result
        try:
            # Call process_single_pdf_page, which returns a DICTIONARY
            page_dict = await process_single_pdf_page(client, page_info, uploaded_files)

            # Ensure essential keys for Pydantic initialization exist
            page_dict.setdefault("page_number", page_num)
            # Add defaults for required fields if not present in the dict returned by Gemini
            # This ensures PageContent(**page_dict) is less likely to fail basic validation
            # (though Gemini *should* return them based on the prompt)
            for field, default_value in PageContent.model_fields.items():
                 if default_value.is_required() and field not in page_dict:
                      # Try to provide a sensible default based on annotation
                      if default_value.annotation == str: page_dict.setdefault(field, "")
                      elif default_value.annotation == int: page_dict.setdefault(field, 0)
                      elif default_value.annotation == bool: page_dict.setdefault(field, False)
                      elif default_value.annotation == list or getattr(default_value.annotation, '__origin__', None) == list: page_dict.setdefault(field, [])
                      else: page_dict.setdefault(field, None) # Fallback

            # Initialize PageContent with the dictionary
            result_page_obj = PageContent(**page_dict)

        except ValidationError as ve:
            st.warning(f"Pydantic validation error creating PageContent for page {page_num} of {doc_name} AFTER initial processing: {ve}")
            # Create error page (pass doc_name in the message)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Pydantic Validation Error in '{doc_name}': {ve}", # Include doc_name in msg
                validation_errors=ve.errors()
            )

        except Exception as e:
            st.error(f"Unexpected error in semaphore task for page {page_num} of {doc_name}: {str(e)}")
            import traceback
            st.error(f"Traceback: {traceback.format_exc()}")
            # Create error page (pass doc_name in the message)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Unexpected processing error in semaphore task for '{doc_name}': {str(e)}" # Include doc_name in msg
            )

        # Return the result packaged with the original info
        return {"info": page_info, "result": result_page_obj}

async def process_single_pdf_page(client, page_info: dict, uploaded_files) -> dict:
    """
    Process a single PDF page image using Gemini with fallback options, returning extracted data as a dictionary.
    Handles errors by returning a dictionary representing an error page structure.
    The returned dictionary does NOT include 'doc_name' as a key.
    """
    page_num = page_info.get("page_num", 0)
    doc_name = page_info.get("doc_name", "Unknown_Document")  # Keep track of it locally
    page_dict = {}  # Initialize the dictionary to be returned

    try:
        # --- Ensure required info is present ---
        if "image_b64" not in page_info or not page_info["image_b64"]:
             # Create error object, convert to dict, and return
             st.warning(f"Missing image data for page {page_num} of {doc_name}.")
             error_page_obj = create_error_page(
                 page_num=page_num,
                 error_msg="Missing or empty image data (image_b64)",
                 validation_errors=None
             )
             return error_page_obj.model_dump()  # Return dict representation

        if "page_num" not in page_info:
             st.warning(f"Missing page number for a page in {doc_name}.")
             # Use 0 or handle as appropriate, create error object/dict
             error_page_obj = create_error_page(
                 page_num=0,  # Or decide on a better default/error indicator
                 error_msg="Missing page number information",
                 validation_errors=None
             )
             return error_page_obj.model_dump()

        # --- Prepare for Gemini Call ---
        img_data = base64.b64decode(page_info["image_b64"])
        image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")  # Assuming JPEG, adjust if needed

        # --- Call Gemini with fallback strategy instead of direct extraction ---
        page_content_obj = await process_page_with_fallback(
            client, image_part, page_num, doc_name  # Pass page_num and doc_name
        )

        # --- Convert successful PageContent object to a dictionary ---
        # This dictionary contains only the fields defined in the PageContent model
        page_dict = page_content_obj.model_dump()

        return page_dict  # Return the clean dictionary

    except Exception as e:
        st.error(f"Error during Gemini processing or data handling for page {page_num} of {doc_name}: {str(e)}")
        # Create an error PageContent object first
        error_page_obj = create_error_page(
            page_num=page_num,
            error_msg=f"Gemini API or processing error: {str(e)}",
            validation_errors=None
        )
        # Convert the error PageContent object to a dictionary to maintain return type consistency
        page_dict = error_page_obj.model_dump()

        return page_dict

async def finalize_document(client, doc_name, pages):
    """Create document summary and final structure with proper model handling and chapter extraction."""
    try:
        # Convert page data to dictionaries first
        validated_pages = []
        for page in pages:
            try:
                # Convert to dict if model
                page_dict = page.model_dump() if hasattr(page, 'model_dump') else page
                
                # Ensure required fields
                page_dict.setdefault("title", f"Page {page_dict.get('page_number', 0)}")
                page_dict.setdefault("text", "")
                
                validated_pages.append(PageContent(**page_dict))
            except Exception as e:
                st.warning(f"Error validating page: {e}")
                validated_pages.append(PageContent(
                    page_number=page.get("page_number", len(validated_pages)+1),
                    text=f"Validation error: {str(e)}",
                    title=f"Page {len(validated_pages)+1} (Invalid)"
                ))
        
        # Process subsections
        
        # 1. Merge cut-off subsections
        merged_pages = await merge_cutoff_subsections(validated_pages)
        
        # 2. Extract chapters from subsections
        chapters = await extract_chapters_from_subsections(client, merged_pages)
        
        # 3. Process each chapter for entity relationships and summaries
        for i, chapter in enumerate(chapters):
            # Analyze entity relationships
            relationships = await analyze_entity_relationships(client, chapter)
            chapter.entity_relationships = relationships
            
            # Generate chapter summary
            updated_chapter = await generate_chapter_summary(client, chapter)
            chapters[i] = updated_chapter

        # Generate document summary as dict
        summary_dict = await generate_financial_summary(merged_pages, doc_name, client)
        
        # Add chapters to summary
        summary_dict['chapters'] = [chapter.model_dump() for chapter in chapters]
        
        # Extract all entity relationships from chapters for document-level relationships
        document_relationships = []
        for chapter in chapters:
            document_relationships.extend([rel.model_dump() for rel in chapter.entity_relationships])
        summary_dict['entity_relationships'] = document_relationships
        
        # Create the final result structure
        result = {
            "raw_extracted_content": {
                "filename": doc_name,
                "pages": [p.model_dump() for p in merged_pages],
                "summary": summary_dict
            }
        }
        
        return result
        
    except Exception as e:
        print(f"Error finalizing document {doc_name}: {str(e)}")
        # Return basic document structure on error
        return {
            "raw_extracted_content": {
                "filename": doc_name,
                "pages": [p.model_dump() if hasattr(p, 'model_dump') else p for p in pages],
                "summary": None
            }
        }

async def extract_pages_from_pdf_bytes(pdf_bytes, file_name):
    """Extract pages from PDF bytes using PyMuPDF."""
    loop = asyncio.get_running_loop()
    pages = []
    try:
        pdf_stream = io.BytesIO(pdf_bytes)
        # Run synchronous PyMuPDF code in an executor thread
        def sync_extract():
            doc = pymupdf.open(stream=pdf_stream, filetype="pdf")
            extracted = []
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    # Moderate DPI for balance
                    pix = page.get_pixmap(dpi=150, alpha=False) # No alpha needed for JPEG
                    img_bytes = pix.tobytes("jpeg", jpg_quality=85) # Control quality
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    extracted.append({
                        "page_num": page_num + 1,
                        "image_b64": img_b64,
                        "doc_name": file_name,
                        # "doc_bytes": None # Avoid storing large bytes here
                    })
                except Exception as page_e:
                     st.warning(f"Error extracting page {page_num+1} from {file_name}: {page_e}")
            doc.close()
            return extracted

        pages = await loop.run_in_executor(None, sync_extract)
        return pages

    except Exception as e:
        st.error(f"Error opening or processing PDF {file_name}: {str(e)}")
        return [] # Return empty list on error

async def generate_financial_summary(pages: List[Union[Dict, PageContent]], filename: str, client) -> Dict:
    """Generate enhanced financial summary using Gemini and return as dict."""
    # Prepare summary input
    summary_input = {
        "page_titles": [],
        "key_metrics": [],
        "financial_terms": [],
        "has_tables": False,
        "has_visuals": False
    }

    # Track if we have any tables or visuals
    has_tables = False
    has_visuals = False

    for page in pages:
        # Handle both dict and PageContent instances
        if isinstance(page, PageContent):
            page_dict = page.model_dump()
        else:
            page_dict = page

        # Add page title
        summary_input["page_titles"].append(page_dict.get("title", f"Page {page_dict.get('page_number', '?')}"))

        # Collect metrics - handle both dict and model instances
        for metric in page_dict.get("key_metrics", []):
            if hasattr(metric, 'model_dump'):  # If it's a Pydantic model
                metric_dict = metric.model_dump()
            else:
                metric_dict = metric
            
            summary_input["key_metrics"].append({
                **metric_dict,
                "page": page_dict.get("page_number", 0)
            })

        # Collect other data
        summary_input["financial_terms"].extend(page_dict.get("financial_terms", []))
        
        # Track if we have tables or visuals
        if page_dict.get("has_tables", False):
            has_tables = True
        if page_dict.get("has_visuals", False):
            has_visuals = True

    summary_input["has_tables"] = has_tables
    summary_input["has_visuals"] = has_visuals

    # Generate summary with Gemini
    try:
        summary_response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_text(text=f"""
                            Create a financial summary for "{filename}" with this structure:
                            {{
                                "title": "concise title",
                                "themes": ["theme1", "theme2"],
                                "questions": ["question1", "question2"],
                                "summary": "detailed summary",
                                "tables_summary": "summary of key tables",
                                "visuals_summary": "summary of key visuals"
                            }}
                            Input Data:
                            {json.dumps(summary_input, indent=2)}
                        """)
                    ]
                )
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json"
            )
        )
        
        if summary_response.candidates:
            json_text = summary_response.candidates[0].content.parts[0].text
            json_text = clean_json_response(json_text)
            summary_data = json.loads(json_text)
            
            # Handle case where API returns a list
            if isinstance(summary_data, list):
                summary_data = summary_data[0] if summary_data else {}
            
            return summary_data
            
    except Exception as e:
        print(f"Summary generation error: {str(e)}")
    
    # Fallback summary as dict
    return {
        "title": f"Summary: {filename}",
        "themes": ["analysis"],
        "questions": ["What are the key points in this document?"],
        "summary": "Generated summary unavailable",
        "tables_summary": "",
        "visuals_summary": ""
    }

# Non-PDF Document Processing Functions

async def process_single_document_memory(client, file, uploaded_files):
    """Process a single non-PDF document from memory with comprehensive extraction."""
    try:
        file_name = file["name"]
        file_type = file["type"].lower()
        file_content = file["content"]
        
        # Process based on file type
        if file_type in ["xlsx", "xls"]:
            # Excel processing
            content = await run_in_executor(process_excel_memory, file_content)
            processed_pages = await process_tabular_data(client, content, file_name, "excel")
        elif file_type == "csv":
            # CSV processing
            content = await run_in_executor(process_csv_memory, file_content)
            processed_pages = await process_tabular_data(client, content, file_name, "csv")
        elif file_type == "docx":
            # Word processing
            processed_pages = await process_word_document(client, file_content, file_name)
        elif file_type == "pptx":
            # PowerPoint processing
            processed_pages = await process_pptx_document(client, file_content, file_name)
        elif file_type in ["txt", "md", "json", "html", "xml"]:
            # Text-based document processing
            processed_pages = await process_text_document(client, {"content": file_content, "name": file_name, "type": file_type})
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        if not processed_pages:
            raise Exception(f"No pages were processed in {file_name}")
            
        # Now process the pages to extract chapters and entity relationships
        # Merge cut-off subsections
        merged_pages = await merge_cutoff_subsections(processed_pages)
        
        # Extract chapters from subsections
        chapters = await extract_chapters_from_subsections(client, merged_pages)
        
        # Process each chapter for entity relationships and summaries
        for i, chapter in enumerate(chapters):
            # Analyze entity relationships using improved method aligned with prompt
            relationships = await analyze_entity_relationships_enhanced(client, chapter)
            chapter.entity_relationships = relationships
            
            # Generate chapter summary
            updated_chapter = await generate_chapter_summary(client, chapter)
            chapters[i] = updated_chapter

        # Generate document summary
        summary_dict = await generate_document_summary(client, merged_pages, file_name)
        
        # Add chapters to summary
        summary_dict['chapters'] = [chapter.model_dump() for chapter in chapters]
        
        # Extract all entity relationships from chapters for document-level relationships
        document_relationships = []
        for chapter in chapters:
            document_relationships.extend([rel.model_dump() for rel in chapter.entity_relationships])
        summary_dict['entity_relationships'] = document_relationships
        
        # Return structured document result
        return {
            "raw_extracted_content": {
                "filename": file_name,
                "pages": [p.model_dump() for p in merged_pages],
                "summary": summary_dict
            }
        }
    except Exception as e:
        st.error(f"Error processing document {file_name}: {str(e)}")
        # Return minimal result with error
        return {
            "raw_extracted_content": {
                "filename": file_name,
                "pages": [],
                "summary": {
                    "title": f"Error Processing {file_name}",
                    "themes": ["error"],
                    "questions": ["What went wrong during processing?"],
                    "summary": f"Error occurred during processing: {str(e)}",
                    "chapters": [],
                    "entity_relationships": []
                },
                "error": str(e)
            }
        }

async def process_tabular_data(client, content, filename, filetype):
    """Process tabular data (Excel, CSV) with enhanced content extraction."""
    pages = []
    
    try:
        # Create basic page structure
        page = PageContent(
            page_number=1,
            text=f"Tabular data from {filename}",
            title=f"Data: {filename}",
            topics=[f"{filetype} data", "tabular data"],
            has_tables=True,
            has_numbers=True,
            subsections=[]
        )
        
        # Process sheet data
        if filetype == "excel" and isinstance(content, dict):
            # Process each sheet as a section
            section_order = 1
            for sheet_name, sheet_data in content.items():
                # Create table content in markdown format
                table_content = create_markdown_table(sheet_data.get("headers", []), sheet_data.get("data", []))
                
                # Create table object
                table_id = f"table_1_{section_order}"
                table = TableData(
                    table_content=table_content,
                    title=f"Sheet: {sheet_name}",
                    summary=f"Data from sheet {sheet_name} in {filename}",
                    page_number=1,
                    table_id=table_id
                )
                page.tables.append(table)
                
                # Create subsection for this sheet
                subsection = Subsection(
                    subsection_id=f"page_1_section_{section_order}",
                    order=section_order,
                    title=f"Sheet: {sheet_name}",
                    content=f"Data from sheet {sheet_name} in {filename}\n\n{table_content[:500]}...",
                    is_cutoff=False,
                    referenced_tables=[table_id],
                    referenced_visuals=[],
                    page_number=1
                )
                page.subsections.append(subsection)
                section_order += 1
                
                # Extract key metrics, entities, etc. from this sheet
                await extract_sheet_insights(client, sheet_data, page)
                
        elif filetype == "csv" and "data" in content:
            # Process CSV data
            headers = content.get("headers", [])
            data = content.get("data", [])
            
            # Create table content
            table_content = create_markdown_table(headers, data)
            
            # Create table object
            table_id = "table_1_1"
            table = TableData(
                table_content=table_content,
                title=f"CSV Data: {filename}",
                summary=f"Data from {filename}",
                page_number=1,
                table_id=table_id
            )
            page.tables.append(table)
            
            # Create a subsection for the CSV data
            subsection = Subsection(
                subsection_id=f"page_1_section_1",
                order=1,
                title=f"CSV Data",
                content=f"Data from {filename}\n\n{table_content[:500]}...",
                is_cutoff=False,
                referenced_tables=[table_id],
                referenced_visuals=[],
                page_number=1
            )
            page.subsections.append(subsection)
            
            # Extract key metrics, entities, etc. from this data
            await extract_sheet_insights(client, {"headers": headers, "data": data}, page)
        
        # Add the page to our pages list
        pages.append(page)
        
        # Use Gemini to generate summary, topic classification, etc.
        enhanced_pages = await enhance_tabular_data_pages(client, pages, filename)
        return enhanced_pages
        
    except Exception as e:
        st.error(f"Error processing tabular data: {str(e)}")
        # Return basic error page
        return [PageContent(
            page_number=1,
            text=f"Error processing tabular data: {str(e)}",
            title=f"Error: {filename}",
            topics=["error"],
            summary=f"Failed to process tabular data due to: {str(e)}",
            subsections=[]
        )]

def create_markdown_table(headers, data):
    """Create a markdown table from headers and data."""
    if not headers or not data:
        return "No data available"
    
    # Format headers
    header_row = "| " + " | ".join(str(h) for h in headers) + " |"
    separator_row = "| " + " | ".join("---" for _ in headers) + " |"
    
    # Format data rows
    data_rows = []
    for row in data:
        if isinstance(row, dict):
            # Handle dictionary data (e.g., from CSV)
            values = [str(row.get(h, "")) for h in headers]
            data_rows.append("| " + " | ".join(values) + " |")
        else:
            # Handle list data
            values = [str(cell) if cell is not None else "" for cell in row] if isinstance(row, (list, tuple)) else [str(row)]
            data_rows.append("| " + " | ".join(values) + " |")
    
    # Combine all rows
    return "\n".join([header_row, separator_row] + data_rows)

async def extract_sheet_insights(client, sheet_data, page):
    """Extract insights from tabular data."""
    try:
        # Extract numerical data points
        if "data" in sheet_data:
            for row in sheet_data["data"]:
                if isinstance(row, dict):
                    for key, value in row.items():
                        if isinstance(value, (int, float)) or (isinstance(value, str) and value.replace('.', '', 1).isdigit()):
                            # Add as numerical data point
                            page.numbers.append(NumericalDataPoint(
                                value=str(value),
                                description=f"Value of {key}",
                                context=f"From {key} column"
                            ))
                            page.has_numbers = True
        
        # Set flags based on content
        page.has_tables = bool(page.tables)
        page.has_numbers = page.has_numbers or bool(page.numbers)
        
    except Exception as e:
        print(f"Error extracting sheet insights: {str(e)}")

async def enhance_tabular_data_pages(client, pages, filename):
    """Use Gemini to enhance tabular data pages with summaries, topics, etc."""
    try:
        enhanced_pages = []
        
        for page in pages:
            # Extract tables and data for context
            tables_data = []
            for table in page.tables:
                tables_data.append({
                    "title": table.title if hasattr(table, 'title') else "Untitled Table",
                    "content": table.table_content if hasattr(table, 'table_content') else ""
                })
            
            # Create a prompt for summarization
            context = f"Filename: {filename}\n\n"
            for i, table in enumerate(tables_data):
                context += f"Table {i+1}: {table['title']}\n"
                # Limit table content to avoid token overflow
                context += f"{table['content'][:1000]}...\n\n" if len(table['content']) > 1000 else f"{table['content']}\n\n"
            
            prompt = f"""
            Analyze this tabular data and return a JSON object with:
            1. A brief title for the document (2-5 words)
            2. Key topics discussed (list of keywords)
            3. A summary of key insights (3-5 sentences)
            4. Entities mentioned (list of strings)
            5. Key numerical findings (list of descriptions)
            
            Data Context:
            {context}
            
            Format your response as a valid JSON object with these exact fields:
            {{
                "title": "Brief title",
                "topics": ["topic1", "topic2"],
                "summary": "Key insights summary",
                "entities": ["entity1", "entity2"],
                "key_findings": ["finding1", "finding2"]
            }}
            """
            
            try:
                # Get enhancement from Gemini
                response = await retry_api_call(
                    client.aio.models.generate_content,
                    model="gemini-2.0-flash",
                    contents=[
                        types.Content(parts=[types.Part.from_text(text=prompt)]),
                    ],
                    config=types.GenerateContentConfig(response_mime_type="application/json")
                )
                
                if response.candidates:
                    json_text = clean_json_response(response.candidates[0].content.parts[0].text)
                    data = json.loads(json_text)
                    
                    # Update page with enhanced data
                    page.title = data.get("title", page.title)
                    page.topics = data.get("topics", page.topics)
                    page.summary = data.get("summary", page.summary)
                    page.entities = data.get("entities", page.entities)
                    
                    # Add numerical findings
                    for finding in data.get("key_findings", []):
                        if not any(n.description == finding for n in page.numbers):
                            page.numbers.append(NumericalDataPoint(
                                value="N/A",
                                description=finding,
                                context="Extracted insight"
                            ))
            except Exception as e:
                print(f"Error enhancing tabular data: {str(e)}")
                # Continue with original page
            
            enhanced_pages.append(page)
        
        return enhanced_pages
        
    except Exception as e:
        print(f"Error in enhance_tabular_data_pages: {str(e)}")
        return pages

async def process_word_document(client, file_content: bytes, filename: str):
    """Process Word document with improved subsection extraction."""
    try:
        # Create in-memory document
        doc = docx.Document(io.BytesIO(file_content))
        
        # Extract paragraphs and headings
        paragraphs = []
        current_heading = None
        sections = []
        current_section = {"title": "Introduction", "content": [], "level": 0}
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
                
            paragraphs.append(text)
            
            # Check if it's a heading
            if paragraph.style.name.startswith('Heading'):
                # Extract heading level
                heading_level = int(paragraph.style.name.replace('Heading', '')) if paragraph.style.name != 'Heading' else 1
                
                # Save previous section if exists
                if current_section["content"]:
                    sections.append(current_section)
                
                # Start new section
                current_section = {"title": text, "content": [], "level": heading_level}
                current_heading = text
            else:
                # Add to current section
                current_section["content"].append(text)
        
        # Add final section
        if current_section["content"]:
            sections.append(current_section)
            
        # Detect page breaks and convert sections to pages
        # Word doesn't store page break info reliably in the docx format,
        # so we'll use sections as a proxy for page content
        
        # Group sections into pages (using heuristics)
        pages = []
        current_page = []
        current_page_section_count = 0
        
        for section in sections:
            current_page.append(section)
            current_page_section_count += 1
            
            # Heuristic: start a new page after several sections or a lot of content
            content_length = sum(len(" ".join(s["content"])) for s in current_page)
            if current_page_section_count >= 3 or content_length > 3000:
                pages.append(current_page)
                current_page = []
                current_page_section_count = 0
        
        # Add any remaining sections
        if current_page:
            pages.append(current_page)
            
        # Now process each "page" (group of sections) into PageContent objects
        processed_pages = []
        for i, page_sections in enumerate(pages, 1):
            # Combine text for this page
            page_text = ""
            subsections = []
            section_order = 1
            
            for section in page_sections:
                title = section["title"]
                content = " ".join(section["content"])
                level = section["level"]
                
                # Add to page text
                if level > 0:  # It's a heading
                    page_text += f"{'#' * level} {title}\n\n"
                page_text += content + "\n\n"
                
                # Create subsection
                subsection = Subsection(
                    subsection_id=f"page_{i}_section_{section_order}",
                    order=section_order,
                    title=title,
                    content=content,
                    is_cutoff=False,  # Word processor files rarely have content cutoff
                    referenced_visuals=[],
                    referenced_tables=[],
                    page_number=i
                )
                subsections.append(subsection)
                section_order += 1
            
            # Extract tables from this page section
            tables = []
            for table_index, table in enumerate(doc.tables):
                # We can't tell exactly which page a table is on, so assign it to the first page
                # In a real implementation, this would need more sophisticated handling
                table_id = f"table_{i}_{table_index+1}"
                table_content = extract_word_table(table)
                
                tables.append(TableData(
                    table_content=table_content,
                    title=f"Table {table_index+1}",
                    summary="",  # Will be filled by Gemini
                    page_number=i,
                    table_id=table_id
                ))
            
            # Create page with basic content
            page = PageContent(
                page_number=i,
                text=page_text,
                title=f"Page {i}" if not page_sections else page_sections[0]["title"],
                topics=[],  # Will be filled by Gemini
                summary="",  # Will be filled by Gemini
                entities=[],
                has_tables=bool(tables),
                has_visuals=False,  # Word processing docs typically don't have visuals
                has_numbers=False,  # Will be determined by content analysis
                tables=tables,
                visuals=[],
                numbers=[],
                dates=[],
                financial_statements=[],
                key_metrics=[],
                financial_terms=[],
                subsections=subsections
            )
            
            # Use Gemini to extract more info
            enhanced_page = await enhance_text_page(client, page, filename)
            processed_pages.append(enhanced_page)
                
        return processed_pages
        
    except Exception as e:
        st.error(f"Error processing Word document {filename}: {str(e)}")
        return [PageContent(
            page_number=1,
            text=f"Error processing Word document: {str(e)}",
            title=f"Error: {filename}",
            topics=["error"],
            summary=f"Failed to process Word document due to: {str(e)}",
            subsections=[]
        )]

def extract_word_table(table):
    """Extract table data from a Word table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    
    if not rows:
        return ""
    
    # Create markdown table
    headers = rows[0]
    data = rows[1:]
    
    return create_markdown_table(headers, data)

async def process_pptx_document(client, file_content: bytes, filename: str):
    """Process PowerPoint files with enhanced slide content extraction."""
    try:
        # Create in-memory presentation
        ppt = Presentation(io.BytesIO(file_content))
        
        # Extract slides
        processed_pages = []
        for i, slide in enumerate(ppt.slides, 1):
            # Extract slide title
            title = "Untitled Slide"
            for shape in slide.shapes:
                if hasattr(shape, "is_title") and shape.is_title and hasattr(shape, "text"):
                    title = shape.text.strip()
                    break
            
            # Extract all text content
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text and text != title:  # Avoid duplicating title
                        slide_text.append(text)
            
            # Extract tables
            tables = []
            table_count = 0
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1
                    table_id = f"table_{i}_{table_count}"
                    
                    # Extract table content
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    
                    # Create markdown table
                    if rows:
                        table_content = create_markdown_table(rows[0], rows[1:])
                        tables.append(TableData(
                            table_content=table_content,
                            title=f"Table on slide {i}",
                            summary="",  # Will be filled by Gemini
                            page_number=i,
                            table_id=table_id
                        ))
            
            # Create subsections (title and content)
            subsections = []
            if slide_text:
                subsection = Subsection(
                    subsection_id=f"page_{i}_section_1",
                    order=1,
                    title=title,
                    content="\n\n".join(slide_text),
                    is_cutoff=False,
                    referenced_visuals=[],
                    referenced_tables=[f"table_{i}_{j+1}" for j in range(table_count)],
                    page_number=i
                )
                subsections.append(subsection)
            
            # Create page with basic content
            page = PageContent(
                page_number=i,
                text="\n\n".join([title] + slide_text),
                title=title,
                topics=[],  # Will be filled by Gemini
                summary="",  # Will be filled by Gemini
                entities=[],
                has_tables=bool(tables),
                has_visuals=False,  # We don't extract visuals yet
                has_numbers=False,  # Will be determined by content analysis
                tables=tables,
                visuals=[],
                numbers=[],
                dates=[],
                financial_statements=[],
                key_metrics=[],
                financial_terms=[],
                subsections=subsections
            )
            
            # Use Gemini to extract more info
            enhanced_page = await enhance_text_page(client, page, filename)
            processed_pages.append(enhanced_page)
        
        return processed_pages
        
    except Exception as e:
        st.error(f"Error processing PowerPoint file {filename}: {str(e)}")
        return [PageContent(
            page_number=1,
            text=f"Error processing PowerPoint file: {str(e)}",
            title=f"Error: {filename}",
            topics=["error"],
            summary=f"Failed to process PowerPoint file due to: {str(e)}",
            subsections=[]
        )]

async def process_text_document(client, file_data: dict):
    """Process text document with intelligent structure extraction."""
    file_buffer = io.BytesIO(file_data["content"])
    filename = file_data["name"]
    file_type = file_data.get("type", "txt")
    
    try:
        # Read content
        text_content = file_buffer.read().decode('utf-8')
        
        # Extract structure based on file type
        if file_type == "md":
            # For markdown, use headers to identify sections
            return await process_markdown_content(client, text_content, filename)
        elif file_type in ["json", "html", "xml"]:
            # For structured text formats, use special handling
            return await process_structured_text(client, text_content, filename, file_type)
        else:
            # For plain text, look for page breaks or section indicators
            return await process_plain_text(client, text_content, filename)
            
    except Exception as e:
        st.error(f"Error processing text document {filename}: {str(e)}")
        return [PageContent(
            page_number=1,
            text=f"Error processing text document: {str(e)}",
            title=f"Error: {filename}",
            topics=["error"],
            summary=f"Failed to process text document due to: {str(e)}",
            subsections=[]
        )]
    finally:
        file_buffer.close()

async def process_markdown_content(client, text_content, filename):
    """Process markdown content with header-based sectioning."""
    # Split by headers
    lines = text_content.split("\n")
    sections = []
    current_section = {"title": "Introduction", "content": [], "level": 0}
    
    for line in lines:
        if line.startswith("#"):
            # Count leading # to determine heading level
            level = 0
            for char in line:
                if char == '#':
                    level += 1
                else:
                    break
            
            title = line[level:].strip()
            
            # Save previous section if it has content
            if current_section["content"]:
                sections.append(current_section)
            
            # Start new section
            current_section = {"title": title, "content": [], "level": level}
        else:
            current_section["content"].append(line)
    
    # Add final section
    if current_section["content"]:
        sections.append(current_section)
    
    # Group sections into pages (similar to Word document processing)
    pages = []
    current_page = []
    current_page_section_count = 0
    
    for section in sections:
        current_page.append(section)
        current_page_section_count += 1
        
        # Heuristic: start a new page after major heading or lots of content
        content_length = sum(len("\n".join(s["content"])) for s in current_page)
        if section["level"] == 1 or content_length > 3000:
            pages.append(current_page)
            current_page = []
            current_page_section_count = 0
    
    # Add any remaining sections
    if current_page:
        pages.append(current_page)
    
    # Process pages
    processed_pages = []
    for i, page_sections in enumerate(pages, 1):
        # Combine text for this page
        page_text = ""
        subsections = []
        section_order = 1
        
        for section in page_sections:
            title = section["title"]
            content = "\n".join(section["content"])
            level = section["level"]
            
            # Add to page text
            if level > 0:  # It's a heading
                page_text += f"{'#' * level} {title}\n\n"
            page_text += content + "\n\n"
            
            # Create subsection
            subsection = Subsection(
                subsection_id=f"page_{i}_section_{section_order}",
                order=section_order,
                title=title,
                content=content,
                is_cutoff=False,
                referenced_visuals=[],
                referenced_tables=[],
                page_number=i
            )
            subsections.append(subsection)
            section_order += 1
        
        # Create page with basic content
        page = PageContent(
            page_number=i,
            text=page_text,
            title=f"Page {i}" if not page_sections else page_sections[0]["title"],
            topics=[],
            summary="",
            entities=[],
            has_tables=False,  # Will be determined by Gemini
            has_visuals=False,
            has_numbers=False,
            tables=[],
            visuals=[],
            numbers=[],
            dates=[],
            financial_statements=[],
            key_metrics=[],
            financial_terms=[],
            subsections=subsections
        )
        
        # Use Gemini to extract more info
        enhanced_page = await enhance_text_page(client, page, filename)
        processed_pages.append(enhanced_page)
    
    return processed_pages

async def process_structured_text(client, text_content, filename, file_type):
    """Process structured text formats (JSON, HTML, XML)."""
    # For structured formats, we'll extract basic structure and then let Gemini help
    try:
        # Create a single page for the structured content
        page = PageContent(
            page_number=1,
            text=text_content[:10000] + ("..." if len(text_content) > 10000 else ""),  # Truncate very long content
            title=f"{file_type.upper()} Document: {filename}",
            topics=[file_type],
            summary="",  # Will be filled by Gemini
            entities=[],
            has_tables=False,
            has_visuals=False,
            has_numbers=False,
            tables=[],
            visuals=[],
            numbers=[],
            dates=[],
            financial_statements=[],
            key_metrics=[],
            financial_terms=[],
            subsections=[]
        )
        
        # For JSON, parse and extract key structures
        if file_type == "json":
            try:
                json_data = json.loads(text_content)
                
                # Extract top-level keys as sections
                section_order = 1
                for key, value in json_data.items():
                    # Format the value based on type
                    if isinstance(value, dict):
                        content = json.dumps(value, indent=2)
                    elif isinstance(value, list):
                        content = json.dumps(value, indent=2)
                    else:
                        content = str(value)
                    
                    # Create subsection
                    subsection = Subsection(
                        subsection_id=f"page_1_section_{section_order}",
                        order=section_order,
                        title=key,
                        content=content,
                        is_cutoff=False,
                        referenced_visuals=[],
                        referenced_tables=[],
                        page_number=1
                    )
                    page.subsections.append(subsection)
                    section_order += 1
            except json.JSONDecodeError:
                # Not valid JSON, create a single section
                subsection = Subsection(
                    subsection_id="page_1_section_1",
                    order=1,
                    title="JSON Content",
                    content=text_content[:5000] + ("..." if len(text_content) > 5000 else ""),
                    is_cutoff=False,
                    referenced_visuals=[],
                    referenced_tables=[],
                    page_number=1
                )
                page.subsections.append(subsection)
        
        # For HTML, extract title, headings, etc.
        elif file_type == "html":
            # Simple regex-based extraction of HTML structure (in real app, use BeautifulSoup)
            title_match = re.search(r"<title>(.*?)</title>", text_content, re.IGNORECASE | re.DOTALL)
            if title_match:
                page.title = title_match.group(1).strip()
            
            # Extract headings to create sections
            heading_matches = re.findall(r"<h(\d)>(.*?)</h\1>", text_content, re.IGNORECASE | re.DOTALL)
            
            if heading_matches:
                # Use headings as section titles
                section_order = 1
                for level, heading in heading_matches:
                    # Create subsection
                    subsection = Subsection(
                        subsection_id=f"page_1_section_{section_order}",
                        order=section_order,
                        title=heading.strip(),
                        content=f"Section with heading: {heading.strip()}",  # Just placeholder
                        is_cutoff=False,
                        referenced_visuals=[],
                        referenced_tables=[],
                        page_number=1
                    )
                    page.subsections.append(subsection)
                    section_order += 1
            else:
                # No headings found, create a single section
                subsection = Subsection(
                    subsection_id="page_1_section_1",
                    order=1,
                    title="HTML Content",
                    content=text_content[:5000] + ("..." if len(text_content) > 5000 else ""),
                    is_cutoff=False,
                    referenced_visuals=[],
                    referenced_tables=[],
                    page_number=1
                )
                page.subsections.append(subsection)
        
        # For XML, extract root elements
        elif file_type == "xml":
            # Simple regex-based extraction (in real app, use proper XML parser)
            root_element_match = re.search(r"<(\w+)[^>]*>", text_content)
            if root_element_match:
                page.title = f"XML Document: {root_element_match.group(1)}"
            
            # Extract top-level elements as sections
            element_matches = re.findall(r"<(\w+)[^>]*>(.*?)</\1>", text_content, re.DOTALL)
            
            if element_matches:
                # Use top-level elements as section titles
                section_order = 1
                for element, content in element_matches[:10]:  # Limit to first 10 to avoid explosion
                    if len(content) > 5000:
                        content = content[:5000] + "..."
                    
                    # Create subsection
                    subsection = Subsection(
                        subsection_id=f"page_1_section_{section_order}",
                        order=section_order,
                        title=f"Element: {element}",
                        content=content,
                        is_cutoff=False,
                        referenced_visuals=[],
                        referenced_tables=[],
                        page_number=1
                    )
                    page.subsections.append(subsection)
                    section_order += 1
            else:
                # No elements found or parsing failed, create a single section
                subsection = Subsection(
                    subsection_id="page_1_section_1",
                    order=1,
                    title="XML Content",
                    content=text_content[:5000] + ("..." if len(text_content) > 5000 else ""),
                    is_cutoff=False,
                    referenced_visuals=[],
                    referenced_tables=[],
                    page_number=1
                )
                page.subsections.append(subsection)
        
        # Use Gemini to enhance the page with insights
        enhanced_page = await enhance_text_page(client, page, filename)
        return [enhanced_page]
        
    except Exception as e:
        print(f"Error processing structured text: {str(e)}")
        # Return basic error page
        return [PageContent(
            page_number=1,
            text=f"Error processing {file_type.upper()} document: {str(e)}",
            title=f"Error: {filename}",
            topics=["error", file_type],
            summary=f"Failed to process {file_type.upper()} document due to: {str(e)}",
            subsections=[]
        )]

async def process_plain_text(client, text_content, filename):
    """Process plain text document with smart section detection."""
    # Split by page breaks (various formats)
    page_breaks = re.compile(r"\n{3,}|=== Page \d+ ===|\f|----+|Page \d+|^\d+\.$", re.MULTILINE)
    pages_text = [p.strip() for p in page_breaks.split(text_content) if p.strip()]
    
    # If no page breaks found and text is very long, split by paragraphs
    if len(pages_text) <= 1 and len(text_content) > 5000:
        paragraphs = [p.strip() for p in re.split(r"\n{2,}", text_content) if p.strip()]
        
        # Group paragraphs into pages of reasonable size
        pages_text = []
        current_page = []
        current_length = 0
        
        for para in paragraphs:
            current_page.append(para)
            current_length += len(para)
            
            if current_length > 3000:
                pages_text.append("\n\n".join(current_page))
                current_page = []
                current_length = 0
        
        if current_page:
            pages_text.append("\n\n".join(current_page))
    
    # If still only one page, use it as is
    if not pages_text:
        pages_text = [text_content]
    
    # Process each page
    processed_pages = []
    
    for i, page_text in enumerate(pages_text, 1):
        # Create basic page
        page = PageContent(
            page_number=i,
            text=page_text,
            title=f"Page {i}" if len(pages_text) > 1 else filename,
            topics=[],
            summary="",
            entities=[],
            has_tables=False,
            has_visuals=False,
            has_numbers=False,
            tables=[],
            visuals=[],
            numbers=[],
            dates=[],
            financial_statements=[],
            key_metrics=[],
            financial_terms=[],
            subsections=[]
        )
        
        # Create default subsection
        subsection = Subsection(
            subsection_id=f"page_{i}_section_1",
            order=1,
            title=f"Content" if len(pages_text) > 1 else filename,
            content=page_text,
            is_cutoff=(i < len(pages_text)),  # Mark as potentially cut off if not last page
            referenced_visuals=[],
            referenced_tables=[],
            page_number=i
        )
        page.subsections.append(subsection)
        
        # Use Gemini to enhance the page
        enhanced_page = await enhance_text_page(client, page, filename)
        processed_pages.append(enhanced_page)
    
    return processed_pages

async def enhance_text_page(client, page, filename):
    """Use Gemini to extract insights and enhance text page content."""
    try:
        # Create a context from page data
        text_sample = page.text[:5000] + ("..." if len(page.text) > 5000 else "")
        subsections_info = ""
        
        if page.subsections:
            subsections_info = "Sections:\n"
            for i, section in enumerate(page.subsections):
                title = section.title if hasattr(section, 'title') else f"Section {i+1}"
                content_preview = (section.content[:200] + "..." if len(section.content) > 200 else section.content) if hasattr(section, 'content') else ""
                subsections_info += f"- {title}: {content_preview}\n"
        
        context = f"""
        Filename: {filename}
        Page number: {page.page_number}
        Current title: {page.title}
        
        Content sample:
        {text_sample}
        
        {subsections_info}
        """
        
        # Create enhancement prompt
        prompt = f"""
        Analyze this text page and return a JSON object with:
        1. A descriptive title for the page (2-5 words)
        2. Key topics discussed (list of 3-5 keywords)
        3. A summary of key points (3-5 sentences)
        4. Entities mentioned (list of important names, organizations, locations, etc.)
        5. Key numerical data points (if any)
        6. Important dates mentioned (if any)
        7. Identify any financial terms (if relevant)
        
        Text Context:
        {context}
        
        Format your response as a valid JSON object with these exact fields:
        {{
            "title": "Descriptive title",
            "topics": ["topic1", "topic2", "topic3"],
            "summary": "Key points summary",
            "entities": ["entity1", "entity2"],
            "numbers": [
                {{
                    "value": "123",
                    "description": "what this number represents",
                    "context": "surrounding context"
                }}
            ],
            "dates": ["date1", "date2"],
            "financial_terms": ["term1", "term2"]
        }}
        """
        
        try:
            # Get enhancement from Gemini
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.0-flash",
                contents=[
                    types.Content(parts=[types.Part.from_text(text=prompt)]),
                ],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            
            if response.candidates:
                json_text = clean_json_response(response.candidates[0].content.parts[0].text)
                data = json.loads(json_text)
                
                # Update page with enhanced data
                if "title" in data and data["title"]:
                    page.title = data["title"]
                
                if "topics" in data and data["topics"]:
                    page.topics = data["topics"]
                
                if "summary" in data and data["summary"]:
                    page.summary = data["summary"]
                
                if "entities" in data and data["entities"]:
                    page.entities = data["entities"]
                
                if "numbers" in data and data["numbers"]:
                    for num_data in data["numbers"]:
                        if isinstance(num_data, dict):
                            page.numbers.append(NumericalDataPoint(
                                value=num_data.get("value", ""),
                                description=num_data.get("description", ""),
                                context=num_data.get("context", "")
                            ))
                    page.has_numbers = bool(page.numbers)
                
                if "dates" in data and data["dates"]:
                    page.dates = data["dates"]
                
                if "financial_terms" in data and data["financial_terms"]:
                    page.financial_terms = data["financial_terms"]
            
        except Exception as e:
            print(f"Error enhancing text page with Gemini: {str(e)}")
            # Continue with basic page info
        
        return page
        
    except Exception as e:
        print(f"Error in enhance_text_page: {str(e)}")
        return page

# Enhanced Entity Relationship Analysis aligned with the provided prompt
async def analyze_entity_relationships_enhanced(client, chapter):
    """
    Analyze entity relationships within a chapter using an enhanced approach
    aligned with the entity extraction prompt structure.
    """
    # Collect content from all subsections
    subsection_texts = []
    for subsection in chapter.subsections:
        subsection_texts.append(f"--- {subsection.title} ---\n{subsection.content}")
    
    # Combine texts while keeping under token limit
    combined_text = "\n\n".join(subsection_texts)
    if len(combined_text) > 100000:
        combined_text = combined_text[:100000] + "..."
    
    # Create prompt for entity extraction using the shared prompt structure
    entity_prompt = f"""
    ---Goal---
    Given the text for chapter "{chapter.title}", identify all entities and relationships among them.
    Use English as output language.

    ---Steps---
    1. Identify all entities. For each identified entity, extract:
       - entity_name: Name of the entity, capitalized if English
       - entity_type: One of the following types: [organization, person, geo, location, event, category, concept, technology, product]
       - entity_description: Comprehensive description of the entity's attributes and activities
    
    2. Identify all pairs of (source_entity, target_entity) that are clearly related to each other.
       For each pair of related entities, extract:
       - source_entity: name of the source entity
       - target_entity: name of the target entity
       - relationship_description: explanation of how entities are related
       - relationship_strength: a numeric score from 1-10 indicating strength
       - relationship_keywords: high-level keywords summarizing the relationship
    
    3. Return output as a JSON object with this structure:
    {{
        "entity_relationships": [
            {{
                "source_entity": "Entity A",
                "target_entity": "Entity B",
                "relationship_description": "Description of relationship",
                "relationship_keywords": ["keyword1", "keyword2"],
                "relationship_strength": 8
            }}
        ],
        "content_keywords": ["keyword1", "keyword2", "keyword3"]
    }}

    Chapter Text:
    {combined_text}
    """
    
    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=entity_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        
        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)
            
            # Convert to EntityRelationship objects
            relationships = []
            for rel_data in data.get("entity_relationships", []):
                try:
                    relationship = EntityRelationship(
                        source_entity=rel_data.get("source_entity", "Unknown"),
                        target_entity=rel_data.get("target_entity", "Unknown"),
                        relationship_description=rel_data.get("relationship_description", ""),
                        relationship_keywords=rel_data.get("relationship_keywords", []),
                        relationship_strength=float(rel_data.get("relationship_strength", 5))
                    )
                    relationships.append(relationship)
                except Exception as rel_error:
                    print(f"Error creating relationship: {rel_error}")
            
            # Update chapter content keywords if available
            if "content_keywords" in data and data["content_keywords"]:
                # Store content keywords in chapter summary
                content_keywords = ", ".join(data["content_keywords"])
                if chapter.summary:
                    chapter.summary += f"\n\nKey themes: {content_keywords}"
                else:
                    chapter.summary = f"Key themes: {content_keywords}"
            
            return relationships
    except Exception as e:
        print(f"Error analyzing entity relationships: {str(e)}")
        return []

async def generate_document_summary(client, pages, filename):
    """Generate comprehensive document summary including theme detection."""
    # Prepare summary input from pages
    summary_input = {
        "filename": filename,
        "page_count": len(pages),
        "page_titles": [],
        "page_summaries": [],
        "entities": set(),
        "topics": set(),
        "has_tables": False,
        "has_visuals": False,
        "has_numbers": False
    }

    # Collect data from pages
    for page in pages:
        # Add page title
        if isinstance(page, dict):
            title = page.get("title", f"Page {page.get('page_number', '?')}")
            summary = page.get("summary", "")
            entities = page.get("entities", [])
            topics = page.get("topics", [])
            has_tables = page.get("has_tables", False)
            has_visuals = page.get("has_visuals", False)
            has_numbers = page.get("has_numbers", False)
        else:
            title = getattr(page, "title", f"Page {getattr(page, 'page_number', '?')}")
            summary = getattr(page, "summary", "")
            entities = getattr(page, "entities", [])
            topics = getattr(page, "topics", [])
            has_tables = getattr(page, "has_tables", False)
            has_visuals = getattr(page, "has_visuals", False)
            has_numbers = getattr(page, "has_numbers", False)
        
        summary_input["page_titles"].append(title)
        summary_input["page_summaries"].append(summary)
        summary_input["entities"].update(entities)
        summary_input["topics"].update(topics)
        summary_input["has_tables"] = summary_input["has_tables"] or has_tables
        summary_input["has_visuals"] = summary_input["has_visuals"] or has_visuals
        summary_input["has_numbers"] = summary_input["has_numbers"] or has_numbers
    
    # Convert sets to lists for JSON serialization
    summary_input["entities"] = list(summary_input["entities"])
    summary_input["topics"] = list(summary_input["topics"])

    # Generate summary with Gemini
    try:
        summary_response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_text(text=f"""
                            Create a comprehensive document summary for "{filename}" with this structure:
                            {{
                                "title": "concise document title",
                                "themes": ["theme1", "theme2", "theme3"],
                                "questions": ["question1", "question2", "question3"],
                                "summary": "detailed multi-paragraph summary",
                                "tables_summary": "summary of key tables",
                                "visuals_summary": "summary of key visuals"
                            }}
                            
                            Document Info:
                            {json.dumps(summary_input, indent=2)}
                        """)
                    ]
                )
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json"
            )
        )
        
        if summary_response.candidates:
            json_text = summary_response.candidates[0].content.parts[0].text
            json_text = clean_json_response(json_text)
            summary_data = json.loads(json_text)
            
            # Handle case where API returns a list
            if isinstance(summary_data, list):
                summary_data = summary_data[0] if summary_data else {}
            
            return summary_data
            
    except Exception as e:
        print(f"Summary generation error: {str(e)}")
    
    # Fallback summary as dict
    return {
        "title": f"Summary: {filename}",
        "themes": list(summary_input["topics"])[:5],
        "questions": ["What are the key points in this document?"],
        "summary": "Generated summary unavailable",
        "tables_summary": "Tables summary unavailable" if summary_input["has_tables"] else "",
        "visuals_summary": "Visuals summary unavailable" if summary_input["has_visuals"] else ""
    }

async def process_all_documents_async(file_data):
    """Process all documents with proper model conversion and error handling."""
    # Set up API client and processing environment
    api_key = get_gemini_api_key()
    client = genai.Client(api_key=api_key,
                          http_options={
                              "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                              "headers": {
                                  "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                              }
                          })
    processed_docs = []
    uploaded_files = []

    # Set up status tracking
    status_container = st.session_state.get("status_container")
    progress_bar = st.session_state.get("progress_bar")
    time_info = st.session_state.get("time_info")

    processing_status = {
        "active": True, "current_step": "", "total_steps": 7, "current_step_num": 0,
        "step_progress": 0.0, "start_time": time.time(), "step_start_time": time.time(),
        "estimated_time_remaining": None, "parallel_count": 1
    }
    steps = {
        "Detecting document types": 1, "Processing PDF documents": 2, "Analyzing PDF pages": 3,
        "Processing non-PDF documents": 4, "Generating document summaries": 5,
        "Finalizing results": 6, "Cleaning up resources": 7
    }

    def update_step_status(step, current=None, total=None, parallel=1):
        processing_status["current_step"] = step
        processing_status["current_step_num"] = steps.get(step, processing_status["current_step_num"])

        step_progress = 0.0
        if current is not None and total is not None and total > 0:
             step_progress = min(1.0, float(current) / float(total))
        processing_status["step_progress"] = step_progress

        overall_progress = (processing_status["current_step_num"] - 1 + processing_status["step_progress"]) / processing_status["total_steps"]
        overall_progress = max(0.0, min(1.0, overall_progress))

        status_msg = f"{processing_status['current_step_num']}/{processing_status['total_steps']}: {step}"
        if current is not None and total is not None:
            status_msg += f" ({current}/{total})"
        if parallel > 1:
            status_msg += f" | Parallelism: {parallel}"

        eta_msg = ""
        if overall_progress > 0.01:
            time_elapsed = time.time() - processing_status["start_time"]
            est_total_time = time_elapsed / overall_progress if overall_progress > 0 else 0
            remaining = max(0, est_total_time - time_elapsed)
            mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
            mins_remaining, secs_remaining = divmod(int(remaining), 60)
            eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: ~{mins_remaining}m {secs_remaining}s"

        if status_container:
             try:
                 status_container.update(label=status_msg, expanded=True)
             except Exception: pass
        if progress_bar:
             try:
                 progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
             except Exception: pass
        if time_info:
             try:
                 time_info.markdown(f"`{eta_msg}`")
             except Exception: pass

    final_docs = []

    try:
        # Step 1: Detect document types
        update_step_status("Detecting document types", 0, len(file_data))
        pdf_docs = []
        other_docs = []
        for i, file in enumerate(file_data):
            update_step_status("Detecting document types", i + 1, len(file_data))
            # Robust file type detection
            file_name = file.get("name", "")
            file_ext = file.get("type", "").lower() or (file_name.split('.')[-1].lower() if '.' in file_name else "")
            file["type"] = file_ext  # Ensure type is set in the dict
            if file_ext == 'pdf':
                pdf_docs.append(file)
            else:
                other_docs.append(file)
        update_step_status("Detecting document types", len(file_data), len(file_data))

        # --- PDF PROCESSING SECTION ---
        if pdf_docs:
            update_step_status("Processing PDF documents", 0, len(pdf_docs))

            MAX_PARALLEL_CONVERSIONS = 5  # Limit concurrent PDF rendering
            all_pdf_pages_info = []  # List to hold page_info dicts from extraction

            # --- Convert PDF pages to images ---
            pdf_conversion_progress = 0
            conversion_tasks = []
            for file in pdf_docs:
                if "content" in file and "name" in file:
                    # Wrap the async call to extract pages
                    conversion_tasks.append(extract_pages_from_pdf_bytes(file["content"], file["name"]))
                else:
                    st.warning(f"Skipping PDF file due to missing 'content' or 'name': {file.get('name', 'Unknown')}")

            # Use a semaphore to limit concurrent PDF rendering if extract_pages_from_pdf_bytes is IO heavy
            pdf_render_semaphore = asyncio.Semaphore(MAX_PARALLEL_CONVERSIONS)
            async def run_with_semaphore(task):
                 async with pdf_render_semaphore:
                     return await task

            tasks_with_semaphore = [run_with_semaphore(task) for task in conversion_tasks]

            # Gather results with progress update
            conversion_results = await asyncio.gather(*tasks_with_semaphore, return_exceptions=True)

            for i, result in enumerate(conversion_results):
                 pdf_conversion_progress += 1  # Increment progress regardless of result
                 if isinstance(result, list):
                     all_pdf_pages_info.extend(result)
                 elif isinstance(result, Exception):
                     st.error(f"PDF page extraction failed for '{pdf_docs[i].get('name', 'Unknown')}': {result}")
                 else:
                     st.warning(f"Unexpected result type during PDF extraction: {type(result)}")

                 update_step_status(
                    "Processing PDF documents",
                    pdf_conversion_progress,
                    len(pdf_docs)  # Total PDFs to convert
                )
            # Final update for this step
            update_step_status("Processing PDF documents", len(pdf_docs), len(pdf_docs))

            # --- Analyze PDF Pages ---
            pdf_page_count = len(all_pdf_pages_info)
            if pdf_page_count > 0:
                update_step_status("Analyzing PDF pages", 0, pdf_page_count)

                MAX_PARALLEL_PAGES = 25  # Gemini API concurrency limit
                semaphore = asyncio.Semaphore(MAX_PARALLEL_PAGES)

                # --- Create tasks ---
                page_tasks = []  # List to hold the Future objects
                valid_page_info_count = 0

                for page_info in all_pdf_pages_info:
                    # Basic validation of the page_info structure
                    if isinstance(page_info, dict) and "doc_name" in page_info and "page_num" in page_info and "image_b64" in page_info:
                        coro = process_single_page_with_semaphore(
                            semaphore, client, page_info, uploaded_files
                        )
                        page_tasks.append(asyncio.create_task(coro))  # Create and store task
                        valid_page_info_count += 1
                    else:
                         st.warning(f"Skipping invalid or incomplete page_info data: {str(page_info)[:100]}...")  # Log truncated data

                # Adjust total page count if some were invalid
                pdf_page_count = valid_page_info_count
                if pdf_page_count == 0:
                    st.warning("No valid PDF page data found to analyze after extraction.")
                else:
                    # --- Collect results using as_completed ---
                    processed_results_with_context = []  # Store dicts: {"doc_name": ..., "page_content": PageContent}
                    processed_count = 0

                    for future in asyncio.as_completed(page_tasks):
                        try:
                            # Get the result dictionary returned by the wrapper
                            task_output = await future  # This is {"info": ..., "result": ...}

                            # Extract the original info and the actual page content/error
                            original_info = task_output.get("info", {})  # Safely get info
                            result_page_obj = task_output.get("result")  # Safely get result (PageContent or Error PageContent)

                            doc_name = original_info.get("doc_name", "Unknown_Context")  # Use fallback
                            page_num = original_info.get("page_num", -1)  # Use fallback

                            if isinstance(result_page_obj, PageContent):
                                processed_results_with_context.append({
                                    "doc_name": doc_name,
                                    "page_content": result_page_obj  # Already has page_num inside
                                })
                            elif result_page_obj is None:
                                 st.warning(f"Page processing task for page {page_num} of {doc_name} returned None result payload.")
                                 # Optionally create an error page here too if needed
                                 error_page = create_error_page(page_num, f"Task returned None payload in '{doc_name}'")
                                 processed_results_with_context.append({"doc_name": doc_name, "page_content": error_page})
                            else:
                                # This case should ideally not happen if the wrapper always returns PageContent or None
                                st.warning(f"Unexpected item in result payload for {doc_name} page {page_num}: {type(result_page_obj)}. Storing as error.")
                                error_page = create_error_page(page_num, f"Unexpected result type {type(result_page_obj)} in '{doc_name}'")
                                processed_results_with_context.append({"doc_name": doc_name, "page_content": error_page})

                        except Exception as e:
                            # Handle exceptions raised by 'await future' itself (e.g., task cancellation)
                            # We don't have reliable context (doc_name/page_num) here if the task itself failed critically
                            st.error(f"Error awaiting page processing task result: {e}")
                            # Log a generic error page if context is lost
                            processed_count += 1  # Increment even on failure to keep progress moving

                        # Increment progress AFTER processing the result or handling the exception
                        processed_count += 1
                        update_step_status(
                            "Analyzing PDF pages",
                            processed_count,
                            pdf_page_count,
                            min(MAX_PARALLEL_PAGES, pdf_page_count - processed_count + 1)  # Active tasks
                        )

                    # Final update for this step
                    update_step_status("Analyzing PDF pages", processed_count, pdf_page_count)

                    # --- Organize PDF pages using collected context ---
                    doc_pages = defaultdict(list)
                    for item in processed_results_with_context:
                        # item is {"doc_name": ..., "page_content": PageContent_object}
                        doc_name = item["doc_name"]
                        page_obj = item["page_content"]  # This is the PageContent object
                        if doc_name != "Unknown_Context":  # Avoid grouping unknown pages
                            doc_pages[doc_name].append(page_obj)
                        else:
                            st.warning(f"Could not determine document name for processed page: {getattr(page_obj, 'page_number', '?')}")

                    # --- Generate Summaries for PDF documents ---
                    doc_count = len(doc_pages)
                    if doc_count > 0:
                        update_step_status("Generating document summaries", 0, doc_count)
                        summary_tasks = []
                        doc_name_order = []  # Keep track of the order for results

                        for doc_name, pages_list in doc_pages.items():
                            # Ensure pages are PageContent objects and sort them
                            valid_pages = sorted(
                                [p for p in pages_list if isinstance(p, PageContent)],
                                key=lambda p: getattr(p, 'page_number', 0)  # Safe access to page_number
                            )
                            if not valid_pages:
                                st.warning(f"No valid pages found for document '{doc_name}' during summary phase.")
                                continue  # Skip summary if no valid pages

                            # Pass the list of PageContent objects to finalize_document
                            summary_tasks.append(finalize_document(client, doc_name, valid_pages))
                            doc_name_order.append(doc_name)

                        if summary_tasks:  # Only run gather if there are tasks
                            summary_results = await asyncio.gather(*summary_tasks, return_exceptions=True)

                            for i, result in enumerate(summary_results):
                                update_step_status("Generating document summaries", i + 1, len(doc_name_order))  # Use len(doc_name_order) as total
                                doc_name = doc_name_order[i]  # Get doc_name based on original order
                                if isinstance(result, Exception):
                                    st.error(f"Failed to finalize summary for {doc_name}: {result}")
                                    # Convert original pages (from doc_pages) to dicts for error output
                                    pages_as_dicts = []
                                    for p in doc_pages.get(doc_name, []):  # Use .get for safety
                                        try:
                                            pages_as_dicts.append(p.model_dump())
                                        except Exception as dump_err:
                                            pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})

                                    processed_docs.append({
                                        "raw_extracted_content": {
                                            "filename": doc_name,
                                            "pages": pages_as_dicts,
                                            "summary": None,
                                            "error": f"Summary generation failed: {result}"
                                        }
                                    })
                                elif isinstance(result, dict) and "raw_extracted_content" in result:
                                    processed_docs.append(result)  # Add the fully processed doc structure
                                else:
                                    st.warning(f"Unexpected result type during summary finalization for {doc_name}: {type(result)}")
                                    # Convert original pages to dicts for error output
                                    pages_as_dicts = []
                                    for p in doc_pages.get(doc_name, []):
                                        try:
                                             pages_as_dicts.append(p.model_dump())
                                        except Exception as dump_err:
                                             pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})
                                    processed_docs.append({
                                        "raw_extracted_content": {
                                            "filename": doc_name,
                                            "pages": pages_as_dicts,
                                            "summary": None,
                                            "error": "Unexpected finalization result type"
                                        }
                                   })
                        else:
                            st.info("No valid documents found to summarize.")

                        # Final status update for summary generation
                        update_step_status("Generating document summaries", len(doc_name_order), len(doc_name_order))
                    else:
                         st.info("No documents found after page processing to summarize.")

            else:  # pdf_page_count == 0 or no valid page info
                 st.info("No PDF pages were successfully extracted or processed, skipping PDF analysis and summary.")

        # --- Non-PDF Document Processing ---
        if other_docs:
            update_step_status("Processing non-PDF documents", 0, len(other_docs))

            MAX_PARALLEL_NON_PDF = 4
            semaphore_non_pdf = asyncio.Semaphore(MAX_PARALLEL_NON_PDF)
            non_pdf_tasks = []
            valid_non_pdf_files = []  # Keep track of files we are actually processing

            for file in other_docs:
                 # Ensure file type and content are present
                 if "type" not in file or not file["type"]:
                      file_name = file.get("name", "Unknown")
                      if '.' in file_name:
                           file["type"] = file_name.split('.')[-1].lower()
                      else:
                           st.warning(f"Cannot determine type for non-PDF file: {file_name}. Skipping.")
                           continue
                 if "content" not in file or not file["content"]:
                     st.warning(f"Missing content for non-PDF file: {file.get('name', 'Unknown')}. Skipping.")
                     continue

                 valid_non_pdf_files.append(file)  # Add to list of files to process

                 # Wrapper to apply semaphore to the processing function
                 async def process_non_pdf_wrapper(f):
                      async with semaphore_non_pdf:
                           try:
                               # process_single_document_memory should return the nested dict { "raw_extracted_content": ... }
                               result_dict = await process_single_document_memory(client, f, uploaded_files)

                               # --- Ensure pages are dicts in the final result ---
                               if result_dict and "raw_extracted_content" in result_dict:
                                   raw_content = result_dict["raw_extracted_content"]
                                   if "pages" in raw_content and isinstance(raw_content["pages"], list):
                                       pages_as_dicts = []
                                       for p in raw_content["pages"]:
                                           try:
                                               pages_as_dicts.append(p.model_dump() if hasattr(p, 'model_dump') else p)
                                           except Exception as dump_err:
                                                pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Dump failed: {dump_err}"})
                                       raw_content["pages"] = pages_as_dicts
                               return result_dict  # Return the modified dict
                           except Exception as e:
                                st.error(f"Error processing non-PDF {f.get('name', 'Unknown')}: {e}")
                                # Return an error structure consistent with successful returns
                                return {
                                    "raw_extracted_content": {
                                        "filename": f.get('name', 'Unknown_NonPDF_Error'),
                                        "pages": [],
                                        "summary": None,
                                        "error": f"Processing failed: {e}"
                                    }
                                }
                 non_pdf_tasks.append(process_non_pdf_wrapper(file))  # Pass the valid file

            non_pdf_results = []
            non_pdf_progress = 0
            total_non_pdf_to_process = len(non_pdf_tasks)  # Based on actual tasks created

            if total_non_pdf_to_process > 0:
                # Use asyncio.gather for potentially better performance if tasks are independent
                gathered_results = await asyncio.gather(*non_pdf_tasks, return_exceptions=True)

                for result in gathered_results:
                    non_pdf_progress += 1
                    if isinstance(result, Exception):
                        # This catches exceptions from the gather itself or unhandled ones in the wrapper
                        st.error(f"Non-PDF processing task failed critically: {result}")
                        # Optionally add an error placeholder if needed, but the wrapper should handle most
                    elif result and isinstance(result, dict) and "raw_extracted_content" in result:
                        non_pdf_results.append(result)  # Add successful or handled error results
                    else:
                        st.warning(f"Received invalid/unexpected result structure from non-PDF processing: {result}")
                    update_step_status("Processing non-PDF documents", non_pdf_progress, total_non_pdf_to_process, min(MAX_PARALLEL_NON_PDF, total_non_pdf_to_process - non_pdf_progress + 1))

                processed_docs.extend(non_pdf_results)
            else:
                st.info("No valid non-PDF documents found to process.")

            # Update status one last time after loop/gather finishes
            update_step_status("Processing non-PDF documents", total_non_pdf_to_process, total_non_pdf_to_process)

        # --- Final Processing Steps ---
        update_step_status("Finalizing results", 1, 1)
        # Ensure all page lists within processed_docs contain only dictionaries
        # This step might be redundant if conversion happens correctly earlier, but acts as a safeguard
        final_cleaned_docs = []
        for doc_result in processed_docs:
             if doc_result and "raw_extracted_content" in doc_result:  # Check if doc_result is valid
                 raw_content = doc_result["raw_extracted_content"]
                 if "pages" in raw_content and isinstance(raw_content["pages"], list):
                     pages_as_dicts = []
                     for p in raw_content["pages"]:
                          if isinstance(p, dict):
                              pages_as_dicts.append(p)
                          elif hasattr(p, 'model_dump'):
                              try:
                                  pages_as_dicts.append(p.model_dump())
                              except Exception as final_dump_err:
                                   pages_as_dicts.append({"page_number": getattr(p, 'page_number', -1), "error": f"Final dump failed: {final_dump_err}"})
                          else:  # Fallback for unexpected types
                               pages_as_dicts.append({"page_number": -1, "error": "Unknown page format", "data": str(p)[:100]})
                     raw_content["pages"] = pages_as_dicts
                 final_cleaned_docs.append(doc_result)  # Add the cleaned doc result
             else:
                 st.warning(f"Skipping invalid document result during finalization: {doc_result}")

        final_docs = final_cleaned_docs  # Assign the cleaned list
        update_step_status("Cleaning up resources", 1, 1)
        processing_status["active"] = False
        if status_container:
            try:
                 status_container.update(label="✅ Document processing complete!", state="complete", expanded=False)
            except Exception: pass

    except asyncio.CancelledError:
         st.warning("Document processing was cancelled.")
         if status_container:
             try:
                 status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
             except Exception: pass
         processing_status["active"] = False
         return []  # Return empty list on cancellation
    except Exception as e:
        st.error(f"❌ An error occurred during document processing: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        if status_container:
            try:
                 status_container.update(label=f"❌ Error: {str(e)}", state="error", expanded=True)
            except Exception: pass
        processing_status["active"] = False
        # Return potentially partially processed docs
        return final_docs  # Return what was processed so far

    finally:
        processing_status["active"] = False
        st.session_state.processing_active = False

    # Ensure return type consistency (list of dictionaries)
    return final_docs

###############################
# UI DISPLAY FUNCTIONS
###############################

def display_table(table):
    """Display table that could be either TableData model or dict."""
    if hasattr(table, 'model_dump'):
        table_data = table.model_dump()
    else:
        table_data = table
    
    # Display table title if available
    if table_data.get("title"):
        st.subheader(table_data["title"])
    
    # Display the table content
    if table_data.get("table_content"):
        try:
            df = pd.read_csv(
                io.StringIO(table_data["table_content"]),
                sep="|",
                engine="python",
                skipinitialspace=True
            )
            df = df.dropna(axis=1, how="all")
            st.dataframe(df)
        except Exception:
            st.markdown(f"```markdown\n{table_data['table_content']}\n```")
        
def display_visual_element(visual):
    """Display a visual element."""
    if isinstance(visual, dict):
        visual_type = visual.get('type', 'Unknown')
        description = visual.get('description', '')
        data_summary = visual.get('data_summary', '')
    else:
        visual_type = getattr(visual, 'type', 'Unknown')
        description = getattr(visual, 'description', '')
        data_summary = getattr(visual, 'data_summary', '')
        
    st.markdown(f"**{visual_type.capitalize()}**")
    st.markdown(f"*Description:* {description}")
    if data_summary:
        st.markdown(f"*Data Summary:* {data_summary}")

def display_entity_relationship_network(relationships, title="Entity Relationship Network"):
    """Display entity relationships as a network graph."""
    if not relationships:
        st.info("No entity relationships to display.")
        return
    
    # Create a network graph
    G = nx.DiGraph()
    
    # Add nodes and edges
    for rel in relationships:
        if isinstance(rel, dict):
            source = rel.get('source_entity', 'Unknown')
            target = rel.get('target_entity', 'Unknown')
            strength = rel.get('relationship_strength', 5)
            description = rel.get('relationship_description', '')
        else:
            source = getattr(rel, 'source_entity', 'Unknown')
            target = getattr(rel, 'target_entity', 'Unknown')
            strength = getattr(rel, 'relationship_strength', 5)
            description = getattr(rel, 'relationship_description', '')
        
        # Add nodes if they don't exist
        if source not in G:
            G.add_node(source)
        if target not in G:
            G.add_node(target)
        
        # Add edge with attributes
        G.add_edge(source, target, weight=strength, description=description)
    
    # Create a figure
    fig, ax = plt.subplots(figsize=(10, 8))
    
    # Set up the layout
    pos = nx.spring_layout(G, seed=42)
    
    # Draw nodes
    nx.draw_networkx_nodes(G, pos, node_size=700, node_color='lightblue', alpha=0.8)
    
    # Draw edges with width based on relationship strength
    edge_width = [G[u][v]['weight']/2 for u, v in G.edges()]
    nx.draw_networkx_edges(G, pos, width=edge_width, alpha=0.7, edge_color='gray', 
                          arrowsize=20, connectionstyle='arc3,rad=0.1')
    
    # Draw labels
    nx.draw_networkx_labels(G, pos, font_size=10, font_weight='bold')
    
    # Add title
    plt.title(title, fontsize=16)
    plt.axis('off')
    
    # Display the plot in Streamlit
    st.pyplot(fig)
    
    # Display edge information as a table
    edge_data = []
    for source, target, data in G.edges(data=True):
        edge_data.append({
            "Source": source,
            "Target": target,
            "Strength": data['weight'],
            "Description": data['description']
        })
    
    if edge_data:
        st.markdown("### Relationship Details")
        df = pd.DataFrame(edge_data)
        st.dataframe(df)

def display_chapter(chapter):
    """Display chapter details including subsections and entity relationships."""
    if isinstance(chapter, dict):
        title = chapter.get('title', 'Untitled Chapter')
        summary = chapter.get('summary', 'No summary available.')
        subsections = chapter.get('subsections', [])
        entity_relationships = chapter.get('entity_relationships', [])
    else:
        title = getattr(chapter, 'title', 'Untitled Chapter')
        summary = getattr(chapter, 'summary', 'No summary available.')
        subsections = getattr(chapter, 'subsections', [])
        entity_relationships = getattr(chapter, 'entity_relationships', [])
    
    st.markdown(f"## {title}")
    st.markdown(summary)
    
    # Display entity relationships if any
    if entity_relationships:
        with st.expander("Entity Relationships", expanded=False):
            display_entity_relationship_network(entity_relationships, f"Entity Relationships in {title}")
    
    # Display subsections
    if subsections:
        st.markdown("### Subsections")
        for subsection in subsections:
            if isinstance(subsection, dict):
                sub_title = subsection.get('title', 'Untitled Subsection')
                sub_content = subsection.get('content', 'No content available.')
                page_num = subsection.get('page_number', '?')
                ref_visuals = subsection.get('referenced_visuals', [])
                ref_tables = subsection.get('referenced_tables', [])
            else:
                sub_title = getattr(subsection, 'title', 'Untitled Subsection')
                sub_content = getattr(subsection, 'content', 'No content available.')
                page_num = getattr(subsection, 'page_number', '?')
                ref_visuals = getattr(subsection, 'referenced_visuals', [])
                ref_tables = getattr(subsection, 'referenced_tables', [])
            
            with st.expander(f"{sub_title} (Page {page_num})"):
                st.markdown(sub_content)
                
                # Display referenced visuals/tables if any
                if ref_visuals:
                    st.markdown("#### Referenced Visuals")
                    st.write(", ".join(ref_visuals))
                
                if ref_tables:
                    st.markdown("#### Referenced Tables")
                    st.write(", ".join(ref_tables))

def display_project_ontology(ontology):
    """Display project-wide ontology."""
    if isinstance(ontology, dict):
        title = ontology.get('title', 'Project Ontology')
        overview = ontology.get('overview', 'No overview available.')
        document_count = ontology.get('document_count', 0)
        documents = ontology.get('documents', [])
        global_themes = ontology.get('global_themes', [])
        entity_relationships = ontology.get('entity_relationships', [])
        key_concepts = ontology.get('key_concepts', [])
    else:
        title = getattr(ontology, 'title', 'Project Ontology')
        overview = getattr(ontology, 'overview', 'No overview available.')
        document_count = getattr(ontology, 'document_count', 0)
        documents = getattr(ontology, 'documents', [])
        global_themes = getattr(ontology, 'global_themes', [])
        entity_relationships = getattr(ontology, 'entity_relationships', [])
        key_concepts = getattr(ontology, 'key_concepts', [])
    
    st.title(title)
    st.markdown(f"**Documents Analyzed:** {document_count}")
    
    # Display overview
    st.markdown("## Overview")
    st.markdown(overview)
    
    # Display global themes
    if global_themes:
        st.markdown("## Global Themes")
        themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" 
                             for tag in global_themes])
        st.markdown(themes_html, unsafe_allow_html=True)
    
    # Display key concepts
    if key_concepts:
        st.markdown("## Key Concepts")
        concepts_html = " ".join([f"<span style='background-color:#f0f0f0; padding:5px; margin:2px; border-radius:5px'>{concept}</span>" 
                               for concept in key_concepts])
        st.markdown(concepts_html, unsafe_allow_html=True)
    
    # Display documents included
    if documents:
        st.markdown("## Documents Included")
        for doc in documents:
            st.markdown(f"- {doc}")
    
    # Display entity relationships
    if entity_relationships:
        st.markdown("## Project-wide Entity Relationships")
        display_entity_relationship_network(entity_relationships, "Project Entity Network")

def display_page_details(document):
    """Display page details, handling both model objects and dictionaries."""
    st.header("📑 Page Level Details")
    
    if hasattr(document, 'pages') and document.pages:
        # Create tab labels that handle both dict and model objects
        tab_labels = []
        for page in document.pages:
            if isinstance(page, dict):
                page_num = page.get('page_number', '?')
            else:
                page_num = getattr(page, 'page_number', '?')
            tab_labels.append(f"Page {page_num}")
        
        tabs = st.tabs(tab_labels)
        
        for i, page in enumerate(document.pages):
            # Convert to dict if it's a model
            if hasattr(page, 'model_dump'):
                page_dict = page.model_dump()
            else:
                page_dict = page
                
            with tabs[i]:
                col1, col2 = st.columns([2, 1])
                
                with col1:
                    st.subheader(page_dict.get('title', f'Page {page_dict.get("page_number", "?")}'))
                    st.markdown("#### Topics")
                    st.write(", ".join(page_dict.get('topics', [])))
                    
                    if page_dict.get('summary'):
                        st.markdown("#### Summary")
                        st.markdown(page_dict.get('summary', ''))
                
                with col2:
                    st.markdown("#### Entities")
                    st.write(", ".join(page_dict.get('entities', [])))
                    
                    st.markdown("#### Content Flags")
                    st.write(f"Tables: {'✅' if page_dict.get('has_tables', False) else '❌'}")
                    st.write(f"Visuals: {'✅' if page_dict.get('has_visuals', False) else '❌'}")
                    st.write(f"Numbers: {'✅' if page_dict.get('has_numbers', False) else '❌'}")
                
                # Display tables
                if page_dict.get('tables'):
                    st.markdown("#### Tables")
                    for table in page_dict.get('tables', []):
                        display_table(table)
                
                # Display visuals
                if page_dict.get('visuals'):
                    st.markdown("#### Visual Elements")
                    for visual in page_dict.get('visuals', []):
                        display_visual_element(visual)
                
                # Display numerical data
                if page_dict.get('numbers'):
                    st.markdown("#### Numerical Data")
                    for num in page_dict.get('numbers', []):
                        value = num.get('value', '') if isinstance(num, dict) else getattr(num, 'value', '')
                        description = num.get('description', '') if isinstance(num, dict) else getattr(num, 'description', '')
                        context = num.get('context', '') if isinstance(num, dict) else getattr(num, 'context', '')
                        
                        st.markdown(f"**{value}**: {description}")
                        if context:
                            st.markdown(f"*Context: {context}*")
                
                # Display subsections
                if page_dict.get('subsections'):
                    st.markdown("#### Subsections")
                    for subsection in page_dict.get('subsections', []):
                        if isinstance(subsection, dict):
                            sub_title = subsection.get('title', 'Untitled')
                            sub_content = subsection.get('content', 'No content')
                            sub_id = subsection.get('subsection_id', 'unknown')
                            is_cutoff = subsection.get('is_cutoff', False)
                        else:
                            sub_title = getattr(subsection, 'title', 'Untitled')
                            sub_content = getattr(subsection, 'content', 'No content')
                            sub_id = getattr(subsection, 'subsection_id', 'unknown')
                            is_cutoff = getattr(subsection, 'is_cutoff', False)
                        
                        with st.expander(f"{sub_title} {'(Cut off)' if is_cutoff else ''}"):
                            st.markdown(sub_content)
                
                # Display page text
                st.markdown("#### Page Content")
                st.text_area("Page Content", page_dict.get('text', ''), height=300, 
                             key=f"page_text_{page_dict.get('page_number', i)}", 
                             label_visibility="collapsed")  # This hides the label visually but keeps it for accessibility

#############################
# STREAMLIT UI COMPONENTS
#############################

def render_unified_document_report(document_data):
    """All document data in one scrollable view with detailed page tabs"""
    # Convert to dict if it's a Pydantic model
    if hasattr(document_data, 'model_dump'):
        doc = document_data.model_dump()
    else:
        doc = document_data
    
    # Safely get raw_content
    raw_content = doc.get('raw_extracted_content', doc)
    
    # --- Header Section ---
    st.title(raw_content.get('filename', 'Document Report'))
    
    # Count tables/visuals (works with both dict and model)
    def count_attr(items, attr):
        return sum(
            1 for item in items 
            if (isinstance(item, dict) and item.get(attr)) 
            or (hasattr(item, attr) and getattr(item, attr))
        )
    
    st.caption(
        f"🔢 {len(raw_content.get('pages', []))} pages | " +
        f"📊 {count_attr(raw_content.get('pages', []), 'has_tables')} tables | " +
        f"📈 {count_attr(raw_content.get('pages', []), 'has_visuals')} visuals"
    )

    # Create tabs for document view modes
    tab1, tab2, tab3 = st.tabs(["Executive Summary", "Chapters", "Detailed Page View"])
    
    with tab1:
        # --- Summary Card ---
        with st.container(border=True):
            summary = raw_content.get('summary', {})
            if isinstance(summary, dict):
                summary_data = summary
            else:
                summary_data = summary.model_dump() if hasattr(summary, 'model_dump') else {}
            
            st.subheader("Executive Summary")
            
            # Display summary text
            if summary_data.get('summary'):
                st.markdown(summary_data.get('summary'))
            else:
                st.info("No summary available for this document.")
            
            # Display themes if available
            if summary_data.get('themes'):
                st.markdown("#### Key Themes")
                themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" 
                                     for tag in summary_data.get('themes', [])])
                st.markdown(themes_html, unsafe_allow_html=True)
            
            # Metrics in columns
            if summary_data.get('key_metrics'):
                st.markdown("#### Key Metrics")
                cols = st.columns(3)
                metrics = summary_data['key_metrics']
                if not isinstance(metrics, list):
                    metrics = list(metrics) if hasattr(metrics, '__iter__') else []
                
                for i, metric in enumerate(metrics[:6]):
                    metric_data = metric if isinstance(metric, dict) else metric.model_dump()
                    cols[i%3].metric(
                        label=metric_data.get('name', 'Metric'),
                        value=metric_data.get('value', 'N/A'),
                        delta=metric_data.get('trend'),
                        help=metric_data.get('context')
                    )
            
            # Display entity relationships if any
            entity_relationships = summary_data.get('entity_relationships', [])
            if entity_relationships:
                st.markdown("#### Entity Relationships")
                with st.expander("View Entity Relationship Network", expanded=False):
                    display_entity_relationship_network(entity_relationships, "Document Entity Relationships")
                    
        # --- Simple Page Navigation (keep this for the summary view) ---
        st.divider()
        pages = raw_content.get('pages', [])
        page_options = []
        for p in pages:
            if isinstance(p, dict):
                title = p.get('title', '')
                num = p.get('page_number', len(page_options)+1)
            else:
                title = getattr(p, 'title', '')
                num = getattr(p, 'page_number', len(page_options)+1)
            page_options.append(f"Page {num} - {title}")
        
        selected_page = st.selectbox(
            "Navigate to page:",
            options=page_options,
            key="page_nav_summary"
        )
        page_idx = int(selected_page.split()[1]) - 1
        page_data = pages[page_idx]
        
        # --- Selected Page Content ---
        with st.container(border=True):
            # Text Content
            st.subheader(page_data.get('title', selected_page))
            st.text_area("Full Text", 
                        page_data.get('text', ''), 
                        height=200,
                        label_visibility="collapsed")

            # Tables
            tables = page_data.get('tables', [])
            if tables:
                st.subheader(f"Tables ({len(tables)})")
                for table in tables:
                    if hasattr(table, 'model_dump'):
                        table = table.model_dump()
                    display_table(table)
    
    with tab2:
        # Display chapters
        chapters = summary_data.get('chapters', [])
        if chapters:
            chapter_titles = [chapter.get('title', f"Chapter {i+1}") for i, chapter in enumerate(chapters)]
            chapter_tabs = st.tabs(chapter_titles)
            
            for i, (tab, chapter) in enumerate(zip(chapter_tabs, chapters)):
                with tab:
                    display_chapter(chapter)
        else:
            st.info("No chapter information available for this document.")
    
    with tab3:
        # Use existing detailed page view but adapt it for the new structure
        pages = raw_content.get('pages', [])
        
        # Convert all pages to PageContent objects or compatible dicts
        processed_pages = []
        for page in pages:
            if hasattr(page, 'model_dump'):
                processed_pages.append(page)  # Already a model
            else:
                # It's a dict, either use as is or convert to PageContent
                processed_pages.append(page)
        
        # Create a wrapper object compatible with display_page_details
        document_wrapper = type('ProcessedDocument', (), {'pages': processed_pages})
        
        # Call your existing function
        display_page_details(document_wrapper)

def display_sidebar_chat():
    """Manages the sidebar content: file upload, chat, and triggers processing."""
    st.sidebar.title("📝 Input & Chat")
    st.sidebar.markdown("Upload documents and chat about their content.")

    # Define supported file types
    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt"]

    # Initialize session state keys if they don't exist
    if "processed_file_names" not in st.session_state:
        st.session_state.processed_file_names = set()  # Store strings only
    if "processing_active" not in st.session_state:
        st.session_state.processing_active = False
    if "last_uploaded_files" not in st.session_state:
        st.session_state.last_uploaded_files = []
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "processed_documents" not in st.session_state:
        st.session_state.processed_documents = []  # Can store both dicts and models
    if "selected_docs" not in st.session_state:
        st.session_state.selected_docs = []
    if "show_results" not in st.session_state:
        st.session_state.show_results = False
    if "selected_doc_to_view" not in st.session_state:
        st.session_state.selected_doc_to_view = None
    if "project_ontology" not in st.session_state:
        st.session_state.project_ontology = None

    # --- File Upload Section ---
    st.sidebar.subheader("📁 Upload Documents")
    
    uploaded_files = st.sidebar.file_uploader(
        "Select files:",
        type=supported_types,
        accept_multiple_files=True,
        key="sidebar_file_uploader_main"
    )

    # Detect *new* files to process automatically
    new_files_to_process = []
    current_files = {f.name for f in uploaded_files} if uploaded_files else set()
    previous_files = {f.name for f in st.session_state.last_uploaded_files} if st.session_state.last_uploaded_files else set()
    
    # If there are new files and we're not already processing
    if uploaded_files and (current_files != previous_files) and not st.session_state.processing_active:
        st.session_state.last_uploaded_files = uploaded_files
        
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                new_files_to_process.append(file)
    
    # If we have new files to process, prepare and trigger processing directly
    if new_files_to_process and not st.session_state.processing_active:
        st.sidebar.info(f"Processing {len(new_files_to_process)} new files...")
        
        # Create the status elements directly in the sidebar
        status_container = st.sidebar.status("Starting document processing...", expanded=True)
        progress_bar = st.sidebar.progress(0)
        time_info = st.sidebar.empty()
        
        # Store in session state for access
        st.session_state.status_container = status_container
        st.session_state.progress_bar = progress_bar
        st.session_state.time_info = time_info
        
        # Prepare file data
        files_data = []
        for file in new_files_to_process:
            try:
                file_content = file.getvalue()
                file_ext = file.name.split('.')[-1].lower()
                files_data.append({
                    "name": file.name,
                    "content": file_content,
                    "type": file_ext
                })
            except Exception as e:
                st.sidebar.error(f"Error reading file {file.name}: {e}")

        if files_data:
            # Process documents directly
            try:
                st.session_state.processing_active = True
                
                # Call the main processing function
                processed_docs = run_async(
                    process_all_documents_async,
                    files_data
                )

                # Update session state with results
                new_processed_docs = [doc for doc in processed_docs if doc is not None]
                st.session_state.processed_documents.extend(new_processed_docs)

                # Update the set of processed file names
                for file_info in files_data:
                    st.session_state.processed_file_names.add(file_info["name"])

                # Set flag to show results
                st.session_state.show_results = True
                
                status_container.update(label="✅ Document processing finished!", state="complete")
                
            except Exception as e:
                st.sidebar.error(f"Error during document processing: {e}")
                status_container.update(label=f"❌ Processing Error: {e}", state="error")
            finally:
                # Ensure processing flag is reset
                st.session_state.processing_active = False
    
    # Display already processed files
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---")
        st.sidebar.write("Processed files:")
        for filename in sorted(st.session_state.processed_file_names):
            st.sidebar.caption(f"✓ {filename}")
        
        # Add a button to view results
        if st.sidebar.button("View Document Reports"):
            st.session_state.show_results = True
            st.rerun()
        
        # Add a button to generate project ontology
        if st.sidebar.button("Generate Project Ontology"):
            if st.session_state.processed_documents:
                with st.spinner("Generating project ontology..."):
                    # Get API client
                    api_key = get_gemini_api_key()
                    client = genai.Client(api_key=api_key)
                    
                    # Generate ontology
                    ontology_data = run_async(
                        generate_project_ontology,
                        client,
                        st.session_state.processed_documents
                    )
                    st.session_state.project_ontology = ontology_data
                    st.success("Project ontology generated successfully!")
                    st.session_state.show_ontology = True
                    st.rerun()
            else:
                st.sidebar.warning("No processed documents available for ontology generation.")

    # --- Chat Section --- 
    st.sidebar.divider()
    st.sidebar.subheader("💬 Document Chat")

    # Document selection for chat context
    if st.session_state.processed_documents:
        doc_options = []
        for doc_result in st.session_state.processed_documents:
            filename = None
            
            if isinstance(doc_result, dict):
                raw_content = doc_result.get("raw_extracted_content", {})
                filename = raw_content.get("filename")
            elif hasattr(doc_result, 'filename'):
                filename = doc_result.filename
            elif hasattr(doc_result, 'raw_extracted_content'):
                filename = doc_result.raw_extracted_content.filename
                
            if filename:
                doc_options.append(filename)

        if doc_options:
            st.session_state.selected_docs = st.sidebar.multiselect(
                "Select documents for chat context:",
                options=sorted(list(set(doc_options))),
                default=st.session_state.selected_docs,
                key="doc_context_selector_sidebar"
            )
        else:
            st.sidebar.caption("No processed documents available to reference.")
    else:
        st.sidebar.caption("No documents processed yet. You can still chat, but I won't be able to reference any document content.")

    # Create a container with border in the sidebar for chat history
    with st.sidebar.container(border=True, height=600):
        st.subheader("Chat History")
        # Display chat history
        if not st.session_state.get("messages", []):
            st.info("No messages yet. Start a conversation below!")
        else:
            for message in st.session_state.messages:
                role = message.get("role", "")
                content = message.get("content", "")
                timestamp = message.get("timestamp", "")
                
                if role == "user":
                    with st.chat_message("user"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Sent: {timestamp}")
                elif role == "assistant":
                    with st.chat_message("assistant"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Received: {timestamp}")
                elif role == "system":
                    with st.chat_message("system"):
                        st.markdown(f"*{content}*")
                        if timestamp:
                            st.caption(f"System: {timestamp}")

        # Chat input
        if prompt := st.sidebar.chat_input("Ask a question...", key="sidebar_chat_input_main"):
            current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages.append({"role": "user", "content": prompt, "timestamp": current_time})
            run_async(process_chat_message, prompt)
            st.rerun()
        
async def process_chat_message(message):
    """Process a chat message with clear user feedback."""
    # Make sure the message list exists
    if "messages" not in st.session_state:
        st.session_state.messages = []
        
    # Add a thinking indicator that will be updated with the real response
    current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.session_state.messages.append({"role": "assistant", "content": "Thinking...", "timestamp": current_time})
    
    try:
        # Get API client
        api_key = get_gemini_api_key()
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                }
            })
        
        # Build context from documents (if any)
        context = ""
        selected_docs = st.session_state.get("selected_docs", [])
        processed_docs = st.session_state.get("processed_documents", [])
        
        if selected_docs and processed_docs:
            # Extract content from selected documents for context
            context = "Using document context from: " + ", ".join(selected_docs) + "\n\n"
            
            for doc_name in selected_docs:
                # Find the document in processed_documents
                selected_doc = None
                for doc in processed_docs:
                    if isinstance(doc, dict) and "raw_extracted_content" in doc:
                        filename = doc["raw_extracted_content"].get("filename", "")
                        if filename == doc_name:
                            selected_doc = doc
                            break
                
                if selected_doc and "raw_extracted_content" in selected_doc:
                    # Extract summary and key pages
                    tech_data = selected_doc["raw_extracted_content"]
                    
                    # Add summary if available
                    if "summary" in tech_data and tech_data["summary"]:
                        summary = tech_data["summary"]
                        context += f"Document: {doc_name}\n"
                        context += f"Title: {summary.get('title', '')}\n"
                        context += f"Summary: {summary.get('summary', '')}\n\n"
                    
                    # Add content from up to 3 key pages
                    if "pages" in tech_data and tech_data["pages"]:
                        pages = tech_data["pages"]
                        # Sort pages by importance (tables, numbers, length)
                        sorted_pages = sorted(pages, 
                                            key=lambda p: (p.get("has_tables", False), 
                                                          p.get("has_numbers", False), 
                                                          len(p.get("text", ""))), 
                                            reverse=True)[:3]
                        
                        for page in sorted_pages:
                            context += f"Page {page.get('page_number', '')}: {page.get('text', '')[:500]}...\n\n"
        
        # Simple prompt for demonstration
        prompt = f"""
        You are a helpful assistant that analyzes documents. 
        
        Document context:
        {context}
        
        Please respond to: {message}
        """
        
        # Generate content - using the correct async pattern
        response = await client.aio.models.generate_content(
            model="gemini-2.0-flash",
            contents=[
                types.Content(
                    role="user",
                    parts=[types.Part.from_text(text=prompt)]
                )
            ]
        )
        
        # Update the thinking message with the actual response
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            if response and response.candidates:
                content = response.candidates[0].content.parts[0].text
                st.session_state.messages[-1]["content"] = content
                st.session_state.messages[-1]["timestamp"] = timestamp
            else:
                st.session_state.messages[-1]["content"] = "Sorry, I couldn't generate a response."
                st.session_state.messages[-1]["timestamp"] = timestamp
        
    except Exception as e:
        # Update thinking message with the error
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages[-1]["content"] = f"Sorry, an error occurred: {str(e)}"
            st.session_state.messages[-1]["timestamp"] = timestamp
    
    # Trigger UI update
    st.rerun()

##############
# Main Application
##############

def main():
    """Main function to run the Streamlit application."""
    # Configure page
    st.set_page_config(
        page_title="Multimodal Document Processor",
        page_icon="📊",
        layout="wide"
    )
    
    # Initialize session state
    if "processed_documents" not in st.session_state:
        st.session_state.processed_documents = []
    
    # --- SIDEBAR ---
    display_sidebar_chat()
    
    # --- MAIN CONTENT ---
    st.title("📑 Advanced Document Processor")
    
    if st.session_state.get("processing_active"):
        with st.status("Processing documents..."):
            st.write("This may take a few minutes")
    elif st.session_state.get("show_ontology") and st.session_state.project_ontology:
        # Display project ontology
        display_project_ontology(st.session_state.project_ontology)
    elif st.session_state.processed_documents:
        # DOCUMENT SELECTION LOGIC
        if len(st.session_state.processed_documents) > 1:
            try:
                # Get list of valid document names
                doc_options = []
                for d in st.session_state.processed_documents:
                    if isinstance(d, dict):
                        name = d.get("filename") or d.get("raw_extracted_content", {}).get("filename")
                    else:
                        name = getattr(d, "filename", None)
                    if name:
                        doc_options.append(name)
                
                # Default to first document if no selection yet
                default_idx = 0
                if "selected_doc_index" in st.session_state:
                    default_idx = min(st.session_state.selected_doc_index, len(doc_options)-1)
                
                selected_name = st.selectbox(
                    "Choose document:",
                    options=doc_options,
                    index=default_idx,
                    key="doc_selector"
                )
                
                # Find matching document
                doc = None
                for d in st.session_state.processed_documents:
                    current_name = None
                    if isinstance(d, dict):
                        current_name = d.get("filename") or d.get("raw_extracted_content", {}).get("filename")
                    else:
                        current_name = getattr(d, "filename", None)
                    
                    if current_name == selected_name:
                        doc = d
                        break
                
                if not doc:
                    st.warning("Document not found, showing first available")
                    doc = st.session_state.processed_documents[0]
                    
            except Exception as e:
                st.error(f"Document selection error: {e}")
                doc = st.session_state.processed_documents[0]
        else:
            doc = st.session_state.processed_documents[0]
        
        # Render the selected document
        render_unified_document_report(doc)
    else:
        # Display welcome screen
        st.info("Upload documents in the sidebar to begin")
        
        # Project description
        st.markdown("""
        ## Enhanced Multimodal Document Processor
        
        This application provides advanced document processing with:
        
        - **Multimodal Processing**: Extract content from PDFs, Word, Excel, PowerPoint, and more
        - **Subsection Extraction**: Intelligently divide documents into logical sections
        - **Chapter Generation**: Group related subsections into coherent chapters
        - **Entity Relationship Analysis**: Identify connections between key concepts
        - **Document Chat**: Interact with your documents through natural language
        - **Project-wide Ontology**: Generate a comprehensive overview across all documents
        
        ### Getting Started
        
        1. Upload one or more documents using the sidebar
        2. Wait for processing to complete
        3. Explore the extracted content, chapters, and relationships
        4. Use the document chat to ask questions about your content
        5. Generate a project-wide ontology to see connections across documents
        """)
        
        st.image("https://via.placeholder.com/800x400.png?text=Enhanced+Document+Processor", caption="Document Processing Pipeline")

if __name__ == "__main__":
    main()
import streamlit as st
import os
import json
import asyncio
import base64
import pandas as pd
import pymupdf  # PyMuPDF
from typing import List, Dict, Optional, Union, Any, Set, Tuple
import pydantic
from pydantic import BaseModel, Field, ValidationError, field_validator, ConfigDict
from google import genai
from google.genai import types
import docx
from pptx import Presentation
import csv
import openpyxl
import re
from collections import Counter, defaultdict
import time
import io
import datetime
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import logging
import traceback

###############################
# DATA MODELS
###############################

class VisualElement(BaseModel):
    """Model for visual elements like charts, graphs, etc."""
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data, or N/A") # Updated description to reflect N/A possibility
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")
    source_url: Optional[str] = Field(None, description="Source URL of the visual")
    alt_text: Optional[str] = Field(None, description="Alternative text for the visual")
    visual_id: str = Field(..., description="Unique identifier for the visual") # Made required as per prompt

class NumericalDataPoint(BaseModel):
    """Model for numerical data points extracted from text."""
    value: str = Field(..., description="The numerical value as a string") # Clarified string type
    description: str = Field(..., description="What the number represents, with units/context") # Made required
    context: Optional[str] = Field(None, description="Surrounding text context, if available") # Changed to Optional[str] and default to None

class TableData(BaseModel):
    """Model for tables extracted from documents."""
    table_content: str = Field(..., description="Markdown formatted table content") # Made required
    title: Optional[str] = Field(None, description="Optional table title") # Kept optional
    summary: Optional[str] = Field(None, description="Optional table summary") # Kept optional
    page_number: int = Field(..., description="Page number where the table appears") # Made required
    table_id: str = Field(..., description="Unique identifier for the table") # Made required

class Subsection(BaseModel):
    """Model for subsections extracted from pages."""
    subsection_id: str = Field(..., description="Unique identifier for the subsection")
    order: int = Field(..., description="Order of the subsection within the page")
    title: str = Field(..., description="Title of the subsection (less than 7 words)") # Added constraint hint
    text: str = Field(..., description="Full text content of the subsection") # Renamed from content to match prompt output
    description: str = Field(..., description="One line description summarizing the main point")  # Made required as per schema
    is_cutoff: bool = Field(..., description="True if content appears to be cut off by page break") # Made required
    referenced_visuals: List[str] = Field(default_factory=list, description="IDs of referenced visuals")
    referenced_tables: List[str] = Field(default_factory=list, description="IDs of referenced tables")

# --- Simplified PageContent Model ---
class PageContent(BaseModel):
    """
    Model for structured content extracted from a single page.
    Stage 1 populates raw_text, tables, visuals, numbers.
    Stage 2 populates subsections from raw_text.
    """
    page_number: int = Field(..., description="Page number in the original document")
    has_tables: bool = Field(..., description="Indicates if any tables were extracted from this page")
    has_visuals: bool = Field(..., description="Indicates if any visuals (charts, images, etc.) were extracted")
    has_numbers: bool = Field(..., description="Indicates if any key numerical data points were extracted")
    tables: List[TableData] = Field(default_factory=list, description="List of tables extracted from the page")
    visuals: List[VisualElement] = Field(default_factory=list, description="List of visual elements extracted from the page")
    numbers: List[NumericalDataPoint] = Field(default_factory=list, description="List of numerical data points extracted from the page")
    # NEW: Store raw text from initial extraction
    raw_text: Optional[str] = Field(None, description="Full raw text extracted initially from the page/chunk")
    # Subsections are now populated in a second step
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections extracted from raw_text")

    model_config = ConfigDict(arbitrary_types_allowed=True, extra='forbid')

# --- New Models for Qdrant Structure ---
class LinkedNode(BaseModel):
    """Structure for linked nodes metadata in Qdrant."""
    target_node_id: str = Field(..., description="ID of the target node (e.g., type_name_chX)")
    relationship_description: str = Field(..., description="Nature of the relationship from source to target")
    relationship_keywords: List[str] = Field(default_factory=list, description="Keywords summarizing the relationship")
    relationship_strength: float = Field(..., description="Strength of the relationship (1-10)")

    @field_validator('relationship_strength', mode='before')
    def strength_to_float(cls, v):
        try:
            return float(v)
        except (ValueError, TypeError):
            return 5.0 # Default strength on conversion error

class QdrantNode(BaseModel):
    """Structure for nodes to be stored in Qdrant."""
    node_id: str = Field(..., description="Unique node ID (e.g., type_sanitizedname_chX)")
    node_type: str = Field(..., description="'entity' or 'concept'")
    name: str = Field(..., description="Entity or concept name")
    description: str = Field(..., description="Comprehensive description from text")
    chapter_id_context: str = Field(..., description="Chapter ID where this node was primarily identified")
    # --- ADDED FIELD ---
    document_context: str = Field(..., description="Identifier (filename, URL, article name, etc.) of the source document where this node was primarily identified")
    # --------------------
    linked_nodes: List[LinkedNode] = Field(default_factory=list, description="List of outgoing relationships")

    model_config = ConfigDict(extra='ignore') # Allow ignoring extra fields if needed during validation

class Chapter(BaseModel):
    """Model for chapters composed of subsections, including Qdrant node data."""
    chapter_id: str = Field(..., description="Unique identifier for the chapter")
    title: str = Field(..., description="Title of the chapter")
    summary: str = Field(..., description="Summary of the chapter")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections in this chapter")
    # REMOVE THIS LINE:
    # entity_relationships: List[EntityRelationship] = Field(default_factory=list, description="Legacy entity relationships within the chapter")
    order: int = Field(..., description="Order of the chapter in the document")
    # New field to store structured node data for Qdrant
    qdrant_nodes: Optional[List[QdrantNode]] = Field(None, description="Structured nodes and relationships for Qdrant")
    
class DocumentSummaryDetails(BaseModel):
    """Structure for summary details generated from node analysis."""
    title: str = Field(..., description="Concise title for the document based on nodes")
    themes: List[str] = Field(default_factory=list, description="Main themes/topics derived from nodes")
    questions: List[str] = Field(default_factory=list, description="Sample questions reflecting node content")
    summary: str = Field(..., description="Comprehensive summary synthesizing node information")
    # Note: No chapters, entity_relationships, table/visual summaries here

class KeyTopic(BaseModel):
    name: str = Field(..., description="Short topic name (1-3 words)")
    description: str = Field(..., description="Brief description of the topic")
    relevance: str = Field(..., description="Relevance level (High/Medium/Low)")
    sentiment: str = Field(..., description="Sentiment analysis (Positive/Neutral/Mixed/Cautionary/Negative)")
    analysis: str = Field(..., description="Brief analysis of the topic")

class QuotedStatement(BaseModel):
    speaker: str = Field(..., description="Person or entity who made the statement")
    quote: str = Field(..., description="The quoted text")
    page: int = Field(..., description="Page number where the quote appears")

class DocumentReport(BaseModel):
    file_name: str = Field(..., description="Original filename")
    page_count: int = Field(..., description="Number of pages in document")
    title: str = Field(..., description="Document title")
    title_summary: str = Field(..., description="Brief summary/subtitle")
    concept_theme_hashtags: List[str] = Field(default_factory=list, description="Theme tags as hashtags")
    date_published: Optional[str] = Field(None, description="Publication date if available")
    source: Optional[str] = Field(None, description="Document source")
    confidence: str = Field("Medium", description="Analysis confidence (High/Medium/Low)")
    document_summary: str = Field(..., description="Comprehensive summary")
    key_insights: List[str] = Field(default_factory=list, description="Key takeaways and insights")
    key_topics: List[KeyTopic] = Field(default_factory=list, description="Detailed topic analysis")
    quoted_statements: List[QuotedStatement] = Field(default_factory=list, description="Important quotes from document")
    content_excerpt: Optional[str] = Field(None, description="Highlighted document content")
    chapters: List[Chapter] = Field(default_factory=list, description="Document chapters")

class ProcessedDocument(BaseModel):
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages") # This now uses the simplified PageContent
    summary: Optional[DocumentSummaryDetails] = None

class ProjectOntology(BaseModel):
    title: str = Field(..., description="Project title")
    overview: str = Field(..., description="Project overview based on synthesized nodes")
    document_count: int = Field(..., description="Number of documents analyzed")
    documents: List[str] = Field(..., description="List of document filenames included in analysis")
    global_themes: List[str] = Field(..., description="High-level project themes derived from nodes")
    key_concepts: List[str] = Field(..., description="Key concepts identified across all documents from nodes")
    # NEW FIELD: Stores the aggregated nodes including inter-document relationships
    project_graph_nodes: Optional[List[QdrantNode]] = Field(None,
                                description="Aggregated nodes and their relationships (including inter-document links) across the project.")

    model_config = ConfigDict(arbitrary_types_allowed=True) # If QdrantNode uses non-standard types


###############################
# API AND UTILITY FUNCTIONS
###############################

def convert_to_model(data, model_class):
    """Convert dictionary data to a Pydantic model with error handling."""
    if not isinstance(data, dict):
        print(f"Warning: Expected dict for {model_class.__name__}, got {type(data)}")
        return None
    
    try:
        # Extract only the fields that the model expects
        model_fields = model_class.__annotations__.keys()
        filtered_data = {k: v for k, v in data.items() if k in model_fields}
        return model_class(**filtered_data)
    except Exception as e:
        print(f"Error converting to {model_class.__name__}: {str(e)}")
        return None

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
        st.stop()
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)
            
            # Validate JSON if applicable
            config = kwargs.get('config')
            if (config and hasattr(config, 'response_mime_type') and 
                config.response_mime_type == 'application/json' and 
                hasattr(response, 'candidates') and response.candidates):
                try:
                    # Try to parse the JSON response
                    json_text = response.candidates[0].content.parts[0].text
                    json.loads(clean_json_response(json_text))  # Use our clean function first
                except json.JSONDecodeError as e:
                    if attempt < max_retries - 1:
                        print(f"Received malformed JSON on attempt {attempt+1}, retrying: {e}")
                        await asyncio.sleep(2 ** attempt)
                        continue
            
            return response
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            print(f"API call failed on attempt {attempt+1}, retrying: {e}")
            await asyncio.sleep(2 ** attempt)

def clean_json_response(json_text: str, extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        return "{}"
    
    # Ensure it's a string
    json_text = str(json_text)
    
    try:
        # Handle markdown code blocks
        if json_text.startswith("```"):
            blocks = json_text.split("```")
            if len(blocks) >= 3:
                json_text = blocks[1]
                if json_text.startswith("json"):
                    json_text = json_text[4:].strip()
            else:
                json_text = json_text.replace("```", "").strip()
        
        elif json_text.startswith("json"):
            json_text = json_text[4:].strip()
        
        # Remove any leading/trailing whitespace
        json_text = json_text.strip()
        
        # Attempt to parse the JSON as is
        try:
            json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # Continue with fixes
            pass
        
        # Common JSON fixes
        # Fix missing quotes around property names
        json_text = re.sub(r'([{,]\s*)(\w+)(\s*:)', r'\1"\2"\3', json_text)
        
        # Fix trailing commas in arrays/objects
        json_text = re.sub(r',(\s*[}\]])', r'\1', json_text)
        
        # Fix missing quotes around string values
        # This is trickier and might cause issues with valid JSON if not careful
        
        # Final attempt to parse with Python's JSON parser
        try:
            data = json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # If we're extracting text on failure and can't fix the JSON
            if extract_text_on_failure:
                print("JSON parsing failed - extracting text content instead")
                
                # Extract what appears to be the main text content
                # Look for patterns that would indicate text fields in a JSON response
                text_content = ""
                
                # Try to find text between quotes after "text": or "content": patterns
                text_matches = re.findall(r'"(?:text|content)"\s*:\s*"([^"]+)"', json_text)
                if text_matches:
                    text_content = " ".join(text_matches)
                
                # If that doesn't work, take any text between quotes as potential content
                if not text_content:
                    # Get all quoted strings that look reasonably long (over 20 chars)
                    quoted_text = re.findall(r'"([^"]{20,})"', json_text)
                    if quoted_text:
                        text_content = " ".join(quoted_text)
                
                # If still nothing, just take the longest line that's not mostly symbols
                if not text_content:
                    lines = json_text.split('\n')
                    content_lines = [line for line in lines if len(line) > 50 and 
                                    sum(c.isalpha() or c.isspace() for c in line) / len(line) > 0.7]
                    if content_lines:
                        text_content = " ".join(content_lines)
                
                # Fallback - create a minimal valid JSON with the extracted text
                fallback_json = {
                    "text": text_content or "Extracted text content not found",
                    "title": "Text Extraction Fallback",
                    "topics": ["text extraction"],
                    "summary": "Structured JSON parsing failed - basic text extraction used as fallback.",
                    "entities": [],
                    "has_tables": False,
                    "has_visuals": False,
                    "has_numbers": False,
                    "dates": [],
                    "tables": [],
                    "visuals": [],
                    "numbers": [],
                    "financial_statements": [],
                    "key_metrics": [],
                    "financial_terms": [],
                    "subsections": []
                }
                return json.dumps(fallback_json)
            else:
                # Return minimal valid JSON as before
                return '{"text":"Error parsing JSON response","error":"Invalid JSON structure"}'
            
    except Exception as e:
        print(f"Error cleaning JSON response: {e}")
        
        if extract_text_on_failure:
            # Create a fallback with whatever raw content we have
            raw_text = json_text.strip()
            if len(raw_text) > 10000:  # Truncate very long responses
                raw_text = raw_text[:10000] + "... (truncated)"
                
            fallback_json = {
                "text": raw_text,
                "title": "Raw Text Fallback",
                "topics": ["raw text"],
                "summary": "Response processing failed - returning raw text.",
                "entities": [],
                "has_tables": False,
                "has_visuals": False,
                "has_numbers": False,
                "dates": [],
                "tables": [],
                "visuals": [],
                "numbers": [],
                "financial_statements": [],
                "key_metrics": [],
                "financial_terms": [],
                "subsections": []
            }
            return json.dumps(fallback_json)
        else:
            return '{"text":"Error cleaning JSON response","error":"' + str(e).replace('"', '\\"') + '"}'

async def run_in_executor(func, *args):
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, func, *args) # Use default executor

def run_async(func, *args, **kwargs):
    """Run an async function from Streamlit."""
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        # If func is already a coroutine, run it directly
        if asyncio.iscoroutine(func):
            return loop.run_until_complete(func)
        # Otherwise call it with the provided args
        return loop.run_until_complete(func(*args, **kwargs))
    finally:
        loop.close()

###############################
# CONTENT EXTRACTION FUNCTIONS
###############################

async def process_page_with_fallback(client, image_part, page_num, filename=None):
    """Process page with multi-level fallback including pure text extraction."""
    doc_name = filename or "Unknown_Document"
    
    try:
        # First attempt: Full structured extraction
        logging.debug(f"Attempting full subsection extraction for page {page_num}")
        page_content = await extract_page_content_from_memory(client, image_part, page_num, doc_name)

        # Check if it failed severely (e.g., returned an error page type, or text is placeholder)
        is_error = any(s.title == "Processing Error" for s in page_content.subsections) if page_content.subsections else False # Check error subsection
        if is_error or page_content.raw_text == "[LLM failed to extract raw text]":
             logging.warning(f"Initial structure extraction failed or yielded no text for page {page_num} of '{doc_name}'. Triggering fallback.")
             raise ValueError("Initial extraction failed or returned no text") # Go to fallback

        logging.info(f"Successfully extracted structure (incl. raw_text) for page {page_num} of '{doc_name}'")
        return page_content

    except Exception as e:
        logging.warning(f"Initial extraction failed for page {page_num} of '{doc_name}': {e}. Trying pure text fallback.")

        # --- Fallback attempt: Pure text extraction ---
        try:
            logging.debug(f"Attempting pure text extraction fallback for page {page_num} of '{doc_name}'")
            text_only_prompt = f"Extract ONLY the plain text content from page {page_num} of this document. Preserve paragraph breaks. Do NOT format as JSON."

            text_response = await retry_api_call(
                 client.aio.models.generate_content,
                 model="gemini-2.5-flash-preview-04-17", # Model suitable for text extraction
                 contents=[
                    types.Content(parts=[image_part, types.Part.from_text(text=text_only_prompt)]),
                 ],
            )

            raw_text = ""
            if text_response.candidates:
                raw_text = text_response.candidates[0].content.parts[0].text.strip()

            if not raw_text:
                 raw_text = "[No text extracted during fallback]"
                 logging.warning(f"Pure text fallback for page {page_num} resulted in empty text.")

            logging.info(f"Successfully extracted pure text fallback for page {page_num} of '{doc_name}'")
            # Create a PageContent object with ONLY raw_text and page_number populated
            return PageContent(
                page_number=page_num,
                raw_text=raw_text,
                # Set defaults for other required fields
                has_tables=False, has_visuals=False, has_numbers=False,
                tables=[], visuals=[], numbers=[], subsections=[]
            )

        except Exception as text_error:
            logging.error(f"Pure text extraction fallback also failed for page {page_num} of '{doc_name}': {text_error}", exc_info=True)
            # Final fallback: return an error page object
            return create_error_page(page_num, f"All extraction methods failed: {text_error}")

# --- Define the NEW prompt for subsection extraction from text ---
subsection_extraction_prompt = """
Analyze the provided text content from a single page ({page_num}) of a document, originating from source '{source_identifier}'. Your goal is to segment this text into logical subsections, providing concise yet **informative descriptions geared towards identifying entities and relationships later**.

Input Text:
--- START TEXT ---
{raw_text_content}
--- END TEXT ---

Instructions:
1.  Read the text and identify logical breaks based on topic shifts, headings (if any within the text), or paragraph structure.
2.  Segment the text into sequential subsections.
3.  For each subsection, determine:
    *   `subsection_id`: Generate a unique ID formatted as `page_{page_num}_section_ORDER` where ORDER is the sequential order number.
    *   `order`: The sequential order number (integer, starting from 1).
    *   `title`: A concise title (less than 7 words). Use headings from the text if present and appropriate, otherwise generate a descriptive title.
    *   `text`: The full text content belonging *only* to this specific subsection. Ensure all original text is captured across subsections.
    *   `description`: **CRITICAL FOR RELATIONSHIP ANALYSIS:** A concise, one-sentence summary (approx. 15-25 words) that **highlights the key entities, concepts, or specific actions/interactions** discussed **within this specific subsection**. This description acts as a vital preview for downstream entity and relationship identification.
        *   **Focus:** Mention the main actors (e.g., 'Parsely', 'GPT-40', 'CafeCorner LLC') and the core action or relationship involving them *in this text block* (e.g., 'developed by', 'utilizes', 'integrates', 'enhances', 'piloted for').
        *   **Example Good Descriptions:** "Explains how Parsely integrates Llamalndex and Qdrant Vector DB.", "Details CafeCorner LLC's role in developing the Parsely system.", "Introduces GPT-40 as the LLM used for internal analysis.", "Discusses the enhancement of Pitchbook data distribution via Parsely.", "Outlines the piloting of Parsely for JPMorgan's ChatIQ assistant."
        *   **Example Less Useful Description:** "This section talks about the system." (Too generic - doesn't mention specific entities or actions).
    *   `is_cutoff`: Set to `true` ONLY if this specific text chunk appears to end mid-sentence or mid-thought, suggesting it was truncated *before* being passed to you. Otherwise `false`.
    *   `referenced_visuals`: An empty list `[]`. (Visuals/Tables extracted previously).
    *   `referenced_tables`: An empty list `[]`. (Visuals/Tables extracted previously).

4.  Return ONLY a valid JSON array containing the subsection objects, structured exactly like this example:
    ```json
    [
      {{
        "subsection_id": "page_{page_num}_section_1",
        "order": 1,
        "title": "Parsely System Overview",
        "text": "The full text content of the first subsection detailing Parsely...",
        "description": "Introduces the Parsely system developed at CafeCorner LLC for data enhancement.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }},
      {{
        "subsection_id": "page_{page_num}_section_2",
        "order": 2,
        "title": "Core Technologies",
        "text": "Details on the integration of various technologies...",
        "description": "Details Parsely's use of GPT-40, Llamalndex, and Qdrant Vector DB for analysis.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }}
      // ... more subsection objects
    ]
    ```
Ensure the output is only the JSON array and adheres strictly to the structure.
"""

# The function call itself was correct, the error was in the prompt string definition
async def extract_subsections_from_text(client: genai.Client, raw_text: str, page_num: int, filename: str, max_structure_retries: int = 2, retry_delay: int = 2) -> List[Subsection]:
    """
    Takes raw text from a page and uses an LLM to segment it into Subsection objects.
    Includes robustness for dict-wrapped lists and retries on structural errors.
    """
    function_name = "extract_subsections_from_text"
    logging.debug(f"[{function_name}] Starting subsection extraction for page {page_num} of '{filename}'")

    if not raw_text or not raw_text.strip():
        logging.warning(f"[{function_name}] No raw text provided for page {page_num} of '{filename}'. Returning empty list.")
        return []

    # Truncate text if needed (same as before)
    MAX_TEXT_LIMIT = 100000
    text_for_prompt = raw_text[:MAX_TEXT_LIMIT] + "... [Content Truncated]" if len(raw_text) > MAX_TEXT_LIMIT else raw_text
    if len(raw_text) > MAX_TEXT_LIMIT:
         logging.warning(f"[{function_name}] Truncating raw text for page {page_num}")


    prompt = subsection_extraction_prompt.format( # Assuming subsection_extraction_prompt is defined correctly now
        page_num=page_num,
        raw_text_content=text_for_prompt
    )

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Attempt {attempt + 1}/{max_structure_retries} for page {page_num} of '{filename}'")
        try:
            response = await retry_api_call( # This handles API-level retries
                client.aio.models.generate_content,
                model="gemini-2.5-flash-preview-04-17", # Use a recent flash model - good balance
                contents=[types.Content(parts=[types.Part.from_text(text=prompt)])],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates from LLM (Attempt {attempt + 1}) for page {page_num}.")
                # Don't retry if API gives no candidates, proceed to fallback outside loop
                subsections_data = None # Signal failure to proceed to fallback
                break # Exit retry loop

            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"[{function_name}] Raw LLM response (Attempt {attempt + 1}, Page {page_num}): {json_text[:500]}...")

            try:
                subsections_data = json.loads(json_text)
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError on attempt {attempt + 1} for page {page_num}: {json_err}. Raw text: {json_text[:200]}...")
                if attempt < max_structure_retries - 1:
                    logging.info(f"[{function_name}] Retrying after delay due to JSON decode error.")
                    await asyncio.sleep(retry_delay * (attempt + 1)) # Exponential backoff for structure retry
                    continue # Go to next attempt
                else:
                    logging.error(f"[{function_name}] Failed to decode JSON after {max_structure_retries} attempts for page {page_num}.")
                    subsections_data = None # Signal failure
                    break # Exit retry loop

            # --- Structure Check and Robustness ---
            final_list_data = None
            if isinstance(subsections_data, list):
                final_list_data = subsections_data # Correct structure received
                logging.debug(f"[{function_name}] Received correct list structure on attempt {attempt + 1} for page {page_num}.")
                break # Success, exit retry loop
            elif isinstance(subsections_data, dict):
                logging.warning(f"[{function_name}] LLM returned a dict instead of a list on attempt {attempt + 1} for page {page_num}. Checking common keys...")
                # Check common keys where the list might be nested
                possible_keys = ["subsections", "data", "items", "result", "page_subsections"]
                found = False
                for key in possible_keys:
                    if key in subsections_data and isinstance(subsections_data[key], list):
                        logging.info(f"[{function_name}] Found list under key '{key}'. Using this list.")
                        final_list_data = subsections_data[key]
                        found = True
                        break # Use the found list
                if found:
                     break # Success (recovered), exit retry loop

                # If not found under common keys, treat as structural failure
                if attempt < max_structure_retries - 1:
                    logging.warning(f"[{function_name}] Could not find list within dict. Retrying after delay.")
                    await asyncio.sleep(retry_delay * (attempt + 1))
                    continue # Go to next attempt
                else:
                    logging.error(f"[{function_name}] LLM returned a dict and list could not be recovered after {max_structure_retries} attempts for page {page_num}.")
                    subsections_data = None # Signal failure
                    break # Exit retry loop
            else:
                # Incorrect type altogether
                logging.error(f"[{function_name}] LLM response was not a list or dict (type: {type(subsections_data)}) on attempt {attempt+1} for page {page_num}.")
                if attempt < max_structure_retries - 1:
                     logging.info(f"[{function_name}] Retrying after delay due to unexpected response type.")
                     await asyncio.sleep(retry_delay * (attempt + 1))
                     continue # Go to next attempt
                else:
                     logging.error(f"[{function_name}] Failed to get valid structure after {max_structure_retries} attempts for page {page_num}.")
                     subsections_data = None # Signal failure
                     break # Exit retry loop


        except Exception as e:
            logging.error(f"[{function_name}] Error during attempt {attempt + 1} for page {page_num}: {e}", exc_info=True)
            if attempt < max_structure_retries - 1:
                logging.info(f"[{function_name}] Retrying after delay due to unexpected error.")
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue
            else:
                 logging.error(f"[{function_name}] Failed after {max_structure_retries} attempts for page {page_num} due to error: {e}")
                 subsections_data = None # Signal failure
                 break # Exit retry loop


    # --- Process the final data (if successful) or fallback ---
    if final_list_data is not None: # Check if we successfully got a list
        validated_subsections = []
        logging.debug(f"[{function_name}] Processing {len(final_list_data)} potential subsection items for page {page_num}.")
        for sub_dict in final_list_data: # Process the list
            try:
                # *** ADDED CHECK FOR NON-EMPTY TEXT ***
                subsection_text = sub_dict.get("text", "")
                if not isinstance(sub_dict, dict) or not subsection_text or not subsection_text.strip():
                    logging.warning(f"[{function_name}] Skipping subsection item with missing or empty text for page {page_num}: {str(sub_dict)[:100]}...")
                    continue # Skip this item entirely if text is missing/empty

                # Proceed with validation if text is present
                sub_dict.setdefault("referenced_visuals", [])
                sub_dict.setdefault("referenced_tables", [])
                sub_dict.setdefault("is_cutoff", False)
                # Ensure order is an int
                if 'order' in sub_dict and not isinstance(sub_dict['order'], int):
                    try: sub_dict['order'] = int(sub_dict['order'])
                    except (ValueError, TypeError): sub_dict['order'] = 0

                subsection_obj = Subsection(**sub_dict)
                validated_subsections.append(subsection_obj) # Add only if text is valid and validation passes
            except ValidationError as ve:
                logging.warning(f"[{function_name}] Validation failed for subsection on page {page_num}: {ve}. Skipping. Data: {sub_dict}")
            except Exception as val_err:
                logging.error(f"[{function_name}] Unexpected validation error for subsection on page {page_num}: {val_err}. Skipping. Data: {sub_dict}", exc_info=True)

        # Sort by order just in case LLM didn't guarantee it
        validated_subsections.sort(key=lambda s: getattr(s, 'order', 0))

        # Re-assign order sequentially after filtering and sorting
        for i, sub in enumerate(validated_subsections):
            sub.order = i + 1
            # Optionally update ID if order changed significantly, though less critical
            # sub.subsection_id = f"page_{page_num}_section_{sub.order}"

        logging.info(f"[{function_name}] Successfully extracted and validated {len(validated_subsections)} non-empty subsections for page {page_num} of '{filename}'.")
        return validated_subsections
    else:
        # Fallback if all retries failed or initial API call failed
        logging.error(f"[{function_name}] All extraction attempts failed for page {page_num}. Using fallback.")
        # Assuming create_error_fallback_subsection exists and creates a single error subsection
        return [create_error_fallback_subsection(page_num, raw_text, "All LLM extraction attempts failed")]

def create_error_fallback_subsection(page_num: int, raw_text: str, error_msg: str) -> Subsection:
     """
     Creates a single fallback Subsection object when subsection extraction fails.
     This ensures the raw text is preserved along with the error context.
     """
     fallback_id = f"page_{page_num}_section_fallback_error"
     fallback_title = "Full Page Content (Subsection Extraction Failed)"
     # Combine error and original text for context
     fallback_text = (
         f"--- ERROR DURING SUBSECTION EXTRACTION ---\n"
         f"{error_msg}\n"
         f"--- END ERROR ---\n\n"
         f"--- ORIGINAL RAW TEXT FOR PAGE {page_num} ---\n"
         f"{raw_text}\n"
         f"--- END RAW TEXT ---"
     )
     fallback_description = f"Could not segment page {page_num} into subsections due to error: {error_msg[:100]}..." # Truncate long errors

     try:
          # Create and return a single Subsection object
          return Subsection(
              subsection_id=fallback_id,
              order=1, # Always the first (and only) subsection in this fallback case
              title=fallback_title,
              text=fallback_text,
              description=fallback_description,
              is_cutoff=False, # Assume not cutoff unless original text was truncated
              referenced_visuals=[], # No references identified
              referenced_tables=[]  # No references identified
          )
     except ValidationError as ve_fb:
         # This should be rare if the inputs (page_num, raw_text, error_msg) are valid types
         logging.error(f"[create_error_fallback_subsection] Critical: Failed to create even the fallback Subsection for page {page_num}: {ve_fb}", exc_info=True)
         # Re-raise the error as we cannot create a valid object, indicating a deeper issue
         raise ve_fb
     except Exception as e:
        # Catch any other unexpected errors during creation
        logging.error(f"[create_error_fallback_subsection] Unexpected error creating fallback Subsection for page {page_num}: {e}", exc_info=True)
        # Depending on desired behavior, could raise or return a dummy/error object
        # Re-raising is often safer to signal the failure clearly
        raise e


# Prompt for initial structure extraction (raw_text + elements)
# --- New Task-Oriented Prompt for Initial Page Structure Extraction ---

page_structure_prompt = """
--- Goal ---
Analyze the input, which represents a single page ({page_num}) of a document (could be an image or text). Your primary task is to extract the complete raw text content AND identify and structure specific elements like tables, visuals, and key numerical data points found on that page.

--- Input ---
- Content of page {page_num} (provided as image or text).

--- Instructions ---
Perform the following tasks based *only* on the provided input content for page {page_num}:

**Task 1: Extract Raw Text**
1a. Extract ALL text content visible on the page.
1b. Preserve original formatting like paragraphs and line breaks as accurately as possible.
1c. Store this complete text in the `raw_text` field of the output JSON.

**Task 2: Extract Tables**
2a. Identify all distinct tables present on the page.
2b. For EACH table found:
    *   Extract the complete cell content accurately.
    *   Format this content STRICTLY as GitHub Flavored Markdown within the `table_content` field (using '|' separators and a '---' header separator line). Do NOT include surrounding paragraphs or explanatory text in `table_content`.
    *   If a clear title or caption is associated with the table, extract it into the `title` field. Otherwise, use `null`.
    *   Write a brief (1-2 sentence) `summary` of the table's main content or purpose, if possible. Otherwise, use `null`.
    *   Generate a unique `table_id` formatted as `page_{page_num}_table_N` (where N is 1, 2, 3...).
    *   Set `page_number` to {page_num}.
    *   Add the complete table object to the `tables` list in the output JSON.
2c. If NO tables are found on the page, the `tables` list MUST be an empty list `[]`.

**Task 3: Extract Visual Elements**
3a. Identify all visual elements (e.g., charts, graphs, diagrams, images, photographs, equations, code blocks).
3b. For EACH visual element found:
    *   Determine its `type` (e.g., "bar chart", "line graph", "diagram", "image", "equation", "code block").
    *   Write a detailed `description` covering what the visual shows, its purpose, and key elements.
    *   Provide a `data_summary` summarizing key data points, trends, or findings presented in the visual. If the visual doesn't present data (e.g., a decorative image) or a summary isn't applicable, use the string "N/A".
    *   Generate a unique `visual_id` formatted as `page_{page_num}_visual_N` (where N is 1, 2, 3...).
    *   Set `page_numbers` to `[{page_num}]`.
    *   Set `source_url` and `alt_text` to `null` unless explicitly available.
    *   Add the complete visual object to the `visuals` list in the output JSON.
3c. If NO visual elements are found, the `visuals` list MUST be an empty list `[]`.

**Task 4: Extract Numerical Data**
4a. Identify significant numerical data points mentioned in the page's text (excluding those already inside tables).
4b. For EACH significant number found:
    *   Extract the `value` as a string (e.g., "123.45", "50%", "$1.2M").
    *   Write a `description` explaining what the number represents, including units or context (e.g., "increase in revenue", "percentage completion").
    *   Extract the surrounding text snippet (~10-20 words) providing immediate context into the `context` field. If no clear surrounding context is available, use `null`.
    *   Add the complete number object to the `numbers` list in the output JSON.
4c. If NO significant numerical data points are found, the `numbers` list MUST be an empty list `[]`.

**Task 5: Set Boolean Flags**
5a. Set `has_tables` to `true` if the `tables` list is NOT empty, otherwise set it to `false`.
5b. Set `has_visuals` to `true` if the `visuals` list is NOT empty, otherwise set it to `false`.
5c. Set `has_numbers` to `true` if the `numbers` list is NOT empty, otherwise set it to `false`.

**Task 6: Assemble Final JSON**
6a. Combine all extracted information (`raw_text`, `has_tables`, `tables`, `has_visuals`, `visuals`, `has_numbers`, `numbers`) into a single JSON object adhering precisely to the schema defined below.

--- Output Format ---
Return ONLY a single, valid JSON object. Do NOT include any text, comments, or markdown formatting before or after the JSON object. It MUST strictly follow this schema:

```json
{{
  "raw_text": "string", // REQUIRED. Full text from the page.
  "has_tables": "boolean", // REQUIRED. True if 'tables' list is not empty.
  "has_visuals": "boolean", // REQUIRED. True if 'visuals' list is not empty.
  "has_numbers": "boolean", // REQUIRED. True if 'numbers' list is not empty.
  "tables": [ // REQUIRED list (can be empty []).
    {{
      "table_id": "string", // REQUIRED. e.g., "page_{page_num}_table_1"
      "table_content": "string", // REQUIRED. Table content as GitHub Flavored Markdown.
      "title": "string | null", // Optional
      "summary": "string | null", // Optional
      "page_number": "integer" // REQUIRED. e.g., {page_num}
    }}
  ],
  "visuals": [ // REQUIRED list (can be empty []).
    {{
      "visual_id": "string", // REQUIRED. e.g., "page_{page_num}_visual_1"
      "type": "string", // REQUIRED. e.g., "bar chart", "image"
      "description": "string", // REQUIRED. Detailed description.
      "data_summary": "string", // REQUIRED. Summary of data or "N/A".
      "page_numbers": ["integer"], // REQUIRED. e.g., [{page_num}]
      "source_url": "string | null", // Optional
      "alt_text": "string | null" // Optional
    }}
  ],
  "numbers": [ // REQUIRED list (can be empty []).
    {{
      "value": "string", // REQUIRED. Numerical value as string.
      "description": "string", // REQUIRED. What the number represents.
      "context": "string | null" // Optional. Surrounding text.
    }}
  ]
}}

--- Critical Rules ---
VALID JSON ONLY: The entire output MUST be a single JSON object starting with {{ and ending with }}.
SCHEMA ADHERENCE: Follow the schema EXACTLY. All REQUIRED fields must be present with the correct data types. Use null for optional string fields if no value is applicable/found. Use [] for empty lists.
MARKDOWN TABLES: Table content MUST be formatted as GitHub Flavored Markdown in the table_content field.
TEXTUAL BASIS: All extracted information must be derived SOLELY from the provided page content. Do not infer or add external knowledge.

Generate the JSON output for page {page_num} now:
"""


async def extract_structure_from_text_content(client: genai.Client, text_content: str, page_num: int, filename: str) -> PageContent:
    """
    Uses Gemini to extract raw text and structural elements (tables, visuals, numbers)
    from the input TEXT content. Populates PageContent WITHOUT subsections initially.
    Includes fallback for pure text preservation.
    """
    function_name = "extract_structure_from_text_content"
    logging.debug(f"[{function_name}] Starting initial structure extraction from text for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    if not text_content or not text_content.strip():
         logging.warning(f"[{function_name}] Input text_content is empty for page {page_num} of '{filename}'. Returning basic PageContent.")
         page_data["raw_text"] = "" # Set raw text to empty string
         try:
              return PageContent(**page_data)
         except ValidationError as ve:
             logging.error(f"[{function_name}] Failed validation even for empty text page {page_num}: {ve}")
             # Return hardcoded error if basic creation fails
             return create_error_page(page_num, f"Validation error on empty text input: {ve}")


    # Limit text content size
    MAX_TEXT_LENGTH = 100000
    if len(text_content) > MAX_TEXT_LENGTH:
        logging.warning(f"[{function_name}] Truncating text input for page {page_num} of '{filename}' from {len(text_content)} to {MAX_TEXT_LENGTH} chars.")
        text_content_for_api = text_content[:MAX_TEXT_LENGTH] + "\n... [Content Truncated] ..."
    else:
        text_content_for_api = text_content

    try:
        prompt_for_page = page_structure_prompt.format(page_num=page_num)

        # --- Attempt structured extraction (raw_text + elements) from text ---
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.5-flash-preview-04-17", # Flash should be fine for this structure extraction
            # Pass the text AS THE PRIMARY CONTENT and the prompt as instructions
            # Check Gemini docs: sometimes it's better to put instructions first.
            # Option 1: Text first (like image)
            # contents=[text_content_for_api, types.Part.from_text(prompt_for_page)],
            # Option 2: Prompt first (often better for text-only models)
            contents=[
                types.Content(parts=[types.Part.from_text(text=prompt_for_page)]),
                types.Content(parts=[types.Part.from_text(text=text_content_for_api)])
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {}
            if not isinstance(data, dict): data = {}

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text") # Get raw text from LLM's perspective
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False) # Less likely but check
            page_data["has_numbers"] = data.get("has_numbers", False)
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])
            page_data["subsections"] = [] # Explicitly ensure subsections are empty here

            # If LLM failed to echo raw_text, use the original input text as fallback
            if not page_data["raw_text"]:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num}. Using original input text.")
                 page_data["raw_text"] = text_content # Use original full text

            # Create PageContent object (validates tables, visuals, numbers)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted initial structure from text for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for text structure extraction on page {page_num}.")
            raise ValueError("LLM response had no candidates") # Trigger fallback

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.warning(f"[{function_name}] Initial structure extraction from text failed for page {page_num}: {e}. Falling back to text preservation.")

        # --- Fallback: Preserve original text ---
        # Create a PageContent object with ONLY raw_text and page_number populated
        page_data["raw_text"] = text_content # Ensure original text is saved
        page_data["has_tables"] = False # Reset flags
        page_data["has_visuals"] = False
        page_data["has_numbers"] = False
        page_data["tables"] = []
        page_data["visuals"] = []
        page_data["numbers"] = []
        page_data["subsections"] = []

        try:
             # Return a valid PageContent object even in fallback
             return PageContent(**page_data)
        except ValidationError as ve_fallback:
            logging.error(f"[{function_name}] Failed to create fallback PageContent for text page {page_num}: {ve_fallback}")
            return create_error_page(page_num, f"Text structure extraction failed AND fallback creation failed: {ve_fallback}")


async def extract_page_content_from_memory(client, image_part, page_num, filename=None):
    """
    Extracts raw text and structural elements (tables, visuals, numbers) from a single page image.
    Populates PageContent WITHOUT subsections initially.
    """
    function_name = "extract_page_structure" # For logging
    logging.debug(f"[{function_name}] Starting extraction for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    page_structure_prompt = f"""
    Analyze page {page_num} from the input document/image. Extract the full raw text content and identify structural elements like tables, visuals, and key numbers. Return ONLY a valid JSON object with the exact structure specified below.

    JSON Structure:
    {{
    "raw_text": "Complete text content extracted from the page. Preserve formatting like paragraphs and line breaks.",
    "has_tables": true/false,
    "has_visuals": true/false,
    "has_numbers": true/false,
    "tables": [ {{ "table_id": "page_{page_num}_table_N", "table_content": "markdown...", "title": "optional", "summary": "optional", "page_number": {page_num} }} ],
    "visuals": [ {{ "visual_id": "page_{page_num}_visual_N", "type": "...", "description": "...", "data_summary": "...", "page_numbers": [{page_num}], "source_url": null, "alt_text": null }} ],
    "numbers": [ {{ "value": "...", "description": "...", "context": "..." }} ]
    }}

    Rules & Guidelines:
    1.  **Raw Text:** Extract ALL text content into the `raw_text` field. Maintain original paragraph structure.
    2.  **Elements:** Accurately identify and extract all tables, visuals (charts, images, diagrams, equations), and significant numbers.
    3.  **IDs:** Generate unique IDs for tables and visuals starting N from 1 for the page.
    4.  **Structure:** Adhere STRICTLY to the specified JSON structure. Use `[]` for empty lists.
    5.  **JSON Validity:** Ensure the output is ONLY a single, valid JSON object.

    EXAMPLE JSON Schema (Conceptual):
    {{
    "type": "object",
    "properties": {{
        "raw_text": {{ "type": ["string", "null"] }},
        "has_tables": {{ "type": "boolean" }},
        "has_visuals": {{ "type": "boolean" }},
        "has_numbers": {{ "type": "boolean" }},
        "tables": {{ "type": "array", "items": {{ "$ref": "#/definitions/TableData" }} }},
        "visuals": {{ "type": "array", "items": {{ "$ref": "#/definitions/VisualElement" }} }},
        "numbers": {{ "type": "array", "items": {{ "$ref": "#/definitions/NumericalDataPoint" }} }}
    }},
    "required": ["raw_text", "has_tables", "has_visuals", "has_numbers", "tables", "visuals", "numbers"]
    // Define TableData, VisualElement, NumericalDataPoint schemas conceptually if needed
    }}
    """

    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.5-flash-preview-04-17",  # Using 2.5 as specified
            contents=[
                types.Content(parts=[image_part, types.Part.from_text(text=page_structure_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )


        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {} # Handle list response
            if not isinstance(data, dict): data = {} # Ensure data is a dict

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text")
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False)
            page_data["has_numbers"] = data.get("has_numbers", False)
            # Pydantic will validate items during PageContent creation
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])

            # Check if raw_text was extracted
            if page_data["raw_text"] is None:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num} of '{filename}'.")
                 page_data["raw_text"] = "[LLM failed to extract raw text]" # Provide placeholder

            # Create PageContent object (will validate tables, visuals, numbers if models defined)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted structure for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for page {page_num} of '{filename}'.")
            # Return error page content (subsections will be empty)
            return create_error_page(page_num, "LLM response had no candidates.")

    except (json.JSONDecodeError, ValidationError, Exception) as e:
        logging.error(f"[{function_name}] Failed for page {page_num} of '{filename}': {e}", exc_info=True)
        # Return error page content (subsections will be empty)
        error_page = create_error_page(page_num, f"Initial extraction failed: {e}")
        # Ensure raw_text is None or set appropriately in error page if needed downstream
        error_page.raw_text = page_data.get("raw_text") # Preserve if extracted before error
        return error_page


def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List[Dict[str, Any]]] = None) -> PageContent:
    """
    Creates a simplified PageContent object representing an error state,
    placing the error details within a single subsection.
    """
    logging.error(f"Creating error page for page {page_num} due to: {error_msg}")
    # --- Format the detailed error text ---
    error_text = f"Error processing page {page_num}: {error_msg}"

    if validation_errors:
        # Ensure validation_errors is a list of dicts before processing
        if isinstance(validation_errors, list) and all(isinstance(item, dict) for item in validation_errors):
            try:
                error_details = "\n\nValidation Errors:\n" + "\n".join(
                    # Safely access keys with .get()
                    f"- Field '{'.'.join(map(str, err.get('loc', ['unknown']))) if err.get('loc') else 'unknown'}': {err.get('msg', 'No message')}"
                    for err in validation_errors
                )
                error_text += error_details
                logging.debug(f"Formatted validation errors for page {page_num}: {error_details}")
            except Exception as format_error:
                 logging.error(f"Could not format validation errors for page {page_num}: {format_error}")
                 error_text += f"\n(Could not format validation errors: {format_error})"
        else:
            logging.warning(f"Received validation_errors in unexpected format for page {page_num}: {type(validation_errors)}. Skipping details.")
            error_text += "\n(Validation error details format incorrect)"

    # --- Create the error subsection ---
    error_subsection = Subsection(
        subsection_id=f"page_{page_num}_section_error",
        order=1,
        title="Processing Error",
        text=error_text, # Place the detailed error message here
        description=f"An error occurred processing page {page_num}.", # Brief summary
        is_cutoff=False,
        referenced_visuals=[],
        referenced_tables=[]
    )

    # --- Initialize the simplified PageContent object ---
    error_page = PageContent(
        page_number=page_num,
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        subsections=[error_subsection] # Put the error subsection in the list
        # Removed old fields: text, title, topics, summary, entities, dates, etc.
    )

    logging.debug(f"Generated error PageContent object for page {page_num}")
    return error_page

async def merge_cutoff_subsections(pages: List[PageContent]) -> List[PageContent]:
    """
    Merge subsections marked as cut off with the first subsection of the next page.
    Operates on the simplified PageContent model.
    """
    if not pages or len(pages) < 2:
        return pages # Nothing to merge if less than 2 pages

    # Ensure pages are sorted by page number
    # Use getattr for safe access in case of malformed objects, though Pydantic should ensure existence
    try:
        sorted_pages = sorted(pages, key=lambda p: getattr(p, 'page_number', float('inf')))
    except Exception as sort_err:
        logging.error(f"Failed to sort pages for merging: {sort_err}. Returning original order.")
        sorted_pages = pages # Proceed with original order if sorting fails


    logging.debug(f"Attempting to merge subsections across {len(sorted_pages)} pages.")
    merged_something = False # Flag to track if any merge happened

    # Iterate through pages up to the second-to-last one
    for i in range(len(sorted_pages) - 1):
        current_page = sorted_pages[i]
        next_page = sorted_pages[i + 1]

        # Basic validation of pages and subsections lists
        if not isinstance(current_page, PageContent) or not isinstance(next_page, PageContent):
             logging.warning(f"Skipping merge check between page {i+1} and {i+2} due to invalid PageContent object types.")
             continue

        if not current_page.subsections or not next_page.subsections:
            # logging.debug(f"Skipping merge check between page {current_page.page_number} and {next_page.page_number}: one has no subsections.")
            continue # Nothing to merge if either page lacks subsections

        # Sort subsections by order for safety, though they should be ordered
        current_subsections = sorted(current_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))
        next_subsections = sorted(next_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))

        # Get the last subsection of the current page and first of the next
        last_subsection = current_subsections[-1]
        first_next_subsection = next_subsections[0]

        # Check if the last subsection is actually marked as cut off
        if getattr(last_subsection, 'is_cutoff', False):
            logging.info(f"Merging cutoff subsection '{last_subsection.subsection_id}' (Page {current_page.page_number}) with '{first_next_subsection.subsection_id}' (Page {next_page.page_number})")

            # --- Perform the merge ---
            # Append text (Use space or newline as appropriate)
            last_subsection.text += "\n" + getattr(first_next_subsection, 'text', '') # Use getattr for safety

            # Mark the merged subsection as not cut off
            last_subsection.is_cutoff = False

            # Merge references (handle potential missing attributes)
            last_subsection.referenced_visuals.extend(getattr(first_next_subsection, 'referenced_visuals', []))
            last_subsection.referenced_tables.extend(getattr(first_next_subsection, 'referenced_tables', []))
            # Ensure uniqueness if needed:
            # last_subsection.referenced_visuals = list(set(last_subsection.referenced_visuals))
            # last_subsection.referenced_tables = list(set(last_subsection.referenced_tables))


            # --- Update the next page ---
            # Remove the first subsection from the next page
            next_page.subsections = next_subsections[1:]

            # Renumber the order of remaining subsections in the next page
            for idx, subsection in enumerate(next_page.subsections):
                 subsection.order = idx + 1

            merged_something = True # Mark that a merge occurred

    if merged_something:
         logging.info("Finished merging cut-off subsections.")
    else:
         logging.debug("No cut-off subsections found to merge.")

    return sorted_pages # Return the list of pages (potentially modified)

async def extract_chapters_from_subsections(client: genai.Client, pages: List[PageContent]) -> List[Chapter]:
    """
    Groups subsections from multiple pages into logical chapters using an LLM.
    Input `pages` should already have cut-off subsections merged.
    """
    all_subsections_with_context = []
    subsection_map = {} # For quick lookup later

    logging.debug(f"Extracting subsections from {len(pages)} pages for chapter generation.")
    for page in pages:
        if not isinstance(page, PageContent) or not page.subsections:
            continue # Skip pages without valid subsections
        for subsection in page.subsections:
             if isinstance(subsection, Subsection):
                # Store the subsection object itself and context for the prompt
                context = {
                    "subsection_id": subsection.subsection_id,
                    "title": subsection.title,
                    "description": subsection.description, # Use the description field
                    "page_number": page.page_number, # Get page number from parent
                    # Optional: Add a very short text preview if description isn't enough context
                    # "text_preview": subsection.text[:50] + "..." if len(subsection.text) > 50 else subsection.text
                }
                all_subsections_with_context.append(context)
                subsection_map[subsection.subsection_id] = subsection # Store the actual object
             else:
                  logging.warning(f"Skipping invalid subsection object on page {page.page_number}")


    if not all_subsections_with_context:
        logging.warning("No valid subsections found across all pages to form chapters.")
        # Return an empty list or a single default chapter if required by downstream logic
        # For now, returning empty list. Adjust if a default chapter is mandatory.
        return []

    # Sort subsections by page number and order (using the context dict)
    # This ensures the LLM sees them in the correct document order.
    sorted_subsections_context = sorted(all_subsections_with_context, key=lambda s: (s["page_number"], s.get("order", 0))) # Use get for order safety

    logging.info(f"Sending {len(sorted_subsections_context)} subsections to LLM for chapter structuring.")

    # Create prompt for chapter extraction using NEW structure
    # Pass subsection_id, title, and description
    chapter_prompt = f"""
    Analyze the following list of subsections extracted sequentially from a document. Each subsection includes its ID, title, and a brief description. Your goal is to group these subsections into logically coherent chapters.

    Subsections List:
    ```json
    {json.dumps(sorted_subsections_context, indent=2)}
    ```

    Instructions:
    1. Group the provided subsections into 3 to 10 logical chapters.
    2. Base the grouping on topic cohesion and narrative flow. Subsections within a chapter should relate to a common theme or part of the document's structure.
    3. Each chapter MUST contain at least one subsection.
    4. Ensure every subsection ID from the input list is assigned to exactly ONE chapter.
    5. Maintain the original relative order of subsections within and across chapters as much as possible.
    6. For each chapter created, provide:
        - A unique `chapter_id` (e.g., "chapter_1", "chapter_2").
        - A concise and descriptive `title` (max 7 words) reflecting the chapter's main theme.
        - A brief `summary` (1-2 sentences) describing the chapter's content or purpose.
        - A list of `subsection_ids` belonging to that chapter, in their original relative order.
    7. Number chapters sequentially starting from 1 using an `order` field.

    Return ONLY a valid JSON array containing the chapter objects, structured exactly like this example:
    ```json
    [
      {{
        "chapter_id": "chapter_1",
        "order": 1,
        "title": "Introduction and Setup",
        "summary": "Provides background information and initial setup instructions.",
        "subsection_ids": ["page_1_section_1", "page_1_section_2", "page_2_section_1"]
      }},
      {{
        "chapter_id": "chapter_2",
        "order": 2,
        "title": "Core Functionality Explained",
        "summary": "Details the main features and how they operate.",
        "subsection_ids": ["page_2_section_2", "page_3_section_1", "page_3_section_2"]
      }}
      // ... more chapters
    ]
    ```
    Ensure the final output is only the JSON array, with no extra text before or after.
    """

    chapters = []
    try:
        logging.debug("Sending chapter extraction request to LLM.")
        # Use the actual client and API call structure
        response = await retry_api_call(
            client.aio.models.generate_content, # Adjusted path might be needed based on client setup
            model="gemini-2.5-flash-preview-04-17", # Or pro, flash might struggle with complex structuring
            contents=[
                types.Content(parts=[types.Part.from_text(text=chapter_prompt)]),
            ], # Prompt first, then text content
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"LLM response for chapters: {json_text[:500]}...") # Log beginning of response
            chapters_data = json.loads(json_text) # Expecting a List[Dict]

            if not isinstance(chapters_data, list):
                 logging.warning(f"LLM returned non-list data for chapters: {type(chapters_data)}. Attempting recovery if dict.")
                 # Handle cases where it might incorrectly return a dict containing the list
                 if isinstance(chapters_data, dict) and "chapters" in chapters_data and isinstance(chapters_data["chapters"], list):
                     chapters_data = chapters_data["chapters"]
                 else:
                     raise ValueError("LLM response for chapters was not a list.")


            processed_subsection_ids = set()
            for i, chapter_data in enumerate(chapters_data):
                # Validate chapter_data structure
                if not isinstance(chapter_data, dict) or not all(k in chapter_data for k in ["title", "summary", "subsection_ids", "order"]):
                     logging.warning(f"Skipping malformed chapter data received from LLM: {chapter_data}")
                     continue

                chapter_subsections = []
                subsection_ids_in_chapter = chapter_data.get("subsection_ids", [])
                if not isinstance(subsection_ids_in_chapter, list):
                     logging.warning(f"Subsection IDs for chapter '{chapter_data.get('title')}' is not a list. Skipping chapter.")
                     continue

                for subsection_id in subsection_ids_in_chapter:
                    if subsection_id in subsection_map:
                        chapter_subsections.append(subsection_map[subsection_id])
                        processed_subsection_ids.add(subsection_id)
                    else:
                        logging.warning(f"Subsection ID '{subsection_id}' from LLM not found in original map for chapter '{chapter_data.get('title')}'.")

                if not chapter_subsections:
                     logging.warning(f"Chapter '{chapter_data.get('title')}' created by LLM has no valid subsections. Skipping.")
                     continue

                # Create the Chapter object (ensure your Chapter model matches)
                chapter = Chapter(
                    chapter_id=chapter_data.get("chapter_id", f"chapter_{chapter_data.get('order', i+1)}"), # Use order for default ID
                    title=chapter_data.get("title", f"Chapter {chapter_data.get('order', i+1)}"),
                    summary=chapter_data.get("summary", "No summary provided."),
                    subsections=chapter_subsections, # List of actual Subsection objects
                    # entity_relationships=[], # Defer relationship extraction
                    order=chapter_data.get("order", i+1)
                )
                chapters.append(chapter)

            # Check if all original subsections were processed
            original_ids = set(subsection_map.keys())
            missing_ids = original_ids - processed_subsection_ids
            if missing_ids:
                 logging.warning(f"LLM failed to assign {len(missing_ids)} subsections to chapters: {missing_ids}. They will be excluded.")
                 # Optionally, create an "Orphaned" chapter for these

            # Sort final chapters list by order
            chapters = sorted(chapters, key=lambda c: c.order)

            if not chapters:
                 logging.warning("LLM processing resulted in zero valid chapters. Creating default chapter.")
                 raise ValueError("LLM returned no valid chapters") # Trigger default chapter creation in except block

            logging.info(f"Successfully extracted {len(chapters)} chapters from LLM.")
            return chapters

        else:
             logging.error("Chapter extraction response from LLM had no candidates.")
             raise ValueError("LLM response had no candidates") # Trigger default

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.error(f"Error processing LLM response or extracting chapters: {e}", exc_info=True)

        # --- Fallback: Create a single default chapter ---
        logging.info("Creating a single default chapter due to chapter extraction error.")
        # Get all original subsections back from the map, sorted
        all_original_subsections = sorted(subsection_map.values(), key=lambda s: (
            # Need page number back for sorting - requires passing it or storing it differently
            # Assuming subsection_map values have access or we modify the storage
            # For now, sort by ID as a fallback sort order
            getattr(s, 'subsection_id', '')
        ))

        default_chapter = Chapter(
            chapter_id="chapter_1_default",
            title="Document Content",
            summary="Content grouped into a single chapter due to an error during automatic chapter structuring.",
            subsections=all_original_subsections, # Assign all subsections here
            # entity_relationships=[],
            order=1
        )
        return [default_chapter]

async def analyze_concept_entity_relationships(
    client: genai.Client,
    chapter: Chapter,
    source_identifier: str, # <-- Added parameter
    max_structure_retries: int = 2,
    retry_delay: int = 3
) -> None:
    """
    Analyzes text within a single chapter to extract key concepts/entities
    and their relationships using an LLM structured prompt. Populates
    chapter.qdrant_nodes in-place. Includes retries on structural errors.

    Args:
        client: The configured GenAI client.
        chapter: The Chapter object containing subsections to analyze.
        source_identifier: Identifier for the source document.
        max_structure_retries: How many times to retry if the LLM returns the wrong JSON structure.
        retry_delay: Base delay (in seconds) between structure retries.
    """
    function_name = "analyze_concept_entity_relationships_refactored"
    chapter_id = getattr(chapter, 'chapter_id', 'unknown_chapter')
    chapter_title = getattr(chapter, 'title', 'Untitled Chapter')
    logging.info(f"[{function_name}] Analyzing relationships for Chapter '{chapter_id}': '{chapter_title}' from source '{source_identifier}'")

    # --- 1. Prepare Input Text ---
    if not chapter.subsections:
        logging.warning(f"[{function_name}] Chapter '{chapter_id}' has no subsections. Skipping analysis.")
        chapter.qdrant_nodes = []
        return


    combined_text = ""
    try:
        subsection_texts = []
        for subsection in chapter.subsections:
            sub_id = getattr(subsection, 'subsection_id', 'unknown_id')
            sub_title = getattr(subsection, 'title', 'Untitled Subsection')
            sub_desc = getattr(subsection, 'description', '')
            sub_text = getattr(subsection, 'text', '')
            subsection_texts.append(f"--- Subsection Start (ID: {sub_id}) ---\nTitle: {sub_title}\nDescription: {sub_desc}\n\n{sub_text}\n--- Subsection End ---")

        combined_text = "\n\n".join(subsection_texts)

        MAX_TEXT_LIMIT = 150000
        if len(combined_text) > MAX_TEXT_LIMIT:
            logging.warning(f"[{function_name}] Truncating combined text for chapter '{chapter_id}' from {len(combined_text)} to {MAX_TEXT_LIMIT} chars.")
            combined_text = combined_text[:MAX_TEXT_LIMIT] + "\n... [Content Truncated] ..."
        elif not combined_text.strip():
            logging.warning(f"[{function_name}] Combined text for chapter '{chapter_id}' is empty after formatting. Skipping analysis.")
            chapter.qdrant_nodes = []
            return
    except Exception as text_prep_error:
        logging.error(f"[{function_name}] Error preparing text for chapter '{chapter_id}': {text_prep_error}", exc_info=True)
        chapter.qdrant_nodes = []
        return

    # --- 2. Extract Chapter Number for Node IDs ---
    chapter_num_str = "1"
    if chapter_id and isinstance(chapter_id, str):
        match = re.search(r'[_-](\d+)$', chapter_id)
        if match: chapter_num_str = match.group(1)
    logging.debug(f"[{function_name}] Using chapter number '{chapter_num_str}' for node IDs in chapter '{chapter_id}'.")

    # --- 3. Define the LLM Prompt *Inside* the Function ---
    # Define the JSON schema clearly within the prompt
    json_schema_description = """
```json
{
  "qdrant_nodes": [
    {
      "node_id": "string", // REQUIRED. Format: "{type}_{sanitized_name}_ch{chapter_num_str}". Example: "concept_machine_learning_ch1". Sanitized name: lowercase, spaces/symbols to underscores, max 50 chars.
      "node_type": "string", // REQUIRED. "entity" or "concept".
      "name": "string", // REQUIRED. Original human-readable name.
      "description": "string", // REQUIRED. 2-4 sentence description based ONLY on chapter text.
      "chapter_id_context": "string", // REQUIRED. Use the provided Chapter ID. Example: "{chapter_id}"
      "linked_nodes": [ // List, REQUIRED (can be empty [] if no links found). Max 7 links per node.
        {
          "target_node_id": "string", // REQUIRED. node_id of another node defined in THIS list.
          "relationship_description": "string", // REQUIRED. Description of link from source to target.
          "relationship_keywords": ["string"], // List, OPTIONAL. 1-3 keywords. Defaults to empty list [].
          "relationship_strength": "float" // REQUIRED. Link strength/importance (1.0-10.0).
        }
      ]
    }
  ]
}
```""" # Note: Made linked_nodes required but can be empty list for better schema validation

    entity_prompt = f"""
--- Goal ---
Analyze the text from the single document chapter provided below. Your task is to identify key specific entities (e.g., organizations, people, products) AND important abstract concepts (e.g., processes, topics, ideas). Extract these as nodes and determine the direct relationships *between these nodes within this chapter*.

--- Input Data ---
Chapter Title: {chapter_title}
Chapter ID: {chapter_id}
Source Identifier: {source_identifier} # Include the source identifier
Chapter Content:
{combined_text}
--- End Input Data ---

--- Instructions ---
Perform the following tasks based *only* on the provided "Chapter Content":

**Task 1: Extract Nodes and Intra-Chapter Relationships**
1a. **Identify Nodes:** Scan the text and identify all distinct key entities and abstract concepts mentioned.
1b. **Define Node Properties:** For each node identified:
    *   Determine its `name` (original, human-readable).
    *   Determine its `node_type` (must be "entity" or "concept").
    *   Write a concise `description` (2-4 sentences) summarizing its role/meaning *within this chapter's context*.
    *   Generate its unique `node_id` using the format: `{{type}}_{{sanitized_name}}_ch{chapter_num_str}`. (Sanitize name: lowercase, spaces/symbols to '_', limit length).
    *   Set `chapter_id_context` to exactly "{chapter_id}".
    *   Set `document_context` to exactly "{source_identifier}". # Set the document context here

1c. **Identify Relationships (Linked Nodes):** For each node created (the *source* node), find its most important (max 5-7) direct relationships *to other nodes also created in this Task* (the *target* nodes). Focus on explicit or clearly implied connections within the text.
    *   For each relationship found, create a `linked_nodes` object containing:
        *   `target_node_id`: The `node_id` of the target node (must be from this chapter's identified nodes).
        *   `relationship_description`: **CRITICAL: Describe the SPECIFIC ACTION or NATURE of the relationship as stated or clearly implied in the text.** Look for verbs connecting the source and target. Explain *how* or *why* they are linked in this context.
            *   **AVOID GENERIC TERMS** like "related to", "associated with", "used in", "interacts with" if more detail exists.
            *   **GOOD EXAMPLES:** "Developed at CafeCorner LLC", "Leverages GPT-40 for text generation", "Stores embeddings in Qdrant Vector DB", "Integrates Llamalndex components", "Enhanced Pitchbook data distribution", "Provided blueprint for JPMorgan's ChatIQ".
            *   **BAD EXAMPLES:** "Used GPT-40", "Related to Qdrant", "Associated with Llamalndex".
        *   `relationship_keywords`: (Optional) 1-3 keywords summarizing the relationship *type* (e.g., ["dependency", "causation", "component", "data source", "integration"]). Default to `[]`.
        *   `relationship_strength`: Estimate the importance/strength of this specific link (float from 1.0 to 10.0, based on textual emphasis).
    *   If a node has no clear relationships to other identified nodes in this chapter, provide an empty list `[]` for its `linked_nodes` field.
1d. **Assemble Output:** Format all extracted node objects (including their `linked_nodes`) into a single JSON object according to the schema specified below.

--- Output Format ---
Return ONLY a single, valid JSON object. It MUST contain exactly one top-level key: "qdrant_nodes". The value associated with "qdrant_nodes" MUST be a JSON array of node objects. Adhere STRICTLY to this schema:

{json_schema_description} # Assuming json_schema_description includes the updated QdrantNode with document_context

--- Critical Rules ---
*   **VALID JSON ONLY:** Your entire output must be a single JSON object starting with `{{` and ending with `}}`. No extra text, comments, or markdown formatting outside the JSON structure.
*   **SCHEMA ADHERENCE:** Follow the provided schema precisely. All REQUIRED fields must be present with the correct data types (string, float, list). `linked_nodes` list is REQUIRED but can be empty `[]`. `relationship_keywords` is OPTIONAL (defaults to `[]`).
*   **NODE ID CONSISTENCY:** Every `target_node_id` in any `linked_nodes` list MUST correspond to a `node_id` that is also defined within the main `qdrant_nodes` list of your response.
*   **RELATIONSHIP LIMIT:** Max 5-7 `linked_nodes` objects per source node.
*   **TEXTUAL BASIS:** All information (`description`, `relationship_description`) must be derived SOLELY from the provided "Chapter Content". Do not add external knowledge. The relationship description must reflect the specific interaction mentioned in the text.

Generate the JSON output now:
"""


    # --- 4. Call LLM and Process Response with Retries (Logic largely unchanged) ---
    chapter.qdrant_nodes = [] # Initialize
    final_data = None

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Attempt {attempt + 1}/{max_structure_retries} for chapter '{chapter_id}'")
        try:
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.5-flash-preview-04-17", # Keep capable model
                contents=[types.Content(parts=[types.Part.from_text(text=entity_prompt)])],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates (Attempt {attempt + 1}) for chapter '{chapter_id}'.")
                final_data = None; break

            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"[{function_name}] Cleaned LLM response (Attempt {attempt + 1}, Chapter {chapter_id}): {json_text[:500]}...")

            try: data = json.loads(json_text)
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError attempt {attempt + 1} for chapter '{chapter_id}': {json_err}.")
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: logging.error(f"[{function_name}] Failed JSON decode after retries for chapter '{chapter_id}'."); final_data = None; break

            # Structure Check
            if isinstance(data, dict) and "qdrant_nodes" in data:
                 logging.info(f"[{function_name}] Received correct structure attempt {attempt + 1} for chapter '{chapter_id}'.")
                 final_data = data; break
            else:
                 logging.warning(f"[{function_name}] Incorrect JSON structure attempt {attempt+1} for chapter '{chapter_id}'. Got: {type(data)}. Structure: {str(data)[:300]}...")
                 if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                 else: logging.error(f"[{function_name}] Failed structure check after retries for chapter '{chapter_id}'."); final_data = None; break

        except Exception as e:
            logging.error(f"[{function_name}] Error attempt {attempt + 1} for chapter '{chapter_id}': {e}", exc_info=True)
            if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else: logging.error(f"[{function_name}] Failed after retries for chapter '{chapter_id}': {e}"); final_data = None; break

    # --- 5. Process Validated Data (Logic unchanged) ---
    if final_data is not None:
        qdrant_nodes_data = final_data.get("qdrant_nodes")
        if isinstance(qdrant_nodes_data, list):
            # --- 6. Validate Individual Nodes (Logic unchanged) ---
            validated_nodes: List[QdrantNode] = []
            node_ids_in_chapter = set()
            for node_dict in qdrant_nodes_data:
                if not isinstance(node_dict, dict): continue
                try:
                    # Set linked_nodes to empty list if missing
                    ln = node_dict.setdefault('linked_nodes', [])
                    if not isinstance(ln, list): node_dict['linked_nodes'] = []
                    
                    # Set keywords to empty list if missing
                    for link_item in node_dict['linked_nodes']:
                        if isinstance(link_item, dict):
                            kw = link_item.setdefault('relationship_keywords', [])
                            if not isinstance(kw, list): link_item['relationship_keywords'] = []

                    # --- ADD document_context BEFORE VALIDATION ---
                    node_dict['document_context'] = source_identifier # Assign the source identifier
                    # ---------------------------------------------

                    node_obj = QdrantNode(**node_dict)
                    validated_nodes.append(node_obj)
                    node_ids_in_chapter.add(node_obj.node_id)
                except ValidationError as node_error:
                    # Use source_identifier in log
                    logging.warning(f"[{function_name}] Pydantic validation failed node in chapter '{chapter_id}' of source '{source_identifier}': {node_error}. Skipping. Data: {node_dict}")
                except Exception as unexpected_err:
                    logging.error(f"[{function_name}] Unexpected error processing node data in chapter '{chapter_id}': {unexpected_err}. Data: {node_dict}", exc_info=True)

            # --- 7. Post-Validation Link Check (Logic unchanged) ---
            final_validated_nodes = []
            MAX_LINKS = 10
            for node in validated_nodes:
                 valid_links = []
                 link_count = 0
                 for link in node.linked_nodes:
                      if link_count >= MAX_LINKS: break
                      if link.target_node_id not in node_ids_in_chapter:
                           logging.warning(f"[{function_name}] Node '{node.node_id}' links to '{link.target_node_id}' (not defined). Dropping link.")
                           continue
                      valid_links.append(link)
                      link_count += 1
                 node.linked_nodes = valid_links
                 final_validated_nodes.append(node)

            chapter.qdrant_nodes = final_validated_nodes
            logging.info(f"[{function_name}] Assigned {len(chapter.qdrant_nodes)} validated nodes for chapter '{chapter_id}'.")
        else:
             logging.error(f"[{function_name}] 'qdrant_nodes' key present but not a list (type: {type(qdrant_nodes_data)}) for chapter '{chapter_id}'.")
             chapter.qdrant_nodes = []
    else:
        logging.error(f"[{function_name}] All relationship extraction attempts failed for chapter '{chapter_id}'.")
        chapter.qdrant_nodes = []

    # Function modifies chapter in-place


async def analyze_inter_chapter_relationships(client: genai.Client, chapters: List[Chapter], source_identifier: Optional[str] = "Document") -> Optional[Dict[str, Any]]:
    """
    Analyzes inter-chapter relationships AND generates a document summary based on aggregated nodes.
    Updates the `linked_nodes` attribute of the QdrantNode objects within the chapters list *in-place*.
    Returns a dictionary containing the document summary details, or None on failure.
    Uses gemini-2.5-flash-preview-04-17.
    """
    function_name = "analyze_inter_chapter_relationships"
    logging.info(f"[{function_name}] Starting analysis for {len(chapters)} chapters of '{source_identifier}'.")

    # --- 1. Aggregate Node Information & Create Lookup Map ---
    all_nodes_for_prompt = []
    node_map: Dict[str, QdrantNode] = {} # node_id -> QdrantNode object

    for chapter in chapters:
        if not chapter.qdrant_nodes: continue
        for node in chapter.qdrant_nodes:
            if not isinstance(node, QdrantNode): continue
            all_nodes_for_prompt.append({
                "node_id": node.node_id, "node_type": node.node_type,
                "name": node.name, "description": node.description,
                "chapter_id_context": node.chapter_id_context
            })
            if node.node_id in node_map:
                 logging.warning(f"[{function_name}] Duplicate node_id '{node.node_id}' found. Overwriting map entry.")
            node_map[node.node_id] = node

    if len(all_nodes_for_prompt) == 0: # Check for zero nodes specifically
        logging.warning(f"[{function_name}] No nodes found across all chapters for '{doc_name}'. Cannot perform analysis or generate summary.")
        return None # No nodes, no analysis, no summary

    logging.info(f"[{function_name}] Aggregated {len(all_nodes_for_prompt)} nodes from {len(chapters)} chapters for '{doc_name}'.")

    # --- 2. Define the Combined LLM Prompt ---
    # Asks for both inter-chapter links AND document summary details
    inter_chapter_and_summary_prompt = f"""
    ---Goal---
    Analyze the provided list of concepts and entities (nodes) aggregated from all chapters of the document '{source_identifier}'. Perform two tasks:
    1. Identify significant relationships *between* nodes originating from *different* chapters.
    2. Generate a concise document-level summary based on the overall collection of nodes.


    ---Input Data---
    List of node definitions:
    ```json
    {json.dumps(all_nodes_for_prompt, indent=2)}
    ```

    ---Instructions---
    **Task 1: Identify Inter-Chapter Relationships**
    1a. Compare nodes across the entire list, focusing on potential connections between nodes from different `chapter_id_context`.
    1b. Identify pairs (`source_node`, `target_node`) with **different** `chapter_id_context` values that have a direct, meaningful relationship based on their names and descriptions.
    1c. For each significant inter-chapter relationship identified, determine: `source_node_id`, `target_node_id`, `relationship_description`, `relationship_keywords` (list), `relationship_strength` (float 1.0-10.0).
    1d. **Prioritize Quality & Limit:** Focus on the 3-10 *strongest* inter-chapter links *per source node*. Do not exceed 10 outgoing links per source.

    **Task 2: Generate Document Summary**
    2a. Synthesize the information from the *entire list* of nodes (names, descriptions, types, chapter context).
    2b. Generate the following summary components:
        *   `title`: A concise title for the document '{doc_name}' based on the aggregated nodes (max 10 words).
        *   `themes`: A list of 3-7 main themes or topics reflected across all nodes.
        *   `questions`: A list of 2-4 insightful questions a reader might have after understanding the key concepts/entities.
        *   `summary`: A comprehensive summary paragraph (4-8 sentences) synthesizing the key nodes and their overall significance or narrative represented by the aggregation.

    **Output Format:**
    Return ONLY a single valid JSON object containing TWO top-level keys: "inter_chapter_links" and "document_summary_details".

    *   `inter_chapter_links`: A JSON array containing objects, where each object represents one identified *inter-chapter* relationship from Task 1 (with keys: `source_node_id`, `target_node_id`, `relationship_description`, `relationship_keywords`, `relationship_strength`).
    *   `document_summary_details`: A JSON object containing the results from Task 2 (with keys: `title`, `themes`, `questions`, `summary`).

    ---Explicit Exclusions---
    *   DO NOT include intra-chapter relationships (source/target in same chapter) in the `inter_chapter_links` list.
    *   DO NOT include the original node definitions in the output.

    Example Output Format:
    ```json
    {{
      "inter_chapter_links": [
        {{
          "source_node_id": "concept_machine_learning_ch2",
          "target_node_id": "entity_gpu_optimization_ch5",
          "relationship_description": "Utilizes concepts from Chapter 2...",
          "relationship_keywords": ["hardware acceleration", "optimization"],
          "relationship_strength": 8.5
        }}
      ],
      "document_summary_details": {{
        "title": "AI Development and Hardware Optimization",
        "themes": ["Machine Learning", "GPU Computing", "Optimization Techniques", "Acme Corp Strategy"],
        "questions": [
          "How does Acme Corp leverage machine learning?",
          "What are the key challenges in GPU optimization mentioned?"
        ],
        "summary": "The document outlines core machine learning concepts introduced in early chapters, detailing Acme Corp's involvement. Later chapters focus on the technical challenges and solutions related to GPU optimization, showing practical application of the earlier concepts."
      }}
    }}
    ```
    Ensure the final output is only the single JSON object described.
    """

    # --- 3. Call LLM and Process Response ---
    document_summary_output = None # Initialize summary return value

    try:
        logging.debug(f"[{function_name}] Sending combined analysis request for '{doc_name}' to LLM ({'gemini-2.5-flash-preview-04-17'}).")
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.5-flash-preview-04-17", # Use the specified powerful model - **NOTE: Ensure correct model name per Gemini API docs**
            contents=[
                types.Content(parts=[types.Part.from_text(text=inter_chapter_and_summary_prompt)]),
            ],
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text) # Expecting the dict with two keys

            if not isinstance(data, dict):
                raise ValueError("LLM response is not a JSON object.")

            # --- Process Inter-Chapter Links (Updates node_map in-place) ---
            newly_found_links_data = data.get("inter_chapter_links", [])
            if isinstance(newly_found_links_data, list):
                links_added_count, links_skipped_count = _merge_links_into_nodes(newly_found_links_data, node_map)
                logging.info(f"[{function_name}] Merged inter-chapter links for '{doc_name}'. Added: {links_added_count}, Skipped: {links_skipped_count}")
            else:
                logging.warning(f"[{function_name}] Key 'inter_chapter_links' missing or not a list in LLM response for '{doc_name}'.")

            # --- Process Document Summary Details ---
            summary_details_data = data.get("document_summary_details")
            if isinstance(summary_details_data, dict):
                 try:
                     # Validate the extracted summary details (optional but recommended)
                     validated_summary = DocumentSummaryDetails(**summary_details_data)
                     document_summary_output = validated_summary.model_dump() # Store as dict
                     logging.info(f"[{function_name}] Successfully extracted document summary details for '{doc_name}'.")
                 except ValidationError as summary_error:
                     logging.error(f"[{function_name}] Validation failed for document_summary_details from LLM for '{doc_name}': {summary_error}")
                     document_summary_output = None # Indicate summary failure
            else:
                logging.warning(f"[{function_name}] Key 'document_summary_details' missing or not a dict in LLM response for '{doc_name}'.")

        else:
            logging.error(f"[{function_name}] Combined analysis response from LLM had no candidates for '{doc_name}'.")

    except (json.JSONDecodeError, ValidationError, Exception) as e:
        logging.error(f"[{function_name}] Error during combined analysis or parsing LLM response for '{doc_name}': {e}", exc_info=True)
        # Summary generation failed

    # The chapters list has been modified in-place by _merge_links_into_nodes via node_map
    logging.info(f"[{function_name}] Finished analysis for '{source_identifier}'. Returning summary.")
    return document_summary_output

# --- Helper function to merge links (extracted from previous version) ---
def _merge_links_into_nodes(link_data_list: List[Dict], node_map: Dict[str, QdrantNode]) -> (int, int):
    """Merges link data into the linked_nodes of QdrantNode objects in node_map."""
    links_added_count = 0
    links_skipped_count = 0
    MAX_OUTGOING_LINKS_PER_NODE = 10

    if not isinstance(link_data_list, list):
        logging.warning("[_merge_links_into_nodes] Input link_data_list is not a list.")
        return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict):
            logging.warning(f"[_merge_links_into_nodes] Skipping invalid item in link data: {link_data}")
            links_skipped_count += 1
            continue

        source_id = link_data.get("source_node_id")
        target_id = link_data.get("target_node_id")
        desc = link_data.get("relationship_description")
        keywords = link_data.get("relationship_keywords", [])
        strength_str = link_data.get("relationship_strength", "5.0")

        if not all([source_id, target_id, desc]) or not isinstance(source_id, str) or not isinstance(target_id, str):
            logging.warning(f"[_merge_links_into_nodes] Skipping incomplete/invalid link data: {link_data}")
            links_skipped_count += 1
            continue

        try: strength = float(strength_str)
        except (ValueError, TypeError): strength = 5.0

        source_node = node_map.get(source_id)
        target_node = node_map.get(target_id)

        if not source_node or not target_node:
            logging.warning(f"[_merge_links_into_nodes] Source ('{source_id}') or Target ('{target_id}') not found. Skipping.")
            links_skipped_count += 1
            continue

        if source_node.chapter_id_context == target_node.chapter_id_context:
            logging.debug(f"[_merge_links_into_nodes] Skipping INTRA-chapter link: {source_id} -> {target_id}")
            links_skipped_count += 1
            continue

        # Add Forward Link
        if source_node.linked_nodes is None: source_node.linked_nodes = []
        if len(source_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == target_id for link in source_node.linked_nodes):
            try:
                forward_link = LinkedNode(**link_data) # Directly use link_data if keys match LinkedNode
                source_node.linked_nodes.append(forward_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added FORWARD link: {source_id} -> {target_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating forward LinkedNode {source_id} -> {target_id}: {ve}. Skipping.")
                 links_skipped_count += 1
        else:
             links_skipped_count += 1 # Skipped due to limit or duplicate

        # Add Reverse Link
        if target_node.linked_nodes is None: target_node.linked_nodes = []
        if len(target_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == source_id for link in target_node.linked_nodes):
            try:
                reverse_desc = f"Related to '{source_node.name}' from {source_node.chapter_id_context} ({desc})"
                reverse_link = LinkedNode(target_node_id=source_id, relationship_description=reverse_desc, relationship_keywords=keywords, relationship_strength=strength)
                target_node.linked_nodes.append(reverse_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added REVERSE link: {target_id} -> {source_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating reverse LinkedNode {target_id} -> {source_id}: {ve}. Skipping.")
                 # links_skipped_count += 1 # Optional separate count
        else:
             # links_skipped_count += 1 # Optional separate count
             pass # Skip reverse if limit reached or duplicate


    return links_added_count, links_skipped_count


def _merge_project_links_into_node_dicts(link_data_list: List[Dict],
                                         node_map: Dict[str, Dict],
                                         max_total_links: int = 15) -> Tuple[int, int]:
    """
    Merges INTER-DOCUMENT link data into the 'linked_nodes' lists of node
    dictionaries stored in node_map.
    (Ensure Args and Returns docstrings are present)
    """
    function_name = "_merge_project_links_into_node_dicts"
    links_added_count = 0
    links_skipped_count = 0

    if not isinstance(link_data_list, list):
        logging.warning(f"[{function_name}] Input link_data_list is not a list.")
        return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict):
            logging.warning(f"[{function_name}] Skipping invalid item in link data: {link_data}")
            links_skipped_count += 1; continue

        source_id = link_data.get("source_node_id")
        target_id = link_data.get("target_node_id")
        desc = link_data.get("relationship_description")

        if not all([source_id, target_id, desc]):
            logging.warning(f"[{function_name}] Skipping incomplete/invalid link data: {link_data}")
            links_skipped_count += 1; continue

        source_node_dict = node_map.get(source_id)
        target_node_dict = node_map.get(target_id)

        if not source_node_dict or not target_node_dict:
            logging.warning(f"[{function_name}] Source ('{source_id}') or Target ('{target_id}') node not found in project map. Skipping link.")
            links_skipped_count += 1; continue

        source_doc = source_node_dict.get('document_context', 'source_unknown')
        target_doc = target_node_dict.get('document_context', 'target_unknown')
        if source_doc == target_doc:
            logging.warning(f"[{function_name}] LLM provided an INTRA-document link ({source_id} in '{source_doc}' -> {target_id} in '{target_doc}'). Skipping.")
            links_skipped_count += 1; continue

        try:
            link_to_add = {
                "target_node_id": target_id,
                "relationship_description": desc,
                "relationship_keywords": link_data.get("relationship_keywords", []),
                "relationship_strength": float(link_data.get("relationship_strength", 5.0))
            }
            if not isinstance(link_to_add["relationship_keywords"], list):
                 link_to_add["relationship_keywords"] = []
        except (ValueError, TypeError) as prep_err:
             logging.warning(f"[{function_name}] Error preparing link data {source_id} -> {target_id}: {prep_err}. Skipping.")
             links_skipped_count += 1; continue

        # Add Forward Link
        source_node_dict.setdefault('linked_nodes', [])
        if isinstance(source_node_dict['linked_nodes'], list):
             if len(source_node_dict['linked_nodes']) < max_total_links and not any(lnk.get('target_node_id') == target_id for lnk in source_node_dict['linked_nodes'] if isinstance(lnk, dict)):
                 source_node_dict['linked_nodes'].append(link_to_add)
                 links_added_count += 1
                 logging.debug(f"[{function_name}] Added FORWARD link: {source_id} ({source_doc}) -> {target_id} ({target_doc})")
             else:
                  logging.debug(f"[{function_name}] Skipped FORWARD link (duplicate or max links reached): {source_id} -> {target_id}")
                  links_skipped_count += 1
        else:
             logging.error(f"[{function_name}] 'linked_nodes' in source node {source_id} is not list. Skipping add.")
             links_skipped_count += 1

        # Optional: Add Reverse Link
        target_node_dict.setdefault('linked_nodes', [])
        if isinstance(target_node_dict['linked_nodes'], list):
             if len(target_node_dict['linked_nodes']) < max_total_links and not any(lnk.get('target_node_id') == source_id for lnk in target_node_dict['linked_nodes'] if isinstance(lnk, dict)):
                 reverse_desc = f"Referenced by '{source_node_dict.get('name', source_id)}' from document '{source_doc}' ({desc})"
                 reverse_link_to_add = {
                     "target_node_id": source_id,
                     "relationship_description": reverse_desc,
                     "relationship_keywords": link_to_add["relationship_keywords"],
                     "relationship_strength": link_to_add["relationship_strength"]
                 }
                 target_node_dict['linked_nodes'].append(reverse_link_to_add)
                 links_added_count += 1
                 logging.debug(f"[{function_name}] Added REVERSE link: {target_id} ({target_doc}) -> {source_id} ({source_doc})")
             else:
                  logging.debug(f"[{function_name}] Skipped REVERSE link (duplicate or max links reached): {target_id} -> {source_id}")
        else:
             logging.error(f"[{function_name}] 'linked_nodes' in target node {target_id} is not list. Skipping reverse add.")

    return links_added_count, links_skipped_count


async def generate_project_ontology(client: genai.Client,
                                  documents: List[Dict], # List of finalized document dicts
                                  max_structure_retries: int = 1, # Less retries needed for summary usually
                                  retry_delay: int = 5
                                  ) -> Optional[ProjectOntology]:
    """
    Generates a project-wide ontology by analyzing aggregated nodes from all
    provided documents. Identifies inter-document relationships and synthesizes
    project-level metadata (overview, themes, concepts).

    Args:
        client: The configured GenAI client.
        documents: A list of finalized document dictionaries, each expected to
                   contain 'raw_extracted_content' -> 'summary' -> 'chapters' -> 'qdrant_nodes'.
        max_structure_retries: Max retries if LLM returns incorrect JSON structure.
        retry_delay: Base delay between retries.

    Returns:
        A ProjectOntology object containing the synthesized information and
        aggregated graph nodes, or None if processing fails significantly.
    """
    function_name = "generate_project_ontology_with_graph"
    logging.info(f"[{function_name}] Starting ontology and graph generation for {len(documents)} documents.")

    # --- 1. Aggregate Node Information & Document Details ---
    all_project_nodes_data: List[Dict] = [] # Store node dicts
    node_map: Dict[str, Dict] = {} # node_id -> node dict (for link merging)
    doc_filenames: List[str] = []
    doc_count = 0

    for doc_dict in documents:
        if not isinstance(doc_dict, dict) or "raw_extracted_content" not in doc_dict:
            logging.warning(f"[{function_name}] Skipping invalid document structure: {type(doc_dict)}")
            continue

        raw_content = doc_dict["raw_extracted_content"]
        if not isinstance(raw_content, dict): continue # Skip if structure invalid

        doc_filename = raw_content.get("filename", f"Unknown_Doc_{doc_count}")
        doc_filenames.append(doc_filename)
        doc_count += 1

        summary = raw_content.get("summary")
        if not isinstance(summary, dict):
            logging.warning(f"[{function_name}] Skipping nodes from '{doc_filename}' due to missing or invalid summary dict.")
            continue

        chapters = summary.get("chapters")
        if not isinstance(chapters, list):
            logging.warning(f"[{function_name}] Skipping nodes from '{doc_filename}' due to missing or invalid chapters list.")
            continue

        nodes_found_in_doc = 0
        for chapter in chapters:
            if not isinstance(chapter, dict): continue
            qdrant_nodes = chapter.get("qdrant_nodes")
            if isinstance(qdrant_nodes, list):
                for node_dict in qdrant_nodes:
                    if isinstance(node_dict, dict) and "node_id" in node_dict:
                        node_id = node_dict["node_id"]
                        # Add document context explicitly IF NOT ALREADY PRESENT
                        # Ideally, this should be added during the per-chapter analysis
                        node_dict.setdefault('document_context', doc_filename)

                        if node_id in node_map:
                            # Handle collision - maybe merge descriptions or log warning
                            logging.warning(f"[{function_name}] Duplicate node_id '{node_id}' encountered across documents ('{node_map[node_id].get('document_context')}' vs '{doc_filename}'). Overwriting map entry; link merging might be affected.")
                        node_map[node_id] = node_dict # Store the dict itself
                        all_project_nodes_data.append(node_dict) # Add to list for prompt
                        nodes_found_in_doc += 1
        logging.debug(f"[{function_name}] Aggregated {nodes_found_in_doc} nodes from '{doc_filename}'.")

    if not all_project_nodes_data:
        logging.warning(f"[{function_name}] No nodes found across any documents. Cannot generate graph or detailed ontology.")
        # Return a basic ontology object indicating failure
        return ProjectOntology(
            title=f"Project Ontology (No Node Data)",
            overview="No node data could be extracted from the documents to perform analysis.",
            document_count=doc_count,
            documents=doc_filenames,
            global_themes=[],
            key_concepts=[],
            project_graph_nodes=[] # Empty list
        )

    logging.info(f"[{function_name}] Aggregated {len(all_project_nodes_data)} total nodes from {doc_count} documents for analysis.")

    # --- 2. Define the LLM Prompt (Task-Oriented) ---
    # Prepare node info for prompt (only essential fields)
    nodes_for_prompt = [{
        "node_id": n.get("node_id"),
        "node_type": n.get("node_type"),
        "name": n.get("name"),
        "description": n.get("description"),
        # Use the document_context we added/ensured earlier
        "document_context": n.get("document_context", "unknown_document")
    } for n in all_project_nodes_data if n.get("node_id")] # Filter out invalid entries

    json_schema_description = """
```json
{
  "project_summary_details": { // REQUIRED object
    "title": "string", // REQUIRED. Concise title for the entire project (max 10 words).
    "overview": "string", // REQUIRED. Comprehensive project overview (3-6 sentences) synthesizing all nodes.
    "global_themes": ["string"], // REQUIRED list. 3-7 main themes across all documents derived from nodes.
    "key_concepts": ["string"] // REQUIRED list. 5-10 key concepts central to the project derived from nodes.
  },
  "inter_document_links": [ // REQUIRED list (can be empty []). Links between nodes from DIFFERENT documents.
    {
      "source_node_id": "string", // REQUIRED. node_id of a source node from the input list.
      "target_node_id": "string", // REQUIRED. node_id of a target node from the input list (must be from a DIFFERENT document_context).
      "relationship_description": "string", // REQUIRED. Description of link from source to target.
      "relationship_keywords": ["string"], // List, OPTIONAL. 1-3 keywords. Defaults to [].
      "relationship_strength": "float" // REQUIRED. Link strength/importance (1.0-10.0).
    }
  ]
}
```"""

    ontology_prompt = f"""
--- Goal ---
Analyze the provided list of concepts and entities (nodes), which have been aggregated from ALL documents in a project. Perform two main tasks:
1.  Identify significant relationships connecting nodes that originate from *different* source documents.
2.  Synthesize the information from *all* nodes to generate high-level project metadata (title, overview, themes, key concepts).

--- Input Data ---
List of node definitions from across the project:
```json
{json.dumps(nodes_for_prompt, indent=2)}
Total Documents Analyzed: {doc_count}
--- Instructions ---
Perform the following tasks based on the entire "Input Data":
Task 1: Identify Inter-Document Relationships
1a. Compare nodes across the entire list. Focus specifically on finding meaningful connections between nodes that have different document_context values.
1b. Identify pairs (source_node, target_node) where source_node.document_context is DIFFERENT from target_node.document_context, and where a direct relationship (e.g., usage, comparison, causality, reference) can be inferred from their names and descriptions.
1c. For each significant inter-document relationship identified, create a link object with: source_node_id, target_node_id, relationship_description, relationship_keywords (optional, default []), and relationship_strength (float 1.0-10.0).
1d. Prioritize & Limit: Focus on the 5-10 most significant inter-document links per source node. Do not list weak or trivial connections. Limit total outgoing links per source node added in this step.
1e. Compile these link objects into the inter_document_links list. If no inter-document links are found, this list MUST be empty [].
Task 2: Generate Project Summary Details
2a. Synthesize information from the entire list of input nodes (consider names, descriptions, types, document context).
2b. Generate the following components and place them within the project_summary_details object:
* title: A concise title for the overall project (max 10 words).
* overview: A comprehensive summary paragraph (3-6 sentences) describing the project's main subjects, scope, purpose, and key findings as represented by the aggregated nodes.
* global_themes: A list of 3-7 high-level themes reflected across multiple documents/nodes.
* key_concepts: A list of the 5-10 most important or frequently occurring concepts identified across the project nodes.
--- Output Format ---
Return ONLY a single, valid JSON object. It MUST contain exactly TWO top-level keys: "project_summary_details" and "inter_document_links". Adhere STRICTLY to this schema:
{json_schema_description}
--- Critical Rules ---
VALID JSON ONLY: Your entire output must be a single JSON object starting with {{ and ending with }}. No extra text or comments.
SCHEMA ADHERENCE: Follow the provided schema precisely. All REQUIRED fields/objects/lists must be present (even if lists are empty []).
INTER-DOCUMENT LINKS ONLY: The inter_document_links list must ONLY contain relationships where the source and target nodes have different document_context values based on the input data.
NODE ID CONSISTENCY: source_node_id and target_node_id in links MUST correspond to node_id values present in the input node list.
TEXTUAL BASIS: Base summary details and relationships on the provided node information.
Generate the JSON output now:
"""

    # --- 3. Call LLM and Process Response with Retries ---
    llm_output_data = None
    final_ontology = None

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Attempt {attempt + 1}/{max_structure_retries} for project ontology.")
        try:
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.5-flash-preview-04-17", # Capable model needed
                contents=[types.Content(parts=[types.Part.from_text(text=ontology_prompt)])],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates (Attempt {attempt + 1}) for project ontology.")
                llm_output_data = None; break

            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"[{function_name}] Cleaned LLM response (Attempt {attempt + 1}): {json_text[:500]}...")

            try: data = json.loads(json_text)
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError attempt {attempt + 1}: {json_err}.")
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: logging.error(f"[{function_name}] Failed JSON decode after retries."); llm_output_data = None; break

            # Structure Check
            if isinstance(data, dict) and "project_summary_details" in data and "inter_document_links" in data:
                logging.info(f"[{function_name}] Received correct structure attempt {attempt + 1}.")
                llm_output_data = data; break
            else:
                logging.warning(f"[{function_name}] Incorrect JSON structure attempt {attempt+1}. Got: {type(data)}. Keys: {list(data.keys()) if isinstance(data,dict) else 'N/A'}")
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: logging.error(f"[{function_name}] Failed structure check after retries."); llm_output_data = None; break

        except Exception as e:
            logging.error(f"[{function_name}] Error attempt {attempt + 1}: {e}", exc_info=True)
            if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else: logging.error(f"[{function_name}] Failed after retries: {e}"); llm_output_data = None; break

    # --- 4. Process LLM Output ---
    if llm_output_data is not None:
        summary_details = llm_output_data.get("project_summary_details", {})
        inter_doc_links = llm_output_data.get("inter_document_links", [])

        # Merge the new inter-document links into the nodes stored in node_map
        if isinstance(inter_doc_links, list) and inter_doc_links:
            logging.info(f"[{function_name}] Merging {len(inter_doc_links)} potential inter-document links.")
            # Use a helper that understands merging into existing dicts and checks context
            added_count, skipped_count = _merge_project_links_into_node_dicts(inter_doc_links, node_map)
            logging.info(f"[{function_name}] Link merge results - Added: {added_count}, Skipped: {skipped_count}")
            # The 'all_project_nodes_data' list now contains dicts potentially updated via node_map

        # --- 5. Validate Final Nodes ---
        validated_project_nodes: List[QdrantNode] = []
        final_node_list_for_ontology = list(node_map.values()) # Get potentially updated nodes
        logging.info(f"[{function_name}] Validating final {len(final_node_list_for_ontology)} project nodes with Pydantic.")
        for node_dict in final_node_list_for_ontology:
            try:
                # Ensure linked_nodes is list (might have been added/modified)
                if 'linked_nodes' not in node_dict or not isinstance(node_dict['linked_nodes'], list):
                    node_dict['linked_nodes'] = []
                # Re-validate the potentially modified node dict
                node_obj = QdrantNode(**node_dict)
                validated_project_nodes.append(node_obj)
            except ValidationError as node_error:
                logging.warning(f"[{function_name}] Pydantic validation failed for final project node '{node_dict.get('node_id')}': {node_error}. Skipping node.")
            except Exception as final_val_err:
                logging.error(f"[{function_name}] Unexpected error during final node validation for '{node_dict.get('node_id')}': {final_val_err}", exc_info=True)


        # --- 6. Create Final ProjectOntology Object ---
        try:
            final_ontology = ProjectOntology(
                title=summary_details.get("title", f"Project Ontology ({doc_count} Docs)"),
                overview=summary_details.get("overview", "Overview not generated."),
                document_count=doc_count,
                documents=doc_filenames, # List of filenames
                global_themes=summary_details.get("global_themes", []),
                key_concepts=summary_details.get("key_concepts", []),
                project_graph_nodes=validated_project_nodes # Assign the validated list
            )
            logging.info(f"[{function_name}] Project ontology created successfully with {len(validated_project_nodes)} nodes.")

        except ValidationError as ve:
            logging.error(f"[{function_name}] Failed to validate final ProjectOntology object: {ve}")
            # Return basic ontology on final validation error
            final_ontology = ProjectOntology(
                title="Project Ontology (Validation Error)",
                overview=f"Error creating final ontology object: {ve}",
                document_count=doc_count, documents=doc_filenames,
                global_themes=[], key_concepts=[], project_graph_nodes=[]
            )

    else:
        # LLM processing failed after retries
        logging.error(f"[{function_name}] LLM processing failed for project ontology after all retries.")
        final_ontology = ProjectOntology(
            title="Project Ontology (LLM Error)",
            overview="Failed to generate project overview and relationships via LLM.",
            document_count=doc_count, documents=doc_filenames,
            global_themes=[], key_concepts=[], project_graph_nodes=[]
        )

    return final_ontology



###############################
# DOCUMENT PROCESSING FUNCTIONS
###############################

async def process_single_page_with_semaphore(semaphore, client, page_info: dict, uploaded_files):
    """
    Process a single page using semaphore for rate limiting.
    Returns a dictionary: {"info": page_info, "result": PageContent object}.
    Handles potential errors during processing. Uses 'source_identifier'.
    """
    # Get source_identifier and page_num early for use in logic and error reporting
    source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
    page_num = page_info.get("page_num", 0)

    async with semaphore:
        result_page_obj = None
        try:
            # Call process_single_pdf_page, which returns a DICTIONARY
            # Pass the whole page_info which now contains source_identifier
            page_dict = await process_single_pdf_page(client, page_info, uploaded_files)

            page_dict.setdefault("page_number", page_num)
            logging.debug(f"Attempting Pydantic validation for page {page_num} dict: {list(page_dict.keys())}")
            result_page_obj = PageContent(**page_dict)
            # Use source_identifier in log
            logging.info(f"Successfully validated PageContent for page {page_num} of {source_identifier}")

        except pydantic.ValidationError as ve:
            # Use source_identifier in log/error
            logging.error(f"Pydantic validation error creating PageContent for page {page_num} of {source_identifier} from returned dictionary: {ve}")
            error_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Pydantic Validation Error in '{source_identifier}': {ve}", # Include identifier in msg
                validation_errors=ve.errors()
            )
            result_page_obj = error_page_obj # Already a PageContent object

        except Exception as e:
            # Use source_identifier in log/error
            logging.error(f"Unexpected error in semaphore task for page {page_num} of {source_identifier}: {str(e)}")
            error_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Unexpected processing error in semaphore task for '{source_identifier}': {str(e)}" # Include identifier in msg
            )
            result_page_obj = error_page_obj

        # Return the result packaged with the original info
        logging.debug(f"Finished processing for page {page_num} of {source_identifier}")
        return {"info": page_info, "result": result_page_obj}


async def process_single_pdf_page(client, page_info: dict, uploaded_files) -> dict:
    """
    Process a single PDF page image using Gemini with fallback options, returning extracted data as a dictionary.
    Uses 'source_identifier' from page_info.
    """
    page_num = page_info.get("page_num", 0)
    source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
    page_dict = {}

    try:
        if "image_b64" not in page_info or not page_info["image_b64"]:
             st.warning(f"Missing image data for page {page_num} of {source_identifier}.")
             error_page_obj = create_error_page(page_num=page_num, error_msg="Missing or empty image data (image_b64)")
             return error_page_obj.model_dump()

        if "page_num" not in page_info:
             st.warning(f"Missing page number for a page in {source_identifier}.")
             error_page_obj = create_error_page(page_num=0, error_msg="Missing page number information")
             return error_page_obj.model_dump()

        img_data = base64.b64decode(page_info["image_b64"])
        image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")

        # Call Gemini with fallback strategy, passing source_identifier
        # Assumes process_page_with_fallback accepts source_identifier now
        page_content_obj = await process_page_with_fallback(
            client, image_part, page_num, source_identifier # <-- Pass identifier
        )

        page_dict = page_content_obj.model_dump()
        return page_dict

    except Exception as e:
        st.error(f"Error during Gemini processing or data handling for page {page_num} of {source_identifier}: {str(e)}")
        error_page_obj = create_error_page(
            page_num=page_num,
            error_msg=f"Gemini API or processing error in '{source_identifier}': {str(e)}",
        )
        page_dict = error_page_obj.model_dump()
        return page_dict

async def finalize_document(client: genai.Client, doc_name: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Merges cut-off subsections.
    2. Extracts chapters based on subsections.
    3. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    4. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    5. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{doc_name}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": doc_name,
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }
    # Attempt to dump incoming pages immediately for error reporting, handle failures
    try:
        initial_pages_dump = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in pages]
        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed to initially dump pages for '{doc_name}' for error reporting: {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{doc_name}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{doc_name}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result

        # --- Step 1: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Merging subsections for '{doc_name}'.")
        merged_pages = await merge_cutoff_subsections(pages) # Operates on PageContent objects
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{doc_name}'. Using original pages.")
             merged_pages = pages

        # --- Step 2: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Extracting chapters for '{doc_name}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{doc_name}'. Cannot proceed with relationship/summary analysis.")
             # Return result with merged pages but no summary/chapter info
             final_result["raw_extracted_content"]["pages"] = [p.model_dump(mode='json') for p in merged_pages]
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = { # Basic placeholder summary
                  "title": f"Processing Incomplete: {doc_name}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result

        # --- Step 3: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Analyzing relationships within {len(chapters)} chapters for '{doc_name}'.")
        # This loop modifies chapter.qdrant_nodes in place
        analysis_tasks = [analyze_concept_entity_relationships(client, chapter) for chapter in chapters]
        await asyncio.gather(*analysis_tasks, return_exceptions=True) # Run in parallel, log errors if any task fails
        # Check results? For now, we assume modification happened or warnings were logged internally.

        # --- Step 4: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Analyzing inter-chapter relationships and generating summary for '{doc_name}'.")
        # This function modifies chapter.qdrant_nodes again (adds inter-chapter links) AND returns the summary dict
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, doc_name)

        # --- Step 5: Assemble Final Result ---
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{doc_name}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {doc_name}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated before dumping
        summary_dict['chapters'] = [ch.model_dump(mode='json') for ch in chapters if isinstance(ch, Chapter)]

        # Add document-level relationships placeholder (could be populated later from qdrant_nodes)
        summary_dict['entity_relationships'] = []

        # Dump the final state of merged pages
        final_pages_as_dicts = [p.model_dump(mode='json') for p in merged_pages if isinstance(p, PageContent)]

        final_result = {
            "raw_extracted_content": {
                "filename": doc_name,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{doc_name}'.")
        return final_result

    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{doc_name}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {doc_name}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": [], "entity_relationships": []
            }
        # Keep the initial pages dump in the error case
        return final_result


async def extract_pages_from_pdf_bytes(pdf_bytes, source_identifier: str): # Renamed parameter
    """Extract pages from PDF bytes using PyMuPDF."""
    loop = asyncio.get_running_loop()
    pages = []
    try:
        pdf_stream = io.BytesIO(pdf_bytes)
        # Run synchronous PyMuPDF code in an executor thread
        def sync_extract():
            # Use source_identifier in logs/errors if needed
            doc = pymupdf.open(stream=pdf_stream, filetype="pdf")
            extracted = []
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=150, alpha=False)
                    img_bytes = pix.tobytes("jpeg", jpg_quality=85)
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    extracted.append({
                        "page_num": page_num + 1,
                        "image_b64": img_b64,
                        "source_identifier": source_identifiedr, # <-- Use new key
                        # "doc_bytes": None
                    })
                except Exception as page_e:
                     # Use source_identifier in warning
                     st.warning(f"Error extracting page {page_num+1} from {source_identifier}: {page_e}")
            doc.close()
            return extracted

        pages = await loop.run_in_executor(None, sync_extract)
        return pages

    except Exception as e:
        # Use source_identifier in error
        st.error(f"Error opening or processing PDF {source_identifier}: {str(e)}")
        return []

# Non-PDF Document Processing Functions

# Assume necessary imports and models are defined

async def process_single_document_memory(client: genai.Client, file_data: dict, uploaded_files) -> List[PageContent]:
    """
    Processes a single non-PDF document from memory.
    Uses 'identifier' from file_data as the source_identifier.
    """
    source_identifier = file_data.get("identifier", "Unknown_Source") # <-- Use 'identifier' key
    file_type = file_data.get("type", "").lower()
    file_content = file_data.get("content")
    processed_pages: List[PageContent] = []

    logging.info(f"Starting page-level processing for non-PDF: '{source_identifier}' (type: {file_type})")

    if not file_content:
        logging.error(f"No content found for file '{source_identifier}'.")
        return [create_error_page(1, f"Source '{source_identifier}' has no content.")]
    if not file_type:
        logging.error(f"Cannot determine file type for '{source_identifier}'.")
        return [create_error_page(1, f"Unknown file type for '{source_identifier}'.")]

    try:
        # --- Dispatch based on file type, passing source_identifier ---
        if file_type in ["xlsx", "xls"]:
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "excel")
        elif file_type == "csv":
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "csv")
        elif file_type == "docx":
            processed_pages = await process_word_document(client, file_content, source_identifier)
        elif file_type == "pptx":
            processed_pages = await process_pptx_document(client, file_content, source_identifier)
        elif file_type in ["txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]:
            # process_text_document takes the whole file_data dict which includes 'identifier'
            processed_pages = await process_text_document(client, file_data)
        else:
            logging.warning(f"Unsupported file type '{file_type}' for '{source_identifier}'.")
            processed_pages = [create_error_page(1, f"Unsupported file type: {file_type}")]

        if not processed_pages:
            logging.warning(f"Processing for '{source_identifier}' resulted in no pages.")
            processed_pages = [create_error_page(1, f"No content pages generated during processing for '{source_identifier}'.")]

        logging.info(f"Finished page-level processing for non-PDF: '{source_identifier}'. Generated {len(processed_pages)} PageContent objects.")
        return processed_pages

    except Exception as e:
        logging.error(f"Core processing failed for document '{source_identifier}': {str(e)}", exc_info=True)
        return [create_error_page(1, f"Core processing error for '{source_identifier}': {e}")]
    
async def finalize_document(client: genai.Client, source_identifier: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Extracts subsections from raw_text for each page. <--- NEW STEP
    2. Merges cut-off subsections.
    3. Extracts chapters based on subsections.
    4. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    5. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    6. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{source_identifier}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": source_identifier, # Store the identifier here
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }
    # Attempt to dump incoming pages immediately for error reporting, handle failures
    initial_pages_dump = []
    try:
        # Dump safely, handle potential model errors during dump
        for i, p in enumerate(pages):
             if isinstance(p, PageContent):
                 try: initial_pages_dump.append(p.model_dump(mode='json'))
                 except Exception as dump_err: initial_pages_dump.append({"page_number": getattr(p, 'page_number', i+1), "error": f"Initial model dump failed: {dump_err}"})
             else: initial_pages_dump.append(p) # Keep as is if not PageContent obj

        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed during initial page dump for '{source_identifier}': {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{source_identifier}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result # Return early with initial page dump

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{source_identifier}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result # Return early

        # --- Step 1: Extract Subsections from Raw Text (NEW STEP) ---
        logging.debug(f"[{function_name}] Step 1: Extracting subsections from raw text for {len(pages)} pages of '{source_identifier}'.")
        subsection_extraction_tasks = []
        pages_needing_subsections = [] # Keep track of pages we are processing

        # Create tasks FIRST to run them concurrently
        for page in pages:
             # Only process if raw_text exists and subsections are currently empty
             if isinstance(page, PageContent) and page.raw_text and not page.subsections:
                 task = extract_subsections_from_text(client, page.raw_text, page.page_number, source_identifier)
                 subsection_extraction_tasks.append(task)
                 pages_needing_subsections.append(page) # Store reference to the page object
             elif not isinstance(page, PageContent):
                  logging.warning(f"[{function_name}] Item in pages list is not PageContent object for '{source_identifier}', page {getattr(page, 'page_number', '?')}. Skipping subsection extraction for it.")

             # else: subsections already exist or no raw text, skip

        # Run tasks concurrently
        if subsection_extraction_tasks:
             logging.info(f"[{function_name}] Running {len(subsection_extraction_tasks)} subsection extraction tasks for '{source_identifier}'.")
             results = await asyncio.gather(*subsection_extraction_tasks, return_exceptions=True)
             logging.info(f"[{function_name}] Finished subsection extraction tasks for '{source_identifier}'.")

             # Assign results back to the corresponding page objects
             if len(results) == len(pages_needing_subsections):
                 for i, result in enumerate(results):
                     page_to_update = pages_needing_subsections[i] # Get the correct page object
                     if isinstance(result, Exception):
                         logging.error(f"[{function_name}] Subsection extraction failed for page {page_to_update.page_number} of '{source_identifier}': {result}")
                         page_to_update.subsections = [] # Ensure empty list on error
                     elif isinstance(result, list): # Expected result is List[Subsection]
                         page_to_update.subsections = result # Assign generated subsections
                     else:
                         logging.error(f"[{function_name}] Unexpected result type {type(result)} for subsection extraction on page {page_to_update.page_number} of '{source_identifier}'.")
                         page_to_update.subsections = []
             else:
                  # This indicates a logic error in tracking tasks/pages
                  logging.error(f"[{function_name}] Mismatch between pages needing processing ({len(pages_needing_subsections)}) and task results ({len(results)}) for '{doc_name}'. Subsections may be missing.")
                  # Attempt to assign based on index anyway, but it might be wrong
                  for idx, page in enumerate(pages_needing_subsections):
                     if idx < len(results):
                         result = results[idx]
                         if isinstance(result, list): page.subsections = result
                         else: page.subsections = []
                     else: page.subsections = [] # No result for this page

        else:
             logging.info(f"[{function_name}] No pages required subsection extraction for '{source_identifier}'.")


        # --- Step 2: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Step 2: Merging subsections for '{source_identifier}'.")
        # Pass the 'pages' list which now has potentially populated subsections
        merged_pages = await merge_cutoff_subsections(pages)
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{source_identifier}'. Using original pages.")
             merged_pages = pages # Fallback to original pages (already updated with subsections)


        # --- Step 3: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Step 3: Extracting chapters for '{source_identifier}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{source_identifier}'. Cannot proceed further.")
             final_pages_as_dicts = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in merged_pages] # Dump state after merge/subsection attempts
             final_result["raw_extracted_content"]["pages"] = final_pages_as_dicts
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = {
                  "title": f"Processing Incomplete: {source_identifier}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result


        # --- Step 4: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Step 4: Analyzing intra-chapter relationships for {len(chapters)} chapters in '{source_identifier}'.")
        # This loop modifies chapter.qdrant_nodes in place
        intra_chapter_tasks = [analyze_concept_entity_relationships(client, chapter) for chapter in chapters]
        intra_results = await asyncio.gather(*intra_chapter_tasks, return_exceptions=True)
        # Log any exceptions from intra-chapter analysis
        for i, res in enumerate(intra_results):
            if isinstance(res, Exception):
                 logging.error(f"[{function_name}] Intra-chapter analysis failed for chapter {chapters[i].chapter_id} of '{source_identifier}': {res}")


        # --- Step 5: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Step 5: Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)


        # --- Step 6: Assemble Final Result ---
        logging.debug(f"[{function_name}] Step 6: Assembling final result for '{source_identifier}'.")
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated (even if empty list) before dumping
        dumped_chapters = []
        for ch in chapters:
            if isinstance(ch, Chapter):
                try: dumped_chapters.append(ch.model_dump(mode='json'))
                except Exception as ch_dump_err:
                     logging.error(f"[{function_name}] Failed to dump chapter {ch.chapter_id}: {ch_dump_err}")
                     dumped_chapters.append({"chapter_id": ch.chapter_id, "title": ch.title, "error": "Failed to serialize chapter"})
            else: logging.warning(f"[{function_name}] Invalid chapter object found during final assembly.")
        summary_dict['chapters'] = dumped_chapters


        # Dump the final state of pages (use merged_pages)
        final_pages_as_dicts = []
        for p in merged_pages:
             if isinstance(p, PageContent):
                 try: final_pages_as_dicts.append(p.model_dump(mode='json'))
                 except Exception as pg_dump_err:
                      logging.error(f"[{function_name}] Failed to dump page {p.page_number}: {pg_dump_err}")
                      final_pages_as_dicts.append({"page_number": p.page_number, "error": "Failed to serialize page"})
             else: final_pages_as_dicts.append(p)


        final_result = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result

    # --- Error Handling ---
    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{source_identifier}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {source_identifier}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": []
                # Removed entity_relationships placeholder here too
            }
        # Keep the initial pages dump in the error case if available
        return final_result



#-------------------------------------
# HELPER: Create Markdown Table
#-------------------------------------
def create_markdown_table(headers, data):
    """Create a markdown table string from headers and list of lists/dicts."""
    if not data or not headers: return "<!-- Empty Table -->"
    header_row = "| " + " | ".join(map(str, headers)) + " |"
    separator_row = "| " + " | ".join("---" * len(headers)) + " |"
    data_rows = []
    for row in data:
        if isinstance(row, dict):
            values = [str(row.get(h, '')) for h in headers]
        elif isinstance(row, (list, tuple)):
            # Ensure row has same length as headers, padding if necessary
            values = [str(row[i]) if i < len(row) else '' for i in range(len(headers))]
        else: # Handle unexpected row types
             values = [str(row)] + [''] * (len(headers) - 1)
        data_rows.append("| " + " | ".join(values) + " |")
    return "\n".join([header_row, separator_row] + data_rows)

#-------------------------------------
# REWRITTEN NON-PDF PROCESSORS
#-------------------------------------

async def process_tabular_data(client: genai.Client, file_content: bytes, source_identifier: str, filetype: str) -> List[PageContent]:
    """Processes Excel/CSV by converting sheets/file to Markdown text and extracting structure."""
    processed_pages = []
    page_counter = 1 # For sheets or single CSV page

    try:
        if filetype == "excel":
            # Use pandas to read all sheets
            try:
                 # Read all sheets into a dictionary of DataFrames
                 excel_data = pd.read_excel(io.BytesIO(file_content), sheet_name=None, engine='openpyxl') # Specify engine
                 if not excel_data:
                     raise ValueError("Excel file is empty or unreadable.")

                 for sheet_name, df in excel_data.items():
                     logging.info(f"Processing Excel sheet '{sheet_name}' for '{source_identifier}' as page {page_counter}")
                     if df.empty:
                          logging.warning(f"Sheet '{sheet_name}' in '{source_identifier}' is empty. Skipping.")
                          # Optionally create an empty page object
                          # processed_pages.append(PageContent(page_number=page_counter, ...))
                          page_counter += 1
                          continue
                     # Convert DataFrame to Markdown text
                     markdown_text = f"# Sheet: {sheet_name}\n\n" + df.to_markdown(index=False)
                     # Extract structure from this markdown text
                     page_content = await extract_structure_from_text_content(client, markdown_text, page_counter, f"{source_identifier} (Sheet: {sheet_name})")
                     processed_pages.append(page_content)
                     page_counter += 1

            except Exception as e:
                logging.error(f"Error reading Excel file '{source_identifier}' with pandas: {e}")
                processed_pages.append(create_error_page(1, f"Error reading Excel file '{source_identifier}': {e}")) # Treat as single page error

        elif filetype == "csv":
            logging.info(f"Processing CSV file '{source_identifier}' as page 1")
            try:
                # Use pandas to read CSV
                df = pd.read_csv(io.BytesIO(file_content))
                if df.empty:
                     raise ValueError("CSV file is empty or unreadable.")
                # Convert DataFrame to Markdown text
                markdown_text = "# CSV Data\n\n" + df.to_markdown(index=False)
                # Extract structure from this markdown text
                page_content = await extract_structure_from_text_content(client, markdown_text, 1, source_identifier)
                processed_pages.append(page_content)
            except Exception as e:
                 logging.error(f"Error reading CSV file '{source_identifier}' with pandas: {e}")
                 processed_pages.append(create_error_page(1, f"Error reading CSV file '{source_identifier}': {e}"))

        else: # Should not happen based on calling logic
            raise ValueError(f"Unsupported tabular file type: {filetype}")

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing tabular data for '{source_identifier}': {e}")
        # Return a single PageContent object indicating the error for the whole file
        return [create_error_page(1, f"Failed to process tabular file '{source_identifier}': {e}")]


async def process_word_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes Word docx file by extracting text and applying structure extraction."""
    processed_pages = []
    try:
        logging.info(f"Processing Word document '{source_identifier}'")
        doc = docx.Document(io.BytesIO(file_content))
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        # Option 1: Treat whole doc as one page (Simpler)
        if not full_text.strip():
             logging.warning(f"Word document '{source_identifier}' contains no text.")
             return [create_error_page(1, "Word document is empty or contains no text.")] # Return error page

        logging.debug(f"Extracted {len(full_text)} chars from '{source_identifier}'. Processing as single page.")
        page_content = await extract_structure_from_text_content(client, full_text, 1, source_identifier)
        processed_pages.append(page_content)

        # Option 2: Heuristic page splitting (More complex, may be less reliable than letting Gemini handle structure)
        # If needed, implement splitting logic here (e.g., by H1 headings) and call extract_structure_from_text_content multiple times.
        # Example (Conceptual):
        # pages_text = split_word_text_into_pages(doc) # Your custom splitting logic
        # for i, text_chunk in enumerate(pages_text, 1):
        #     page_content = await extract_structure_from_text_content(client, text_chunk, i, filename)
        #     processed_pages.append(page_content)

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing Word document '{source_identifier}': {e}")
        return [create_error_page(1, f"Failed to process Word document '{source_identifier}': {e}")]


async def process_pptx_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes PowerPoint pptx file by extracting text per slide and applying structure extraction."""
    processed_pages = []
    try:
        logging.info(f"Processing PowerPoint document '{source_identifier}'")
        ppt = Presentation(io.BytesIO(file_content))

        if not ppt.slides:
             logging.warning(f"PowerPoint document '{source_identifier}' contains no slides.")
             return [create_error_page(1, "PowerPoint document has no slides.")]

        for i, slide in enumerate(ppt.slides, 1):
            slide_texts = []
            # Extract title
            title = f"Slide {i}" # Default title
            has_title_shape = False
            try: # Add try-except for potential shape access errors
                for shape in slide.shapes:
                     if hasattr(shape, "is_title") and shape.is_title and shape.has_text_frame and shape.text.strip():
                         title = shape.text.strip()
                         slide_texts.append(f"# {title}") # Add title as H1 markdown
                         has_title_shape = True
                         break # Found title
            except Exception as shape_error:
                 logging.warning(f"Error accessing title shape on slide {i} of '{source_identifier}': {shape_error}")

            # Extract other text
            try:
                for shape in slide.shapes:
                     if shape.has_text_frame and shape.text and not (has_title_shape and hasattr(shape, "is_title") and shape.is_title):
                         # Add bullet points for structure if possible
                         for para in shape.text_frame.paragraphs:
                              if para.text.strip():
                                   prefix = "- " * (para.level + 1) # Indent based on level
                                   slide_texts.append(f"{prefix}{para.text.strip()}")
            except Exception as shape_error:
                 logging.warning(f"Error accessing text shapes on slide {i} of '{source_identifier}': {shape_error}")


            # Extract notes
            try:
                 if slide.has_notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip():
                     slide_texts.append("\n## Speaker Notes")
                     slide_texts.append(slide.notes_slide.notes_text_frame.text.strip())
            except Exception as notes_error:
                logging.warning(f"Error accessing notes on slide {i} of '{source_identifier}': {notes_error}")


            slide_full_text = "\n".join(slide_texts)

            if not slide_full_text.strip():
                 logging.warning(f"Slide {i} of '{source_identifier}' contains no text content. Skipping.")
                 # Optionally create empty page
                 # processed_pages.append(create_error_page(i, f"Slide {i} is empty."))
                 continue

            logging.debug(f"Processing slide {i} of '{source_identifier}'")
            page_content = await extract_structure_from_text_content(client, slide_full_text, i, f"{source_identifier} (Slide {i})")
            processed_pages.append(page_content)

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing PowerPoint document '{source_identifier}': {e}")
        return [create_error_page(1, f"Failed to process PowerPoint document '{source_identifier}': {e}")]

# --- Simplified process_text_document and helpers ---

async def process_text_document(client: genai.Client, file_data: dict) -> List[PageContent]:
    """Processes plain text, markdown, json, html, xml files by chunking and extracting structure."""
    source_identifier = file_data["name"]
    file_type = file_data.get("type", "txt")
    processed_pages = []
    page_counter = 1

    try:
        logging.info(f"Processing {file_type} document '{source_identifier}'")
        # Decode content safely
        try:
            text_content = file_data["content"].decode('utf-8')
        except UnicodeDecodeError:
             try:
                 text_content = file_data["content"].decode('latin-1') # Try alternative
             except Exception as decode_err:
                 logging.error(f"Failed to decode file '{source_identifier}': {decode_err}")
                 return [create_error_page(1, f"Failed to decode file '{source_identifier}'")]

        if not text_content.strip():
            logging.warning(f"Text document '{source_identifier}' is empty.")
            return [create_error_page(1, f"Text document '{source_identifier}' is empty.")]

        # --- Chunking Strategy (Example: Split by ~3000 words, respecting paragraphs) ---
        # More sophisticated chunking (e.g., LangChain text splitters) could be used here.
        # This is a simple heuristic.
        words = text_content.split()
        target_chunk_size_words = 3000 # Adjust as needed
        text_chunks = []
        current_chunk_words = []

        for word in words:
            current_chunk_words.append(word)
            # Check if chunk size is reached (approximately)
            if len(current_chunk_words) >= target_chunk_size_words:
                 # Find last paragraph break near the end to avoid splitting mid-sentence
                 chunk_text = " ".join(current_chunk_words)
                 last_para_break = chunk_text.rfind("\n\n")
                 if last_para_break > len(chunk_text) * 0.7 and last_para_break > 0: # Found a recent para break
                      final_chunk = chunk_text[:last_para_break].strip()
                      remainder_words = chunk_text[last_para_break:].strip().split()
                 else: # No good break point, just split words
                      final_chunk = chunk_text
                      remainder_words = []

                 if final_chunk: text_chunks.append(final_chunk)
                 current_chunk_words = remainder_words # Start next chunk with remainder
            # Add a newline check - handle cases where paragraphs themselves exceed the limit
            elif word.endswith('\n\n') and len(current_chunk_words) > target_chunk_size_words * 0.8:
                 chunk_text = " ".join(current_chunk_words)
                 text_chunks.append(chunk_text.strip())
                 current_chunk_words = []


        # Add the last remaining chunk
        if current_chunk_words:
            text_chunks.append(" ".join(current_chunk_words).strip())

        if not text_chunks: # Handle case where splitting resulted in nothing
             text_chunks = [text_content]

        logging.info(f"Split '{source_identifier}' into {len(text_chunks)} chunks for processing.")

        # --- Process each chunk ---
        for text_chunk in text_chunks:
            if not text_chunk.strip(): continue # Skip empty chunks
            page_content = await extract_structure_from_text_content(client, text_chunk, page_counter, f"{source_identifier} (Chunk {page_counter})")
            processed_pages.append(page_content)
            page_counter += 1

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing text document '{source_identifier}': {e}")
        return [create_error_page(1, f"Failed to process text document '{source_identifier}': {e}")]

async def process_all_documents_async(file_data):
    """Process all documents with proper model conversion and error handling."""
    # --- Setup ---
    try:
        api_key = get_gemini_api_key()
        # Example basic client init (use your actual setup)
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                }
            })

        logging.info("Gemini client initialized.")
    except Exception as client_err:
        logging.error(f"Failed to initialize Gemini client: {client_err}")
        st.error(f"Failed to initialize Gemini client: {client_err}")
        return []


    final_docs = [] # Stores the FINAL dictionary results from finalize_document
    uploaded_files = [] # Assuming this might be needed by processors

    # Intermediate storage for pages BEFORE finalization
    # Maps: doc_name -> List[PageContent objects]
    all_doc_pages: Dict[str, List[PageContent]] = defaultdict(list)
    doc_order: List[str] = [] # Maintain original processing order

    # --- Status Tracking Setup ---
    status_container = st.session_state.get("status_container")
    progress_bar = st.session_state.get("progress_bar")
    time_info = st.session_state.get("time_info")
    processing_status = {
        "active": True, "current_step": "",
        "total_steps": 5, # Detect, PDFPages, NonPDFPages, Finalize, Cleanup
        "current_step_num": 0, "step_progress": 0.0,
        "start_time": time.time(), "step_start_time": time.time(),
        "estimated_time_remaining": None, "parallel_count": 1
    }
    steps = { # Renamed and renumbered steps
        "Detecting document types": 1,
        "Processing PDF Pages": 2,
        "Processing Non-PDF Pages": 3,
        "Finalizing Documents": 4,
        "Cleaning up resources": 5
    }

    def update_step_status(step, current=None, total=None, parallel=1):
        """Updates the processing status and UI elements."""
        try:
            processing_status["current_step"] = step
            processing_status["current_step_num"] = steps.get(step, processing_status["current_step_num"])
            processing_status["step_start_time"] = time.time() # Reset step timer

            step_progress = 0.0
            if current is not None and total is not None and total > 0:
                step_progress = min(1.0, float(current) / float(total))
            processing_status["step_progress"] = step_progress

            # Ensure current_step_num is at least 1 for progress calculation if processing started
            step_num_for_calc = max(1, processing_status["current_step_num"])
            overall_progress = (step_num_for_calc - 1 + processing_status["step_progress"]) / processing_status["total_steps"]
            overall_progress = max(0.0, min(1.0, overall_progress))

            status_msg = f"{processing_status['current_step_num']}/{processing_status['total_steps']}: {step}"
            if current is not None and total is not None:
                status_msg += f" ({current}/{total})"
            if parallel > 1:
                status_msg += f" | Parallelism: {parallel}"

            eta_msg = ""
            time_elapsed = time.time() - processing_status["start_time"]
            if overall_progress > 0.01: # Only estimate if some progress made
                est_total_time = time_elapsed / overall_progress
                remaining = max(0, est_total_time - time_elapsed)
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                mins_remaining, secs_remaining = divmod(int(remaining), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: ~{mins_remaining}m {secs_remaining}s"
            else:
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: Calculating..."


            if status_container: status_container.update(label=status_msg, expanded=True)
            if progress_bar: progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
            if time_info: time_info.markdown(f"`{eta_msg}`")

        except Exception as ui_err:
             logging.warning(f"Failed to update UI status: {ui_err}") # Log UI errors but don't stop processing

    # --- Main Processing Logic ---
    try:
        # --- Step 1: Detect document types ---
        update_step_status("Detecting document types", 0, len(file_data))
        pdf_docs_info = [] # List of file dicts for PDFs
        other_docs_info = [] # List of file dicts for Non-PDFs
        for i, file in enumerate(file_data):
            if not isinstance(file, dict):
                st.warning(f"Skipping invalid file data item (not a dict): {type(file)}")
                continue
            file_name = file.get("name", f"Unnamed_File_{i}")
            file["name"] = file_name
            if file_name not in doc_order: doc_order.append(file_name)

            file_ext = file.get("type", "").lower() or (file_name.split('.')[-1].lower() if '.' in file_name else "")
            file["type"] = file_ext

            if file_ext == 'pdf':
                pdf_docs_info.append(file)
            else:
                # Basic check for content needed by non-PDF processors
                if "content" not in file or not file["content"]:
                     st.warning(f"Skipping file '{file_name}' due to missing or empty 'content'.")
                else:
                     other_docs_info.append(file)

            update_step_status("Detecting document types", i + 1, len(file_data))
        update_step_status("Detecting document types", len(file_data), len(file_data))


        # --- Step 2: Process PDF Pages -> List[PageContent] ---
        if pdf_docs_info:
            pdf_processing_start_time = time.time()
            update_step_status("Processing PDF Pages", 0, len(pdf_docs_info))

            # --- 2a. Convert PDFs to intermediate page info (images) ---
            MAX_PARALLEL_CONVERSIONS = 5 # Limit concurrent CPU/IO-bound rendering
            pdf_render_semaphore = asyncio.Semaphore(MAX_PARALLEL_CONVERSIONS)
            all_pdf_pages_info = [] # List of dicts: {"doc_name": ..., "page_num": ..., "image_b64": ...}
            conversion_tasks = []

            for file in pdf_docs_info:
                if "content" in file and "name" in file:
                    conversion_tasks.append(extract_pages_from_pdf_bytes(file["content"], file["name"]))
                else:
                    st.warning(f"Cannot convert PDF '{file.get('name', 'Unknown')}' due to missing content or name.")

            async def run_with_render_semaphore(task):
                 async with pdf_render_semaphore: return await task
            tasks_with_sem = [run_with_render_semaphore(task) for task in conversion_tasks]

            if tasks_with_sem:
                conversion_results = await asyncio.gather(*tasks_with_sem, return_exceptions=True)
                for i, result in enumerate(conversion_results):
                     if isinstance(result, list): all_pdf_pages_info.extend(result)
                     elif isinstance(result, Exception): st.error(f"PDF page image extraction failed for '{pdf_docs_info[i].get('name', 'Unknown')}': {result}")
                     else: st.warning(f"Unexpected result type during PDF extraction: {type(result)}")
                logging.info(f"Finished PDF to image conversion for {len(pdf_docs_info)} files in {time.time() - pdf_processing_start_time:.2f}s")
            else:
                 st.info("No valid PDF files to convert to images.")
            # Update status after conversion sub-step
            update_step_status("Processing PDF Pages", 0, len(all_pdf_pages_info), 1) # Reset progress for analysis


            # --- 2b. Analyze Page Images -> PageContent Objects ---
            pdf_page_count = len(all_pdf_pages_info)
            if pdf_page_count > 0:
                logging.info(f"Analyzing {pdf_page_count} PDF pages using Gemini...")
                MAX_PARALLEL_PAGES = 25 # Gemini API concurrency limit
                page_api_semaphore = asyncio.Semaphore(MAX_PARALLEL_PAGES)
                page_tasks = []
                valid_page_info_count = 0

                for page_info in all_pdf_pages_info:
                    if isinstance(page_info, dict) and all(k in page_info for k in ["doc_name", "page_num", "image_b64"]):
                        # Ensure image_b64 is not empty
                        if page_info["image_b64"]:
                            page_tasks.append(asyncio.create_task(process_single_page_with_semaphore(page_api_semaphore, client, page_info, uploaded_files)))
                            valid_page_info_count += 1
                        else:
                             st.warning(f"Skipping page {page_info.get('page_num')} of '{page_info.get('doc_name')}' due to empty image data.")
                    else:
                         st.warning(f"Skipping invalid PDF page_info data: {str(page_info)[:100]}...")

                pdf_page_count = valid_page_info_count # Update actual count

                # --- Collect PageContent results and group by doc_name ---
                processed_pdf_results_collector: Dict[str, List[PageContent]] = defaultdict(list)
                processed_pdf_count = 0
                if page_tasks:
                    for future in asyncio.as_completed(page_tasks):
                        try:
                             task_output = await future # {"info":..., "result": PageContent_obj}
                             original_info = task_output.get("info", {})
                             result_page_obj = task_output.get("result")
                             doc_name = original_info.get("doc_name", "Unknown_PDF_Context")

                             if isinstance(result_page_obj, PageContent):
                                 processed_pdf_results_collector[doc_name].append(result_page_obj)
                             else: # Handle errors / None results
                                 error_page = create_error_page(original_info.get("page_num", -1), f"Invalid result type {type(result_page_obj)} from task for {doc_name}")
                                 processed_pdf_results_collector[doc_name].append(error_page)
                        except Exception as e:
                             logging.error(f"Error awaiting PDF page task result: {e}", exc_info=True)
                             # Consider adding a placeholder error page to a default error list if needed

                        processed_pdf_count += 1
                        # Calculate active tasks more accurately
                        active_tasks = len(page_tasks) - processed_pdf_count # Approximation
                        update_step_status("Processing PDF Pages", processed_pdf_count, pdf_page_count, max(1, min(MAX_PARALLEL_PAGES, active_tasks)))

                    # --- Add collected pages to the main dictionary ---
                    for doc_name, pages_list in processed_pdf_results_collector.items():
                        if doc_name != "Unknown_PDF_Context":
                            sorted_pages = sorted(pages_list, key=lambda p: getattr(p, 'page_number', float('inf')))
                            all_doc_pages[doc_name].extend(sorted_pages)
                        else:
                            st.warning("Some PDF pages processed without known document context.")
                else:
                    st.warning("No valid PDF page processing tasks created.")
            else:
                 st.info("No PDF pages extracted to analyze.")
            # Final update for PDF processing step
            update_step_status("Processing PDF Pages", len(pdf_docs_info), len(pdf_docs_info)) # Show step as complete


        # --- Step 3: Process Non-PDF Pages -> List[PageContent] ---
        if other_docs_info:
            non_pdf_start_time = time.time()
            update_step_status("Processing Non-PDF Pages", 0, len(other_docs_info))
            MAX_PARALLEL_NON_PDF = 4 # Limit concurrency for potentially CPU/IO bound local parsing
            semaphore_non_pdf = asyncio.Semaphore(MAX_PARALLEL_NON_PDF)
            non_pdf_tasks = []

            for file in other_docs_info:
                 # Wrapper expects file dict, calls process_single_document_memory which returns List[PageContent]
                 async def process_non_pdf_wrapper(f):
                      async with semaphore_non_pdf:
                           pages_list: List[PageContent] = await process_single_document_memory(client, f, uploaded_files)
                           return (f.get("name"), pages_list)
                 non_pdf_tasks.append(process_non_pdf_wrapper(file)) # Pass the file dict

            non_pdf_progress = 0
            total_non_pdf_to_process = len(non_pdf_tasks)

            if total_non_pdf_to_process > 0:
                logging.info(f"Processing {total_non_pdf_to_process} non-PDF documents...")
                gathered_results = await asyncio.gather(*non_pdf_tasks, return_exceptions=True)

                for result in gathered_results:
                    non_pdf_progress += 1
                    doc_name_from_result = "Unknown_NonPDF" # Default for logging
                    try:
                        if isinstance(result, Exception):
                            # Log the exception associated with the task if possible (hard without context)
                            logging.error(f"Non-PDF processing task failed critically: {result}")
                            st.error(f"A non-PDF processing task failed: {result}")
                        elif isinstance(result, tuple) and len(result) == 2:
                            doc_name, pages_list = result
                            doc_name_from_result = doc_name or "Unnamed NonPDF"
                            if doc_name and isinstance(pages_list, list) and all(isinstance(p, PageContent) for p in pages_list):
                                sorted_pages = sorted(pages_list, key=lambda p: getattr(p, 'page_number', float('inf')))
                                all_doc_pages[doc_name].extend(sorted_pages) # Add to main collection
                            else:
                                logging.warning(f"Received invalid result tuple from non-PDF wrapper: doc={doc_name}, pages type={type(pages_list)}. Adding error page.")
                                if doc_name: # Add error page if we know the doc name
                                     all_doc_pages[doc_name].append(create_error_page(1, f"Invalid data from non-PDF processing for {doc_name}"))
                                st.warning(f"Received invalid result tuple from non-PDF wrapper for {doc_name}")
                        else:
                            logging.warning(f"Received unexpected result structure from non-PDF processing: {type(result)}")
                            st.warning(f"Received unexpected result structure from non-PDF processing: {type(result)}")
                    except Exception as inner_err:
                         logging.error(f"Error processing result for non-PDF doc '{doc_name_from_result}': {inner_err}", exc_info=True)

                    # Update progress
                    active_tasks = total_non_pdf_to_process - non_pdf_progress
                    update_step_status("Processing Non-PDF Pages", non_pdf_progress, total_non_pdf_to_process, max(1, min(MAX_PARALLEL_NON_PDF, active_tasks)))

                logging.info(f"Finished non-PDF processing in {time.time() - non_pdf_start_time:.2f}s")
            else:
                st.info("No valid non-PDF documents found to process.")
            # Final update for this step
            update_step_status("Processing Non-PDF Pages", len(other_docs_info), len(other_docs_info))


        # --- Step 4: Finalize ALL Documents (Common Finalization) ---
        docs_to_finalize_count = len(all_doc_pages) # Number of unique documents with processed pages
        if docs_to_finalize_count > 0:
            finalize_start_time = time.time()
            update_step_status("Finalizing Documents", 0, docs_to_finalize_count)
            MAX_PARALLEL_FINALIZE = 5 # Limit concurrency for finalization (often involves LLM calls)
            semaphore_finalize = asyncio.Semaphore(MAX_PARALLEL_FINALIZE)
            finalize_tasks = []

            # Use the collected doc names for iteration
            doc_names_to_process = [name for name in doc_order if name in all_doc_pages]
            processed_names = set(doc_names_to_process)
            doc_names_to_process.extend([name for name in all_doc_pages if name not in processed_names])

            logging.info(f"Starting finalization for {docs_to_finalize_count} documents...")

            for doc_name in doc_names_to_process:
                 pages_list = all_doc_pages[doc_name] # This is List[PageContent]
                 if not pages_list:
                     logging.warning(f"Skipping finalization for '{doc_name}' as it has no processed pages.")
                     # Add an error entry to final_docs for this file
                     final_docs.append({
                         "raw_extracted_content": {
                             "filename": doc_name, "pages": [], "summary": None,
                             "error": "No pages were successfully processed for this document."
                         }
                     })
                     continue

                 # Wrapper for finalize_document
                 async def finalize_wrapper(name, pages):
                      async with semaphore_finalize:
                           logging.debug(f"Starting finalize_document task for '{name}'.")
                           result_dict = await finalize_document(client, name, pages)
                           logging.debug(f"Finished finalize_document task for '{name}'.")
                           return result_dict # finalize_document returns the final dictionary

                 finalize_tasks.append(finalize_wrapper(doc_name, pages_list))

            # Collect results from finalization tasks
            finalized_count = 0
            if finalize_tasks:
                 gathered_final_results = await asyncio.gather(*finalize_tasks, return_exceptions=True)
                 for result in gathered_final_results:
                      finalized_count += 1
                      try:
                          if isinstance(result, Exception):
                              # Log the exception, but try to find context if possible (difficult here)
                              logging.error(f"Document finalization task failed critically: {result}", exc_info=True)
                              st.error(f"A document finalization task failed: {result}")
                              # Add generic error doc (can we get the name?)
                              # If gather maintains order, we could map back to doc_names_to_process[index]
                              failed_doc_name = "Unknown_Finalize_Error" # Placeholder
                              final_docs.append({
                                   "raw_extracted_content": {"filename": failed_doc_name, "pages": [], "summary": None, "error": f"Finalization failed: {result}"}
                              })
                          elif isinstance(result, dict) and "raw_extracted_content" in result:
                              final_docs.append(result) # Add the final dict structure
                          else:
                              logging.warning(f"Received unexpected result type from finalize_document: {type(result)}")
                              final_docs.append({
                                 "raw_extracted_content": {"filename": "Unknown_Finalize_Error", "pages": [], "summary": None, "error": f"Unexpected finalize result type {type(result)}"}
                              })
                      except Exception as final_proc_err:
                           logging.error(f"Error processing finalization result: {final_proc_err}", exc_info=True)

                      active_tasks = len(finalize_tasks) - finalized_count
                      update_step_status("Finalizing Documents", finalized_count, docs_to_finalize_count, max(1, min(MAX_PARALLEL_FINALIZE, active_tasks)))
            else:
                 st.info("No documents needed finalization.")

            logging.info(f"Finished document finalization in {time.time() - finalize_start_time:.2f}s")
            # Final update for this step
            update_step_status("Finalizing Documents", docs_to_finalize_count, docs_to_finalize_count)
        else:
             st.info("No documents processed, skipping finalization.")


        # --- Step 5: Cleaning up resources ---
        update_step_status("Cleaning up resources", 1, 1)
        logging.info("Cleaning up resources...")
        # ... (cleanup logic if any, e.g., close files, clear temp data) ...
        processing_status["active"] = False
        if status_container:
             try: status_container.update(label="✅ Document processing complete!", state="complete", expanded=False)
             except Exception as ui_err: logging.warning(f"Failed to update final status: {ui_err}")

    # --- Exception Handling & Finally Block ---
    except asyncio.CancelledError:
         logging.warning("Document processing was cancelled.")
         st.warning("Document processing was cancelled.")
         if status_container:
             try: status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
             except Exception: pass
         processing_status["active"] = False
         return [] # Return empty list on cancellation
    except Exception as e:
        logging.error(f"❌ An error occurred during the main document processing pipeline: {str(e)}", exc_info=True)
        st.error(f"❌ An error occurred during document processing: {str(e)}")
        # st.error(traceback.format_exc()) # Optionally show full traceback in UI
        if status_container:
            try: status_container.update(label=f"❌ Pipeline Error: {str(e)}", state="error", expanded=True)
            except Exception: pass
        processing_status["active"] = False
        # Return partially processed docs (which are the final dictionaries)
        logging.info(f"Returning {len(final_docs)} potentially partially finalized documents due to pipeline error.")
        return final_docs
    finally:
        logging.info("Document processing pipeline finished or exited.")
        processing_status["active"] = False
        if "st" in locals() and hasattr(st, "session_state"):
            st.session_state.processing_active = False # Example if using session state

    # Return the list of final document dictionaries
    logging.info(f"Successfully processed and finalized {len(final_docs)} documents.")
    return final_docs

###############################
# UI DISPLAY FUNCTIONS (Updated)
###############################

def display_table(table):
    """Display table data (accepts model or dict)."""
    if not table: return
    table_data = table.model_dump() if hasattr(table, 'model_dump') else table

    title = table_data.get("title")
    content = table_data.get("table_content")

    if title:
        st.subheader(f"Table: {title}")

    if content:
        # Attempt to display as DataFrame, fallback to Markdown
        try:
            if pd: # Check if pandas is available
                # Preprocess markdown for better parsing: remove leading/trailing whitespace, split lines
                lines = [line.strip() for line in content.strip().split('\n')]
                # Remove separator line if present
                if len(lines) > 1 and all(c in '-:| ' for c in lines[1]):
                    lines.pop(1)
                # Reconstruct cleaned content for StringIO
                cleaned_content = "\n".join(lines)

                df = pd.read_csv(
                    io.StringIO(cleaned_content),
                    sep="|",
                    skipinitialspace=True,
                    index_col=False # Important: Avoid assuming first column is index
                ).iloc[1:] # Skip the header row parsed as data
                # Clean up columns: remove leading/trailing spaces in names, drop empty columns
                df.columns = [col.strip() for col in df.columns]
                df = df.drop(columns=[col for col in df.columns if col == ''], errors='ignore') # Drop unnamed columns from extra separators
                df = df.dropna(axis=1, how='all')
                st.dataframe(df.reset_index(drop=True)) # Show clean index
            else:
                st.markdown(f"```markdown\n{content}\n```")
        except Exception as e:
            logging.warning(f"Pandas failed to parse table, showing markdown. Error: {e}. Content:\n{content}")
            st.markdown(f"```markdown\n{content}\n```")
    else:
        st.caption("Table content not available.")

def display_visual_element(visual):
    """Display visual element data (accepts model or dict)."""
    if not visual: return
    visual_data = visual.model_dump() if hasattr(visual, 'model_dump') else visual

    visual_type = visual_data.get('type', 'Unknown Type')
    description = visual_data.get('description', 'No description available.')
    data_summary = visual_data.get('data_summary') # Can be None or "N/A"

    st.markdown(f"**🖼️ {visual_type.capitalize()}**")
    st.markdown(f"*Description:* {description}")
    if data_summary and data_summary != "N/A":
        st.markdown(f"*Data Summary:* {data_summary}")
    # Consider adding visual_id if useful: st.caption(f"ID: {visual_data.get('visual_id')}")

def display_numerical_data_point(num_data):
    """Displays a single numerical data point."""
    if not num_data: return
    num_dict = num_data.model_dump() if hasattr(num_data, 'model_dump') else num_data

    value = num_dict.get('value', 'N/A')
    description = num_dict.get('description', 'No description.')
    context = num_dict.get('context')

    st.markdown(f"**{value}**: {description}")
    if context:
        st.caption(f"Context: {context}")

def display_qdrant_node_network(qdrant_nodes_data: Optional[List[Dict]], title="Concept/Entity Network"): # Accept List[Dict]
    """Displays relationships from QdrantNode data (as dicts) as a network graph."""
    if not qdrant_nodes_data: # Handles None or empty list
        st.info("No relationship data available to display.")
        return

    # Ensure input is a list of dictionaries
    if not isinstance(qdrant_nodes_data, list) or not all(isinstance(n, dict) for n in qdrant_nodes_data):
         st.error("Invalid data format passed to node network display function (expected list of dicts).")
         logging.error(f"display_qdrant_node_network received invalid data type: {type(qdrant_nodes_data)}")
         return

    if not nx or not plt:
        # ... (Fallback text display as before, using .get() on dicts) ...
        st.warning("NetworkX/Matplotlib not installed. Displaying text list.")
        st.markdown("#### Relationships (Text List):")
        for node_dict in qdrant_nodes_data:
             linked_nodes = node_dict.get('linked_nodes', [])
             if linked_nodes:
                 for link_dict in linked_nodes:
                     st.markdown(f"- **{node_dict.get('name','?')}** → **{link_dict.get('target_node_id','?')}**: {link_dict.get('relationship_description','?')} (Strength: {link_dict.get('relationship_strength','?')})")
        return

    # Create a network graph
    G = nx.DiGraph()
    node_labels = {}
    edge_weights = {}
    edge_labels = {}
    node_types_map = {} # Store node type for coloring

    # Add nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        node_id = node_dict.get('node_id')
        node_name = node_dict.get('name', 'Unknown Node')
        node_type = node_dict.get('node_type', 'unknown')
        if node_id and node_id not in G: # Check if node_id exists and is not None
            G.add_node(node_id)
            node_labels[node_id] = f"{node_name}\n({node_type})"
            node_types_map[node_id] = node_type # Store type

    # Add edges from linked_nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        source_id = node_dict.get('node_id')
        linked_nodes = node_dict.get('linked_nodes', [])
        if not source_id or not G.has_node(source_id) or not isinstance(linked_nodes, list):
             continue # Skip if source_id invalid or no links

        for link_dict in linked_nodes:
             if not isinstance(link_dict, dict): continue # Skip invalid links
             target_id = link_dict.get('target_node_id')

             if not target_id: continue # Skip links without target

             # Ensure target node exists in the graph
             if target_id not in G:
                 G.add_node(target_id)
                 # Try to create a basic label from ID if target details aren't in current node list
                 target_name_match = re.match(r"entity_(.*?)_ch\d+|concept_(.*?)_ch\d+", target_id)
                 target_name_part = target_name_match.group(1) or target_name_match.group(2) if target_name_match else target_id
                 node_labels[target_id] = target_name_part.replace('_',' ').title() # Basic label
                 node_types_map[target_id] = 'unknown' # Mark as unknown type

             # Add edge if not already present
             if not G.has_edge(source_id, target_id):
                 strength = float(link_dict.get('relationship_strength', 5.0)) # Use get with default
                 description = link_dict.get('relationship_description', '')
                 G.add_edge(source_id, target_id, weight=strength)
                 edge_weights[(source_id, target_id)] = strength
                 edge_labels[(source_id, target_id)] = description

    # --- Graph drawing logic remains largely the same ---
    if not G.nodes():
         st.info("No nodes found to build the graph.")
         return
    if not G.edges():
         st.info("Nodes found, but no links between them to display.")
         st.write("Nodes identified:", list(node_labels.values()))
         return

    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(12, 10))
    try:
        pos = nx.spring_layout(G, k=0.5, iterations=50, seed=42)
    except Exception as layout_err:
        logging.error(f"NetworkX layout failed: {layout_err}. Using random layout.")
        pos = nx.random_layout(G, seed=42)

    # Draw nodes with color based on type
    node_colors = ['skyblue' if node_types_map.get(n) == 'entity' else 'lightgreen' if node_types_map.get(n) == 'concept' else 'lightgray' for n in G.nodes()]
    nx.draw_networkx_nodes(G, pos, node_size=1500, node_color=node_colors, alpha=0.9, ax=ax)

    # Draw edges
    if edge_weights: # Check if edges exist before calculating weights
        min_strength = min(edge_weights.values()) if edge_weights else 1
        max_strength = max(edge_weights.values()) if edge_weights else 10
        if max_strength == min_strength: max_strength += 1
        edge_width = [1 + 4 * (edge_weights.get(edge, 5.0) - min_strength) / (max_strength - min_strength) for edge in G.edges()]
        nx.draw_networkx_edges(G, pos, width=edge_width, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)
    else: # Draw default thin edges if no weights somehow
         nx.draw_networkx_edges(G, pos, width=1.0, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)


    # Draw labels
    nx.draw_networkx_labels(G, pos, labels=node_labels, font_size=9, font_weight='normal', ax=ax)

    ax.set_title(title, fontsize=16)
    fig.tight_layout()
    plt.axis('off')
    st.pyplot(fig)
    plt.close(fig)

    # Display edge details table
    edge_data = []
    # Use edge_labels which was populated correctly
    for (source_id, target_id), desc in edge_labels.items():
        edge_data.append({
            "Source": node_labels.get(source_id, source_id),
            "Target": node_labels.get(target_id, target_id),
            "Strength": edge_weights.get((source_id, target_id), 5.0),
            "Description": desc
        })

    if edge_data:
        st.markdown("### Relationship Details")
        try:
             if pd: # Check if pandas is available
                 df = pd.DataFrame(edge_data)
                 st.dataframe(df)
             else:
                 for item in edge_data: st.write(item) # Fallback display
        except Exception as df_err:
            logging.error(f"Failed to create DataFrame for edge data: {df_err}")
            for item in edge_data: st.write(item) # Fallback display
            
def display_chapter(chapter):
    """Display chapter details including subsections and Qdrant node relationships."""
    if not chapter: return
    chapter_data = chapter.model_dump() if hasattr(chapter, 'model_dump') else chapter

    title = chapter_data.get('title', 'Untitled Chapter')
    summary = chapter_data.get('summary', 'No summary available.')
    subsections = chapter_data.get('subsections', [])
    qdrant_nodes = chapter_data.get('qdrant_nodes') # This holds QdrantNode objects/dicts

    st.markdown(f"## {title}")
    st.markdown(f"_{summary}_")

    # Display Qdrant Node relationships if available
    if qdrant_nodes:
        with st.expander("View Concepts/Entities and Relationships within this Chapter", expanded=False):
            # Pass the actual list of nodes (could be dicts or models)
            display_qdrant_node_network(qdrant_nodes, title=f"Network for '{title}'")
    else:
        st.caption("No concept/entity relationship data available for this chapter.")

    # Display subsections
    if subsections:
        st.markdown("### Subsections Contained")
        for sub_idx, subsection_data in enumerate(subsections):
            if not subsection_data: continue
            # Handle dict or model
            sub_obj = subsection_data
            if isinstance(subsection_data, dict):
                 # Minimal validation or create Subsection obj if needed for consistency
                 # sub_obj = Subsection(**subsection_data) # Optional
                 pass

            sub_title = getattr(sub_obj, 'title', f'Subsection {sub_idx + 1}')
            sub_text = getattr(sub_obj, 'text', 'No text available.')
            sub_description = getattr(sub_obj, 'description', '')
            sub_id = getattr(sub_obj, 'subsection_id', 'N/A')
            page_num = "?" # Page number not stored on subsection anymore
            is_cutoff = getattr(sub_obj, 'is_cutoff', False)
            ref_visuals = getattr(sub_obj, 'referenced_visuals', [])
            ref_tables = getattr(sub_obj, 'referenced_tables', [])

            expander_title = f"{sub_idx + 1}. {sub_title}"
            if is_cutoff: expander_title += " (Continues...)"

            with st.expander(expander_title):
                if sub_description:
                    st.caption(f"Description: {sub_description}")
                st.markdown(sub_text) # Display the actual text
                st.caption(f"Subsection ID: {sub_id}")

                # Display referenced visuals/tables if any
                if ref_visuals: st.write(f"**Ref Visuals:** {', '.join(ref_visuals)}")
                if ref_tables: st.write(f"**Ref Tables:** {', '.join(ref_tables)}")
    else:
        st.caption("No subsections listed for this chapter.")

def display_node_links(all_nodes: List[Dict | Any], # List of node dicts or Pydantic objects
                       selected_node_id: str,
                       node_map: Optional[Dict[str, Dict | Any]] = None):
    """
    Finds a specific node by ID and displays its outgoing relationships (linked_nodes),
    including the origin context (document, chapter) of the target node.

    Args:
        all_nodes: A list containing all nodes (as dictionaries or Pydantic objects).
        selected_node_id: The node_id of the node to inspect.
        node_map: Optional pre-built dictionary mapping node_id to node object/dict
                  for faster lookups, especially for target node info.
    """
    function_name = "display_node_links"
    if not selected_node_id:
        st.caption("No node selected.")
        return

    # --- Find the selected node (Source Node) ---
    selected_node = None
    if node_map and selected_node_id in node_map:
        selected_node = node_map[selected_node_id]
    else:
        # Fallback search if map not provided
        for node in all_nodes:
            is_model = not isinstance(node, dict)
            node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
            if node_id == selected_node_id:
                selected_node = node
                break

    if selected_node is None:
        st.warning(f"Node with ID '{selected_node_id}' not found.")
        return

    is_source_model = not isinstance(selected_node, dict)
    source_name = getattr(selected_node, 'name', selected_node_id) if is_source_model else selected_node.get('name', selected_node_id)
    source_doc = getattr(selected_node, 'document_context', '?') if is_source_model else selected_node.get('document_context', '?')
    source_chapter = getattr(selected_node, 'chapter_id_context', '?') if is_source_model else selected_node.get('chapter_id_context', '?')

    st.markdown(f"#### Relationships FROM: **{source_name}**")
    st.caption(f"*(Source Origin: Document '{source_doc}', Chapter '{source_chapter}')*")
    st.markdown("---")


    linked_nodes = getattr(selected_node, 'linked_nodes', []) if is_source_model else selected_node.get('linked_nodes', [])

    if not linked_nodes or not isinstance(linked_nodes, list):
        st.caption("No outgoing links found for this node.")
        return

    # --- Build Node Map if needed (for target node lookups) ---
    if node_map is None:
        logging.debug(f"[{function_name}] Building temporary node_map for target info lookup.")
        node_map = {}
        for node in all_nodes:
             inner_is_model = not isinstance(node, dict)
             inner_node_id = getattr(node, 'node_id', None) if inner_is_model else node.get('node_id')
             if inner_node_id:
                 node_map[inner_node_id] = node

    # --- Display Links with Target Context ---
    if not node_map: # Check if map is still empty after trying to build
         st.warning("Could not build node map to look up target details.")

    for i, link in enumerate(linked_nodes):
        link_is_model = not isinstance(link, dict)
        target_id = getattr(link, 'target_node_id', 'Unknown Target ID') if link_is_model else link.get('target_node_id', 'Unknown Target ID')
        description = getattr(link, 'relationship_description', 'No description') if link_is_model else link.get('relationship_description', 'No description')
        strength = getattr(link, 'relationship_strength', 0.0) if link_is_model else link.get('relationship_strength', 0.0)
        keywords = getattr(link, 'relationship_keywords', []) if link_is_model else link.get('relationship_keywords', [])
        if not isinstance(keywords, list): keywords = []

        # --- Find Target Node Info using the map ---
        target_node = node_map.get(target_id) if node_map else None
        target_name = target_id # Default to ID if not found
        target_type = "unknown"
        target_doc = "unknown"
        target_chapter = "unknown"

        if target_node:
             target_is_model = not isinstance(target_node, dict)
             target_name = getattr(target_node, 'name', target_id) if target_is_model else target_node.get('name', target_id)
             target_type = getattr(target_node, 'node_type', 'unknown') if target_is_model else target_node.get('node_type', 'unknown')
             target_doc = getattr(target_node, 'document_context', 'unknown') if target_is_model else target_node.get('document_context', 'unknown')
             target_chapter = getattr(target_node, 'chapter_id_context', 'unknown') if target_is_model else target_node.get('chapter_id_context', 'unknown')
        elif target_id != 'Unknown Target ID':
             logging.warning(f"[{function_name}] Target node '{target_id}' linked from '{selected_node_id}' not found in node map.")


        # Display formatted link info including target context
        st.markdown(f"**{i+1}. Target Node:** {target_name} (`{target_type}`)")
        # Add target origin using italics or caption
        st.markdown(f"    *Origin: Document '{target_doc}', Chapter '{target_chapter}'*")
        st.markdown(f"    **Link Description:** {description}")
        st.markdown(f"    **Strength:** {strength:.1f}/10.0")
        if keywords:
            st.markdown(f"    **Keywords:** " + ", ".join(keywords))
        st.markdown(f"    *Target Node ID:* `{target_id}`") # Keep ID easily visible
        st.markdown("---") # Separator

    if not linked_nodes: # Redundant check, but safe
        st.caption("No valid outgoing links found to display.")



def display_project_ontology(ontology: Optional[ProjectOntology]): # Type hint remains Optional[ProjectOntology]
    """Displays project-wide ontology including the aggregated graph.
       Expects a ProjectOntology Pydantic model object as input.
    """
    function_name = "display_project_ontology"
    if not ontology:
        st.warning("No project ontology data available to display.")
        return

    # --- Treat input strictly as a ProjectOntology object ---
    # Remove the is_model check and the dictionary fallback logic
    if not hasattr(ontology, '__class__') or ontology.__class__.__name__ != 'ProjectOntology':
        st.error(f"[display_project_ontology] Error: Expected a ProjectOntology object based on name, but received type {type(ontology)}. Cannot display.")
        logging.error(f"[display_project_ontology] display_project_ontology received unexpected type: {type(ontology)}")
        return

    try:
        # Access attributes directly using getattr for safety or direct access
        title = getattr(ontology, 'title', 'Project Ontology [Default]')
        overview = getattr(ontology, 'overview', 'No overview available.')
        doc_count = getattr(ontology, 'document_count', 0)
        documents = getattr(ontology, 'documents', [])
        global_themes = getattr(ontology, 'global_themes', [])
        key_concepts = getattr(ontology, 'key_concepts', [])
        project_nodes = getattr(ontology, 'project_graph_nodes', None) # Access the graph nodes field

        st.title(f"🌐 {title}")
        st.markdown(f"**Documents Analyzed:** {doc_count}")

        # Display documents included
        if documents:
            st.markdown("## Documents Included In Analysis")
            if hasattr(st, 'chip'):
                for doc_title in documents: st.chip(doc_title, icon="📄")
            else:  # Fallback for older Streamlit versions
                docs_html = " ".join([f"<span style='background-color:#c5e1ff; padding:5px; margin:2px; border-radius:5px'>{doc}</span>" for doc in documents])
                st.markdown(docs_html, unsafe_allow_html=True)

        # Display overview
        st.markdown("## Overview")
        st.markdown(overview)

        # Display global themes
        if global_themes:
            st.markdown("## Global Themes")
            if hasattr(st, 'chip'):
                for theme in global_themes: st.chip(theme, icon="🏷️")
            else: # Fallback
                themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" for tag in global_themes])
                st.markdown(themes_html, unsafe_allow_html=True)

        # Display key concepts
        if key_concepts:
            st.markdown("## Key Concepts")
            if hasattr(st, 'chip'):
                 for concept in key_concepts: st.chip(concept, icon="💡")
            else: # Fallback
                concepts_html = " ".join([f"<span style='background-color:#f0f0f0; padding:5px; margin:2px; border-radius:5px'>{concept}</span>" for concept in key_concepts])
                st.markdown(concepts_html, unsafe_allow_html=True)

        # Display Project-Level Graph
        st.divider()
        st.markdown("## Project Knowledge Graph")
        if project_nodes and isinstance(project_nodes, list):
            logging.info(f"[{function_name}] Rendering project graph with {len(project_nodes)} nodes.")
            with st.expander("View Full Project Concept/Entity Network", expanded=False):
                 try:
                     # Convert Pydantic nodes back to dicts IF display_qdrant_node_network expects dicts
                     # If it accepts Pydantic objects, you can pass project_nodes directly
                     nodes_as_dicts = [node.model_dump(mode='json') if hasattr(node, 'model_dump') else node for node in project_nodes]
                     display_qdrant_node_network(nodes_as_dicts, title="Project-Wide Concept/Entity Network")
                 except NameError:
                      st.error("`display_qdrant_node_network` function not found.")
                      logging.error("`display_qdrant_node_network` function not found.")
                 except Exception as graph_err:
                      st.error(f"Failed to render project network graph: {graph_err}")
                      logging.error(f"Failed to render project network graph: {graph_err}", exc_info=True)
        else:
            st.caption("No cross-document relationship data was generated or processed to display a project graph.")


        # --- NEW: Node Relationship Explorer ---
        st.divider()
        st.markdown("## Node Relationship Explorer")

        if project_nodes and isinstance(project_nodes, list):
            # Create map for faster lookups and consistent data access
            node_map_for_display: Dict[str, Any] = {}
            node_options: Dict[str, str] = {} # Map display label -> node_id
            # Sort nodes alphabetically by name for the dropdown
            sorted_nodes = sorted(project_nodes, key=lambda n: getattr(n, 'name', '') if not isinstance(n, dict) else n.get('name', ''))

            for node in sorted_nodes:
                is_model = not isinstance(node, dict)
                node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
                node_name = getattr(node, 'name', node_id) if is_model else node.get('name', node_id)
                node_type = getattr(node, 'node_type', 'unknown') if is_model else node.get('node_type', 'unknown')

                if node_id:
                    display_label = f"{node_name} ({node_type})"
                    # Handle potential duplicate display labels (though unlikely with type added)
                    count = 1
                    original_label = display_label
                    while display_label in node_options:
                        count += 1
                        display_label = f"{original_label} ({count})"

                    node_options[display_label] = node_id
                    node_map_for_display[node_id] = node # Store original object/dict

            if node_options:
                selected_label = st.selectbox(
                    "Select a Node to see its connections:",
                    options=["Select a Node..."] + list(node_options.keys()),
                    index=0,
                    key="project_node_selector"
                )

                if selected_label != "Select a Node...":
                    selected_id = node_options.get(selected_label)
                    if selected_id:
                        # Call the helper function to display links
                        display_node_links(
                            project_nodes, # Pass the original list
                            selected_id,
                            node_map=node_map_for_display # Pass the map for efficiency
                         )
                    else:
                        st.warning("Could not find the ID for the selected node.")
            else:
                st.info("No valid nodes available to select.")
        else:
            st.caption("No nodes available to explore relationships.")
        # --- END: Node Relationship Explorer ---

    except AttributeError as attr_err:
        # Catch cases where expected attributes are missing from the model instance
        st.error(f"Error accessing project ontology data: Missing attribute - {attr_err}")
        logging.error(f"[{function_name}] Missing attribute in ProjectOntology object: {attr_err}", exc_info=True)
    except Exception as e:
        # Catch any other unexpected errors during display
        st.error(f"An unexpected error occurred while displaying the project ontology: {e}")
        logging.error(f"[{function_name}] Unexpected error displaying ontology: {e}", exc_info=True)


def display_page_details(page_content_data):
    """Displays details for a single page based on the simplified PageContent model."""
    if not page_content_data: return
    page_dict = page_content_data.model_dump() if hasattr(page_content_data, 'model_dump') else page_content_data

    page_num = page_dict.get("page_number", "?")
    st.header(f"📑 Page {page_num} Details")

    # Display Content Flags
    st.caption(f"Detected Content: "
               f"{'Tables ✅' if page_dict.get('has_tables') else 'Tables ❌'} | "
               f"{'Visuals ✅' if page_dict.get('has_visuals') else 'Visuals ❌'} | "
               f"{'Numbers ✅' if page_dict.get('has_numbers') else 'Numbers ❌'}")
    st.divider()

    # Display Subsections first as they contain the core text
    subsections = page_dict.get('subsections', [])
    if subsections:
        st.subheader("Content Subsections")
        for sub_idx, sub_data in enumerate(subsections):
             sub_title = sub_data.get('title', f'Subsection {sub_idx+1}')
             sub_desc = sub_data.get('description', '')
             sub_text = sub_data.get('text', 'No text.')
             sub_id = sub_data.get('subsection_id', 'N/A')
             is_cutoff = sub_data.get('is_cutoff', False)
             ref_v = sub_data.get('referenced_visuals', [])
             ref_t = sub_data.get('referenced_tables', [])

             exp_title = f"{sub_idx+1}. {sub_title}"
             if is_cutoff: exp_title += " (Continues...)"

             with st.expander(exp_title, expanded=True): # Expand by default
                 if sub_desc:
                     st.caption(f"Description: {sub_desc}")
                 st.markdown(sub_text)
                 # Optionally show refs if needed
                 if ref_v or ref_t:
                      refs = []
                      if ref_v: refs.append(f"Visuals: {', '.join(ref_v)}")
                      if ref_t: refs.append(f"Tables: {', '.join(ref_t)}")
                      st.caption(f"References: {'; '.join(refs)} | ID: {sub_id}")
                 else:
                      st.caption(f"ID: {sub_id}")

        st.divider()
    else:
         st.info("No subsections were extracted for this page.")

    # Display Tables
    tables = page_dict.get('tables', [])
    if tables:
        st.subheader(f"Tables on Page {page_num}")
        for table in tables:
            display_table(table) # Use the existing helper
        st.divider()

    # Display Visuals
    visuals = page_dict.get('visuals', [])
    if visuals:
        st.subheader(f"Visuals on Page {page_num}")
        for visual in visuals:
            display_visual_element(visual) # Use the existing helper
        st.divider()

    # Display Numerical Data
    numbers = page_dict.get('numbers', [])
    if numbers:
        st.subheader(f"Numerical Data on Page {page_num}")
        for num_data in numbers:
            display_numerical_data_point(num_data) # Use new helper
        st.divider()

    # --- REMOVED display of old fields: title, topics, summary, entities, page text area ---


#############################
# STREAMLIT UI COMPONENTS (Updated)
#############################

def render_unified_document_report(document_data: Optional[Dict]):
    """
    Displays the full report for a single processed document, including
    Executive Summary (with full graph), Chapters, and Detailed Page View tabs.

    Args:
        document_data: A dictionary containing the finalized document structure,
                       typically the output from finalize_document. Expected structure:
                       {
                           "raw_extracted_content": {
                               "filename": str,
                               "pages": List[Dict], # List of PageContent dicts
                               "summary": Dict | None, # Dict from DocumentSummaryDetails + chapters
                               "error": str | None
                           }
                       }
                       Can also accept the ProcessedDocument model object directly.
    """
    function_name = "render_unified_document_report"
    if not document_data:
        st.error(f"[{function_name}] No document data provided to display.")
        return

    # --- Data Extraction and Validation ---
    raw_content = {}
    # Handle both dict input and potential Pydantic model input
    if hasattr(document_data, 'model_dump'): # Check if it's a Pydantic model
        logging.debug(f"[{function_name}] Input is a Pydantic model, dumping to dict.")
        try:
            doc = document_data.model_dump(mode='json') # Use mode='json' for potential complex types
        except Exception as dump_err:
            st.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}")
            logging.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}", exc_info=True)
            return
    elif isinstance(document_data, dict):
        doc = document_data
    else:
        st.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        logging.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        return

    # Safely get raw_content, checking structure
    if isinstance(doc, dict) and "raw_extracted_content" in doc and isinstance(doc["raw_extracted_content"], dict):
        raw_content = doc["raw_extracted_content"]
        logging.debug(f"[{function_name}] Successfully extracted raw_extracted_content.")
    # Add fallback if top-level dict *is* raw_content (less likely with finalize_document output)
    elif isinstance(doc, dict) and "filename" in doc and "pages" in doc:
        logging.warning(f"[{function_name}] Input dict seems to be raw_extracted_content directly. Using it.")
        raw_content = doc
    else:
        st.error(f"[{function_name}] Could not find 'raw_extracted_content' dictionary in the provided data.")
        logging.error(f"[{function_name}] Invalid data structure, missing 'raw_extracted_content'. Data keys: {list(doc.keys()) if isinstance(doc, dict) else 'N/A'}")
        return

    # Extract key components, providing defaults
    filename = raw_content.get('filename', 'Unknown Document')
    processing_error = raw_content.get('error')
    pages = raw_content.get('pages', []) # Should be list of dicts
    summary_data = raw_content.get('summary', {}) # Should be a dict
    if not isinstance(summary_data, dict):
        logging.warning(f"[{function_name}] 'summary' data for '{filename}' is not a dict ({type(summary_data)}). Resetting to empty dict.")
        summary_data = {} # Ensure it's a dict for safe access later
    chapters_list = summary_data.get('chapters', []) # Get chapter list

    # --- Display Processing Error (if any) ---
    if processing_error:
        st.error(f"🚨 **Processing Error for {filename}:** {processing_error}")
        st.warning("Displayed data below might be incomplete or relate to a previous successful run.")
        # Decide whether to continue rendering or stop here
        # return # Option: Stop rendering if there was a critical error

    # --- Header Section ---
    st.header(f"📄 Report: {filename}")

    # Count actual elements present in pages for the caption
    num_pages = len(pages) if isinstance(pages, list) else 0
    num_tables = 0
    num_visuals = 0
    num_numbers = 0
    if isinstance(pages, list):
        for p in pages:
            if isinstance(p, dict):
                 num_tables += len(p.get('tables', [])) if isinstance(p.get('tables'), list) else 0
                 num_visuals += len(p.get('visuals', [])) if isinstance(p.get('visuals'), list) else 0
                 num_numbers += len(p.get('numbers', [])) if isinstance(p.get('numbers'), list) else 0

    caption_parts = [f"📄 {num_pages} pages"]
    if num_tables > 0: caption_parts.append(f"📊 {num_tables} tables")
    if num_visuals > 0: caption_parts.append(f"🖼️ {num_visuals} visuals")
    if num_numbers > 0: caption_parts.append(f"🔢 {num_numbers} numbers")
    st.caption(" | ".join(caption_parts))

    # Create tabs for document view modes
    tab_titles = ["Executive Summary", "Chapters", "Detailed Page View"]
    tab1, tab2, tab3 = st.tabs(tab_titles)

    # --- Tab 1: Executive Summary ---
    with tab1:
        # Display Summary, Themes, Questions
        with st.container(border=True):
            st.subheader("Executive Summary")
            summary_title = summary_data.get('title')
            if summary_title and summary_title != f"Summary Failed: {filename}" and summary_title != f"Error Finalizing {filename}":
                 st.markdown(f"**Title:** {summary_title}") # Display LLM generated title if valid

            summary_text = summary_data.get('summary')
            if summary_text and not summary_text.startswith("Failed to generate"):
                st.markdown(summary_text)
            else:
                st.info("No summary text available or summary generation failed.")

            themes = summary_data.get('themes')
            if themes and themes != ["error"]:
                st.markdown("#### Key Themes")
                if hasattr(st, 'chip'):
                    for theme in themes: st.chip(theme, icon="🏷️")
                else: # Fallback display
                     st.markdown("- " + "\n- ".join(themes))

            questions = summary_data.get('questions')
            if questions:
                 st.markdown("#### Key Questions Answered")
                 for q in questions: st.markdown(f"- {q}")

        # Display Document-Level Graph
        st.divider()
        st.subheader("Document Concept & Entity Network")

        # 1. Aggregate all nodes from all chapters
        all_document_nodes = []
        if isinstance(chapters_list, list):
            for chapter_dict in chapters_list:
                # Ensure chapter_dict is a dict and contains qdrant_nodes list
                if isinstance(chapter_dict, dict):
                    nodes_in_chapter = chapter_dict.get('qdrant_nodes', [])
                    if isinstance(nodes_in_chapter, list):
                        # Nodes should already be dicts after model_dump in finalize_document
                        all_document_nodes.extend(nodes_in_chapter)
                    else:
                        logging.warning(f"[{function_name}] Chapter '{chapter_dict.get('chapter_id')}' has non-list 'qdrant_nodes'.")
                else:
                    logging.warning(f"[{function_name}] Found non-dictionary item in chapters list for '{filename}'.")

        # 2. Display the aggregated network using the helper function
        if all_document_nodes:
            logging.info(f"[{function_name}] Rendering overall network graph with {len(all_document_nodes)} nodes for {filename}")
            with st.expander("View Full Document Network Graph", expanded=False):
                 try:
                      # Ensure display_qdrant_node_network exists and handles list of dicts
                      display_qdrant_node_network(
                          all_document_nodes,
                          title=f"Overall Concept/Entity Network for {filename}"
                      )
                 except NameError:
                      st.error("`display_qdrant_node_network` function not found.")
                      logging.error("`display_qdrant_node_network` function not found.")
                 except Exception as graph_err:
                      st.error(f"Failed to render document network graph: {graph_err}")
                      logging.error(f"Failed to render document network graph for {filename}: {graph_err}", exc_info=True)
        else:
            st.caption("No concept or entity data was extracted or processed correctly across chapters to build a document network.")

        # Page Content Preview Section
        st.divider()
        st.subheader("Page Content Preview")
        if pages and isinstance(pages, list):
            # Create options for the selectbox
            page_options = {} # Map label to index
            for i, p in enumerate(pages):
                page_num = p.get('page_number', i + 1) if isinstance(p, dict) else i + 1
                page_options[f"Page {page_num}"] = i

            if page_options:
                selected_page_label = st.selectbox(
                    "Select page to preview content:",
                    options=["Select a page..."] + list(page_options.keys()),
                    index=0, # Default to placeholder
                    key=f"page_nav_summary_{filename}" # Unique key per document
                )
                if selected_page_label != "Select a page...":
                    try:
                        page_idx = page_options[selected_page_label]
                        page_to_preview = pages[page_idx]

                        if isinstance(page_to_preview, dict):
                             with st.container(border=True):
                                  st.markdown(f"**Previewing Page {page_to_preview.get('page_number')}**")
                                  # Show subsection titles and descriptions for the selected page
                                  page_subsections = page_to_preview.get('subsections', [])
                                  if page_subsections and isinstance(page_subsections, list):
                                       st.markdown("**Subsections on this page:**")
                                       for sub_idx, sub in enumerate(page_subsections):
                                           if isinstance(sub, dict):
                                                sub_title = sub.get('title', f'Subsection {sub_idx+1}')
                                                sub_desc = sub.get('description', '')
                                                st.markdown(f"- **{sub_title}**: {sub_desc}")
                                           else: st.caption("- Invalid subsection format")
                                  else:
                                       st.caption("No subsections extracted or available for this page.")
                                  # Optionally show raw text preview if helpful
                                  raw_page_text = page_to_preview.get('raw_text')
                                  if raw_page_text:
                                       with st.expander("View Raw Text (preview)"):
                                             st.text(raw_page_text[:1000] + ("..." if len(raw_page_text) > 1000 else ""))

                        else:
                            st.warning(f"Invalid data format for page index {page_idx}.")

                    except (IndexError, KeyError, ValueError) as e:
                        st.warning(f"Could not select page for preview: {e}")
            else:
                st.info("No pages available in this document.")
        else:
            st.info("No page data available for preview.")


    # --- Tab 2: Chapters ---
    with tab2:
        if chapters_list and isinstance(chapters_list, list):
            chapter_titles = [ch.get('title', f"Chapter {ch.get('order', i+1)}") for i, ch in enumerate(chapters_list) if isinstance(ch, dict)]
            if chapter_titles:
                try:
                    chapter_tabs = st.tabs(chapter_titles)
                    for i, (tab, chapter_dict) in enumerate(zip(chapter_tabs, chapters_list)):
                        # Ensure we only process valid chapter dicts
                        if isinstance(chapter_dict, dict):
                            with tab:
                                try:
                                    # Ensure display_chapter exists and handles dicts
                                    display_chapter(chapter_dict)
                                except NameError:
                                    st.error("`display_chapter` function not found.")
                                    logging.error("`display_chapter` function not found.")
                                    break # Stop trying if function is missing
                                except Exception as chapter_disp_err:
                                     st.error(f"Error displaying chapter '{chapter_dict.get('title')}': {chapter_disp_err}")
                                     logging.error(f"Error displaying chapter {chapter_dict.get('chapter_id')}: {chapter_disp_err}", exc_info=True)

                        else:
                             logging.warning(f"[{function_name}] Skipping non-dictionary item found in chapters list for '{filename}'.")

                except Exception as tabs_err: # Catch errors from st.tabs if titles are invalid etc.
                     st.error(f"Could not create chapter tabs: {tabs_err}")
                     logging.error(f"Could not create chapter tabs for {filename}: {tabs_err}", exc_info=True)
            else:
                 st.info("No valid chapter titles found in the summary data.")
        else:
            st.info("No chapter information available or processed for this document.")
            if summary_data.get("error") and "chapter extraction" in summary_data.get("error","").lower():
                 st.warning(f"Note: Chapter extraction previously failed for this document. Reason: {summary_data['error']}")


    # --- Tab 3: Detailed Page View ---
    with tab3:
        if pages and isinstance(pages, list):
             # Create options mapping Page Number (int) -> Index (int)
             page_num_options: Dict[int, int] = {}
             valid_page_nums: List[int] = [] # Store the actual page numbers (integers)
             for i, p in enumerate(pages):
                 # Ensure page data is dict and page_number is an int
                 if isinstance(p, dict) and isinstance(p.get('page_number'), int):
                      page_num = p['page_number']
                      page_num_options[page_num] = i # Map the integer page number to its list index
                      valid_page_nums.append(page_num)
                 else:
                      logging.warning(f"[{function_name}] Skipping invalid page data at index {i} (page_number missing or not int) for page view selector in '{filename}'. Page data keys: {list(p.keys()) if isinstance(p,dict) else 'N/A'}")

             if valid_page_nums:
                  # Sort page numbers numerically for the selector options
                  sorted_page_nums = sorted(valid_page_nums)

                  # Use format_func to display "Page X", but the options are the integers
                  selected_page_num = st.selectbox(
                       "Select Page Number:",
                       options=sorted_page_nums,       # Pass the list of INT page numbers
                       format_func=lambda x: f"Page {x}", # Display "Page X" to the user
                       index=0,                         # Default selection to the first page
                       key=f"page_detail_selector_{filename}" # Unique key for this selectbox
                  )
                  # 'selected_page_num' now directly holds the chosen INTEGER page number (e.g., 5)

                  # Find the corresponding page dictionary using the selected integer page number
                  try:
                       # REMOVED: selected_num = int(selected_page_num.split()[-1])
                       # USE selected_page_num directly as the key
                       selected_index = page_num_options.get(selected_page_num)

                       if selected_index is not None:
                           selected_page_dict = pages[selected_index]
                           try:
                               # Ensure display_page_details exists and handles dicts
                               display_page_details(selected_page_dict)
                           except NameError:
                                st.error("`display_page_details` function not found.")
                                logging.error("`display_page_details` function not found.")
                           except Exception as page_disp_err:
                                st.error(f"Error displaying details for page {selected_page_num}: {page_disp_err}")
                                logging.error(f"Error displaying details for page {selected_page_num} of {filename}: {page_disp_err}", exc_info=True)
                       else:
                           # This should be unlikely if valid_page_nums was populated correctly
                           st.warning(f"Internal error: Could not find index for selected page number {selected_page_num}.")
                  except Exception as e: # Catch any unexpected errors during index lookup/display
                       st.error(f"Error processing page selection or display: {e}")
                       logging.error(f"Error processing page selection {selected_page_num} or display for {filename}: {e}", exc_info=True)

             else:
                  st.info("No valid page data with page numbers found to display details.")
        else:
             st.info("No page data available to display details.")


# --- Sidebar function (Mostly unchanged, ensure it calls correct ontology function if using simplified) ---

def display_sidebar_chat():
    """Manages the sidebar content: file upload, chat, and triggers processing."""
    st.sidebar.title("📝 Input & Chat")
    st.sidebar.markdown("Upload documents and chat about their content.")

    # Define supported file types
    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt"]

    # Initialize session state keys if they don't exist
    if "processed_file_names" not in st.session_state:
        st.session_state.processed_file_names = set()  # Store strings only
    if "processing_active" not in st.session_state:
        st.session_state.processing_active = False
    if "last_uploaded_files" not in st.session_state:
        st.session_state.last_uploaded_files = []
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "processed_documents" not in st.session_state:
        st.session_state.processed_documents = []  # Can store both dicts and models
    if "selected_docs" not in st.session_state:
        st.session_state.selected_docs = []
    if "show_results" not in st.session_state:
        st.session_state.show_results = False
    if "selected_doc_to_view" not in st.session_state:
        st.session_state.selected_doc_to_view = None
    if "project_ontology" not in st.session_state:
        st.session_state.project_ontology = None

    # --- File Upload Section ---
    st.sidebar.subheader("📁 Upload Documents")
    
    uploaded_files = st.sidebar.file_uploader(
        "Select files:",
        type=supported_types,
        accept_multiple_files=True,
        key="sidebar_file_uploader_main"
    )

    # Detect *new* files to process automatically
    new_files_to_process = []
    current_files = {f.name for f in uploaded_files} if uploaded_files else set()
    previous_files = {f.name for f in st.session_state.last_uploaded_files} if st.session_state.last_uploaded_files else set()
    
    # If there are new files and we're not already processing
    if uploaded_files and (current_files != previous_files) and not st.session_state.processing_active:
        st.session_state.last_uploaded_files = uploaded_files
        
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                new_files_to_process.append(file)
    
    # If we have new files to process, prepare and trigger processing directly
    if new_files_to_process and not st.session_state.processing_active:
        st.sidebar.info(f"Processing {len(new_files_to_process)} new files...")
        
        # Create the status elements directly in the sidebar
        status_container = st.sidebar.status("Starting document processing...", expanded=True)
        progress_bar = st.sidebar.progress(0)
        time_info = st.sidebar.empty()
        
        # Store in session state for access
        st.session_state.status_container = status_container
        st.session_state.progress_bar = progress_bar
        st.session_state.time_info = time_info
        
        # Prepare file data
        files_data = []
        for file in new_files_to_process:
            try:
                file_content = file.getvalue()
                file_ext = file.name.split('.')[-1].lower()
                files_data.append({
                    "name": file.name,
                    "content": file_content,
                    "type": file_ext
                })
            except Exception as e:
                st.sidebar.error(f"Error reading file {file.name}: {e}")

        if files_data:
            # Process documents directly
            try:
                st.session_state.processing_active = True
                
                # Call the main processing function
                processed_docs = run_async(
                    process_all_documents_async,
                    files_data
                )

                # Update session state with results
                new_processed_docs = [doc for doc in processed_docs if doc is not None]
                st.session_state.processed_documents.extend(new_processed_docs)

                # Update the set of processed file names
                for file_info in files_data:
                    st.session_state.processed_file_names.add(file_info["name"])

                # Set flag to show results
                st.session_state.show_results = True
                
                status_container.update(label="✅ Document processing finished!", state="complete")
                
            except Exception as e:
                st.sidebar.error(f"Error during document processing: {e}")
                status_container.update(label=f"❌ Processing Error: {e}", state="error")
            finally:
                # Ensure processing flag is reset
                st.session_state.processing_active = False

    # Display already processed files
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---")
        st.sidebar.write("Processed files:")
        for filename in sorted(st.session_state.processed_file_names):
            st.sidebar.caption(f"✓ {filename}")

        col1, col2 = st.sidebar.columns(2) # Use columns for buttons

        # Button to view document reports
        with col1:
            if st.button("View Reports"):
                st.session_state.show_results = True
                st.session_state.show_ontology = False # Deactivate ontology view
                st.rerun() # Rerun to update main area

        # Button to generate/view project ontology
        with col2:
            # Label changes depending on whether ontology exists
            ontology_button_label = "View Ontology" if st.session_state.project_ontology else "Gen. Ontology"
            if st.button(ontology_button_label):
                if st.session_state.processed_documents:
                    # If ontology doesn't exist yet, generate it
                    if not st.session_state.project_ontology:
                        with st.spinner("Generating project ontology..."):
                            try: # Add try/except for generation
                                # Get API client (ensure get_gemini_api_key and client setup work)
                                api_key = get_gemini_api_key()
                                client = genai.Client(api_key=api_key,
                                    http_options={
                                        "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                                        "headers": {
                                            "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                                            "helicone-target-url": 'https://generativelanguage.googleapis.com'
                                        }
                                    })

                                # Generate ontology
                                ontology_data = run_async(
                                    generate_project_ontology, # Use the correct function name
                                    client,
                                    st.session_state.processed_documents
                                )
                                st.session_state.project_ontology = ontology_data
                                st.success("Project ontology generated!")
                            except Exception as gen_err:
                                st.error(f"Ontology generation failed: {gen_err}")
                                st.session_state.project_ontology = None # Ensure it's None on failure

                    # If ontology exists (or was just generated), show it
                    if st.session_state.project_ontology:
                         st.session_state.show_ontology = True
                         st.session_state.show_results = False # Deactivate report view
                         st.rerun() # Rerun to update main area
                    # else: generation failed, message shown above

                else:
                    st.sidebar.warning("No processed documents available for ontology generation.")


    # --- Chat Section --- 
    st.sidebar.divider()
    st.sidebar.subheader("💬 Document Chat")

    # Document selection for chat context
    if st.session_state.processed_documents:
        doc_options = []
        for doc_result in st.session_state.processed_documents:
            filename = None
            
            if isinstance(doc_result, dict):
                raw_content = doc_result.get("raw_extracted_content", {})
                filename = raw_content.get("filename")
            elif hasattr(doc_result, 'filename'):
                filename = doc_result.filename
            elif hasattr(doc_result, 'raw_extracted_content'):
                filename = doc_result.raw_extracted_content.filename
                
            if filename:
                doc_options.append(filename)

        if doc_options:
            st.session_state.selected_docs = st.sidebar.multiselect(
                "Select documents for chat context:",
                options=sorted(list(set(doc_options))),
                default=st.session_state.selected_docs,
                key="doc_context_selector_sidebar"
            )
        else:
            st.sidebar.caption("No processed documents available to reference.")
    else:
        st.sidebar.caption("No documents processed yet. You can still chat, but I won't be able to reference any document content.")

    # Create a container with border in the sidebar for chat history
    with st.sidebar.container(border=True, height=600):
        st.subheader("Chat History")
        # Display chat history
        if not st.session_state.get("messages", []):
            st.info("No messages yet. Start a conversation below!")
        else:
            for message in st.session_state.messages:
                role = message.get("role", "")
                content = message.get("content", "")
                timestamp = message.get("timestamp", "")
                
                if role == "user":
                    with st.chat_message("user"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Sent: {timestamp}")
                elif role == "assistant":
                    with st.chat_message("assistant"):
                        st.markdown(content)
                        if timestamp:
                            st.caption(f"Received: {timestamp}")
                elif role == "system":
                    with st.chat_message("system"):
                        st.markdown(f"*{content}*")
                        if timestamp:
                            st.caption(f"System: {timestamp}")

        # Chat input
        if prompt := st.sidebar.chat_input("Ask a question...", key="sidebar_chat_input_main"):
            current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages.append({"role": "user", "content": prompt, "timestamp": current_time})
            run_async(process_chat_message, prompt)
            st.rerun()

async def process_chat_message(message):
    """Process a chat message with clear user feedback."""
    # Make sure the message list exists
    if "messages" not in st.session_state:
        st.session_state.messages = []
        
    # Add a thinking indicator that will be updated with the real response
    current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.session_state.messages.append({"role": "assistant", "content": "Thinking...", "timestamp": current_time})
    
    try:
        # Get API client
        api_key = get_gemini_api_key()
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                }
            })
        
        # Build context from documents (if any)
        context = ""
        selected_docs = st.session_state.get("selected_docs", [])
        processed_docs = st.session_state.get("processed_documents", [])
        
        if selected_docs and processed_docs:
            # Extract content from selected documents for context
            context = "Using document context from: " + ", ".join(selected_docs) + "\n\n"
            
            for doc_name in selected_docs:
                # Find the document in processed_documents
                selected_doc = None
                for doc in processed_docs:
                    if isinstance(doc, dict) and "raw_extracted_content" in doc:
                        filename = doc["raw_extracted_content"].get("filename", "")
                        if filename == doc_name:
                            selected_doc = doc
                            break
                
                if selected_doc and "raw_extracted_content" in selected_doc:
                    # Extract summary and key pages
                    tech_data = selected_doc["raw_extracted_content"]
                    
                    # Add summary if available
                    if "summary" in tech_data and tech_data["summary"]:
                        summary = tech_data["summary"]
                        context += f"Document: {doc_name}\n"
                        context += f"Title: {summary.get('title', '')}\n"
                        context += f"Summary: {summary.get('summary', '')}\n\n"
                    
                    # Add content from up to 3 key pages
                    if "pages" in tech_data and tech_data["pages"]:
                        pages = tech_data["pages"]
                        # Sort pages by importance (tables, numbers, length)
                        sorted_pages = sorted(pages, 
                                            key=lambda p: (p.get("has_tables", False), 
                                                          p.get("has_numbers", False), 
                                                          len(p.get("text", ""))), 
                                            reverse=True)[:3]
                        
                        for page in sorted_pages:
                            context += f"Page {page.get('page_number', '')}: {page.get('text', '')[:500]}...\n\n"
        
        # Simple prompt for demonstration
        prompt = f"""
        You are a helpful assistant that analyzes documents. 
        
        Document context:
        {context}
        
        Please respond to: {message}
        """
        
        # Generate content - using the correct async pattern
        response = await client.aio.models.generate_content(
            model="gemini-2.5-flash-preview-04-17",
            contents=[
                types.Content(
                    role="user",
                    parts=[types.Part.from_text(text=prompt)]
                )
            ]
        )
        
        # Update the thinking message with the actual response
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            if response and response.candidates:
                content = response.candidates[0].content.parts[0].text
                st.session_state.messages[-1]["content"] = content
                st.session_state.messages[-1]["timestamp"] = timestamp
            else:
                st.session_state.messages[-1]["content"] = "Sorry, I couldn't generate a response."
                st.session_state.messages[-1]["timestamp"] = timestamp
        
    except Exception as e:
        # Update thinking message with the error
        if st.session_state.messages and st.session_state.messages[-1].get("role") == "assistant":
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages[-1]["content"] = f"Sorry, an error occurred: {str(e)}"
            st.session_state.messages[-1]["timestamp"] = timestamp
    
    # Trigger UI update
    st.rerun()

##############
# Main Application
##############

# --- Main Application ---
def main():
    """Main function to run the Streamlit application."""
    st.set_page_config(
        page_title="Multimodal Document Processor",
        page_icon="📊",
        layout="wide"
    )

    # Initialize session state keys robustly
    default_session_state = {
        "processed_documents": [],
        "processed_file_names": set(),
        "processing_active": False,
        "last_uploaded_files": [],
        "messages": [],
        "selected_docs": [], # For chat context
        "show_results": False, # To show document reports
        "selected_doc_to_view": None, # Specific doc name (if needed)
        "project_ontology": None, # Holds ProjectOntology object
        "show_ontology": False, # To show ontology view
        "status_container": None, # UI elements
        "progress_bar": None,
        "time_info": None
    }
    for key, default_value in default_session_state.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

    # --- SIDEBAR ---
    display_sidebar_chat() # Contains file upload, buttons, chat

    # --- MAIN CONTENT AREA ---
    st.title("📊 Advanced Document Processor")
    st.markdown("Upload PDFs, Office documents, text files, and more to extract structure, relationships, and summaries.")
    st.divider()

    if st.session_state.get("processing_active"):
        st.info("Processing documents... Check sidebar for progress.")
        # Optionally display a general spinner in the main area too
        with st.spinner("Processing..."):
            time.sleep(2) # Keep UI responsive
    elif st.session_state.get("show_ontology") and st.session_state.project_ontology:
        # Display project ontology if requested and available
        display_project_ontology(st.session_state.project_ontology)
    elif st.session_state.get("show_results") and st.session_state.processed_documents:
        # Display document reports if requested and available
        st.header("Document Reports")
        # Document Selection Logic
        doc_options = {} # Map display name to index
        for i, doc_result in enumerate(st.session_state.processed_documents):
            filename = "Unknown Document"
            if isinstance(doc_result, dict) and "raw_extracted_content" in doc_result:
                 filename = doc_result["raw_extracted_content"].get("filename", f"Doc_{i+1}")
            doc_options[f"{i+1}: {filename}"] = i # Map display label to original index

        if len(st.session_state.processed_documents) > 1:
             selected_display_name = st.selectbox(
                 "Choose document report to view:",
                 options=list(doc_options.keys()),
                 index=0, # Default to first
                 key="doc_report_selector"
             )
             selected_index = doc_options.get(selected_display_name, 0)
        elif len(st.session_state.processed_documents) == 1:
            selected_index = 0 # Only one doc
        else:
            selected_index = -1 # Should not happen if check above passes

        if selected_index != -1:
            selected_doc_data = st.session_state.processed_documents[selected_index]
            # Render the selected document report
            render_unified_document_report(selected_doc_data)
        else:
            # This case should ideally not be reached if processed_documents is not empty
            st.info("Select a document or process files.")
    else:
        # Display welcome screen
        st.info("Upload documents using the sidebar to begin analysis.")

        # Project description - Updated
        st.markdown("""
        ## 📊 Advanced Multimodal Document Analyzer

        Leverage the power of Google's Gemini models to unlock insights from your documents. This application offers:

        - **Broad File Support**: Process PDFs (image-based), Office documents (Word, PowerPoint, Excel), plain text, Markdown, code files, and more.
        - **Intelligent Structure Extraction**: Automatically segment documents into meaningful **subsections**, identifying titles, descriptions, and content flow. Detects and extracts **tables**, **visuals** (charts, images), and key **numerical data points** within each page or chunk.
        - **Logical Chapter Generation**: Automatically group related subsections across pages into coherent **chapters** based on thematic content, providing a structured overview.
        - **Concept & Entity Graphing**: Identify key **concepts** and **entities** within each chapter and analyze their relationships, including connections **between chapters**, building a knowledge graph foundation. (Visualized per chapter).
        - **Automated Summarization**: Generate a concise **document-level summary**, including key themes and potential questions, derived from the extracted concepts and entities.
        - **Project-Wide Ontology (Optional)**: Create a high-level overview synthesizing the main themes and concepts across *all* uploaded documents.
        - **Interactive Chat**: Engage in conversation with your documents; ask questions based on the extracted content and summaries (context-aware based on selected documents).

        ### Getting Started:

        1.  **Upload:** Use the sidebar to upload one or more supported documents. Processing will begin automatically for new files.
        2.  **Monitor:** Track the processing progress in the sidebar.
        3.  **Explore Reports:** Once finished, click "View Document Reports" in the sidebar. Use the tabs:
            *   **Executive Summary:** Read the AI-generated summary and key themes.
            *   **Chapters:** Navigate through the structured chapters and view their concept/entity networks.
            *   **Detailed Page View:** Examine the extracted subsections, tables, visuals, and numbers for each page.
        4.  **Chat:** Select documents in the sidebar and use the chat interface to ask specific questions.
        5.  **Synthesize:** Click "Generate Project Ontology" (optional) to get a combined overview of all processed documents.
        """)

        # Placeholder Image (Consider replacing with a relevant architecture diagram or feature collage)
        st.image("https://via.placeholder.com/800x400.png?text=Document+Structure+and+Analysis",
                 caption="From Raw Files to Structured Insights and Knowledge Graphs")

if __name__ == "__main__":
    # Basic logging setup
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
    main()

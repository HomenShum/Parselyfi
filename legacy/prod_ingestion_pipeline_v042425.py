# --- Standard Library Imports ---
import asyncio
import base64
import io
import json
import logging
import os
import re
import sys
import time
import uuid
import hashlib
import datetime
import random
import pathlib
import traceback
import csv
from enum import Enum
from collections import defaultdict, Counter
from typing import (
    Any, Dict, List, Optional, Tuple, Union, Literal, Type, Callable, Set
)

# --- Third-Party Imports ---
# Attempt to import essential libraries first
try:
    from pydantic import (
        BaseModel, Field, ConfigDict, field_validator, model_validator, ValidationError
    )
except ImportError:
    print("CRITICAL ERROR: Pydantic not found. Please install it: pip install pydantic")
    sys.exit(1)

try:
    import streamlit as st
except ImportError:
    st = None # Allow running without streamlit for backend testing
    print("Streamlit not found, UI features will be unavailable.")

try:
    from google import genai
    from google.generativeai import types as genai_types
    from google.api_core import exceptions as google_exceptions
except ImportError:
    print("CRITICAL ERROR: Google Generative AI SDK not found. Please install it: pip install google-generativeai")
    # Stop execution if core dependency is missing
    sys.exit(1)

# PDF/Office/Media/ML Processing Libraries (Optional, handled gracefully)
try:
    import pymupdf as fitz # PyMuPDF
except ImportError:
    fitz = None
    print("PyMuPDF not found (pip install pymupdf). PDF image extraction will be limited.")

try:
    import pandas as pd
except ImportError:
    pd = None
    print("Pandas not found (pip install pandas). Tabular data processing will be limited.")

try:
    import docx
except ImportError:
    docx = None
    print("python-docx not found (pip install python-docx). DOCX processing will be unavailable.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("python-pptx not found (pip install python-pptx). PPTX processing will be unavailable.")

try:
    import openpyxl # Needed by pandas for xlsx
except ImportError:
    openpyxl = None
    print("openpyxl not found (pip install openpyxl). XLSX processing might fail.")

# Qdrant/Embedding/Graph Libraries (Optional, handled gracefully)
try:
    from qdrant_client import AsyncQdrantClient, models
    from qdrant_client.http.models import ( # Use http models for clarity in definition
        PointStruct, Distance, VectorParams, SparseVectorParams, SparseIndexParams,
        Filter, FieldCondition, MatchValue, Range, CollectionStatus, SparseVector
    )
except ImportError:
    AsyncQdrantClient, models, PointStruct, Distance, VectorParams, SparseVectorParams, SparseIndexParams, Filter, FieldCondition, MatchValue, Range, CollectionStatus, SparseVector = None, None, None, None, None, None, None, None, None, None, None, None, None
    print("qdrant-client not found (pip install qdrant-client). Qdrant features will be unavailable.")

try:
    from fastembed import TextEmbedding, SparseTextEmbedding
except ImportError:
    TextEmbedding, SparseTextEmbedding = None, None
    print("fastembed not found (pip install fastembed[qdrant]). Embedding/Qdrant ingestion will be unavailable.")

try:
    import networkx as nx
except ImportError:
    nx = None
    print("NetworkX not found (pip install networkx). Graph visualization will be unavailable.")

try:
    import matplotlib.pyplot as plt
except ImportError:
    plt = None
    print("Matplotlib not found (pip install matplotlib). Graph visualization will be unavailable.")

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    import numpy as np
except ImportError:
    TfidfVectorizer, cosine_similarity, np = None, None, None
    print("Scikit-learn/Numpy not found (pip install scikit-learn numpy). Some advanced text analysis features might be limited.")

# --- Helper Function for Safe Import Access ---
def _check_import(module, name=""):
    if module is None:
        module_name = name or (module.__name__ if hasattr(module, '__name__') else 'Required library')
        logging.error(f"{module_name} is not installed. Please install the necessary dependencies.")
        if st:
            st.error(f"Missing Library: {module_name}. Please install it. Some features may be disabled.")
        # Don't raise error immediately, allow app to load partially maybe?
        # raise ImportError(f"{module_name} is not installed.")
        return False
    return True


############################################################
# Logging Configuration
############################################################

def setup_logging(log_level_str: str = "INFO", log_file: str = "document_processor.log"):
    """Configure logging."""
    log_level = getattr(logging, log_level_str.upper(), logging.INFO)
    # Ensure handlers are not added multiple times if called again
    root_logger = logging.getLogger()
    if root_logger.hasHandlers():
        # Optionally remove existing handlers if re-configuration is needed
        # for handler in root_logger.handlers[:]:
        #     root_logger.removeHandler(handler)
        pass # Or just skip setup if already configured
    else:
        log_formatter = logging.Formatter('%(asctime)s - %(levelname)s - [%(name)s:%(funcName)s:%(lineno)d] %(message)s')
        # Stream Handler (stdout)
        stream_handler = logging.StreamHandler(sys.stdout)
        stream_handler.setFormatter(log_formatter)
        root_logger.addHandler(stream_handler)
        # File Handler
        try:
            file_handler = logging.FileHandler(log_file, mode="a", encoding='utf-8')
            file_handler.setFormatter(log_formatter)
            root_logger.addHandler(file_handler)
        except Exception as e:
            print(f"Warning: Could not setup file logging to {log_file}: {e}")

        root_logger.setLevel(log_level)
        # Set lower levels for noisy libraries
        logging.getLogger("httpx").setLevel(logging.WARNING)
        logging.getLogger("httpcore").setLevel(logging.WARNING)
        logging.getLogger("fastembed").setLevel(logging.INFO) # Adjust as needed
        logging.getLogger("qdrant_client").setLevel(logging.INFO) # Adjust as needed

        logging.info(f"Logging initialized. Level: {log_level_str.upper()}. File: {log_file}")

# Call setup_logging early
setup_logging()
_logger = logging.getLogger(__name__) # General logger for the file


############################################################
# Qdrant Configuration (Add these near the top)
############################################################

# Fetch configuration from Streamlit secrets safely
QDRANT_URL = None
QDRANT_API_KEY = None
QDRANT_COLLECTION_NAME = "project_ontology_nodes_v3" # Choose a suitable name, maybe versioned
EMB_MODEL_NAME = "BAAI/bge-base-en-v1.5" # Ensure this matches the model loaded
SPARSE_MODEL_NAME = "Qdrant/bm25"          # Ensure this matches the model loaded
EMB_DIM = 768                             # Dimension for bge-base-en-v1.5

# Configuration for upsert batching
BATCH_SIZE = 100 # How many points to send in one upsert call
INGEST_CONCURRENCY = 5 # How many batches to process per UI update cycle
HNSW_EF_SEARCH = 64
OVERSAMPLING = 4.0
DENSE_M_PROD = 16
SEARCH_PREFETCH_LIMIT = 50

# Attempt to load secrets only if Streamlit is available
if st:
    try:
        QDRANT_URL = st.secrets.get("qdrant_url")
        QDRANT_API_KEY = st.secrets.get("qdrant_api_key") # Optional: For Qdrant Cloud
        if not QDRANT_URL:
            _logger.warning("QDRANT_URL not found in Streamlit secrets.")
            st.sidebar.warning("Qdrant URL not configured in secrets.", icon="⚠️")
        else:
            _logger.info(f"Qdrant URL loaded from secrets: {QDRANT_URL}")
        if QDRANT_API_KEY:
            _logger.info("Qdrant API Key loaded from secrets.")
        else:
            _logger.info("Qdrant API Key not found or not needed.")

    except Exception as e:
        _logger.error(f"Error accessing Streamlit secrets: {e}")
        st.sidebar.error("Error loading Qdrant configuration from secrets.")
else:
    # Fallback for non-Streamlit environment (e.g., load from environment variables)
    QDRANT_URL = os.environ.get("QDRANT_URL")
    QDRANT_API_KEY = os.environ.get("QDRANT_API_KEY")
    if not QDRANT_URL:
        _logger.warning("QDRANT_URL not found in environment variables.")
    else:
        _logger.info(f"Qdrant URL loaded from env: {QDRANT_URL}")


############################################################
# DATA MODELS (Using models from the user's combined `main.py` for compatibility)
############################################################
class VisualElement(BaseModel):
    """Model for visual elements like charts, graphs, etc."""
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data, or N/A")
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")
    source_url: Optional[str] = Field(None, description="Source URL of the visual")
    alt_text: Optional[str] = Field(None, description="Alternative text for the visual")
    visual_id: str = Field(..., description="Unique identifier for the visual")

class NumericalDataPoint(BaseModel):
    """Model for numerical data points extracted from text."""
    value: str = Field(..., description="The numerical value as a string")
    description: str = Field(..., description="What the number represents, with units/context")
    context: Optional[str] = Field(None, description="Surrounding text context, if available")

class TableData(BaseModel):
    """Model for tables extracted from documents."""
    table_content: str = Field(..., description="Markdown formatted table content")
    title: Optional[str] = Field(None, description="Optional table title")
    summary: Optional[str] = Field(None, description="Optional table summary")
    page_number: int = Field(..., description="Page number where the table appears")
    table_id: str = Field(..., description="Unique identifier for the table")

class Subsection(BaseModel):
    """Model for subsections extracted from pages."""
    subsection_id: str = Field(..., description="Unique identifier for the subsection")
    order: int = Field(..., description="Order of the subsection within the page")
    title: str = Field(..., description="Title of the subsection (less than 7 words)")
    text: str = Field(..., description="Full text content of the subsection")
    description: str = Field(..., description="One line description summarizing the main point")
    is_cutoff: bool = Field(..., description="True if content appears to be cut off by page break")
    referenced_visuals: List[str] = Field(default_factory=list, description="IDs of referenced visuals")
    referenced_tables: List[str] = Field(default_factory=list, description="IDs of referenced tables")

class PageContent(BaseModel):
    """
    Model for structured content extracted from a single page.
    Stage 1 populates raw_text, tables, visuals, numbers.
    Stage 2 populates subsections from raw_text.
    """
    page_number: int = Field(..., description="Page number in the original document")
    has_tables: bool = Field(..., description="Indicates if any tables were extracted from this page")
    has_visuals: bool = Field(..., description="Indicates if any visuals (charts, images, etc.) were extracted")
    has_numbers: bool = Field(..., description="Indicates if any key numerical data points were extracted")
    tables: List[TableData] = Field(default_factory=list, description="List of tables extracted from the page")
    visuals: List[VisualElement] = Field(default_factory=list, description="List of visual elements extracted from the page")
    numbers: List[NumericalDataPoint] = Field(default_factory=list, description="List of numerical data points extracted from the page")
    raw_text: Optional[str] = Field(None, description="Full raw text extracted initially from the page/chunk")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections extracted from raw_text")

    model_config = ConfigDict(arbitrary_types_allowed=True, extra='forbid')

class LinkedNode(BaseModel):
    """Structure for linked nodes metadata in Qdrant."""
    target_node_id: str = Field(..., description="ID of the target node (e.g., type_name_chX)")
    relationship_description: str = Field(..., description="Nature of the relationship from source to target")
    relationship_keywords: List[str] = Field(default_factory=list, description="Keywords summarizing the relationship")
    relationship_strength: float = Field(..., description="Strength of the relationship (1-10)")

    @field_validator('relationship_strength', mode='before')
    def strength_to_float(cls, v):
        try:
            val = float(v)
            # Clamp value between 1.0 and 10.0
            return max(1.0, min(val, 10.0))
        except (ValueError, TypeError):
            return 5.0 # Default strength on conversion error

class QdrantNode(BaseModel):
    """Structure for nodes to be stored in Qdrant, using direct fields."""
    node_id: str = Field(..., description="Unique node ID (e.g., UUID generated from type/name/context)")
    node_type: str = Field(..., description="'entity', 'concept', 'document', or 'meta_summary'")
    name: str = Field(..., description="Entity, concept, or document name")
    chapter_id_context: str = Field(..., description="Chapter ID or other context identifier (e.g., 'doc_level::filename')")
    title: str = Field(..., description="Concise title for this node's content focus within its context (5-10 words)")
    core_claim: str = Field(..., description="1-3 sentence summary capturing the node's main point, assertion, or function.")
    information_type: str = Field(..., description="Classification of the main information type (e.g., 'Definition', 'Explanation', 'Process Description').")
    key_entities: str = Field(..., description="JSON stringified list of key entity names mentioned in context. Example: '[\"Acme Corp\", \"Project X\"]'. Use '[]' if none.")
    tags: str = Field(..., description="JSON stringified list of relevant keywords/tags (1-5 tags, lowercase, underscores). Example: '[\"ai\", \"optimization\"]'. Use '[]' if none.")
    sentiment_tone: str = Field(..., description="Overall sentiment/tone (e.g., 'Positive', 'Neutral', 'Objective').")
    source_credibility: str = Field(..., description="Inferred credibility (e.g., 'Company Report', 'Technical Documentation', 'Unknown').")
    source: Optional[str] = Field(None, description="Source identifier (e.g., filename/URL), added programmatically.")
    stored_at: Optional[str] = Field(None, description="ISO timestamp of when the node was stored, added programmatically.")
    original_node_id: Optional[str] = Field(None, description="Original ID before UUID generation, if applicable.")
    executive_summary: Optional[str] = Field(None, description="Executive summary, primarily for doc/meta nodes.")
    linked_nodes: List[LinkedNode] = Field(default_factory=list, description="List of outgoing relationships")

    model_config = ConfigDict(extra='ignore')

    @field_validator('key_entities', 'tags', mode='before')
    def check_json_string(cls, value):
        if isinstance(value, list):
            try:
                return json.dumps(value)
            except TypeError:
                 raise ValueError("List provided for key_entities/tags could not be JSON serialized")
        if not isinstance(value, str):
            raise ValueError('Must be a string')
        try:
            loaded = json.loads(value)
            if not isinstance(loaded, list):
                 raise ValueError('JSON string must represent a list')
        except json.JSONDecodeError:
            raise ValueError('Must be a valid JSON string representation of a list')
        return value

class Chapter(BaseModel):
    """Model for chapters composed of subsections, including Qdrant node data."""
    chapter_id: str = Field(..., description="Unique identifier for the chapter")
    title: str = Field(..., description="Title of the chapter")
    summary: str = Field(..., description="Summary of the chapter")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections in this chapter")
    order: int = Field(..., description="Order of the chapter in the document")
    qdrant_nodes: Optional[List[QdrantNode]] = Field(None, description="Structured nodes and relationships for Qdrant")

class DocumentSummaryDetails(BaseModel):
    """Structure for summary details generated from node analysis."""
    title: str = Field(..., description="Concise title for the document based on nodes")
    themes: List[str] = Field(default_factory=list, description="Main themes/topics derived from nodes")
    questions: List[str] = Field(default_factory=list, description="Sample questions reflecting node content")
    summary: str = Field(..., description="Comprehensive summary synthesizing node information")
    # Adding chapters for compatibility with UI display logic
    chapters: List[Chapter] = Field(default_factory=list, description="List of chapters in the document")

# Removed DocumentReport and related models as they are superseded by the pipeline output

class ProcessedDocument(BaseModel): # Retained for potential internal use, but pipeline output is ProcessedDocumentOutput
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages")
    summary: Optional[DocumentSummaryDetails] = None

class ProjectOntology(BaseModel):
    """Represents the aggregated knowledge graph across documents."""
    title: str = Field(..., description="Project title")
    overview: str = Field(..., description="Project overview based on synthesized nodes")
    document_count: int = Field(..., description="Number of documents analyzed")
    documents: List[str] = Field(..., description="List of document filenames included in analysis")
    global_themes: List[str] = Field(..., description="High-level project themes derived from nodes")
    key_concepts: List[str] = Field(..., description="Key concepts identified across all documents from nodes")
    project_graph_nodes: Optional[List[QdrantNode]] = Field(None, description="Aggregated nodes and their relationships across the project.")

    model_config = ConfigDict(arbitrary_types_allowed=True)


############################################################
# Qdrant/Embedding Client Factory & Utilities
############################################################

# Qdrant Client Factory (Cached)
@st.cache_resource(show_spinner="Connecting to Qdrant...")
def get_qdrant_client():
    """Initializes and caches the AsyncQdrantClient."""
    if not _check_import(AsyncQdrantClient, "qdrant-client"):
        return None
    if not QDRANT_URL:
        _logger.error("Qdrant URL is not configured.")
        if st: st.error("Qdrant URL is not configured. Cannot initialize client.")
        return None

    _logger.info(f"Initializing Qdrant client for {QDRANT_URL}")
    try:
        client = AsyncQdrantClient(
            url=QDRANT_URL,
            api_key=QDRANT_API_KEY, # Handles None if not provided
            prefer_grpc=False, # Set to True if using gRPC
            timeout=60 # Adjust timeout as needed
        )
        # Optional: Perform a quick check if connection works (in async context)
        # Needs to be called from an async function, e.g., during button press
        _logger.info("Qdrant client object created (connection not verified here).")
        return client
    except Exception as e:
        _logger.error(f"Qdrant client initialization failed: {e}", exc_info=True)
        if st: st.error(f"Failed to initialize Qdrant client: {e}")
        return None

# Embedding models (Cached)
@st.cache_resource(show_spinner="Loading embedding models...")
def load_models():
    """Loads and caches the dense and sparse embedding models."""
    if not _check_import(TextEmbedding, "fastembed") or not _check_import(SparseTextEmbedding, "fastembed"):
        return None, None
    try:
        _logger.info(f"Loading embedding models ({EMB_MODEL_NAME}, {SPARSE_MODEL_NAME})...")
        start = time.time()
        # Consider adding cache_dir parameter if needed
        dense_model = TextEmbedding(EMB_MODEL_NAME)
        sparse_model = SparseTextEmbedding(SPARSE_MODEL_NAME)
        _logger.info(f"Embedding models loaded in {time.time() - start:.2f} seconds.")
        return dense_model, sparse_model
    except Exception as e:
        _logger.error(f"Failed to load embedding models: {e}", exc_info=True)
        if st: st.error(f"Failed to load embedding models: {e}")
        return None, None

def get_or_create_eventloop():
    """Gets the current event loop or creates a new one if none exists."""
    try:
        return asyncio.get_running_loop()
    except RuntimeError:
        _logger.info("No running event loop found, creating a new one.")
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        return loop

def run_async(func, *args, **kwargs):
    """Runs an async function, creating a new event loop if necessary."""
    try:
        loop = asyncio.get_running_loop()
        if loop.is_running():
             # If a loop is running (like in Streamlit context sometimes),
             # schedule the coroutine and wait for it.
             # This might require nest_asyncio in some complex scenarios.
             _logger.warning("Detected running event loop. Attempting to schedule task.")
             future = asyncio.ensure_future(func(*args, **kwargs))
             # This approach might still block if the outer loop isn't processed.
             # Using a dedicated thread might be safer but adds complexity.
             # For simplicity in Streamlit, nest_asyncio or threadpool might be needed
             # if this pattern causes issues. Let's try direct run first.
             # return loop.run_until_complete(future) # This can fail if loop is running
             # Fallback: Run in new loop (might cause issues if interacting with existing async state)
             new_loop = asyncio.new_event_loop()
             return new_loop.run_until_complete(func(*args, **kwargs))
        else:
             # If no loop is running, run it normally
             return loop.run_until_complete(func(*args, **kwargs))
    except RuntimeError:
        # No loop exists, create one
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            return loop.run_until_complete(func(*args, **kwargs))
        finally:
            # loop.close() # Closing the loop can cause issues in some environments
            asyncio.set_event_loop(None) # Detach the loop


############################################################
# Gemini API & Utility Functions
############################################################

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or (st.secrets.get("GEMINI_API_KEY") if st else None)
    if not api_key:
        _logger.error("Gemini API key not found.")
        if st:
            st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
            st.stop()
        else:
            raise ValueError("Gemini API key not found.")
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    _check_import(genai, "google.generativeai")
    _check_import(google_exceptions, "google.api_core.exceptions")
    last_exception = None
    base_delay = 1.5 # Base delay in seconds

    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)

            # Gemini Specific Checks (Safety, Content Parts)
            if isinstance(response, genai_types.GenerateContentResponse):
                if not response.candidates:
                    if response.prompt_feedback and response.prompt_feedback.block_reason:
                        reason = response.prompt_feedback.block_reason.name
                        safety_ratings = {rating.category.name: rating.probability.name for rating in response.prompt_feedback.safety_ratings}
                        err_msg = f"API call blocked: {reason}. Ratings: {safety_ratings}"
                        _logger.error(err_msg)
                        # Raise a specific, non-retryable error
                        raise google_exceptions.PermissionDenied(err_msg)
                    else:
                        # No candidates, no block reason - potentially transient?
                        err_msg = f"API call returned no candidates (attempt {attempt+1})"
                        _logger.warning(err_msg)
                        # Raise an error that *might* be retryable depending on the exception handling below
                        raise ValueError(err_msg)

                # Check for empty content parts in the first candidate
                if response.candidates:
                    first_candidate = response.candidates[0]
                    if not first_candidate.content or not first_candidate.content.parts:
                         err_msg = f"API call candidate has no content parts (attempt {attempt+1})"
                         _logger.warning(err_msg)
                         raise ValueError(err_msg) # Treat as potentially retryable


            # Validate JSON if applicable (check response_mime_type in config)
            gen_config = kwargs.get('generation_config') or kwargs.get('config') # Check both possible names
            is_json_response = False
            if isinstance(gen_config, genai_types.GenerationConfig) and gen_config.response_mime_type == 'application/json':
                is_json_response = True
            elif isinstance(gen_config, dict) and gen_config.get('response_mime_type') == 'application/json':
                is_json_response = True

            if is_json_response and response.candidates:
                try:
                    json_text = response.candidates[0].content.parts[0].text
                    clean_json_response(json_text) # Use clean function (which includes parse check)
                except (json.JSONDecodeError, ValueError) as json_e:
                    _logger.warning(f"Malformed JSON detected on attempt {attempt+1}: {json_e}")
                    if attempt < max_retries - 1:
                        last_exception = json_e
                        # Use exponential backoff with jitter
                        delay = base_delay * (2 ** attempt) + random.uniform(0, 0.5)
                        _logger.info(f"Retrying after delay: {delay:.2f}s")
                        await asyncio.sleep(delay)
                        continue # Retry
                    else:
                        _logger.error(f"Malformed JSON after {max_retries} attempts.")
                        raise json_e from None # Raise the JSON error after retries

            return response # Success

        except (google_exceptions.ResourceExhausted, google_exceptions.ServiceUnavailable, google_exceptions.DeadlineExceeded, asyncio.TimeoutError) as e:
            last_exception = e
            if attempt == max_retries - 1:
                _logger.error(f"API call failed after {max_retries} attempts due to retryable error: {e}")
                raise last_exception from e
            # Exponential backoff with jitter
            delay = base_delay * (2 ** attempt) + random.uniform(0, 0.5)
            _logger.warning(f"API call failed (Attempt {attempt+1}/{max_retries}, retryable: {type(e).__name__}). Retrying in {delay:.2f}s...")
            await asyncio.sleep(delay)

        except google_exceptions.PermissionDenied as e:
             # Non-retryable safety block
             _logger.error(f"API call permanently blocked by safety settings: {e}")
             raise # Re-raise immediately

        except Exception as e:
            # Catch other unexpected errors
            _logger.error(f"API call failed unexpectedly (Attempt {attempt+1}/{max_retries}): {type(e).__name__} - {e}", exc_info=True)
            last_exception = e
            # Decide whether to retry based on type? For now, assume non-retryable
            raise # Re-raise immediately

    # Fallback if loop completes without returning or raising (shouldn't happen)
    _logger.critical("Exited retry loop unexpectedly.")
    raise last_exception or Exception(f"API call failed after {max_retries} attempts with unknown error")


def clean_json_response(json_text: Optional[str], extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        _logger.warning("clean_json_response received None input, returning empty JSON object.")
        return "{}"

    # Ensure it's a string
    json_text = str(json_text).strip()
    if not json_text:
        _logger.warning("clean_json_response received empty string input, returning empty JSON object.")
        return "{}"

    # Handle markdown code blocks
    if json_text.startswith("```"):
        cleaned_text = re.sub(r"^```(?:json)?\s*", "", json_text, flags=re.IGNORECASE)
        cleaned_text = re.sub(r"\s*```$", "", cleaned_text)
        cleaned_text = cleaned_text.strip()
    elif json_text.lower().startswith("json"):
         cleaned_text = json_text[4:].strip()
    else:
         cleaned_text = json_text

    # First attempt: Check if it's already valid JSON
    try:
        json.loads(cleaned_text)
        return cleaned_text # It's valid, return as is
    except json.JSONDecodeError:
        _logger.debug("Initial JSON parse failed, attempting fixes...")
        pass # Continue to fixes

    # Common JSON fixes (Add more if needed)
    # Fix trailing commas in arrays/objects
    fixed_text = re.sub(r',(\s*[}\]])', r'\1', cleaned_text)
    # Add more robust fixes here if necessary (e.g., missing quotes, escaped chars)
    # Be cautious with regex fixes, they can break valid JSON.

    # Final attempt to parse with fixes
    try:
        json.loads(fixed_text)
        _logger.info("JSON successfully parsed after applying fixes.")
        return fixed_text
    except json.JSONDecodeError as e:
        _logger.warning(f"JSON parsing failed even after fixes: {e}. Original text (start): '{json_text[:200]}...' Fixed attempt: '{fixed_text[:200]}...'")

        # Fallback Strategy (Optional)
        if extract_text_on_failure:
            _logger.info("Attempting text extraction fallback from failed JSON.")
            # Simple fallback: return the cleaned/fixed text wrapped in a basic JSON structure
            fallback_data = {
                "error": f"JSON parsing failed: {e}",
                "raw_llm_output_preview": fixed_text[:1000] + ("..." if len(fixed_text) > 1000 else ""),
                # Set default values for expected keys to prevent downstream errors
                "raw_text": fixed_text, # Best guess for text content
                "tables": [], "visuals": [], "numbers": [], "subsections": [],
                "has_tables": False, "has_visuals": False, "has_numbers": False,
                "title": "Processing Error", "summary": "Could not parse LLM response.",
                "themes": ["error"], "questions": [], "chapters": [],
                "qdrant_nodes": [], "inter_chapter_links": [], "document_summary_details": None,
                "node_analysis": [], "document_summary": {}
            }
            try:
                return json.dumps(fallback_data)
            except TypeError: # Should not happen with this structure
                 return '{"error": "JSON fallback creation failed"}'
        else:
            # If not extracting text, return a minimal error JSON
             return '{"error": "Invalid JSON structure", "details": "' + str(e).replace('"', '\\"') + '"}'

    except Exception as e:
        _logger.error(f"Unexpected error during JSON cleaning: {e}", exc_info=True)
        return '{"error": "Error cleaning JSON response", "details": "' + str(e).replace('"', '\\"') + '"}'


############################################################
# Content Extraction Functions (Adapted from main.py)
############################################################

# --- Prompt Templates ---

# Prompt for initial structure extraction (raw_text + elements) - STAGE 1
page_structure_prompt_template = """
--- Goal ---
Analyze the input, which represents a single page ({page_num}) of a document '{source_identifier}'. Your primary task is to extract the complete raw text content AND identify and structure specific elements like tables, visuals, and key numerical data points found on that page.

--- Input ---
- Content of page {page_num} (provided as image or text).

--- Instructions ---
Perform the following tasks based *only* on the provided input content for page {page_num}:

**Task 1: Extract Raw Text**
1a. Extract ALL text content visible on the page.
1b. Preserve original formatting like paragraphs and line breaks as accurately as possible. Handle lists, headings, and code blocks appropriately.
1c. Store this complete text in the `raw_text` field of the output JSON. If no text is found, use an empty string "".

**Task 2: Extract Tables**
2a. Identify all distinct tables present on the page.
2b. For EACH table found:
    *   Generate a unique `table_id` formatted as `page_{page_num}_table_N` (where N is 1, 2, 3...).
    *   Extract the complete cell content accurately.
    *   Format this content STRICTLY as GitHub Flavored Markdown within the `table_content` field (using '|' separators and a '---' header separator line). Do NOT include surrounding paragraphs or explanatory text in `table_content`.
    *   If a clear title or caption is associated with the table, extract it into the `title` field. Otherwise, use `null`.
    *   Write a brief (1-2 sentence) `summary` of the table's main content or purpose, if possible. Otherwise, use `null`.
    *   Set `page_number` to {page_num}.
    *   Add the complete table object to the `tables` list in the output JSON.
2c. If NO tables are found on the page, the `tables` list MUST be an empty list `[]`.

**Task 3: Extract Visual Elements**
3a. Identify all visual elements (e.g., charts, graphs, diagrams, images, photographs, equations, maps).
3b. For EACH visual element found:
    *   Generate a unique `visual_id` formatted as `page_{page_num}_visual_N` (where N is 1, 2, 3...).
    *   Determine its `type` (e.g., "bar chart", "line graph", "diagram", "image", "photo", "map").
    *   Write a detailed `description` covering what the visual shows, its purpose, and key elements visible.
    *   Provide a `data_summary` summarizing key data points, trends, or findings presented in the visual. If the visual doesn't present data (e.g., a decorative image) or a summary isn't applicable/possible, use the string "N/A".
    *   Set `page_numbers` to `[{page_num}]`.
    *   Set `source_url` and `alt_text` to `null` unless explicitly available.
    *   Add the complete visual object to the `visuals` list in the output JSON.
3c. If NO visual elements are found, the `visuals` list MUST be an empty list `[]`.

**Task 4: Extract Numerical Data**
4a. Identify significant numerical data points mentioned in the page's text (excluding those already inside tables). Look for percentages, currencies, quantities, measurements etc.
4b. For EACH significant number found:
    *   Extract the `value` as a string (e.g., "123.45", "50%", "$1.2M", "3/4").
    *   Write a `description` explaining what the number represents, including units or context (e.g., "increase in revenue", "percentage completion", "total users").
    *   Extract the surrounding text snippet (~10-20 words) providing immediate context into the `context` field. If no clear surrounding context is available, use `null`.
    *   Add the complete number object to the `numbers` list in the output JSON.
4c. If NO significant numerical data points are found, the `numbers` list MUST be an empty list `[]`.

**Task 5: Set Boolean Flags**
5a. Set `has_tables` to `true` if the `tables` list is NOT empty, otherwise set it to `false`.
5b. Set `has_visuals` to `true` if the `visuals` list is NOT empty, otherwise set it to `false`.
5c. Set `has_numbers` to `true` if the `numbers` list is NOT empty, otherwise set it to `false`.

**Task 6: Assemble Final JSON**
6a. Combine all extracted information (`raw_text`, `has_tables`, `tables`, `has_visuals`, `visuals`, `has_numbers`, `numbers`) into a single JSON object adhering precisely to the schema defined below.

--- Output Format ---
Return ONLY a single, valid JSON object. Do NOT include any text, comments, or markdown formatting like ```json before or after the JSON object. It MUST strictly follow this schema:

```json
{{
  "raw_text": "string | null",
  "has_tables": true,
  "has_visuals": false,
  "has_numbers": true,
  "tables": [
    {{
      "table_id": "page_{page_num}_table_1",
      "table_content": "| Header 1 | Header 2 |\\n|---|---|\\n| Data A | Data B |",
      "title": "Example Table Title",
      "summary": "This table shows example data.",
      "page_number": {page_num}
    }}
  ],
  "visuals": [],
  "numbers": [
    {{
      "value": "50%",
      "description": "Completion rate of project X",
      "context": "... reached a 50% completion rate..."
    }}
  ]
}}
```

--- Critical Rules ---
*   VALID JSON ONLY: The entire output MUST be a single JSON object starting with { and ending with }.
*   SCHEMA ADHERENCE: Follow the schema EXACTLY. Ensure all REQUIRED fields are present (see models above for hints, though follow THIS specific structure for the JSON output). Use `null` for optional string fields if no value is applicable/found. Use `[]` for empty lists.
*   MARKDOWN TABLES: Table content MUST be formatted as GitHub Flavored Markdown in the `table_content` field.
*   TEXTUAL BASIS: All extracted information must be derived SOLELY from the provided page content. Do not infer or add external knowledge.

Generate the JSON output for page {page_num} now:
"""

# Prompt for subsection extraction from raw text (STAGE 2)
subsection_extraction_prompt_template = """
Analyze the provided text content from a single page ({page_num}) of document '{source_identifier}'. Your goal is to segment this text into logical subsections, providing concise yet **informative descriptions geared towards identifying entities and relationships later**.

Input Text:
--- START TEXT ---
{raw_text_content}
--- END TEXT ---

Instructions:
1.  Read the text and identify logical breaks based on topic shifts, headings (if any within the text), or paragraph structure. Aim for subsections that represent coherent units of meaning.
2.  Segment the text into sequential subsections. Ensure the entire input text is covered without loss or overlap.
3.  For each subsection, determine:
    *   `subsection_id`: Generate a unique ID formatted as `page_{page_num}_section_ORDER` where ORDER is the sequential order number (1-based).
    *   `order`: The sequential order number (integer, starting from 1).
    *   `title`: A concise title (less than 7 words). Use headings from the text if present and appropriate, otherwise generate a descriptive title.
    *   `text`: The full text content belonging *only* to this specific subsection. Ensure correct segmentation.
    *   `description`: **CRITICAL FOR ANALYSIS:** A concise, one-sentence summary (approx. 15-25 words) that **highlights the key entities, concepts, or specific actions/interactions** discussed **within this specific subsection**. Focus on specific nouns and verbs.
    *   `is_cutoff`: Set to `true` ONLY if this specific text chunk appears to end mid-sentence or mid-thought, suggesting it might be truncated *before* being passed to you. Otherwise `false`.
    *   `referenced_visuals`: List of visual IDs (e.g., `page_{page_num}_visual_N`) mentioned or clearly referred to within this subsection's text. Empty list `[]` if none.
    *   `referenced_tables`: List of table IDs (e.g., `page_{page_num}_table_N`) mentioned or clearly referred to within this subsection's text. Empty list `[]` if none.

4.  Return ONLY a valid JSON array containing the subsection objects, structured exactly like this example:
    ```json
    [
      {{
        "subsection_id": "page_{page_num}_section_1",
        "order": 1,
        "title": "System Overview",
        "text": "The full text content of the first subsection...",
        "description": "Introduces the Parsely system developed for data enhancement.",
        "is_cutoff": false,
        "referenced_visuals": ["page_{page_num}_visual_1"],
        "referenced_tables": []
      }},
      {{
        "subsection_id": "page_{page_num}_section_2",
        "order": 2,
        "title": "Core Technologies",
        "text": "Details on the integration of various technologies...",
        "description": "Details use of GPT-40, Llamalndex, and Qdrant for analysis.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": ["page_{page_num}_table_1"]
      }}
    ]
    ```
CRITICAL: Your entire response MUST be ONLY the JSON array, starting with `[` and ending with `]`. Do NOT include ```json markdown fences, introductory text, explanations, or any other characters outside the JSON structure. Ensure all string values are enclosed in double quotes and properly escaped. Validate your JSON structure carefully.
"""

# Prompt for JSON repair (if initial LLM calls fail parsing)
json_repair_prompt_template = """
The following text was generated by an AI attempting to create a valid JSON array or object based on previous instructions. However, the text has syntax errors and is not valid JSON.

Your task is to analyze the provided text, correct the JSON syntax errors, and return ONLY the corrected, valid JSON (either an array `[...]` or an object `{{...}}` as appropriate based on the content). Do NOT add any introductory text, explanations, or markdown formatting like ```json.

Ensure:
- All strings are enclosed in double quotes.
- Internal quotes within strings are properly escaped (e.g., \\").
- Objects within arrays are correctly separated by commas.
- Array elements are correctly separated by commas.
- Key-value pairs within objects are correctly separated by commas.
- There are no trailing commas after the last element in an array or the last property in an object.
- All brackets (`[]`) and braces (`{{}}`) are correctly paired and closed.

Faulty Text to Repair:
--- START FAULTY TEXT ---
{faulty_json_text}
--- END FAULTY TEXT ---

Return ONLY the valid JSON:
"""

# Prompt for Chapter generation (STAGE 3)
chapter_extraction_prompt_template = """
Analyze the following list of subsections extracted sequentially from document '{source_identifier}'. Each subsection includes its ID, title, and description. Group these subsections into logically coherent chapters (typically 3-10 chapters per document, unless very short).

Subsections List:
```json
{subsection_context_json}
```

Instructions:
1. Group the provided subsections into chapters based on topic cohesion and narrative flow.
2. Each chapter MUST contain at least one subsection.
3. Every subsection ID from the input list must be assigned to exactly ONE chapter.
4. Maintain the original relative order of subsections within and across chapters.
5. For each chapter created, provide:
    - A unique `chapter_id` (e.g., "chapter_1", "chapter_2").
    - A concise and descriptive `title` (max 7 words).
    - A brief `summary` (1-2 sentences) describing the chapter's content.
    - `subsection_ids`: A list of subsection IDs belonging to that chapter, in order.
    - `order`: The sequential chapter number (1-based).

Return ONLY a valid JSON array containing the chapter objects, structured exactly like this example:
```json
[
  {{
    "chapter_id": "chapter_1",
    "order": 1,
    "title": "Introduction and Setup",
    "summary": "Provides background information and initial setup.",
    "subsection_ids": ["page_1_section_1", "page_1_section_2", "page_2_section_1"]
  }},
  {{
    "chapter_id": "chapter_2",
    "order": 2,
    "title": "Core Functionality",
    "summary": "Details the main features and operations.",
    "subsection_ids": ["page_2_section_2", "page_3_section_1"]
  }}
]
```
Ensure the output is only the JSON array, starting with `[` and ending with `]`. No extra text.
"""

# Prompt for Node & Relationship analysis (STAGE 3)
node_relationship_prompt_template = """
--- Goal ---
Analyze the text content within a single document chapter ('{chapter_id}': '{chapter_title}') from source '{source_identifier}'. Identify key specific entities (e.g., organizations, products) AND important abstract concepts (e.g., processes, topics). Extract these as nodes. For each node, define its properties based *only* on the chapter text. Also, identify direct relationships *between these nodes within this chapter*.

--- Input Data ---
Chapter Title: {chapter_title}
Chapter ID: {chapter_id}
Source Identifier (Filename/URL): {source_identifier}
Chapter Content (Combined Subsection Text):
{combined_text}
--- End Input Data ---

--- Instructions ---
**Task: Extract Nodes and Intra-Chapter Relationships**
1.  **Identify Nodes:** Scan the text and identify distinct key entities and abstract concepts (~5-15 significant nodes per chapter recommended).
2.  **Define Node Structure:** For each node identified, structure its information:
    *   `node_type`: "entity" or "concept".
    *   `name`: Original, human-readable name (e.g., "Acme Corp", "Machine Learning Workflow").
    *   `chapter_id_context`: Set to exactly "{chapter_id}".
    *   `title`: Concise title (5-10 words) for the node's discussion *in this chapter*.
    *   `core_claim`: 1-3 sentence summary capturing the node's central point/function *in this chapter*. Be specific.
    *   `information_type`: Primary info type (e.g., "Definition", "Explanation", "Process Step", "Example").
    *   `key_entities`: JSON string list of *other specific entity names* explicitly mentioned in the node's context (e.g., `"[\\"NVIDIA GPUs\\", \\"TensorFlow\\"]"`). Use `"[]"` if none.
    *   `tags`: JSON string list of 1-5 relevant keywords (lowercase, underscores, e.g., `"[\\"rag_framework\\", \\"llm_enhancement\\"]"`). Use `"[]"` if none.
    *   `sentiment_tone`: Overall sentiment/tone for this node in the chapter (e.g., "Positive", "Neutral", "Objective").
    *   `source_credibility`: Inferred credibility based on source type/tone (e.g., "Company Report", "Technical Doc", "Unknown").
    *   `linked_nodes`: An initially empty list `[]`.
3.  **Identify Relationships (Linked Nodes):** For each node created (source), find its most important (max 5-7) direct relationships to *other nodes created in this Task* (target).
    *   For each relationship, create a `linked_nodes` object: `target_node_id` (must be another node's generated ID from *this* response), `relationship_description` (SPECIFIC verbs, e.g., "calculates cost based on", "is a component of"), `relationship_keywords` (optional list), `relationship_strength` (float 1.0-10.0).
    *   Populate the `linked_nodes` list of the source node. If no links, use `[]`.
4.  **Generate Node ID:** For each node, create a unique `node_id` using the format `{{type}}_{{sanitized_name}}_ch{{chapter_num_str}}`. Ensure uniqueness within this response.
5.  **Assemble Output:** Format all extracted node objects into a single JSON object.

--- Output Format ---
Return ONLY a single, valid JSON object containing ONE top-level key: "qdrant_nodes". The value MUST be a JSON array of node objects. Adhere STRICTLY to this schema:
```json
{{
  "qdrant_nodes": [
    {{
      "node_id": "string", // Generated: e.g., "concept_machine_learning_ch1"
      "node_type": "string", // "entity" or "concept"
      "name": "string",
      "chapter_id_context": "string", // e.g., "{chapter_id}"
      "title": "string",
      "core_claim": "string",
      "information_type": "string",
      "key_entities": "string", // JSON list e.g., "[\\"Other Entity\\"]"
      "tags": "string", // JSON list e.g., "[\\"tag1\\", \\"tag2\\"]"
      "sentiment_tone": "string",
      "source_credibility": "string",
      "linked_nodes": [
        {{
          "target_node_id": "string", // ID of another node in THIS list
          "relationship_description": "string",
          "relationship_keywords": ["string"], // Optional
          "relationship_strength": 10.0 // Float 1.0-10.0
        }}
      ]
      // DO NOT include 'source' or 'stored_at' fields
    }}
  ]
}}
```
--- Critical Rules ---
*   VALID JSON ONLY: Output must be a single JSON object `{{...}}`. No extra text.
*   SCHEMA ADHERENCE: Follow the schema precisely. All required fields must be present.
*   JSON STRINGS: Ensure `key_entities` and `tags` are valid JSON strings representing lists. Double-check escaping.
*   NODE ID CONSISTENCY: `target_node_id` MUST correspond to a `node_id` defined in the SAME `qdrant_nodes` list.
*   TEXTUAL BASIS: All information must be derived SOLELY from the provided "Chapter Content".

Generate the JSON output now:
"""

# Prompt for Inter-Chapter Links and Document Summary (STAGE 3)
inter_chapter_summary_prompt_template = """
---Goal---
Analyze the provided list of nodes (concepts/entities) aggregated from all chapters of document '{source_identifier}'. Perform two tasks:
1. Identify significant relationships *between* nodes originating from *different* chapters.
2. Generate a concise document-level summary based on the overall collection of nodes.

---Input Data---
List of node definitions (node_id is the unique UUID):
```json
{nodes_json_for_prompt}
```

---Instructions---
**Task 1: Identify Inter-Chapter Relationships**
1a. Compare nodes across the entire list, focusing on connections between nodes with different `chapter_id_context`.
1b. Identify pairs (`source_node`, `target_node`) with **different** `chapter_id_context` values that have a direct, meaningful relationship based on names/descriptions/tags. Use the provided `node_id` (UUIDs) for linking.
1c. For each significant inter-chapter relationship identified, determine: `source_node_id` (source UUID), `target_node_id` (target UUID), `relationship_description` (specific, verb-driven), `relationship_keywords` (list, optional), `relationship_strength` (float 1.0-10.0).
1d. **Limit:** Max 5-7 *strongest* inter-chapter links *per source node*.

**Task 2: Generate Document Summary**
2a. Synthesize information from the *entire list* of nodes.
2b. Generate:
    *   `title`: Concise title for the document '{source_identifier}' (max 10 words).
    *   `themes`: List of 3-7 main themes/topics.
    *   `questions`: List of 2-4 insightful questions the document might answer.
    *   `summary`: Comprehensive summary paragraph (4-8 sentences) synthesizing key nodes and narrative.

**Output Format:**
Return ONLY a single valid JSON object with TWO top-level keys: "inter_chapter_links" and "document_summary_details".

*   `inter_chapter_links`: JSON array of relationship objects from Task 1 (using UUIDs).
*   `document_summary_details`: JSON object with summary results from Task 2 (`title`, `themes`, `questions`, `summary`).

Example:
```json
{{
  "inter_chapter_links": [
    {{
      "source_node_id": "uuid-...",
      "target_node_id": "uuid-...",
      "relationship_description": "Utilizes concept from Chapter X...",
      "relationship_keywords": ["optimization"],
      "relationship_strength": 8.5
    }}
  ],
  "document_summary_details": {{
    "title": "AI Development Overview",
    "themes": ["Machine Learning", "GPU Optimization"],
    "questions": ["How is ML applied?", "What are optimization challenges?"],
    "summary": "The document outlines core ML concepts... applied to GPU optimization..."
  }}
}}
```
CRITICAL: Ensure the output is ONLY the single valid JSON object. Use the provided `node_id` UUIDs in links.
"""

# Prompt for Project-Level Ontology Generation (Optional - called separately)
project_ontology_prompt_template = """
Act as Knowledge Graph Architect. Create an initial project-level ontology by analyzing summaries from multiple documents.

Input: List of Document/Meta Summaries:
```json
{doc_nodes_json}
```

Tasks:
1. Find significant inter-document relationships between the provided nodes (using their `node_id`, e.g., 'doc::filename' or 'meta::idx'). Base relationships on the content of their `executive_summary`.
2. Generate project metadata: `title`, `overview` (synthesized from all summaries), `global_themes` (list), `key_concepts` (list of the most important original document names or meta-summary themes).

Output Schema (JSON ONLY):
```json
{{
  "project_summary_details": {{
    "title": "string",
    "overview": "string",
    "global_themes": ["string"],
    "key_concepts": ["string"]
  }},
  "inter_document_links": [
    {{
      "source_node_id": "string", // e.g., "doc::report.pdf" or "meta::group_1"
      "target_node_id": "string",
      "relationship_description": "string",
      "relationship_keywords": ["string"],
      "relationship_strength": 7.5 // Float 1.0-10.0
    }}
  ]
}}
```
Rules: Valid JSON ONLY. Use the provided `node_id`s for links. Base analysis SOLELY on the executive summaries provided. `key_concepts` should reflect the most central document topics or meta-summary themes.
Generate JSON:
"""

# --- Extraction and Processing Functions ---

def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List[Dict[str, Any]]] = None) -> PageContent:
    """Creates a PageContent object representing an error state."""
    _logger.error(f"Creating error page for page {page_num} due to: {error_msg}")
    error_text = f"Error processing page {page_num}: {error_msg}"
    if validation_errors:
        try:
            error_details = "\n\nValidation Errors:\n" + "\n".join(
                f"- Field '{'.'.join(map(str, err.get('loc', ['N/A'])))}': {err.get('msg', 'No message')}"
                for err in validation_errors
            )
            error_text += error_details
        except Exception as format_error:
            _logger.warning(f"Could not format validation errors for page {page_num}: {format_error}")
            error_text += "\n(Could not format validation error details)"

    # Create minimal error subsection
    error_subsection = Subsection(
        subsection_id=f"page_{page_num}_section_error",
        order=1,
        title="Processing Error",
        text=error_text,
        description=f"Error processing page {page_num}.",
        is_cutoff=False,
        referenced_visuals=[],
        referenced_tables=[]
    )

    error_page = PageContent(
        page_number=page_num,
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        subsections=[error_subsection], # Include the error subsection
        raw_text=f"[ERROR: {error_msg}]" # Store error in raw_text as well
    )
    return error_page

def create_error_fallback_subsection(page_num: int, raw_text: Optional[str], error_msg: str) -> Subsection:
     """Creates a single fallback Subsection when extraction fails, preserving raw text."""
     fallback_id = f"page_{page_num}_section_fallback_error"
     fallback_title = "Full Page Content (Subsection Extraction Failed)"
     fallback_text = (
         f"--- ERROR DURING SUBSECTION EXTRACTION ---\n"
         f"{error_msg}\n"
         f"--- END ERROR ---\n\n"
         f"--- ORIGINAL RAW TEXT FOR PAGE {page_num} ---\n"
         f"{raw_text or '[Raw text not available]'}\n" # Handle None raw_text
         f"--- END RAW TEXT ---"
     )
     fallback_description = f"Could not segment page {page_num} into subsections: {error_msg[:100]}..."

     try:
          return Subsection(
              subsection_id=fallback_id, order=1, title=fallback_title,
              text=fallback_text, description=fallback_description,
              is_cutoff=False, referenced_visuals=[], referenced_tables=[]
          )
     except ValidationError as ve_fb:
         _logger.error(f"Critical: Failed to create fallback Subsection for page {page_num}: {ve_fb}", exc_info=True)
         # Return a minimal valid Subsection even if data is compromised
         return Subsection(
              subsection_id=fallback_id, order=1, title="CRITICAL FALLBACK ERROR",
              text=f"Error creating fallback object: {ve_fb}", description="Critical error.",
              is_cutoff=False, referenced_visuals=[], referenced_tables=[]
          )


async def extract_page_content_from_image_or_text(
    client: genai.GenerativeModel,
    page_num: int,
    source_identifier: str,
    image_part: Optional[genai_types.Part] = None,
    text_content: Optional[str] = None
) -> PageContent:
    """
    Extracts raw text and structural elements (tables, visuals, numbers) from a single page image OR text.
    Populates PageContent WITHOUT subsections initially (Stage 1).
    """
    function_name = "extract_page_content_from_image_or_text"
    _logger.debug(f"[{function_name}] Starting Stage 1 extraction for page/unit {page_num} of '{source_identifier}'")

    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    # --- Prepare Prompt and Content ---
    try:
        prompt = page_structure_prompt_template.format(page_num=page_num, source_identifier=source_identifier)
        content_parts = [genai_types.Part.from_text(prompt)]

        if image_part:
            content_parts.append(image_part)
            _logger.debug(f"[{function_name}] Using image input for page {page_num}.")
        elif text_content:
            # Limit text input size if necessary
            MAX_TEXT_LENGTH = 150000 # Adjust based on model limits
            if len(text_content) > MAX_TEXT_LENGTH:
                 _logger.warning(f"[{function_name}] Truncating text input for page {page_num} ({len(text_content)} -> {MAX_TEXT_LENGTH} chars).")
                 text_content = text_content[:MAX_TEXT_LENGTH] + "\n... [Content Truncated]"
            content_parts.append(genai_types.Part.from_text(text_content))
            _logger.debug(f"[{function_name}] Using text input for page {page_num}.")
        else:
            raise ValueError("No image or text content provided for extraction.")

        # --- Call LLM (Stage 1) ---
        response = await retry_api_call(
            client.generate_content,
            content_parts, # Pass the list of parts
            generation_config=genai_types.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.1 # Lower temp for structured extraction
            )
        )

        if not response.candidates:
            raise ValueError("LLM response had no candidates.")

        json_text = clean_json_response(response.candidates[0].content.parts[0].text)
        data = json.loads(json_text)

        if isinstance(data, list): data = data[0] if data else {}
        if not isinstance(data, dict):
             _logger.warning(f"[{function_name}] LLM returned non-dict JSON for page {page_num}. Attempting fallback.")
             # Use the fallback mechanism within clean_json_response if enabled
             # Or create a basic error structure here
             raise ValueError(f"LLM returned non-dictionary JSON: {type(data)}")

        # Populate page_data from the LLM response, validating types
        page_data["raw_text"] = data.get("raw_text") if isinstance(data.get("raw_text"), str) else None
        page_data["has_tables"] = bool(data.get("has_tables", False))
        page_data["has_visuals"] = bool(data.get("has_visuals", False))
        page_data["has_numbers"] = bool(data.get("has_numbers", False))

        # Pydantic will validate list items during PageContent creation
        page_data["tables"] = data.get("tables", []) if isinstance(data.get("tables"), list) else []
        page_data["visuals"] = data.get("visuals", []) if isinstance(data.get("visuals"), list) else []
        page_data["numbers"] = data.get("numbers", []) if isinstance(data.get("numbers"), list) else []

        # Basic fallback if raw text wasn't extracted by LLM but text was input
        if page_data["raw_text"] is None and text_content:
             _logger.warning(f"[{function_name}] LLM failed to return 'raw_text' for page {page_num}. Using original input text as fallback.")
             page_data["raw_text"] = text_content

        # Create PageContent object (will validate internal models)
        page_content_obj = PageContent(**page_data)
        _logger.info(f"[{function_name}] Successfully extracted Stage 1 structure for page {page_num} of '{source_identifier}'.")
        return page_content_obj

    except (json.JSONDecodeError, ValidationError, ValueError) as e:
        _logger.error(f"[{function_name}] Failed Stage 1 processing for page {page_num}: {e}", exc_info=True)
        # If text_content was the input, try to preserve it in the error page
        error_raw_text = text_content if text_content else f"[ERROR during Stage 1: {e}]"
        error_page = create_error_page(page_num, f"Stage 1 extraction failed: {e}")
        error_page.raw_text = error_raw_text # Override default error text if input text available
        return error_page
    except Exception as e:
        _logger.error(f"[{function_name}] Unexpected error during Stage 1 for page {page_num}: {e}", exc_info=True)
        error_page = create_error_page(page_num, f"Unexpected Stage 1 error: {e}")
        error_page.raw_text = text_content if text_content else f"[UNEXPECTED ERROR: {e}]"
        return error_page


async def extract_subsections_from_text_content(
    client: genai.GenerativeModel,
    raw_text: Optional[str],
    page_num: int,
    source_identifier: str,
    max_retries: int = 2,
    retry_delay: int = 2
) -> List[Subsection]:
    """
    Takes raw text from a page and uses an LLM to segment it into Subsection objects (Stage 2).
    Includes JSON repair mechanism.
    """
    function_name = "extract_subsections_from_text_content"
    _logger.debug(f"[{function_name}] Starting subsection extraction (Stage 2) for page {page_num} of '{source_identifier}'")

    if not raw_text or not raw_text.strip():
        _logger.warning(f"[{function_name}] No raw text provided for page {page_num}. Returning empty list.")
        return []

    # Limit text size for the prompt
    MAX_TEXT_LIMIT = 150000 # Adjust as needed
    text_for_prompt = raw_text
    if len(raw_text) > MAX_TEXT_LIMIT:
        _logger.warning(f"[{function_name}] Truncating raw text for page {page_num} from {len(raw_text)} to {MAX_TEXT_LIMIT} chars.")
        text_for_prompt = raw_text[:MAX_TEXT_LIMIT] + "... [Content Truncated]"

    initial_prompt = subsection_extraction_prompt_template.format(
        page_num=page_num,
        raw_text_content=text_for_prompt,
        source_identifier=source_identifier
    )

    faulty_json_text_to_repair = None
    parsed_data = None

    # --- Initial LLM Call with Retries ---
    for attempt in range(max_retries):
        _logger.debug(f"[{function_name}] Initial LLM call attempt {attempt + 1}/{max_retries} for page {page_num}")
        try:
            response = await retry_api_call(
                client.generate_content,
                [initial_prompt],
                generation_config=genai_types.GenerationConfig(
                    response_mime_type="application/json",
                    temperature=0.2
                )
            )
            if not response.candidates: raise ValueError("LLM response had no candidates.")

            raw_json_text = response.candidates[0].content.parts[0].text
            cleaned_json_text = clean_json_response(raw_json_text, extract_text_on_failure=False) # Don't extract text yet
            _logger.debug(f"[{function_name}] Initial Cleaned JSON Response (Attempt {attempt + 1}, Page {page_num}): {cleaned_json_text[:500]}...")

            parsed_data = json.loads(cleaned_json_text) # Attempt to parse

            # Structure check (expecting a list)
            if isinstance(parsed_data, list):
                _logger.info(f"[{function_name}] Initial call successful (got list) on attempt {attempt + 1} for page {page_num}.")
                faulty_json_text_to_repair = None # Success, no repair needed
                break # Exit retry loop
            else:
                 _logger.warning(f"[{function_name}] Initial call response was not a list (type: {type(parsed_data)}) on attempt {attempt+1}, Page {page_num}.")
                 # Treat as parsing failure for retry/repair logic
                 raise json.JSONDecodeError(f"Expected list, got {type(parsed_data)}", cleaned_json_text, 0)

        except (json.JSONDecodeError, ValueError) as e:
            _logger.warning(f"[{function_name}] Initial JSON parse/validation failed (Attempt {attempt + 1}, Page {page_num}): {e}")
            faulty_json_text_to_repair = raw_json_text # Store the original faulty text for repair
            if attempt < max_retries - 1:
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue
            else:
                _logger.error(f"[{function_name}] Failed initial JSON parse/validation after {max_retries} attempts for page {page_num}.")
                break # Exit retry loop, proceed to repair attempt

        except Exception as e:
            _logger.error(f"[{function_name}] Error during initial LLM call (Attempt {attempt + 1}, Page {page_num}): {e}", exc_info=True)
            faulty_json_text_to_repair = f"LLM call failed: {e}" # Store error info
            if attempt < max_retries - 1:
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue
            else:
                 _logger.error(f"[{function_name}] Initial call failed definitively after {max_retries} attempts for page {page_num}.")
                 break

    # --- JSON Repair Attempt ---
    if faulty_json_text_to_repair and parsed_data is None: # Only repair if initial call failed AND didn't parse
        _logger.info(f"[{function_name}] Attempting JSON repair for page {page_num}.")
        try:
            repair_prompt = json_repair_prompt_template.format(faulty_json_text=faulty_json_text_to_repair)
            repair_response = await retry_api_call(
                client.generate_content,
                [repair_prompt],
                generation_config=genai_types.GenerationConfig(
                    response_mime_type="application/json", # Expect repair to return JSON
                    temperature=0.1
                ),
                max_retries=1 # Only one retry for repair attempt
            )
            if repair_response.candidates:
                repaired_json_text = clean_json_response(repair_response.candidates[0].content.parts[0].text, extract_text_on_failure=False)
                _logger.debug(f"[{function_name}] Repaired JSON Response (Page {page_num}): {repaired_json_text[:500]}...")
                try:
                    repaired_data = json.loads(repaired_json_text)
                    if isinstance(repaired_data, list):
                        _logger.info(f"[{function_name}] Successfully parsed JSON after repair call for page {page_num}.")
                        parsed_data = repaired_data # Use the repaired data
                    else:
                        _logger.error(f"[{function_name}] Repaired JSON was not a list for page {page_num}. Type: {type(repaired_data)}")
                except json.JSONDecodeError as repair_err:
                    _logger.error(f"[{function_name}] Failed to decode JSON even after repair attempt for page {page_num}: {repair_err}")
            else:
                 _logger.error(f"[{function_name}] Repair LLM call returned no candidates for page {page_num}.")

        except Exception as repair_call_err:
            _logger.error(f"[{function_name}] Error during repair LLM call for page {page_num}: {repair_call_err}", exc_info=True)


    # --- Process Final Data (if successful) or Fallback ---
    validated_subsections: List[Subsection] = []
    if parsed_data is not None and isinstance(parsed_data, list):
        _logger.debug(f"[{function_name}] Validating {len(parsed_data)} potential subsection items for page {page_num}.")
        items_skipped = 0
        for i, sub_dict in enumerate(parsed_data):
            if not isinstance(sub_dict, dict):
                _logger.warning(f"[{function_name}] Skipping item {i+1} (not a dict) for page {page_num}.")
                items_skipped += 1
                continue
            if not sub_dict.get("text", "").strip():
                _logger.warning(f"[{function_name}] Skipping item {i+1} (empty text) for page {page_num}.")
                items_skipped += 1
                continue

            try:
                # Ensure required fields exist with defaults if possible
                sub_dict.setdefault("referenced_visuals", [])
                sub_dict.setdefault("referenced_tables", [])
                sub_dict.setdefault("is_cutoff", False)
                # Ensure order is int, default to index
                try: sub_dict['order'] = int(sub_dict.get('order', i + 1))
                except (ValueError, TypeError): sub_dict['order'] = i + 1

                # Create Subsection object (Pydantic validation)
                subsection_obj = Subsection(**sub_dict)
                validated_subsections.append(subsection_obj)
            except ValidationError as ve:
                _logger.warning(f"[{function_name}] Pydantic validation failed for subsection {i+1} on page {page_num}: {ve}. Skipping. Data: {sub_dict}")
                items_skipped += 1
            except Exception as val_err:
                _logger.error(f"[{function_name}] Unexpected validation error for subsection {i+1} on page {page_num}: {val_err}. Skipping.", exc_info=True)
                items_skipped += 1

        # Post-process: Sort and re-number order
        if validated_subsections:
            validated_subsections.sort(key=lambda s: s.order)
            for i, sub in enumerate(validated_subsections): sub.order = i + 1 # Ensure sequential order
            _logger.info(f"[{function_name}] Successfully extracted and validated {len(validated_subsections)} subsections for page {page_num} (skipped {items_skipped}).")
            return validated_subsections
        else:
             _logger.warning(f"[{function_name}] No valid subsections remained after validation for page {page_num}. Using fallback.")
             return [create_error_fallback_subsection(page_num, raw_text, "LLM response processed, but no valid subsections remained.")]
    else:
        # Fallback if all initial attempts AND repair failed
        fail_reason = "Initial/repair LLM calls failed or returned invalid structure."
        _logger.error(f"[{function_name}] All subsection extraction attempts failed for page {page_num}. Using fallback. Reason: {fail_reason}")
        return [create_error_fallback_subsection(page_num, raw_text, fail_reason)]


async def merge_cutoff_subsections(pages: List[PageContent]) -> List[PageContent]:
    """Merges subsections marked as cut off with the first subsection of the next page."""
    if len(pages) < 2: return pages

    try:
        # Ensure pages are sorted by page number for correct merging
        sorted_pages = sorted(pages, key=lambda p: p.page_number)
    except Exception as sort_err:
        _logger.error(f"Failed to sort pages for merging: {sort_err}. Returning original order.")
        sorted_pages = pages # Proceed with original order if sorting fails

    _logger.debug(f"Attempting to merge subsections across {len(sorted_pages)} pages.")
    merged_something = False

    # Iterate through pages up to the second-to-last one
    for i in range(len(sorted_pages) - 1):
        current_page = sorted_pages[i]
        next_page = sorted_pages[i + 1]

        # Basic validation
        if not current_page.subsections or not next_page.subsections: continue

        # Check if the last subsection of the current page is cut off
        # Need to ensure subsections are sorted by order within the page first
        current_subs_sorted = sorted(current_page.subsections, key=lambda s: s.order)
        next_subs_sorted = sorted(next_page.subsections, key=lambda s: s.order)

        if not current_subs_sorted or not next_subs_sorted: continue # Skip if sorting failed or list empty

        last_subsection = current_subs_sorted[-1]
        first_next_subsection = next_subs_sorted[0]

        if last_subsection.is_cutoff:
            _logger.info(f"Merging cutoff subsection '{last_subsection.subsection_id}' (Page {current_page.page_number}) with '{first_next_subsection.subsection_id}' (Page {next_page.page_number})")

            # Append text
            last_subsection.text += "\n" + first_next_subsection.text
            # Merge descriptions (simple concat, could be smarter)
            last_subsection.description += " " + first_next_subsection.description
            # Mark the merged subsection as not cut off (unless the *next* one was also cut off)
            last_subsection.is_cutoff = first_next_subsection.is_cutoff
            # Merge references
            last_subsection.referenced_visuals = list(set(last_subsection.referenced_visuals + first_next_subsection.referenced_visuals))
            last_subsection.referenced_tables = list(set(last_subsection.referenced_tables + first_next_subsection.referenced_tables))

            # Remove the first subsection from the next page
            next_page.subsections = next_subs_sorted[1:]
            # Renumber the order of remaining subsections in the next page
            for idx, subsection in enumerate(next_page.subsections):
                 subsection.order = idx + 1

            merged_something = True

    if merged_something: _logger.info("Finished merging cut-off subsections.")
    else: _logger.debug("No cut-off subsections found to merge.")

    return sorted_pages # Return the potentially modified list


async def extract_chapters_from_subsections(client: genai.GenerativeModel, pages: List[PageContent], source_identifier: str) -> List[Chapter]:
    """Groups subsections from multiple pages into logical chapters using an LLM."""
    all_subsections_with_context = []
    subsection_map = {} # subsection_id -> Subsection object

    _logger.debug(f"Extracting subsections from {len(pages)} pages for chapter generation for '{source_identifier}'.")
    for page in pages:
        if not page.subsections: continue
        for subsection in page.subsections:
            context = {
                "subsection_id": subsection.subsection_id,
                "title": subsection.title,
                "description": subsection.description,
                "page_number": page.page_number,
                "order": subsection.order
            }
            all_subsections_with_context.append(context)
            subsection_map[subsection.subsection_id] = subsection

    if not all_subsections_with_context:
        _logger.warning(f"No valid subsections found to form chapters for '{source_identifier}'.")
        return []

    # Sort by page number and then order within page
    sorted_subsections_context = sorted(all_subsections_with_context, key=lambda s: (s["page_number"], s["order"]))
    _logger.info(f"Sending {len(sorted_subsections_context)} subsections to LLM for chapter structuring for '{source_identifier}'.")

    # --- Create prompt ---
    subsection_context_json = "[]"
    try: subsection_context_json = json.dumps(sorted_subsections_context, indent=2)
    except TypeError: _logger.error("Failed to serialize subsection context for chapter prompt."); return []

    chapter_prompt = chapter_extraction_prompt_template.format(
        source_identifier=source_identifier,
        subsection_context_json=subsection_context_json
    )

    chapters_result: List[Chapter] = []
    try:
        response = await retry_api_call(
            client.generate_content,
            [chapter_prompt],
            generation_config=genai_types.GenerationConfig(response_mime_type="application/json")
        )
        json_text = clean_json_response(response.candidates[0].content.parts[0].text)
        chapters_data = json.loads(json_text)

        if not isinstance(chapters_data, list):
             raise ValueError("LLM response for chapters was not a list.")

        processed_subsection_ids = set()
        validated_chapters = []
        for i, chapter_data in enumerate(chapters_data):
             if not isinstance(chapter_data, dict) or not all(k in chapter_data for k in ["title", "summary", "subsection_ids", "order"]):
                 _logger.warning(f"Skipping malformed chapter data: {chapter_data}")
                 continue

             chapter_subsections = []
             subsection_ids_in_chapter = chapter_data.get("subsection_ids", [])
             if not isinstance(subsection_ids_in_chapter, list):
                 _logger.warning(f"Subsection IDs for chapter '{chapter_data.get('title')}' not a list. Skipping chapter.")
                 continue

             valid_chapter = True
             temp_seen_ids = set()
             for subsection_id in subsection_ids_in_chapter:
                 if subsection_id not in subsection_map:
                      _logger.warning(f"Subsection ID '{subsection_id}' from LLM not found in map for chapter '{chapter_data.get('title')}'. Skipping chapter.")
                      valid_chapter = False; break
                 if subsection_id in processed_subsection_ids:
                      _logger.warning(f"Subsection ID '{subsection_id}' reused in chapter '{chapter_data.get('title')}'. Skipping chapter.")
                      valid_chapter = False; break
                 chapter_subsections.append(subsection_map[subsection_id])
                 temp_seen_ids.add(subsection_id)

             if not valid_chapter: continue
             if not chapter_subsections: _logger.warning(f"Chapter '{chapter_data.get('title')}' has no valid subsections. Skipping."); continue

             try:
                 # Ensure 'order' is int
                 chapter_data['order'] = int(chapter_data.get('order', i + 1))
                 # Add chapter_id if missing
                 chapter_data.setdefault('chapter_id', f"chap_{chapter_data['order']}_{uuid.uuid4().hex[:4]}")
                 # Add the actual Subsection objects
                 chapter_data['subsections'] = chapter_subsections
                 # Validate with Pydantic
                 chapter_obj = Chapter(**chapter_data)
                 validated_chapters.append(chapter_obj)
                 processed_subsection_ids.update(temp_seen_ids) # Mark IDs as used
             except (ValidationError, ValueError, TypeError) as ch_val_err:
                  _logger.warning(f"Validation failed for chapter '{chapter_data.get('title')}': {ch_val_err}. Skipping.")

        # Check coverage
        original_ids = set(subsection_map.keys())
        missing_ids = original_ids - processed_subsection_ids
        if missing_ids: _logger.warning(f"LLM failed to assign {len(missing_ids)} subsections to chapters: {missing_ids}")

        chapters_result = sorted(validated_chapters, key=lambda c: c.order)
        if not chapters_result and all_subsections_with_context:
             _logger.warning("LLM processing resulted in zero valid chapters. Creating default chapter.")
             raise ValueError("LLM returned no valid chapters")

    except (json.JSONDecodeError, ValidationError, ValueError) as e:
        _logger.error(f"Error processing chapter LLM response for '{source_identifier}': {e}", exc_info=True)
    except Exception as e:
        _logger.error(f"LLM call failed during chapter extraction for '{source_identifier}': {e}", exc_info=True)

    # Fallback: Create a single default chapter if extraction failed or yielded nothing
    if not chapters_result and all_subsections_with_context:
        _logger.info(f"Creating a single default chapter for '{source_identifier}'.")
        all_original_subsections = sorted(subsection_map.values(), key=lambda s: (
            next((ctx['page_number'] for ctx in sorted_subsections_context if ctx['subsection_id'] == s.subsection_id), 0),
            s.order
        ))
        default_chapter = Chapter(
            chapter_id="chapter_1_default",
            title="Document Content",
            summary="Content grouped into a single chapter due to an error during automatic chapter structuring.",
            subsections=all_original_subsections,
            order=1
        )
        chapters_result = [default_chapter]

    _logger.info(f"Successfully generated {len(chapters_result)} chapters for '{source_identifier}'.")
    return chapters_result


async def analyze_concept_entity_relationships(
    client: genai.GenerativeModel,
    chapter: Chapter, # Takes the Chapter object
    source_identifier: str # Now required
) -> None: # Modifies chapter.qdrant_nodes in-place
    """
    Analyzes text within a single chapter to extract key concepts/entities
    and their relationships using an LLM structured prompt. Populates
    chapter.qdrant_nodes in-place. Includes retries and fallback logic.
    """
    function_name = "analyze_concept_entity_relationships"
    chapter_id = chapter.chapter_id
    chapter_title = chapter.title
    _logger.info(f"[{function_name}] Analyzing relationships for Chapter '{chapter_id}': '{chapter_title}' from '{source_identifier}'")

    if not chapter.subsections:
        _logger.warning(f"[{function_name}] Chapter '{chapter_id}' has no subsections. Skipping analysis.")
        chapter.qdrant_nodes = []
        return

    # --- Prepare Input Text ---
    combined_text = ""
    subsection_texts = []
    for subsection in chapter.subsections:
        sub_text = getattr(subsection, 'text', '')
        if sub_text and sub_text.strip():
             subsection_texts.append(
                 f"--- Subsection Start (ID: {getattr(subsection, 'subsection_id', 'unknown')}) ---\n"
                 f"Title: {getattr(subsection, 'title', 'Untitled')}\n"
                 f"Description: {getattr(subsection, 'description', '')}\n\n"
                 f"{sub_text.strip()}\n"
                 f"--- Subsection End ---"
             )
    combined_text = "\n\n".join(subsection_texts)

    MAX_TEXT_LIMIT = 150000 # Adjust if needed
    if len(combined_text) > MAX_TEXT_LIMIT:
         _logger.warning(f"[{function_name}] Truncating combined text for chapter '{chapter_id}' ({len(combined_text)} -> {MAX_TEXT_LIMIT} chars).")
         combined_text = combined_text[:MAX_TEXT_LIMIT] + "\n... [Content Truncated]"
    elif not combined_text.strip():
        _logger.warning(f"[{function_name}] Combined text for chapter '{chapter_id}' is empty. Skipping analysis.")
        chapter.qdrant_nodes = []
        return

    # Extract chapter number for node ID generation
    chapter_num_str = "1"; match = re.search(r'[_-](\d+)$', chapter_id)
    if match: chapter_num_str = match.group(1)

    # --- Define LLM Prompt ---
    entity_prompt = node_relationship_prompt_template.format(
        chapter_title=chapter_title,
        chapter_id=chapter_id,
        source_identifier=source_identifier,
        combined_text=combined_text,
        chapter_num_str=chapter_num_str
    )

    # --- Call LLM and Process Response (with Fallback) ---
    max_structure_retries = 2; retry_delay = 3
    final_data = None
    last_raw_response_text = "N/A"; last_cleaned_response_text = "N/A"; last_parsed_data = None
    initial_call_failed_parsing = False

    for attempt in range(max_structure_retries):
        _logger.debug(f"[{function_name}] LLM Call Attempt {attempt + 1}/{max_structure_retries} for chapter '{chapter_id}'")
        try:
            response = await retry_api_call(
                client.generate_content,
                [entity_prompt],
                generation_config=genai_types.GenerationConfig(response_mime_type="application/json")
            )
            if not response or not response.candidates: raise ValueError("LLM response had no candidates")

            last_raw_response_text = response.candidates[0].content.parts[0].text
            json_text = clean_json_response(last_raw_response_text, extract_text_on_failure=False) # Don't use fallback here yet
            last_cleaned_response_text = json_text
            _logger.debug(f"[{function_name}] Cleaned LLM response (Attempt {attempt + 1}, Chapter {chapter_id}): {json_text[:500]}...")

            try: data = json.loads(json_text); last_parsed_data = data
            except json.JSONDecodeError as json_err:
                _logger.warning(f"[{function_name}] JSONDecodeError attempt {attempt + 1} for chapter '{chapter_id}': {json_err}.")
                # Store faulty text for potential repair
                faulty_json_text_to_repair = last_raw_response_text
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: initial_call_failed_parsing = True; break # Failed parsing after retries

            # Structure Check (expecting {"qdrant_nodes": [...]})
            if isinstance(data, dict) and "qdrant_nodes" in data and isinstance(data.get("qdrant_nodes"), list):
                _logger.info(f"[{function_name}] Received correct structure attempt {attempt + 1} for chapter '{chapter_id}'.")
                final_data = data; initial_call_failed_parsing = False; break
            else:
                _logger.warning(f"[{function_name}] Incorrect JSON structure attempt {attempt+1} for chapter '{chapter_id}'. Got: {type(data)}. Structure: {str(data)[:300]}...")
                faulty_json_text_to_repair = last_raw_response_text # Store faulty text
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: initial_call_failed_parsing = True; break # Failed structure check after retries

        except Exception as e:
            _logger.error(f"[{function_name}] Error during LLM call attempt {attempt + 1} for chapter '{chapter_id}': {e}", exc_info=True)
            faulty_json_text_to_repair = f"LLM call failed: {e}" # Store error info
            if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else: initial_call_failed_parsing = True; break # Call failed after retries

    # --- JSON Repair Attempt (if initial parsing failed) ---
    if initial_call_failed_parsing and faulty_json_text_to_repair:
        _logger.info(f"[{function_name}] Attempting JSON repair for chapter '{chapter_id}'.")
        try:
            repair_prompt = json_repair_prompt_template.format(faulty_json_text=faulty_json_text_to_repair)
            repair_response = await retry_api_call(
                client.generate_content, [repair_prompt],
                generation_config=genai_types.GenerationConfig(response_mime_type="application/json", temperature=0.1),
                max_retries=1
            )
            if repair_response and repair_response.candidates:
                repaired_json_text = clean_json_response(repair_response.candidates[0].content.parts[0].text, extract_text_on_failure=False)
                _logger.debug(f"[{function_name}] Repaired JSON Response (Chapter {chapter_id}): {repaired_json_text[:500]}...")
                try:
                    repaired_data = json.loads(repaired_json_text)
                    if isinstance(repaired_data, dict) and "qdrant_nodes" in repaired_data and isinstance(repaired_data.get("qdrant_nodes"), list):
                        _logger.info(f"[{function_name}] Fallback repair SUCCEEDED for chapter '{chapter_id}'.")
                        final_data = repaired_data
                    else: _logger.error(f"[{function_name}] Repaired JSON structure still incorrect for chapter '{chapter_id}'.")
                except json.JSONDecodeError as repair_err: _logger.error(f"[{function_name}] Failed to decode JSON even after repair for chapter '{chapter_id}': {repair_err}")
            else: _logger.error(f"[{function_name}] Repair LLM call returned no candidates for chapter '{chapter_id}'.")
        except Exception as fallback_err: _logger.error(f"[{function_name}] Error during fallback/repair LLM call for chapter '{chapter_id}': {fallback_err}", exc_info=True)

    # --- Process Final Data ---
    chapter.qdrant_nodes = [] # Initialize/reset
    if final_data is not None:
        qdrant_nodes_data = final_data.get("qdrant_nodes")
        if isinstance(qdrant_nodes_data, list):
            validated_nodes: List[QdrantNode] = []
            temp_node_map_by_llm_id: Dict[str, str] = {} # Map LLM ID -> generated UUID
            nodes_to_process_later: List[Dict] = [] # Store raw dicts for link resolution

            current_iso_ts = datetime.datetime.now(datetime.timezone.utc).isoformat()

            # --- First Pass: Prepare nodes, generate UUIDs ---
            for node_dict in qdrant_nodes_data:
                if not isinstance(node_dict, dict): continue
                try:
                    # Validate required fields from LLM prompt
                    if not all(k in node_dict for k in ["node_type", "name", "chapter_id_context", "title", "core_claim", "information_type", "key_entities", "tags", "sentiment_tone", "source_credibility"]):
                         _logger.warning(f"Node dict missing required fields generated by LLM. Skipping. Data: {str(node_dict)[:200]}")
                         continue

                    # Generate deterministic UUID based on type, name, context
                    node_type = node_dict["node_type"]; name = node_dict["name"]; chapter_ctx = node_dict["chapter_id_context"]
                    node_uuid = generate_qdrant_id(node_type, name, chapter_ctx)
                    original_llm_node_id = node_dict.get('node_id') # Get the ID generated by LLM (type_name_chX)
                    if original_llm_node_id: temp_node_map_by_llm_id[original_llm_node_id] = node_uuid

                    # Prepare the dictionary for QdrantNode validation
                    node_data_for_pydantic = node_dict.copy()
                    node_data_for_pydantic['node_id'] = node_uuid # Assign the generated UUID
                    node_data_for_pydantic['source'] = source_identifier
                    node_data_for_pydantic['stored_at'] = current_iso_ts
                    node_data_for_pydantic['original_node_id'] = original_llm_node_id

                    # Ensure linked_nodes is a list (initially contains dicts from LLM)
                    node_data_for_pydantic.setdefault('linked_nodes', [])
                    if not isinstance(node_data_for_pydantic['linked_nodes'], list):
                        node_data_for_pydantic['linked_nodes'] = []

                    nodes_to_process_later.append(node_data_for_pydantic)

                except Exception as node_prep_err:
                     _logger.error(f"Error preparing node during first pass: {node_prep_err}. Data: {str(node_dict)[:200]}", exc_info=True)

            # --- Second Pass: Resolve links and validate QdrantNode ---
            final_node_ids_created = set()
            for node_dict_prepared in nodes_to_process_later:
                node_uuid = node_dict_prepared['node_id']
                raw_links = node_dict_prepared.get('linked_nodes', []) # List[Dict] from LLM
                resolved_link_models: List[LinkedNode] = [] # Store validated Pydantic models

                if isinstance(raw_links, list):
                     for link_dict in raw_links:
                         if not isinstance(link_dict, dict): continue
                         target_original_id = link_dict.get("target_node_id") # LLM returns original format ID
                         target_uuid = temp_node_map_by_llm_id.get(target_original_id) if target_original_id else None

                         if target_uuid: # If lookup successful
                             try:
                                 # Create LinkedNode model (validation happens here)
                                 resolved_link = LinkedNode(
                                     target_node_id=target_uuid, # Use the resolved UUID
                                     relationship_description=link_dict.get('relationship_description', '?'),
                                     relationship_keywords=link_dict.get('relationship_keywords', []),
                                     relationship_strength=link_dict.get('relationship_strength', 5.0)
                                 )
                                 resolved_link_models.append(resolved_link)
                             except (ValidationError, ValueError, TypeError) as link_val_err:
                                  _logger.warning(f"Failed to create/validate LinkedNode for target UUID '{target_uuid}' in node '{node_uuid}': {link_val_err}")
                         else:
                             _logger.warning(f"Could not resolve target UUID for link target '{target_original_id}' in node '{node_uuid}'. Dropping link.")

                # Update the node dict with the resolved LIST OF LINKEDNODE MODELS
                node_dict_prepared['linked_nodes'] = resolved_link_models

                # Final validation of the entire node dict against QdrantNode model
                try:
                    node_obj = QdrantNode(**node_dict_prepared)
                    validated_nodes.append(node_obj)
                    final_node_ids_created.add(node_obj.node_id)
                except ValidationError as node_error:
                    _logger.warning(f"Pydantic validation failed for node UUID '{node_uuid}' after link resolution: {node_error}. Skipping node. Data: {json.dumps(node_dict_prepared, indent=2, default=str)}")
                except Exception as final_val_err:
                    _logger.error(f"Unexpected error during final node validation UUID '{node_uuid}': {final_val_err}", exc_info=True)

            # Assign the list of validated Pydantic objects
            chapter.qdrant_nodes = validated_nodes
            _logger.info(f"Assigned {len(chapter.qdrant_nodes)} validated nodes (using UUIDs) for chapter '{chapter_id}'.")

        else: # qdrant_nodes data was not a list
            _logger.error(f"Final data for chapter '{chapter_id}' had 'qdrant_nodes' but not a list.")
    else: # Initial or fallback LLM call failed completely
        _logger.error(f"All node extraction attempts failed for chapter '{chapter_id}'. Assigning empty list.")

    # Function modifies chapter in-place, no return value needed.


async def analyze_inter_chapter_relationships(
    client: genai.GenerativeModel, # Use model instance
    chapters: List[Chapter],
    source_identifier: str
) -> Optional[Dict]: # Returns DocumentSummaryDetails structure as dict
    """
    Analyzes inter-chapter relationships AND generates a document summary based on aggregated nodes.
    Updates the `linked_nodes` attribute of the QdrantNode objects within the chapters list *in-place*.
    Returns a dictionary conforming to DocumentSummaryDetails or None on failure.
    """
    function_name = "analyze_inter_chapter_relationships"
    _logger.info(f"[{function_name}] Starting analysis for {len(chapters)} chapters of '{source_identifier}'.")

    # --- Aggregate Node Information & Create Lookup Map ---
    all_nodes_for_prompt = []
    node_map: Dict[str, QdrantNode] = {} # node_id (UUID) -> QdrantNode object

    for chapter in chapters:
        if not hasattr(chapter, 'qdrant_nodes') or not isinstance(chapter.qdrant_nodes, list):
            _logger.warning(f"[{function_name}] Chapter '{getattr(chapter, 'chapter_id', '?')}' missing 'qdrant_nodes'. Skipping.")
            continue
        for node in chapter.qdrant_nodes:
             if not isinstance(node, QdrantNode) or not node.node_id: continue

             node_description = getattr(node, 'core_claim', 'No description.')
             node_tags = getattr(node, 'tags', '[]')

             nodes_for_prompt_entry = {
                 "node_id": node.node_id, # Use the UUID
                 "node_type": getattr(node, 'node_type', 'unknown'),
                 "name": getattr(node, 'name', 'Unknown Name'),
                 "description": node_description,
                 "chapter_id_context": getattr(node, 'chapter_id_context', 'unknown_context'),
                 "tags": node_tags
             }
             all_nodes_for_prompt.append(nodes_for_prompt_entry)
             if node.node_id in node_map: _logger.warning(f"[{function_name}] Duplicate node_id '{node.node_id}'. Overwriting entry.")
             node_map[node.node_id] = node

    if not all_nodes_for_prompt:
        _logger.warning(f"[{function_name}] No valid nodes found across chapters for '{source_identifier}'. Cannot analyze/summarize.")
        return None

    _logger.info(f"[{function_name}] Aggregated {len(all_nodes_for_prompt)} nodes from {len(chapters)} chapters for '{source_identifier}'.")

    nodes_json_for_prompt = "[]"; MAX_PROMPT_CHARS = 150000 # Limit prompt size
    try:
        temp_json = json.dumps(all_nodes_for_prompt, indent=2)
        if len(temp_json) > MAX_PROMPT_CHARS:
             _logger.warning(f"Truncating node context for prompt ({len(temp_json)} -> {MAX_PROMPT_CHARS} chars).")
             # Simple truncation, could be smarter
             nodes_json_for_prompt = temp_json[:MAX_PROMPT_CHARS] + "... ]"
        else:
             nodes_json_for_prompt = temp_json
    except TypeError as json_err:
         _logger.error(f"[{function_name}] Failed to serialize nodes for prompt: {json_err}"); return None

    # --- Define the LLM Prompt ---
    inter_chapter_prompt = inter_chapter_summary_prompt_template.format(
        source_identifier=source_identifier,
        nodes_json_for_prompt=nodes_json_for_prompt
    )

    # --- Call LLM and Process Response (with Retries/Fallback implicitly handled by retry_api_call) ---
    final_llm_output = None
    try:
        response = await retry_api_call(
            client.generate_content,
            [inter_chapter_prompt],
            generation_config=genai_types.GenerationConfig(response_mime_type="application/json")
        )
        if not response or not response.candidates: raise ValueError("LLM response had no candidates")

        json_text = clean_json_response(response.candidates[0].content.parts[0].text)
        final_llm_output = json.loads(json_text)

        # --- Structure Check ---
        if not (isinstance(final_llm_output, dict) and
                "inter_chapter_links" in final_llm_output and isinstance(final_llm_output.get("inter_chapter_links"), list) and
                "document_summary_details" in final_llm_output and isinstance(final_llm_output.get("document_summary_details"), dict)):
            _logger.warning(f"[{function_name}] Incorrect JSON structure received from LLM for '{source_identifier}'.")
            # Attempt to use fallback JSON from clean_json_response if it exists
            if isinstance(final_llm_output, dict) and final_llm_output.get("error"):
                 _logger.warning(f"Using fallback error JSON due to parsing failure.")
            else:
                 raise ValueError(f"Incorrect JSON structure: {str(final_llm_output)[:300]}...")

    except (json.JSONDecodeError, ValueError) as e:
         _logger.error(f"[{function_name}] Failed to process/parse LLM response for inter-chapter analysis '{source_identifier}': {e}", exc_info=True)
         final_llm_output = {"error": str(e), "inter_chapter_links": [], "document_summary_details": None} # Ensure keys exist for logic below
    except Exception as e:
         _logger.error(f"[{function_name}] LLM call failed for inter-chapter analysis '{source_identifier}': {e}", exc_info=True)
         final_llm_output = {"error": str(e), "inter_chapter_links": [], "document_summary_details": None}


    # --- Process LLM Output ---
    document_summary_output = None

    # Process Inter-Chapter Links (modify node_map in place)
    newly_found_links_data = final_llm_output.get("inter_chapter_links", [])
    if isinstance(newly_found_links_data, list):
        links_added_count, links_skipped_count = _merge_links_into_nodes(newly_found_links_data, node_map)
        _logger.info(f"[{function_name}] Merged inter-chapter links for '{source_identifier}'. Added: {links_added_count}, Skipped: {links_skipped_count}")
    else:
        _logger.warning(f"[{function_name}] Key 'inter_chapter_links' missing or not a list in LLM output for '{source_identifier}'.")

    # Process Document Summary Details
    summary_details_data = final_llm_output.get("document_summary_details")
    if isinstance(summary_details_data, dict):
         try:
             # Validate against the Pydantic model
             # Need to add chapters back for validation/consistency if needed by model
             summary_details_data['chapters'] = chapters # Add the list of Chapter objects
             validated_summary = DocumentSummaryDetails(**summary_details_data)
             document_summary_output = validated_summary.model_dump(mode='json') # Store as dict
             _logger.info(f"[{function_name}] Successfully extracted and validated document summary details for '{source_identifier}'.")
         except ValidationError as summary_error:
             _logger.error(f"[{function_name}] Validation failed for document_summary_details from LLM for '{source_identifier}': {summary_error}")
             document_summary_output = {
                "title": f"Summary Validation Failed: {source_identifier}",
                "themes": ["error"], "questions": [],
                "summary": f"LLM returned summary data, but it failed validation: {summary_error}",
                "chapters": [ch.model_dump(mode='json') for ch in chapters] # Include chapter data in error dict
             }
         except Exception as unexpected_summary_err:
             _logger.error(f"[{function_name}] Unexpected error processing summary details for '{source_identifier}': {unexpected_summary_err}", exc_info=True)
             document_summary_output = {
                 "title": f"Summary Processing Error: {source_identifier}",
                 "themes": ["error"], "questions": [],
                 "summary": f"An unexpected error occurred during summary processing: {unexpected_summary_err}",
                 "chapters": [ch.model_dump(mode='json') for ch in chapters]
             }
    else:
        _logger.warning(f"[{function_name}] Key 'document_summary_details' missing or not a dict in LLM output for '{source_identifier}'.")
        document_summary_output = {
            "title": f"Summary Missing: {source_identifier}", "themes": [], "questions": [],
            "summary": "LLM did not provide document summary details.",
            "chapters": [ch.model_dump(mode='json') for ch in chapters]
        }

    # The chapters list (containing QdrantNode objects, accessible via node_map) has been modified in-place
    _logger.info(f"[{function_name}] Finished inter-chapter analysis for '{source_identifier}'. Returning summary details.")
    return document_summary_output # Returns summary dict

# Helper for merging links (modifies node_map)
def _merge_links_into_nodes(link_data_list: List[Dict], node_map: Dict[str, QdrantNode]) -> Tuple[int, int]:
    """Merges link data into the linked_nodes of QdrantNode objects in node_map."""
    links_added_count = 0; links_skipped_count = 0
    MAX_OUTGOING_LINKS_PER_NODE = 10 # Limit links per node

    if not isinstance(link_data_list, list): return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict): links_skipped_count += 1; continue

        source_id = link_data.get("source_node_id") # Expecting UUID
        target_id = link_data.get("target_node_id") # Expecting UUID
        desc = link_data.get("relationship_description")

        if not all([source_id, target_id, desc]): links_skipped_count += 1; continue

        source_node = node_map.get(source_id)
        target_node = node_map.get(target_id)

        if not source_node or not target_node:
            _logger.debug(f"[_merge_links] Source ('{source_id}') or Target ('{target_id}') node not found. Skipping link.")
            links_skipped_count += 1; continue

        # Ensure it's an INTER-chapter link (or handle doc/meta nodes)
        if source_node.chapter_id_context == target_node.chapter_id_context and not source_node.chapter_id_context.startswith("doc_level"):
             _logger.debug(f"[_merge_links] Skipping intra-chapter link: {source_id} -> {target_id}")
             links_skipped_count += 1; continue

        # Add Forward Link
        if source_node.linked_nodes is None: source_node.linked_nodes = []
        # Check if link already exists and respect limit
        if (len(source_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and
            not any(link.target_node_id == target_id for link in source_node.linked_nodes if hasattr(link, 'target_node_id'))):
            try:
                # Create LinkedNode directly from dict if keys match
                # Strength validation happens in LinkedNode model
                forward_link = LinkedNode(**link_data)
                # Update target_node_id again just in case (should be redundant)
                forward_link.target_node_id = target_id
                source_node.linked_nodes.append(forward_link)
                links_added_count += 1
            except (ValidationError, TypeError) as ve:
                 _logger.warning(f"Validation failed creating forward LinkedNode {source_id} -> {target_id}: {ve}")
                 links_skipped_count += 1
        else:
             links_skipped_count += 1 # Skipped due to limit or duplicate

    return links_added_count, links_skipped_count


# --- Project Ontology Generation (Adapted from main.py) ---
# NOTE: This uses a separate LLM call and assumes document processing is complete.
# It now works with the QdrantNode structure directly.

# Helper function for consistent ID generation
def generate_qdrant_id(node_type: str, name: str, context_id: str) -> str:
    """Generates a deterministic UUID v5 for Qdrant based on node info."""
    unique_string = f"{node_type}:{name}:{context_id}".lower().strip()
    # Generate UUID using NAMESPACE_DNS (ensures same input -> same output)
    return str(uuid.uuid5(uuid.NAMESPACE_DNS, unique_string))

async def incrementally_generate_doc_level_ontology(
    client: genai.GenerativeModel,
    documents: List[Dict], # Expects list of dicts from finalize_document results
    *,
    ontology_path: Union[str, pathlib.Path] = "doc_level_ontology.json",
    char_budget: int = 100_000,
    max_new_links: int = 50,
    branch_size: int = 20,
    llm_retries: int = 2,
    llm_delay: int = 5
) -> Optional[ProjectOntology]:
    """
    Builds or incrementally updates project ontology using document results.
    Outputs nodes with QdrantNode structure.
    Uses internal tree summarizer if needed.
    """
    log = logging.getLogger("incremental_doc_ontology")
    path = pathlib.Path(ontology_path)
    CHAR_BUDGET = char_budget; MAX_NEW_LINKS = max_new_links; BRANCH_SIZE = branch_size
    LLM_RETRIES = llm_retries; LLM_DELAY = llm_delay

    # --- Nested Tree Summarizer ---
    async def _tree_summarise_nested(input_docs: List[Dict], *, target_chars: int) -> List[Dict]:
        """NESTED: Recursively compress summaries until JSON size <= target_chars."""
        tree_log = logging.getLogger("tree_summarise_nested")
        def blob_size(nodes_list: List[Dict]) -> int:
            try: return len(json.dumps(nodes_list, ensure_ascii=False))
            except TypeError: return float('inf')

        current_size = blob_size(input_docs); tree_log.debug(f"Input: {len(input_docs)} nodes, Size: {current_size}, Target: {target_chars}")
        if current_size <= target_chars or len(input_docs) <= 1: return input_docs

        groups = [input_docs[i:i + BRANCH_SIZE] for i in range(0, len(input_docs), BRANCH_SIZE)]
        tree_log.info(f"Compressing {len(input_docs)} nodes into {len(groups)} groups (branch size {BRANCH_SIZE}).")
        compressed_results: List[Dict] = []; tasks = []

        async def call_gemini_for_summary(prompt: str, group_idx: int) -> Optional[Dict]:
            """Inner function for single API call."""
            for attempt in range(LLM_RETRIES):
                try:
                    tree_log.debug(f"Calling Gemini for meta-summary group {group_idx} (Attempt {attempt+1})")
                    response = await retry_api_call(
                        client.generate_content, [prompt],
                        generation_config=genai_types.GenerationConfig(response_mime_type="application/json"),
                        max_retries=1 # Retry handled externally by retry_api_call
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text
                    cleaned_json = clean_json_response(raw_text)
                    parsed_json = json.loads(cleaned_json)
                    # Expecting structure {"meta_node": {...}}
                    if isinstance(parsed_json, dict) and "meta_node" in parsed_json and isinstance(parsed_json["meta_node"], dict):
                        return parsed_json["meta_node"] # Return just the node dict
                    else:
                        tree_log.warning(f"Unexpected meta-summary structure group {group_idx}: {parsed_json}")
                        return None
                except Exception as e:
                    tree_log.warning(f"Meta-summary Gemini call/parse failed group {group_idx} (Attempt {attempt+1}): {e}")
                    if attempt == LLM_RETRIES - 1: tree_log.error(f"Meta-summary failed group {group_idx} after retries."); return None
                    await asyncio.sleep(LLM_DELAY * (attempt + 1))
            return None

        for idx, chunk in enumerate(groups, 1):
            if not chunk: continue
            chunk_for_prompt = [{"node_id": d.get("node_id"), "document_name": d.get("document_name"), "executive_summary_preview": d.get("executive_summary","")[:300]+"..."} for d in chunk]
            chunk_json = json.dumps(chunk_for_prompt, indent=2); MAX_CHUNK_JSON = int(CHAR_BUDGET * 0.8)
            if len(chunk_json) > MAX_CHUNK_JSON: chunk_json = chunk_json[:MAX_CHUNK_JSON] + "\n... ]"

            prompt = f"""Synthesize the key ideas from the following {len(chunk)} document/meta summaries into **one single, concise meta-summary node**.
                Input Summaries (Previews): ```json\n{chunk_json}\n```
                Goal: Create a meta-summary representing this group.
                Instructions: Generate an `executive_summary` (<800 words). Assign `node_id` "meta::group_{idx}", `document_name` "meta_group_{idx}".
                Output Format: **Valid JSON ONLY**: `{{"meta_node": {{"node_id": "meta::group_{idx}", "document_name": "meta_group_{idx}", "executive_summary": "<summary>"}}}}`
                Generate JSON:"""
            tasks.append(asyncio.create_task(call_gemini_for_summary(prompt, idx)))

        task_results = await asyncio.gather(*tasks)
        for i, result in enumerate(task_results, 1):
            if result and isinstance(result, dict): compressed_results.append(result)
            else: tree_log.error(f"Skipping group {i} due to API/JSON failure in meta-summary generation.")

        if not compressed_results: tree_log.error("All summarization tasks failed."); return input_docs # Fallback
        return await _tree_summarise_nested(compressed_results, target_chars=target_chars) # Recurse

    # --- Main Ontology Logic ---
    ontology = load_ontology_json(path) # Uses external helper
    first_run = ontology is None
    if first_run:
        log.info(f"No ontology file found at '{path}' – creating new.")
        ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
    else:
        if not isinstance(ontology, ProjectOntology):
            log.error(f"Loaded data from '{path}' invalid. Starting fresh."); first_run = True
            ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
        else: log.info(f"Loaded existing ontology with {ontology.document_count} documents from '{path}'.")

    existing_files = set(ontology.documents or [])
    target_docs_data = [] # List of result dicts from finalize_document
    new_doc_count = 0
    for d_result in documents: # documents is List[Dict] from finalize_document
        try:
            filename = d_result.get("raw_extracted_content", {}).get("filename")
            processing_error = d_result.get("raw_extracted_content", {}).get("error")
            # Process only new, successfully processed documents
            if filename and filename not in existing_files and not processing_error:
                target_docs_data.append(d_result)
                new_doc_count += 1
        except Exception as e: log.warning(f"Error checking document result: {e}")
    log.info(f"Incremental run: Found {new_doc_count} new, successfully processed documents.")
    if not target_docs_data: log.info("No new valid documents found for ontology update."); return ontology

    # Prepare nodes from NEW document summaries for LLM input
    doc_nodes_for_llm: List[Dict] = []
    processed_filenames = []
    for d_result in target_docs_data:
        try:
            raw = d_result.get("raw_extracted_content", {}); fname = raw.get("filename")
            summ_block = raw.get("summary") # This is a dict from DocumentSummaryDetails
            summary_text = ""
            if isinstance(summ_block, dict): summary_text = summ_block.get("summary", "")
            if not fname or not summary_text: log.warning(f"Doc '{fname or 'unknown'}' missing filename or summary text. Skipping."); continue

            doc_nodes_for_llm.append({"node_id": f"doc::{fname}", "document_name": fname, "executive_summary": summary_text})
            processed_filenames.append(fname)
        except Exception as e: log.warning(f"Error processing doc result data '{raw.get('filename', 'unknown')}': {e}")
    if not doc_nodes_for_llm: log.error("No usable summaries found in new documents."); return ontology

    # Apply Tree Summarization if needed
    final_doc_nodes_for_llm = doc_nodes_for_llm; current_char_size = len(json.dumps(doc_nodes_for_llm))
    if current_char_size > CHAR_BUDGET:
        log.info(f"Summary size ({current_char_size:,} chars) > budget ({CHAR_BUDGET:,}). Summarizing...")
        final_doc_nodes_for_llm = await _tree_summarise_nested(doc_nodes_for_llm, target_chars=CHAR_BUDGET)
        final_size = len(json.dumps(final_doc_nodes_for_llm))
        log.info(f"Summarization complete. Final size: {final_size:,} chars ({len(final_doc_nodes_for_llm)} nodes).")

    # Prepare final JSON for LLM prompts
    final_doc_nodes_json_str = "[]"; max_json_chars_in_prompt = int(CHAR_BUDGET * 0.85)
    try:
        temp_json = json.dumps(final_doc_nodes_for_llm, indent=2)
        final_doc_nodes_json_str = temp_json[:max_json_chars_in_prompt] + "... ]" if len(temp_json) > max_json_chars_in_prompt else temp_json
    except Exception as e: log.error(f"Failed to serialize final doc nodes for prompt: {e}"); final_doc_nodes_json_str = "[]"


    # --- Call LLM for Ontology Generation / Update ---
    llm_links = []; llm_summary_details = {}; llm_call_failed = False
    try:
        if first_run:
            log.info("Performing initial ontology generation LLM call...")
            init_prompt = project_ontology_prompt_template.format(doc_nodes_json=final_doc_nodes_json_str)
            response = await retry_api_call(client.generate_content, [init_prompt], generation_config=genai_types.GenerationConfig(response_mime_type="application/json"), max_retries=LLM_RETRIES)
            if not response or not response.candidates: raise ValueError("LLM call failed (no candidates)")
            llm_out = json.loads(clean_json_response(response.candidates[0].content.parts[0].text))
            llm_links = llm_out.get("inter_document_links", [])
            llm_summary_details = llm_out.get("project_summary_details", {})
            log.info("Initial ontology LLM call successful.")
        else: # Incremental update
            log.info("Performing incremental ontology update LLM call...")
            # Sample existing nodes (if any) for context - use QdrantNode structure
            existing_nodes_for_context = []
            if ontology.project_graph_nodes:
                 sample_size = min(len(ontology.project_graph_nodes), max(len(final_doc_nodes_for_llm)*2, 30))
                 sampled_nodes_pydantic = random.sample(ontology.project_graph_nodes, min(sample_size, len(ontology.project_graph_nodes)))
                 max_summary_preview=800
                 for node_pydantic in sampled_nodes_pydantic:
                     summary_text = node_pydantic.executive_summary or node_pydantic.core_claim or ""
                     existing_nodes_for_context.append({
                         "node_id": node_pydantic.node_id, # Use the existing UUID
                         "document_name": node_pydantic.name, # Use name (might be doc name or concept name)
                         "executive_summary_preview": summary_text[:max_summary_preview] + ("..." if len(summary_text) > max_summary_preview else "")
                     })

            existing_context_json_str = "[]"; max_context_json_chars = int(CHAR_BUDGET * 0.4)
            try:
                temp_json = json.dumps(existing_nodes_for_context, indent=2)
                existing_context_json_str = temp_json[:max_context_json_chars] + "... ]" if len(temp_json) > max_context_json_chars else temp_json
            except Exception as e: log.error(f"Failed to serialize existing node context: {e}")

            # Link Discovery Prompt
            link_prompt = f"""Act as KG Analyst. Update ontology. Find links involving NEW docs/summaries.
                Input 1 (New Docs/Summaries): ```json\n{final_doc_nodes_json_str}\n```
                Input 2 (Existing Node Sample - UUIDs used): ```json\n{existing_context_json_str}\n```
                Task: Find links connecting New nodes (Input 1, IDs like 'doc::fname' or 'meta::id') to other New nodes OR Existing nodes (Input 2, UUIDs). Use the `node_id` provided for linking. Base relationships on `executive_summary` content. Limit: {MAX_NEW_LINKS} links.
                Output (JSON ONLY): `{{"new_inter_document_links": [ {{ "source_node_id": "...", "target_node_id": "...", "relationship_description": "...", "relationship_keywords": ["..."], "relationship_strength": 7.5 }} ]}}`
                Rules: Valid JSON. Use provided `node_id`s. At least one node in link MUST be from Input 1. Empty list [] if none.
                Generate JSON:"""
            try:
                response = await retry_api_call(client.generate_content, [link_prompt], generation_config=genai_types.GenerationConfig(response_mime_type="application/json"), max_retries=LLM_RETRIES)
                if response and response.candidates:
                     link_json = json.loads(clean_json_response(response.candidates[0].content.parts[0].text))
                     llm_links = link_json.get("new_inter_document_links", [])
                     log.info(f"Incremental link discovery call successful. Found {len(llm_links)} potential links.")
                else: log.warning("Incremental link discovery call failed or yielded no candidates.")
            except Exception as e: log.warning(f"Incremental link discovery failed: {e}")

            # Summary Revision Prompt
            current_summary = {"title": ontology.title, "overview": ontology.overview, "global_themes": ontology.global_themes, "key_concepts": ontology.key_concepts}
            current_summary_json_str = json.dumps(current_summary, indent=2)
            rev_prompt = f"""Act as Content Strategist. Update project summary with new docs/summaries.
                Input 1 (Current Summary): ```json\n{current_summary_json_str}\n```
                Input 2 (New Docs/Summaries): ```json\n{final_doc_nodes_json_str}\n```
                Task: Revise title, overview, global_themes, key_concepts. Integrate info from Input 2. `key_concepts` should list important document topics or meta-summary themes.
                Output (JSON ONLY): `{{"revised_project_summary_details": {{ "title": "...", "overview": "...", "global_themes": ["..."], "key_concepts": ["..."] }}}}`
                Rules: Valid JSON. Integrate new info coherently.
                Generate JSON:"""
            try:
                response = await retry_api_call(client.generate_content, [rev_prompt], generation_config=genai_types.GenerationConfig(response_mime_type="application/json"), max_retries=LLM_RETRIES)
                if response and response.candidates:
                    summary_json = json.loads(clean_json_response(response.candidates[0].content.parts[0].text))
                    llm_summary_details = summary_json.get("revised_project_summary_details", {})
                    log.info("Summary revision call successful.")
                else: log.warning("Summary revision call failed or yielded no candidates."); llm_summary_details = current_summary # Fallback
            except Exception as e: log.warning(f"Summary revision failed: {e}"); llm_summary_details = current_summary # Fallback

    except Exception as e: log.error(f"Critical error during ontology LLM processing: {e}", exc_info=True); llm_call_failed = True

    if llm_call_failed and first_run: log.error("Initial LLM calls failed critically. Aborting ontology generation."); return None
    elif llm_call_failed: log.warning("One+ incremental LLM calls failed. Proceeding with partial ontology update.")

    # --- Merge LLM Results into Ontology ---
    # node_map holds UUID -> QdrantNode object
    node_map: Dict[str, QdrantNode] = {node.node_id: node for node in (ontology.project_graph_nodes or [])}
    new_node_uuid_map: Dict[str, str] = {} # Map original doc/meta ID -> generated UUID
    current_iso_ts = datetime.datetime.now(datetime.timezone.utc).isoformat()

    # Process NEW doc/meta summary nodes from LLM input (final_doc_nodes_for_llm)
    for new_node_data in final_doc_nodes_for_llm:
         original_node_id = new_node_data.get("node_id") # e.g., "doc::fname" or "meta::id"
         doc_name = new_node_data.get("document_name")
         exec_summary = new_node_data.get("executive_summary")
         if not original_node_id or not doc_name or exec_summary is None: continue

         node_type = "document" if original_node_id.startswith("doc::") else "meta_summary"
         context_for_id = original_node_id
         node_uuid = generate_qdrant_id(node_type, doc_name, context_for_id)
         new_node_uuid_map[original_node_id] = node_uuid

         # Create or update QdrantNode object in the map
         if node_uuid not in node_map:
             try:
                 node_map[node_uuid] = QdrantNode(
                     node_id=node_uuid, node_type=node_type, name=doc_name,
                     chapter_id_context=f"doc_level::{doc_name}", # Context for doc/meta nodes
                     title=f"Summary: {doc_name}",
                     core_claim=exec_summary[:800] + ("..." if len(exec_summary) > 800 else ""),
                     information_type="Summary",
                     key_entities="[]", tags="[]", sentiment_tone="Neutral",
                     source_credibility="Generated Summary" if node_type == "meta_summary" else "Document Summary",
                     source=doc_name, stored_at=current_iso_ts,
                     original_node_id=original_node_id, executive_summary=exec_summary,
                     linked_nodes=[]
                 )
             except ValidationError as ve:
                 log.error(f"Validation Error creating new doc/meta node {node_uuid}: {ve}")
         else: # Update existing node (e.g., if meta-summary ID collides)
              node_map[node_uuid].executive_summary = exec_summary
              node_map[node_uuid].stored_at = current_iso_ts


    # Resolve and Merge LLM Links into the QdrantNode objects within node_map
    resolved_llm_links = []
    if isinstance(llm_links, list):
        for link_dict in llm_links:
             source_original_id = link_dict.get("source_node_id")
             target_original_id = link_dict.get("target_node_id")
             source_uuid = new_node_uuid_map.get(source_original_id) or (source_original_id if source_original_id in node_map else None)
             target_uuid = new_node_uuid_map.get(target_original_id) or (target_original_id if target_original_id in node_map else None)
             if source_uuid and target_uuid:
                 resolved_link_dict = link_dict.copy()
                 resolved_link_dict["source_node_id"] = source_uuid
                 resolved_link_dict["target_node_id"] = target_uuid
                 resolved_llm_links.append(resolved_link_dict)
             else: log.warning(f"Could not resolve UUIDs for link: {source_original_id} -> {target_original_id}")

    # Merge resolved links into the QdrantNode objects in the map
    added_links, skipped_links = _merge_links_into_nodes(resolved_llm_links, node_map) # Modifies node_map in place
    log.info(f"Merged project-level LLM links: {added_links} added, {skipped_links} skipped.")

    # --- Update Ontology Metadata ---
    ontology.title = llm_summary_details.get("title", ontology.title or "Project Document Ontology")
    ontology.overview = llm_summary_details.get("overview", ontology.overview or "Analysis of project documents.")
    ontology.global_themes = llm_summary_details.get("global_themes", ontology.global_themes or [])
    ontology.key_concepts = llm_summary_details.get("key_concepts", ontology.key_concepts or [])
    # Update list of nodes
    ontology.project_graph_nodes = list(node_map.values())
    # Update list of documents
    final_doc_filenames_in_graph = {node.name for node in ontology.project_graph_nodes if node.node_type == "document"}
    updated_doc_list = sorted(list(existing_files.union(final_doc_filenames_in_graph)))
    ontology.documents = updated_doc_list
    ontology.document_count = len(ontology.documents)
    log.info(f"Ontology update complete. State: {ontology.document_count} docs, {len(ontology.project_graph_nodes)} nodes.")

    # Persist Ontology & Return
    try: save_ontology_json(ontology, path)
    except Exception as e: log.error(f"Failed to save ontology to '{path}': {e}", exc_info=True)
    return ontology

# Helper to save ontology
def save_ontology_json(ont: ProjectOntology, path: pathlib.Path):
    """Saves Pydantic ontology model to JSON file."""
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("w", encoding="utf-8") as f:
            # Use model_dump for Pydantic v2
            json.dump(ont.model_dump(mode="json"), f, ensure_ascii=False, indent=2)
        _logger.info(f"Ontology saved to {path}")
    except Exception as e:
        _logger.error(f"Error saving ontology to {path}: {e}", exc_info=True)
        if st: st.error(f"Failed to save ontology: {e}")

# Helper to load ontology
def load_ontology_json(path: pathlib.Path) -> Optional[ProjectOntology]:
    """Loads ontology from JSON file, validates with Pydantic."""
    if not path.exists(): return None
    try:
        with path.open("r", encoding="utf-8") as f: data = json.load(f)
        # Use model_validate for Pydantic v2
        return ProjectOntology.model_validate(data)
    except (json.JSONDecodeError, ValidationError, IOError, Exception) as e:
        _logger.error(f"Error loading or validating ontology from {path}: {e}", exc_info=True)
        return None


############################################################
# DOCUMENT PROCESSING WORKFLOW FUNCTION (Using Stage 1 & Finalize)
############################################################

async def process_and_finalize_single_document(
    client: genai.GenerativeModel, # Pass model instance
    file_data: dict,
    page_api_semaphore: asyncio.Semaphore,
    stage2_semaphore: asyncio.Semaphore, # Semaphore for subsection/node extraction
    stage3_semaphore: asyncio.Semaphore  # Semaphore for chapter/relationship/summary analysis
) -> Dict: # Returns dict compatible with UI { "raw_extracted_content": { ... } }
    """
    Handles the processing pipeline for a single document using multiple stages.
    Stage 1: Initial content/element extraction per page/chunk.
    Stage 2: Subsection extraction (run concurrently).
    Stage 3: Chapter extraction, Node/Relationship Analysis, Summary Generation (run sequentially for now).
    """
    function_name = "process_and_finalize_single_document"
    source_identifier = file_data.get("name", "Unknown_Document")
    file_type = file_data.get("type", "").lower()
    _logger.info(f"[{function_name}] Starting workflow for '{source_identifier}' (type: {file_type})")

    initial_pages: List[PageContent] = []
    final_result_dict = { # Default error structure
        "raw_extracted_content": { "filename": source_identifier, "pages": [], "summary": None, "error": "Processing did not complete." }
    }

    try:
        # --- Stage 1: Initial Page/Chunk Extraction ---
        stage1_start = time.time()
        _logger.debug(f"[{function_name}] --- Stage 1: Initial Extraction START --- '{source_identifier}'")

        # --- Dispatch based on file type ---
        if file_type == 'pdf':
            pdf_content = file_data.get("content")
            if not pdf_content: raise ValueError("Missing PDF content")
            if not _check_import(fitz, "PyMuPDF"): raise ImportError("PyMuPDF needed for PDF processing.")

            page_infos = await extract_pages_from_pdf_bytes(pdf_content, source_identifier) # Uses PyMuPDF
            if not page_infos:
                initial_pages = [create_error_page(1, "Failed to extract any pages from PDF")]
            else:
                page_tasks = [
                    extract_page_content_from_image_or_text(client, pi.get("page_num", i+1), source_identifier, image_part=genai_types.Part.from_data(data=base64.b64decode(pi["image_b64"]), mime_type=pi.get("image_mime_type", "image/png")))
                    for i, pi in enumerate(page_infos) if pi.get("image_b64")
                ]
                if page_tasks:
                     # Use semaphore from argument for concurrency control
                     async def run_with_semaphore(task):
                         async with page_api_semaphore: return await task
                     tasks_with_semaphore = [run_with_semaphore(task) for task in page_tasks]
                     gathered_results = await asyncio.gather(*tasks_with_semaphore, return_exceptions=True)

                     # Process results, creating error pages for exceptions
                     processed_pages = []
                     for i, res in enumerate(gathered_results):
                          page_num_ctx = page_infos[i].get('page_num', i + 1)
                          if isinstance(res, PageContent): processed_pages.append(res)
                          elif isinstance(res, Exception): processed_pages.append(create_error_page(page_num_ctx, f"Stage 1 PDF page error: {res}"))
                          else: processed_pages.append(create_error_page(page_num_ctx, f"Unexpected Stage 1 PDF page result type: {type(res)}"))
                     initial_pages = sorted(processed_pages, key=lambda p: p.page_number)
                else: initial_pages = [create_error_page(1, "No valid page images found in PDF")]

        # Handle other file types (TEXTUAL PROCESSING ONLY FOR NOW)
        # Reuse the extract_page_content_from_image_or_text for text input
        elif file_type in ["txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp", ""]:
             text_content_bytes = file_data.get("content")
             if not text_content_bytes: raise ValueError("Missing text content")
             try: text_content = text_content_bytes.decode('utf-8')
             except UnicodeDecodeError: text_content = text_content_bytes.decode('latin-1', errors='ignore')
             # Treat whole text doc as one "page" for Stage 1 analysis for simplicity now
             # Could implement chunking here if needed
             async with page_api_semaphore: # Use semaphore
                initial_pages = [await extract_page_content_from_image_or_text(client, 1, source_identifier, text_content=text_content)]

        # Handle Office/Tabular types (Using simplified text extraction)
        elif file_type == 'docx':
            if not _check_import(docx, "python-docx"): raise ImportError("python-docx needed.")
            doc = docx.Document(io.BytesIO(file_data["content"]))
            full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            async with page_api_semaphore:
                initial_pages = [await extract_page_content_from_image_or_text(client, 1, source_identifier, text_content=full_text)]
        elif file_type == 'pptx':
             if not _check_import(Presentation, "python-pptx"): raise ImportError("python-pptx needed.")
             prs = Presentation(io.BytesIO(file_data["content"]))
             ppt_pages = []
             for i, slide in enumerate(prs.slides):
                 slide_texts = []
                 for shape in slide.shapes:
                      if hasattr(shape, "text"): slide_texts.append(shape.text)
                 slide_full_text = "\n".join(slide_texts)
                 async with page_api_semaphore:
                     ppt_pages.append(await extract_page_content_from_image_or_text(client, i + 1, f"{source_identifier} (Slide {i+1})", text_content=slide_full_text))
             initial_pages = ppt_pages
        elif file_type in ['xlsx', 'xls', 'csv']:
            if not _check_import(pd, "pandas"): raise ImportError("pandas needed.")
            if file_type == 'csv': dfs = {'Sheet1': pd.read_csv(io.BytesIO(file_data["content"]))}
            else: dfs = pd.read_excel(io.BytesIO(file_data["content"]), sheet_name=None)
            excel_pages = []
            page_num = 1
            for sheet_name, df in dfs.items():
                 sheet_text = f"# Sheet: {sheet_name}\n\n{df.to_markdown(index=False)}"
                 async with page_api_semaphore:
                     excel_pages.append(await extract_page_content_from_image_or_text(client, page_num, f"{source_identifier} (Sheet: {sheet_name})", text_content=sheet_text))
                 page_num += 1
            initial_pages = excel_pages
        else:
             initial_pages = [create_error_page(1, f"Unsupported file type for processing: {file_type}")]

        stage1_duration = time.time() - stage1_start
        _logger.info(f"[{function_name}] --- Stage 1: Initial Extraction END --- '{source_identifier}' ({stage1_duration:.2f}s). Pages generated: {len(initial_pages)}")

        # Check if Stage 1 produced usable content
        if not initial_pages or all(p.raw_text and p.raw_text.startswith("[ERROR:") for p in initial_pages):
             _logger.error(f"[{function_name}] Stage 1 failed to produce usable content for '{source_identifier}'. Skipping further stages.")
             final_result_dict["raw_extracted_content"]["error"] = "Initial page extraction (Stage 1) failed or produced only errors."
             final_result_dict["raw_extracted_content"]["pages"] = [p.model_dump(mode='json') for p in initial_pages] # Include error pages
             return final_result_dict

        # --- Stage 2: Subsection Extraction (Concurrent per page) ---
        stage2_start = time.time()
        _logger.debug(f"[{function_name}] --- Stage 2: Subsection Extraction START --- '{source_identifier}'")
        subsection_tasks = []
        pages_processed_stage2 = [] # Keep track of pages successfully processed
        for page in initial_pages:
             if page.raw_text and not page.subsections: # Process only if raw text exists and no subsections yet
                 async def run_subsection_task(page_obj):
                     async with stage2_semaphore: # Use stage-specific semaphore
                         return await extract_subsections_from_text_content(client, page_obj.raw_text, page_obj.page_number, source_identifier)
                 subsection_tasks.append(run_subsection_task(page))
                 pages_processed_stage2.append(page)
             else:
                  pages_processed_stage2.append(page) # Keep page even if not processed

        if subsection_tasks:
             subsection_results = await asyncio.gather(*subsection_tasks, return_exceptions=True)
             # Assign results back to pages
             task_idx = 0
             for i, page in enumerate(pages_processed_stage2):
                  # Only update pages that were actually processed in this stage
                  if page.raw_text and not page.subsections:
                       if task_idx < len(subsection_results):
                           res = subsection_results[task_idx]
                           if isinstance(res, list): page.subsections = res
                           elif isinstance(res, Exception):
                               _logger.error(f"Subsection extraction failed for page {page.page_number} of '{source_identifier}': {res}")
                               # Create a fallback subsection indicating the error
                               page.subsections = [create_error_fallback_subsection(page.page_number, page.raw_text, f"Subsection Extraction Error: {res}")]
                           task_idx += 1
                       else: # Should not happen if logic is correct
                            _logger.error(f"Result mismatch for subsection task for page {page.page_number}")
                            page.subsections = [create_error_fallback_subsection(page.page_number, page.raw_text, "Result mismatch error")]

        # Update initial_pages list with modified pages
        initial_pages = pages_processed_stage2
        stage2_duration = time.time() - stage2_start
        _logger.info(f"[{function_name}] --- Stage 2: Subsection Extraction END --- '{source_identifier}' ({stage2_duration:.2f}s).")

        # --- Stage 3: Merging, Chapters, Node Analysis, Summary (Sequential for now) ---
        stage3_start = time.time()
        _logger.debug(f"[{function_name}] --- Stage 3: Final Analysis START --- '{source_identifier}'")
        final_summary_dict = None

        async with stage3_semaphore: # Acquire semaphore for the entire stage 3 block
            # 3a. Merge cut-off subsections
            merged_pages = await merge_cutoff_subsections(initial_pages)

            # 3b. Extract chapters
            chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages, source_identifier)
            if not chapters:
                 _logger.warning(f"[{function_name}] No chapters extracted for '{source_identifier}'. Node analysis/summary might be affected.")
                 # Proceed without chapters, analysis will run on all subsections

            # 3c. Analyze nodes WITHIN each chapter
            _logger.debug(f"[{function_name}] Analyzing relationships within {len(chapters)} chapters for '{source_identifier}'.")
            intra_chapter_tasks = [analyze_concept_entity_relationships(client, chapter, source_identifier) for chapter in chapters]
            await asyncio.gather(*intra_chapter_tasks, return_exceptions=True) # Run in parallel, errors logged internally

            # 3d. Analyze relationships BETWEEN chapters & Generate Summary
            _logger.debug(f"[{function_name}] Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
            final_summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)

        stage3_duration = time.time() - stage3_start
        _logger.info(f"[{function_name}] --- Stage 3: Final Analysis END --- '{source_identifier}' ({stage3_duration:.2f}s).")

        # --- Assemble Final Result ---
        if final_summary_dict is None:
             _logger.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             final_summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 "chapters": [ch.model_dump(mode='json') for ch in chapters] # Include chapters anyway
             }

        # Use merged_pages for the final output
        final_pages_as_dicts = [p.model_dump(mode='json') for p in merged_pages]

        final_result_dict = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": final_summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing reached here
            }
        }
        _logger.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result_dict

    except ImportError as e:
         _logger.error(f"[{function_name}] Processing failed for '{source_identifier}' due to missing library: {e}")
         final_result_dict["raw_extracted_content"]["error"] = f"Processing failed: Missing library {e}"
         return final_result_dict
    except Exception as e:
        _logger.error(f"[{function_name}] Unhandled error during workflow for '{source_identifier}': {e}", exc_info=True)
        final_result_dict["raw_extracted_content"]["error"] = f"Workflow failed unexpectedly: {str(e)}"
        # Attempt to dump initial pages if available
        try: final_result_dict["raw_extracted_content"]["pages"] = [p.model_dump(mode='json') for p in initial_pages]
        except: pass # Ignore errors during error reporting
        return final_result_dict


async def process_all_documents_async(client: genai.GenerativeModel, file_data_list: List[Dict], status_container=None, progress_bar=None, time_info=None) -> List[Dict]:
    """
    Process all documents concurrently using the new workflow function.
    Uses semaphores to control concurrency at different stages.
    """
    final_doc_results = []
    original_doc_count = len(file_data_list)
    docs_to_process_count = 0

    processing_status = {
        "active": True, "current_step": "Initializing...",
        "total_steps": original_doc_count, "current_step_num": 0,
        "start_time": time.time(), "eta_msg": "Calculating ETA..."
    }

    def update_doc_progress(completed_count, total_count, last_completed_name=None):
        """Updates the processing status."""
        try:
            processing_status["total_steps"] = total_count
            processing_status["current_step_num"] = completed_count
            overall_progress = float(completed_count) / float(total_count) if total_count > 0 else 0.0
            overall_progress = max(0.0, min(1.0, overall_progress))

            status_msg = f"Processing Documents ({completed_count}/{total_count})"
            if last_completed_name: status_msg = f"Completed: {last_completed_name} ({completed_count}/{total_count})"

            eta_msg = "Calculating ETA..."
            time_elapsed = time.time() - processing_status["start_time"]
            if overall_progress > 0.01:
                est_total_time = time_elapsed / overall_progress
                remaining = max(0, est_total_time - time_elapsed)
                mins_e, secs_e = divmod(int(time_elapsed), 60)
                mins_r, secs_r = divmod(int(remaining), 60)
                eta_msg = f"Elapsed: {mins_e}m {secs_e}s | ETA: ~{mins_r}m {secs_r}s"
            else:
                mins_e, secs_e = divmod(int(time_elapsed), 60)
                eta_msg = f"Elapsed: {mins_e}m {secs_e}s | ETA: Calculating..."
            processing_status["eta_msg"] = eta_msg

            if st and status_container: status_container.update(label=status_msg, expanded=True)
            if st and progress_bar: progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
            if st and time_info: time_info.markdown(f"`{eta_msg}`")

        except Exception as ui_err: _logger.warning(f"Failed to update UI status: {ui_err}")

    # --- Main Concurrent Processing ---
    try:
        if st: update_doc_progress(0, original_doc_count) # Initial UI update

        # --- Concurrency Limits & Semaphores ---
        # Example limits - adjust based on API limits and resources
        MAX_PARALLEL_STAGE1_API = 8 # Concurrent page/unit LLM calls across all docs
        MAX_PARALLEL_STAGE2_API = 5 # Concurrent subsection LLM calls
        MAX_PARALLEL_STAGE3_API = 3 # Concurrent chapter/analysis LLM calls
        MAX_PARALLEL_DOCS = 4       # Max documents processed *overall* at the same time

        page_api_semaphore = asyncio.Semaphore(MAX_PARALLEL_STAGE1_API)
        stage2_semaphore = asyncio.Semaphore(MAX_PARALLEL_STAGE2_API)
        stage3_semaphore = asyncio.Semaphore(MAX_PARALLEL_STAGE3_API)
        doc_workflow_semaphore = asyncio.Semaphore(MAX_PARALLEL_DOCS)

        # --- Create Tasks ---
        doc_tasks = []
        valid_file_infos = []
        for file_info in file_data_list:
            if not isinstance(file_info, dict) or "name" not in file_info or "content" not in file_info or not file_info["content"]:
                if st: st.warning(f"Skipping invalid file data: {str(file_info.get('name', 'N/A'))}...")
                _logger.warning(f"Skipping invalid file data: {str(file_info.get('name', 'N/A'))}")
                continue
            valid_file_infos.append(file_info)
            docs_to_process_count += 1

        processing_status["total_steps"] = docs_to_process_count
        if st: update_doc_progress(0, docs_to_process_count) # Update total for progress bar

        for file_info_arg in valid_file_infos:
            async def run_workflow_with_overall_semaphore(f_info):
                async with doc_workflow_semaphore: # Limit overall doc concurrency
                    return await process_and_finalize_single_document(
                        client, f_info, page_api_semaphore, stage2_semaphore, stage3_semaphore
                    )
            doc_tasks.append(run_workflow_with_overall_semaphore(file_info_arg))

        # --- Execute and Gather ---
        completed_docs = 0
        if not doc_tasks:
            _logger.warning("No valid documents to process."); processing_status["active"] = False
            if st: st.warning("No valid documents found to process."); status_container.update(label="⏹️ No valid documents found.", state="error", expanded=False)
            return []

        _logger.info(f"Launching {len(doc_tasks)} document processing workflows (max docs: {MAX_PARALLEL_DOCS}, stage1: {MAX_PARALLEL_STAGE1_API}, stage2: {MAX_PARALLEL_STAGE2_API}, stage3: {MAX_PARALLEL_STAGE3_API})...")

        for future in asyncio.as_completed(doc_tasks):
            last_completed_name = "Unknown"
            try:
                result_dict = await future
                if isinstance(result_dict, dict) and "raw_extracted_content" in result_dict:
                     final_doc_results.append(result_dict)
                     last_completed_name = result_dict.get("raw_extracted_content", {}).get("filename", "Unknown")
                     error_msg = result_dict.get("raw_extracted_content", {}).get("error")
                     if error_msg: _logger.warning(f"Workflow for '{last_completed_name}' completed with error: {error_msg}")
                     else: _logger.info(f"Workflow for '{last_completed_name}' completed successfully.")
                else:
                     _logger.error(f"Workflow task returned unexpected data type: {type(result_dict)}")
                     final_doc_results.append({"raw_extracted_content": {"filename": "Unknown_Workflow_Error", "error": f"Workflow task failed unexpectedly, returned type {type(result_dict)}"}})
            except asyncio.CancelledError:
                 _logger.warning("A document processing task was cancelled.")
                 final_doc_results.append({"raw_extracted_content": {"filename": "Cancelled_Task", "error": "Processing was cancelled."}})
                 last_completed_name = "Cancelled Task"
            except Exception as e:
                 _logger.error(f"Error awaiting document task result: {e}", exc_info=True)
                 final_doc_results.append({"raw_extracted_content": {"filename": "Unknown_Workflow_Error", "error": f"Workflow task failed: {str(e)}"}})
                 last_completed_name = "Errored Task"

            completed_docs += 1
            if st: update_doc_progress(completed_docs, docs_to_process_count, last_completed_name)

        # Final status update
        if st: update_doc_progress(completed_docs, docs_to_process_count)
        processing_status["active"] = False
        if st and status_container:
             final_state = "complete" if completed_docs == docs_to_process_count and not any(d.get('raw_extracted_content',{}).get('error') for d in final_doc_results) else "error"
             final_label = "✅ Document processing complete!" if final_state == "complete" else "⚠️ Processing finished with errors."
             try: status_container.update(label=final_label, state=final_state, expanded=False)
             except Exception: pass

    except asyncio.CancelledError:
         _logger.warning("Document processing orchestration cancelled."); processing_status["active"] = False
         if st: st.warning("Document processing cancelled."); status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
         return []
    except Exception as e:
        _logger.error(f"Error during main document processing orchestration: {e}", exc_info=True)
        processing_status["active"] = False
        if st: st.error(f"Processing pipeline error: {e}"); status_container.update(label=f"❌ Pipeline Error: {str(e)}", state="error", expanded=True)
        return final_doc_results # Return partial results if any
    finally:
        _logger.info("Document processing orchestration finished or exited.")
        processing_status["active"] = False
        if st and "st" in locals() and hasattr(st, "session_state"): st.session_state.processing_active = False

    _logger.info(f"Completed processing. Returning {len(final_doc_results)} document results.")
    return final_doc_results


############################################################
# QDRANT INGESTION WORKFLOW (2-Phase Foreground)
############################################################

# --- Phase 1a: Create Collection (m=0) ---
async def create_collection_phase1a(client: AsyncQdrantClient):
    """Phase 1a: Ensures collection exists with HNSW disabled (m=0)."""
    if not _check_import(AsyncQdrantClient, "qdrant-client"): return False
    if not client: _logger.error("Qdrant client is None."); return False
    if not QDRANT_COLLECTION_NAME: _logger.error("QDRANT_COLLECTION_NAME is not set."); return False

    collection_name = QDRANT_COLLECTION_NAME
    try:
        collection_exists = False
        try:
             # Check if collection exists efficiently
             await client.get_collection(collection_name=collection_name)
             collection_exists = True
             _logger.warning(f"Collection '{collection_name}' already exists. Phase 1a assumes it's usable or will be managed.")
             if st: st.sidebar.warning(f"Collection '{collection_name}' exists. Proceeding.")
             return True # Treat as success if exists
        except Exception as e:
             # Handle expected "not found" type errors gracefully
             if 'NotFoundError' in str(type(e)) or ' TonicHTTPError' in str(type(e)) and 'status=404' in str(e) or 'status_code=404' in str(e):
                  _logger.info(f"Collection '{collection_name}' does not exist. Will create.")
                  collection_exists = False
             else:
                  # Re-raise unexpected errors during check
                  raise e

        # If it doesn't exist, create it
        _logger.info(f"Phase 1a: Creating collection '{collection_name}' (m=0)...")
        # Prepare vector params based on configuration
        vectors_params = {
            EMB_MODEL_NAME: models.VectorParams(
                size=EMB_DIM,
                distance=models.Distance.COSINE,
                on_disk=True,
                hnsw_config=models.HnswConfigDiff(m=0) # HNSW disabled initially
            )
        }
        sparse_vectors_params = {
             SPARSE_MODEL_NAME: models.SparseVectorParams(
                 index=models.SparseIndexParams(on_disk=True)
             )
        } if SPARSE_MODEL_NAME else None # Only add if sparse model is defined

        create_kwargs = {
            "collection_name": collection_name,
            "vectors_config": vectors_params,
            "optimizers_config": models.OptimizersConfigDiff(memmap_threshold=20000)
            # Optional: Binary quantization (adjust always_ram as needed)
            # "quantization_config": models.BinaryQuantization(
            #     binary=models.BinaryQuantizationConfig(always_ram=True)
            # ),
        }
        if sparse_vectors_params: create_kwargs["sparse_vectors_config"] = sparse_vectors_params

        await client.create_collection(**create_kwargs)

        _logger.info(f"Collection '{collection_name}' created (HNSW disabled).")
        await asyncio.sleep(1) # Give Qdrant a moment
        info = await client.get_collection(collection_name=collection_name)
        _logger.info(f"Collection status after creation: {info.status}")
        return True

    except Exception as e:
        _logger.error(f"Failed during Phase 1a for collection '{collection_name}': {e}", exc_info=True)
        if st: st.error(f"Phase 1a Error: Failed to ensure collection '{collection_name}': {e}")
        return False


# --- Phase 1b: Upload Single Batch ---
async def upload_batch_nodes(
    client: AsyncQdrantClient,
    batch_nodes: List[QdrantNode], # Expects list of QdrantNode objects
    batch_num: int,
    total_batches: int,
    dense_model: TextEmbedding,
    sparse_model: Optional[SparseTextEmbedding] # Sparse model is optional
):
    """Processes and upserts a single batch of QdrantNode objects."""
    if not _check_import(AsyncQdrantClient, "qdrant-client"): return 0
    if not client: _logger.error("Qdrant client is None."); return 0
    if not batch_nodes: _logger.debug(f"Batch {batch_num}: Empty, skipping."); return 0

    _logger.debug(f"Processing batch {batch_num}/{total_batches} ({len(batch_nodes)} nodes)")
    texts = []
    valid_nodes_in_batch = []

    # 1. Prepare texts for embedding
    for node in batch_nodes:
        if not isinstance(node, QdrantNode) or not node.node_id: continue

        # Use name and core_claim for embedding text (adjust as needed)
        node_name = getattr(node, 'name', '')
        core_claim = getattr(node, 'core_claim', '')
        exec_summary = getattr(node, 'executive_summary', '') # Use if available (for doc/meta nodes)

        # Prioritize summary > claim > title > name
        content_text = exec_summary or core_claim or getattr(node, 'title', '') or node_name
        text = f"{node_name}: {content_text}".strip() if node_name else content_text.strip()
        if not text or text == ":": text = f"Placeholder content for node {node.node_id}"

        texts.append(text)
        valid_nodes_in_batch.append(node)

    if not texts: _logger.warning(f"Batch {batch_num}: No valid text extracted."); return 0

    try:
        # 2. Generate embeddings
        if not dense_model: raise ValueError("Dense embedding model not loaded")
        dense_embs_iter = dense_model.embed(texts, batch_size=len(texts)) # Embed whole batch
        dense_embs = [emb.tolist() for emb in dense_embs_iter]

        sparse_vectors = []
        if sparse_model:
            sparse_embs_iter = sparse_model.embed(texts, batch_size=len(texts))
            for emb_obj in sparse_embs_iter:
                if hasattr(emb_obj, 'indices') and hasattr(emb_obj, 'values'):
                    sparse_vectors.append(SparseVector(indices=emb_obj.indices.tolist(), values=emb_obj.values.tolist()))
                else: sparse_vectors.append(SparseVector(indices=[], values=[])) # Fallback

        if len(dense_embs) != len(texts): raise ValueError(f"Dense embedding count mismatch")
        if sparse_model and len(sparse_vectors) != len(texts): raise ValueError(f"Sparse embedding count mismatch")

        # 3. Create PointStructs
        points = []
        processed_in_batch = 0
        for i, node in enumerate(valid_nodes_in_batch):
            try:
                payload_dict = node.model_dump(mode='json', exclude_none=True)
                vector_payload = { EMB_MODEL_NAME: dense_embs[i] }
                if sparse_model and sparse_vectors:
                    vector_payload[SPARSE_MODEL_NAME] = sparse_vectors[i]

                points.append(PointStruct(
                    id=node.node_id, # Use the UUID
                    vector=vector_payload,
                    payload=payload_dict
                ))
                processed_in_batch += 1
            except Exception as point_err:
                _logger.error(f"Error creating PointStruct for node {node.node_id} in batch {batch_num}: {point_err}", exc_info=True)

        # 4. Upsert
        if points:
            await client.upsert(collection_name=QDRANT_COLLECTION_NAME, points=points, wait=False) # wait=False for potentially faster ingest
            _logger.debug(f"Upsert command for batch {batch_num} sent ({len(points)} points).")
            return processed_in_batch
        else:
            _logger.warning(f"Batch {batch_num}: No points generated after processing.")
            return 0
    except Exception as e:
        _logger.error(f"Error processing/upserting batch {batch_num}: {e}", exc_info=True)
        return 0 # Indicate failure for this batch

# --- Phase 1b: Process Group of Batches ---
async def process_batch_node_group(
    client: AsyncQdrantClient,
    group_node_slice: List[QdrantNode], # List of QdrantNode objects
    group_start_batch_num: int,
    total_batches: int,
    dense_model: TextEmbedding,
    sparse_model: Optional[SparseTextEmbedding]
):
    """Processes a specific group of node batches concurrently."""
    tasks = []
    num_nodes_in_group = len(group_node_slice)
    num_batches_in_group = (num_nodes_in_group + BATCH_SIZE - 1) // BATCH_SIZE
    _logger.debug(f"Group starting at batch {group_start_batch_num}: Processing {num_batches_in_group} batches ({num_nodes_in_group} nodes).")

    for i in range(0, num_nodes_in_group, BATCH_SIZE):
        batch_num = group_start_batch_num + (i // BATCH_SIZE)
        batch_node_list = group_node_slice[i : i + BATCH_SIZE]
        if not batch_node_list: continue
        tasks.append(
            upload_batch_nodes(client, batch_node_list, batch_num, total_batches, dense_model, sparse_model)
        )

    if not tasks: _logger.warning(f"Group {group_start_batch_num}: No tasks generated."); return True

    try:
        results = await asyncio.gather(*tasks) # Returns list of processed counts per batch
        total_processed_in_group = sum(results)
        _logger.debug(f"Group {group_start_batch_num}: asyncio.gather completed. Processed points in group: {total_processed_in_group}")
        return True # Group processed (individual batch errors handled within upload_batch_nodes)
    except Exception as e:
        _logger.error(f"Error gathering batch tasks for group {group_start_batch_num}: {e}", exc_info=True)
        return False # Indicate failure for this group

# --- Phase 1b: Foreground Orchestrator ---
def run_bulk_ingest_foreground_nodes_with_ui(
    client: AsyncQdrantClient,
    nodes: List[QdrantNode], # List of QdrantNode objects
    dense_model: TextEmbedding,
    sparse_model: Optional[SparseTextEmbedding]
) -> Tuple[bool, float, bool]:
    """Phase 1b: Uploads QdrantNode data in groups, updating UI, returns (success, time, has_data)."""
    start_ingest_time = time.perf_counter()
    _logger.info("Phase 1b: Starting FOREGROUND bulk node ingest process...")

    progress_bar = None; status_placeholder = None
    if st: # Setup UI elements if Streamlit is running
        progress_bar = st.progress(0.0, text="Initializing ingest...")
        status_placeholder = st.empty()
        status_placeholder.info("Initiating node ingest process...")

    # Get Event Loop
    current_loop = get_or_create_eventloop()
    overall_successful = True
    ingest_time = 0.0
    has_data = False # Track if any data was successfully processed

    try:
        if not nodes:
            if status_placeholder: status_placeholder.warning("No nodes provided to ingest.")
            _logger.warning("No nodes provided for Qdrant ingest.")
            return False, 0.0, False

        total_nodes = len(nodes)
        total_batches = (total_nodes + BATCH_SIZE - 1) // BATCH_SIZE
        num_groups = (total_batches + INGEST_CONCURRENCY - 1) // INGEST_CONCURRENCY
        _logger.info(f"Total nodes: {total_nodes}, Batch Size: {BATCH_SIZE}, Total Batches: {total_batches}, Concurrency Groups: {num_groups}")

        processed_nodes_count = 0

        if status_placeholder: status_placeholder.info(f"Starting upload of {total_nodes} nodes in {total_batches} batches ({num_groups} groups)...")
        if st: time.sleep(0.1) # Slight delay for UI update

        for group_idx in range(num_groups):
            group_start_batch_num = group_idx * INGEST_CONCURRENCY + 1
            group_end_batch_num = min((group_idx + 1) * INGEST_CONCURRENCY, total_batches)
            node_start_index = group_idx * INGEST_CONCURRENCY * BATCH_SIZE
            node_end_index = min((group_idx + 1) * INGEST_CONCURRENCY * BATCH_SIZE, total_nodes)
            group_data_slice = nodes[node_start_index : node_end_index]
            nodes_in_group = len(group_data_slice)

            if nodes_in_group == 0: _logger.warning(f"Group {group_idx+1}/{num_groups}: Slice is empty, skipping."); continue

            log_text = f"Processing group {group_idx+1}/{num_groups} (Batches {group_start_batch_num}-{group_end_batch_num}, {nodes_in_group} nodes)..."
            _logger.info(log_text)
            if status_placeholder: status_placeholder.info(log_text + " Processing...")

            # --- Run async processing for THIS GROUP (Blocks UI thread until group done) ---
            group_success = current_loop.run_until_complete(
                process_batch_node_group(
                    client, group_data_slice, group_start_batch_num, total_batches,
                    dense_model, sparse_model
                )
            )

            if not group_success:
                error_msg = f"Group {group_idx+1} failed. Aborting ingest."
                _logger.error(error_msg)
                if status_placeholder: status_placeholder.error(error_msg)
                overall_successful = False
                break # Stop

            # Update Progress (based on index, simpler than summing batch results)
            processed_nodes_count = min(node_end_index, total_nodes)
            progress_percentage = float(processed_nodes_count) / float(total_nodes) if total_nodes > 0 else 0.0
            progress_text = f"Ingested ~{processed_nodes_count}/{total_nodes} nodes ({progress_percentage:.1%})"
            if progress_bar: progress_bar.progress(progress_percentage, text=progress_text)
            if status_placeholder: status_placeholder.info(f"Completed group {group_idx+1}/{num_groups}. {progress_text}")
            if st: time.sleep(0.05) # Allow UI to potentially update

        # --- After the loop ---
        if overall_successful:
            _logger.info("All groups processed. Waiting for final Qdrant operations...")
            if status_placeholder: status_placeholder.info("Waiting for final operations...")
            current_loop.run_until_complete(asyncio.sleep(5)) # Wait for potential background indexing

            # Verify approximate count
            final_count = 0
            try:
                final_count_res = current_loop.run_until_complete(client.count(collection_name=QDRANT_COLLECTION_NAME, exact=False))
                final_count = final_count_res.count
                _logger.info(f"Final approximate count in Qdrant '{QDRANT_COLLECTION_NAME}': {final_count}")
                if final_count > 0:
                    has_data = True; _logger.info("Bulk ingest completed successfully.")
                    if status_placeholder: status_placeholder.success(f"Phase 1b Ingest complete. Final count: ~{final_count}")
                else: # Count is 0, might indicate an issue
                    _logger.warning(f"Ingest finished, but Qdrant count is {final_count}. Check logs.")
                    if status_placeholder: status_placeholder.warning(f"Ingest finished, but Qdrant count is {final_count}. Check logs.")
                    # Consider setting overall_successful = False if count=0 is unexpected
            except Exception as count_err:
                _logger.error(f"Failed to get final count from Qdrant: {count_err}")
                if status_placeholder: status_placeholder.warning("Could not verify final count in Qdrant.")

        else: # Failure already logged
            if status_placeholder: status_placeholder.error("Phase 1b Ingest failed.")


    except Exception as e:
        error_msg = f"Bulk ingest coordination failed: {type(e).__name__}"
        _logger.error(error_msg, exc_info=True)
        if status_placeholder: status_placeholder.error(f"{error_msg} - See console logs.")
        overall_successful = False; has_data = False
    finally:
        end_ingest_time = time.perf_counter()
        ingest_time = end_ingest_time - start_ingest_time
        _logger.info(f"Bulk node ingest coordination ended after {ingest_time:.2f}s. Success: {overall_successful}")
        # Clear temporary UI elements after a delay
        if st:
            time.sleep(3)
            if status_placeholder: status_placeholder.empty()
            if progress_bar: progress_bar.empty()

    return overall_successful, ingest_time, has_data


# --- Phase 2: Enable HNSW ---
async def enable_hnsw_phase2(client: AsyncQdrantClient, m_val=DENSE_M_PROD):
    """Phase 2: Enables HNSW index build."""
    if not _check_import(AsyncQdrantClient, "qdrant-client"): return False
    if not client: _logger.error("Qdrant client is None."); return False
    if not QDRANT_COLLECTION_NAME: _logger.error("QDRANT_COLLECTION_NAME is not set."); return False

    _logger.info(f"Phase 2: Enabling HNSW (m={m_val}) for '{QDRANT_COLLECTION_NAME}'...")
    if st: st.info(f"Phase 2: Sending HNSW update (m={m_val}). Waiting for optimization...")

    try:
        # Update HNSW config specifically for the dense vector model name
        await client.update_collection(
            collection_name=QDRANT_COLLECTION_NAME,
            vectors_config={
                EMB_MODEL_NAME: models.VectorParamsDiff(
                    hnsw_config=models.HnswConfigDiff(m=m_val)
                )
            }
            # Note: Sparse index config usually doesn't need changing here
        )
        _logger.info("HNSW update command sent. Waiting for optimizer status GREEN...")

        # Wait for optimization to finish
        start_wait = time.time()
        timeout = 300 # 5 minutes timeout
        while time.time() - start_wait < timeout:
            collection_info = await client.get_collection(collection_name=QDRANT_COLLECTION_NAME)
            current_status = collection_info.status
            optim_status = collection_info.optimizer_status
            index_count = collection_info.points_count
            _logger.debug(f"Checking status: Collection={current_status}, Optimizer OK={optim_status.ok}, Error={optim_status.error}, Indexed Count={index_count}")

            if current_status == CollectionStatus.GREEN and optim_status.ok:
                _logger.info("Status GREEN & Optimizer OK. HNSW build likely complete.")
                if st: st.success(f"Phase 2: HNSW index ready (Status: GREEN, Points: {index_count}).")
                return True
            elif current_status == CollectionStatus.YELLOW or not optim_status.ok:
                wait_msg = f"Status {current_status} (Optimizing... Indexed: {index_count or '?'}). Waiting..."
                _logger.info(wait_msg + " (5s)")
                if st: st.info(wait_msg)
                await asyncio.sleep(5)
            elif current_status == CollectionStatus.RED:
                _logger.error(f"Status RED. Optimization failed: {optim_status.error}")
                if st: st.error(f"Optimization failed (Status: RED): {optim_status.error}")
                return False
            else: # Should not happen
                _logger.warning(f"Unexpected status: {current_status}. Waiting 5s.")
                if st: st.warning(f"Unexpected status: {current_status}. Waiting...")
                await asyncio.sleep(5)

        _logger.error("Timeout waiting for HNSW optimization to complete.")
        if st: st.error("Timeout waiting for HNSW optimization.")
        return False

    except Exception as e:
        _logger.error(f"Failed to enable HNSW: {e}", exc_info=True)
        if st: st.error(f"Failed to enable HNSW: {e}")
        return False

# --- Qdrant Search Function (Example) ---
async def search_qdrant_nodes(
    qdrant_client: AsyncQdrantClient,
    query_text: str,
    dense_model: TextEmbedding,
    sparse_model: Optional[SparseTextEmbedding],
    top_k: int = 5
) -> List[Dict]:
    """Performs hybrid search in Qdrant and returns results."""
    if not all([qdrant_client, query_text, dense_model]):
        _logger.error("Missing client, query, or dense model for Qdrant search.")
        return []
    if not _check_import(AsyncQdrantClient, "qdrant-client"): return []

    try:
        _logger.info(f"Performing hybrid search for query: '{query_text[:50]}...'")
        # 1. Embed query
        dense_query_vector = next(dense_model.embed([query_text])).tolist()
        sparse_query_vector = None
        if sparse_model:
            sparse_emb = next(sparse_model.embed([query_text]))
            sparse_query_vector = SparseVector(
                indices=sparse_emb.indices.tolist(),
                values=sparse_emb.values.tolist()
            )

        # 2. Build search requests
        search_requests = []
        # Dense search request
        search_requests.append(
            models.SearchRequest(
                vector=models.NamedVector(name=EMB_MODEL_NAME, vector=dense_query_vector),
                limit=top_k,
                with_payload=True, # Get payload data
                # Optional: Add hnsw_ef search parameter if needed
                # params=models.SearchParams(hnsw_ef=HNSW_EF_SEARCH)
            )
        )
        # Sparse search request (if applicable)
        if sparse_query_vector:
             search_requests.append(
                 models.SearchRequest(
                     vector=models.NamedVector(name=SPARSE_MODEL_NAME, vector=sparse_query_vector),
                     limit=top_k,
                     with_payload=True,
                 )
             )

        # 3. Perform search using batch interface for hybrid search
        search_result = await qdrant_client.search_batch(
            collection_name=QDRANT_COLLECTION_NAME,
            requests=search_requests
        )

        # 4. Process and combine results (Simple RRF or just return top dense/sparse)
        # For simplicity, let's just take top K dense results here.
        # Real hybrid would involve re-ranking.
        processed_results = []
        if search_result and len(search_result) > 0:
            dense_results = search_result[0] # Assuming dense is first request
            for hit in dense_results:
                 result_dict = {
                     "id": hit.id,
                     "score": hit.score,
                     "payload": hit.payload
                 }
                 processed_results.append(result_dict)
            _logger.info(f"Qdrant search returned {len(processed_results)} results.")
        else:
             _logger.warning("Qdrant search returned no results.")

        return processed_results

    except Exception as e:
        _logger.error(f"Error during Qdrant search: {e}", exc_info=True)
        if st: st.error(f"Search failed: {e}")
        return []


############################################################
# UI DISPLAY FUNCTIONS (Adapted for integrated models)
############################################################

def display_table(table_data: Union[TableData, Dict]):
    """Display table data (accepts model or dict)."""
    if not table_data: return
    if not st: return # No Streamlit

    try:
        table_dict = table_data if isinstance(table_data, dict) else table_data.model_dump()

        title = table_dict.get("title")
        content = table_dict.get("table_content")
        summary = table_dict.get("summary")
        page_num = table_dict.get("page_number", "?")
        table_id = table_dict.get("table_id", "N/A")

        st.markdown(f"**Table (Page {page_num}, ID: `{table_id}`)**")
        if title: st.markdown(f"*{title}*")
        if summary: st.caption(f"Summary: {summary}")

        if content:
            # Attempt to display as DataFrame, fallback to Markdown
            if _check_import(pd, "pandas"):
                try:
                    lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
                    if len(lines) > 1 and all(c in '-:| ' for c in lines[1]): lines.pop(1) # Remove separator
                    cleaned_content = "\n".join(lines)
                    # Use StringIO for pandas compatibility
                    df = pd.read_csv(io.StringIO(cleaned_content), sep='|', skipinitialspace=True, index_col=False).iloc[1:] # Skip header row parsed as data
                    df.columns = [col.strip() for col in df.columns] # Clean column names
                    df = df.drop(columns=[col for col in df.columns if col == ''], errors='ignore') # Drop empty cols from extra separators
                    df = df.iloc[:, 1:-1] # Drop first/last empty columns from leading/trailing '|'
                    df = df.dropna(axis=1, how='all') # Drop fully empty columns
                    st.dataframe(df.reset_index(drop=True)) # Show clean index
                except Exception as e:
                    _logger.warning(f"Pandas failed to parse table {table_id}, showing markdown. Error: {e}. Content:\n{content[:200]}...")
                    st.markdown(f"```markdown\n{content}\n```")
            else:
                st.markdown(f"```markdown\n{content}\n```") # Fallback if pandas not installed
        else:
            st.caption("Table content not available.")
    except Exception as e:
         _logger.error(f"Error displaying table: {e}", exc_info=True)
         st.error("Error displaying table.")

def display_visual_element(visual_data: Union[VisualElement, Dict]):
    """Display visual element data (accepts model or dict)."""
    if not visual_data: return
    if not st: return

    try:
        visual_dict = visual_data if isinstance(visual_data, dict) else visual_data.model_dump()

        visual_type = visual_dict.get('type', 'Unknown Type')
        description = visual_dict.get('description', 'No description available.')
        data_summary = visual_dict.get('data_summary', 'N/A')
        page_nums = visual_dict.get('page_numbers', ['?'])
        visual_id = visual_dict.get('visual_id', 'N/A')

        st.markdown(f"**🖼️ {visual_type.capitalize()} (Page(s): {', '.join(map(str, page_nums))}, ID: `{visual_id}`)**")
        st.markdown(f"*Description:* {description}")
        if data_summary and data_summary.lower() != "n/a":
            st.markdown(f"*Data Summary:* {data_summary}")
        # Consider adding source_url or alt_text if useful
    except Exception as e:
        _logger.error(f"Error displaying visual element: {e}", exc_info=True)
        st.error("Error displaying visual.")


def display_numerical_data_point(num_data: Union[NumericalDataPoint, Dict]):
    """Displays a single numerical data point."""
    if not num_data: return
    if not st: return

    try:
        num_dict = num_data if isinstance(num_data, dict) else num_data.model_dump()

        value = num_dict.get('value', 'N/A')
        description = num_dict.get('description', 'No description.')
        context = num_dict.get('context')

        st.markdown(f"**{value}**: {description}")
        if context: st.caption(f"Context: ...{context}...")
    except Exception as e:
        _logger.error(f"Error displaying numerical data: {e}", exc_info=True)
        st.error("Error displaying number.")


def display_qdrant_node_network(qdrant_nodes_data: Optional[List[Union[QdrantNode, Dict]]], title="Concept/Entity Network"):
    """Displays relationships from QdrantNode data (dicts or models) as a network graph."""
    if not qdrant_nodes_data:
        if st: st.info("No node data available to display network.")
        return
    if not st: return

    if not _check_import(nx, "NetworkX") or not _check_import(plt, "Matplotlib"):
        st.warning("NetworkX/Matplotlib not installed. Displaying relationships as text list.")
        st.markdown("#### Relationships (Text List):")
        rel_count = 0
        for node_data in qdrant_nodes_data:
            try:
                node_dict = node_data if isinstance(node_data, dict) else node_data.model_dump()
                linked_nodes = node_dict.get('linked_nodes', [])
                source_name = node_dict.get('name','Unknown Source')
                if linked_nodes and isinstance(linked_nodes, list):
                    for link_dict in linked_nodes:
                         if isinstance(link_dict, dict):
                            target_id = link_dict.get('target_node_id','Unknown Target')
                            # Try to find target name in the provided list (basic lookup)
                            target_name = target_id
                            for target_data in qdrant_nodes_data:
                                target_data_dict = target_data if isinstance(target_data, dict) else target_data.model_dump()
                                if target_data_dict.get('node_id') == target_id:
                                    target_name = target_data_dict.get('name', target_id); break
                            st.markdown(f"- **{source_name}** → **{target_name}**: {link_dict.get('relationship_description','?')} (Strength: {link_dict.get('relationship_strength','?')})")
                            rel_count += 1
            except Exception as text_disp_err:
                _logger.error(f"Error displaying text relationship: {text_disp_err}")
        if rel_count == 0: st.caption("No relationships found in data.")
        return

    # --- NetworkX Graphing ---
    G = nx.DiGraph()
    node_labels = {}
    edge_weights = {}
    edge_labels = {}
    node_types_map = {}

    # Add nodes
    node_ids_added = set()
    for node_data in qdrant_nodes_data:
        try:
            node_dict = node_data if isinstance(node_data, dict) else node_data.model_dump()
            node_id = node_dict.get('node_id')
            if node_id and node_id not in node_ids_added:
                node_name = node_dict.get('name', 'Unknown')
                node_type = node_dict.get('node_type', 'unknown')
                G.add_node(node_id)
                # Make label slightly smaller for readability
                node_labels[node_id] = f"{node_name[:25]}{'...' if len(node_name)>25 else ''}\n({node_type})"
                node_types_map[node_id] = node_type
                node_ids_added.add(node_id)
        except Exception as node_add_err:
            _logger.error(f"Error adding node to graph: {node_add_err}")

    # Add edges
    edges_added = set()
    for node_data in qdrant_nodes_data:
        try:
            node_dict = node_data if isinstance(node_data, dict) else node_data.model_dump()
            source_id = node_dict.get('node_id')
            linked_nodes = node_dict.get('linked_nodes', [])
            if not source_id or not G.has_node(source_id) or not isinstance(linked_nodes, list): continue

            for link_data in linked_nodes:
                 link_dict = link_data if isinstance(link_data, dict) else link_data.model_dump()
                 target_id = link_dict.get('target_node_id')

                 if not target_id: continue
                 # Ensure target exists (add if missing - signifies inter-document or incomplete graph)
                 if not G.has_node(target_id):
                      _logger.warning(f"Target node '{target_id}' not found in main list. Adding placeholder.")
                      G.add_node(target_id)
                      # Basic label from ID
                      target_name_part = target_id.split('_')[1] if '_' in target_id else target_id
                      node_labels[target_id] = f"{target_name_part[:25]}{'...' if len(target_name_part)>25 else ''}\n(external?)"
                      node_types_map[target_id] = 'unknown'

                 edge_tuple = (source_id, target_id)
                 if edge_tuple not in edges_added:
                     try: strength = float(link_dict.get('relationship_strength', 5.0))
                     except: strength = 5.0
                     description = link_dict.get('relationship_description', '')

                     G.add_edge(source_id, target_id, weight=strength)
                     edge_weights[edge_tuple] = strength
                     edge_labels[edge_tuple] = description[:30] + '...' if len(description)>30 else description # Shorten label
                     edges_added.add(edge_tuple)
        except Exception as edge_add_err:
            _logger.error(f"Error adding edges for node {node_dict.get('node_id')}: {edge_add_err}")


    if not G.nodes(): st.info("No nodes found to build the graph."); return
    if not G.edges():
         st.info("Nodes found, but no links between them to display.")
         node_list_str = ", ".join([lbl.split('\n')[0] for lbl in node_labels.values()][:10])
         if len(node_labels) > 10: node_list_str += ", ..."
         st.write(f"Nodes identified: {node_list_str}")
         return

    # --- Draw Graph ---
    try:
        plt.style.use('seaborn-v0_8-whitegrid')
        fig, ax = plt.subplots(figsize=(14, 11)) # Adjusted size
        pos = nx.spring_layout(G, k=0.6, iterations=50, seed=42) # Adjusted k

        # Node colors
        type_colors = {'entity': 'skyblue', 'concept': 'lightgreen', 'document': 'lightcoral', 'meta_summary': 'lightsalmon'}
        node_colors = [type_colors.get(node_types_map.get(n, 'unknown'), 'lightgray') for n in G.nodes()]
        nx.draw_networkx_nodes(G, pos, node_size=1200, node_color=node_colors, alpha=0.9, ax=ax) # Slightly smaller nodes

        # Edges with width based on strength
        min_w = 1.0; max_w = 5.0
        weights = list(edge_weights.values())
        if weights:
             min_s = min(weights) if weights else 1.0
             max_s = max(weights) if weights else 10.0
             if max_s == min_s: max_s += 1 # Avoid division by zero
             edge_width = [min_w + (max_w - min_w) * (edge_weights.get(edge, 5.0) - min_s) / (max_s - min_s) for edge in G.edges()]
        else: edge_width = [min_w] * len(G.edges())

        nx.draw_networkx_edges(G, pos, width=edge_width, alpha=0.5, edge_color='gray',
                              arrowsize=12, connectionstyle='arc3,rad=0.1', ax=ax)

        # Labels
        nx.draw_networkx_labels(G, pos, labels=node_labels, font_size=8, font_weight='normal', ax=ax) # Smaller font

        # Optional: Edge labels (can get crowded)
        # nx.draw_networkx_edge_labels(G, pos, edge_labels=edge_labels, font_size=7, font_color='red', ax=ax)

        ax.set_title(title, fontsize=16)
        plt.axis('off')
        st.pyplot(fig)
        plt.close(fig) # Close plot to free memory

        # Display edge details table (Optional)
        if _check_import(pd, "pandas"):
            edge_data = [{
                "Source": node_labels.get(s, s).split('\n')[0],
                "Target": node_labels.get(t, t).split('\n')[0],
                "Strength": edge_weights.get((s, t), '?'),
                "Description": edge_labels.get((s, t), '?')
            } for s, t in edge_labels.keys()]
            if edge_data:
                with st.expander("View Relationship Details (Table)"):
                    st.dataframe(pd.DataFrame(edge_data))

    except Exception as draw_err:
         _logger.error(f"Error drawing graph: {draw_err}", exc_info=True)
         st.error(f"Could not draw graph: {draw_err}")


def display_chapter(chapter_data: Union[Chapter, Dict]):
    """Display chapter details including subsections and Qdrant node relationships."""
    if not chapter_data: return
    if not st: return

    try:
        chapter_dict = chapter_data if isinstance(chapter_data, dict) else chapter_data.model_dump()

        title = chapter_dict.get('title', 'Untitled Chapter')
        summary = chapter_dict.get('summary', 'No summary available.')
        subsections = chapter_dict.get('subsections', [])
        qdrant_nodes = chapter_dict.get('qdrant_nodes') # List of QdrantNode models/dicts

        st.markdown(f"### {title}")
        st.markdown(f"_{summary}_")

        # Display Qdrant Node relationships if available
        if qdrant_nodes and isinstance(qdrant_nodes, list):
            with st.expander("View Concepts/Entities & Relationships within this Chapter", expanded=False):
                display_qdrant_node_network(qdrant_nodes, title=f"Network for Chapter '{title}'")
        else:
            st.caption("No concept/entity relationship data available for this chapter.")

        # Display subsections
        st.markdown("#### Subsections")
        if subsections and isinstance(subsections, list):
            for sub_idx, sub_data in enumerate(subsections):
                sub_dict = sub_data if isinstance(sub_data, dict) else sub_data.model_dump()
                sub_title = sub_dict.get('title', f'Subsection {sub_idx + 1}')
                sub_desc = sub_dict.get('description', '')
                sub_text = sub_dict.get('text', 'No text available.')
                sub_id = sub_dict.get('subsection_id', 'N/A')

                with st.expander(f"{sub_idx + 1}. {sub_title}"):
                    if sub_desc: st.caption(f"Description: {sub_desc}")
                    st.markdown(sub_text)
                    st.caption(f"ID: {sub_id}")
        else:
            st.caption("No subsections listed for this chapter.")

    except Exception as e:
        _logger.error(f"Error displaying chapter: {e}", exc_info=True)
        st.error("Error displaying chapter details.")


async def display_node_links_qdrant(
    qdrant_client: AsyncQdrantClient,
    collection_name: str,
    selected_node_id: str
):
    """Fetches node and links from Qdrant and displays them."""
    if not selected_node_id: st.caption("No node selected."); return
    if not _check_import(AsyncQdrantClient, "qdrant-client") or not qdrant_client:
        st.error("Qdrant client is not available."); return
    if not st: return

    loading_msg = st.info(f"Fetching details for node `{selected_node_id}` from Qdrant...")
    selected_payload = None; target_points_map = {}

    try:
        # Fetch Source Node
        retrieved_points = await qdrant_client.retrieve(
            collection_name=collection_name, ids=[selected_node_id], with_payload=True, with_vectors=False
        )
        if not retrieved_points: st.warning(f"Node ID '{selected_node_id}' not found."); loading_msg.empty(); return
        selected_payload = retrieved_points[0].payload
        if not selected_payload: st.warning(f"Node '{selected_node_id}' has no payload."); loading_msg.empty(); return

        source_name = selected_payload.get('name', selected_node_id)
        source_type = selected_payload.get('node_type', 'unknown')
        source_doc = selected_payload.get('source', '?')
        source_context = selected_payload.get('chapter_id_context', '?')
        source_title = selected_payload.get('title', 'N/A')
        source_claim = selected_payload.get('core_claim', 'N/A')

        linked_nodes_data = selected_payload.get('linked_nodes', [])
        target_ids = list(set(link.get('target_node_id') for link in linked_nodes_data if isinstance(link, dict) and link.get('target_node_id')))

        # Fetch Target Nodes (if any links exist)
        if target_ids:
            _logger.debug(f"Retrieving {len(target_ids)} target node details from Qdrant.")
            target_points = await qdrant_client.retrieve(
                collection_name=collection_name, ids=target_ids, with_payload=True, with_vectors=False
            )
            for point in target_points:
                if point.payload: target_points_map[point.id] = point.payload
            _logger.debug(f"Retrieved details for {len(target_points_map)} targets.")

        loading_msg.empty() # Clear loading message

        # --- Display Node Info ---
        st.markdown(f"#### Relationships FROM: **{source_name}** (`{source_type}`) ")
        st.caption(f"*(Origin: Doc '{source_doc}', Context '{source_context}')*")
        with st.expander("View Source Node Details"):
            st.markdown(f"**Title:** {source_title}"); st.markdown(f"**Core Claim:** {source_claim}")
            st.markdown(f"**Info Type:** {selected_payload.get('information_type', 'N/A')}")
            st.markdown(f"**Sentiment:** {selected_payload.get('sentiment_tone', 'N/A')}")
            try: tags_list = json.loads(selected_payload.get('tags', '[]'))
            except: tags_list = []
            if tags_list: st.markdown(f"**Tags:** {', '.join(tags_list)}")
            st.caption(f"ID: `{selected_node_id}`")
        st.markdown("---")

        # --- Display Links ---
        if not linked_nodes_data or not isinstance(linked_nodes_data, list):
            st.caption("No outgoing links found for this node.")
            return

        links_displayed = 0
        for i, link_data in enumerate(linked_nodes_data):
            if not isinstance(link_data, dict): continue
            target_id = link_data.get('target_node_id'); description = link_data.get('relationship_description', '?')
            keywords = link_data.get('relationship_keywords', []); strength = link_data.get('relationship_strength', 0.0)
            if not target_id: continue

            target_payload = target_points_map.get(target_id)
            target_name = target_id; target_type = "unknown"; target_doc = "?"; target_context = "?"

            if target_payload:
                target_name = target_payload.get('name', target_id)
                target_type = target_payload.get('node_type', 'unknown')
                target_doc = target_payload.get('source', '?')
                target_context = target_payload.get('chapter_id_context', '?')

            st.markdown(f"**{i+1}. Target Node:** {target_name} (`{target_type}`)")
            st.markdown(f"    *Origin: Doc '{target_doc}', Context '{target_context}'*")
            st.markdown(f"    **Link:** {description}")
            st.markdown(f"    **Strength:** {strength:.1f}/10.0")
            if keywords and isinstance(keywords, list): st.markdown(f"    **Keywords:** {', '.join(keywords)}")
            st.markdown(f"    *Target Node ID:* `{target_id}`")
            st.markdown("---")
            links_displayed += 1

        if links_displayed == 0: st.caption("No valid outgoing links could be displayed.")

    except Exception as e:
        loading_msg.empty()
        st.error(f"Error displaying node links from Qdrant: {e}")
        _logger.error(f"Error displaying node links for {selected_node_id}: {e}", exc_info=True)


def display_page_details(page_content_data: Union[PageContent, Dict]):
    """Displays details for a single page."""
    if not page_content_data: return
    if not st: return

    try:
        page_dict = page_content_data if isinstance(page_content_data, dict) else page_content_data.model_dump()
        page_num = page_dict.get("page_number", "?")
        st.markdown(f"#### 📄 Page {page_num}")

        flags = []
        if page_dict.get('has_tables'): flags.append("Tables ✅")
        if page_dict.get('has_visuals'): flags.append("Visuals ✅")
        if page_dict.get('has_numbers'): flags.append("Numbers ✅")
        st.caption("Detected: " + " | ".join(flags) if flags else "Detected: None")
        st.divider()

        # Display Subsections
        subsections = page_dict.get('subsections', [])
        if subsections and isinstance(subsections, list):
            st.markdown("**Subsections**")
            for sub_idx, sub_data in enumerate(subsections):
                 sub_dict = sub_data if isinstance(sub_data, dict) else sub_data.model_dump()
                 sub_title = sub_dict.get('title', f'Section {sub_idx+1}')
                 exp_title = f"{sub_idx+1}. {sub_title}"
                 if sub_dict.get('is_cutoff'): exp_title += " (Continues...)"
                 with st.expander(exp_title, expanded=True):
                     st.caption(f"Desc: {sub_dict.get('description', 'N/A')}")
                     st.markdown(sub_dict.get('text', 'No text.'))
                     refs = []
                     if sub_dict.get('referenced_visuals'): refs.append(f"Visuals: {', '.join(sub_dict['referenced_visuals'])}")
                     if sub_dict.get('referenced_tables'): refs.append(f"Tables: {', '.join(sub_dict['referenced_tables'])}")
                     if refs: st.caption(f"References: {'; '.join(refs)}")
                     st.caption(f"ID: {sub_dict.get('subsection_id', 'N/A')}")
            st.divider()
        else:
             st.caption("No subsections extracted for this page.")
             # Show raw text if no subsections
             if page_dict.get('raw_text'):
                  with st.expander("View Raw Page Text", expanded=False):
                       st.text(page_dict['raw_text'])
                  st.divider()


        # Display Tables, Visuals, Numbers
        if page_dict.get('tables'):
            st.markdown("**Tables**")
            for table in page_dict['tables']: display_table(table)
            st.divider()
        if page_dict.get('visuals'):
            st.markdown("**Visuals**")
            for visual in page_dict['visuals']: display_visual_element(visual)
            st.divider()
        if page_dict.get('numbers'):
            st.markdown("**Numerical Data**")
            for num_data in page_dict['numbers']: display_numerical_data_point(num_data)
            st.divider()

    except Exception as e:
        _logger.error(f"Error displaying page details: {e}", exc_info=True)
        st.error("Error displaying page details.")


def render_unified_document_report(document_data: Optional[Dict]):
    """Displays the full report for a single processed document."""
    if not document_data: st.error("No document data provided."); return
    if not st: return

    try:
        # Assuming input is the dict structure from process_all_documents_async
        raw_content = document_data.get("raw_extracted_content", {})
        if not raw_content: st.error("Invalid document data format."); return

        filename = raw_content.get('filename', 'Unknown Document')
        processing_error = raw_content.get('error')
        pages = raw_content.get('pages', []) # List of PageContent dicts
        summary_data = raw_content.get('summary', {}) # Dict from DocumentSummaryDetails + chapters
        chapters_list = summary_data.get('chapters', []) if isinstance(summary_data, dict) else []

        st.header(f"📄 Report: {filename}")
        if processing_error: st.error(f"🚨 **Processing Error:** {processing_error}")

        num_pages = len(pages) if isinstance(pages, list) else 0
        st.caption(f"📄 {num_pages} pages/units processed") # Simplified caption

        tab_summary, tab_chapters, tab_pages, tab_graph = st.tabs([
            "Executive Summary", "Chapters", "Detailed Pages", "Document Graph"
        ])

        # --- Summary Tab ---
        with tab_summary:
            if isinstance(summary_data, dict):
                 st.subheader("Document Summary")
                 summary_title = summary_data.get('title')
                 if summary_title and "Failed" not in summary_title: st.markdown(f"**Title:** {summary_title}")
                 summary_text = summary_data.get('summary')
                 if summary_text and "Failed" not in summary_text: st.markdown(summary_text)
                 else: st.info("No summary text available or generation failed.")

                 themes = summary_data.get('themes')
                 if themes and themes != ["error"]:
                     st.markdown("#### Key Themes"); st.markdown("- " + "\n- ".join(themes))
                 questions = summary_data.get('questions')
                 if questions:
                      st.markdown("#### Key Questions Answered"); st.markdown("- " + "\n- ".join(questions))
            else: st.warning("Summary data is missing or invalid.")

        # --- Chapters Tab ---
        with tab_chapters:
             if chapters_list and isinstance(chapters_list, list):
                 chapter_titles = [ch.get('title', f"Chapter {ch.get('order', i+1)}") for i, ch in enumerate(chapters_list) if isinstance(ch, dict)]
                 if chapter_titles:
                     try:
                         chapter_tabs = st.tabs(chapter_titles)
                         for tab, chapter_dict in zip(chapter_tabs, chapters_list):
                             if isinstance(chapter_dict, dict):
                                  with tab: display_chapter(chapter_dict)
                     except Exception as tabs_err: st.error(f"Could not create chapter tabs: {tabs_err}")
                 else: st.info("No valid chapter titles found.")
             else: st.info("No chapter information available.")

        # --- Pages Tab ---
        with tab_pages:
             if pages and isinstance(pages, list):
                  page_num_options = {} # Map Page Num -> Index
                  valid_page_nums = []
                  for i, p_dict in enumerate(pages):
                      if isinstance(p_dict, dict) and isinstance(p_dict.get('page_number'), int):
                           pg_num = p_dict['page_number']
                           page_num_options[pg_num] = i
                           valid_page_nums.append(pg_num)
                      else: _logger.warning(f"Invalid page data at index {i} for page selector.")

                  if valid_page_nums:
                       sorted_page_nums = sorted(valid_page_nums)
                       selected_page_num = st.selectbox(
                           "Select Page Number to View Details:", options=sorted_page_nums,
                           format_func=lambda x: f"Page {x}", index=0,
                           key=f"page_detail_selector_{filename}"
                       )
                       selected_index = page_num_options.get(selected_page_num)
                       if selected_index is not None: display_page_details(pages[selected_index])
                       else: st.warning("Could not find selected page index.")
                  else: st.info("No valid page data with page numbers found.")
             else: st.info("No page data available.")

        # --- Graph Tab ---
        with tab_graph:
             st.subheader("Document Concept & Entity Network")
             all_document_nodes = []
             if isinstance(chapters_list, list):
                 for chapter_dict in chapters_list:
                      if isinstance(chapter_dict, dict):
                          nodes_in_chapter = chapter_dict.get('qdrant_nodes', [])
                          if isinstance(nodes_in_chapter, list): all_document_nodes.extend(nodes_in_chapter)

             if all_document_nodes:
                  display_qdrant_node_network(all_document_nodes, title=f"Overall Network for {filename}")
                  st.divider()
                  # Add Document-Level Node Explorer
                  st.subheader("Explore Node Relationships (This Document)")
                  qdrant_client = get_qdrant_client()
                  if qdrant_client is None: st.error("Qdrant client unavailable.")
                  else:
                      doc_node_options: Dict[str, str] = {} # label -> node_id
                      try: sorted_doc_nodes = sorted(all_document_nodes, key=lambda n: n.get('name', '') if isinstance(n, dict) else '')
                      except: sorted_doc_nodes = all_document_nodes
                      for node_data in sorted_doc_nodes:
                           node_dict = node_data if isinstance(node_data, dict) else node_data.model_dump()
                           node_id = node_dict.get('node_id'); node_name = node_dict.get('name', node_id)
                           node_type = node_dict.get('node_type', 'unknown')
                           if node_id:
                                label = f"{node_name} ({node_type})"; count=1
                                while label in doc_node_options: label = f"{node_name} ({node_type}) ({count})"; count+=1
                                doc_node_options[label] = node_id

                      if doc_node_options:
                           selected_label = st.selectbox("Select Node:", options=["Select..."] + list(doc_node_options.keys()), key=f"doc_node_explorer_{filename}")
                           if selected_label != "Select...":
                                selected_id = doc_node_options.get(selected_label)
                                if selected_id: run_async(display_node_links_qdrant, qdrant_client, QDRANT_COLLECTION_NAME, selected_id)
                      else: st.info("No valid nodes found in this document.")
             else:
                  st.info("No node data available to build document network.")

    except Exception as e:
        _logger.error(f"Error rendering document report for {filename}: {e}", exc_info=True)
        st.error(f"Error displaying report: {e}")


def display_project_ontology(ontology: Optional[ProjectOntology]):
    """Displays project-wide ontology including the aggregated graph."""
    if not ontology: st.warning("No project ontology data available."); return
    if not st: return

    try:
        st.header(f"🌐 Project Ontology: {ontology.title}")
        st.markdown(f"**Documents Analyzed:** {ontology.document_count}")

        with st.expander("Documents Included"):
             st.markdown("- " + "\n- ".join(ontology.documents or []))

        st.markdown("**Overview:**"); st.markdown(ontology.overview)
        if ontology.global_themes: st.markdown("**Global Themes:**"); st.markdown("- " + "\n- ".join(ontology.global_themes))
        if ontology.key_concepts: st.markdown("**Key Concepts/Docs:**"); st.markdown("- " + "\n- ".join(ontology.key_concepts))

        st.divider()
        st.subheader("Project Knowledge Graph")
        if ontology.project_graph_nodes and isinstance(ontology.project_graph_nodes, list):
             _logger.info(f"Rendering project graph with {len(ontology.project_graph_nodes)} nodes.")
             with st.expander("View Full Project Network", expanded=False):
                  display_qdrant_node_network(ontology.project_graph_nodes, title="Project-Wide Network")

             # --- Project-Level Node Relationship Explorer ---
             st.divider()
             st.subheader("Explore Node Relationships (Project-Wide)")
             qdrant_client = get_qdrant_client()
             if qdrant_client is None: st.error("Qdrant client unavailable.")
             else:
                 project_node_options: Dict[str, str] = {} # label -> node_id
                 try: sorted_proj_nodes = sorted(ontology.project_graph_nodes, key=lambda n: n.name)
                 except: sorted_proj_nodes = ontology.project_graph_nodes or []

                 for node in sorted_proj_nodes:
                      node_id = node.node_id; node_name = node.name; node_type = node.node_type
                      if node_id:
                           label = f"{node_name} ({node_type})"; count=1
                           while label in project_node_options: label = f"{node_name} ({node_type}) ({count})"; count+=1
                           project_node_options[label] = node_id

                 if project_node_options:
                      selected_label = st.selectbox("Select Node:", options=["Select..."] + list(project_node_options.keys()), key="project_node_explorer_main")
                      if selected_label != "Select...":
                           selected_id = project_node_options.get(selected_label)
                           if selected_id: run_async(display_node_links_qdrant, qdrant_client, QDRANT_COLLECTION_NAME, selected_id)
                 else: st.info("No valid nodes found in the project ontology.")
        else:
            st.caption("No cross-document relationship data generated or available.")

    except Exception as e:
        _logger.error(f"Error displaying project ontology: {e}", exc_info=True)
        st.error(f"Error displaying project ontology: {e}")


############################################################
# Streamlit Application Setup
############################################################

def display_sidebar_chat(): # Takes client model instance
    """Manages sidebar: file upload, Qdrant control, status, chat."""
    if not st: return
    st.sidebar.title("📝 Input & Control")
    st.sidebar.markdown("Upload docs, process, generate ontology, manage Qdrant, chat.")

    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]

    # --- Initialize Session State ---
    default_state = {
        "processed_file_names": set(), "processing_active": False, "processed_documents": [],
        "show_results": False, "project_ontology": None, "show_ontology": False,
        "messages": [], "selected_docs": [],
        "qdrant_client": None, "dense_model": None, "sparse_model": None,
        "collection_created_phase1a": False, "indexing_complete_phase1b": False,
        "hnsw_enabled_phase2": False, "collection_has_data": False,
        "status_container": None, "progress_bar": None, "time_info": None, # UI elements
        "gemini_client_model": None # Store the Gemini model instance
    }
    for key, default in default_state.items():
        if key not in st.session_state: st.session_state[key] = default

    # --- Ensure Clients/Models are loaded ---
    if st.session_state.gemini_client_model is None:
         api_key = get_gemini_api_key()
         if api_key:
             try:
                  genai.configure(api_key=api_key)
                  # Use 1.5 Flash as default for general processing
                  st.session_state.gemini_client_model = genai.GenerativeModel('gemini-1.5-flash-preview-0514')
                  _logger.info("Gemini client model initialized.")
             except Exception as e:
                  _logger.error(f"Failed to initialize Gemini client model: {e}", exc_info=True)
                  st.sidebar.error(f"Gemini init error: {e}")
         else:
              _logger.error("Gemini API Key unavailable.") # Error already shown by get_gemini_api_key

    if st.session_state.qdrant_client is None: st.session_state.qdrant_client = get_qdrant_client()
    if st.session_state.dense_model is None or st.session_state.sparse_model is None:
        st.session_state.dense_model, st.session_state.sparse_model = load_models()

    gemini_ready = st.session_state.gemini_client_model is not None
    qdrant_ready = st.session_state.qdrant_client is not None
    models_ready = st.session_state.dense_model is not None and st.session_state.sparse_model is not None

    # --- File Upload ---
    st.sidebar.subheader("📁 Upload & Process Documents")
    uploaded_files = st.sidebar.file_uploader("Select files:", type=supported_types, accept_multiple_files=True, key="sidebar_file_uploader_main")

    files_data_for_processing = []
    new_files_count = 0
    if uploaded_files:
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                try:
                    file_content = file.getvalue(); file_ext = file.name.split('.')[-1].lower() if '.' in file.name else ''
                    if not file_content: st.sidebar.warning(f"'{file.name}' is empty."); continue
                    files_data_for_processing.append({"name": file.name, "content": file_content, "type": file_ext})
                    new_files_count += 1
                except Exception as e: st.sidebar.error(f"Error reading {file.name}: {e}")

    # --- Document Processing Button ---
    if new_files_count > 0 and not st.session_state.processing_active and gemini_ready:
        if st.sidebar.button(f"Process {new_files_count} New File(s)", key="process_files_btn"):
            if files_data_for_processing:
                with st.sidebar: # Create status elements dynamically
                    st.session_state.status_container = st.status("Starting document processing...", expanded=True)
                    # Need to access the container object to add elements inside
                    with st.session_state.status_container:
                        st.session_state.progress_bar = st.progress(0)
                        st.session_state.time_info = st.empty()

                try:
                    st.session_state.processing_active = True; st.sidebar.warning("Processing started...")
                    # Run the async processing function
                    processed_results_list = run_async(
                        process_all_documents_async,
                        st.session_state.gemini_client_model,
                        files_data_for_processing,
                        status_container=st.session_state.status_container,
                        progress_bar=st.session_state.progress_bar,
                        time_info=st.session_state.time_info
                    )

                    # Update state with results
                    new_processed_docs = []; processed_names = set()
                    for r in processed_results_list:
                        if isinstance(r, dict) and "raw_extracted_content" in r:
                            new_processed_docs.append(r); fname = r["raw_extracted_content"].get("filename")
                            if fname: processed_names.add(fname)
                    st.session_state.processed_documents.extend(new_processed_docs)
                    st.session_state.processed_file_names.update(processed_names)
                    st.session_state.show_results = True; st.session_state.show_ontology = False # Show results tab
                    if st.session_state.status_container: st.session_state.status_container.update(label="✅ Processing finished!", state="complete", expanded=False)
                    st.sidebar.success(f"Processed {len(processed_names)} file(s).")
                except Exception as e:
                     st.sidebar.error(f"Processing error: {e}")
                     _logger.error(f"Doc processing error: {e}", exc_info=True)
                     if st.session_state.status_container: st.session_state.status_container.update(label=f"Error: {e}", state="error")
                finally:
                    st.session_state.processing_active = False
                    # Clean up UI element references from session state
                    for key in ["status_container", "progress_bar", "time_info"]:
                         if key in st.session_state: del st.session_state[key]
                    st.rerun() # Rerun to reflect changes
            else: st.sidebar.warning("No valid new files selected.")
    elif not gemini_ready:
         st.sidebar.warning("Gemini client not ready. Cannot process documents.")


    # --- Display Processed Files & View Buttons ---
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---"); st.sidebar.markdown("**Processed Files:**")
        for filename in sorted(list(st.session_state.processed_file_names)):
             has_error = any(
                 isinstance(d, dict) and d.get("raw_extracted_content", {}).get("filename") == filename and d.get("raw_extracted_content", {}).get("error")
                 for d in st.session_state.processed_documents
             )
             st.sidebar.caption(f"{'⚠️' if has_error else '✓'} {filename}")
        st.sidebar.markdown("---")

        col1, col2 = st.sidebar.columns(2)
        with col1:
            if st.button("View Reports", key="view_reports_btn", disabled=st.session_state.processing_active or not st.session_state.processed_documents):
                st.session_state.show_results = True; st.session_state.show_ontology = False; st.rerun()
        with col2:
            ontology_label = "View Ontology" if st.session_state.project_ontology else "Gen. Ontology"
            if st.button(ontology_label, key="ontology_btn", disabled=st.session_state.processing_active or not gemini_ready):
                # Filter for successfully processed docs
                valid_docs = [d for d in st.session_state.processed_documents if isinstance(d, dict) and not d.get("raw_extracted_content", {}).get("error")]
                if valid_docs:
                    if not st.session_state.project_ontology: # Generate if not exists
                        with st.spinner("Generating project ontology..."):
                            try:
                                ontology_data = run_async(
                                    incrementally_generate_doc_level_ontology,
                                    st.session_state.gemini_client_model,
                                    valid_docs
                                )
                                st.session_state.project_ontology = ontology_data
                                if ontology_data: st.sidebar.success("Ontology generated!")
                                else: st.sidebar.error("Ontology generation failed.")
                            except Exception as gen_err:
                                st.sidebar.error(f"Ontology error: {gen_err}")
                                _logger.error(f"Ontology gen failed: {gen_err}", exc_info=True)
                                st.session_state.project_ontology = None

                    # Show ontology if available (either pre-existing or just generated)
                    if st.session_state.project_ontology:
                        st.session_state.show_ontology = True; st.session_state.show_results = False; st.rerun()
                    elif not st.session_state.project_ontology and ontology_label == "View Ontology":
                         st.sidebar.warning("Ontology not available.") # Warn if trying to view non-existent

                else: st.sidebar.warning("No valid documents processed successfully for ontology generation.")


    # --- Qdrant 2-Phase Control ---
    st.sidebar.divider()
    st.sidebar.subheader("⬆️ Qdrant Indexing")
    st.sidebar.caption(f"Collection: `{QDRANT_COLLECTION_NAME}`")

    # Prerequisites for Qdrant operations
    can_start_qdrant = (
        qdrant_ready and models_ready and
        st.session_state.project_ontology is not None and
        isinstance(st.session_state.project_ontology.project_graph_nodes, list) and # Check it's a list
        st.session_state.project_ontology.project_graph_nodes and # Check it's not empty
        not st.session_state.processing_active # Ensure doc processing is not running
    )
    nodes_to_index = st.session_state.project_ontology.project_graph_nodes if can_start_qdrant else []

    # Phase 1a Button
    phase1a_tooltip = "Creates/Ensures collection exists (HNSW disabled)."
    if not qdrant_ready: phase1a_tooltip = "Qdrant client unavailable."
    elif st.session_state.processing_active: phase1a_tooltip = "Wait for document processing."
    phase1a_disabled = not qdrant_ready or st.session_state.processing_active
    if st.sidebar.button("Phase 1a: Create Collection", key="qdrant_phase1a_btn", disabled=phase1a_disabled, help=phase1a_tooltip):
        with st.spinner("Phase 1a: Ensuring collection exists..."):
            st.session_state.indexing_complete_phase1b = False; st.session_state.hnsw_enabled_phase2 = False; st.session_state.collection_has_data = False # Reset state
            success = run_async(create_collection_phase1a, st.session_state.qdrant_client)
            st.session_state.collection_created_phase1a = success
            if success: st.sidebar.success("Phase 1a Complete.")
            else: st.sidebar.error("Phase 1a Failed.")
        st.rerun()

    # Phase 1b Button
    phase1b_tooltip = "Uploads ontology nodes to Qdrant (UI blocks per group)."
    if not st.session_state.collection_created_phase1a: phase1b_tooltip = "Complete Phase 1a first."
    elif not can_start_qdrant: phase1b_tooltip = "Generate ontology with nodes first."
    elif not nodes_to_index: phase1b_tooltip = "No nodes found in ontology to ingest."
    phase1b_disabled = (not st.session_state.collection_created_phase1a or not can_start_qdrant or not nodes_to_index or st.session_state.processing_active)
    if st.sidebar.button("Phase 1b: Bulk Ingest Nodes", key="qdrant_phase1b_btn", disabled=phase1b_disabled, help=phase1b_tooltip):
        try:
            # Call the synchronous orchestrator which handles UI updates internally
            success, time_taken, has_data = run_bulk_ingest_foreground_nodes_with_ui(
                st.session_state.qdrant_client, nodes_to_index,
                st.session_state.dense_model, st.session_state.sparse_model
            )
            st.session_state.indexing_complete_phase1b = success
            st.session_state.collection_has_data = has_data
            st.session_state.hnsw_enabled_phase2 = False # Reset Phase 2
            if success: st.sidebar.success(f"Phase 1b Complete ({time_taken:.1f}s).")
            else: st.sidebar.error("Phase 1b Failed.")
        except Exception as e:
            st.sidebar.error(f"Phase 1b Error: {e}")

    # --- Phase 2: Enable HNSW (MISSING IMPLEMENTATION IN UI) ---
    phase2_tooltip = "Optimizes the collection by enabling the HNSW index (takes time)."
    if not st.session_state.indexing_complete_phase1b: phase2_tooltip = "Complete Phase 1b (Ingest) first."
    elif not st.session_state.collection_has_data: phase2_tooltip = "No data was ingested in Phase 1b to index."
    elif not qdrant_ready: phase2_tooltip = "Qdrant client unavailable."
    elif st.session_state.processing_active: phase2_tooltip = "Wait for document processing."

    # Disable if Phase 1b didn't succeed, no data, client not ready, or other processing active
    phase2_disabled = (
        not st.session_state.indexing_complete_phase1b or
        not st.session_state.collection_has_data or
        not qdrant_ready or
        st.session_state.processing_active
    )

    if st.sidebar.button("Phase 2: Enable HNSW Index", key="qdrant_phase2_btn", disabled=phase2_disabled, help=phase2_tooltip):
        with st.spinner("Phase 2: Enabling HNSW index... This can take time depending on data size."):
            # Run the async function to enable HNSW
            success = run_async(enable_hnsw_phase2, st.session_state.qdrant_client)
            st.session_state.hnsw_enabled_phase2 = success
            if success:
                st.sidebar.success("Phase 2 Complete: HNSW index enabled.")
            else:
                st.sidebar.error("Phase 2 Failed: Could not enable HNSW index.")
        st.rerun() # Rerun to update UI state (e.g., enable search)

    # Display Status Indicators for Qdrant Phases
    st.sidebar.markdown("---")
    st.sidebar.markdown("**Indexing Status:**")
    status_1a = "✅ Done" if st.session_state.collection_created_phase1a else "⏳ Pending"
    status_1b = "✅ Done" if st.session_state.indexing_complete_phase1b else "⏳ Pending"
    status_2 = "✅ Done" if st.session_state.hnsw_enabled_phase2 else "⏳ Pending"
    st.sidebar.caption(f"Phase 1a (Create): {status_1a}")
    st.sidebar.caption(f"Phase 1b (Ingest): {status_1b}")
    st.sidebar.caption(f"Phase 2 (Index): {status_2}")
    st.sidebar.markdown("---")

    # --- End of Qdrant Control Section ---
    # (Chat interface follows, which is present in the original code)
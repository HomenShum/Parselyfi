import streamlit as st
import os
import json
import asyncio
import base64
import pandas as pd
import pymupdf  # PyMuPDF
from typing import List, Dict, Optional, Union, Any, Set, Tuple
import pydantic
from pydantic import BaseModel, Field, ValidationError, field_validator, ConfigDict
from google import genai
from google.genai import types
import docx
from pptx import Presentation
import csv
import openpyxl
import re
from collections import Counter, defaultdict
import time
import io
import datetime
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import logging
import traceback
import pathlib
import random
from qdrant_client import AsyncQdrantClient, models
from qdrant_client.http.models import ( # Use http models for clarity in definition
    PointStruct, Distance, VectorParams, SparseVectorParams, SparseIndexParams,
    Filter, FieldCondition, MatchValue, Range, CollectionStatus, SparseVector
)
from fastembed import TextEmbedding, SparseTextEmbedding
import uuid

################# --- Qdrant Configuration (Add these near the top) ---
QDRANT_URL = st.secrets["qdrant_url"] # Or your Qdrant URL
QDRANT_API_KEY = st.secrets["qdrant_api_key"] # Optional: For Qdrant Cloud
QDRANT_COLLECTION_NAME = "project_ontology_nodes_v2" # Choose a suitable name, maybe versioned
EMB_MODEL_NAME = "BAAI/bge-base-en-v1.5" # Ensure this matches the model loaded
SPARSE_MODEL_NAME = "Qdrant/bm25"          # Ensure this matches the model loaded
EMB_DIM = 768                             # Dimension for bge-base-en-v1.5
# Configuration for upsert batching
BATCH_SIZE = 100 # How many points to send in one upsert call
INGEST_CONCURRENCY = 5 # How many batches to process per UI update cycle
HNSW_EF_SEARCH = 64
OVERSAMPLING = 4.0
DENSE_M_PROD = 16
SEARCH_PREFETCH_LIMIT = 50

################# --- Qdrant Client Factory (Cached) ---
@st.cache_resource
def get_qdrant_client():
    """Initializes and caches the AsyncQdrantClient."""
    logging.info(f"Initializing Qdrant client for {QDRANT_URL}")
    try:
        client = AsyncQdrantClient(
            url=QDRANT_URL,
            api_key=QDRANT_API_KEY,
            prefer_grpc=False, # Set to True if using gRPC
            timeout=60 # Adjust timeout as needed
        )
        # Optional: Perform a quick check if connection works
        # asyncio.run(client.get_collections()) # This blocks, better check elsewhere or handle async properly
        logging.info("Qdrant client initialized (connection not verified here).")
        return client
    except Exception as e:
        st.error(f"Failed to initialize Qdrant client: {e}")
        logging.error(f"Qdrant client initialization failed: {e}", exc_info=True)
        return None # Return None on failure

def get_or_create_eventloop():
    try:
        return asyncio.get_running_loop()
    except RuntimeError:
        logging.info("No running event loop found, creating a new one for this thread.")
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        return loop

###############################
# DATA MODELS
###############################

class VisualElement(BaseModel):
    """Model for visual elements like charts, graphs, etc."""
    type: str = Field(..., description="Type of visual element")
    description: str = Field(..., description="Description of the visual")
    data_summary: Optional[str] = Field(None, description="Summary of the data, or N/A") # Updated description to reflect N/A possibility
    page_numbers: List[int] = Field(default_factory=list, description="Pages where this appears")
    source_url: Optional[str] = Field(None, description="Source URL of the visual")
    alt_text: Optional[str] = Field(None, description="Alternative text for the visual")
    visual_id: str = Field(..., description="Unique identifier for the visual") # Made required as per prompt

class NumericalDataPoint(BaseModel):
    """Model for numerical data points extracted from text."""
    value: str = Field(..., description="The numerical value as a string") # Clarified string type
    description: str = Field(..., description="What the number represents, with units/context") # Made required
    context: Optional[str] = Field(None, description="Surrounding text context, if available") # Changed to Optional[str] and default to None

class TableData(BaseModel):
    """Model for tables extracted from documents."""
    table_content: str = Field(..., description="Markdown formatted table content") # Made required
    title: Optional[str] = Field(None, description="Optional table title") # Kept optional
    summary: Optional[str] = Field(None, description="Optional table summary") # Kept optional
    page_number: int = Field(..., description="Page number where the table appears") # Made required
    table_id: str = Field(..., description="Unique identifier for the table") # Made required

class Subsection(BaseModel):
    """Model for subsections extracted from pages."""
    subsection_id: str = Field(..., description="Unique identifier for the subsection")
    order: int = Field(..., description="Order of the subsection within the page")
    title: str = Field(..., description="Title of the subsection (less than 7 words)") # Added constraint hint
    text: str = Field(..., description="Full text content of the subsection") # Renamed from content to match prompt output
    description: str = Field(..., description="One line description summarizing the main point")  # Made required as per schema
    is_cutoff: bool = Field(..., description="True if content appears to be cut off by page break") # Made required
    referenced_visuals: List[str] = Field(default_factory=list, description="IDs of referenced visuals")
    referenced_tables: List[str] = Field(default_factory=list, description="IDs of referenced tables")

# --- Simplified PageContent Model ---
class PageContent(BaseModel):
    """
    Model for structured content extracted from a single page.
    Stage 1 populates raw_text, tables, visuals, numbers.
    Stage 2 populates subsections from raw_text.
    """
    page_number: int = Field(..., description="Page number in the original document")
    has_tables: bool = Field(..., description="Indicates if any tables were extracted from this page")
    has_visuals: bool = Field(..., description="Indicates if any visuals (charts, images, etc.) were extracted")
    has_numbers: bool = Field(..., description="Indicates if any key numerical data points were extracted")
    tables: List[TableData] = Field(default_factory=list, description="List of tables extracted from the page")
    visuals: List[VisualElement] = Field(default_factory=list, description="List of visual elements extracted from the page")
    numbers: List[NumericalDataPoint] = Field(default_factory=list, description="List of numerical data points extracted from the page")
    # NEW: Store raw text from initial extraction
    raw_text: Optional[str] = Field(None, description="Full raw text extracted initially from the page/chunk")
    # Subsections are now populated in a second step
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections extracted from raw_text")

    model_config = ConfigDict(arbitrary_types_allowed=True, extra='forbid')

# --- New Models for Qdrant Structure ---
class LinkedNode(BaseModel):
    """Structure for linked nodes metadata in Qdrant."""
    target_node_id: str = Field(..., description="ID of the target node (e.g., type_name_chX)")
    relationship_description: str = Field(..., description="Nature of the relationship from source to target")
    relationship_keywords: List[str] = Field(default_factory=list, description="Keywords summarizing the relationship")
    relationship_strength: float = Field(..., description="Strength of the relationship (1-10)")

    @field_validator('relationship_strength', mode='before')
    def strength_to_float(cls, v):
        try:
            return float(v)
        except (ValueError, TypeError):
            return 5.0 # Default strength on conversion error

class QdrantNode(BaseModel):
    """Structure for nodes to be stored in Qdrant, using direct fields."""
    # Core Identifiers
    node_id: str = Field(..., description="Unique node ID (e.g., UUID generated from type/name/context)")
    node_type: str = Field(..., description="'entity', 'concept', 'document', or 'meta_summary'") # Expanded types
    name: str = Field(..., description="Entity, concept, or document name")
    chapter_id_context: str = Field(..., description="Chapter ID or other context identifier (e.g., 'doc_level::filename')")

    # Extracted Content & Metadata Fields (Previously in properties)
    title: str = Field(..., description="Concise title for this node's content focus within its context (5-10 words)")
    core_claim: str = Field(..., description="1-3 sentence summary capturing the node's main point, assertion, or function.")
    information_type: str = Field(..., description="Classification of the main information type (e.g., 'Definition', 'Explanation', 'Process Description').")
    key_entities: str = Field(..., description="JSON stringified list of key entity names mentioned in context. Example: '[\"Acme Corp\", \"Project X\"]'. Use '[]' if none.")
    tags: str = Field(..., description="JSON stringified list of relevant keywords/tags (1-5 tags, lowercase, underscores). Example: '[\"ai\", \"optimization\"]'. Use '[]' if none.")
    sentiment_tone: str = Field(..., description="Overall sentiment/tone (e.g., 'Positive', 'Neutral', 'Objective').")
    source_credibility: str = Field(..., description="Inferred credibility (e.g., 'Company Report', 'Technical Documentation', 'Unknown').")

    # Fields added programmatically
    source: Optional[str] = Field(None, description="Source identifier (e.g., filename/URL), added programmatically.")
    stored_at: Optional[str] = Field(None, description="ISO timestamp of when the node was stored, added programmatically.")
    # Added field for original ID if needed (e.g. for doc/meta nodes)
    original_node_id: Optional[str] = Field(None, description="Original ID before UUID generation, if applicable.")
    # Added field for summary text specific to doc/meta nodes
    executive_summary: Optional[str] = Field(None, description="Executive summary, primarily for doc/meta nodes.")

    # Relationships
    linked_nodes: List[LinkedNode] = Field(default_factory=list, description="List of outgoing relationships")

    # Removed description @property
    # Removed properties dict

    model_config = ConfigDict(extra='ignore') # Allow ignoring extra fields if needed during validation

    # Add validator for key_entities and tags to ensure they are valid JSON strings
    @field_validator('key_entities', 'tags', mode='before')
    def check_json_string(cls, value):
        if isinstance(value, list): # Handle if LLM returns list instead of string
            try:
                return json.dumps(value)
            except TypeError:
                 raise ValueError("List provided for key_entities/tags could not be JSON serialized")
        if not isinstance(value, str):
            raise ValueError('Must be a string')
        try:
            json.loads(value) # Check if it's valid JSON
        except json.JSONDecodeError:
            raise ValueError('Must be a valid JSON string representation of a list')
        return value

class Chapter(BaseModel):
    """Model for chapters composed of subsections, including Qdrant node data."""
    chapter_id: str = Field(..., description="Unique identifier for the chapter")
    title: str = Field(..., description="Title of the chapter")
    summary: str = Field(..., description="Summary of the chapter")
    subsections: List[Subsection] = Field(default_factory=list, description="List of subsections in this chapter")
    # REMOVE THIS LINE:
    # entity_relationships: List[EntityRelationship] = Field(default_factory=list, description="Legacy entity relationships within the chapter")
    order: int = Field(..., description="Order of the chapter in the document")
    # New field to store structured node data for Qdrant
    qdrant_nodes: Optional[List[QdrantNode]] = Field(None, description="Structured nodes and relationships for Qdrant")
    
class DocumentSummaryDetails(BaseModel):
    """Structure for summary details generated from node analysis."""
    title: str = Field(..., description="Concise title for the document based on nodes")
    themes: List[str] = Field(default_factory=list, description="Main themes/topics derived from nodes")
    questions: List[str] = Field(default_factory=list, description="Sample questions reflecting node content")
    summary: str = Field(..., description="Comprehensive summary synthesizing node information")
    # Note: No chapters, entity_relationships, table/visual summaries here

class KeyTopic(BaseModel):
    name: str = Field(..., description="Short topic name (1-3 words)")
    description: str = Field(..., description="Brief description of the topic")
    relevance: str = Field(..., description="Relevance level (High/Medium/Low)")
    sentiment: str = Field(..., description="Sentiment analysis (Positive/Neutral/Mixed/Cautionary/Negative)")
    analysis: str = Field(..., description="Brief analysis of the topic")

class QuotedStatement(BaseModel):
    speaker: str = Field(..., description="Person or entity who made the statement")
    quote: str = Field(..., description="The quoted text")
    page: int = Field(..., description="Page number where the quote appears")

class DocumentReport(BaseModel):
    file_name: str = Field(..., description="Original filename")
    page_count: int = Field(..., description="Number of pages in document")
    title: str = Field(..., description="Document title")
    title_summary: str = Field(..., description="Brief summary/subtitle")
    concept_theme_hashtags: List[str] = Field(default_factory=list, description="Theme tags as hashtags")
    date_published: Optional[str] = Field(None, description="Publication date if available")
    source: Optional[str] = Field(None, description="Document source")
    confidence: str = Field("Medium", description="Analysis confidence (High/Medium/Low)")
    document_summary: str = Field(..., description="Comprehensive summary")
    key_insights: List[str] = Field(default_factory=list, description="Key takeaways and insights")
    key_topics: List[KeyTopic] = Field(default_factory=list, description="Detailed topic analysis")
    quoted_statements: List[QuotedStatement] = Field(default_factory=list, description="Important quotes from document")
    content_excerpt: Optional[str] = Field(None, description="Highlighted document content")
    chapters: List[Chapter] = Field(default_factory=list, description="Document chapters")

class ProcessedDocument(BaseModel):
    filename: str = Field(..., description="Original filename")
    pages: List[PageContent] = Field(..., description="Processed pages") # This now uses the simplified PageContent
    summary: Optional[DocumentSummaryDetails] = None

class ProjectOntology(BaseModel):
    title: str = Field(..., description="Project title")
    overview: str = Field(..., description="Project overview based on synthesized nodes")
    document_count: int = Field(..., description="Number of documents analyzed")
    documents: List[str] = Field(..., description="List of document filenames included in analysis")
    global_themes: List[str] = Field(..., description="High-level project themes derived from nodes")
    key_concepts: List[str] = Field(..., description="Key concepts identified across all documents from nodes")
    # NEW FIELD: Stores the aggregated nodes including inter-document relationships
    project_graph_nodes: Optional[List[QdrantNode]] = Field(None,
                                description="Aggregated nodes and their relationships (including inter-document links) across the project.")

    model_config = ConfigDict(arbitrary_types_allowed=True) # If QdrantNode uses non-standard types


###############################
# API AND UTILITY FUNCTIONS
###############################

def convert_to_model(data, model_class):
    """Convert dictionary data to a Pydantic model with error handling."""
    if not isinstance(data, dict):
        print(f"Warning: Expected dict for {model_class.__name__}, got {type(data)}")
        return None
    
    try:
        # Extract only the fields that the model expects
        model_fields = model_class.__annotations__.keys()
        filtered_data = {k: v for k, v in data.items() if k in model_fields}
        return model_class(**filtered_data)
    except Exception as e:
        print(f"Error converting to {model_class.__name__}: {str(e)}")
        return None

def get_gemini_api_key():
    """Get the Gemini API key from environment or secrets."""
    api_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        st.error("🚫 Gemini API key not found. Please set the GEMINI_API_KEY environment variable or in Streamlit secrets.")
        st.stop()
    return api_key

async def retry_api_call(func, *args, max_retries=3, **kwargs):
    """Retry API call with exponential backoff and JSON validation."""
    for attempt in range(max_retries):
        try:
            response = await func(*args, **kwargs)
            
            # Validate JSON if applicable
            config = kwargs.get('config')
            if (config and hasattr(config, 'response_mime_type') and 
                config.response_mime_type == 'application/json' and 
                hasattr(response, 'candidates') and response.candidates):
                try:
                    # Try to parse the JSON response
                    json_text = response.candidates[0].content.parts[0].text
                    json.loads(clean_json_response(json_text))  # Use our clean function first
                except json.JSONDecodeError as e:
                    if attempt < max_retries - 1:
                        print(f"Received malformed JSON on attempt {attempt+1}, retrying: {e}")
                        await asyncio.sleep(2 ** attempt)
                        continue
            
            return response
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            print(f"API call failed on attempt {attempt+1}, retrying: {e}")
            await asyncio.sleep(2 ** attempt)

def clean_json_response(json_text: str, extract_text_on_failure=True) -> str:
    """Clean Gemini JSON response with improved error handling and text extraction fallback."""
    if json_text is None:
        return "{}"
    
    # Ensure it's a string
    json_text = str(json_text)
    
    try:
        # Handle markdown code blocks
        if json_text.startswith("```"):
            blocks = json_text.split("```")
            if len(blocks) >= 3:
                json_text = blocks[1]
                if json_text.startswith("json"):
                    json_text = json_text[4:].strip()
            else:
                json_text = json_text.replace("```", "").strip()
        
        elif json_text.startswith("json"):
            json_text = json_text[4:].strip()
        
        # Remove any leading/trailing whitespace
        json_text = json_text.strip()
        
        # Attempt to parse the JSON as is
        try:
            json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # Continue with fixes
            pass
        
        # Common JSON fixes
        # Fix missing quotes around property names
        json_text = re.sub(r'([{,]\s*)(\w+)(\s*:)', r'\1"\2"\3', json_text)
        
        # Fix trailing commas in arrays/objects
        json_text = re.sub(r',(\s*[}\]])', r'\1', json_text)
        
        # Fix missing quotes around string values
        # This is trickier and might cause issues with valid JSON if not careful
        
        # Final attempt to parse with Python's JSON parser
        try:
            data = json.loads(json_text)
            return json_text
        except json.JSONDecodeError:
            # If we're extracting text on failure and can't fix the JSON
            if extract_text_on_failure:
                logging.warning(f"Initial JSON parsing failed. Raw Response Text: '{json_text[:2000]}...'") # Log raw response
                
                # Extract what appears to be the main text content
                # Look for patterns that would indicate text fields in a JSON response
                text_content = ""
                
                # Try to find text between quotes after "text": or "content": patterns
                text_matches = re.findall(r'"(?:text|content)"\s*:\s*"([^"]+)"', json_text)
                if text_matches:
                    text_content = " ".join(text_matches)
                
                # If that doesn't work, take any text between quotes as potential content
                if not text_content:
                    # Get all quoted strings that look reasonably long (over 20 chars)
                    quoted_text = re.findall(r'"([^"]{20,})"', json_text)
                    if quoted_text:
                        text_content = " ".join(quoted_text)
                
                # If still nothing, just take the longest line that's not mostly symbols
                if not text_content:
                    lines = json_text.split('\n')
                    content_lines = [line for line in lines if len(line) > 50 and 
                                    sum(c.isalpha() or c.isspace() for c in line) / len(line) > 0.7]
                    if content_lines:
                        text_content = " ".join(content_lines)
                
                # Fallback - create a minimal valid JSON with the extracted text
                fallback_json = {
                    "text": text_content or "Extracted text content not found",
                    "title": "Text Extraction Fallback",
                    "topics": ["text extraction"],
                    "summary": "Structured JSON parsing failed - basic text extraction used as fallback.",
                    "entities": [],
                    "has_tables": False,
                    "has_visuals": False,
                    "has_numbers": False,
                    "dates": [],
                    "tables": [],
                    "visuals": [],
                    "numbers": [],
                    "financial_statements": [],
                    "key_metrics": [],
                    "financial_terms": [],
                    "subsections": []
                }
                return json.dumps(fallback_json)
            else:
                # Return minimal valid JSON as before
                return '{"text":"Error parsing JSON response","error":"Invalid JSON structure"}'
            
    except Exception as e:
        logging.error(f"Error cleaning JSON response: {e}")
        
        if extract_text_on_failure:
            # Create a fallback with whatever raw content we have
            raw_text = json_text.strip()
            if len(raw_text) > 10000:  # Truncate very long responses
                raw_text = raw_text[:10000] + "... (truncated)"
                
            fallback_json = {
                "text": raw_text,
                "title": "Raw Text Fallback",
                "topics": ["raw text"],
                "summary": "Response processing failed - returning raw text.",
                "entities": [],
                "has_tables": False,
                "has_visuals": False,
                "has_numbers": False,
                "dates": [],
                "tables": [],
                "visuals": [],
                "numbers": [],
                "financial_statements": [],
                "key_metrics": [],
                "financial_terms": [],
                "subsections": []
            }
            return json.dumps(fallback_json)
        else:
            return '{"text":"Error cleaning JSON response","error":"' + str(e).replace('"', '\\"') + '"}'

async def run_in_executor(func, *args):
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, func, *args) # Use default executor

def run_async(func, *args, **kwargs):
    """Run an async function from Streamlit."""
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        # If func is already a coroutine, run it directly
        if asyncio.iscoroutine(func):
            return loop.run_until_complete(func)
        # Otherwise call it with the provided args
        return loop.run_until_complete(func(*args, **kwargs))
    finally:
        loop.close()

###############################
# CONTENT EXTRACTION FUNCTIONS
###############################

async def process_page_with_fallback(client, image_part, page_num, filename=None):
    """Process page with multi-level fallback including pure text extraction."""
    source_identifier = filename or "Unknown_Document"
    
    try:
        # First attempt: Full structured extraction
        logging.debug(f"Attempting full subsection extraction for page {page_num}")
        page_content = await extract_page_content_from_memory(client, image_part, page_num, source_identifier)

        # Check if it failed severely (e.g., returned an error page type, or text is placeholder)
        is_error = any(s.title == "Processing Error" for s in page_content.subsections) if page_content.subsections else False # Check error subsection
        if is_error or page_content.raw_text == "[LLM failed to extract raw text]":
             logging.warning(f"Initial structure extraction failed or yielded no text for page {page_num} of '{source_identifier}'. Triggering fallback.")
             raise ValueError("Initial extraction failed or returned no text") # Go to fallback

        logging.info(f"Successfully extracted structure (incl. raw_text) for page {page_num} of '{source_identifier}'")
        return page_content

    except Exception as e:
        logging.warning(f"Initial extraction failed for page {page_num} of '{source_identifier}': {e}. Trying pure text fallback.")

        # --- Fallback attempt: Pure text extraction ---
        try:
            logging.debug(f"Attempting pure text extraction fallback for page {page_num} of '{source_identifier}'")
            text_only_prompt = f"Extract ONLY the plain text content from page {page_num} of this document. Preserve paragraph breaks. Do NOT format as JSON."

            text_response = await retry_api_call(
                 client.aio.models.generate_content,
                 model="gemini-2.0-flash", # Model suitable for text extraction
                 contents=[
                    types.Content(parts=[image_part, types.Part.from_text(text=text_only_prompt)]),
                 ],
            )

            raw_text = ""
            if text_response.candidates:
                raw_text = text_response.candidates[0].content.parts[0].text.strip()

            if not raw_text:
                 raw_text = "[No text extracted during fallback]"
                 logging.warning(f"Pure text fallback for page {page_num} resulted in empty text.")

            logging.info(f"Successfully extracted pure text fallback for page {page_num} of '{source_identifier}'")
            # Create a PageContent object with ONLY raw_text and page_number populated
            return PageContent(
                page_number=page_num,
                raw_text=raw_text,
                # Set defaults for other required fields
                has_tables=False, has_visuals=False, has_numbers=False,
                tables=[], visuals=[], numbers=[], subsections=[]
            )

        except Exception as text_error:
            logging.error(f"Pure text extraction fallback also failed for page {page_num} of '{source_identifier}': {text_error}", exc_info=True)
            # Final fallback: return an error page object
            return create_error_page(page_num, f"All extraction methods failed: {text_error}")

# --- Define the NEW prompt for subsection extraction from text ---
subsection_extraction_prompt = """
Analyze the provided text content from a single page ({page_num}) of a document, originating from source '{source_identifier}'. Your goal is to segment this text into logical subsections, providing concise yet **informative descriptions geared towards identifying entities and relationships later**.

Input Text:
--- START TEXT ---
{raw_text_content}
--- END TEXT ---

Instructions:
1.  Read the text and identify logical breaks based on topic shifts, headings (if any within the text), or paragraph structure.
2.  Segment the text into sequential subsections.
3.  For each subsection, determine:
    *   `subsection_id`: Generate a unique ID formatted as `page_{page_num}_section_ORDER` where ORDER is the sequential order number.
    *   `order`: The sequential order number (integer, starting from 1).
    *   `title`: A concise title (less than 7 words). Use headings from the text if present and appropriate, otherwise generate a descriptive title.
    *   `text`: The full text content belonging *only* to this specific subsection. Ensure all original text is captured across subsections.
    *   `description`: **CRITICAL FOR RELATIONSHIP ANALYSIS:** A concise, one-sentence summary (approx. 15-25 words) that **highlights the key entities, concepts, or specific actions/interactions** discussed **within this specific subsection**. This description acts as a vital preview for downstream entity and relationship identification.
        *   **Focus:** Mention the main actors (e.g., 'Parsely', 'GPT-40', 'CafeCorner LLC') and the core action or relationship involving them *in this text block* (e.g., 'developed by', 'utilizes', 'integrates', 'enhances', 'piloted for').
        *   **Example Good Descriptions:** "Explains how Parsely integrates Llamalndex and Qdrant Vector DB.", "Details CafeCorner LLC's role in developing the Parsely system.", "Introduces GPT-40 as the LLM used for internal analysis.", "Discusses the enhancement of Pitchbook data distribution via Parsely.", "Outlines the piloting of Parsely for JPMorgan's ChatIQ assistant."
        *   **Example Less Useful Description:** "This section talks about the system." (Too generic - doesn't mention specific entities or actions).
    *   `is_cutoff`: Set to `true` ONLY if this specific text chunk appears to end mid-sentence or mid-thought, suggesting it was truncated *before* being passed to you. Otherwise `false`.
    *   `referenced_visuals`: An empty list `[]`. (Visuals/Tables extracted previously).
    *   `referenced_tables`: An empty list `[]`. (Visuals/Tables extracted previously).

4.  Return ONLY a valid JSON array containing the subsection objects, structured exactly like this example:
    ```json
    [
      {{
        "subsection_id": "page_{page_num}_section_1",
        "order": 1,
        "title": "Parsely System Overview",
        "text": "The full text content of the first subsection detailing Parsely...",
        "description": "Introduces the Parsely system developed at CafeCorner LLC for data enhancement.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }},
      {{
        "subsection_id": "page_{page_num}_section_2",
        "order": 2,
        "title": "Core Technologies",
        "text": "Details on the integration of various technologies...",
        "description": "Details Parsely's use of GPT-40, Llamalndex, and Qdrant Vector DB for analysis.",
        "is_cutoff": false,
        "referenced_visuals": [],
        "referenced_tables": []
      }}
      // ... more subsection objects
    ]
    ```
CRITICAL: Your entire response MUST be ONLY the JSON array, starting with `[` and ending with `]`. Do NOT include ```json markdown fences, introductory text, explanations, or any other characters outside the JSON structure. Ensure all string values are enclosed in double quotes and properly escaped. Validate your JSON structure carefully before outputting.
"""

# --- NEW Prompt for JSON Repair ---
json_repair_prompt_template = """
The following text was generated by an AI attempting to create a valid JSON array of subsection objects based on previous instructions. However, the text has syntax errors and is not valid JSON.

Your task is to analyze the provided text, correct the JSON syntax errors, and return ONLY the corrected, valid JSON array. Do NOT add any introductory text, explanations, or markdown formatting like ```json. The output must start with `[` and end with `]`.

Ensure:
- All strings are enclosed in double quotes.
- Internal quotes within strings are properly escaped (e.g., \\").
- Objects within the array are correctly separated by commas.
- There are no trailing commas after the last element in the array or the last property in an object.
- All brackets (`[]`) and braces (`{{}}`) are correctly paired and closed.

Faulty Text to Repair:
--- START FAULTY TEXT ---
{faulty_json_text}
--- END FAULTY TEXT ---

Return ONLY the valid JSON array:
"""


async def extract_subsections_from_text(
    client: genai.Client,
    raw_text: str,
    page_num: int,
    source_identifier: str,
    max_initial_retries: int = 2, # Retries for the initial API call
    retry_delay: int = 2
) -> List[Subsection]:
    """
    Takes raw text from a page and uses an LLM to segment it into Subsection objects.
    Includes robustness for dict-wrapped lists and a secondary LLM call to attempt
    JSON repair if initial parsing fails.
    Incorporates enhanced logging for debugging.

    Args:
        client: The configured GenAI client.
        raw_text: The raw text content of the page.
        page_num: The page number.
        source_identifier: Identifier of the source document (e.g., filename).
        max_initial_retries: Max attempts for the first LLM call.
        retry_delay: Base delay (seconds) between retries.

    Returns:
        A list of validated Subsection objects, or a list containing a single
        error fallback subsection if processing fails after all attempts.
    """
    function_name = "extract_subsections_from_text_with_repair" # Updated name
    logging.debug(f"[{function_name}] Starting subsection extraction for page {page_num} of '{source_identifier}'")

    # --- Input Validation ---
    if not raw_text or not raw_text.strip():
        logging.warning(f"[{function_name}] No raw text provided for page {page_num} of '{source_identifier}'. Returning empty list.")
        return []

    # --- Prepare Initial Prompt ---
    MAX_TEXT_LIMIT = 100000 # Adjust as needed
    text_for_prompt = raw_text
    if len(raw_text) > MAX_TEXT_LIMIT:
        logging.warning(f"[{function_name}] Truncating raw text for page {page_num} of '{source_identifier}' from {len(raw_text)} to {MAX_TEXT_LIMIT} chars.")
        text_for_prompt = raw_text[:MAX_TEXT_LIMIT] + "... [Content Truncated]"

    try:
        initial_prompt = subsection_extraction_prompt.format(
            page_num=page_num,
            raw_text_content=text_for_prompt,
            source_identifier=source_identifier
        )
    except KeyError as fmt_err:
         logging.error(f"[{function_name}] Error formatting initial prompt for page {page_num}: Missing key {fmt_err}. Using default prompt.")
         initial_prompt = f"Extract subsections as JSON for page {page_num}. Text: {text_for_prompt}"

    # --- Initial LLM Call and Retry Logic ---
    final_list_data: Optional[List[Dict]] = None # Variable to hold the successfully parsed list data
    last_exception: Optional[Exception] = None # Store the last exception encountered
    initial_call_failed_parsing = False # Flag if the first call resulted in parse failure

    for attempt in range(max_initial_retries):
        logging.debug(f"[{function_name}] Initial LLM Call Attempt {attempt + 1}/{max_initial_retries} for page {page_num}")
        last_exception = None # Reset exception for this attempt
        try:
            response = await retry_api_call( # Handles API-level retries (timeouts, server errors)
                client.aio.models.generate_content,
                # Consider using a more robust model here if issues persist
                model="gemini-2.5-flash-preview-04-17", 
                contents=[types.Content(parts=[types.Part.from_text(text=initial_prompt)])],
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.2 # Lower temperature for consistency
                    ),
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates from initial LLM call (Attempt {attempt + 1}) for page {page_num}.")
                # If no candidates after retries, we can't proceed to repair
                if attempt == max_initial_retries - 1:
                    final_list_data = None
                    initial_call_failed_parsing = False # Call itself failed, not parsing
                    last_exception = ValueError("LLM response had no candidates after retries.")
                # Continue to next retry if not last attempt
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue

            # --- Clean and Attempt Initial Parse ---
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"[{function_name}] Initial Cleaned JSON Response (Attempt {attempt + 1}, Page {page_num}): {json_text[:1000]}...") # Log more for debug

            try:
                # Check if it looks like a list before parsing fully
                if not isinstance(json_text, str) or not json_text.strip().startswith('['):
                     logging.warning(f"[{function_name}] Initial response doesn't start with '[' (Attempt {attempt+1}, Page {page_num}). Raw: {json_text[:200]}...")
                     # Treat this as a potential parsing failure for retry/repair logic
                     raise json.JSONDecodeError("Response does not start with '['", json_text, 0)

                subsections_data = json.loads(json_text)

                # --- Initial Structure Check ---
                if isinstance(subsections_data, list):
                    final_list_data = subsections_data # Correct structure received
                    initial_call_failed_parsing = False # Success!
                    logging.debug(f"[{function_name}] Initial call successful with list structure on attempt {attempt + 1} for page {page_num}.")
                    break # Success, exit retry loop

                elif isinstance(subsections_data, dict):
                    logging.warning(f"[{function_name}] Initial call returned a dict (Attempt {attempt + 1}, Page {page_num}). Checking common keys...")
                    possible_keys = ["subsections", "data", "items", "result", "page_subsections", "sections"]
                    found = False
                    for key in possible_keys:
                        if key in subsections_data and isinstance(subsections_data.get(key), list):
                            logging.info(f"[{function_name}] Found list under key '{key}' in dict. Using this list.")
                            final_list_data = subsections_data[key]
                            initial_call_failed_parsing = False # Recovered successfully
                            found = True
                            break
                    if found:
                        break # Success (recovered), exit retry loop
                    else:
                         # Could not find list within dict, treat as structural failure for retry/repair
                         logging.warning(f"[{function_name}] Could not find list within dict structure (Attempt {attempt+1}, Page {page_num}).")
                         raise json.JSONDecodeError("Dict received, but known list keys not found", json_text, 0)
                else:
                    # Incorrect type altogether
                    logging.warning(f"[{function_name}] Initial call response was not list or dict (type: {type(subsections_data)}) on attempt {attempt+1}, Page {page_num}.")
                    raise json.JSONDecodeError(f"Expected list or dict, got {type(subsections_data)}", json_text, 0)

            except json.JSONDecodeError as json_err:
                last_exception = json_err # Store the parsing error
                logging.warning(f"[{function_name}] Initial JSONDecodeError on attempt {attempt + 1} for page {page_num}: {json_err}. Raw text starts with: {json_text[:200]}...")
                if attempt < max_initial_retries - 1:
                    logging.info(f"[{function_name}] Retrying initial call after delay due to JSON decode error.")
                    await asyncio.sleep(retry_delay * (attempt + 1))
                    continue # Go to next initial call attempt
                else:
                    logging.error(f"[{function_name}] Failed to decode JSON from initial call after {max_initial_retries} attempts for page {page_num}.")
                    initial_call_failed_parsing = True # Mark that parsing failed on the last attempt
                    final_list_data = None # Ensure we don't proceed with bad data
                    # Store the faulty text for potential repair attempt
                    faulty_json_text_to_repair = json_text
                    break # Exit initial retry loop, proceed to repair attempt

        except Exception as e:
            last_exception = e # Store any other exception
            logging.error(f"[{function_name}] Error during initial LLM call/processing (Attempt {attempt + 1}) for page {page_num}: {e}", exc_info=True)
            if attempt < max_initial_retries - 1:
                logging.info(f"[{function_name}] Retrying initial call after delay due to unexpected error.")
                await asyncio.sleep(retry_delay * (attempt + 1))
                continue
            else:
                 logging.error(f"[{function_name}] Initial call failed after {max_initial_retries} attempts for page {page_num} due to error: {e}")
                 final_list_data = None # Signal failure
                 initial_call_failed_parsing = False # The call itself failed, not parsing
                 break # Exit retry loop

    # --- JSON Repair Attempt (if initial parsing failed on last try) ---
    if initial_call_failed_parsing and faulty_json_text_to_repair:
        logging.info(f"[{function_name}] Initial JSON parsing failed for page {page_num}. Attempting self-repair LLM call.")
        try:
            repair_prompt = json_repair_prompt_template.format(faulty_json_text=faulty_json_text_to_repair)

            repair_response = await retry_api_call( # Use retry for the repair call too (1 attempt usually enough)
                client.aio.models.generate_content,
                # Use a capable model for repair
                model="gemini-2.0-flash", 
                contents=[types.Content(parts=[types.Part.from_text(text=repair_prompt)])],
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.1 # Very low temperature for repair
                ),
                max_retries=1 # Only retry the repair call once on API errors
            )

            if repair_response and repair_response.candidates:
                repaired_json_text = clean_json_response(repair_response.candidates[0].content.parts[0].text)
                logging.debug(f"[{function_name}] Repaired JSON Response (Page {page_num}): {repaired_json_text[:1000]}...")
                try:
                    # Attempt to parse the REPAIRED text
                    repaired_data = json.loads(repaired_json_text)
                    if isinstance(repaired_data, list):
                        logging.info(f"[{function_name}] Successfully parsed JSON after repair call for page {page_num}.")
                        final_list_data = repaired_data # Use the repaired data
                    elif isinstance(repaired_data, dict) and "subsections" in repaired_data and isinstance(repaired_data["subsections"], list):
                        # Handle if repair somehow wrapped it again
                        logging.warning(f"[{function_name}] Repaired JSON was dict with 'subsections' list for page {page_num}. Using list.")
                        final_list_data = repaired_data["subsections"]
                    else:
                        logging.error(f"[{function_name}] Repaired JSON was not a list or expected dict structure for page {page_num}. Type: {type(repaired_data)}")
                        final_list_data = None # Repair failed structurally
                        last_exception = json.JSONDecodeError("Repaired JSON structure invalid", repaired_json_text, 0)

                except json.JSONDecodeError as repair_err:
                    logging.error(f"[{function_name}] Failed to decode JSON even after repair attempt for page {page_num}: {repair_err}")
                    final_list_data = None # Repair failed parsing
                    last_exception = repair_err # Store the repair parsing error
            else:
                logging.error(f"[{function_name}] Repair LLM call returned no candidates for page {page_num}.")
                final_list_data = None # Repair call failed
                last_exception = ValueError("Repair LLM call had no candidates.")

        except Exception as repair_call_err:
            logging.error(f"[{function_name}] Error during repair LLM call for page {page_num}: {repair_call_err}", exc_info=True)
            final_list_data = None # Repair call failed
            last_exception = repair_call_err


    # --- Process the final data (if successful) or fallback ---
    if final_list_data is not None and isinstance(final_list_data, list): # Check if we successfully got a list
        validated_subsections: List[Subsection] = []
        logging.debug(f"[{function_name}] Processing {len(final_list_data)} potential subsection items for page {page_num}.")
        items_skipped_validation = 0
        for i, sub_dict in enumerate(final_list_data):
            # Skip items that are not dicts or have empty text
            if not isinstance(sub_dict, dict):
                logging.warning(f"[{function_name}] Skipping item {i+1} (not a dict) for page {page_num}: {str(sub_dict)[:100]}...")
                items_skipped_validation += 1
                continue
            subsection_text = sub_dict.get("text", "")
            if not subsection_text or not subsection_text.strip():
                logging.warning(f"[{function_name}] Skipping item {i+1} with missing/empty 'text' for page {page_num}: {str(sub_dict)[:100]}...")
                items_skipped_validation += 1
                continue

            # Proceed with validation
            try:
                # Set defaults before validation
                sub_dict.setdefault("referenced_visuals", [])
                sub_dict.setdefault("referenced_tables", [])
                sub_dict.setdefault("is_cutoff", False)

                # Ensure order is int, default to index
                try:
                    sub_dict['order'] = int(sub_dict.get('order', i + 1))
                except (ValueError, TypeError):
                    logging.warning(f"[{function_name}] Converting order '{sub_dict.get('order')}' to int failed for item {i+1} page {page_num}. Using index {i+1}.")
                    sub_dict['order'] = i + 1

                subsection_obj = Subsection(**sub_dict)
                validated_subsections.append(subsection_obj)

            except ValidationError as ve:
                logging.warning(f"[{function_name}] Pydantic validation failed for subsection item {i+1} on page {page_num}: {ve}. Skipping. Data: {sub_dict}")
                items_skipped_validation += 1
            except Exception as val_err:
                logging.error(f"[{function_name}] Unexpected error during subsection validation {i+1} on page {page_num}: {val_err}. Skipping. Data: {sub_dict}", exc_info=True)
                items_skipped_validation += 1

        # Post-process validated subsections
        if validated_subsections:
            # Sort by order
            validated_subsections.sort(key=lambda s: getattr(s, 'order', float('inf')))

            # Re-assign sequential order
            for i, sub in enumerate(validated_subsections):
                if sub.order != i + 1:
                    logging.debug(f"[{function_name}] Re-assigning order for subsection '{sub.subsection_id}' from {sub.order} to {i+1} on page {page_num}.")
                    sub.order = i + 1

            logging.info(f"[{function_name}] Successfully extracted and validated {len(validated_subsections)} non-empty subsections for page {page_num} (skipped {items_skipped_validation} during validation).")
            return validated_subsections
        else:
            # Fallback if validation filtered everything out
            logging.warning(f"[{function_name}] Processing resulted in zero valid subsections after validation for page {page_num}. Using fallback.")
            return [create_error_fallback_subsection(page_num, raw_text, "LLM response processed, but no valid/non-empty subsections remained after validation.")]

    else:
        # Fallback if all initial attempts AND the repair attempt failed
        fail_reason = "Initial LLM calls failed or returned invalid structure."
        if initial_call_failed_parsing:
             fail_reason = "Initial JSON parsing failed, and subsequent repair attempt also failed."
        if last_exception:
             fail_reason += f" Last known error: {type(last_exception).__name__}: {str(last_exception)[:100]}..." # Add error context

        logging.error(f"[{function_name}] All subsection extraction attempts failed for page {page_num}. Using fallback. Reason: {fail_reason}")
        logging.debug(f"[{function_name}] Input raw_text for failed page {page_num}:\n--- START RAW TEXT ---\n{raw_text[:1000]}...\n--- END RAW TEXT ---")
        return [create_error_fallback_subsection(page_num, raw_text, fail_reason)]


def create_error_fallback_subsection(page_num: int, raw_text: str, error_msg: str) -> Subsection:
     """
     Creates a single fallback Subsection object when subsection extraction fails.
     This ensures the raw text is preserved along with the error context.
     """
     fallback_id = f"page_{page_num}_section_fallback_error"
     fallback_title = "Full Page Content (Subsection Extraction Failed)"
     # Combine error and original text for context
     fallback_text = (
         f"--- ERROR DURING SUBSECTION EXTRACTION ---\n"
         f"{error_msg}\n"
         f"--- END ERROR ---\n\n"
         f"--- ORIGINAL RAW TEXT FOR PAGE {page_num} ---\n"
         f"{raw_text}\n"
         f"--- END RAW TEXT ---"
     )
     fallback_description = f"Could not segment page {page_num} into subsections due to error: {error_msg[:100]}..." # Truncate long errors

     try:
          # Create and return a single Subsection object
          return Subsection(
              subsection_id=fallback_id,
              order=1, # Always the first (and only) subsection in this fallback case
              title=fallback_title,
              text=fallback_text,
              description=fallback_description,
              is_cutoff=False, # Assume not cutoff unless original text was truncated
              referenced_visuals=[], # No references identified
              referenced_tables=[]  # No references identified
          )
     except ValidationError as ve_fb:
         # This should be rare if the inputs (page_num, raw_text, error_msg) are valid types
         logging.error(f"[create_error_fallback_subsection] Critical: Failed to create even the fallback Subsection for page {page_num}: {ve_fb}", exc_info=True)
         # Re-raise the error as we cannot create a valid object, indicating a deeper issue
         raise ve_fb
     except Exception as e:
        # Catch any other unexpected errors during creation
        logging.error(f"[create_error_fallback_subsection] Unexpected error creating fallback Subsection for page {page_num}: {e}", exc_info=True)
        # Depending on desired behavior, could raise or return a dummy/error object
        # Re-raising is often safer to signal the failure clearly
        raise e


# Prompt for initial structure extraction (raw_text + elements)
# --- New Task-Oriented Prompt for Initial Page Structure Extraction ---

page_structure_prompt = """
--- Goal ---
Analyze the input, which represents a single page ({page_num}) of a document (could be an image or text). Your primary task is to extract the complete raw text content AND identify and structure specific elements like tables, visuals, and key numerical data points found on that page.

--- Input ---
- Content of page {page_num} (provided as image or text).

--- Instructions ---
Perform the following tasks based *only* on the provided input content for page {page_num}:

**Task 1: Extract Raw Text**
1a. Extract ALL text content visible on the page.
1b. Preserve original formatting like paragraphs and line breaks as accurately as possible.
1c. Store this complete text in the `raw_text` field of the output JSON.

**Task 2: Extract Tables**
2a. Identify all distinct tables present on the page.
2b. For EACH table found:
    *   Extract the complete cell content accurately.
    *   Format this content STRICTLY as GitHub Flavored Markdown within the `table_content` field (using '|' separators and a '---' header separator line). Do NOT include surrounding paragraphs or explanatory text in `table_content`.
    *   If a clear title or caption is associated with the table, extract it into the `title` field. Otherwise, use `null`.
    *   Write a brief (1-2 sentence) `summary` of the table's main content or purpose, if possible. Otherwise, use `null`.
    *   Generate a unique `table_id` formatted as `page_{page_num}_table_N` (where N is 1, 2, 3...).
    *   Set `page_number` to {page_num}.
    *   Add the complete table object to the `tables` list in the output JSON.
2c. If NO tables are found on the page, the `tables` list MUST be an empty list `[]`.

**Task 3: Extract Visual Elements**
3a. Identify all visual elements (e.g., charts, graphs, diagrams, images, photographs, equations, code blocks).
3b. For EACH visual element found:
    *   Determine its `type` (e.g., "bar chart", "line graph", "diagram", "image", "equation", "code block").
    *   Write a detailed `description` covering what the visual shows, its purpose, and key elements.
    *   Provide a `data_summary` summarizing key data points, trends, or findings presented in the visual. If the visual doesn't present data (e.g., a decorative image) or a summary isn't applicable, use the string "N/A".
    *   Generate a unique `visual_id` formatted as `page_{page_num}_visual_N` (where N is 1, 2, 3...).
    *   Set `page_numbers` to `[{page_num}]`.
    *   Set `source_url` and `alt_text` to `null` unless explicitly available.
    *   Add the complete visual object to the `visuals` list in the output JSON.
3c. If NO visual elements are found, the `visuals` list MUST be an empty list `[]`.

**Task 4: Extract Numerical Data**
4a. Identify significant numerical data points mentioned in the page's text (excluding those already inside tables).
4b. For EACH significant number found:
    *   Extract the `value` as a string (e.g., "123.45", "50%", "$1.2M").
    *   Write a `description` explaining what the number represents, including units or context (e.g., "increase in revenue", "percentage completion").
    *   Extract the surrounding text snippet (~10-20 words) providing immediate context into the `context` field. If no clear surrounding context is available, use `null`.
    *   Add the complete number object to the `numbers` list in the output JSON.
4c. If NO significant numerical data points are found, the `numbers` list MUST be an empty list `[]`.

**Task 5: Set Boolean Flags**
5a. Set `has_tables` to `true` if the `tables` list is NOT empty, otherwise set it to `false`.
5b. Set `has_visuals` to `true` if the `visuals` list is NOT empty, otherwise set it to `false`.
5c. Set `has_numbers` to `true` if the `numbers` list is NOT empty, otherwise set it to `false`.

**Task 6: Assemble Final JSON**
6a. Combine all extracted information (`raw_text`, `has_tables`, `tables`, `has_visuals`, `visuals`, `has_numbers`, `numbers`) into a single JSON object adhering precisely to the schema defined below.

--- Output Format ---
Return ONLY a single, valid JSON object. Do NOT include any text, comments, or markdown formatting before or after the JSON object. It MUST strictly follow this schema:

```json
{{
  "raw_text": "string", // REQUIRED. Full text from the page.
  "has_tables": "boolean", // REQUIRED. True if 'tables' list is not empty.
  "has_visuals": "boolean", // REQUIRED. True if 'visuals' list is not empty.
  "has_numbers": "boolean", // REQUIRED. True if 'numbers' list is not empty.
  "tables": [ // REQUIRED list (can be empty []).
    {{
      "table_id": "string", // REQUIRED. e.g., "page_{page_num}_table_1"
      "table_content": "string", // REQUIRED. Table content as GitHub Flavored Markdown.
      "title": "string | null", // Optional
      "summary": "string | null", // Optional
      "page_number": "integer" // REQUIRED. e.g., {page_num}
    }}
  ],
  "visuals": [ // REQUIRED list (can be empty []).
    {{
      "visual_id": "string", // REQUIRED. e.g., "page_{page_num}_visual_1"
      "type": "string", // REQUIRED. e.g., "bar chart", "image"
      "description": "string", // REQUIRED. Detailed description.
      "data_summary": "string", // REQUIRED. Summary of data or "N/A".
      "page_numbers": ["integer"], // REQUIRED. e.g., [{page_num}]
      "source_url": "string | null", // Optional
      "alt_text": "string | null" // Optional
    }}
  ],
  "numbers": [ // REQUIRED list (can be empty []).
    {{
      "value": "string", // REQUIRED. Numerical value as string.
      "description": "string", // REQUIRED. What the number represents.
      "context": "string | null" // Optional. Surrounding text.
    }}
  ]
}}

--- Critical Rules ---
VALID JSON ONLY: The entire output MUST be a single JSON object starting with {{ and ending with }}.
SCHEMA ADHERENCE: Follow the schema EXACTLY. All REQUIRED fields must be present with the correct data types. Use null for optional string fields if no value is applicable/found. Use [] for empty lists.
MARKDOWN TABLES: Table content MUST be formatted as GitHub Flavored Markdown in the table_content field.
TEXTUAL BASIS: All extracted information must be derived SOLELY from the provided page content. Do not infer or add external knowledge.

Generate the JSON output for page {page_num} now:
"""


async def extract_structure_from_text_content(client: genai.Client, text_content: str, page_num: int, filename: str) -> PageContent:
    """
    Uses Gemini to extract raw text and structural elements (tables, visuals, numbers)
    from the input TEXT content. Populates PageContent WITHOUT subsections initially.
    Includes fallback for pure text preservation.
    """
    function_name = "extract_structure_from_text_content"
    logging.debug(f"[{function_name}] Starting initial structure extraction from text for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    if not text_content or not text_content.strip():
         logging.warning(f"[{function_name}] Input text_content is empty for page {page_num} of '{filename}'. Returning basic PageContent.")
         page_data["raw_text"] = "" # Set raw text to empty string
         try:
              return PageContent(**page_data)
         except ValidationError as ve:
             logging.error(f"[{function_name}] Failed validation even for empty text page {page_num}: {ve}")
             # Return hardcoded error if basic creation fails
             return create_error_page(page_num, f"Validation error on empty text input: {ve}")


    # Limit text content size
    MAX_TEXT_LENGTH = 100000
    if len(text_content) > MAX_TEXT_LENGTH:
        logging.warning(f"[{function_name}] Truncating text input for page {page_num} of '{filename}' from {len(text_content)} to {MAX_TEXT_LENGTH} chars.")
        text_content_for_api = text_content[:MAX_TEXT_LENGTH] + "\n... [Content Truncated] ..."
    else:
        text_content_for_api = text_content

    try:
        prompt_for_page = page_structure_prompt.format(page_num=page_num)

        full_prompt_with_data = f"""
        {prompt_for_page}
        
        {text_content_for_api}
        """

        # --- Attempt structured extraction (raw_text + elements) from text ---
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[types.Part.from_text(text=full_prompt_with_data)]),
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {}
            if not isinstance(data, dict): data = {}

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text") # Get raw text from LLM's perspective
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False) # Less likely but check
            page_data["has_numbers"] = data.get("has_numbers", False)
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])
            page_data["subsections"] = [] # Explicitly ensure subsections are empty here

            # If LLM failed to echo raw_text, use the original input text as fallback
            if not page_data["raw_text"]:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num}. Using original input text.")
                 page_data["raw_text"] = text_content # Use original full text

            # Create PageContent object (validates tables, visuals, numbers)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted initial structure from text for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for text structure extraction on page {page_num}.")
            raise ValueError("LLM response had no candidates") # Trigger fallback

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.warning(f"[{function_name}] Initial structure extraction from text failed for page {page_num}: {e}. Falling back to text preservation.")

        # --- Fallback: Preserve original text ---
        # Create a PageContent object with ONLY raw_text and page_number populated
        page_data["raw_text"] = text_content # Ensure original text is saved
        page_data["has_tables"] = False # Reset flags
        page_data["has_visuals"] = False
        page_data["has_numbers"] = False
        page_data["tables"] = []
        page_data["visuals"] = []
        page_data["numbers"] = []
        page_data["subsections"] = []

        try:
             # Return a valid PageContent object even in fallback
             return PageContent(**page_data)
        except ValidationError as ve_fallback:
            logging.error(f"[{function_name}] Failed to create fallback PageContent for text page {page_num}: {ve_fallback}")
            return create_error_page(page_num, f"Text structure extraction failed AND fallback creation failed: {ve_fallback}")


async def extract_page_content_from_memory(client, image_part, page_num, filename=None):
    """
    Extracts raw text and structural elements (tables, visuals, numbers) from a single page image.
    Populates PageContent WITHOUT subsections initially.
    """
    function_name = "extract_page_structure" # For logging
    logging.debug(f"[{function_name}] Starting extraction for page {page_num} of '{filename}'")

    # Initialize with defaults, note subsections is empty
    page_data = {
        "page_number": page_num, "has_tables": False, "has_visuals": False, "has_numbers": False,
        "tables": [], "visuals": [], "numbers": [], "raw_text": None, "subsections": []
    }

    page_structure_prompt = f"""
    Analyze page {page_num} from the input document/image. Extract the full raw text content and identify structural elements like tables, visuals, and key numbers. Return ONLY a valid JSON object with the exact structure specified below.

    JSON Structure:
    {{
    "raw_text": "Complete text content extracted from the page. Preserve formatting like paragraphs and line breaks.",
    "has_tables": true/false,
    "has_visuals": true/false,
    "has_numbers": true/false,
    "tables": [ {{ "table_id": "page_{page_num}_table_N", "table_content": "markdown...", "title": "optional", "summary": "optional", "page_number": {page_num} }} ],
    "visuals": [ {{ "visual_id": "page_{page_num}_visual_N", "type": "...", "description": "...", "data_summary": "...", "page_numbers": [{page_num}], "source_url": null, "alt_text": null }} ],
    "numbers": [ {{ "value": "...", "description": "...", "context": "..." }} ]
    }}

    Rules & Guidelines:
    1.  **Raw Text:** Extract ALL text content into the `raw_text` field. Maintain original paragraph structure.
    2.  **Elements:** Accurately identify and extract all tables, visuals (charts, images, diagrams, equations), and significant numbers.
    3.  **IDs:** Generate unique IDs for tables and visuals starting N from 1 for the page.
    4.  **Structure:** Adhere STRICTLY to the specified JSON structure. Use `[]` for empty lists.
    5.  **JSON Validity:** Ensure the output is ONLY a single, valid JSON object.

    EXAMPLE JSON Schema (Conceptual):
    {{
    "type": "object",
    "properties": {{
        "raw_text": {{ "type": ["string", "null"] }},
        "has_tables": {{ "type": "boolean" }},
        "has_visuals": {{ "type": "boolean" }},
        "has_numbers": {{ "type": "boolean" }},
        "tables": {{ "type": "array", "items": {{ "$ref": "#/definitions/TableData" }} }},
        "visuals": {{ "type": "array", "items": {{ "$ref": "#/definitions/VisualElement" }} }},
        "numbers": {{ "type": "array", "items": {{ "$ref": "#/definitions/NumericalDataPoint" }} }}
    }},
    "required": ["raw_text", "has_tables", "has_visuals", "has_numbers", "tables", "visuals", "numbers"]
    // Define TableData, VisualElement, NumericalDataPoint schemas conceptually if needed
    }}
    """

    try:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.0-flash",
            contents=[
                types.Content(parts=[image_part, types.Part.from_text(text=page_structure_prompt)]),
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )



        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            data = json.loads(json_text)

            if isinstance(data, list): data = data[0] if data else {} # Handle list response
            if not isinstance(data, dict): data = {} # Ensure data is a dict

            # Populate page_data from the LLM response
            page_data["raw_text"] = data.get("raw_text")
            page_data["has_tables"] = data.get("has_tables", False)
            page_data["has_visuals"] = data.get("has_visuals", False)
            page_data["has_numbers"] = data.get("has_numbers", False)
            # Pydantic will validate items during PageContent creation
            page_data["tables"] = data.get("tables", [])
            page_data["visuals"] = data.get("visuals", [])
            page_data["numbers"] = data.get("numbers", [])

            # Check if raw_text was extracted
            if page_data["raw_text"] is None:
                 logging.warning(f"[{function_name}] LLM did not return 'raw_text' for page {page_num} of '{filename}'.")
                 page_data["raw_text"] = "[LLM failed to extract raw text]" # Provide placeholder

            # Create PageContent object (will validate tables, visuals, numbers if models defined)
            page_content_obj = PageContent(**page_data)
            logging.info(f"[{function_name}] Successfully extracted structure for page {page_num} of '{filename}'.")
            return page_content_obj

        else:
            logging.error(f"[{function_name}] No candidates from LLM for page {page_num} of '{filename}'.")
            # Return error page content (subsections will be empty)
            return create_error_page(page_num, "LLM response had no candidates.")

    except (json.JSONDecodeError, ValidationError, Exception) as e:
        logging.error(f"[{function_name}] Failed for page {page_num} of '{filename}': {e}", exc_info=True)
        # Return error page content (subsections will be empty)
        error_page = create_error_page(page_num, f"Initial extraction failed: {e}")
        # Ensure raw_text is None or set appropriately in error page if needed downstream
        error_page.raw_text = page_data.get("raw_text") # Preserve if extracted before error
        return error_page


def create_error_page(page_num: int, error_msg: str, validation_errors: Optional[List[Dict[str, Any]]] = None) -> PageContent:
    """
    Creates a simplified PageContent object representing an error state,
    placing the error details within a single subsection.
    """
    logging.error(f"Creating error page for page {page_num} due to: {error_msg}")
    # --- Format the detailed error text ---
    error_text = f"Error processing page {page_num}: {error_msg}"

    if validation_errors:
        # Ensure validation_errors is a list of dicts before processing
        if isinstance(validation_errors, list) and all(isinstance(item, dict) for item in validation_errors):
            try:
                error_details = "\n\nValidation Errors:\n" + "\n".join(
                    # Safely access keys with .get()
                    f"- Field '{'.'.join(map(str, err.get('loc', ['unknown']))) if err.get('loc') else 'unknown'}': {err.get('msg', 'No message')}"
                    for err in validation_errors
                )
                error_text += error_details
                logging.debug(f"Formatted validation errors for page {page_num}: {error_details}")
            except Exception as format_error:
                 logging.error(f"Could not format validation errors for page {page_num}: {format_error}")
                 error_text += f"\n(Could not format validation errors: {format_error})"
        else:
            logging.warning(f"Received validation_errors in unexpected format for page {page_num}: {type(validation_errors)}. Skipping details.")
            error_text += "\n(Validation error details format incorrect)"

    # --- Create the error subsection ---
    error_subsection = Subsection(
        subsection_id=f"page_{page_num}_section_error",
        order=1,
        title="Processing Error",
        text=error_text, # Place the detailed error message here
        description=f"An error occurred processing page {page_num}.", # Brief summary
        is_cutoff=False,
        referenced_visuals=[],
        referenced_tables=[]
    )

    # --- Initialize the simplified PageContent object ---
    error_page = PageContent(
        page_number=page_num,
        has_tables=False,
        has_visuals=False,
        has_numbers=False,
        tables=[],
        visuals=[],
        numbers=[],
        subsections=[error_subsection] # Put the error subsection in the list
        # Removed old fields: text, title, topics, summary, entities, dates, etc.
    )

    logging.debug(f"Generated error PageContent object for page {page_num}")
    return error_page

async def merge_cutoff_subsections(pages: List[PageContent]) -> List[PageContent]:
    """
    Merge subsections marked as cut off with the first subsection of the next page.
    Operates on the simplified PageContent model.
    """
    if not pages or len(pages) < 2:
        return pages # Nothing to merge if less than 2 pages

    # Ensure pages are sorted by page number
    # Use getattr for safe access in case of malformed objects, though Pydantic should ensure existence
    try:
        sorted_pages = sorted(pages, key=lambda p: getattr(p, 'page_number', float('inf')))
    except Exception as sort_err:
        logging.error(f"Failed to sort pages for merging: {sort_err}. Returning original order.")
        sorted_pages = pages # Proceed with original order if sorting fails


    logging.debug(f"Attempting to merge subsections across {len(sorted_pages)} pages.")
    merged_something = False # Flag to track if any merge happened

    # Iterate through pages up to the second-to-last one
    for i in range(len(sorted_pages) - 1):
        current_page = sorted_pages[i]
        next_page = sorted_pages[i + 1]

        # Basic validation of pages and subsections lists
        if not isinstance(current_page, PageContent) or not isinstance(next_page, PageContent):
             logging.warning(f"Skipping merge check between page {i+1} and {i+2} due to invalid PageContent object types.")
             continue

        if not current_page.subsections or not next_page.subsections:
            # logging.debug(f"Skipping merge check between page {current_page.page_number} and {next_page.page_number}: one has no subsections.")
            continue # Nothing to merge if either page lacks subsections

        # Sort subsections by order for safety, though they should be ordered
        current_subsections = sorted(current_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))
        next_subsections = sorted(next_page.subsections, key=lambda s: getattr(s, 'order', float('inf')))

        # Get the last subsection of the current page and first of the next
        last_subsection = current_subsections[-1]
        first_next_subsection = next_subsections[0]

        # Check if the last subsection is actually marked as cut off
        if getattr(last_subsection, 'is_cutoff', False):
            logging.info(f"Merging cutoff subsection '{last_subsection.subsection_id}' (Page {current_page.page_number}) with '{first_next_subsection.subsection_id}' (Page {next_page.page_number})")

            # --- Perform the merge ---
            # Append text (Use space or newline as appropriate)
            last_subsection.text += "\n" + getattr(first_next_subsection, 'text', '') # Use getattr for safety

            # Mark the merged subsection as not cut off
            last_subsection.is_cutoff = False

            # Merge references (handle potential missing attributes)
            last_subsection.referenced_visuals.extend(getattr(first_next_subsection, 'referenced_visuals', []))
            last_subsection.referenced_tables.extend(getattr(first_next_subsection, 'referenced_tables', []))
            # Ensure uniqueness if needed:
            # last_subsection.referenced_visuals = list(set(last_subsection.referenced_visuals))
            # last_subsection.referenced_tables = list(set(last_subsection.referenced_tables))


            # --- Update the next page ---
            # Remove the first subsection from the next page
            next_page.subsections = next_subsections[1:]

            # Renumber the order of remaining subsections in the next page
            for idx, subsection in enumerate(next_page.subsections):
                 subsection.order = idx + 1

            merged_something = True # Mark that a merge occurred

    if merged_something:
         logging.info("Finished merging cut-off subsections.")
    else:
         logging.debug("No cut-off subsections found to merge.")

    return sorted_pages # Return the list of pages (potentially modified)

async def extract_chapters_from_subsections(client: genai.Client, pages: List[PageContent]) -> List[Chapter]:
    """
    Groups subsections from multiple pages into logical chapters using an LLM.
    Input `pages` should already have cut-off subsections merged.
    """
    all_subsections_with_context = []
    subsection_map = {} # For quick lookup later

    logging.debug(f"Extracting subsections from {len(pages)} pages for chapter generation.")
    for page in pages:
        if not isinstance(page, PageContent) or not page.subsections:
            continue # Skip pages without valid subsections
        for subsection in page.subsections:
             if isinstance(subsection, Subsection):
                # Store the subsection object itself and context for the prompt
                context = {
                    "subsection_id": subsection.subsection_id,
                    "title": subsection.title,
                    "description": subsection.description, # Use the description field
                    "page_number": page.page_number, # Get page number from parent
                    # Optional: Add a very short text preview if description isn't enough context
                    # "text_preview": subsection.text[:50] + "..." if len(subsection.text) > 50 else subsection.text
                }
                all_subsections_with_context.append(context)
                subsection_map[subsection.subsection_id] = subsection # Store the actual object
             else:
                  logging.warning(f"Skipping invalid subsection object on page {page.page_number}")


    if not all_subsections_with_context:
        logging.warning("No valid subsections found across all pages to form chapters.")
        # Return an empty list or a single default chapter if required by downstream logic
        # For now, returning empty list. Adjust if a default chapter is mandatory.
        return []

    # Sort subsections by page number and order (using the context dict)
    # This ensures the LLM sees them in the correct document order.
    sorted_subsections_context = sorted(all_subsections_with_context, key=lambda s: (s["page_number"], s.get("order", 0))) # Use get for order safety

    logging.info(f"Sending {len(sorted_subsections_context)} subsections to LLM for chapter structuring.")

    # Create prompt for chapter extraction using NEW structure
    # Pass subsection_id, title, and description
    chapter_prompt = f"""
    Analyze the following list of subsections extracted sequentially from a document. Each subsection includes its ID, title, and a brief description. Your goal is to group these subsections into logically coherent chapters.

    Subsections List:
    ```json
    {json.dumps(sorted_subsections_context, indent=2)}
    ```

    Instructions:
    1. Group the provided subsections into 3 to 10 logical chapters.
    2. Base the grouping on topic cohesion and narrative flow. Subsections within a chapter should relate to a common theme or part of the document's structure.
    3. Each chapter MUST contain at least one subsection.
    4. Ensure every subsection ID from the input list is assigned to exactly ONE chapter.
    5. Maintain the original relative order of subsections within and across chapters as much as possible.
    6. For each chapter created, provide:
        - A unique `chapter_id` (e.g., "chapter_1", "chapter_2").
        - A concise and descriptive `title` (max 7 words) reflecting the chapter's main theme.
        - A brief `summary` (1-2 sentences) describing the chapter's content or purpose.
        - A list of `subsection_ids` belonging to that chapter, in their original relative order.
    7. Number chapters sequentially starting from 1 using an `order` field.

    Return ONLY a valid JSON array containing the chapter objects, structured exactly like this example:
    ```json
    [
      {{
        "chapter_id": "chapter_1",
        "order": 1,
        "title": "Introduction and Setup",
        "summary": "Provides background information and initial setup instructions.",
        "subsection_ids": ["page_1_section_1", "page_1_section_2", "page_2_section_1"]
      }},
      {{
        "chapter_id": "chapter_2",
        "order": 2,
        "title": "Core Functionality Explained",
        "summary": "Details the main features and how they operate.",
        "subsection_ids": ["page_2_section_2", "page_3_section_1", "page_3_section_2"]
      }}
      // ... more chapters
    ]
    ```
    Ensure the final output is only the JSON array, with no extra text before or after.
    """

    chapters = []
    try:
        logging.debug("Sending chapter extraction request to LLM.")
        # Use the actual client and API call structure
        response = await retry_api_call(
            client.aio.models.generate_content, # Adjusted path might be needed based on client setup
            model="gemini-2.0-flash", # Or pro, flash might struggle with complex structuring
            contents=[
                types.Content(parts=[types.Part.from_text(text=chapter_prompt)]),
            ], # Prompt first, then text content
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )

        if response.candidates:
            json_text = clean_json_response(response.candidates[0].content.parts[0].text)
            logging.debug(f"LLM response for chapters: {json_text[:500]}...") # Log beginning of response
            chapters_data = json.loads(json_text) # Expecting a List[Dict]

            if not isinstance(chapters_data, list):
                 logging.warning(f"LLM returned non-list data for chapters: {type(chapters_data)}. Attempting recovery if dict.")
                 # Handle cases where it might incorrectly return a dict containing the list
                 if isinstance(chapters_data, dict) and "chapters" in chapters_data and isinstance(chapters_data["chapters"], list):
                     chapters_data = chapters_data["chapters"]
                 else:
                     raise ValueError("LLM response for chapters was not a list.")


            processed_subsection_ids = set()
            for i, chapter_data in enumerate(chapters_data):
                # Validate chapter_data structure
                if not isinstance(chapter_data, dict) or not all(k in chapter_data for k in ["title", "summary", "subsection_ids", "order"]):
                     logging.warning(f"Skipping malformed chapter data received from LLM: {chapter_data}")
                     continue

                chapter_subsections = []
                subsection_ids_in_chapter = chapter_data.get("subsection_ids", [])
                if not isinstance(subsection_ids_in_chapter, list):
                     logging.warning(f"Subsection IDs for chapter '{chapter_data.get('title')}' is not a list. Skipping chapter.")
                     continue

                for subsection_id in subsection_ids_in_chapter:
                    if subsection_id in subsection_map:
                        chapter_subsections.append(subsection_map[subsection_id])
                        processed_subsection_ids.add(subsection_id)
                    else:
                        logging.warning(f"Subsection ID '{subsection_id}' from LLM not found in original map for chapter '{chapter_data.get('title')}'.")

                if not chapter_subsections:
                     logging.warning(f"Chapter '{chapter_data.get('title')}' created by LLM has no valid subsections. Skipping.")
                     continue

                # Create the Chapter object (ensure your Chapter model matches)
                chapter = Chapter(
                    chapter_id=chapter_data.get("chapter_id", f"chapter_{chapter_data.get('order', i+1)}"), # Use order for default ID
                    title=chapter_data.get("title", f"Chapter {chapter_data.get('order', i+1)}"),
                    summary=chapter_data.get("summary", "No summary provided."),
                    subsections=chapter_subsections, # List of actual Subsection objects
                    # entity_relationships=[], # Defer relationship extraction
                    order=chapter_data.get("order", i+1)
                )
                chapters.append(chapter)

            # Check if all original subsections were processed
            original_ids = set(subsection_map.keys())
            missing_ids = original_ids - processed_subsection_ids
            if missing_ids:
                 logging.warning(f"LLM failed to assign {len(missing_ids)} subsections to chapters: {missing_ids}. They will be excluded.")
                 # Optionally, create an "Orphaned" chapter for these

            # Sort final chapters list by order
            chapters = sorted(chapters, key=lambda c: c.order)

            if not chapters:
                 logging.warning("LLM processing resulted in zero valid chapters. Creating default chapter.")
                 raise ValueError("LLM returned no valid chapters") # Trigger default chapter creation in except block

            logging.info(f"Successfully extracted {len(chapters)} chapters from LLM.")
            return chapters

        else:
             logging.error("Chapter extraction response from LLM had no candidates.")
             raise ValueError("LLM response had no candidates") # Trigger default

    except (json.JSONDecodeError, ValidationError, ValueError, Exception) as e:
        logging.error(f"Error processing LLM response or extracting chapters: {e}", exc_info=True)

        # --- Fallback: Create a single default chapter ---
        logging.info("Creating a single default chapter due to chapter extraction error.")
        # Get all original subsections back from the map, sorted
        all_original_subsections = sorted(subsection_map.values(), key=lambda s: (
            # Need page number back for sorting - requires passing it or storing it differently
            # Assuming subsection_map values have access or we modify the storage
            # For now, sort by ID as a fallback sort order
            getattr(s, 'subsection_id', '')
        ))

        default_chapter = Chapter(
            chapter_id="chapter_1_default",
            title="Document Content",
            summary="Content grouped into a single chapter due to an error during automatic chapter structuring.",
            subsections=all_original_subsections, # Assign all subsections here
            # entity_relationships=[],
            order=1
        )
        return [default_chapter]

# Helper function for consistent ID generation
def generate_qdrant_id(node_type: str, name: str, context_id: str) -> str:
    """Generates a deterministic UUID for Qdrant based on node info."""
    # Combine type, name, and context for uniqueness
    unique_string = f"{node_type}:{name}:{context_id}".lower()
    # Generate UUID using NAMESPACE_DNS (ensures same input -> same output)
    return str(uuid.uuid5(uuid.NAMESPACE_DNS, unique_string))


async def analyze_concept_entity_relationships(
    client: genai.Client,
    chapter: Chapter,
    source_identifier: str,
    max_structure_retries: int = 2,
    retry_delay: int = 3
) -> None:
    """
    Analyzes text within a single chapter to extract key concepts/entities
    and their relationships using an LLM structured prompt. Populates
    chapter.qdrant_nodes in-place. Includes retries on structural errors,
    AND a fallback mechanism using extracted text from failed responses.

    Args:
        client: The configured GenAI client/model instance.
        chapter: The Chapter object containing subsections to analyze.
        source_identifier: Identifier for the source document.
        max_structure_retries: How many times to retry the INITIAL prompt if the LLM returns the wrong JSON structure.
        retry_delay: Base delay (in seconds) between structure retries.
    """
    function_name = "analyze_concept_entity_relationships_fallback" # Renamed
    chapter_id = getattr(chapter, 'chapter_id', 'unknown_chapter')
    chapter_title = getattr(chapter, 'title', 'Untitled Chapter')
    logging.info(f"[{function_name}] Analyzing relationships for Chapter '{chapter_id}': '{chapter_title}' from source '{source_identifier}'")

    # --- 1. Prepare Input Text ---
    if not chapter.subsections:
        logging.warning(f"[{function_name}] Chapter '{chapter_id}' has no subsections. Skipping analysis.")
        chapter.qdrant_nodes = []
        return

    combined_text = ""
    try:
        # --- Start Expansion: combined_text generation ---
        subsection_texts = []
        for subsection in chapter.subsections:
             # Use getattr for safe access to attributes
             sub_id = getattr(subsection, 'subsection_id', 'unknown_id')
             sub_title = getattr(subsection, 'title', 'Untitled Subsection')
             sub_desc = getattr(subsection, 'description', '') # Use the description for context
             sub_text = getattr(subsection, 'text', '')

             # Only include non-empty text
             if sub_text and sub_text.strip():
                 # Format clearly for the LLM
                 subsection_texts.append(
                     f"--- Subsection Start (ID: {sub_id}) ---\n"
                     f"Title: {sub_title}\n"
                     f"Description: {sub_desc}\n\n" # Description provides LLM context
                     f"{sub_text.strip()}\n" # Add the actual text
                     f"--- Subsection End ---"
                 )
             else:
                 logging.debug(f"[{function_name}] Skipping empty subsection {sub_id} in chapter {chapter_id}.")

        combined_text = "\n\n".join(subsection_texts)
        # --- End Expansion: combined_text generation ---

        MAX_TEXT_LIMIT = 150000 # Adjust based on model context window limits
        if len(combined_text) > MAX_TEXT_LIMIT:
             logging.warning(f"[{function_name}] Truncating combined text for chapter '{chapter_id}' from {len(combined_text)} to {MAX_TEXT_LIMIT} chars.")
             combined_text = combined_text[:MAX_TEXT_LIMIT] + "\n... [Content Truncated] ..."
        elif not combined_text.strip():
            logging.warning(f"[{function_name}] Combined text for chapter '{chapter_id}' is empty after processing subsections. Skipping analysis.")
            chapter.qdrant_nodes = []
            return
    except Exception as text_prep_error:
        logging.error(f"[{function_name}] Error preparing text for chapter '{chapter_id}': {text_prep_error}", exc_info=True)
        chapter.qdrant_nodes = []
        return

    # --- 2. Extract Chapter Number ---
    chapter_num_str = "1" # Default
    if chapter_id and isinstance(chapter_id, str):
        match = re.search(r'[_-](\d+)$', chapter_id)
        if match:
            chapter_num_str = match.group(1)

    # --- 3. Define the *ORIGINAL* LLM Prompt ---

    # --- Start Expansion: new_json_schema_description ---
    new_json_schema_description = """
```json
{
  "qdrant_nodes": [
    {
      // --- Core Identifiers (Generated by LLM) ---
      "node_id": "string", // REQUIRED. Format: "{type}_{sanitized_name}_ch{chapter_num_str}". Example: "concept_machine_learning_ch1".
      "node_type": "string", // REQUIRED. "entity" or "concept".
      "name": "string", // REQUIRED. Original human-readable name.
      "chapter_id_context": "string", // REQUIRED. Use the provided Chapter ID. Example: "{chapter_id}"

      // --- Extracted Content & Metadata Fields (Generated by LLM) ---
      "title": "string", // REQUIRED. Concise title for this node's content (5-10 words).
      "core_claim": "string", // REQUIRED. 1-3 sentence summary of the central point/function in context. Be specific.
      "information_type": "string", // REQUIRED. Classification of info (e.g., "Definition", "Explanation", "Process Description").
      "key_entities": "string", // REQUIRED. JSON string list of specific entities in context. Example: "[\\"NVIDIA GPUs\\", \\"TensorFlow\\"]". Use "[]" if none.
      "tags": "string", // REQUIRED. JSON string list of 1-5 relevant keywords. Example: "[\\"rag\\", \\"llm\\"]". Use "[]" if none.
      "sentiment_tone": "string", // REQUIRED. Overall sentiment/tone (e.g., "Positive", "Neutral", "Objective").
      "source_credibility": "string", // REQUIRED. Inferred credibility (e.g., "Company Report", "Technical Documentation", "Unknown").

      // --- Relationships (Generated by LLM) ---
      "linked_nodes": [ // List, REQUIRED (can be empty []). Max 7 links per node.
        {
          "target_node_id": "string", // REQUIRED. node_id of another node defined in THIS response's "qdrant_nodes" list.
          "relationship_description": "string", // REQUIRED. Specific description (Use verbs). Example: 'calculates cost based on'.
          "relationship_keywords": ["string"], // List, OPTIONAL. 1-3 keywords. Defaults to []. Example: ["calculation", "dependency"].
          "relationship_strength": "float" // REQUIRED. Link strength/importance (1.0-10.0).
        }
      ]
      // NOTE: 'source' and 'stored_at' fields MUST NOT be included here. They are added later by code.
    }
  ]
}
```"""

    # --- End Expansion: new_json_schema_description ---
    initial_entity_prompt = f"""
        --- Goal ---
        Analyze the text from the single document chapter provided below. Your task is to identify key specific entities (e.g., organizations, people, products) AND important abstract concepts (e.g., processes, topics, ideas). Extract these as nodes. For each node, determine its intrinsic properties based *only* on the chapter text. Also, identify direct relationships *between these nodes within this chapter*.

        --- Input Data ---
        Chapter Title: {chapter_title}
        Chapter ID: {chapter_id}
        Source Identifier (Filename/URL): {source_identifier}
        Chapter Content:
        {combined_text}
        --- End Input Data ---

        --- Instructions ---
        Perform the following tasks based *only* on the provided "Chapter Content":

        **Task 1: Extract Nodes and Intra-Chapter Relationships**
        1a. **Identify Nodes:** Scan the text and identify all distinct key entities and abstract concepts. Aim for ~5-15 significant nodes per chapter.
        1b. **Define Node Structure:** For each node identified, structure its information directly within the node object (do NOT use a nested 'properties' object):
            *   `node_id`: Generate using format `{{type}}_{{sanitized_name}}_ch{{chapter_num_str}}`. Ensure uniqueness within this chapter's output.
            *   `node_type`: Must be exactly "entity" or "concept".
            *   `name`: The original, human-readable name.
            *   `chapter_id_context`: Set to exactly "{chapter_id}".
            *   `title`: Generate a concise title (5-10 words) capturing the node's discussion within the chapter.
            *   `core_claim`: Write a 1-3 sentence summary capturing the central point/function in this chapter's context. Be specific.
            *   `information_type`: Classify the primary info type (e.g., "Definition", "Explanation", "Process Step").
            *   `key_entities`: Provide a JSON stringified list of specific entity names explicitly mentioned in the node's context. Example: `"[\\"NVIDIA GPUs\\", \\"TensorFlow Framework\\"]"`. Use `"[]"` if none.
            *   `tags`: Provide a JSON stringified list of 1-5 relevant keywords (lowercase, underscores). Example: `"[\\"rag_framework\\", \\"llm_enhancement\\"]"`. Use `"[]"` if none.
            *   `sentiment_tone`: Assess the overall sentiment/tone (e.g., "Positive", "Neutral", "Objective").
            *   `source_credibility`: Infer credibility based on source '{source_identifier}' and text tone (e.g., "Company Report", "Technical Documentation", "Unknown").
            *   `linked_nodes`: An initially empty list `[]`.
        1c. **Identify Relationships (Linked Nodes):** For each node created (source), find its most important (max 5-7) direct relationships to *other nodes created in this Task* (target).
            *   For each relationship, create a `linked_nodes` object containing: `target_node_id`, `relationship_description` (CRITICAL: specific verbs), `relationship_keywords` (optional), `relationship_strength` (float 1.0-10.0).
            *   Populate the `linked_nodes` list of the source node. If no links, use empty `[]`.
        1d. **Assemble Output:** Format all extracted node objects into a single JSON object according to the schema specified below.

        --- End Instructions ---

        --- Output Format ---
        Return ONLY a single, valid JSON object. It MUST contain exactly one top-level key: "qdrant_nodes". The value MUST be a JSON array of node objects. Adhere STRICTLY to this schema:

        {new_json_schema_description}
        --- End Output Format ---

        --- Critical Rules ---
        *   VALID JSON ONLY: Output must be a single JSON object `{{...}}`. No extra text or markdown.
        *   SCHEMA ADHERENCE: Follow the schema precisely. All REQUIRED fields must be present.
        *   **NO NESTED PROPERTIES:** Do NOT create a nested "properties" dictionary. All fields like 'title', 'core_claim' must be at the top level of the node object.
        *   **DO NOT GENERATE:** The fields `source` and `stored_at` MUST NOT be generated by you; they are added later by code.
        *   JSON STRINGS: Ensure `key_entities` and `tags` values are valid JSON strings representing lists (e.g., `"[]"`, `"["item1"]"`). Double-check escaping.
        *   NODE ID CONSISTENCY: `target_node_id` MUST correspond to a `node_id` defined in the SAME `qdrant_nodes` list.
        *   RELATIONSHIP SPECIFICITY: Use specific, verb-driven descriptions.
        *   TEXTUAL BASIS: All information must be derived SOLELY from the provided "Chapter Content".
        *   UNICODE: Handle Unicode characters correctly.

        --- End Critical Rules ---

        Generate the JSON output now:
        """


    # --- 4. Call LLM and Process Response (Modified with Fallback) ---

    chapter.qdrant_nodes = [] # Initialize
    final_data = None
    last_raw_response_text = "N/A" # Store last raw response for logging/fallback
    last_cleaned_response_text = "N/A" # Store last cleaned for logging/fallback
    last_parsed_data = None # Store last parsed data for fallback check

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Initial Prompt Attempt {attempt + 1}/{max_structure_retries} for chapter '{chapter_id}'")
        response = None # Reset response for this attempt
        json_text = ""
        data = None

        try:
            response = await retry_api_call(
                client.aio.models.generate_content, # Correct method assuming 'client' is the model instance
                model="gemini-2.5-flash-preview-04-17", # Use model name
                contents=[types.Content(parts=[types.Part.from_text(text=initial_entity_prompt)])], # Pass prompt directly as content
                config=types.GenerateContentConfig(response_mime_type="application/json") # Use GenerateContentConfig
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates (Initial Attempt {attempt + 1}) for chapter '{chapter_id}'.")
                last_raw_response_text = "No candidates received"
                last_cleaned_response_text = "{}"
                last_parsed_data = {}
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: break

            last_raw_response_text = response.candidates[0].content.parts[0].text
            json_text = clean_json_response(last_raw_response_text)
            last_cleaned_response_text = json_text
            logging.debug(f"[{function_name}] Cleaned LLM response (Initial Attempt {attempt + 1}, Chapter {chapter_id}): {json_text[:500]}...")

            try:
                data = json.loads(json_text)
                last_parsed_data = data
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError initial attempt {attempt + 1} for chapter '{chapter_id}': {json_err}.")
                # Attempt to extract text even from decode error for fallback
                if isinstance(json_text, str) and len(json_text.strip()) > 50:
                     last_parsed_data = {"error": "JSONDecodeError", "text": json_text}
                else:
                     last_parsed_data = {"error": "JSONDecodeError", "text": ""}
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: break

            # --- Structure Check ---
            if isinstance(data, dict) and "qdrant_nodes" in data and isinstance(data.get("qdrant_nodes"), list):
                logging.info(f"[{function_name}] Received correct structure initial attempt {attempt + 1} for chapter '{chapter_id}'.")
                final_data = data
                break
            else:
                logging.warning(f"[{function_name}] Incorrect JSON structure initial attempt {attempt+1} for chapter '{chapter_id}'. Got: {type(data)}. Structure: {str(data)[:300]}...")
                # Log the failing structure for debugging
                logging.debug(f"[{function_name}] Failing structure details (Attempt {attempt+1}, Chapter {chapter_id}): Type={type(data)}, Content='{str(data)[:1000]}...'")
                if attempt < max_structure_retries - 1:
                    await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else:
                    logging.info(f"[{function_name}] Initial prompt failed structure check after {max_structure_retries} attempts for chapter '{chapter_id}'. Proceeding to fallback attempt.")
                    break

        except Exception as e:
            logging.error(f"[{function_name}] Error during initial attempt {attempt + 1} for chapter '{chapter_id}': {e}", exc_info=True)
            last_raw_response_text = f"Exception: {e}"
            last_cleaned_response_text = "{}"
            last_parsed_data = {"error": str(e)}
            if attempt < max_structure_retries - 1:
                 await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else:
                 logging.error(f"[{function_name}] Initial attempts failed completely after {max_structure_retries} tries for chapter '{chapter_id}'. Proceeding to fallback attempt.")
                 break

    # --- 5. Fallback Attempt (if final_data is still None after initial retries) ---
    if final_data is None:
        logging.info(f"[{function_name}] Entering fallback: Attempting extraction using text from last failed response for chapter '{chapter_id}'.")

        extracted_text_for_fallback = None
        # Check if the last parsed data (from failed attempt) contains usable text
        if isinstance(last_parsed_data, dict) and "text" in last_parsed_data and isinstance(last_parsed_data["text"], str) and len(last_parsed_data["text"].strip()) > 50: # Basic check for usable text
            extracted_text_for_fallback = last_parsed_data["text"].strip()
            logging.debug(f"[{function_name}] Found extracted text for fallback (length {len(extracted_text_for_fallback)}) for chapter '{chapter_id}'.")
        else:
            logging.warning(f"[{function_name}] Could not find usable extracted text in the last failed response for chapter '{chapter_id}'. Cannot perform text-based fallback.")
            logging.debug(f"[{function_name}] Last parsed data type for fallback check: {type(last_parsed_data)}, keys: {list(last_parsed_data.keys()) if isinstance(last_parsed_data, dict) else 'N/A'}")


        if extracted_text_for_fallback:
            # --- Construct Fallback Prompt ---
            # Uses the same instructions, schema, and rules, but different input text
            fallback_entity_prompt = f"""
--- Goal ---
Analyze the provided text, which was extracted from a previous failed attempt to structure chapter content. Your task is to identify key specific entities AND important abstract concepts from THIS EXTRACTED TEXT ONLY. Extract these as nodes with properties and relationships.

--- Input Data ---
Chapter Title: {chapter_title} (Context Only)
Chapter ID: {chapter_id} (Context Only)
Source Identifier (Filename/URL): {source_identifier} (Context Only)
Extracted Chapter Content (Potentially Incomplete):
{extracted_text_for_fallback}
--- End Input Data ---

--- Instructions ---
Perform the following tasks based *only* on the provided "Extracted Chapter Content":
{'''
**Task 1: Extract Nodes, Properties, and Intra-Chapter Relationships**
1a. **Identify Nodes:** Scan the text and identify all distinct key entities and abstract concepts. An entity is typically a concrete noun (person, place, organization, product). A concept is an abstract idea, process, topic, or technique. Aim for ~5-15 significant nodes per chapter, focusing on items central to the chapter's narrative.
1b. **Define Node Structure:** For each node identified, structure its information as follows:
    *   `node_id`: Generate using the format `{type}_{sanitized_name}_ch{chapter_num_str}`. '{type}' is 'entity' or 'concept'. '{sanitized_name}' is the lowercase node name with spaces/symbols replaced by underscores `_`, truncated to ~50 characters. Ensure uniqueness within this chapter's output.
    *   `node_type`: Must be exactly "entity" or "concept".
    *   `name`: The original, human-readable name as found in the text or a standard representation (e.g., "Machine Learning" not "machine learning").
    *   `chapter_id_context`: Set to exactly "{chapter_id}".
    *   `properties`: A nested JSON object containing the following fields derived from the text segment most relevant to *this specific node*:
        *   `title`: Generate a concise title (5-10 words) capturing the essence of this node's discussion within the chapter. If a clear heading exists for the relevant section, adapt it. Otherwise, create a descriptive title. Example: "Function of the RAG System Retriever".
        *   `source_credibility`: Infer the credibility level based on the source identifier '{source_identifier}' and the text content/tone. Choose from: "Peer-Reviewed", "Company Report", "News Article", "Blog Post", "Technical Documentation", "User Manual", "Internal Memo", "Marketing Material", "Personal Communication", "Unknown".
        *   `information_type`: Classify the primary type of information this node represents within its context. Choose from: "Definition", "Explanation", "Description", "Example", "Comparison", "Contrast", "Data Point", "Metric", "Process Step", "Workflow", "Case Study", "Methodology", "Algorithm", "API Specification", "Configuration Detail", "Requirement", "Constraint", "Benefit", "Limitation", "Finding", "Claim", "Hypothesis", "Opinion", "Prediction", "Guideline", "Best Practice".
        *   `core_claim`: Write a 1-3 sentence summary capturing the central point, assertion, or function of this node *in this chapter's context*. Be specific and action-oriented where possible. Avoid generic statements. Example: "The retrieval module (R) searches external knowledge sources to find relevant context for the generation module." NOT "This node describes the retrieval module."
        *   `key_entities`: Provide a JSON stringified list of specific entity names (people, organizations, product names, locations, specific systems mentioned by name) explicitly mentioned in the immediate context of this node's primary discussion. Example: "["NVIDIA GPUs", "TensorFlow Framework"]". If none are mentioned in the direct context, use "[]".
        *   `sentiment_tone`: Assess the overall sentiment or tone of the text discussing this node. Choose from: "Positive", "Neutral", "Negative", "Mixed", "Objective", "Subjective", "Promotional", "Cautionary", "Speculative", "Instructional", "Analytical".
        *   `tags`: Provide a JSON stringified list of 1-5 relevant keywords or tags that categorize this node. Use lowercase and underscores. Example: "["rag_framework", "llm_enhancement", "vector_database"]". If no specific tags apply, use "[]".
    *   `linked_nodes`: An initially empty list `[]` where relationship objects will be added (see step 1c).

1c. **Identify Relationships (Linked Nodes):** For each node created (the *source* node), scan the text again to find its most important (max 5-7) direct relationships *to other nodes also created in this Task* (the *target* nodes). Relationships should represent actions, dependencies, comparisons, compositions, etc.
    *   For each relationship found, create a `linked_nodes` object containing:
        *   `target_node_id`: The `node_id` (generated in step 1b) of the target node. This target node MUST be one of the other nodes defined in this response's `qdrant_nodes` list. Do NOT link to nodes outside this list.
        *   `relationship_description`: **CRITICAL: Describe the SPECIFIC ACTION or NATURE of the relationship using precise verbs.** Avoid generic terms like 'related to', 'associated with', 'discusses'. Examples: "implements algorithm", "evaluates performance of", "optimizes parameters for", "is a component of", "contrasts with", "cites findings from".
        *   `relationship_keywords`: (Optional) 1-3 lowercase keywords summarizing the relationship type (e.g., ["dependency", "causation", "comparison", "optimization"]). Default to `[]` if none clearly apply.
        *   `relationship_strength`: Estimate the importance or centrality of this link within the chapter's narrative (float 1.0-10.0). A core dependency might be 9.0, a brief mention or comparison might be 3.0.
    *   Populate the `linked_nodes` list of the source node with these relationship objects. If a node has no direct links to *other nodes in this list*, its `linked_nodes` list remains empty `[]`.

1d. **Assemble Output:** Format all extracted node objects (including their `properties` and `linked_nodes`) into a single JSON object according to the schema specified below. Ensure the entire output is one valid JSON structure.
'''}
--- End Instructions ---

--- Output Format ---
Return ONLY a single, valid JSON object. It MUST contain exactly one top-level key: "qdrant_nodes". Adhere STRICTLY to this schema:

{new_json_schema_description}
--- End Output Format ---

--- Critical Rules ---
{'''
*   **VALID JSON ONLY:** Entire output must be a single JSON object starting with `{` and ending with `}`. No introductory text, explanations, apologies, or markdown formatting outside the JSON.
*   **SCHEMA ADHERENCE:** Follow the provided schema precisely. All REQUIRED fields/objects must be present (e.g., `properties` dict must exist, `linked_nodes` list must exist even if empty `[]`). Check data types carefully (strings, floats, lists).
*   **PROPERTIES:** All fields within the `properties` object are REQUIRED, except `source` and `stored_at` which MUST NOT be generated here (they are added later by code). Ensure default/placeholder values are provided if info isn't in text (e.g., "Unknown" for credibility).
*   **JSON STRINGS:** Ensure `key_entities` and `tags` values within `properties` are valid JSON strings representing lists (e.g., `"[]"`, `"["item1", "item2"]"`). Double-check escaping.
*   **NODE ID CONSISTENCY:** Every `target_node_id` in any `linked_nodes` list MUST correspond to a `node_id` defined earlier in the main `qdrant_nodes` list within THIS SAME RESPONSE. Do not hallucinate target IDs or link outside the generated list.
*   **RELATIONSHIP SPECIFICITY:** Relationship descriptions MUST be specific and verb-driven.
*   **RELATIONSHIP LIMIT:** Max 5-7 `linked_nodes` objects per source node. Focus on the most direct and important links mentioned in the text.
*   **TEXTUAL BASIS:** All information (`properties`, `relationship_description`) must be derived SOLELY from the provided "Extracted Chapter Content". Do not infer external knowledge or details not present in the text.
*   **UNICODE:** Ensure correct handling of Unicode characters within the JSON output.
'''}
--- End Critical Rules ---

Generate the JSON output now based *only* on the 'Extracted Chapter Content':
"""
            # --- Execute Fallback LLM Call ---
            logging.debug(f"[{function_name}] Making fallback LLM call for chapter '{chapter_id}'.")
            fallback_response = None
            fallback_json_text = ""
            fallback_data = None
            try:
                fallback_response = await retry_api_call(
                    client.aio.models.generate_content, # Use retry_api_call for robustness
                    model="gemini-2.0-flash", # Use Gemini 2.0 Flash for speed
                    contents=[types.Content(parts=[types.Part.from_text(text=fallback_entity_prompt)])],
                    config=types.GenerateContentConfig(response_mime_type="application/json"),
                )

                if fallback_response and fallback_response.candidates:
                    fallback_json_text = clean_json_response(fallback_response.candidates[0].content.parts[0].text)
                    logging.debug(f"[{function_name}] Cleaned Fallback Response (Chapter {chapter_id}): {fallback_json_text[:500]}...")
                    try:
                        fallback_data = json.loads(fallback_json_text)
                        # --- Final Structure Check (Fallback) ---
                        if isinstance(fallback_data, dict) and "qdrant_nodes" in fallback_data and isinstance(fallback_data.get("qdrant_nodes"), list):
                            logging.info(f"[{function_name}] Fallback attempt SUCCEEDED with correct structure for chapter '{chapter_id}'.")
                            final_data = fallback_data # Assign data from successful fallback
                        else:
                            logging.error(f"[{function_name}] Fallback attempt FAILED structure check for chapter '{chapter_id}'. Got type: {type(fallback_data)}")
                            # Optionally log the failing structure: logging.debug(f"Fallback failing structure: {str(fallback_data)[:1000]}")
                    except json.JSONDecodeError as fallback_json_err:
                        logging.error(f"[{function_name}] Fallback JSONDecodeError for chapter '{chapter_id}': {fallback_json_err}")
                else:
                    logging.error(f"[{function_name}] Fallback LLM call returned no candidates for chapter '{chapter_id}'.")

            except Exception as fallback_err:
                logging.error(f"[{function_name}] Error during fallback LLM call for chapter '{chapter_id}': {fallback_err}", exc_info=True)
        # End of fallback execution block (if extracted_text_for_fallback was found)

    # --- 6. Process Final Data (if successful from initial or fallback) ---
    if final_data is not None:
        qdrant_nodes_data = final_data.get("qdrant_nodes")
        if isinstance(qdrant_nodes_data, list):
            validated_nodes: List[QdrantNode] = []
            temp_node_map_by_name_ctx: Dict[Tuple[str, str, str], str] = {} # (type, name, ctx) -> uuid
            nodes_to_process_later: List[Dict] = [] # Store dicts needing link resolution

            # --- First Pass: Prepare nodes, generate UUIDs, build temporary map ---
            current_iso_ts = datetime.datetime.now(datetime.timezone.utc).isoformat()
            # Removed unused variable: required_fields_from_llm

            for node_dict in qdrant_nodes_data:
                if not isinstance(node_dict, dict): continue
                try:
                    # --- A. Prepare and Validate Node Dict from LLM ---
                    # Check for essential fields generated by LLM
                    if not all(k in node_dict for k in ["node_type", "name", "chapter_id_context"]):
                         logging.warning(f"Node missing core type/name/chapter_id_context. Skipping. Data: {str(node_dict)[:200]}")
                         continue

                    # Assign Deterministic UUID (same as before)
                    node_type = node_dict["node_type"]; name = node_dict["name"]; chapter_ctx = node_dict["chapter_id_context"]
                    node_uuid = generate_qdrant_id(node_type, name, chapter_ctx)
                    node_dict['node_id'] = node_uuid # Assign the generated UUID

                    # Store mapping for link resolution
                    temp_node_map_by_name_ctx[(node_type, name, chapter_ctx)] = node_uuid

                    # Add fields managed by Python code DIRECTLY to node_dict
                    node_dict['source'] = source_identifier
                    node_dict['stored_at'] = current_iso_ts

                    # Ensure other required fields exist (provide defaults if LLM missed them)
                    llm_required_fields = ["title", "core_claim", "information_type", "key_entities", "tags", "sentiment_tone", "source_credibility"]
                    for key in llm_required_fields:
                         if key not in node_dict or node_dict[key] is None:
                              logging.debug(f"[{function_name}] Node '{node_uuid}' missing LLM field '{key}'. Setting default.")
                              # Set appropriate defaults based on expected type
                              if key in ["key_entities", "tags"]: node_dict[key] = "[]"
                              else: node_dict[key] = f"Missing: {key}" # Or more specific defaults like "Unknown"

                    # Ensure key_entities and tags are valid JSON strings (validator might handle this, but double-check)
                    for key in ["key_entities", "tags"]:
                        val = node_dict.get(key)
                        if isinstance(val, list): # If LLM gave a list instead of string
                            try: node_dict[key] = json.dumps(val)
                            except TypeError: node_dict[key] = "[]"; logging.warning(f"Node {node_uuid} {key} list couldn't be serialized.")
                        elif not isinstance(val, str): node_dict[key] = "[]"; logging.warning(f"Node {node_uuid} {key} is not string or list: {type(val)}.")
                        else: # It's a string, validate JSON
                            try: json.loads(val)
                            except json.JSONDecodeError: node_dict[key] = "[]"; logging.warning(f"Node {node_uuid} {key} invalid JSON string: '{val[:50]}...'")

                    # Remove fields that should not be at the top level (these don't exist anymore)
                    # node_dict.pop('description', None) # Unnecessary
                    # node_dict.pop('document_context', None) # Unnecessary

                    # Ensure linked_nodes is present and is a list
                    linked_nodes_list = node_dict.setdefault('linked_nodes', [])
                    if not isinstance(linked_nodes_list, list):
                        logging.warning(f"[{function_name}] Node '{node_uuid}' had non-list 'linked_nodes'. Resetting to [].")
                        node_dict['linked_nodes'] = []

                    # Ensure keywords/strength within links are valid (before second pass)
                    for link_item in node_dict['linked_nodes']:
                         if isinstance(link_item, dict):
                              kw = link_item.setdefault('relationship_keywords', [])
                              if not isinstance(kw, list): link_item['relationship_keywords'] = []
                              strength = link_item.get('relationship_strength', 5.0)
                              try: link_item['relationship_strength'] = float(strength)
                              except (ValueError, TypeError): link_item['relationship_strength'] = 5.0

                    # Add to list for second pass (link resolution)
                    nodes_to_process_later.append(node_dict) # node_dict is now flat

                except Exception as node_prep_err:
                     logging.error(f"Error preparing node during first pass: {node_prep_err}. Data: {str(node_dict)[:200]}", exc_info=True)


            # --- Second Pass: Resolve links and validate QdrantNode ---
            final_node_ids = set()
            validated_nodes: List[QdrantNode] = [] # Define here

            for node_dict in nodes_to_process_later: # node_dict already has UUID and source/stored_at
                node_uuid = node_dict['node_id']
                raw_links = node_dict.get('linked_nodes', []) # This is still List[Dict] from LLM
                resolved_links_models: List[LinkedNode] = [] # Store validated Pydantic models

                if isinstance(raw_links, list):
                     # --- Link Resolution Logic (remains the same, using temp_node_map_by_name_ctx) ---
                     for link_dict in raw_links:
                         if not isinstance(link_dict, dict): continue
                         target_original_id = link_dict.get("target_node_id") # LLM still returns original format ID
                         target_uuid = None
                         if target_original_id:
                             # Try parsing: assumes format "type_name_chX"
                             parts = target_original_id.split('_')
                             if len(parts) >= 3:
                                  target_type = parts[0]
                                  # Robust chapter context extraction
                                  ch_match = re.search(r"_ch(\d+)$", target_original_id)
                                  if ch_match:
                                      target_chapter_ctx_suffix = ch_match.group(0) # like "_ch1"
                                      target_chapter_ctx = ch_match.group(1) # like "1" -> Use the actual number part maybe? Or keep "ch1"? Let's keep "ch1" format consistent with generation
                                      target_chapter_ctx_full = "ch" + target_chapter_ctx # Reconstruct "chX"

                                      name_end_index = target_original_id.rfind(target_chapter_ctx_suffix)
                                      name_start_index = len(target_type) + 1
                                      target_name_underscores = target_original_id[name_start_index:name_end_index]
                                      target_name = target_name_underscores.replace('_', ' ').title()

                                      # Lookup UUID using reconstructed info
                                      target_uuid = temp_node_map_by_name_ctx.get((target_type, target_name, target_chapter_ctx_full)) # Use chX format for lookup
                                      if not target_uuid: logging.warning(f"UUID lookup failed for ({target_type}, {target_name}, {target_chapter_ctx_full}) from {target_original_id}")
                                  else:
                                      logging.warning(f"Could not extract chapter context from target_original_id '{target_original_id}' for node '{node_uuid}'.")
                             else:
                                  logging.warning(f"Could not parse target_original_id '{target_original_id}' for node '{node_uuid}'.")


                         if target_uuid: # If lookup successful
                             try:
                                 # Validate link data and create Pydantic LinkedNode model instance
                                 # Strength conversion happens in LinkedNode validator now
                                 resolved_link = LinkedNode(
                                     target_node_id=target_uuid, # Use the resolved UUID
                                     relationship_description=link_dict.get('relationship_description', '?'),
                                     relationship_keywords=link_dict.get('relationship_keywords', []), # Validator handles if not list? No, ensure it is list here.
                                     relationship_strength=link_dict.get('relationship_strength', 5.0)
                                 )
                                 resolved_links_models.append(resolved_link)
                             except (ValidationError, ValueError, TypeError) as link_val_err:
                                  logging.warning(f"Failed to create/validate LinkedNode for target UUID '{target_uuid}' in node '{node_uuid}': {link_val_err}")
                         else:
                             logging.warning(f"Could not resolve target UUID for link target '{target_original_id}' in node '{node_uuid}'. Dropping link.")

                # Update the node dict with the resolved LIST OF LINKEDNODE MODELS
                # The QdrantNode model expects List[LinkedNode]
                node_dict['linked_nodes'] = resolved_links_models

                try:
                    # Validate the whole node dict with Pydantic
                    # QdrantNode expects linked_nodes to be List[LinkedNode] which we now have
                    node_obj = QdrantNode(**node_dict)
                    validated_nodes.append(node_obj) # Append the validated Pydantic object
                    final_node_ids.add(node_obj.node_id) # Add final UUID to set
                except ValidationError as node_error:
                    # Log the problematic dictionary content for debugging
                    logging.warning(f"Pydantic validation failed for node UUID '{node_uuid}' after link resolution: {node_error}. Skipping node. Data: {json.dumps(node_dict, indent=2, default=str)}") # Use default=str for models in dict
                except Exception as final_val_err:
                    logging.error(f"Unexpected error during final node validation UUID '{node_uuid}': {final_val_err}", exc_info=True)

            # Assign the list of validated Pydantic objects
            chapter.qdrant_nodes = validated_nodes
            logging.info(f"Assigned {len(chapter.qdrant_nodes)} validated nodes (using UUIDs) for chapter '{chapter_id}'.")

        else: # qdrant_nodes data was not a list
            logging.error(f"Final data for chapter '{chapter_id}' had 'qdrant_nodes' but not a list.")
            chapter.qdrant_nodes = []
    else: # Initial or fallback LLM call failed completely
        logging.error(f"All node extraction attempts failed for chapter '{chapter_id}'. Assigning empty list.")
        chapter.qdrant_nodes = []

    # Function modifies chapter in-place, no return value needed.

max_structure_retries = 2 # Example value, adjust as needed
retry_delay = 3           # Example value, adjust as needed

async def analyze_inter_chapter_relationships(
    client: genai.Client,
    chapters: List[Chapter],
    source_identifier: Optional[str] = "Document",
    # Added max_structure_retries and retry_delay as parameters for flexibility
    max_structure_retries: int = max_structure_retries,
    retry_delay: int = retry_delay
) -> Optional[Dict[str, Any]]:
    """
    Analyzes inter-chapter relationships AND generates a document summary based on aggregated nodes.
    Updates the `linked_nodes` attribute of the QdrantNode objects within the chapters list *in-place*.
    Returns a dictionary containing the document summary details, or None on failure.
    Uses the specified Gemini client/model.

    Args:
        client: The configured GenAI client/model instance.
        chapters: List of Chapter objects, expected to contain QdrantNode objects.
        source_identifier: Identifier for the source document (e.g., filename).
        max_structure_retries: How many times to retry LLM call on structural errors.
        retry_delay: Base delay (in seconds) between retries.

    Returns:
        A dictionary conforming to DocumentSummaryDetails or None on failure.
    """
    function_name = "analyze_inter_chapter_relationships"
    logging.info(f"[{function_name}] Starting analysis for {len(chapters)} chapters of '{source_identifier}'.")

    # --- 1. Aggregate Node Information & Create Lookup Map ---
    all_nodes_for_prompt = []
    node_map: Dict[str, QdrantNode] = {} # node_id (UUID) -> QdrantNode object

    for chapter in chapters:
        # Ensure chapter.qdrant_nodes exists and is a list
        if not hasattr(chapter, 'qdrant_nodes') or not isinstance(chapter.qdrant_nodes, list):
            logging.warning(f"[{function_name}] Chapter '{getattr(chapter, 'chapter_id', '?')}' is missing 'qdrant_nodes' list. Skipping its nodes.")
            continue

        for node in chapter.qdrant_nodes:
            # Ensure node is a QdrantNode instance and has a node_id
            if not isinstance(node, QdrantNode) or not hasattr(node, 'node_id') or not node.node_id:
                logging.warning(f"[{function_name}] Skipping invalid node object in chapter '{getattr(chapter, 'chapter_id', '?')}': Type {type(node)}")
                continue

            # Prepare simplified node info for the prompt
            # *** THE FIX IS HERE: Use getattr for safe access to direct attributes ***
            node_description = getattr(node, 'core_claim', 'No description available.')
            node_tags = getattr(node, 'tags', '[]') # Get tags, default to empty JSON list string

            nodes_for_prompt_entry = {
                "node_id": node.node_id, # Should be the UUID
                "node_type": getattr(node, 'node_type', 'unknown'),
                "name": getattr(node, 'name', 'Unknown Name'),
                "description": node_description, # Use the safely retrieved core_claim
                "chapter_id_context": getattr(node, 'chapter_id_context', 'unknown_context'),
                "tags": node_tags # Use the safely retrieved tags string
            }
            all_nodes_for_prompt.append(nodes_for_prompt_entry)

            if node.node_id in node_map:
                 # This should ideally not happen if UUIDs are generated correctly
                 logging.warning(f"[{function_name}] Duplicate node_id '{node.node_id}' found while building map. Overwriting entry.")
            node_map[node.node_id] = node # Store the actual QdrantNode object

    if not all_nodes_for_prompt:
        logging.warning(f"[{function_name}] No valid nodes found across all chapters for '{source_identifier}'. Cannot analyze inter-chapter relationships or summarize.")
        # Return None or perhaps a default empty summary structure
        return None

    logging.info(f"[{function_name}] Aggregated {len(all_nodes_for_prompt)} nodes from {len(chapters)} chapters for '{source_identifier}'.")

    # Limit size if needed before dumping (optional)
    # MAX_PROMPT_NODES = 500 # Example limit
    # if len(all_nodes_for_prompt) > MAX_PROMPT_NODES:
    #    logging.warning(f"[{function_name}] Truncating nodes for prompt from {len(all_nodes_for_prompt)} to {MAX_PROMPT_NODES}")
    #    all_nodes_for_prompt = all_nodes_for_prompt[:MAX_PROMPT_NODES] # Or sample strategically

    nodes_json_for_prompt = "[]"
    try:
        nodes_json_for_prompt = json.dumps(all_nodes_for_prompt, indent=2)
    except TypeError as json_err:
         logging.error(f"[{function_name}] Failed to serialize nodes for prompt: {json_err}")
         return None # Cannot proceed without nodes in prompt

    # --- 2. Define the Initial LLM Prompt ---
    # (Using the prompt structure from your provided code)
    initial_inter_chapter_and_summary_prompt = f"""
    ---Goal---
    Analyze the provided list of concepts and entities (nodes) aggregated from all chapters of the document '{source_identifier}'. Perform two tasks:
    1. Identify significant relationships *between* nodes originating from *different* chapters.
    2. Generate a concise document-level summary based on the overall collection of nodes.

    ---Input Data---
    List of node definitions (node_id is the unique identifier):
    ```json
    {nodes_json_for_prompt}
    ```

    ---Instructions---
    **Task 1: Identify Inter-Chapter Relationships**
    1a. Compare nodes across the entire list, focusing on potential connections between nodes from different `chapter_id_context`.
    1b. Identify pairs (`source_node`, `target_node`) with **different** `chapter_id_context` values that have a direct, meaningful relationship based on their names and descriptions/tags. Use the provided `node_id` values for linking.
    1c. For each significant inter-chapter relationship identified, determine: `source_node_id` (the UUID of the source), `target_node_id` (the UUID of the target), `relationship_description` (specific, verb-driven), `relationship_keywords` (list, optional), `relationship_strength` (float 1.0-10.0).
    1d. **Prioritize Quality & Limit:** Focus on the 3-10 *strongest* inter-chapter links *per source node*. Do not exceed 10 outgoing links per source node.

    **Task 2: Generate Document Summary**
    2a. Synthesize the information from the *entire list* of nodes (names, descriptions, types, chapter context, tags).
    2b. Generate the following summary components:
        *   `title`: A concise title for the document '{source_identifier}' based on the aggregated nodes (max 10 words).
        *   `themes`: A list of 3-7 main themes or topics reflected across all nodes.
        *   `questions`: A list of 2-4 insightful questions a reader might have after understanding the key concepts/entities.
        *   `summary`: A comprehensive summary paragraph (4-8 sentences) synthesizing the key nodes and their overall significance or narrative represented by the aggregation.

    **Output Format:**
    Return ONLY a single valid JSON object containing TWO top-level keys: "inter_chapter_links" and "document_summary_details".

    *   `inter_chapter_links`: A JSON array containing objects, where each object represents one identified *inter-chapter* relationship from Task 1 (with keys: `source_node_id`, `target_node_id`, `relationship_description`, `relationship_keywords`, `relationship_strength`). Use the UUIDs provided in the input.
    *   `document_summary_details`: A JSON object containing the results from Task 2 (with keys: `title`, `themes`, `questions`, `summary`).

    Example Output Format:
    ```json
    {{
      "inter_chapter_links": [
        {{
          "source_node_id": "uuid-for-concept_machine_learning_ch2",
          "target_node_id": "uuid-for-entity_gpu_optimization_ch5",
          "relationship_description": "Utilizes machine learning concepts from Chapter 2 for GPU optimization.",
          "relationship_keywords": ["hardware acceleration", "optimization"],
          "relationship_strength": 8.5
        }}
      ],
      "document_summary_details": {{
        "title": "AI Development and Hardware Optimization",
        "themes": ["Machine Learning", "GPU Computing", "Optimization Techniques", "Acme Corp Strategy"],
        "questions": [
          "How does Acme Corp leverage machine learning?",
          "What are the key challenges in GPU optimization mentioned?"
        ],
        "summary": "The document outlines core machine learning concepts introduced in early chapters, detailing Acme Corp's involvement. Later chapters focus on the technical challenges and solutions related to GPU optimization, showing practical application of the earlier concepts."
      }}
    }}
    ```
    CRITICAL: Ensure the entire output is ONLY the single valid JSON object described, starting with `{{` and ending with `}}`. Do not include ```json markdown fences or any other text. Use the provided `node_id` values (which are UUIDs) in the `inter_chapter_links`.
    """

    # --- 3. Initial Call LLM and Process Response ---
    final_llm_output = None
    initial_call_failed_parsing = False
    last_exception = None
    last_raw_response_text = "N/A" # Store for potential fallback

    for attempt in range(max_structure_retries):
        logging.debug(f"[{function_name}] Initial Call Attempt {attempt + 1}/{max_structure_retries} for '{source_identifier}'")
        response = None
        json_text = ""
        data = None
        last_exception = None

        try:
            # Assuming retry_api_call and client setup are correct from provided code
            response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.5-flash-preview-04-17", # Or appropriate model
                contents=[types.Content(parts=[types.Part.from_text(text=initial_inter_chapter_and_summary_prompt)])],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )

            if not response or not response.candidates:
                logging.error(f"[{function_name}] No candidates (Initial Attempt {attempt + 1}) for '{source_identifier}'.")
                last_exception = ValueError("LLM response had no candidates")
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: break # Failed after retries

            # Store raw text before cleaning for fallback
            last_raw_response_text = response.candidates[0].content.parts[0].text
            json_text = clean_json_response(last_raw_response_text) # Assuming clean_json_response exists
            logging.debug(f"[{function_name}] Cleaned Initial Response (Attempt {attempt + 1}, {source_identifier}): {json_text[:500]}...")

            try: data = json.loads(json_text)
            except json.JSONDecodeError as json_err:
                logging.warning(f"[{function_name}] JSONDecodeError initial attempt {attempt + 1} for '{source_identifier}': {json_err}.")
                last_exception = json_err
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: initial_call_failed_parsing = True; break # Failed parsing after retries

            # --- Structure Check ---
            if (isinstance(data, dict) and
                "inter_chapter_links" in data and isinstance(data.get("inter_chapter_links"), list) and
                "document_summary_details" in data and isinstance(data.get("document_summary_details"), dict)):
                logging.info(f"[{function_name}] Received correct structure initial attempt {attempt + 1} for '{source_identifier}'.")
                final_llm_output = data # Success!
                initial_call_failed_parsing = False
                break # Exit retry loop
            else:
                logging.warning(f"[{function_name}] Incorrect JSON structure initial attempt {attempt+1} for '{source_identifier}'. Got: {type(data)}. Structure: {str(data)[:300]}...")
                last_exception = ValueError(f"Incorrect JSON structure received (type: {type(data)})")
                if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
                else: initial_call_failed_parsing = True; break # Failed structure check after retries

        except Exception as e:
            logging.error(f"[{function_name}] Error during initial attempt {attempt + 1} for '{source_identifier}': {e}", exc_info=True)
            last_exception = e
            if attempt < max_structure_retries - 1: await asyncio.sleep(retry_delay * (attempt + 1)); continue
            else: initial_call_failed_parsing = False; break # API call itself failed

    # --- 4. Fallback Attempt (if initial parsing failed on last try) ---
    # (Using the fallback logic from your provided code)
    if initial_call_failed_parsing:
        logging.warning(f"[{function_name}] Initial call failed structure/parsing for '{source_identifier}'. Attempting fallback using input nodes.")
        # Construct Fallback Prompt (using the same node list as input)
        fallback_prompt = f"""
        ---Goal---
        A previous attempt to analyze the provided nodes failed to produce the correct JSON structure. Please try again based *only* on the input node data below.
        Analyze the provided list of concepts and entities (nodes) aggregated from all chapters of the document '{source_identifier}'. Perform two tasks:
        1. Identify significant relationships *between* nodes originating from *different* chapters. Use the provided `node_id` values for linking.
        2. Generate a concise document-level summary based on the overall collection of nodes.

        ---Input Data (Same as before)---
        List of node definitions:
        ```json
        {nodes_json_for_prompt}
        ```

        ---Instructions (Same as before)---
        **Task 1: Identify Inter-Chapter Relationships**
        (Instructions identical to initial prompt...)
        1d. **Prioritize Quality & Limit:** Focus on the 3-10 *strongest* inter-chapter links *per source node*. Do not exceed 10 outgoing links per source node.

        **Task 2: Generate Document Summary**
        (Instructions identical to initial prompt...)

        **Output Format (Same as before):**
        Return ONLY a single valid JSON object containing TWO top-level keys: "inter_chapter_links" and "document_summary_details". Ensure `node_id` values used in links are the UUIDs provided in the input.
        ```json
        {{
          "inter_chapter_links": [ ... ],
          "document_summary_details": {{ ... }}
        }}
        ```
        Ensure the final output is only the single JSON object described.
        """

        # Execute Fallback LLM Call (1 attempt usually sufficient)
        try:
            logging.debug(f"[{function_name}] Making fallback LLM call for '{source_identifier}'.")
            fallback_response = await retry_api_call(
                client.aio.models.generate_content,
                model="gemini-2.0-flash", # Can use flash for retry
                contents=[types.Content(parts=[types.Part.from_text(text=fallback_prompt)])],
                config=types.GenerateContentConfig(response_mime_type="application/json"),
                max_retries=1 # Only retry once on API errors for fallback
            )

            if fallback_response and fallback_response.candidates:
                fb_json_text = clean_json_response(fallback_response.candidates[0].content.parts[0].text)
                logging.debug(f"[{function_name}] Cleaned Fallback Response ({source_identifier}): {fb_json_text[:500]}...")
                try:
                    fb_data = json.loads(fb_json_text)
                    # Final Structure Check (Fallback)
                    if (isinstance(fb_data, dict) and
                        "inter_chapter_links" in fb_data and isinstance(fb_data.get("inter_chapter_links"), list) and
                        "document_summary_details" in fb_data and isinstance(fb_data.get("document_summary_details"), dict)):
                        logging.info(f"[{function_name}] Fallback attempt SUCCEEDED with correct structure for '{source_identifier}'.")
                        final_llm_output = fb_data # Assign data from successful fallback
                    else:
                        logging.error(f"[{function_name}] Fallback attempt FAILED structure check for '{source_identifier}'. Got type: {type(fb_data)}")
                except json.JSONDecodeError as fallback_json_err:
                    logging.error(f"[{function_name}] Fallback JSONDecodeError for '{source_identifier}': {fallback_json_err}")
            else:
                logging.error(f"[{function_name}] Fallback LLM call returned no candidates for '{source_identifier}'.")

        except Exception as fallback_err:
            logging.error(f"[{function_name}] Error during fallback LLM call for '{source_identifier}': {fallback_err}", exc_info=True)
            # Fallback failed, final_llm_output remains None

    # --- 5. Process Final LLM Output (if successful from initial or fallback) ---
    document_summary_output = None
    if final_llm_output is not None:
        logging.info(f"[{function_name}] Processing final LLM output for '{source_identifier}'.")
        # --- Process Inter-Chapter Links ---
        newly_found_links_data = final_llm_output.get("inter_chapter_links", [])
        if isinstance(newly_found_links_data, list):
            # Assuming _merge_links_into_nodes exists and modifies node_map correctly
            # Ensure _merge_links_into_nodes correctly handles QdrantNode objects
            links_added_count, links_skipped_count = _merge_links_into_nodes(newly_found_links_data, node_map)
            logging.info(f"[{function_name}] Merged inter-chapter links for '{source_identifier}'. Added: {links_added_count}, Skipped: {links_skipped_count}")
        else:
            logging.warning(f"[{function_name}] Key 'inter_chapter_links' missing or not a list in final LLM output for '{source_identifier}'.")

        # --- Process Document Summary Details ---
        summary_details_data = final_llm_output.get("document_summary_details")
        if isinstance(summary_details_data, dict):
             try:
                 # Validate against the Pydantic model
                 validated_summary = DocumentSummaryDetails(**summary_details_data)
                 document_summary_output = validated_summary.model_dump() # Store as dict
                 logging.info(f"[{function_name}] Successfully extracted and validated document summary details for '{source_identifier}'.")
             except ValidationError as summary_error:
                 logging.error(f"[{function_name}] Validation failed for document_summary_details from LLM for '{source_identifier}': {summary_error}")
                 # Create a default error summary structure
                 document_summary_output = {
                    "title": f"Summary Validation Failed: {source_identifier}",
                    "themes": ["error"], "questions": [],
                    "summary": f"LLM returned summary data, but it failed validation: {summary_error}",
                 }
             except Exception as unexpected_summary_err:
                 logging.error(f"[{function_name}] Unexpected error processing summary details for '{source_identifier}': {unexpected_summary_err}", exc_info=True)
                 document_summary_output = {
                     "title": f"Summary Processing Error: {source_identifier}",
                     "themes": ["error"], "questions": [],
                     "summary": f"An unexpected error occurred during summary processing: {unexpected_summary_err}",
                 }
        else:
            logging.warning(f"[{function_name}] Key 'document_summary_details' missing or not a dict in final LLM output for '{source_identifier}'.")
            # Optionally create a default summary if it's missing entirely
            document_summary_output = {
                "title": f"Summary Missing: {source_identifier}", "themes": [], "questions": [],
                "summary": "LLM did not provide document summary details.",
            }
    else:
        # All attempts (initial + fallback) failed to produce valid structure
        logging.error(f"[{function_name}] All LLM attempts failed for '{source_identifier}'. Last error: {last_exception}")
        # Create a default error summary
        document_summary_output = {
            "title": f"Summary Generation Failed: {source_identifier}", "themes": ["error"], "questions": [],
            "summary": f"Failed to generate document summary via LLM after multiple attempts. Last error: {last_exception}",
        }

    # The chapters list (containing QdrantNode objects, accessible via node_map) has been modified in-place by _merge_links_into_nodes
    logging.info(f"[{function_name}] Finished analysis for '{source_identifier}'. Returning summary details.")
    return document_summary_output # Returns summary dict (validated or default/error)


# --- Helper function to merge links (extracted from previous version) ---
def _merge_links_into_nodes(link_data_list: List[Dict], node_map: Dict[str, QdrantNode]) -> (int, int):
    """Merges link data into the linked_nodes of QdrantNode objects in node_map."""
    links_added_count = 0
    links_skipped_count = 0
    MAX_OUTGOING_LINKS_PER_NODE = 10

    if not isinstance(link_data_list, list):
        logging.warning("[_merge_links_into_nodes] Input link_data_list is not a list.")
        return 0, 0

    for link_data in link_data_list:
        if not isinstance(link_data, dict):
            logging.warning(f"[_merge_links_into_nodes] Skipping invalid item in link data: {link_data}")
            links_skipped_count += 1
            continue

        source_id = link_data.get("source_node_id")
        target_id = link_data.get("target_node_id")
        desc = link_data.get("relationship_description")
        keywords = link_data.get("relationship_keywords", [])
        strength_str = link_data.get("relationship_strength", "5.0")

        if not all([source_id, target_id, desc]) or not isinstance(source_id, str) or not isinstance(target_id, str):
            logging.warning(f"[_merge_links_into_nodes] Skipping incomplete/invalid link data: {link_data}")
            links_skipped_count += 1
            continue

        try: strength = float(strength_str)
        except (ValueError, TypeError): strength = 5.0

        source_node = node_map.get(source_id)
        target_node = node_map.get(target_id)

        if not source_node or not target_node:
            logging.warning(f"[_merge_links_into_nodes] Source ('{source_id}') or Target ('{target_id}') not found. Skipping.")
            links_skipped_count += 1
            continue

        if source_node.chapter_id_context == target_node.chapter_id_context:
            logging.debug(f"[_merge_links_into_nodes] Skipping INTRA-chapter link: {source_id} -> {target_id}")
            links_skipped_count += 1
            continue

        # Add Forward Link
        if source_node.linked_nodes is None: source_node.linked_nodes = []
        if len(source_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == target_id for link in source_node.linked_nodes):
            try:
                forward_link = LinkedNode(**link_data) # Directly use link_data if keys match LinkedNode
                source_node.linked_nodes.append(forward_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added FORWARD link: {source_id} -> {target_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating forward LinkedNode {source_id} -> {target_id}: {ve}. Skipping.")
                 links_skipped_count += 1
        else:
             links_skipped_count += 1 # Skipped due to limit or duplicate

        # Add Reverse Link
        if target_node.linked_nodes is None: target_node.linked_nodes = []
        if len(target_node.linked_nodes) < MAX_OUTGOING_LINKS_PER_NODE and not any(link.target_node_id == source_id for link in target_node.linked_nodes):
            try:
                reverse_desc = f"Related to '{source_node.name}' from {source_node.chapter_id_context} ({desc})"
                reverse_link = LinkedNode(target_node_id=source_id, relationship_description=reverse_desc, relationship_keywords=keywords, relationship_strength=strength)
                target_node.linked_nodes.append(reverse_link)
                links_added_count += 1
                logging.debug(f"[_merge_links_into_nodes] Added REVERSE link: {target_id} -> {source_id}")
            except ValidationError as ve:
                 logging.warning(f"[_merge_links_into_nodes] Validation failed creating reverse LinkedNode {target_id} -> {source_id}: {ve}. Skipping.")
                 # links_skipped_count += 1 # Optional separate count
        else:
             # links_skipped_count += 1 # Optional separate count
             pass # Skip reverse if limit reached or duplicate


    return links_added_count, links_skipped_count

async def incrementally_generate_doc_level_ontology(
    client: genai.Client,
    documents: List[Dict],
    *,
    ontology_path: str | pathlib.Path = "doc_level_ontology.json",
    char_budget: int = 100_000,
    max_new_links: int = 50,
    branch_size: int = 20,
    llm_retries: int = 2,
    llm_delay: int = 5
) -> Optional[ProjectOntology]:
    """
    Builds or incrementally updates ontology using DIRECT Gemini API calls.
    Includes nested tree summarizer and internal constants.
    Outputs nodes with FLATTENED structure.

    Args:
        client: The configured Generative AI client instance.
        documents: List of *all* document dictionaries.
        ontology_path: Path to load/save the ontology JSON file.
        char_budget: Max characters for prompts (affects summarization).
        max_new_links: Max new links to generate in incremental step.
        branch_size: Fan-out for tree summarization.
        llm_retries: Number of retries for Gemini API calls.
        llm_delay: Base delay (seconds) between retries.

    Returns:
        The generated/updated ProjectOntology object, or None on critical failure.
    """
    # --- Setup Logging and Constants ---
    log = logging.getLogger("incremental_doc_ontology")
    path = pathlib.Path(ontology_path)
    # Constants are now local to this function scope
    CHAR_BUDGET = char_budget
    MAX_NEW_LINKS = max_new_links
    BRANCH_SIZE = branch_size
    LLM_RETRIES = llm_retries
    LLM_DELAY = llm_delay

    # ----------------------------------------------------------------------- #
    # ❹ Tree Summarizer (Nested Function)                                     #
    # ----------------------------------------------------------------------- #
    async def _tree_summarise_nested(
        input_docs: List[Dict[str, Any]],
        *,
        target_chars: int # Target chars passed explicitly
    ) -> List[Dict[str, Any]]:
        """
        NESTED: Recursively compress docs until JSON size <= target_chars.
        Uses client, BRANCH_SIZE, CHAR_BUDGET, LLM_RETRIES, LLM_DELAY from outer scope.
        """
        tree_log = logging.getLogger("tree_summarise_nested") # Use specific logger

        def blob_size(nodes_list: List[Dict]) -> int:
            try: return len(json.dumps(nodes_list, ensure_ascii=False))
            except TypeError: return float('inf')

        current_size = blob_size(input_docs)
        tree_log.debug(f"Input docs: {len(input_docs)}, Current size: {current_size}, Target size: {target_chars}")

        if current_size <= target_chars: return input_docs
        if not input_docs: return []
        if len(input_docs) == 1 and current_size > target_chars:
            tree_log.warning(f"Single document exceeds target size ({current_size} > {target_chars}). Cannot compress further.")
            return input_docs

        groups = [input_docs[i:i + BRANCH_SIZE] for i in range(0, len(input_docs), BRANCH_SIZE)]
        tree_log.info(f"Compressing {len(input_docs)} docs into {len(groups)} groups (branch size {BRANCH_SIZE}).")

        compressed_results: List[Dict] = []
        tasks = []

        # --- Direct API Call Task Creation within Nested Tree Summarizer ---
        async def call_gemini_for_summary(prompt: str, group_idx: int) -> Optional[Dict]:
            """Inner function for single API call with retry & JSON parsing"""
            # Uses LLM_RETRIES, LLM_DELAY, client from outer scope
            for attempt in range(LLM_RETRIES):
                try:
                    tree_log.debug(f"Calling Gemini for group {group_idx} (Attempt {attempt+1}/{LLM_RETRIES})")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text
                    raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    parsed_json = json.loads(raw_text)
                    return parsed_json
                except Exception as e:
                    tree_log.warning(f"Gemini call/parse failed group {group_idx} (Attempt {attempt+1}/{LLM_RETRIES}): {e}")
                    if attempt == LLM_RETRIES - 1:
                        tree_log.error(f"Gemini call failed group {group_idx} after {LLM_RETRIES} attempts.")
                        return None
                    await asyncio.sleep(LLM_DELAY * (attempt + 1))
            return None

        for idx, chunk in enumerate(groups, 1):
            if not chunk: continue
            chunk_for_prompt = []
            estimated_prompt_size = 1000
            # Use CHAR_BUDGET from outer scope
            chars_per_doc_limit = (CHAR_BUDGET - estimated_prompt_size) // len(chunk) if len(chunk) > 0 else CHAR_BUDGET
            chars_per_doc_limit = max(500, chars_per_doc_limit)
            for doc_item in chunk:
                 summary = doc_item.get("executive_summary", "")
                 truncated_summary = summary[:chars_per_doc_limit] + "..." if len(summary) > chars_per_doc_limit else summary
                 chunk_for_prompt.append({"node_id": doc_item.get("node_id"), "document_name": doc_item.get("document_name"), "executive_summary_preview": truncated_summary})

            chunk_json = json.dumps(chunk_for_prompt, indent=2, ensure_ascii=False)
            # Use CHAR_BUDGET from outer scope
            if len(chunk_json) > (CHAR_BUDGET * 0.8): chunk_json = chunk_json[:int(CHAR_BUDGET * 0.8)] + "\n... ]\n```"

            prompt = f"""
                You are an expert summarizer consolidating information from multiple document summaries.
                Input: {len(chunk)} summaries/previews:
                ```json
                {chunk_json}
                ```
                Goal: Create **one single, concise meta-summary node** representing this group.
                Instructions: Synthesize key ideas into an executive summary (<800 words). Assign node_id "meta::{idx}", document_name "meta_group_{idx}".
                Output Format: **Valid JSON only**: `{{"meta_node": {{"node_id": "meta::{idx}", "document_name": "meta_group_{idx}", "executive_summary": "<summary>"}}}}`
                Group Index: {idx}
                Generate JSON:"""
            tasks.append(asyncio.create_task(call_gemini_for_summary(prompt, idx), name=f"SummarizeGroup_{idx}"))

        task_results = await asyncio.gather(*tasks, return_exceptions=False)

        for i, result in enumerate(task_results, 1):
            if result is None: tree_log.error(f"Skipping group {i} due to API/JSON failure."); continue
            if isinstance(result, dict) and "meta_node" in result and isinstance(result["meta_node"], dict):
                meta_node = result["meta_node"]
                if ("node_id" in meta_node and "document_name" in meta_node and "executive_summary" in meta_node): compressed_results.append(meta_node)
                else: tree_log.warning(f"Invalid meta_node structure group {i}: {meta_node}")
            else: tree_log.warning(f"Unexpected result type group {i}: {type(result)}")

        if not compressed_results:
             tree_log.error("All summarization tasks failed. Returning original docs.")
             return input_docs # Return original to avoid losing data

        tree_log.debug(f"Recursion step. Input docs: {len(compressed_results)}, Target chars: {target_chars}")
        # Recursive call to the *nested* function
        return await _tree_summarise_nested(compressed_results, target_chars=target_chars)
    # ----------------------------------------------------------------------- #
    # End of Nested Tree Summarizer Definition                                #
    # ----------------------------------------------------------------------- #


    # --- 1. Load Existing Ontology or Initialize ---
    ontology = load_ontology_json(path) # Uses external helper
    first_run = ontology is None
    if first_run:
        log.info(f"No ontology file found at '{path}' – creating new.")
        ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
    else:
        if not isinstance(ontology, ProjectOntology):
             log.error(f"Loaded data from '{path}' invalid. Starting fresh.")
             first_run = True
             ontology = ProjectOntology(title="", overview="", document_count=0, documents=[], global_themes=[], key_concepts=[], project_graph_nodes=[])
        else:
            log.info(f"Loaded existing ontology with {ontology.document_count} documents from '{path}'.")


    # --- 2. Identify New Documents ---
    existing_files = set(ontology.documents)
    target_docs_data = []
    if first_run: target_docs_data = documents
    else:
        new_doc_count = 0
        for d in documents:
            try:
                filename = d.get("raw_extracted_content", {}).get("filename")
                if filename and isinstance(filename, str) and filename not in existing_files: target_docs_data.append(d); new_doc_count += 1
            except Exception as e: log.warning(f"Error accessing filename: {e} - Skipping.")
        log.info(f"Incremental run: Found {new_doc_count} new documents.")

    if not target_docs_data: log.info("No new documents found."); return ontology


    # --- 3. Prepare Nodes from New Document Summaries ---
    doc_nodes_for_llm: List[Dict] = []
    processed_filenames = []
    for d in target_docs_data:
        try:
            raw = d.get("raw_extracted_content", {}); fname = raw.get("filename", f"unknown_doc_{random.randint(1000,9999)}")
            # --- Extract summary information robustly ---
            summ_block = raw.get("summary")
            summary_text = ""
            # Handle case where summary is a DocumentSummaryDetails model (or dict)
            if isinstance(summ_block, dict):
                summary_text = summ_block.get("summary") or summ_block.get("overview", "")
            # Handle if summary is somehow just a string (less likely now)
            elif isinstance(summ_block, str):
                summary_text = summ_block

            summary_text = summary_text.strip()

            if not summary_text:
                log.warning(f"Document '{fname}' lacks summary text – skipping.")
                continue

            doc_nodes_for_llm.append({
                "node_id": f"doc::{fname}",
                "document_name": fname,
                "executive_summary": summary_text # Use the extracted summary
            })
            processed_filenames.append(fname)
        except Exception as e: log.warning(f"Error processing doc data '{raw.get('filename', 'unknown')}': {e} - Skipping.")

    if not doc_nodes_for_llm:
        log.error("No usable summaries found in new documents to process for ontology.")
        return ontology if not first_run else None # Return existing or None if first run
    log.info(f"Prepared {len(doc_nodes_for_llm)} nodes from new summaries.")


    # --- 4. Apply Tree Summarization (using nested function) if Needed ---
    summaries_json_blob = json.dumps(doc_nodes_for_llm, ensure_ascii=False)
    current_char_size = len(summaries_json_blob)
    final_doc_nodes = doc_nodes_for_llm # Default

    if current_char_size > CHAR_BUDGET: # Use local constant
        log.info(f"Summary size ({current_char_size:,} chars) > budget ({CHAR_BUDGET:,}). Summarizing...")
        try:
            # Call the NESTED summarizer
            final_doc_nodes = await _tree_summarise_nested(
                doc_nodes_for_llm,
                target_chars=CHAR_BUDGET # Pass target explicitly
            )
            final_size = len(json.dumps(final_doc_nodes, ensure_ascii=False))
            log.info(f"Summarization complete. Final size: {final_size:,} chars ({len(final_doc_nodes)} nodes).")
            if final_size > CHAR_BUDGET: log.warning(f"Summarization didn't meet budget.")
        except Exception as e:
            log.error(f"Error during tree summarization: {e}. Proceeding with original nodes.", exc_info=True)
            final_doc_nodes = doc_nodes_for_llm # Fallback
    else:
        log.info(f"Summary size ({current_char_size:,} chars) within budget.")


    # --- 5. Call LLM for Ontology Generation / Update (Direct Calls) ---
    llm_links = []
    llm_summary_details = {}
    llm_call_failed = False

    max_json_chars_in_prompt = int(CHAR_BUDGET * 0.85) # Use local constant
    final_doc_nodes_json_str = json.dumps(final_doc_nodes, indent=2, ensure_ascii=False)
    if len(final_doc_nodes_json_str) > max_json_chars_in_prompt:
        log.warning(f"Final doc nodes JSON ({len(final_doc_nodes_json_str)} chars) truncated for prompt.")
        final_doc_nodes_json_str = final_doc_nodes_json_str[:max_json_chars_in_prompt] + "\n... ]\n```"

    try:
        if first_run:
            log.info("Performing initial ontology generation LLM call (direct)...")
            # NOTE: Prompt needs updating to reflect flat node input for existing nodes if sampled
            init_prompt = f"""Act as KG Architect. Create initial ontology from document/meta summaries.
                Input: Document/Meta Summaries:
                ```json
                {final_doc_nodes_json_str}
                ```
                Tasks:
                1. Find significant inter-document relationships between the nodes provided in the input (using their `node_id` which is like 'doc::filename' or 'meta::idx'). Base relationships on the content of their `executive_summary`.
                2. Generate project metadata (title, overview, global_themes, key_documents). `key_documents` should be a list of the original document names (e.g., "report.pdf") deemed most important.

                Output Schema (JSON ONLY): `{{"project_summary_details": {{"title": "...", "overview": "...", "global_themes": ["..."], "key_documents": ["doc_name_1.pdf", "..."]}}, "inter_document_links": [ {{ "source_node_id": "doc::source_doc_name.pdf", "target_node_id": "doc::target_doc_name.pdf", "relationship_description": "...", "relationship_keywords": ["..."], "relationship_strength": 7.5 }} ]}}`
                Rules: Valid JSON. Use the provided `node_id`s for links. Base summary/links SOLELY on the executive summaries provided. `key_documents` MUST be original filenames.
                Generate JSON:"""
            llm_out = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Initial Ontology: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=init_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    llm_out = json.loads(raw_text); log.info(f"Initial ontology LLM call successful."); break
                except Exception as e:
                    log.warning(f"Initial Ontology: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Initial Ontology: Call failed after retries."); llm_call_failed = True; break
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if llm_out: llm_links = llm_out.get("inter_document_links", []); llm_summary_details = llm_out.get("project_summary_details", {})
            else: llm_links = []; llm_summary_details = {}

        else: # Incremental update
            log.info("Performing incremental ontology update LLM calls (direct)...")
            # --- 5a. Link Discovery ---
            existing_nodes_for_context = []
            if ontology.project_graph_nodes:
                 # --- CORRECTED: Sample using direct attributes ---
                 key_doc_node_ids = {f"doc::{kd}" for kd in ontology.key_concepts if isinstance(kd, str)} # key_concepts are filenames
                 # Build map using UUID -> Pydantic object
                 existing_node_map_obj: Dict[str, QdrantNode] = {n.node_id: n for n in ontology.project_graph_nodes}

                 # Get priority nodes based on key concepts (if node exists)
                 priority_nodes_obj: List[QdrantNode] = []
                 for node_uuid, node_obj in existing_node_map_obj.items():
                     # Check if the node represents a key document/concept
                     # This might need refinement - how are key_concepts linked back to nodes? Assuming name is the filename for doc nodes.
                     if node_obj.node_type == "document" and node_obj.name in ontology.key_concepts:
                          priority_nodes_obj.append(node_obj)

                 other_nodes_obj = [n for n in ontology.project_graph_nodes if n not in priority_nodes_obj]

                 sample_size=min(len(ontology.project_graph_nodes), max(len(final_doc_nodes)*2, 30))
                 sample_count=min(len(ontology.project_graph_nodes), sample_size)

                 sample_nodes_pydantic = priority_nodes_obj # Start with priority nodes
                 remaining_sample_count = sample_count - len(sample_nodes_pydantic)

                 if remaining_sample_count > 0 and other_nodes_obj:
                      sample_nodes_pydantic.extend(random.sample(other_nodes_obj, min(remaining_sample_count, len(other_nodes_obj))))

                 max_summary_preview=800
                 for n in sample_nodes_pydantic: # n is QdrantNode object
                      # Access direct attributes for summary preview
                      summary_text = getattr(n, 'executive_summary', '') or getattr(n, 'core_claim', '')
                      summary_preview = summary_text[:max_summary_preview] + ("..." if len(summary_text) > max_summary_preview else "")
                      # Use node_id (UUID) and name directly
                      existing_nodes_for_context.append({
                          "node_id": n.node_id,
                          "document_name": getattr(n, 'name', 'Unknown'),
                          "executive_summary_preview": summary_preview
                      })
                 # --- End CORRECTED sampling ---

            existing_context_json_str = json.dumps(existing_nodes_for_context, indent=2, ensure_ascii=False)
            max_context_json_chars = int(CHAR_BUDGET * 0.4) # Use local constant
            if len(existing_context_json_str) > max_context_json_chars: existing_context_json_str = existing_context_json_str[:max_context_json_chars] + "\n... ]\n```"

            # Note: Prompt needs updating to reflect flat node input for existing nodes
            link_prompt = f"""Act as KG Analyst. Update ontology. Find links involving NEW docs/summaries.
                Input 1 (New Docs/Summaries):
                ```json
                {final_doc_nodes_json_str}
                ```
                Input 2 (Existing Node Sample - UUIDs used):
                ```json
                {existing_context_json_str}
                ```
                Task: Find links connecting New nodes (from Input 1) to either other New nodes OR Existing nodes (from Input 2). Use the `node_id` provided in BOTH inputs for linking. Base relationships on `executive_summary` content. Limit: {MAX_NEW_LINKS} links.
                Output (JSON ONLY): `{{"new_inter_document_links": [ {{ "source_node_id": "...", "target_node_id": "...", "relationship_description": "...", "relationship_keywords": ["..."], "relationship_strength": 7.5 }} ]}}`
                Rules: Valid JSON. Use the provided `node_id`s (which are UUIDs for existing, doc::/meta:: for new). At least one node in each link MUST be from Input 1. Empty list [] if none.
                Generate JSON:"""
            link_json = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Incr Links: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=link_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    link_json = json.loads(raw_text); log.info(f"Incr link discovery call successful."); break
                except Exception as e:
                    log.warning(f"Incr Links: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Incr Links: Call failed after retries."); break # Continue even if links fail
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if link_json: llm_links = link_json.get("new_inter_document_links", [])
            else: llm_links = []
            log.info(f"Found {len(llm_links)} potential new links.")

            # --- 5b. Summary Revision ---
            current_summary = {"title": ontology.title, "overview": ontology.overview, "global_themes": ontology.global_themes, "key_documents": ontology.key_concepts} # Use key_concepts which are filenames
            current_summary_json_str = json.dumps(current_summary, indent=2)
            rev_prompt = f"""Act as Content Strategist. Update project summary with new docs/summaries.
                Input 1 (Current Summary):
                ```json
                {current_summary_json_str}
                ```
                Input 2 (New Docs/Summaries):
                ```json
                {final_doc_nodes_json_str}
                ```
                Task: Revise title, overview, global_themes, key_documents. Integrate info from Input 2. `key_documents` MUST list original document filenames.
                Output (JSON ONLY): `{{"revised_project_summary_details": {{ "title": "...", "overview": "...", "global_themes": ["..."], "key_documents": ["doc_name_1.pdf", "..."] }}}}`
                Rules: Valid JSON. Integrate new info coherently.
                Generate JSON:"""
            summary_json = None
            for attempt in range(LLM_RETRIES): # Use local constant
                try:
                    log.debug(f"Summary Rev: Call Attempt {attempt+1}/{LLM_RETRIES}")
                    response = await retry_api_call(
                        client.aio.models.generate_content,
                        model="gemini-2.5-flash-preview-04-17",
                        contents=[types.Content(parts=[types.Part.from_text(text=rev_prompt)])],
                        config=types.GenerateContentConfig(response_mime_type="application/json")
                    )
                    if not response or not response.candidates: raise ValueError("Missing candidates")
                    raw_text = response.candidates[0].content.parts[0].text; raw_text = raw_text.strip().lstrip("```json").rstrip("```").strip()
                    summary_json = json.loads(raw_text); log.info("Summary revision call successful."); break
                except Exception as e:
                    log.warning(f"Summary Rev: Attempt {attempt+1}/{LLM_RETRIES} failed: {e}")
                    if attempt == LLM_RETRIES - 1: log.error("Summary Rev: Call failed after retries."); llm_call_failed = True; break
                    await asyncio.sleep(LLM_DELAY * (attempt + 1)) # Use local constant

            if summary_json: llm_summary_details = summary_json.get("revised_project_summary_details", {})
            else: llm_summary_details = current_summary # Fallback to current

    except Exception as e: log.error(f"Critical error during LLM processing: {e}", exc_info=True); llm_call_failed = True

    if llm_call_failed and first_run: log.error("Initial LLM calls failed critically. Aborting."); return None
    elif llm_call_failed: log.warning("One+ incremental LLM calls failed. Proceeding with partial update.")

    # --- 6. Merge LLM Results into Ontology ---
    node_map: Dict[str, Dict] = {} # Map UUID -> node dict (FLAT STRUCTURE)
    if not first_run and ontology and ontology.project_graph_nodes:
         for node in ontology.project_graph_nodes: # These are QdrantNode objects
              try:
                  # Dump the existing Pydantic object to a flat dictionary
                  node_map[node.node_id] = node.model_dump(mode='json')
              except Exception as dump_err:
                  log.warning(f"Could not dump existing node {node.node_id}: {dump_err}")

    # Process the NEW nodes (doc/meta summaries) from the current run (final_doc_nodes)
    new_node_uuid_map: Dict[str, str] = {} # Map original doc/meta ID -> generated UUID
    current_iso_ts = datetime.datetime.now(datetime.timezone.utc).isoformat() # Get timestamp once

    for new_node_data in final_doc_nodes: # new_node_data is like {"node_id": "doc::fname", ...}
         original_node_id = new_node_data.get("node_id")
         doc_name = new_node_data.get("document_name")
         exec_summary = new_node_data.get("executive_summary")
         if not original_node_id or not doc_name or exec_summary is None: continue

         node_type = "document" if original_node_id.startswith("doc::") else "meta_summary"
         context_for_id = original_node_id
         node_uuid = generate_qdrant_id(node_type, doc_name, context_for_id)
         new_node_uuid_map[original_node_id] = node_uuid

         chapter_context_value = f"doc_level::{doc_name}"

         # *** CORRECTED: Create/update node dictionary with FLAT structure ***
         node_map[node_uuid] = {
             "node_id": node_uuid,
             "node_type": node_type,
             "name": doc_name,
             "chapter_id_context": chapter_context_value,
             # --- Direct fields ---
             "title": f"Summary: {doc_name}",
             "core_claim": exec_summary[:800] + ("..." if exec_summary and len(exec_summary) > 800 else ""), # Use summary as core claim
             "information_type": "Summary",
             "key_entities": "[]", # Defaults for summary nodes
             "tags": "[]",
             "sentiment_tone": "Neutral",
             "source_credibility": "Generated Summary" if node_type == "meta_summary" else "Document Summary",
             # --- Programmatic fields ---
             "source": doc_name,
             "stored_at": current_iso_ts,
             "original_node_id": original_node_id,
             "executive_summary": exec_summary, # Store the full summary here as well
             # --- Relationships ---
             "linked_nodes": node_map.get(node_uuid, {}).get("linked_nodes", []) # Preserve existing links
             # Removed properties dict entirely
         }


    # --- Resolve and Merge LLM Links ---
    resolved_llm_links = []
    if isinstance(llm_links, list):
        for link_dict in llm_links:
             source_original_id = link_dict.get("source_node_id")
             target_original_id = link_dict.get("target_node_id")

             # Resolve source UUID (could be new doc::/meta:: or existing UUID)
             source_uuid = new_node_uuid_map.get(source_original_id)
             if not source_uuid: source_uuid = source_original_id if source_original_id in node_map else None

             # Resolve target UUID (could be new doc::/meta:: or existing UUID)
             target_uuid = new_node_uuid_map.get(target_original_id)
             if not target_uuid: target_uuid = target_original_id if target_original_id in node_map else None

             if source_uuid and target_uuid:
                 resolved_link_dict = link_dict.copy()
                 resolved_link_dict["source_node_id"] = source_uuid
                 resolved_link_dict["target_node_id"] = target_uuid
                 resolved_llm_links.append(resolved_link_dict)
             else:
                 log.warning(f"Could not resolve UUIDs for link: {source_original_id} -> {target_original_id}. Skipping.")

    # *** CORRECTED: Call the RENAMED helper function ***
    added_links, skipped_links = _merge_project_links_into_node_dicts(resolved_llm_links, node_map)
    log.info(f"Merged LLM links: {added_links} added, {skipped_links} skipped.")


    # --- Validate Nodes before Final Assignment ---
    validated_nodes_list: List[QdrantNode] = []
    final_node_ids = set()
    for node_id, node_dict in node_map.items(): # node_dict should now be FLAT
        try:
            # Ensure linked_nodes is list of dicts suitable for Pydantic validation
            raw_links = node_dict.setdefault('linked_nodes', [])
            valid_link_dicts = []
            if isinstance(raw_links, list):
                 for link in raw_links:
                     if isinstance(link, LinkedNode): # Convert Pydantic models back to dicts
                         valid_link_dicts.append(link.model_dump(mode='json'))
                     elif isinstance(link, dict):
                         # Basic check/fix for strength before validation
                         try: link['relationship_strength'] = float(link.get('relationship_strength', 5.0))
                         except (ValueError, TypeError): link['relationship_strength'] = 5.0
                         valid_link_dicts.append(link)
            node_dict['linked_nodes'] = valid_link_dicts # Assign cleaned list of dicts

            # Validate the flat node dict against the QdrantNode model
            validated_node = QdrantNode.model_validate(node_dict) # Use model_validate for dicts
            validated_nodes_list.append(validated_node)
            final_node_ids.add(node_id)
        except ValidationError as e:
            log.error(f"Pydantic validation failed node '{node_id}'. Skipping. Error: {e} Data: {json.dumps(node_dict, indent=2, default=str)}") # Log the dict
        except Exception as final_val_err:
            log.error(f"Unexpected validation error node '{node_id}': {final_val_err}", exc_info=True)

    # --- Filter links pointing to non-validated nodes ---
    final_validated_nodes_list: List[QdrantNode] = []
    for node in validated_nodes_list:
        # linked_nodes are now dicts, need to convert to LinkedNode to check target_node_id easily?
        # Or just check the dict structure
        valid_final_links = []
        for link_item in node.linked_nodes: # node.linked_nodes is now List[LinkedNode]
             if link_item.target_node_id in final_node_ids:
                  valid_final_links.append(link_item)
        node.linked_nodes = valid_final_links
        final_validated_nodes_list.append(node)


    log.info(f"Validated {len(final_validated_nodes_list)} final nodes.")


    # --- Update Ontology Metadata ---
    ontology.title = llm_summary_details.get("title", ontology.title or "Project Document Ontology")
    ontology.overview = llm_summary_details.get("overview", ontology.overview or "Analysis of project documents.")
    raw_themes = llm_summary_details.get("global_themes", ontology.global_themes)
    ontology.global_themes = [str(t) for t in raw_themes if isinstance(t, str)] if isinstance(raw_themes, list) else ontology.global_themes
    # Key concepts should be filenames from llm_summary_details['key_documents']
    raw_concepts = llm_summary_details.get("key_documents", ontology.key_concepts) # Use key_documents from LLM output
    ontology.key_concepts = [str(c) for c in raw_concepts if isinstance(c, str)] if isinstance(raw_concepts, list) else ontology.key_concepts

    ontology.project_graph_nodes = final_validated_nodes_list # Assign final list of Pydantic objects
    # Update list of documents based on validated nodes
    final_doc_filenames_in_graph = {node.name for node in final_validated_nodes_list if node.node_type == "document"}
    updated_doc_list = sorted(list(existing_files.union(final_doc_filenames_in_graph)))
    ontology.documents = updated_doc_list
    ontology.document_count = len(ontology.documents)

    log.info(f"Ontology update complete. State: {ontology.document_count} docs, {len(ontology.project_graph_nodes)} nodes.")

    # --- 7. Persist Ontology & Return ---
    try:
        save_ontology_json(ontology, path) # Uses external helper
        log.info(f"Successfully saved updated ontology to '{path}'.")
    except Exception as e: log.error(f"Failed to save ontology to '{path}': {e}", exc_info=True)

    return ontology # Return the final object (or None if aborted)

def _merge_project_links_into_node_dicts(
    link_data_list: List[Dict],
    node_map: Dict[str, Dict], # Expects UUID -> node dict (FLAT structure)
    max_total_links: int = 15
) -> Tuple[int, int]:
    """
    Merges INTER-DOCUMENT link data into the 'linked_nodes' lists of node
    dictionaries stored in node_map. Assumes link_data_list contains RESOLVED UUIDs.
    Accesses 'source' directly on the node dict.
    """
    links_added_count = 0
    links_skipped_count = 0
    if not isinstance(link_data_list, list):
        logging.warning("Link merging received non-list input for link data.")
        return 0, 0

    for link_data in link_data_list: # Assume link_data now has UUIDs
        if not isinstance(link_data, dict): continue

        source_uuid = link_data.get("source_node_id") # Expecting UUID
        target_uuid = link_data.get("target_node_id") # Expecting UUID
        desc = link_data.get("relationship_description")

        if not all([source_uuid, target_uuid, desc]):
            links_skipped_count += 1; continue

        source_node_dict = node_map.get(source_uuid)
        target_node_dict = node_map.get(target_uuid)

        if not source_node_dict or not target_node_dict:
            logging.warning(f"Link Merge: Source UUID '{source_uuid}' or Target UUID '{target_uuid}' not found in node_map. Skipping.")
            links_skipped_count += 1; continue

        # *** Access source directly ***
        source_doc = source_node_dict.get('source', 'source_unknown')
        target_doc = target_node_dict.get('source', 'target_unknown')

        if source_doc == target_doc:
            links_skipped_count += 1; continue # Skip intra-doc links

        try:
            # Prepare the LinkedNode structure as a dictionary for adding to the list
            link_to_add = {
                "target_node_id": target_uuid,
                "relationship_description": desc,
                "relationship_keywords": link_data.get("relationship_keywords", []),
                "relationship_strength": float(link_data.get("relationship_strength", 5.0))
            }
            if not isinstance(link_to_add["relationship_keywords"], list):
                link_to_add["relationship_keywords"] = []
        except (ValueError, TypeError) as prep_err:
            logging.warning(f"Could not prepare link data {source_uuid} -> {target_uuid}: {prep_err}")
            links_skipped_count += 1; continue

        # Add Forward Link (to the dictionary in node_map)
        source_node_dict.setdefault('linked_nodes', [])
        if isinstance(source_node_dict['linked_nodes'], list):
             link_exists = False
             for existing_link in source_node_dict['linked_nodes']:
                 # Check if existing link is dict or Pydantic model
                 existing_target_id = existing_link.get('target_node_id') if isinstance(existing_link, dict) else getattr(existing_link, 'target_node_id', None)
                 if existing_target_id == target_uuid:
                     link_exists = True; break

             if len(source_node_dict['linked_nodes']) < max_total_links and not link_exists:
                 source_node_dict['linked_nodes'].append(link_to_add) # Add the dictionary
                 links_added_count += 1
             else: links_skipped_count += 1
        else: links_skipped_count += 1 # linked_nodes wasn't a list

        # Optional: Add Reverse Link (to the dictionary in node_map)
        target_node_dict.setdefault('linked_nodes', [])
        if isinstance(target_node_dict['linked_nodes'], list):
            reverse_link_exists = False
            for existing_link in target_node_dict['linked_nodes']:
                 existing_target_id = existing_link.get('target_node_id') if isinstance(existing_link, dict) else getattr(existing_link, 'target_node_id', None)
                 if existing_target_id == source_uuid:
                     reverse_link_exists = True; break

            if len(target_node_dict['linked_nodes']) < max_total_links and not reverse_link_exists:
                reverse_desc = f"Referenced by '{source_node_dict.get('name', source_uuid)}' from document '{source_doc}' ({desc})"
                reverse_link_to_add = {
                    "target_node_id": source_uuid,
                    "relationship_description": reverse_desc,
                    "relationship_keywords": link_to_add["relationship_keywords"],
                    "relationship_strength": link_to_add["relationship_strength"]
                }
                target_node_dict['linked_nodes'].append(reverse_link_to_add) # Add the dictionary
                # Optionally count reverse links separately
                # links_added_count += 1
            # else: links_skipped_count += 1
        # else: links_skipped_count += 1

    return links_added_count, links_skipped_count

def load_ontology_json(path: pathlib.Path) -> Optional[ProjectOntology]:
    """Loads ontology from JSON file, validates with Pydantic."""
    if not path.exists():
        return None
    try:
        with path.open("r", encoding="utf-8") as f:
            data = json.load(f)
        return ProjectOntology.model_validate(data)
    except (json.JSONDecodeError, ValidationError, IOError) as e:
        logging.error(f"Error loading or validating ontology from {path}: {e}", exc_info=True)
        return None

def save_ontology_json(ont: ProjectOntology, path: pathlib.Path):
    """Saves Pydantic ontology model to JSON file."""
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("w", encoding="utf-8") as f:
            json.dump(ont.model_dump(mode="json"), f, ensure_ascii=False, indent=2)
    except (IOError, TypeError) as e:
        logging.error(f"Error saving ontology to {path}: {e}", exc_info=True)

###############################
# DOCUMENT PROCESSING FUNCTIONS
###############################

# async def process_single_page_with_semaphore(semaphore, client, page_info: dict, uploaded_files):
#     """
#     Process a single page using semaphore for rate limiting.
#     Returns a dictionary: {"info": page_info, "result": PageContent object}.
#     Handles potential errors during processing. Uses 'source_identifier'.
#     """
#     # Get source_identifier and page_num early for use in logic and error reporting
#     source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
#     page_num = page_info.get("page_num", 0)

#     async with semaphore:
#         result_page_obj = None
#         try:
#             # Call process_single_pdf_page, which returns a DICTIONARY
#             # Pass the whole page_info which now contains source_identifier
#             page_dict = await process_single_pdf_page(client, page_info, uploaded_files)

#             page_dict.setdefault("page_number", page_num)
#             logging.debug(f"Attempting Pydantic validation for page {page_num} dict: {list(page_dict.keys())}")
#             result_page_obj = PageContent(**page_dict)
#             # Use source_identifier in log
#             logging.info(f"Successfully validated PageContent for page {page_num} of {source_identifier}")

#         except pydantic.ValidationError as ve:
#             # Use source_identifier in log/error
#             logging.error(f"Pydantic validation error creating PageContent for page {page_num} of {source_identifier} from returned dictionary: {ve}")
#             error_page_obj = create_error_page(
#                 page_num=page_num,
#                 error_msg=f"Pydantic Validation Error in '{source_identifier}': {ve}", # Include identifier in msg
#                 validation_errors=ve.errors()
#             )
#             result_page_obj = error_page_obj # Already a PageContent object

#         except Exception as e:
#             # Use source_identifier in log/error
#             logging.error(f"Unexpected error in semaphore task for page {page_num} of {source_identifier}: {str(e)}")
#             error_page_obj = create_error_page(
#                 page_num=page_num,
#                 error_msg=f"Unexpected processing error in semaphore task for '{source_identifier}': {str(e)}" # Include identifier in msg
#             )
#             result_page_obj = error_page_obj

#         # Return the result packaged with the original info
#         logging.debug(f"Finished processing for page {page_num} of {source_identifier}")
#         return {"info": page_info, "result": result_page_obj}


# async def process_single_pdf_page(client, page_info: dict, uploaded_files) -> dict:
    # """
    # Process a single PDF page image using Gemini with fallback options, returning extracted data as a dictionary.
    # Uses 'source_identifier' from page_info.
    # """
    # page_num = page_info.get("page_num", 0)
    # source_identifier = page_info.get("source_identifier", "Unknown_Source") # <-- Use new key
    # page_dict = {}

    # try:
    #     if "image_b64" not in page_info or not page_info["image_b64"]:
    #          st.warning(f"Missing image data for page {page_num} of {source_identifier}.")
    #          error_page_obj = create_error_page(page_num=page_num, error_msg="Missing or empty image data (image_b64)")
    #          return error_page_obj.model_dump()

    #     if "page_num" not in page_info:
    #          st.warning(f"Missing page number for a page in {source_identifier}.")
    #          error_page_obj = create_error_page(page_num=0, error_msg="Missing page number information")
    #          return error_page_obj.model_dump()

    #     img_data = base64.b64decode(page_info["image_b64"])
    #     image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")

    #     # Call Gemini with fallback strategy, passing source_identifier
    #     # Assumes process_page_with_fallback accepts source_identifier now
    #     page_content_obj = await process_page_with_fallback(
    #         client, image_part, page_num, source_identifier # <-- Pass identifier
    #     )

    #     page_dict = page_content_obj.model_dump()
    #     return page_dict

    # except Exception as e:
    #     st.error(f"Error during Gemini processing or data handling for page {page_num} of {source_identifier}: {str(e)}")
    #     error_page_obj = create_error_page(
    #         page_num=page_num,
    #         error_msg=f"Gemini API or processing error in '{source_identifier}': {str(e)}",
    #     )
    #     page_dict = error_page_obj.model_dump()
    #     return page_dict
    
# --- New/Renamed Function for PDF Page Image Processing (Stage 1) ---
async def process_single_pdf_page_image(semaphore: asyncio.Semaphore, client: genai.Client, page_info: dict) -> PageContent:
    """
    Processes a single PDF page image using Gemini (Stage 1 extraction only).
    Uses semaphore for rate limiting. Returns a PageContent object.
    Handles potential errors during processing. Uses 'source_identifier'.
    """
    source_identifier = page_info.get("source_identifier", "Unknown_Source")
    page_num = page_info.get("page_num", 0)
    result_page_obj = None

    async with semaphore:
        try:
            if "image_b64" not in page_info or not page_info["image_b64"]:
                 logging.warning(f"Missing image data for page {page_num} of {source_identifier}.")
                 result_page_obj = create_error_page(page_num=page_num, error_msg="Missing or empty image data (image_b64)")
            else:
                img_data = base64.b64decode(page_info["image_b64"])
                image_part = types.Part.from_bytes(data=img_data, mime_type="image/jpeg")

                # Call Gemini (Stage 1 extraction - NO subsections yet)
                # Uses process_page_with_fallback internally for robustness
                # Assumes extract_page_content_from_memory performs Stage 1 extraction
                result_page_obj = await extract_page_content_from_memory(
                    client, image_part, page_num, source_identifier
                )

                # Basic check if extraction failed severely (e.g. returned error type)
                # This check depends on how create_error_page marks errors
                if isinstance(result_page_obj, PageContent) and \
                   hasattr(result_page_obj, 'subsections') and \
                   any(getattr(s,'title', '') == "Processing Error" for s in getattr(result_page_obj, 'subsections', [])):
                     logging.warning(f"Initial structure extraction seems to have failed for page {page_num} of '{source_identifier}'.")
                     # Keep the error page object

                elif not isinstance(result_page_obj, PageContent):
                     logging.error(f"Unexpected type {type(result_page_obj)} returned from page processing for page {page_num} of {source_identifier}")
                     result_page_obj = create_error_page(page_num=page_num, error_msg=f"Unexpected type from processing: {type(result_page_obj)}")

        except pydantic.ValidationError as ve:
            logging.error(f"Pydantic validation error creating PageContent for page {page_num} of {source_identifier}: {ve}")
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Pydantic Validation Error in '{source_identifier}': {ve}",
                validation_errors=ve.errors()
            )
        except Exception as e:
            logging.error(f"Unexpected error in page processing task for page {page_num} of {source_identifier}: {str(e)}", exc_info=True)
            result_page_obj = create_error_page(
                page_num=page_num,
                error_msg=f"Unexpected processing error in page task for '{source_identifier}': {str(e)}"
            )

    # Ensure we always return a PageContent object, even if it's an error one
    if result_page_obj is None:
        logging.error(f"Processing yielded None for page {page_num} of {source_identifier}")
        result_page_obj = create_error_page(page_num=page_num, error_msg=f"Processing yielded None for page {page_num} of {source_identifier}")

    # logging.debug(f"Finished initial processing for page {page_num} of {source_identifier}")
    return result_page_obj


async def finalize_document(client: genai.Client, source_identifier: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Merges cut-off subsections.
    2. Extracts chapters based on subsections.
    3. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    4. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    5. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{source_identifier}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": source_identifier, # Store the identifier here
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }

    # Attempt to dump incoming pages immediately for error reporting, handle failures
    try:
        initial_pages_dump = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in pages]
        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed to initially dump pages for '{source_identifier}' for error reporting: {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{source_identifier}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{source_identifier}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result

        # --- Step 1: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Merging subsections for '{source_identifier}'.")
        merged_pages = await merge_cutoff_subsections(pages) # Operates on PageContent objects
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{source_identifier}'. Using original pages.")
             merged_pages = pages

        # --- Step 2: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Extracting chapters for '{source_identifier}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{source_identifier}'. Cannot proceed with relationship/summary analysis.")
             # Return result with merged pages but no summary/chapter info
             final_result["raw_extracted_content"]["pages"] = [p.model_dump(mode='json') for p in merged_pages]
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = { # Basic placeholder summary
                  "title": f"Processing Incomplete: {source_identifier}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result

        # --- Step 3: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Analyzing relationships within {len(chapters)} chapters for '{source_identifier}'.")
        # This loop modifies chapter.qdrant_nodes in place
        analysis_tasks = [analyze_concept_entity_relationships(client, chapter) for chapter in chapters]
        await asyncio.gather(*analysis_tasks, return_exceptions=True) # Run in parallel, log errors if any task fails
        # Check results? For now, we assume modification happened or warnings were logged internally.

        # --- Step 4: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
        # This function modifies chapter.qdrant_nodes again (adds inter-chapter links) AND returns the summary dict
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)

        # --- Step 5: Assemble Final Result ---
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated before dumping
        summary_dict['chapters'] = [ch.model_dump(mode='json') for ch in chapters if isinstance(ch, Chapter)]

        # Add document-level relationships placeholder (could be populated later from qdrant_nodes)
        summary_dict['entity_relationships'] = []

        # Dump the final state of merged pages
        final_pages_as_dicts = [p.model_dump(mode='json') for p in merged_pages if isinstance(p, PageContent)]

        final_result = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result

    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{source_identifier}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {source_identifier}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": [], "entity_relationships": []
            }
        # Keep the initial pages dump in the error case
        return final_result


async def extract_pages_from_pdf_bytes(pdf_bytes, source_identifier: str) -> List[Dict]:
    """Extract pages from PDF bytes as image info dicts."""
    loop = asyncio.get_running_loop()
    pages_info = []
    try:
        pdf_stream = io.BytesIO(pdf_bytes)
        # Run synchronous PyMuPDF code in an executor thread
        def sync_extract():
            doc = pymupdf.open(stream=pdf_stream, filetype="pdf")
            extracted = []
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    # Reduce DPI slightly to potentially save memory/time if quality allows
                    pix = page.get_pixmap(dpi=150, alpha=False)
                    img_bytes = pix.tobytes("jpeg", jpg_quality=80) # Slightly lower quality
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    extracted.append({
                        "page_num": page_num + 1,
                        "image_b64": img_b64,
                        "source_identifier": source_identifier # Use the passed identifier
                    })
                except Exception as page_e:
                     logging.warning(f"Error extracting page {page_num+1} from {source_identifier}: {page_e}")
            doc.close()
            return extracted

        pages_info = await loop.run_in_executor(None, sync_extract)
        return pages_info

    except Exception as e:
        logging.error(f"Error opening or processing PDF {source_identifier}: {str(e)}")
        # Return empty list, caller should handle this
        return []


# Non-PDF Document Processing Functions

async def process_single_document_memory(client: genai.Client, file_data: dict) -> List[PageContent]:
    """
    Processes a single non-PDF document from memory (Stage 1 extraction only).
    Calls specific handlers (Word, PPTX, Tabular, Text) which perform the
    initial structure extraction.
    Returns a list of PageContent objects (one per page/chunk).
    """
    source_identifier = file_data.get("identifier", file_data.get("name", "Unknown_Source")) # Get name robustly
    file_type = file_data.get("type", "").lower()
    file_content = file_data.get("content")
    processed_pages: List[PageContent] = []

    logging.info(f"Starting initial page-level processing for non-PDF: '{source_identifier}' (type: {file_type})")

    if not file_content:
        logging.error(f"No content found for file '{source_identifier}'.")
        return [create_error_page(1, f"Source '{source_identifier}' has no content.")]
    if not file_type:
        logging.error(f"Cannot determine file type for '{source_identifier}'.")
        return [create_error_page(1, f"Unknown file type for '{source_identifier}'.")]

    try:
        # Dispatch based on file type
        if file_type in ["xlsx", "xls"]:
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "excel")
        elif file_type == "csv":
            processed_pages = await process_tabular_data(client, file_content, source_identifier, "csv")
        elif file_type == "docx":
            processed_pages = await process_word_document(client, file_content, source_identifier)
        elif file_type == "pptx":
            processed_pages = await process_pptx_document(client, file_content, source_identifier)
        # Ensure all expected text types are handled
        elif file_type in ["txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]:
            # Pass the dict containing name, content, type
            processed_pages = await process_text_document(client, file_data)
        else:
            logging.warning(f"Unsupported file type '{file_type}' for '{source_identifier}'.")
            processed_pages = [create_error_page(1, f"Unsupported file type: {file_type}")]

        if not processed_pages:
            logging.warning(f"Initial processing for '{source_identifier}' resulted in no pages.")
            # Return an error page to indicate failure for this doc
            processed_pages = [create_error_page(1, f"No content pages generated during initial processing for '{source_identifier}'.")]

        logging.info(f"Finished initial page-level processing for non-PDF: '{source_identifier}'. Generated {len(processed_pages)} PageContent objects.")
        return processed_pages

    except Exception as e:
        logging.error(f"Core initial processing failed for document '{source_identifier}': {str(e)}", exc_info=True)
        return [create_error_page(1, f"Core initial processing error for '{source_identifier}': {e}")]

# --- NEW: Workflow function for a single document ---
async def process_and_finalize_single_document(
    client: genai.Client,
    file_data: dict,
    page_api_semaphore: asyncio.Semaphore, # Pass semaphore for page-level API calls
    # Optional: Add another semaphore if finalize_document also needs rate limiting
    # finalize_semaphore: asyncio.Semaphore
) -> Dict:
    """
    Handles the complete processing pipeline for a single document asynchronously.
    1. Initial page/chunk extraction (Stage 1)
    2. Finalization (Stage 2: subsections, chapters, relationships, summary)
    Returns the final dictionary structure produced by finalize_document.
    """
    function_name = "process_and_finalize_single_document"
    source_identifier = file_data.get("name", "Unknown_Document")
    file_type = file_data.get("type", "").lower()
    logging.info(f"[{function_name}] Starting workflow for '{source_identifier}' (type: {file_type})")

    initial_pages: List[PageContent] = []
    final_result_dict: Optional[Dict] = None

    try:
        # --- Stage 1: Initial Page/Chunk Extraction ---
        stage1_start = time.time()
        if file_type == 'pdf':
            pdf_content = file_data.get("content")
            if not pdf_content:
                 logging.error(f"[{function_name}] Missing PDF content for '{source_identifier}'.")
                 raise ValueError("Missing PDF content")

            # 1a. Extract page images
            logging.debug(f"[{function_name}] Extracting PDF page images for '{source_identifier}'")
            page_infos = await extract_pages_from_pdf_bytes(pdf_content, source_identifier)

            if not page_infos:
                 logging.warning(f"[{function_name}] No pages extracted from PDF '{source_identifier}'.")
                 # Create a single error page for this document
                 initial_pages = [create_error_page(1, f"Failed to extract any pages from PDF: {source_identifier}")]
            else:
                # 1b. Process page images concurrently (Stage 1 LLM)
                logging.debug(f"[{function_name}] Processing {len(page_infos)} PDF pages for '{source_identifier}' (Stage 1)")
                page_tasks = []
                valid_page_info_indices = [] # Keep track of indices corresponding to tasks
                for idx, pi in enumerate(page_infos):
                    if isinstance(pi, dict) and pi.get("image_b64"):
                        page_tasks.append(process_single_pdf_page_image(page_api_semaphore, client, pi))
                        valid_page_info_indices.append(idx)
                    else:
                        logging.warning(f"[{function_name}] Skipping invalid page info at index {idx} for PDF '{source_identifier}'")


                if page_tasks:
                     # Gather results - these are PageContent objects
                     gathered_results = await asyncio.gather(*page_tasks, return_exceptions=True)

                     # Filter out potential exceptions, replace with error pages
                     processed_pages = []
                     for i, res in enumerate(gathered_results):
                         # Get original page info index for context
                         original_info_index = valid_page_info_indices[i]
                         page_num_ctx = page_infos[original_info_index].get('page_num', original_info_index + 1)

                         if isinstance(res, PageContent):
                             processed_pages.append(res)
                         elif isinstance(res, Exception):
                             logging.error(f"[{function_name}] Error processing PDF page {page_num_ctx} of '{source_identifier}': {res}")
                             processed_pages.append(create_error_page(page_num_ctx, f"PDF page processing error: {res}"))
                         else:
                             logging.error(f"[{function_name}] Unexpected result type {type(res)} processing PDF page {page_num_ctx} of '{source_identifier}'")
                             processed_pages.append(create_error_page(page_num_ctx, f"Unexpected PDF page result type: {type(res)}"))
                     initial_pages = processed_pages
                     # Sort pages by page number
                     initial_pages.sort(key=lambda p: getattr(p, 'page_number', float('inf')))
                else:
                     logging.warning(f"[{function_name}] No valid page processing tasks created for PDF '{source_identifier}'.")
                     initial_pages = [create_error_page(1, f"No valid pages found or processed in PDF: {source_identifier}")]

        else: # Non-PDF types
            logging.debug(f"[{function_name}] Processing non-PDF '{source_identifier}' (Stage 1)")
            # This function handles different non-PDF types internally and calls Stage 1 LLM
            initial_pages = await process_single_document_memory(client, file_data)

        stage1_duration = time.time() - stage1_start
        logging.info(f"[{function_name}] Stage 1 (Initial Extraction) for '{source_identifier}' completed in {stage1_duration:.2f}s. Pages generated: {len(initial_pages)}")

        # Check if Stage 1 produced *any* valid PageContent objects (even if they contain errors internally)
        if not isinstance(initial_pages, list) or not initial_pages:
            logging.error(f"[{function_name}] Stage 1 failed catastrophically for '{source_identifier}'. No PageContent objects returned. Skipping finalization.")
            final_result_dict = {
                "raw_extracted_content": { "filename": source_identifier, "pages": [], "summary": None, "error": "Initial page extraction (Stage 1) failed completely." }
            }
            return final_result_dict

        # Further check: Did any page actually get usable text? (More robust check)
        has_content = any(isinstance(p, PageContent) and (p.raw_text or p.subsections) for p in initial_pages)
        if not has_content:
            logging.error(f"[{function_name}] Stage 1 produced PageContent objects but none contained usable text/subsections for '{source_identifier}'. Skipping finalization.")
            # Dump the potentially error-filled pages
            dumped_pages = []
            for p in initial_pages:
                 try: dumped_pages.append(p.model_dump(mode='json'))
                 except Exception: dumped_pages.append({"error": "Failed to dump page object"})

            final_result_dict = {
                "raw_extracted_content": { "filename": source_identifier, "pages": dumped_pages, "summary": None, "error": "Initial page extraction (Stage 1) failed to produce usable content." }
            }
            return final_result_dict # Return early


        # --- Stage 2: Finalization ---
        stage2_start = time.time()
        logging.debug(f"[{function_name}] Starting Stage 2 (Finalization) for '{source_identifier}'")

        # Use the existing finalize_document function
        # Optional: Add semaphore here if needed: async with finalize_semaphore: ...
        final_result_dict = await finalize_document(client, source_identifier, initial_pages)

        stage2_duration = time.time() - stage2_start
        logging.info(f"[{function_name}] Stage 2 (Finalization) for '{source_identifier}' completed in {stage2_duration:.2f}s.")

        # Log if finalization itself reported an error
        if isinstance(final_result_dict, dict) and final_result_dict.get("raw_extracted_content", {}).get("error"):
            logging.warning(f"[{function_name}] Finalization step reported an error for '{source_identifier}': {final_result_dict['raw_extracted_content']['error']}")

        return final_result_dict

    except Exception as e:
        logging.error(f"[{function_name}] Unhandled error during workflow for '{source_identifier}': {e}", exc_info=True)
        # Return an error dictionary
        # Attempt to dump initial pages if available
        dumped_initial_pages = []
        if initial_pages:
             for p in initial_pages:
                  try: dumped_initial_pages.append(p.model_dump(mode='json'))
                  except Exception: dumped_initial_pages.append({"error": "Failed to dump initial page object"})

        return {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": dumped_initial_pages,
                "summary": None,
                "error": f"Workflow failed unexpectedly: {str(e)}"
            }
        }

    
async def finalize_document(client: genai.Client, source_identifier: str, pages: List[PageContent]) -> Dict:
    """
    Finalizes document processing after page-level extraction following the new flow.
    1. Extracts subsections from raw_text for each page. <--- NEW STEP
    2. Merges cut-off subsections.
    3. Extracts chapters based on subsections.
    4. Analyzes nodes/relationships WITHIN each chapter (populates chapter.qdrant_nodes).
    5. Analyzes relationships BETWEEN chapters (updates chapter.qdrant_nodes) AND generates document summary.
    6. Assembles and returns the final dictionary structure.
    """
    function_name = "finalize_document"
    logging.info(f"[{function_name}] Finalizing document '{source_identifier}' with {len(pages)} initial pages.")

    # --- Default Error Structure ---
    final_result = {
        "raw_extracted_content": {
            "filename": source_identifier, # Store the identifier here
            "pages": [], # Will be populated later or with error pages
            "summary": None,
            "error": None
        }
    }
    # Attempt to dump incoming pages immediately for error reporting, handle failures
    initial_pages_dump = []
    try:
        # Dump safely, handle potential model errors during dump
        for i, p in enumerate(pages):
             if isinstance(p, PageContent):
                 try: initial_pages_dump.append(p.model_dump(mode='json'))
                 except Exception as dump_err: initial_pages_dump.append({"page_number": getattr(p, 'page_number', i+1), "error": f"Initial model dump failed: {dump_err}"})
             else: initial_pages_dump.append(p) # Keep as is if not PageContent obj

        final_result["raw_extracted_content"]["pages"] = initial_pages_dump
    except Exception as initial_dump_err:
         logging.error(f"[{function_name}] Failed during initial page dump for '{source_identifier}': {initial_dump_err}")
         final_result["raw_extracted_content"]["pages"] = [{"page_number": i+1, "error": "Original page data could not be serialized"} for i in range(len(pages))]


    try:
        # --- Input Validation ---
        if not isinstance(pages, list) or not all(isinstance(p, PageContent) for p in pages):
            error_msg = "Invalid input: 'pages' must be a list of PageContent objects."
            logging.error(f"[{function_name}] {error_msg} for '{source_identifier}'.")
            final_result["raw_extracted_content"]["error"] = error_msg
            return final_result # Return early with initial page dump

        if not pages:
             logging.warning(f"[{function_name}] No pages provided for '{source_identifier}'. Finalization cannot proceed.")
             final_result["raw_extracted_content"]["error"] = "No pages provided for finalization."
             return final_result # Return early

        # --- Step 1: Extract Subsections from Raw Text (NEW STEP) ---
        logging.debug(f"[{function_name}] Step 1: Extracting subsections from raw text for {len(pages)} pages of '{source_identifier}'.")
        subsection_extraction_tasks = []
        pages_needing_subsections = [] # Keep track of pages we are processing

        # Create tasks FIRST to run them concurrently
        for page in pages:
             # Only process if raw_text exists and subsections are currently empty
             if isinstance(page, PageContent) and page.raw_text and not page.subsections:
                 task = extract_subsections_from_text(client, page.raw_text, page.page_number, source_identifier)
                 subsection_extraction_tasks.append(task)
                 pages_needing_subsections.append(page) # Store reference to the page object
             elif not isinstance(page, PageContent):
                  logging.warning(f"[{function_name}] Item in pages list is not PageContent object for '{source_identifier}', page {getattr(page, 'page_number', '?')}. Skipping subsection extraction for it.")

             # else: subsections already exist or no raw text, skip

        # Run tasks concurrently
        if subsection_extraction_tasks:
             logging.info(f"[{function_name}] Running {len(subsection_extraction_tasks)} subsection extraction tasks for '{source_identifier}'.")
             results = await asyncio.gather(*subsection_extraction_tasks, return_exceptions=True)
             logging.info(f"[{function_name}] Finished subsection extraction tasks for '{source_identifier}'.")

             # Assign results back to the corresponding page objects
             if len(results) == len(pages_needing_subsections):
                 for i, result in enumerate(results):
                     page_to_update = pages_needing_subsections[i] # Get the correct page object
                     if isinstance(result, Exception):
                         logging.error(f"[{function_name}] Subsection extraction failed for page {page_to_update.page_number} of '{source_identifier}': {result}")
                         page_to_update.subsections = [] # Ensure empty list on error
                     elif isinstance(result, list): # Expected result is List[Subsection]
                         page_to_update.subsections = result # Assign generated subsections
                     else:
                         logging.error(f"[{function_name}] Unexpected result type {type(result)} for subsection extraction on page {page_to_update.page_number} of '{source_identifier}'.")
                         page_to_update.subsections = []
             else:
                  # This indicates a logic error in tracking tasks/pages
                  logging.error(f"[{function_name}] Mismatch between pages needing processing ({len(pages_needing_subsections)}) and task results ({len(results)}) for '{source_identifier}'. Subsections may be missing.")
                  # Attempt to assign based on index anyway, but it might be wrong
                  for idx, page in enumerate(pages_needing_subsections):
                     if idx < len(results):
                         result = results[idx]
                         if isinstance(result, list): page.subsections = result
                         else: page.subsections = []
                     else: page.subsections = [] # No result for this page

        else:
             logging.info(f"[{function_name}] No pages required subsection extraction for '{source_identifier}'.")


        # --- Step 2: Merge cut-off subsections ---
        logging.debug(f"[{function_name}] Step 2: Merging subsections for '{source_identifier}'.")
        # Pass the 'pages' list which now has potentially populated subsections
        merged_pages = await merge_cutoff_subsections(pages)
        if not merged_pages:
             logging.warning(f"[{function_name}] Subsection merging resulted in empty page list for '{source_identifier}'. Using original pages.")
             merged_pages = pages # Fallback to original pages (already updated with subsections)


        # --- Step 3: Extract chapters from subsections ---
        logging.debug(f"[{function_name}] Step 3: Extracting chapters for '{source_identifier}'.")
        chapters: List[Chapter] = await extract_chapters_from_subsections(client, merged_pages)
        if not chapters:
             logging.warning(f"[{function_name}] No chapters were extracted for '{source_identifier}'. Cannot proceed further.")
             final_pages_as_dicts = [p.model_dump(mode='json') if isinstance(p, PageContent) else p for p in merged_pages] # Dump state after merge/subsection attempts
             final_result["raw_extracted_content"]["pages"] = final_pages_as_dicts
             final_result["raw_extracted_content"]["error"] = "Chapter extraction failed or yielded no chapters."
             final_result["raw_extracted_content"]["summary"] = {
                  "title": f"Processing Incomplete: {source_identifier}", "themes": [], "questions": [], "summary": "Failed to structure document into chapters.", "chapters": []
             }
             return final_result


        # --- Step 4: Analyze nodes/relationships WITHIN each chapter ---
        logging.debug(f"[{function_name}] Step 4: Analyzing intra-chapter relationships for {len(chapters)} chapters in '{source_identifier}'.")
        # This loop modifies chapter.qdrant_nodes in place
        intra_chapter_tasks = [
            analyze_concept_entity_relationships(client, chapter, source_identifier) # <-- ADD source_identifier HERE
            for chapter in chapters
        ]
        intra_results = await asyncio.gather(*intra_chapter_tasks, return_exceptions=True)

        # Log any exceptions from intra-chapter analysis
        for i, res in enumerate(intra_results):
            if isinstance(res, Exception):
                 logging.error(f"[{function_name}] Intra-chapter analysis failed for chapter {chapters[i].chapter_id} of '{source_identifier}': {res}")


        # --- Step 5: Analyze relationships BETWEEN chapters & Generate Summary ---
        logging.debug(f"[{function_name}] Step 5: Analyzing inter-chapter relationships and generating summary for '{source_identifier}'.")
        summary_dict = await analyze_inter_chapter_relationships(client, chapters, source_identifier)


        # --- Step 6: Assemble Final Result ---
        logging.debug(f"[{function_name}] Step 6: Assembling final result for '{source_identifier}'.")
        if summary_dict is None:
             logging.warning(f"[{function_name}] Failed to generate document summary details for '{source_identifier}'. Using placeholder.")
             summary_dict = {
                 "title": f"Summary Failed: {source_identifier}", "themes": ["error"], "questions": [],
                 "summary": "Failed to generate document summary via LLM.",
                 # Chapters will be added below anyway
             }

        # Add chapter data (dumped) to the summary dictionary
        # Ensure chapters have qdrant_nodes populated (even if empty list) before dumping
        dumped_chapters = []
        for ch in chapters:
            if isinstance(ch, Chapter):
                try: dumped_chapters.append(ch.model_dump(mode='json'))
                except Exception as ch_dump_err:
                     logging.error(f"[{function_name}] Failed to dump chapter {ch.chapter_id}: {ch_dump_err}")
                     dumped_chapters.append({"chapter_id": ch.chapter_id, "title": ch.title, "error": "Failed to serialize chapter"})
            else: logging.warning(f"[{function_name}] Invalid chapter object found during final assembly.")
        summary_dict['chapters'] = dumped_chapters


        # Dump the final state of pages (use merged_pages)
        final_pages_as_dicts = []
        for p in merged_pages:
             if isinstance(p, PageContent):
                 try: final_pages_as_dicts.append(p.model_dump(mode='json'))
                 except Exception as pg_dump_err:
                      logging.error(f"[{function_name}] Failed to dump page {p.page_number}: {pg_dump_err}")
                      final_pages_as_dicts.append({"page_number": p.page_number, "error": "Failed to serialize page"})
             else: final_pages_as_dicts.append(p)


        final_result = {
            "raw_extracted_content": {
                "filename": source_identifier,
                "pages": final_pages_as_dicts,
                "summary": summary_dict, # Assign the final summary dictionary
                "error": None # Clear error state if processing succeeded
            }
        }
        logging.info(f"[{function_name}] Successfully finalized document '{source_identifier}'.")
        return final_result

    # --- Error Handling ---
    except Exception as e:
        error_msg = f"Unexpected error during finalize_document for '{source_identifier}': {str(e)}"
        logging.error(error_msg, exc_info=True)
        final_result["raw_extracted_content"]["error"] = error_msg
        # Ensure summary is None or a basic error dict on critical failure
        if final_result["raw_extracted_content"]["summary"] is None:
            final_result["raw_extracted_content"]["summary"] = {
                "title": f"Error Finalizing {source_identifier}", "themes": ["error"], "questions": [],
                "summary": f"Critical error during final processing: {e}", "chapters": []
                # Removed entity_relationships placeholder here too
            }
        # Keep the initial pages dump in the error case if available
        return final_result



#-------------------------------------
# HELPER: Create Markdown Table
#-------------------------------------
def create_markdown_table(headers, data):
    """Create a markdown table string from headers and list of lists/dicts."""
    if not data or not headers: return "<!-- Empty Table -->"
    header_row = "| " + " | ".join(map(str, headers)) + " |"
    separator_row = "| " + " | ".join("---" * len(headers)) + " |"
    data_rows = []
    for row in data:
        if isinstance(row, dict):
            values = [str(row.get(h, '')) for h in headers]
        elif isinstance(row, (list, tuple)):
            # Ensure row has same length as headers, padding if necessary
            values = [str(row[i]) if i < len(row) else '' for i in range(len(headers))]
        else: # Handle unexpected row types
             values = [str(row)] + [''] * (len(headers) - 1)
        data_rows.append("| " + " | ".join(values) + " |")
    return "\n".join([header_row, separator_row] + data_rows)

#-------------------------------------
# REWRITTEN NON-PDF PROCESSORS
#-------------------------------------

async def process_tabular_data(client: genai.Client, file_content: bytes, source_identifier: str, filetype: str) -> List[PageContent]:
    """Processes Excel/CSV by converting sheets/file to Markdown text and extracting structure (Stage 1)."""
    processed_pages = []
    page_counter = 1 # For sheets or single CSV page

    try:
        if filetype == "excel":
            try:
                 excel_data = pd.read_excel(io.BytesIO(file_content), sheet_name=None, engine='openpyxl')
                 if not excel_data:
                     raise ValueError("Excel file is empty or unreadable.")

                 for sheet_name, df in excel_data.items():
                     logging.info(f"Processing Excel sheet '{sheet_name}' for '{source_identifier}' as page {page_counter}")
                     if df.empty:
                          logging.warning(f"Sheet '{sheet_name}' in '{source_identifier}' is empty. Skipping.")
                          page_counter += 1
                          continue
                     markdown_text = f"# Sheet: {sheet_name}\n\n" + df.to_markdown(index=False)
                     # --- Calls Stage 1 Extraction ---
                     page_content = await extract_structure_from_text_content(client, markdown_text, page_counter, f"{source_identifier} (Sheet: {sheet_name})")
                     processed_pages.append(page_content)
                     page_counter += 1

            except Exception as e:
                logging.error(f"Error reading Excel file '{source_identifier}' with pandas: {e}", exc_info=True)
                processed_pages.append(create_error_page(1, f"Error reading Excel file '{source_identifier}': {e}"))

        elif filetype == "csv":
            logging.info(f"Processing CSV file '{source_identifier}' as page 1")
            try:
                df = pd.read_csv(io.BytesIO(file_content))
                if df.empty:
                     raise ValueError("CSV file is empty or unreadable.")
                markdown_text = "# CSV Data\n\n" + df.to_markdown(index=False)
                # --- Calls Stage 1 Extraction ---
                page_content = await extract_structure_from_text_content(client, markdown_text, 1, source_identifier)
                processed_pages.append(page_content)
            except Exception as e:
                 logging.error(f"Error reading CSV file '{source_identifier}' with pandas: {e}", exc_info=True)
                 processed_pages.append(create_error_page(1, f"Error reading CSV file '{source_identifier}': {e}"))
        else:
            raise ValueError(f"Unsupported tabular file type: {filetype}")

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing tabular data for '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process tabular file '{source_identifier}': {e}")]

async def process_word_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes Word docx file by extracting text and applying structure extraction (Stage 1)."""
    processed_pages = []
    try:
        logging.info(f"Processing Word document '{source_identifier}'")
        doc = docx.Document(io.BytesIO(file_content))
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        if not full_text.strip():
             logging.warning(f"Word document '{source_identifier}' contains no text.")
             return [create_error_page(1, "Word document is empty or contains no text.")]

        logging.debug(f"Extracted {len(full_text)} chars from '{source_identifier}'. Processing as single page.")
        # --- Calls Stage 1 Extraction ---
        page_content = await extract_structure_from_text_content(client, full_text, 1, source_identifier)
        processed_pages.append(page_content)

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing Word document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process Word document '{source_identifier}': {e}")]

async def process_pptx_document(client: genai.Client, file_content: bytes, source_identifier: str) -> List[PageContent]:
    """Processes PowerPoint pptx file by extracting text per slide and applying structure extraction (Stage 1)."""
    processed_pages = []
    try:
        logging.info(f"Processing PowerPoint document '{source_identifier}'")
        ppt = Presentation(io.BytesIO(file_content))

        if not ppt.slides:
             logging.warning(f"PowerPoint document '{source_identifier}' contains no slides.")
             return [create_error_page(1, "PowerPoint document has no slides.")]

        for i, slide in enumerate(ppt.slides, 1):
            slide_texts = []
            title = f"Slide {i}"
            has_title_shape = False
            try:
                for shape in slide.shapes:
                     if hasattr(shape, "is_title") and shape.is_title and shape.has_text_frame and shape.text.strip():
                         title = shape.text.strip()
                         slide_texts.append(f"# {title}")
                         has_title_shape = True
                         break
            except Exception as shape_error:
                 logging.warning(f"Error accessing title shape on slide {i} of '{source_identifier}': {shape_error}")

            try:
                for shape in slide.shapes:
                     if shape.has_text_frame and shape.text and not (has_title_shape and hasattr(shape, "is_title") and shape.is_title):
                         for para in shape.text_frame.paragraphs:
                              if para.text.strip():
                                   prefix = "- " * (para.level + 1)
                                   slide_texts.append(f"{prefix}{para.text.strip()}")
            except Exception as shape_error:
                 logging.warning(f"Error accessing text shapes on slide {i} of '{source_identifier}': {shape_error}")

            try:
                 if slide.has_notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip():
                     slide_texts.append("\n## Speaker Notes")
                     slide_texts.append(slide.notes_slide.notes_text_frame.text.strip())
            except Exception as notes_error:
                logging.warning(f"Error accessing notes on slide {i} of '{source_identifier}': {notes_error}")

            slide_full_text = "\n".join(slide_texts)

            if not slide_full_text.strip():
                 logging.warning(f"Slide {i} of '{source_identifier}' contains no text content. Skipping.")
                 continue

            logging.debug(f"Processing slide {i} of '{source_identifier}'")
            # --- Calls Stage 1 Extraction ---
            page_content = await extract_structure_from_text_content(client, slide_full_text, i, f"{source_identifier} (Slide {i})")
            processed_pages.append(page_content)

        # Ensure at least one page object is returned if slides existed but were all empty/failed
        if ppt.slides and not processed_pages:
            logging.warning(f"PPTX '{source_identifier}' had slides but none yielded processable text.")
            return [create_error_page(1, "PPTX had slides but no text could be processed.")]

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing PowerPoint document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process PowerPoint document '{source_identifier}': {e}")]

async def process_text_document(client: genai.Client, file_data: dict) -> List[PageContent]:
    """Processes plain text, markdown, json, html, xml files by chunking and extracting structure (Stage 1)."""
    source_identifier = file_data.get("name", "Unknown_Text_File")
    file_type = file_data.get("type", "txt")
    processed_pages = []
    page_counter = 1

    try:
        logging.info(f"Processing {file_type} document '{source_identifier}'")
        try:
            text_content = file_data["content"].decode('utf-8')
        except UnicodeDecodeError:
             try:
                 text_content = file_data["content"].decode('latin-1')
             except Exception as decode_err:
                 logging.error(f"Failed to decode file '{source_identifier}': {decode_err}")
                 return [create_error_page(1, f"Failed to decode file '{source_identifier}'")]

        if not text_content.strip():
            logging.warning(f"Text document '{source_identifier}' is empty.")
            return [create_error_page(1, f"Text document '{source_identifier}' is empty.")]

        # Simple chunking by character count, aiming for rough page equivalents
        MAX_CHUNK_CHARS = 10000 # Adjust based on typical LLM context window and desired granularity
        text_chunks = [text_content[i:i + MAX_CHUNK_CHARS] for i in range(0, len(text_content), MAX_CHUNK_CHARS)]

        if not text_chunks:
             text_chunks = [text_content] # Ensure at least one chunk if splitting failed

        logging.info(f"Split '{source_identifier}' into {len(text_chunks)} chunks for processing.")

        for text_chunk in text_chunks:
            if not text_chunk.strip(): continue
            # --- Calls Stage 1 Extraction ---
            page_content = await extract_structure_from_text_content(client, text_chunk, page_counter, f"{source_identifier} (Chunk {page_counter})")
            processed_pages.append(page_content)
            page_counter += 1

        return processed_pages

    except Exception as e:
        logging.error(f"Error processing text document '{source_identifier}': {e}", exc_info=True)
        return [create_error_page(1, f"Failed to process text document '{source_identifier}': {e}")]


async def process_all_documents_async(file_data: List[Dict]) -> List[Dict]:
    """
    Process all documents concurrently, including asynchronous finalization per document.
    """
    # --- Setup ---
    try:
        api_key = get_gemini_api_key()
        # Ensure Helicone proxy details are correct or remove if not used
        client = genai.Client(api_key=api_key,
            http_options={
                "base_url": st.secrets.get("HELICONE_PROXY_URL", 'https://generativelanguage.googleapis.com'),
                "headers": {
                    "helicone-auth": f'Bearer {st.secrets.get("HELICONE_API_KEY", "")}',
                    "helicone-target-url": 'https://generativelanguage.googleapis.com'
                } if st.secrets.get("HELICONE_PROXY_URL") and st.secrets.get("HELICONE_API_KEY") else {} # Conditionally add headers
            })
        logging.info("Gemini client initialized for document processing.")
    except Exception as client_err:
        logging.error(f"Failed to initialize Gemini client: {client_err}")
        st.error(f"Failed to initialize Gemini client: {client_err}")
        return []

    final_doc_results = [] # Stores the FINAL dictionary results from each document's workflow
    original_doc_count = len(file_data)
    docs_to_process_count = 0 # Count valid files

    # --- Status Tracking Setup ---
    status_container = st.session_state.get("status_container")
    progress_bar = st.session_state.get("progress_bar")
    time_info = st.session_state.get("time_info")
    processing_status = {
        "active": True, "current_step": "Initializing...",
        "total_steps": original_doc_count, # Track documents as steps initially
        "current_step_num": 0, "step_progress": 0.0, # Not used per-step
        "start_time": time.time(), "step_start_time": time.time(),
        "estimated_time_remaining": None, "parallel_count": 1
    }
    # We no longer have distinct pipeline phases in the status update logic

    def update_doc_progress(completed_count, total_count, current_source_identifier=None):
        """Updates the processing status based on completed documents."""
        try:
            # Update total steps if some files were skipped early
            processing_status["total_steps"] = total_count
            processing_status["current_step_num"] = completed_count
            overall_progress = 0.0
            if total_count > 0:
                overall_progress = min(1.0, float(completed_count) / float(total_count))
            overall_progress = max(0.0, min(1.0, overall_progress))

            status_msg = f"Processing Documents ({completed_count}/{total_count})"
            # Display last completed document name might be more reliable than current
            if current_source_identifier:
                status_msg = f"Completed: {current_source_identifier} ({completed_count}/{total_count})"


            eta_msg = ""
            time_elapsed = time.time() - processing_status["start_time"]
            if overall_progress > 0.01:
                est_total_time = time_elapsed / overall_progress
                remaining = max(0, est_total_time - time_elapsed)
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                mins_remaining, secs_remaining = divmod(int(remaining), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: ~{mins_remaining}m {secs_remaining}s"
            else:
                mins_elapsed, secs_elapsed = divmod(int(time_elapsed), 60)
                eta_msg = f"Elapsed: {mins_elapsed}m {secs_elapsed}s | ETA: Calculating..."

            if status_container: status_container.update(label=status_msg, expanded=True)
            if progress_bar: progress_bar.progress(overall_progress, text=f"{int(overall_progress*100)}%")
            if time_info: time_info.markdown(f"`{eta_msg}`")

        except Exception as ui_err:
             logging.warning(f"Failed to update UI status: {ui_err}")

    # --- Main Concurrent Processing Logic ---
    try:
        update_doc_progress(0, original_doc_count) # Initial status

        # --- Concurrency Limits ---
        # Limit for page-level API calls (within each doc's Stage 1)
        MAX_PARALLEL_PAGES = 100
        page_api_semaphore = asyncio.Semaphore(MAX_PARALLEL_PAGES)

        # Limit for how many documents process concurrently overall
        MAX_PARALLEL_DOCS = 10
        doc_workflow_semaphore = asyncio.Semaphore(MAX_PARALLEL_DOCS)

        # --- Create Tasks for each Document Workflow ---
        doc_tasks = []
        valid_file_infos = []
        for file_info in file_data:
            if not isinstance(file_info, dict) or "name" not in file_info or "content" not in file_info or not file_info["content"]:
                 st.warning(f"Skipping invalid file data item (missing name, content, or empty): {str(file_info.get('name', 'N/A'))}...")
                 continue # Skip this file

            valid_file_infos.append(file_info)
            docs_to_process_count += 1 # Increment count of actual files to process

            async def run_workflow_with_semaphore(file_info_arg):
                async with doc_workflow_semaphore:
                    # This allows MAX_PARALLEL_DOCS workflows to run concurrently
                    # Each workflow uses page_api_semaphore internally for its page processing
                    logging.debug(f"Acquired doc semaphore for: {file_info_arg.get('name')}")
                    result = await process_and_finalize_single_document(
                        client,
                        file_info_arg,
                        page_api_semaphore
                        # Add finalize_semaphore if needed
                    )
                    logging.debug(f"Releasing doc semaphore for: {file_info_arg.get('name')}")
                    return result

            doc_tasks.append(run_workflow_with_semaphore(file_info))

        # Update total steps based on valid files
        processing_status["total_steps"] = docs_to_process_count
        update_doc_progress(0, docs_to_process_count) # Update progress bar total

        # --- Execute and Gather Results ---
        completed_docs = 0
        if not doc_tasks:
             logging.warning("No valid documents found to process.")
             st.warning("No valid documents found to process.")
             processing_status["active"] = False
             if status_container:
                 try: status_container.update(label="⏹️ No valid documents found.", state="error", expanded=False)
                 except Exception: pass
             return []


        logging.info(f"Launching {len(doc_tasks)} document processing workflows concurrently (max {MAX_PARALLEL_DOCS})...")
        processing_status["parallel_count"] = MAX_PARALLEL_DOCS

        # Use asyncio.as_completed to process results as they finish and update progress
        for future in asyncio.as_completed(doc_tasks):
            last_completed_name = "Unknown"
            try:
                result_dict = await future # result_dict is the final dict from finalize_document or an error dict
                if isinstance(result_dict, dict) and "raw_extracted_content" in result_dict:
                     final_doc_results.append(result_dict)
                     last_completed_name = result_dict.get("raw_extracted_content", {}).get("filename", "Unknown")
                     # Log success or failure for this doc
                     if result_dict.get("raw_extracted_content", {}).get("error"):
                          logging.warning(f"Workflow for '{last_completed_name}' completed with error: {result_dict['raw_extracted_content']['error']}")
                     else:
                          logging.info(f"Workflow for '{last_completed_name}' completed successfully.")
                else:
                     # This case indicates a severe failure within the workflow task itself
                     logging.error(f"Workflow task returned unexpected data type: {type(result_dict)}")
                     last_completed_name = "Workflow Error"
                     final_doc_results.append({
                         "raw_extracted_content": {"filename": "Unknown_Workflow_Error", "pages": [], "summary": None, "error": f"Workflow task failed unexpectedly, returned type {type(result_dict)}"}
                     })

                completed_docs += 1
                update_doc_progress(completed_docs, docs_to_process_count, last_completed_name)

            except asyncio.CancelledError:
                 logging.warning("A document processing task was cancelled.")
                 completed_docs += 1 # Count it as "done" for progress calculation
                 update_doc_progress(completed_docs, docs_to_process_count, "Cancelled Task")
                 # Optionally add a specific marker for cancelled tasks if needed
                 final_doc_results.append({
                     "raw_extracted_content": {"filename": "Cancelled_Task", "pages": [], "summary": None, "error": "Processing was cancelled."}
                 })
            except Exception as e:
                 logging.error(f"Error awaiting document task result: {e}", exc_info=True)
                 completed_docs += 1 # Count error task as "done"
                 last_completed_name = "Errored Task"
                 update_doc_progress(completed_docs, docs_to_process_count, last_completed_name)
                 # Add a generic error dict to results
                 final_doc_results.append({
                     "raw_extracted_content": {"filename": "Unknown_Workflow_Error", "pages": [], "summary": None, "error": f"Workflow task failed: {str(e)}"}
                 })

        # Final status update
        update_doc_progress(completed_docs, docs_to_process_count)
        processing_status["active"] = False
        if status_container:
             final_state = "complete" if completed_docs == docs_to_process_count else "error"
             final_label = "✅ Document processing complete!" if final_state == "complete" else "⚠️ Processing finished with errors."
             try: status_container.update(label=final_label, state=final_state, expanded=False)
             except Exception as ui_err: logging.warning(f"Failed to update final status: {ui_err}")

    # --- Exception Handling & Finally Block ---
    except asyncio.CancelledError:
         logging.warning("Document processing orchestration was cancelled.")
         st.warning("Document processing was cancelled.")
         if status_container:
             try: status_container.update(label="⏹️ Processing Cancelled", state="error", expanded=False)
             except Exception: pass
         processing_status["active"] = False
         return [] # Return empty list on cancellation
    except Exception as e:
        logging.error(f"❌ An error occurred during the main document processing orchestration: {str(e)}", exc_info=True)
        st.error(f"❌ An error occurred during document processing: {str(e)}")
        if status_container:
            try: status_container.update(label=f"❌ Pipeline Error: {str(e)}", state="error", expanded=True)
            except Exception: pass
        processing_status["active"] = False
        # Return whatever was collected before the error
        logging.info(f"Returning {len(final_doc_results)} potentially partially processed documents due to pipeline error.")
        return final_doc_results
    finally:
        logging.info("Document processing orchestration finished or exited.")
        processing_status["active"] = False
        if "st" in locals() and hasattr(st, "session_state"):
            st.session_state.processing_active = False # Ensure flag is reset

    # Return the list of final document dictionaries
    logging.info(f"Successfully processed {len(final_doc_results)} documents out of {docs_to_process_count} attempted.")
    return final_doc_results

###############################
# QDRANT INGESTION WORKFLOW (REPLACED with 2-Phase Foreground Method)
###############################

# --- Phase 1a: Create Collection (m=0) ---
# --- Phase 1a: Create Collection (m=0) - UPDATED ---
async def create_collection_phase1a(client: AsyncQdrantClient):
    """
    Phase 1a: Ensures the collection exists with HNSW disabled (m=0) and binary quantization.
    If the collection already exists, it logs a warning but proceeds, assuming the user
    intends to use the existing one or will manage its state manually if needed.
    Consider adding config checks or recreation logic if strict state reset is required.
    """
    if not client:
        logging.error("Qdrant client is None.")
        return False
    collection_name = QDRANT_COLLECTION_NAME # Use the constant

    try:
        # --- Check if collection exists ---
        collection_exists = await client.collection_exists(collection_name=collection_name)

        if collection_exists:
            logging.warning(f"Collection '{collection_name}' already exists. Phase 1a assumes it's usable or will be managed.")
            # Optional: Add checks here to see if the *existing* collection's config
            # matches the desired m=0 state. If not, you could try to update it
            # using client.update_collection, or raise an error/warning.
            # For simplicity now, we just proceed.
            st.sidebar.warning(f"Collection '{collection_name}' already exists. Proceeding.")
            return True # Treat as success for this phase if it exists

        # --- If it doesn't exist, create it ---
        logging.info(f"Phase 1a: Creating collection '{collection_name}' (m=0)...")
        await client.create_collection(
            collection_name=collection_name,
            vectors_config={
                EMB_MODEL_NAME: models.VectorParams(
                    size=EMB_DIM,
                    distance=models.Distance.COSINE,
                    on_disk=True,
                    hnsw_config=models.HnswConfigDiff(m=0) # HNSW disabled initially
                )
            },
            sparse_vectors_config={
                SPARSE_MODEL_NAME: models.SparseVectorParams(
                    index=models.SparseIndexParams(on_disk=True)
                )
            },
            quantization_config=models.BinaryQuantization(
                binary=models.BinaryQuantizationConfig(always_ram=True)
            ),
            optimizers_config=models.OptimizersConfigDiff(memmap_threshold=20000)
        )
        logging.info(f"Collection '{collection_name}' created (HNSW disabled).")
        await asyncio.sleep(1) # Give Qdrant a moment
        info = await client.get_collection(collection_name=collection_name)
        logging.info(f"Collection status after creation: {info.status}")
        return True

    except Exception as e:
        # Catch potential errors during check or creation
        logging.error(f"Failed during Phase 1a for collection '{collection_name}': {e}", exc_info=True)
        # Display specific Qdrant conflict error if available
        if isinstance(e, models.UnexpectedResponse) and e.status_code == 409:
             st.error(f"Phase 1a Error: Collection '{collection_name}' exists but encountered an issue. Check Qdrant logs.")
        else:
             st.error(f"Phase 1a Error: Failed to ensure collection '{collection_name}': {e}")
        return False

# --- Phase 1b - Step 3: Upload Single Batch (Adapted) ---
async def upload_batch_nodes(
    client: AsyncQdrantClient,
    batch_nodes: List[QdrantNode], # Takes a list of QdrantNode objects
    batch_num: int,
    total_batches: int,
    dense_model: TextEmbedding,
    sparse_model: SparseTextEmbedding
):
    """Processes and upserts a single batch of QdrantNode objects."""
    if not client: logging.error("Qdrant client is None."); return 0
    if not batch_nodes: logging.debug(f"Batch {batch_num}: Empty, skipping."); return 0

    logging.debug(f"Processing batch {batch_num}/{total_batches} ({len(batch_nodes)} nodes)")
    texts = []
    valid_nodes_in_batch = []

    # 1. Prepare texts for embedding
    for node in batch_nodes: # batch_nodes are QdrantNode Pydantic objects here
        if not isinstance(node, QdrantNode) or not node.node_id: continue

        node_name = getattr(node, 'name', '')
        # --- FIXED: Access direct attributes ---
        core_claim = getattr(node, 'core_claim', '')
        title = getattr(node, 'title', '')
        # executive_summary might exist for doc/meta nodes
        exec_summary = getattr(node, 'executive_summary', '')

        # Prioritize exec_summary > core_claim > title > name
        content_text = exec_summary or core_claim or title or "" # Ensure string

        if node_name and content_text:
            text = f"{node_name}: {content_text}".strip()
        elif content_text:
            text = content_text.strip() # Use content if name is missing
        elif node_name:
            text = node_name.strip() # Use name if content is missing
        else:
            logging.warning(f"Node {node.node_id}: No name or content fields found. Using placeholder.")
            text = f"Placeholder content for node {node.node_id}"

        # Final check for effectively empty string
        if not text or text == ":":
             logging.warning(f"Node {node.node_id}: Text extraction resulted in empty string after checks. Using placeholder.")
             text = f"Placeholder content for node {node.node_id}"

        texts.append(text)
        valid_nodes_in_batch.append(node)

    if not texts: logging.warning(f"Batch {batch_num}: No valid text extracted."); return 0

    try:
        # 2. Generate embeddings
        if not dense_model or not sparse_model: raise ValueError("Models not loaded")
        dense_embs_iter = dense_model.embed(texts, batch_size=len(texts))
        sparse_embs_iter = sparse_model.embed(texts, batch_size=len(texts))
        dense_embs = [emb.tolist() for emb in dense_embs_iter]
        sparse_vectors = []
        for emb_obj in sparse_embs_iter:
            if hasattr(emb_obj, 'indices') and hasattr(emb_obj, 'values'):
                sparse_vectors.append(SparseVector(indices=emb_obj.indices.tolist(), values=emb_obj.values.tolist()))
            else: sparse_vectors.append(SparseVector(indices=[], values=[]))

        if len(dense_embs) != len(texts) or len(sparse_vectors) != len(texts):
            raise ValueError(f"Embedding count mismatch for batch {batch_num}")

        # 3. Create PointStructs
        points = []
        processed_in_batch = 0
        for i, node in enumerate(valid_nodes_in_batch):
            try:
                payload_dict = node.model_dump(mode='json', exclude_none=True)
                points.append(PointStruct(
                    id=node.node_id, # Use the UUID
                    vector={
                        EMB_MODEL_NAME: dense_embs[i],        # Use dense model name key
                        SPARSE_MODEL_NAME: sparse_vectors[i] # Use sparse model name key
                    },
                    payload=payload_dict
                ))
                processed_in_batch += 1
            except Exception as point_err:
                logging.error(f"Error creating PointStruct for node {node.node_id} in batch {batch_num}: {point_err}", exc_info=True)

        # 4. Upsert
        if points:
            await client.upsert(collection_name=QDRANT_COLLECTION_NAME, points=points, wait=False) # wait=False for speed
            logging.debug(f"Upsert for batch {batch_num} sent ({len(points)} points).")
            return processed_in_batch
        else:
            logging.warning(f"Batch {batch_num}: No points generated after processing.")
            return 0
    except Exception as e:
        logging.error(f"Error processing/upserting batch {batch_num}: {e}", exc_info=True)
        return 0 # Indicate failure

# --- Phase 1b - Step 2: Process Group of Batches (Adapted) ---
async def process_batch_node_group(
    client: AsyncQdrantClient,
    group_node_slice: List[QdrantNode], # Takes a slice of QdrantNode objects
    group_start_batch_num: int,
    total_batches: int,
    dense_model: TextEmbedding,
    sparse_model: SparseTextEmbedding
):
    """Processes a specific group of node batches concurrently."""
    tasks = []
    num_nodes_in_group = len(group_node_slice)
    num_batches_in_group = (num_nodes_in_group + BATCH_SIZE - 1) // BATCH_SIZE
    logging.debug(f"Group {group_start_batch_num}: Processing {num_batches_in_group} batches ({num_nodes_in_group} nodes).")

    for i in range(0, num_nodes_in_group, BATCH_SIZE):
        batch_num = group_start_batch_num + (i // BATCH_SIZE)
        batch_node_list = group_node_slice[i : i + BATCH_SIZE]
        if not batch_node_list: continue # Skip empty slice

        tasks.append(
            upload_batch_nodes(client, batch_node_list, batch_num, total_batches, dense_model, sparse_model)
        )

    if not tasks:
        logging.warning(f"Group {group_start_batch_num}: No tasks generated.")
        return True # Nothing to do

    try:
        # Gather results (number processed per batch)
        results = await asyncio.gather(*tasks)
        total_processed_in_group = sum(results)
        logging.debug(f"Group {group_start_batch_num}: asyncio.gather completed. Processed points: {total_processed_in_group}")
        # Could add checks here if needed (e.g., if sum(results) == 0 despite tasks)
        return True # Group processed successfully (individual batch errors handled in upload_batch_nodes)
    except Exception as e:
        logging.error(f"Error gathering batch tasks for group {group_start_batch_num}: {e}", exc_info=True)
        return False # Indicate failure for this group

# --- Phase 1b - Step 1: Foreground Orchestrator (Adapted) ---
def run_bulk_ingest_foreground_nodes_with_ui(
    client: AsyncQdrantClient,
    nodes: List[QdrantNode], # Takes the list of QdrantNode objects
    dense_model: TextEmbedding,
    sparse_model: SparseTextEmbedding
):
    """Phase 1b: Uploads QdrantNode data in groups, updating UI progress."""
    start_ingest_time = time.perf_counter()
    logging.info("Phase 1b: Starting FOREGROUND bulk node ingest process...")

    # --- UI Elements ---
    progress_bar = st.progress(0.0)
    status_placeholder = st.empty()
    status_placeholder.info("Initiating node ingest process...")

    # --- Get Event Loop ---
    current_loop = get_or_create_eventloop()
    overall_successful = True
    ingest_time = 0.0
    has_data = False # Track if any data was successfully processed

    try:
        if not nodes:
            status_placeholder.warning("No nodes provided to ingest.")
            return False, 0.0, False

        total_docs = len(nodes)
        total_batches = (total_docs + BATCH_SIZE - 1) // BATCH_SIZE
        num_groups = (total_batches + INGEST_CONCURRENCY - 1) // INGEST_CONCURRENCY
        logging.info(f"Total nodes: {total_docs}, Batch Size: {BATCH_SIZE}, Total Batches: {total_batches}, Concurrency Groups: {num_groups}")

        processed_docs_count = 0

        # --- Loop through groups ---
        status_placeholder.info(f"Starting upload of {total_docs} nodes in {total_batches} batches ({num_groups} groups)...")
        time.sleep(0.1) # Slight delay for UI

        for group_idx in range(num_groups):
            group_start_batch_num = group_idx * INGEST_CONCURRENCY + 1
            group_end_batch_num = min((group_idx + 1) * INGEST_CONCURRENCY, total_batches)

            # Calculate node indices for the slice
            node_start_index = group_idx * INGEST_CONCURRENCY * BATCH_SIZE
            node_end_index = min((group_idx + 1) * INGEST_CONCURRENCY * BATCH_SIZE, total_docs)

            # Get the data slice for the current group (list of QdrantNode objects)
            group_data_slice = nodes[node_start_index : node_end_index]
            nodes_in_group = len(group_data_slice)

            if nodes_in_group == 0:
                logging.warning(f"Group {group_idx+1}/{num_groups}: Calculated slice is empty, skipping.")
                continue

            log_text = f"Processing group {group_idx+1}/{num_groups} (Batches {group_start_batch_num}-{group_end_batch_num}, {nodes_in_group} nodes)..."
            logging.info(log_text)
            status_placeholder.info(log_text + " Processing...") # UI update

            # --- Run async processing for THIS GROUP ---
            # This BLOCKS the UI thread until the group is done
            group_success = current_loop.run_until_complete(
                process_batch_node_group(
                    client, group_data_slice, group_start_batch_num, total_batches,
                    dense_model, sparse_model
                )
            )

            if not group_success:
                error_msg = f"Group {group_idx+1} failed. Aborting ingest."
                logging.error(error_msg)
                status_placeholder.error(error_msg)
                overall_successful = False
                break # Stop processing

            # --- Update Progress ---
            # Estimate processed docs based on index, not batch results to simplify
            processed_docs_count = min(node_end_index, total_docs)

            progress_percentage = processed_docs_count / total_docs if total_docs > 0 else 0.0
            progress_bar.progress(progress_percentage)
            status_placeholder.info(f"Completed ~{processed_docs_count}/{total_docs} nodes ({progress_percentage:.1%}). Processed group {group_idx+1}/{num_groups}.")
            time.sleep(0.05) # Allow UI to potentially update

        # --- After the loop ---
        if overall_successful:
            logging.info("All groups processed. Waiting for final upserts...")
            status_placeholder.info("Waiting for final upserts...")
            current_loop.run_until_complete(asyncio.sleep(5)) # Wait

            # Verify count (Approximate)
            final_count_res = current_loop.run_until_complete(client.count(collection_name=QDRANT_COLLECTION_NAME, exact=False))
            final_count = final_count_res.count
            logging.info(f"Final approximate count in Qdrant: {final_count}")
            if final_count > 0:
                has_data = True; logging.info("Bulk ingest completed successfully.")
                status_placeholder.success(f"Ingest complete. Final count: ~{final_count}")
            else:
                logging.warning(f"Ingest finished, but Qdrant count is {final_count}. Check logs.")
                status_placeholder.warning(f"Ingest finished, but Qdrant count is {final_count}. Check logs.")
                overall_successful = False
        # else: Failure already logged

    except Exception as e:
        error_msg = f"Bulk ingest coordination failed: {type(e).__name__}"
        logging.error(error_msg, exc_info=True)
        status_placeholder.error(f"{error_msg} - See console logs.")
        overall_successful = False; has_data = False
    finally:
        end_ingest_time = time.perf_counter()
        ingest_time = end_ingest_time - start_ingest_time
        logging.info(f"Bulk node ingest coordination ended after {ingest_time:.2f}s. Success: {overall_successful}")
        # Optionally clear UI elements after a delay
        # time.sleep(5)
        # status_placeholder.empty()
        # progress_bar.empty()

    return overall_successful, ingest_time, has_data

# --- Phase 2: Enable HNSW ---
async def enable_hnsw_phase2(client: AsyncQdrantClient, m_val=DENSE_M_PROD):
    """Phase 2: Enables HNSW index build."""
    if not client: logging.error("Qdrant client is None."); return False
    logging.info(f"Phase 2: Enabling HNSW (m={m_val}) for '{QDRANT_COLLECTION_NAME}'...")
    try:
        # Update HNSW config for the specific named vector
        update_op = await client.update_collection(
            collection_name=QDRANT_COLLECTION_NAME,
            # Specify the vector name to update
            vectors_config={
                EMB_MODEL_NAME: models.VectorParamsDiff(
                    hnsw_config=models.HnswConfigDiff(m=m_val)
                )
            }
        )
        if not update_op: logging.warning("Update collection op returned False (might be expected).")

        logging.info("HNSW update sent. Waiting for optimizer status GREEN...")
        # Wait for optimization to finish (check collection status)
        start_wait = time.time()
        timeout = 300 # 5 minutes timeout for optimization
        while time.time() - start_wait < timeout:
            collection_info = await client.get_collection(collection_name=QDRANT_COLLECTION_NAME)
            current_status = collection_info.status
            optim_status = collection_info.optimizer_status
            logging.debug(f"Collection status: {current_status}, Optimizer Status: {optim_status}")
            if current_status == CollectionStatus.GREEN and optim_status.ok:
                logging.info("Status GREEN & Optimizer OK. HNSW ready.")
                return True
            elif current_status == CollectionStatus.YELLOW or not optim_status.ok:
                logging.info(f"Status {current_status} (Optimizing: {optim_status.ok}, Err: {optim_status.error})... waiting 5s.")
                await asyncio.sleep(5)
            elif current_status == CollectionStatus.RED:
                logging.error(f"Status RED. Optimization failed: {optim_status.error}")
                st.error(f"Optimization failed (Status: RED): {optim_status.error}")
                return False
            else: # Should not happen
                logging.warning(f"Unexpected status: {current_status}. Waiting 5s.")
                await asyncio.sleep(5)

        logging.error("Timeout waiting for HNSW optimization to complete.")
        st.error("Timeout waiting for HNSW optimization.")
        return False

    except Exception as e:
        logging.error(f"Failed to enable HNSW: {e}", exc_info=True)
        st.error(f"Failed to enable HNSW: {e}")
        return False


###############################
# UI DISPLAY FUNCTIONS (Updated)
###############################

def display_table(table):
    """Display table data (accepts model or dict)."""
    if not table: return
    table_data = table.model_dump() if hasattr(table, 'model_dump') else table

    title = table_data.get("title")
    content = table_data.get("table_content")

    if title:
        st.subheader(f"Table: {title}")

    if content:
        # Attempt to display as DataFrame, fallback to Markdown
        try:
            if pd: # Check if pandas is available
                # Preprocess markdown for better parsing: remove leading/trailing whitespace, split lines
                lines = [line.strip() for line in content.strip().split('\n')]
                # Remove separator line if present
                if len(lines) > 1 and all(c in '-:| ' for c in lines[1]):
                    lines.pop(1)
                # Reconstruct cleaned content for StringIO
                cleaned_content = "\n".join(lines)

                df = pd.read_csv(
                    io.StringIO(cleaned_content),
                    sep="|",
                    skipinitialspace=True,
                    index_col=False # Important: Avoid assuming first column is index
                ).iloc[1:] # Skip the header row parsed as data
                # Clean up columns: remove leading/trailing spaces in names, drop empty columns
                df.columns = [col.strip() for col in df.columns]
                df = df.drop(columns=[col for col in df.columns if col == ''], errors='ignore') # Drop unnamed columns from extra separators
                df = df.dropna(axis=1, how='all')
                st.dataframe(df.reset_index(drop=True)) # Show clean index
            else:
                st.markdown(f"```markdown\n{content}\n```")
        except Exception as e:
            logging.warning(f"Pandas failed to parse table, showing markdown. Error: {e}. Content:\n{content}")
            st.markdown(f"```markdown\n{content}\n```")
    else:
        st.caption("Table content not available.")

def display_visual_element(visual):
    """Display visual element data (accepts model or dict)."""
    if not visual: return
    visual_data = visual.model_dump() if hasattr(visual, 'model_dump') else visual

    visual_type = visual_data.get('type', 'Unknown Type')
    description = visual_data.get('description', 'No description available.')
    data_summary = visual_data.get('data_summary') # Can be None or "N/A"

    st.markdown(f"**🖼️ {visual_type.capitalize()}**")
    st.markdown(f"*Description:* {description}")
    if data_summary and data_summary != "N/A":
        st.markdown(f"*Data Summary:* {data_summary}")
    # Consider adding visual_id if useful: st.caption(f"ID: {visual_data.get('visual_id')}")

def display_numerical_data_point(num_data):
    """Displays a single numerical data point."""
    if not num_data: return
    num_dict = num_data.model_dump() if hasattr(num_data, 'model_dump') else num_data

    value = num_dict.get('value', 'N/A')
    description = num_dict.get('description', 'No description.')
    context = num_dict.get('context')

    st.markdown(f"**{value}**: {description}")
    if context:
        st.caption(f"Context: {context}")

def display_qdrant_node_network(qdrant_nodes_data: Optional[List[Dict]], title="Concept/Entity Network"): # Accept List[Dict]
    """Displays relationships from QdrantNode data (as dicts) as a network graph."""
    if not qdrant_nodes_data: # Handles None or empty list
        st.info("No relationship data available to display.")
        return

    # Ensure input is a list of dictionaries
    if not isinstance(qdrant_nodes_data, list) or not all(isinstance(n, dict) for n in qdrant_nodes_data):
         st.error("Invalid data format passed to node network display function (expected list of dicts).")
         logging.error(f"display_qdrant_node_network received invalid data type: {type(qdrant_nodes_data)}")
         return

    if not nx or not plt:
        # ... (Fallback text display as before, using .get() on dicts) ...
        st.warning("NetworkX/Matplotlib not installed. Displaying text list.")
        st.markdown("#### Relationships (Text List):")
        for node_dict in qdrant_nodes_data:
             linked_nodes = node_dict.get('linked_nodes', [])
             if linked_nodes:
                 for link_dict in linked_nodes:
                     st.markdown(f"- **{node_dict.get('name','?')}** → **{link_dict.get('target_node_id','?')}**: {link_dict.get('relationship_description','?')} (Strength: {link_dict.get('relationship_strength','?')})")
        return

    # Create a network graph
    G = nx.DiGraph()
    node_labels = {}
    edge_weights = {}
    edge_labels = {}
    node_types_map = {} # Store node type for coloring

    # Add nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        node_id = node_dict.get('node_id')
        node_name = node_dict.get('name', 'Unknown Node')
        node_type = node_dict.get('node_type', 'unknown')
        if node_id and node_id not in G: # Check if node_id exists and is not None
            G.add_node(node_id)
            node_labels[node_id] = f"{node_name}\n({node_type})"
            node_types_map[node_id] = node_type # Store type

    # Add edges from linked_nodes using data from dictionaries
    for node_dict in qdrant_nodes_data:
        source_id = node_dict.get('node_id')
        linked_nodes = node_dict.get('linked_nodes', [])
        if not source_id or not G.has_node(source_id) or not isinstance(linked_nodes, list):
             continue # Skip if source_id invalid or no links

        for link_dict in linked_nodes:
             if not isinstance(link_dict, dict): continue # Skip invalid links
             target_id = link_dict.get('target_node_id')

             if not target_id: continue # Skip links without target

             # Ensure target node exists in the graph
             if target_id not in G:
                 G.add_node(target_id)
                 # Try to create a basic label from ID if target details aren't in current node list
                 target_name_match = re.match(r"entity_(.*?)_ch\d+|concept_(.*?)_ch\d+", target_id)
                 target_name_part = target_name_match.group(1) or target_name_match.group(2) if target_name_match else target_id
                 node_labels[target_id] = target_name_part.replace('_',' ').title() # Basic label
                 node_types_map[target_id] = 'unknown' # Mark as unknown type

             # Add edge if not already present
             if not G.has_edge(source_id, target_id):
                 strength = float(link_dict.get('relationship_strength', 5.0)) # Use get with default
                 description = link_dict.get('relationship_description', '')
                 G.add_edge(source_id, target_id, weight=strength)
                 edge_weights[(source_id, target_id)] = strength
                 edge_labels[(source_id, target_id)] = description

    # --- Graph drawing logic remains largely the same ---
    if not G.nodes():
         st.info("No nodes found to build the graph.")
         return
    if not G.edges():
         st.info("Nodes found, but no links between them to display.")
         st.write("Nodes identified:", list(node_labels.values()))
         return

    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(12, 10))
    try:
        pos = nx.spring_layout(G, k=0.5, iterations=50, seed=42)
    except Exception as layout_err:
        logging.error(f"NetworkX layout failed: {layout_err}. Using random layout.")
        pos = nx.random_layout(G, seed=42)

    # Draw nodes with color based on type
    node_colors = ['skyblue' if node_types_map.get(n) == 'entity' else 'lightgreen' if node_types_map.get(n) == 'concept' else 'lightgray' for n in G.nodes()]
    nx.draw_networkx_nodes(G, pos, node_size=1500, node_color=node_colors, alpha=0.9, ax=ax)

    # Draw edges
    if edge_weights: # Check if edges exist before calculating weights
        min_strength = min(edge_weights.values()) if edge_weights else 1
        max_strength = max(edge_weights.values()) if edge_weights else 10
        if max_strength == min_strength: max_strength += 1
        edge_width = [1 + 4 * (edge_weights.get(edge, 5.0) - min_strength) / (max_strength - min_strength) for edge in G.edges()]
        nx.draw_networkx_edges(G, pos, width=edge_width, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)
    else: # Draw default thin edges if no weights somehow
         nx.draw_networkx_edges(G, pos, width=1.0, alpha=0.6, edge_color='gray',
                              arrowsize=15, connectionstyle='arc3,rad=0.1', ax=ax)


    # Draw labels
    nx.draw_networkx_labels(G, pos, labels=node_labels, font_size=9, font_weight='normal', ax=ax)

    ax.set_title(title, fontsize=16)
    fig.tight_layout()
    plt.axis('off')
    st.pyplot(fig)
    plt.close(fig)

    # Display edge details table
    edge_data = []
    # Use edge_labels which was populated correctly
    for (source_id, target_id), desc in edge_labels.items():
        edge_data.append({
            "Source": node_labels.get(source_id, source_id),
            "Target": node_labels.get(target_id, target_id),
            "Strength": edge_weights.get((source_id, target_id), 5.0),
            "Description": desc
        })

    if edge_data:
        st.markdown("### Relationship Details")
        try:
             if pd: # Check if pandas is available
                 df = pd.DataFrame(edge_data)
                 st.dataframe(df)
             else:
                 for item in edge_data: st.write(item) # Fallback display
        except Exception as df_err:
            logging.error(f"Failed to create DataFrame for edge data: {df_err}")
            for item in edge_data: st.write(item) # Fallback display
            
def display_chapter(chapter):
    """Display chapter details including subsections and Qdrant node relationships."""
    if not chapter: return
    chapter_data = chapter.model_dump() if hasattr(chapter, 'model_dump') else chapter

    title = chapter_data.get('title', 'Untitled Chapter')
    summary = chapter_data.get('summary', 'No summary available.')
    subsections = chapter_data.get('subsections', [])
    qdrant_nodes = chapter_data.get('qdrant_nodes') # This holds QdrantNode objects/dicts

    st.markdown(f"## {title}")
    st.markdown(f"_{summary}_")

    # Display Qdrant Node relationships if available
    if qdrant_nodes:
        with st.expander("View Concepts/Entities and Relationships within this Chapter", expanded=False):
            # Pass the actual list of nodes (could be dicts or models)
            display_qdrant_node_network(qdrant_nodes, title=f"Network for '{title}'")
    else:
        st.caption("No concept/entity relationship data available for this chapter.")

    # Display subsections
    if subsections:
        st.markdown("### Subsections Contained")
        for sub_idx, subsection_data in enumerate(subsections):
            if not subsection_data: continue
            # Handle dict or model
            sub_obj = subsection_data
            if isinstance(subsection_data, dict):
                 # Minimal validation or create Subsection obj if needed for consistency
                 # sub_obj = Subsection(**subsection_data) # Optional
                 pass

            sub_title = getattr(sub_obj, 'title', f'Subsection {sub_idx + 1}')
            sub_text = getattr(sub_obj, 'text', 'No text available.')
            sub_description = getattr(sub_obj, 'description', '')
            sub_id = getattr(sub_obj, 'subsection_id', 'N/A')
            page_num = "?" # Page number not stored on subsection anymore
            is_cutoff = getattr(sub_obj, 'is_cutoff', False)
            ref_visuals = getattr(sub_obj, 'referenced_visuals', [])
            ref_tables = getattr(sub_obj, 'referenced_tables', [])

            expander_title = f"{sub_idx + 1}. {sub_title}"
            if is_cutoff: expander_title += " (Continues...)"

            with st.expander(expander_title):
                if sub_description:
                    st.caption(f"Description: {sub_description}")
                st.markdown(sub_text) # Display the actual text
                st.caption(f"Subsection ID: {sub_id}")

                # Display referenced visuals/tables if any
                if ref_visuals: st.write(f"**Ref Visuals:** {', '.join(ref_visuals)}")
                if ref_tables: st.write(f"**Ref Tables:** {', '.join(ref_tables)}")
    else:
        st.caption("No subsections listed for this chapter.")

# --- Qdrant Version of Node Link Display ---
async def display_node_links_qdrant(
    qdrant_client: AsyncQdrantClient,
    collection_name: str,
    selected_node_id: str
):
    """
    Fetches a specific node and its linked targets directly from Qdrant
    and displays its outgoing relationships.

    Args:
        qdrant_client: An initialized AsyncQdrantClient instance.
        collection_name: The name of the Qdrant collection.
        selected_node_id: The node_id of the node to inspect.
    """
    function_name = "display_node_links_qdrant"
    if not selected_node_id:
        st.caption("No node selected.")
        return
    if not qdrant_client:
        st.error("Qdrant client is not available.")
        return

    st.markdown("*(Fetching data from Qdrant...)*") # Indicate loading

    # --- Fetch the Selected Node (Source) ---
    selected_payload = None
    try:
        retrieved_points = await qdrant_client.retrieve(
            collection_name=collection_name,
            ids=[selected_node_id],
            with_payload=True,
            with_vectors=False
        )
        if not retrieved_points:
            st.warning(f"Node ID '{selected_node_id}' not found in Qdrant collection '{collection_name}'.")
            return
        selected_point = retrieved_points[0]
        if not selected_point.payload:
            st.warning(f"Node '{selected_node_id}' found but has no payload data in Qdrant.")
            return
        selected_payload = selected_point.payload # This is the dictionary stored in Qdrant
        logging.debug(f"[{function_name}] Retrieved source node '{selected_node_id}' payload.")

    except Exception as e:
        st.error(f"Error retrieving node '{selected_node_id}' from Qdrant: {e}")
        logging.error(f"[{function_name}] Qdrant retrieve failed for {selected_node_id}: {e}", exc_info=True)
        return

    # --- Extract Links & Target IDs ---
    source_name = selected_payload.get('name', selected_node_id)
    # *** FIXED: Access direct fields from payload ***
    source_doc = selected_payload.get('source', '?')
    source_chapter = selected_payload.get('chapter_id_context', '?')
    source_title = selected_payload.get('title', 'N/A')
    source_claim = selected_payload.get('core_claim', 'N/A')
    source_type = selected_payload.get('node_type', 'unknown')

    linked_nodes_data = selected_payload.get('linked_nodes', [])

    st.markdown(f"#### Relationships FROM: **{source_name}** (`{source_type}`) ")
    st.caption(f"*(Origin: Doc '{source_doc}', Context '{source_chapter}')*")
    with st.expander("View Source Node Details"):
        st.markdown(f"**Title:** {source_title}")
        st.markdown(f"**Core Claim:** {source_claim}")
        # Display other relevant fields directly
        st.markdown(f"**Info Type:** {selected_payload.get('information_type', 'N/A')}")
        st.markdown(f"**Sentiment:** {selected_payload.get('sentiment_tone', 'N/A')}")
        tags_str = selected_payload.get('tags', '[]')
        try: tags_list = json.loads(tags_str)
        except: tags_list = []
        if tags_list: st.markdown(f"**Tags:** {', '.join(tags_list)}")
        st.caption(f"ID: `{selected_node_id}`")
    st.markdown("---")


    if not linked_nodes_data or not isinstance(linked_nodes_data, list):
        st.caption("No outgoing links found for this node in Qdrant data.")
        return

    target_ids = list(set([ # Get unique target IDs
        link.get('target_node_id')
        for link in linked_nodes_data if isinstance(link, dict) and link.get('target_node_id')
    ]))

    if not target_ids:
        st.caption("Links found, but no valid target IDs specified.")
        return

    # --- Fetch Target Nodes (Batch) ---
    target_points_map: Dict[str, Dict] = {} # Map ID -> Payload Dict
    try:
        logging.debug(f"[{function_name}] Retrieving details for {len(target_ids)} target nodes: {target_ids}")
        target_points = await qdrant_client.retrieve(
            collection_name=collection_name,
            ids=target_ids,
            with_payload=True,
            with_vectors=False
        )
        for point in target_points:
            if point.payload: # Only store if payload exists
                target_points_map[point.id] = point.payload
        logging.debug(f"[{function_name}] Successfully retrieved payloads for {len(target_points_map)} target nodes.")

    except Exception as e:
        st.error(f"Error retrieving target node details from Qdrant: {e}")
        logging.error(f"[{function_name}] Qdrant retrieve failed for targets {target_ids}: {e}", exc_info=True)
        # Continue, but target details will be missing

    # --- Display Links with Target Context ---
    links_displayed_count = 0
    for i, link_data in enumerate(linked_nodes_data):
        if not isinstance(link_data, dict): continue

        target_id = link_data.get('target_node_id')
        if not target_id: continue

        description = link_data.get('relationship_description', 'No description')
        strength = 0.0
        try: strength = float(link_data.get('relationship_strength', 0.0))
        except (ValueError, TypeError): pass # Keep strength 0.0
        keywords = link_data.get('relationship_keywords', [])
        if not isinstance(keywords, list): keywords = []

        # Get target details from the map
        target_payload = target_points_map.get(target_id)
        target_name = target_id
        target_type = "unknown"
        target_doc = "unknown"
        target_chapter = "unknown"
        target_title = "" # For context
        # target_claim = "" # Optional

        if target_payload:
            target_name = target_payload.get('name', target_id)
            target_type = target_payload.get('node_type', 'unknown')
            # *** FIXED: Access direct fields from target payload ***
            target_doc = target_payload.get('source', 'unknown')
            target_chapter = target_payload.get('chapter_id_context', 'unknown')
            target_title = target_payload.get('title', '')
            # target_claim = target_payload.get('core_claim', '')
        else:
            # Target not found or had no payload
            st.caption(f"_(Warning: Details for target node '{target_id}' could not be retrieved from Qdrant.)_")

        st.markdown(f"**{i+1}. Target Node:** {target_name} (`{target_type}`)")
        st.markdown(f"    *Origin: Doc '{target_doc}', Context '{target_chapter}'*")
        if target_title: st.markdown(f"    *Target Title:* {target_title}")
        st.markdown(f"    **Link Description:** {description}")
        st.markdown(f"    **Strength:** {strength:.1f}/10.0")
        if keywords:
            st.markdown(f"    **Keywords:** " + ", ".join(keywords))
        st.markdown(f"    *Target Node ID:* `{target_id}`")
        st.markdown("---")
        links_displayed_count += 1


    if links_displayed_count == 0:
        st.caption("No valid outgoing links could be displayed (check target IDs and Qdrant data).")


def display_node_links_from_list(all_nodes: List[Dict | Any], # List of node dicts or Pydantic objects
                       selected_node_id: str,
                       node_map: Optional[Dict[str, Dict | Any]] = None):
    """
    Finds a specific node by ID and displays its outgoing relationships (linked_nodes),
    including the origin context (document, chapter) of the target node.

    Args:
        all_nodes: A list containing all nodes (as dictionaries or Pydantic objects).
        selected_node_id: The node_id of the node to inspect.
        node_map: Optional pre-built dictionary mapping node_id to node object/dict
                  for faster lookups, especially for target node info.
    """
    function_name = "display_node_links_from_list"
    if not selected_node_id:
        st.caption("No node selected.")
        return

    # --- Find the selected node (Source Node) ---
    selected_node = None
    if node_map and selected_node_id in node_map:
        selected_node = node_map[selected_node_id]
    else:
        # Fallback search if map not provided
        for node in all_nodes:
            is_model = not isinstance(node, dict)
            node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
            if node_id == selected_node_id:
                selected_node = node
                break

    if selected_node is None:
        st.warning(f"Node with ID '{selected_node_id}' not found.")
        return

    is_source_model = not isinstance(selected_node, dict)
    source_name = getattr(selected_node, 'name', selected_node_id) if is_source_model else selected_node.get('name', selected_node_id)
    source_doc = getattr(selected_node, 'document_context', '?') if is_source_model else selected_node.get('document_context', '?')
    source_chapter = getattr(selected_node, 'chapter_id_context', '?') if is_source_model else selected_node.get('chapter_id_context', '?')

    st.markdown(f"#### Relationships FROM: **{source_name}**")
    st.caption(f"*(Source Origin: Document '{source_doc}', Chapter '{source_chapter}')*")
    st.markdown("---")


    linked_nodes = getattr(selected_node, 'linked_nodes', []) if is_source_model else selected_node.get('linked_nodes', [])

    if not linked_nodes or not isinstance(linked_nodes, list):
        st.caption("No outgoing links found for this node.")
        return

    # --- Build Node Map if needed (for target node lookups) ---
    if node_map is None:
        logging.debug(f"[{function_name}] Building temporary node_map for target info lookup.")
        node_map = {}
        for node in all_nodes:
             inner_is_model = not isinstance(node, dict)
             inner_node_id = getattr(node, 'node_id', None) if inner_is_model else node.get('node_id')
             if inner_node_id:
                 node_map[inner_node_id] = node

    # --- Display Links with Target Context ---
    if not node_map: # Check if map is still empty after trying to build
         st.warning("Could not build node map to look up target details.")

    for i, link in enumerate(linked_nodes):
        link_is_model = not isinstance(link, dict)
        target_id = getattr(link, 'target_node_id', 'Unknown Target ID') if link_is_model else link.get('target_node_id', 'Unknown Target ID')
        description = getattr(link, 'relationship_description', 'No description') if link_is_model else link.get('relationship_description', 'No description')
        strength = getattr(link, 'relationship_strength', 0.0) if link_is_model else link.get('relationship_strength', 0.0)
        keywords = getattr(link, 'relationship_keywords', []) if link_is_model else link.get('relationship_keywords', [])
        if not isinstance(keywords, list): keywords = []

        # --- Find Target Node Info using the map ---
        target_node = node_map.get(target_id) if node_map else None
        target_name = target_id # Default to ID if not found
        target_type = "unknown"
        target_doc = "unknown"
        target_chapter = "unknown"

        if target_node:
             target_is_model = not isinstance(target_node, dict)
             target_name = getattr(target_node, 'name', target_id) if target_is_model else target_node.get('name', target_id)
             target_type = getattr(target_node, 'node_type', 'unknown') if target_is_model else target_node.get('node_type', 'unknown')
             target_doc = getattr(target_node, 'document_context', 'unknown') if target_is_model else target_node.get('document_context', 'unknown')
             target_chapter = getattr(target_node, 'chapter_id_context', 'unknown') if target_is_model else target_node.get('chapter_id_context', 'unknown')
        elif target_id != 'Unknown Target ID':
             logging.warning(f"[{function_name}] Target node '{target_id}' linked from '{selected_node_id}' not found in node map.")


        # Display formatted link info including target context
        st.markdown(f"**{i+1}. Target Node:** {target_name} (`{target_type}`)")
        # Add target origin using italics or caption
        st.markdown(f"    *Origin: Document '{target_doc}', Chapter '{target_chapter}'*")
        st.markdown(f"    **Link Description:** {description}")
        st.markdown(f"    **Strength:** {strength:.1f}/10.0")
        if keywords:
            st.markdown(f"    **Keywords:** " + ", ".join(keywords))
        st.markdown(f"    *Target Node ID:* `{target_id}`") # Keep ID easily visible
        st.markdown("---") # Separator

    if not linked_nodes: # Redundant check, but safe
        st.caption("No valid outgoing links found to display.")

# --- Qdrant Ingestion Helper Functions ---

async def create_collection_if_not_exists(
    client: AsyncQdrantClient,
    collection_name: str,
    vector_model_name: str,
    vector_dim: int,
    sparse_vector_model_name: str
):
    """Checks if the collection exists and creates it if not."""
    try:
        collections = await client.get_collections()
        existing_names = {col.name for col in collections.collections}

        if collection_name in existing_names:
            logging.info(f"Collection '{collection_name}' already exists.")
            # Optional: Check if config matches and warn/recreate if needed
            # collection_info = await client.get_collection(collection_name=collection_name)
            # Compare vector params, etc.
            return True

        logging.info(f"Collection '{collection_name}' not found. Creating...")
        await client.create_collection(
            collection_name=collection_name,
            vectors_config={
                # Dense vector config
                vector_model_name: models.VectorParams(
                    size=vector_dim,
                    distance=models.Distance.COSINE,
                    on_disk=True # Recommended for potentially large vectors
                )
            },
            sparse_vectors_config={
                # Sparse vector config
                sparse_vector_model_name: models.SparseVectorParams(
                    index=models.SparseIndexParams(
                        on_disk=True # Store sparse index on disk
                    )
                )
            }
            # Add payload indexing here later if needed for performance:
            # payload_schema={
            #     "node_type": models.PayloadSchemaType.KEYWORD,
            #     "properties.source": models.PayloadSchemaType.KEYWORD,
            #     "chapter_id_context": models.PayloadSchemaType.KEYWORD,
            #     "properties.report_year": models.PayloadSchemaType.INTEGER # Example
            # }
        )
        logging.info(f"Collection '{collection_name}' created successfully.")
        return True
    except Exception as e:
        logging.error(f"Failed to create or check collection '{collection_name}': {e}", exc_info=True)
        st.error(f"Failed to setup Qdrant collection '{collection_name}': {e}")
        return False

async def upsert_qdrant_nodes(
    client: AsyncQdrantClient,
    collection_name: str,
    nodes: List[QdrantNode], # Expects list of Pydantic models
    dense_model: TextEmbedding,
    sparse_model: SparseTextEmbedding,
    vector_model_name: str, # e.g., "BAAI/bge-base-en-v1.5"
    sparse_vector_model_name: str, # e.g., "bm25"
    batch_size: int = BATCH_SIZE
):
    """
    Generates embeddings for QdrantNodes and upserts them into the collection.

    Args:
        client: Initialized AsyncQdrantClient.
        collection_name: Name of the target Qdrant collection.
        nodes: List of QdrantNode Pydantic objects.
        dense_model: Loaded fastembed TextEmbedding model.
        sparse_model: Loaded fastembed SparseTextEmbedding model.
        vector_model_name: Name used for the dense vector in Qdrant config.
        sparse_vector_model_name: Name used for the sparse vector in Qdrant config.
        batch_size: Number of points to upsert in each API call.
    """
    if not nodes:
        logging.warning("No nodes provided for Qdrant upsert.")
        st.warning("No node data found in the project ontology to upsert.")
        return True # Nothing to do, considered success

    logging.info(f"Starting Qdrant upsert for {len(nodes)} nodes into '{collection_name}'.")

    # --- 1. Prepare Texts for Embedding ---
    texts_to_embed = []
    valid_nodes_for_upsert = [] # Store nodes corresponding to texts_to_embed
    node_ids_processed = set()

    for node in nodes:
        if not isinstance(node, QdrantNode) or not node.node_id:
            logging.warning(f"Skipping invalid node object during Qdrant prep: {type(node)}")
            continue
        # Avoid processing duplicate node IDs if they somehow exist in the input list
        if node.node_id in node_ids_processed:
            logging.warning(f"Skipping duplicate node ID '{node.node_id}' during Qdrant prep.")
            continue

        # Combine name and core_claim for embedding (adjust if needed)
        node_name = getattr(node, 'name', '')
        core_claim = node.properties.get('core_claim', '') if isinstance(getattr(node, 'properties', None), dict) else ''
        text = f"{node_name}: {core_claim}".strip()
        if not text or text == ":": # Handle cases where name/claim are empty
             logging.warning(f"Node '{node.node_id}' has empty name/claim. Using name only or placeholder.")
             text = node_name if node_name else f"Placeholder for node {node.node_id}"

        texts_to_embed.append(text)
        valid_nodes_for_upsert.append(node)
        node_ids_processed.add(node.node_id)

    if not texts_to_embed:
        logging.error("No valid text could be extracted from nodes for embedding.")
        st.error("Could not prepare any node data for Qdrant ingestion.")
        return False

    logging.info(f"Prepared {len(texts_to_embed)} texts for embedding.")

    # --- 2. Generate Embeddings (Batched) ---
    try:
        logging.debug(f"Generating dense embeddings using '{vector_model_name}'...")
        # Fastembed handles internal batching, but provide list directly
        dense_embeddings = list(dense_model.embed(texts_to_embed))

        logging.debug(f"Generating sparse embeddings using '{sparse_vector_model_name}'...")
        sparse_embeddings = list(sparse_model.embed(texts_to_embed))

        if len(dense_embeddings) != len(texts_to_embed) or len(sparse_embeddings) != len(texts_to_embed):
            raise ValueError("Mismatch between number of texts and generated embeddings.")
        logging.info("Embeddings generated successfully.")

    except Exception as e:
        logging.error(f"Failed to generate embeddings: {e}", exc_info=True)
        st.error(f"Embedding generation failed: {e}")
        return False

    # --- 3. Create PointStructs and Upsert in Batches ---
    points_batch: List[PointStruct] = []
    total_upserted = 0
    total_failed = 0

    status_placeholder = st.empty() # For progress updates

    for i, node in enumerate(valid_nodes_for_upsert):
        try:
            # Get corresponding embeddings
            dense_vec = dense_embeddings[i].tolist()
            sparse_emb_obj = sparse_embeddings[i]
            sparse_indices = sparse_emb_obj.indices.tolist() if hasattr(sparse_emb_obj, 'indices') else []
            sparse_values = sparse_emb_obj.values.tolist() if hasattr(sparse_emb_obj, 'values') else []

            # Convert Pydantic node to dict for payload
            payload_dict = node.model_dump(mode='json', exclude_none=True) # Exclude None values if desired

            # Create PointStruct
            point = PointStruct(
                id=node.node_id,
                vector={
                    vector_model_name: dense_vec,
                    sparse_vector_model_name: SparseVector(indices=sparse_indices, values=sparse_values)
                },
                payload=payload_dict
            )
            points_batch.append(point)

            # Upsert when batch is full or it's the last item
            if len(points_batch) >= batch_size or i == len(valid_nodes_for_upsert) - 1:
                logging.debug(f"Upserting batch of {len(points_batch)} points (Total processed: {i+1}/{len(valid_nodes_for_upsert)})...")
                status_placeholder.info(f"Upserting nodes {total_upserted + 1} - {i + 1} / {len(valid_nodes_for_upsert)}...")
                try:
                    await client.upsert(
                        collection_name=collection_name,
                        points=points_batch,
                        wait=True # Wait for confirmation, helps manage rate limits
                    )
                    total_upserted += len(points_batch)
                    logging.debug(f"Batch upsert successful.")
                except Exception as upsert_err:
                    logging.error(f"Failed to upsert batch ending at index {i}: {upsert_err}", exc_info=True)
                    st.warning(f"Failed to upsert batch of {len(points_batch)} nodes. Check Qdrant logs.")
                    total_failed += len(points_batch)
                finally:
                    points_batch = [] # Clear batch

        except Exception as point_prep_err:
            logging.error(f"Failed to prepare PointStruct for node '{node.node_id}': {point_prep_err}", exc_info=True)
            total_failed += 1
            if points_batch: # Also try to upsert remaining points in batch before error
                 logging.warning("Attempting to upsert partial batch before skipping errored node.")
                 try:
                     await client.upsert(collection_name=collection_name, points=points_batch, wait=True)
                     total_upserted += len(points_batch)
                 except Exception as partial_upsert_err:
                      logging.error(f"Failed to upsert partial batch: {partial_upsert_err}")
                      total_failed += len(points_batch) # Count these as failed too
                 finally:
                      points_batch = [] # Clear batch

    status_placeholder.empty() # Clear progress message

    # --- Final Report ---
    logging.info(f"Qdrant upsert finished. Successfully upserted: {total_upserted}, Failed/Skipped: {total_failed}")
    if total_failed > 0:
        st.error(f"Qdrant Ingestion: {total_upserted} nodes successfully upserted, {total_failed} failed or were skipped.")
        return False
    elif total_upserted == 0 and len(valid_nodes_for_upsert) > 0:
        st.error("Qdrant Ingestion: No nodes were successfully upserted, although data was present.")
        return False
    else:
        st.success(f"Qdrant Ingestion: Successfully upserted {total_upserted} nodes into '{collection_name}'.")
        return True

def display_project_ontology(ontology: Optional[ProjectOntology]): # Type hint remains Optional[ProjectOntology]
    """Displays project-wide ontology including the aggregated graph.
       Expects a ProjectOntology Pydantic model object as input.
    """
    function_name = "display_project_ontology"
    if not ontology:
        st.warning("No project ontology data available to display.")
        return

    # --- Treat input strictly as a ProjectOntology object ---
    # Remove the is_model check and the dictionary fallback logic
    if not hasattr(ontology, '__class__') or ontology.__class__.__name__ != 'ProjectOntology':
        st.error(f"[display_project_ontology] Error: Expected a ProjectOntology object based on name, but received type {type(ontology)}. Cannot display.")
        logging.error(f"[display_project_ontology] display_project_ontology received unexpected type: {type(ontology)}")
        return

    try:
        # Access attributes directly using getattr for safety or direct access
        title = getattr(ontology, 'title', 'Project Ontology [Default]')
        overview = getattr(ontology, 'overview', 'No overview available.')
        doc_count = getattr(ontology, 'document_count', 0)
        documents = getattr(ontology, 'documents', [])
        global_themes = getattr(ontology, 'global_themes', [])
        key_concepts = getattr(ontology, 'key_concepts', [])
        project_nodes = getattr(ontology, 'project_graph_nodes', None) # Access the graph nodes field

        st.title(f"🌐 {title}")
        st.markdown(f"**Documents Analyzed:** {doc_count}")

        # Display documents included
        if documents:
            st.markdown("## Documents Included In Analysis")
            if hasattr(st, 'chip'):
                for doc_title in documents: st.chip(doc_title, icon="📄")
            else:  # Fallback for older Streamlit versions
                docs_html = " ".join([f"<span style='background-color:#c5e1ff; padding:5px; margin:2px; border-radius:5px'>{doc}</span>" for doc in documents])
                st.markdown(docs_html, unsafe_allow_html=True)

        # Display overview
        st.markdown("## Overview")
        st.markdown(overview)

        # Display global themes
        if global_themes:
            st.markdown("## Global Themes")
            if hasattr(st, 'chip'):
                for theme in global_themes: st.chip(theme, icon="🏷️")
            else: # Fallback
                themes_html = " ".join([f"<span style='background-color:#e6f3ff; padding:5px; margin:2px; border-radius:5px'>#{tag}</span>" for tag in global_themes])
                st.markdown(themes_html, unsafe_allow_html=True)

        # Display key concepts
        if key_concepts:
            st.markdown("## Key Concepts")
            if hasattr(st, 'chip'):
                 for concept in key_concepts: st.chip(concept, icon="💡")
            else: # Fallback
                concepts_html = " ".join([f"<span style='background-color:#f0f0f0; padding:5px; margin:2px; border-radius:5px'>{concept}</span>" for concept in key_concepts])
                st.markdown(concepts_html, unsafe_allow_html=True)

        # Display Project-Level Graph
        st.divider()
        st.markdown("## Project Knowledge Graph")
        if project_nodes and isinstance(project_nodes, list):
            logging.info(f"[{function_name}] Rendering project graph with {len(project_nodes)} nodes.")
            with st.expander("View Full Project Concept/Entity Network", expanded=False):
                 try:
                     # Convert Pydantic nodes back to dicts IF display_qdrant_node_network expects dicts
                     # If it accepts Pydantic objects, you can pass project_nodes directly
                     nodes_as_dicts = [node.model_dump(mode='json') if hasattr(node, 'model_dump') else node for node in project_nodes]
                     display_qdrant_node_network(nodes_as_dicts, title="Project-Wide Concept/Entity Network")
                 except NameError:
                      st.error("`display_qdrant_node_network` function not found.")
                      logging.error("`display_qdrant_node_network` function not found.")
                 except Exception as graph_err:
                      st.error(f"Failed to render project network graph: {graph_err}")
                      logging.error(f"Failed to render project network graph: {graph_err}", exc_info=True)
        else:
            st.caption("No cross-document relationship data was generated or processed to display a project graph.")


        # --- NEW: Node Relationship Explorer ---
        st.divider()
        st.markdown("## Node Relationship Explorer")

        # Use the Qdrant client instance
        qdrant_client = get_qdrant_client() # Get cached client

        if qdrant_client is None:
            st.error("Qdrant client is unavailable. Cannot explore relationships.")
            return # Stop if client failed

        # --- Use project_graph_nodes to populate the SELECTBOX only ---
        project_nodes_list = getattr(ontology, 'project_graph_nodes', None)

        if project_nodes_list and isinstance(project_nodes_list, list):
            node_options: Dict[str, str] = {} # Map display label -> node_id
            # Sort nodes alphabetically by name for the dropdown
            # Ensure sorting handles both Pydantic objects and potential dicts safely
            try:
                sorted_nodes = sorted(
                    project_nodes_list,
                    key=lambda n: getattr(n, 'name', '') if hasattr(n, 'name') else (n.get('name', '') if isinstance(n, dict) else '')
                )
            except Exception as sort_err:
                logging.warning(f"Failed to sort nodes for dropdown: {sort_err}. Using original order.")
                sorted_nodes = project_nodes_list

            for node in sorted_nodes:
                is_model = not isinstance(node, dict)
                node_id = getattr(node, 'node_id', None) if is_model else node.get('node_id')
                node_name = getattr(node, 'name', node_id) if is_model else node.get('name', node_id)
                node_type = getattr(node, 'node_type', 'unknown') if is_model else node.get('node_type', 'unknown')

                if node_id:
                    display_label = f"{node_name} ({node_type})"
                    count = 1
                    original_label = display_label
                    while display_label in node_options:
                        count += 1
                        display_label = f"{original_label} ({count})"
                    node_options[display_label] = node_id
            # --- End of populating node_options ---

            if node_options:
                selected_label = st.selectbox(
                    "Select a Node to see its connections:",
                    options=["Select a Node..."] + list(node_options.keys()),
                    index=0,
                    key="project_node_selector_qdrant" # Ensure unique key
                )

                if selected_label != "Select a Node...":
                    selected_id = node_options.get(selected_label)
                    if selected_id:
                        # --- !!! CHANGE IS HERE !!! ---
                        try:
                            logging.info(f"Running Qdrant link display for node: {selected_id} using asyncio.run")
                            # Use asyncio.run() instead of the custom run_async helper
                            # This is the standard way to run an async function from sync code.
                            asyncio.run(display_node_links_qdrant(
                                qdrant_client,
                                QDRANT_COLLECTION_NAME, # Pass collection name
                                selected_id
                            ))
                            logging.info(f"Finished Qdrant link display for node: {selected_id}")

                        except RuntimeError as e:
                            # Handle potential issue if a loop is already running
                            if "cannot run event loop while another event loop is running" in str(e):
                                st.error("Async Error: Cannot run nested event loops. This may happen in some Streamlit setups.")
                                logging.error("Cannot run asyncio.run; loop already running.", exc_info=True)
                                # If this happens, more complex solutions like nest_asyncio might be needed
                                # import nest_asyncio
                                # nest_asyncio.apply()
                                # try:
                                #    asyncio.run(display_node_links_qdrant(...)) # Try again
                                # except Exception as nested_e: ...
                            elif "Event loop is closed" in str(e):
                                st.error("Async Runtime Error: The event loop was closed unexpectedly.")
                                logging.error("Event loop closed error persisted even with asyncio.run.", exc_info=True)
                            else:
                                st.error(f"Async Runtime Error: {e}")
                                logging.error(f"Async Runtime Error: {e}", exc_info=True)
                        except Exception as q_err:
                            st.error(f"Failed to display node links from Qdrant: {q_err}")
                            logging.error(f"Error running display_node_links_qdrant via asyncio.run: {q_err}", exc_info=True)
                        # --- !!! END OF CHANGE !!! ---
                    else:
                        st.warning("Could not find the ID for the selected node.")
            else:
                st.info("No valid nodes available to select.")
        else:
            st.caption("No project nodes available to explore relationships.")
        # --- END: Node Relationship Explorer ---

    except AttributeError as attr_err:
        # Catch cases where expected attributes are missing from the model instance
        st.error(f"Error accessing project ontology data: Missing attribute - {attr_err}")
        logging.error(f"[{function_name}] Missing attribute in ProjectOntology object: {attr_err}", exc_info=True)
    except Exception as e:
        # Catch any other unexpected errors during display
        st.error(f"An unexpected error occurred while displaying the project ontology: {e}")
        logging.error(f"[{function_name}] Unexpected error displaying ontology: {e}", exc_info=True)


def display_page_details(page_content_data):
    """Displays details for a single page based on the simplified PageContent model."""
    if not page_content_data: return
    page_dict = page_content_data.model_dump() if hasattr(page_content_data, 'model_dump') else page_content_data

    page_num = page_dict.get("page_number", "?")
    st.header(f"📑 Page {page_num} Details")

    # Display Content Flags
    st.caption(f"Detected Content: "
               f"{'Tables ✅' if page_dict.get('has_tables') else 'Tables ❌'} | "
               f"{'Visuals ✅' if page_dict.get('has_visuals') else 'Visuals ❌'} | "
               f"{'Numbers ✅' if page_dict.get('has_numbers') else 'Numbers ❌'}")
    st.divider()

    # Display Subsections first as they contain the core text
    subsections = page_dict.get('subsections', [])
    if subsections:
        st.subheader("Content Subsections")
        for sub_idx, sub_data in enumerate(subsections):
             sub_title = sub_data.get('title', f'Subsection {sub_idx+1}')
             sub_desc = sub_data.get('description', '')
             sub_text = sub_data.get('text', 'No text.')
             sub_id = sub_data.get('subsection_id', 'N/A')
             is_cutoff = sub_data.get('is_cutoff', False)
             ref_v = sub_data.get('referenced_visuals', [])
             ref_t = sub_data.get('referenced_tables', [])

             exp_title = f"{sub_idx+1}. {sub_title}"
             if is_cutoff: exp_title += " (Continues...)"

             with st.expander(exp_title, expanded=True): # Expand by default
                 if sub_desc:
                     st.caption(f"Description: {sub_desc}")
                 st.markdown(sub_text)
                 # Optionally show refs if needed
                 if ref_v or ref_t:
                      refs = []
                      if ref_v: refs.append(f"Visuals: {', '.join(ref_v)}")
                      if ref_t: refs.append(f"Tables: {', '.join(ref_t)}")
                      st.caption(f"References: {'; '.join(refs)} | ID: {sub_id}")
                 else:
                      st.caption(f"ID: {sub_id}")

        st.divider()
    else:
         st.info("No subsections were extracted for this page.")

    # Display Tables
    tables = page_dict.get('tables', [])
    if tables:
        st.subheader(f"Tables on Page {page_num}")
        for table in tables:
            display_table(table) # Use the existing helper
        st.divider()

    # Display Visuals
    visuals = page_dict.get('visuals', [])
    if visuals:
        st.subheader(f"Visuals on Page {page_num}")
        for visual in visuals:
            display_visual_element(visual) # Use the existing helper
        st.divider()

    # Display Numerical Data
    numbers = page_dict.get('numbers', [])
    if numbers:
        st.subheader(f"Numerical Data on Page {page_num}")
        for num_data in numbers:
            display_numerical_data_point(num_data) # Use new helper
        st.divider()

    # --- REMOVED display of old fields: title, topics, summary, entities, page text area ---


#############################
# STREAMLIT UI COMPONENTS (Updated)
#############################

def render_unified_document_report(document_data: Optional[Dict]):
    """
    Displays the full report for a single processed document, including
    Executive Summary (with full graph), Chapters, and Detailed Page View tabs.

    Args:
        document_data: A dictionary containing the finalized document structure,
                       typically the output from finalize_document. Expected structure:
                       {
                           "raw_extracted_content": {
                               "filename": str,
                               "pages": List[Dict], # List of PageContent dicts
                               "summary": Dict | None, # Dict from DocumentSummaryDetails + chapters
                               "error": str | None
                           }
                       }
                       Can also accept the ProcessedDocument model object directly.
    """
    function_name = "render_unified_document_report"
    if not document_data:
        st.error(f"[{function_name}] No document data provided to display.")
        return

    # --- Data Extraction and Validation ---
    raw_content = {}
    # Handle both dict input and potential Pydantic model input
    if hasattr(document_data, 'model_dump'): # Check if it's a Pydantic model
        logging.debug(f"[{function_name}] Input is a Pydantic model, dumping to dict.")
        try:
            doc = document_data.model_dump(mode='json') # Use mode='json' for potential complex types
        except Exception as dump_err:
            st.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}")
            logging.error(f"[{function_name}] Error dumping Pydantic model: {dump_err}", exc_info=True)
            return
    elif isinstance(document_data, dict):
        doc = document_data
    else:
        st.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        logging.error(f"[{function_name}] Invalid document data type provided: {type(document_data)}")
        return

    # Safely get raw_content, checking structure
    if isinstance(doc, dict) and "raw_extracted_content" in doc and isinstance(doc["raw_extracted_content"], dict):
        raw_content = doc["raw_extracted_content"]
        logging.debug(f"[{function_name}] Successfully extracted raw_extracted_content.")
    # Add fallback if top-level dict *is* raw_content (less likely with finalize_document output)
    elif isinstance(doc, dict) and "filename" in doc and "pages" in doc:
        logging.warning(f"[{function_name}] Input dict seems to be raw_extracted_content directly. Using it.")
        raw_content = doc
    else:
        st.error(f"[{function_name}] Could not find 'raw_extracted_content' dictionary in the provided data.")
        logging.error(f"[{function_name}] Invalid data structure, missing 'raw_extracted_content'. Data keys: {list(doc.keys()) if isinstance(doc, dict) else 'N/A'}")
        return

    # Extract key components, providing defaults
    filename = raw_content.get('filename', 'Unknown Document')
    processing_error = raw_content.get('error')
    pages = raw_content.get('pages', []) # Should be list of dicts
    summary_data = raw_content.get('summary', {}) # Should be a dict
    if not isinstance(summary_data, dict):
        logging.warning(f"[{function_name}] 'summary' data for '{filename}' is not a dict ({type(summary_data)}). Resetting to empty dict.")
        summary_data = {} # Ensure it's a dict for safe access later
    chapters_list = summary_data.get('chapters', []) # Get chapter list

    # --- Display Processing Error (if any) ---
    if processing_error:
        st.error(f"🚨 **Processing Error for {filename}:** {processing_error}")
        st.warning("Displayed data below might be incomplete or relate to a previous successful run.")
        # Decide whether to continue rendering or stop here
        # return # Option: Stop rendering if there was a critical error

    # --- Header Section ---
    st.header(f"📄 Report: {filename}")

    # Count actual elements present in pages for the caption
    num_pages = len(pages) if isinstance(pages, list) else 0
    num_tables = 0
    num_visuals = 0
    num_numbers = 0
    if isinstance(pages, list):
        for p in pages:
            if isinstance(p, dict):
                 num_tables += len(p.get('tables', [])) if isinstance(p.get('tables'), list) else 0
                 num_visuals += len(p.get('visuals', [])) if isinstance(p.get('visuals'), list) else 0
                 num_numbers += len(p.get('numbers', [])) if isinstance(p.get('numbers'), list) else 0

    caption_parts = [f"📄 {num_pages} pages"]
    if num_tables > 0: caption_parts.append(f"📊 {num_tables} tables")
    if num_visuals > 0: caption_parts.append(f"🖼️ {num_visuals} visuals")
    if num_numbers > 0: caption_parts.append(f"🔢 {num_numbers} numbers")
    st.caption(" | ".join(caption_parts))

    # Create tabs for document view modes
    tab_titles = ["Executive Summary", "Chapters", "Detailed Page View"]
    tab1, tab2, tab3 = st.tabs(tab_titles)

    # --- Tab 1: Executive Summary ---
    with tab1:
        # Display Summary, Themes, Questions
        with st.container(border=True):
            st.subheader("Executive Summary")
            summary_title = summary_data.get('title')
            if summary_title and summary_title != f"Summary Failed: {filename}" and summary_title != f"Error Finalizing {filename}":
                 st.markdown(f"**Title:** {summary_title}") # Display LLM generated title if valid

            summary_text = summary_data.get('summary')
            if summary_text and not summary_text.startswith("Failed to generate"):
                st.markdown(summary_text)
            else:
                st.info("No summary text available or summary generation failed.")

            themes = summary_data.get('themes')
            if themes and themes != ["error"]:
                st.markdown("#### Key Themes")
                if hasattr(st, 'chip'):
                    for theme in themes: st.chip(theme, icon="🏷️")
                else: # Fallback display
                     st.markdown("- " + "\n- ".join(themes))

            questions = summary_data.get('questions')
            if questions:
                 st.markdown("#### Key Questions Answered")
                 for q in questions: st.markdown(f"- {q}")

        # --- Display Document-Level Graph (From Session State) ---
        st.divider()
        st.subheader("Document Concept & Entity Network (Overview)")

        all_document_nodes = []
        if isinstance(chapters_list, list):
            for chapter_dict in chapters_list:
                 if isinstance(chapter_dict, dict):
                     nodes_in_chapter = chapter_dict.get('qdrant_nodes', [])
                     if isinstance(nodes_in_chapter, list):
                         all_document_nodes.extend(nodes_in_chapter)

        if all_document_nodes:
            logging.info(f"[render_unified_document_report] Rendering network graph with {len(all_document_nodes)} nodes for {filename}")
            with st.expander("View Full Document Network Graph", expanded=False):
                 try:
                     # Assuming display_qdrant_node_network accepts list of dicts
                     display_qdrant_node_network(
                         all_document_nodes,
                         title=f"Overall Concept/Entity Network for {filename}"
                     )
                 except Exception as graph_err:
                      st.error(f"Failed to render document network graph: {graph_err}")
                      logging.error(f"Failed to render document network graph for {filename}: {graph_err}", exc_info=True)
        else:
            st.caption("No concept or entity data available to build a document network.")

        # --- NEW: Qdrant Node Relationship Explorer (Document Level) ---
        st.divider()
        st.subheader("Explore Node Relationships (This Document)")

        qdrant_client = get_qdrant_client() # Get cached client

        if qdrant_client is None:
            st.error("Qdrant client is unavailable. Cannot explore relationships.")
        elif not all_document_nodes:
            st.info("No nodes available in this document to explore.")
        else:
            # Populate selectbox using nodes from THIS document only
            doc_node_options: Dict[str, str] = {} # label -> node_id
            try:
                # Sort nodes from this document for the dropdown
                sorted_doc_nodes = sorted(
                    all_document_nodes,
                    key=lambda n: n.get('name', '') if isinstance(n, dict) else ''
                )
            except Exception as sort_err:
                logging.warning(f"Failed to sort document nodes for dropdown: {sort_err}. Using original order.")
                sorted_doc_nodes = all_document_nodes

            for node_dict in sorted_doc_nodes:
                node_id = node_dict.get('node_id')
                node_name = node_dict.get('name', node_id)
                node_type = node_dict.get('node_type', 'unknown')
                if node_id:
                    display_label = f"{node_name} ({node_type})"
                    # Handle duplicates (unlikely within single doc if IDs are UUIDs)
                    count = 1; original_label = display_label
                    while display_label in doc_node_options: display_label = f"{original_label} ({count})"; count += 1
                    doc_node_options[display_label] = node_id

            if doc_node_options:
                selected_doc_node_label = st.selectbox(
                    "Select a node within this document:",
                    options=["Select a Node..."] + list(doc_node_options.keys()),
                    index=0,
                    key=f"doc_node_explorer_{filename}" # Unique key
                )

                if selected_doc_node_label != "Select a Node...":
                    selected_doc_node_id = doc_node_options.get(selected_doc_node_label)
                    if selected_doc_node_id:
                        # Call the async Qdrant display function
                        try:
                            logging.info(f"Running Qdrant link display for document node: {selected_doc_node_id}")
                            run_async(
                                display_node_links_qdrant,
                                qdrant_client,
                                QDRANT_COLLECTION_NAME,
                                selected_doc_node_id
                            )
                        except Exception as q_disp_err:
                             st.error(f"Failed to display node links from Qdrant: {q_disp_err}")
                             logging.error(f"Error running display_node_links_qdrant for doc node: {q_disp_err}", exc_info=True)
                    else:
                        st.warning("Could not find ID for selected document node.")
            else:
                st.info("No valid nodes found in this document to explore.")
        # --- End Document-Level Qdrant Explorer ---


        # Page Content Preview Section
        st.divider()
        st.subheader("Page Content Preview")
        if pages and isinstance(pages, list):
            # Create options for the selectbox
            page_options = {} # Map label to index
            for i, p in enumerate(pages):
                page_num = p.get('page_number', i + 1) if isinstance(p, dict) else i + 1
                page_options[f"Page {page_num}"] = i

            if page_options:
                selected_page_label = st.selectbox(
                    "Select page to preview content:",
                    options=["Select a page..."] + list(page_options.keys()),
                    index=0, # Default to placeholder
                    key=f"page_nav_summary_{filename}" # Unique key per document
                )
                if selected_page_label != "Select a page...":
                    try:
                        page_idx = page_options[selected_page_label]
                        page_to_preview = pages[page_idx]

                        if isinstance(page_to_preview, dict):
                             with st.container(border=True):
                                  st.markdown(f"**Previewing Page {page_to_preview.get('page_number')}**")
                                  # Show subsection titles and descriptions for the selected page
                                  page_subsections = page_to_preview.get('subsections', [])
                                  if page_subsections and isinstance(page_subsections, list):
                                       st.markdown("**Subsections on this page:**")
                                       for sub_idx, sub in enumerate(page_subsections):
                                           if isinstance(sub, dict):
                                                sub_title = sub.get('title', f'Subsection {sub_idx+1}')
                                                sub_desc = sub.get('description', '')
                                                st.markdown(f"- **{sub_title}**: {sub_desc}")
                                           else: st.caption("- Invalid subsection format")
                                  else:
                                       st.caption("No subsections extracted or available for this page.")
                                  # Optionally show raw text preview if helpful
                                  raw_page_text = page_to_preview.get('raw_text')
                                  if raw_page_text:
                                       with st.expander("View Raw Text (preview)"):
                                             st.text(raw_page_text[:1000] + ("..." if len(raw_page_text) > 1000 else ""))

                        else:
                            st.warning(f"Invalid data format for page index {page_idx}.")

                    except (IndexError, KeyError, ValueError) as e:
                        st.warning(f"Could not select page for preview: {e}")
            else:
                st.info("No pages available in this document.")
        else:
            st.info("No page data available for preview.")


    # --- Tab 2: Chapters ---
    with tab2:
        if chapters_list and isinstance(chapters_list, list):
            chapter_titles = [ch.get('title', f"Chapter {ch.get('order', i+1)}") for i, ch in enumerate(chapters_list) if isinstance(ch, dict)]
            if chapter_titles:
                try:
                    chapter_tabs = st.tabs(chapter_titles)
                    for i, (tab, chapter_dict) in enumerate(zip(chapter_tabs, chapters_list)):
                        # Ensure we only process valid chapter dicts
                        if isinstance(chapter_dict, dict):
                            with tab:
                                try:
                                    # Ensure display_chapter exists and handles dicts
                                    display_chapter(chapter_dict)
                                except NameError:
                                    st.error("`display_chapter` function not found.")
                                    logging.error("`display_chapter` function not found.")
                                    break # Stop trying if function is missing
                                except Exception as chapter_disp_err:
                                     st.error(f"Error displaying chapter '{chapter_dict.get('title')}': {chapter_disp_err}")
                                     logging.error(f"Error displaying chapter {chapter_dict.get('chapter_id')}: {chapter_disp_err}", exc_info=True)

                        else:
                             logging.warning(f"[{function_name}] Skipping non-dictionary item found in chapters list for '{filename}'.")

                except Exception as tabs_err: # Catch errors from st.tabs if titles are invalid etc.
                     st.error(f"Could not create chapter tabs: {tabs_err}")
                     logging.error(f"Could not create chapter tabs for {filename}: {tabs_err}", exc_info=True)
            else:
                 st.info("No valid chapter titles found in the summary data.")
        else:
            st.info("No chapter information available or processed for this document.")
            if summary_data.get("error") and "chapter extraction" in summary_data.get("error","").lower():
                 st.warning(f"Note: Chapter extraction previously failed for this document. Reason: {summary_data['error']}")


    # --- Tab 3: Detailed Page View ---
    with tab3:
        if pages and isinstance(pages, list):
             # Create options mapping Page Number (int) -> Index (int)
             page_num_options: Dict[int, int] = {}
             valid_page_nums: List[int] = [] # Store the actual page numbers (integers)
             for i, p in enumerate(pages):
                 # Ensure page data is dict and page_number is an int
                 if isinstance(p, dict) and isinstance(p.get('page_number'), int):
                      page_num = p['page_number']
                      page_num_options[page_num] = i # Map the integer page number to its list index
                      valid_page_nums.append(page_num)
                 else:
                      logging.warning(f"[{function_name}] Skipping invalid page data at index {i} (page_number missing or not int) for page view selector in '{filename}'. Page data keys: {list(p.keys()) if isinstance(p,dict) else 'N/A'}")

             if valid_page_nums:
                  # Sort page numbers numerically for the selector options
                  sorted_page_nums = sorted(valid_page_nums)

                  # Use format_func to display "Page X", but the options are the integers
                  selected_page_num = st.selectbox(
                       "Select Page Number:",
                       options=sorted_page_nums,       # Pass the list of INT page numbers
                       format_func=lambda x: f"Page {x}", # Display "Page X" to the user
                       index=0,                         # Default selection to the first page
                       key=f"page_detail_selector_{filename}" # Unique key for this selectbox
                  )
                  # 'selected_page_num' now directly holds the chosen INTEGER page number (e.g., 5)

                  # Find the corresponding page dictionary using the selected integer page number
                  try:
                       # REMOVED: selected_num = int(selected_page_num.split()[-1])
                       # USE selected_page_num directly as the key
                       selected_index = page_num_options.get(selected_page_num)

                       if selected_index is not None:
                           selected_page_dict = pages[selected_index]
                           try:
                               # Ensure display_page_details exists and handles dicts
                               display_page_details(selected_page_dict)
                           except NameError:
                                st.error("`display_page_details` function not found.")
                                logging.error("`display_page_details` function not found.")
                           except Exception as page_disp_err:
                                st.error(f"Error displaying details for page {selected_page_num}: {page_disp_err}")
                                logging.error(f"Error displaying details for page {selected_page_num} of {filename}: {page_disp_err}", exc_info=True)
                       else:
                           # This should be unlikely if valid_page_nums was populated correctly
                           st.warning(f"Internal error: Could not find index for selected page number {selected_page_num}.")
                  except Exception as e: # Catch any unexpected errors during index lookup/display
                       st.error(f"Error processing page selection or display: {e}")
                       logging.error(f"Error processing page selection {selected_page_num} or display for {filename}: {e}", exc_info=True)

             else:
                  st.info("No valid page data with page numbers found to display details.")
        else:
             st.info("No page data available to display details.")


# --- Helper function to run the async Qdrant steps sequentially ---
async def _run_qdrant_ingestion_workflow(qdrant_client, nodes, dense_model, sparse_model):
    """Wrapper to run collection creation and upsert."""
    # Ensure collection exists
    collection_ok = await create_collection_if_not_exists(
        qdrant_client,
        QDRANT_COLLECTION_NAME,
        EMB_MODEL_NAME,
        EMB_DIM,
        SPARSE_MODEL_NAME
    )
    if not collection_ok:
        return False # Stop if collection setup fails

    # Perform the upsert
    upsert_ok = await upsert_qdrant_nodes(
        qdrant_client,
        QDRANT_COLLECTION_NAME,
        nodes,
        dense_model,
        sparse_model,
        EMB_MODEL_NAME,
        SPARSE_MODEL_NAME
        # batch_size can be passed if needed, defaults to BATCH_SIZE
    )
    return upsert_ok


# --- Sidebar function (Mostly unchanged, ensure it calls correct ontology function if using simplified) ---
# --- Embedding models ---
@st.cache_resource
def load_models():
    logging.info(f"Loading embedding models ({EMB_MODEL_NAME}, {SPARSE_MODEL_NAME})...")
    start = time.time()
    dense_model = TextEmbedding(EMB_MODEL_NAME)
    sparse_model = SparseTextEmbedding(SPARSE_MODEL_NAME)
    logging.info(f"Embedding models loaded in {time.time() - start:.2f} seconds.")
    return dense_model, sparse_model

def display_sidebar_chat():
    """Manages sidebar: file upload, 3-phase Qdrant control, status."""
    st.sidebar.title("📝 Input & Control")
    st.sidebar.markdown("Upload documents, process them, generate ontology, and manage Qdrant indexing.")

    supported_types = ["pdf", "xlsx", "xls", "docx", "pptx", "csv", "txt", "md", "json", "html", "xml", "py", "js", "css", "java", "c", "cpp", "h", "hpp"]

    # --- Initialize Session State Robustly ---
    default_state = {
        "processed_file_names": set(), "processing_active": False, "processed_documents": [],
        "show_results": False, "project_ontology": None, "show_ontology": False,
        "messages": [], "selected_docs": [],
        # Qdrant Client/Models
        "qdrant_client": None, "dense_model": None, "sparse_model": None,
        # NEW State for 3-Phase Qdrant Ingestion
        "collection_created_phase1a": False,
        "indexing_complete_phase1b": False,
        "hnsw_enabled_phase2": False,
        "collection_has_data": False, # Track if Phase 1b likely succeeded
        # UI element references (will be managed dynamically)
    }
    for key, default in default_state.items():
        if key not in st.session_state: st.session_state[key] = default

    # --- Ensure Qdrant Client and Models are loaded ---
    if st.session_state.qdrant_client is None: st.session_state.qdrant_client = get_qdrant_client()
    if st.session_state.dense_model is None or st.session_state.sparse_model is None:
        st.session_state.dense_model, st.session_state.sparse_model = load_models()

    qdrant_ready = st.session_state.qdrant_client is not None
    models_ready = st.session_state.dense_model is not None and st.session_state.sparse_model is not None

    # --- File Upload ---
    st.sidebar.subheader("📁 Upload & Process Documents")
    uploaded_files = st.sidebar.file_uploader("Select files:", type=supported_types, accept_multiple_files=True, key="sidebar_file_uploader_main")

    files_data_for_processing = []
    new_files_count = 0
    if uploaded_files:
        for file in uploaded_files:
            if file.name not in st.session_state.processed_file_names:
                try:
                    file_content = file.getvalue(); file_ext = file.name.split('.')[-1].lower() if '.' in file.name else ''
                    if not file_content: st.sidebar.warning(f"'{file.name}' is empty."); continue
                    files_data_for_processing.append({"name": file.name, "content": file_content, "type": file_ext})
                    new_files_count += 1
                except Exception as e: st.sidebar.error(f"Error reading {file.name}: {e}")

    # --- Document Processing Button ---
    if new_files_count > 0 and not st.session_state.processing_active:
        if st.sidebar.button(f"Process {new_files_count} New File(s)", key="process_files_btn"):
            if files_data_for_processing:
                with st.sidebar: # Status elements
                    st.session_state.status_container = st.status("Starting document processing...", expanded=True)
                    st.session_state.progress_bar = st.progress(0)
                    st.session_state.time_info = st.empty()
                try:
                    st.session_state.processing_active = True; st.sidebar.warning("Processing started...")
                    processed_results_list = run_async(process_all_documents_async, files_data_for_processing)
                    # Update state (same as before)
                    new_processed_docs=[]; processed_names=set()
                    for r in processed_results_list:
                        if isinstance(r,dict) and "raw_extracted_content" in r:
                            new_processed_docs.append(r); fname=r["raw_extracted_content"].get("filename")
                            if fname: processed_names.add(fname)
                    st.session_state.processed_documents.extend(new_processed_docs)
                    st.session_state.processed_file_names.update(processed_names)
                    st.session_state.show_results = True
                    if 'status_container' in st.session_state and st.session_state.status_container: st.session_state.status_container.update(label="✅ Processing finished!", state="complete", expanded=False)
                    st.sidebar.success(f"Processed {len(processed_names)} file(s).")
                except Exception as e: st.sidebar.error(f"Processing error: {e}"); logging.error(f"Proc error: {e}", exc_info=True)
                finally:
                    st.session_state.processing_active = False
                    for key in ["status_container", "progress_bar", "time_info"]: # Cleanup UI refs
                         if key in st.session_state:
                              try: ui_el = st.session_state.pop(key, None); # ... (potential cleanup)
                              except Exception: pass
                    st.rerun()
            else: st.sidebar.warning("No valid new files.")

    # --- Display Processed Files & View Buttons ---
    if st.session_state.processed_file_names:
        st.sidebar.markdown("---"); st.sidebar.markdown("**Processed Files:**")
        for filename in sorted(list(st.session_state.processed_file_names)):
             has_error = any(d.get("raw_extracted_content", {}).get("filename") == filename and d.get("raw_extracted_content", {}).get("error") for d in st.session_state.processed_documents if isinstance(d, dict))
             st.sidebar.caption(f"{'⚠️' if has_error else '✓'} {filename}")
        st.sidebar.markdown("---")
        col1, col2 = st.sidebar.columns(2)
        with col1:
            if st.button("View Reports", key="view_reports_btn", disabled=st.session_state.processing_active):
                if st.session_state.processed_documents: st.session_state.show_results = True; st.session_state.show_ontology = False; st.rerun()
                else: st.sidebar.warning("No processed data.")
        with col2:
            ontology_exists = st.session_state.project_ontology is not None
            ontology_label = "View Ontology" if ontology_exists else "Gen. Ontology"
            if st.button(ontology_label, key="ontology_btn", disabled=st.session_state.processing_active):
                valid_docs = [d for d in st.session_state.processed_documents if isinstance(d,dict) and not d.get("raw_extracted_content",{}).get("error")]
                if valid_docs:
                    if not ontology_exists:
                        with st.spinner("Generating project ontology..."):
                            try:
                                gemini_client = get_gemini_client()
                                ontology_data = run_async(incrementally_generate_doc_level_ontology, gemini_client, valid_docs)
                                st.session_state.project_ontology = ontology_data
                                if ontology_data: st.sidebar.success("Ontology generated!")
                                else: st.sidebar.error("Ontology generation failed.")
                            except Exception as gen_err: st.sidebar.error(f"Ontology error: {gen_err}"); logging.error(f"Ontology gen failed: {gen_err}", exc_info=True); st.session_state.project_ontology = None
                    if st.session_state.project_ontology: st.session_state.show_ontology = True; st.session_state.show_results = False; st.rerun()
                    elif not ontology_exists: pass
                    else: st.sidebar.warning("Ontology not available.")
                else: st.sidebar.warning("No valid documents for ontology.")

    # --- Qdrant 3-Phase Control ---
    st.sidebar.divider()
    st.sidebar.subheader("⬆️ Qdrant 2-Phase Indexing")
    st.sidebar.caption(f"Target Collection: `{QDRANT_COLLECTION_NAME}`")

    # Check prerequisites for the entire Qdrant workflow
    can_start_qdrant = (
        qdrant_ready and models_ready and
        st.session_state.project_ontology is not None and
        st.session_state.project_ontology.project_graph_nodes is not None and
        not st.session_state.processing_active # Ensure doc processing is not running
    )
    nodes_to_index = st.session_state.project_ontology.project_graph_nodes if can_start_qdrant else []

    # --- Phase 1a Button ---
    phase1a_tooltip = "Creates/recreates the Qdrant collection (m=0, binary quant)."
    if not qdrant_ready: phase1a_tooltip = "Qdrant client unavailable."
    elif not models_ready: phase1a_tooltip = "Embedding models unavailable."

    if st.sidebar.button("Phase 1a: Create Collection", key="qdrant_phase1a_btn",
                          disabled=not qdrant_ready or not models_ready or st.session_state.processing_active,
                          help=phase1a_tooltip):
        with st.spinner("Phase 1a: Creating collection..."):
            try:
                # Reset subsequent phase states
                st.session_state.indexing_complete_phase1b = False
                st.session_state.hnsw_enabled_phase2 = False
                st.session_state.collection_has_data = False
                # Run async function
                # success = run_async(create_collection_phase1a, st.session_state.qdrant_client)
                success = asyncio.run(create_collection_phase1a(st.session_state.qdrant_client))
                st.session_state.collection_created_phase1a = success
                if success: st.sidebar.success("Phase 1a Complete.")
                else: st.sidebar.error("Phase 1a Failed.")
            except Exception as e: st.sidebar.error(f"Phase 1a Error: {e}"); st.session_state.collection_created_phase1a = False
        st.rerun()

    # --- Phase 1b Button ---
    phase1b_disabled = (not st.session_state.collection_created_phase1a or
                        st.session_state.indexing_complete_phase1b or # Don't re-run if complete
                        not can_start_qdrant or # Need ontology data
                        st.session_state.processing_active)
    phase1b_tooltip = "Uploads ontology nodes to Qdrant. This blocks UI per group."
    if not st.session_state.collection_created_phase1a: phase1b_tooltip = "Complete Phase 1a first."
    elif not can_start_qdrant: phase1b_tooltip = "Generate ontology with nodes first."

    if st.sidebar.button("Phase 1b: Bulk Ingest Nodes (FG)", key="qdrant_phase1b_btn",
                          disabled=phase1b_disabled, help=phase1b_tooltip):
        if nodes_to_index:
            # The function run_bulk_ingest_foreground_nodes_with_ui handles its own UI updates (spinner not needed here)
            try:
                # Call the synchronous orchestrator directly
                success, time_taken, has_data = run_bulk_ingest_foreground_nodes_with_ui(
                    st.session_state.qdrant_client,
                    nodes_to_index,
                    st.session_state.dense_model,
                    st.session_state.sparse_model
                )
                st.session_state.indexing_complete_phase1b = success
                st.session_state.collection_has_data = has_data
                st.session_state.hnsw_enabled_phase2 = False # Reset Phase 2 status
                if success: st.sidebar.success(f"Phase 1b Complete ({time_taken:.1f}s).")
                else: st.sidebar.error("Phase 1b Failed.")
            except Exception as e: st.sidebar.error(f"Phase 1b Run Error: {e}"); st.session_state.indexing_complete_phase1b = False; st.session_state.collection_has_data = False
            st.rerun()
        else: st.sidebar.warning("No ontology nodes available to ingest.")

    # --- Phase 2 Button ---
    phase2_disabled = (not st.session_state.indexing_complete_phase1b or
                       st.session_state.hnsw_enabled_phase2 or # Don't re-run if complete
                       not can_start_qdrant or # Need ontology data
                       st.session_state.processing_active)
    phase2_tooltip = f"Enables HNSW index (m={DENSE_M_PROD}) for faster search."
    if not st.session_state.indexing_complete_phase1b: phase2_tooltip = "Complete Phase 1b first."

    if st.sidebar.button(f"Phase 2: Enable HNSW (m={DENSE_M_PROD})", key="qdrant_phase2_btn",
                          disabled=phase2_disabled, help=phase2_tooltip):
        with st.spinner(f"Phase 2: Enabling HNSW (m={DENSE_M_PROD})... This may take time."):
            try:
                success = run_async(enable_hnsw_phase2, st.session_state.qdrant_client, DENSE_M_PROD)
                st.session_state.hnsw_enabled_phase2 = success
                if success: st.sidebar.success("Phase 2 Complete.")
                else: st.sidebar.error("Phase 2 Failed.")
            except Exception as e: st.sidebar.error(f"Phase 2 Error: {e}"); st.session_state.hnsw_enabled_phase2 = False
        st.rerun()

    # --- Qdrant Status Display ---
    st.sidebar.markdown("---")
    st.sidebar.markdown("**Qdrant Index Status:**")
    if st.session_state.collection_created_phase1a: st.sidebar.success("✅ Phase 1a: Collection Created")
    else: st.sidebar.warning("⚪ Phase 1a: Collection Not Created")

    if st.session_state.indexing_complete_phase1b: st.sidebar.success("✅ Phase 1b: Nodes Ingested")
    elif st.session_state.collection_created_phase1a: st.sidebar.warning("⚪ Phase 1b: Nodes Not Ingested")
    else: st.sidebar.caption("⚪ Phase 1b: (Requires Collection)")

    if st.session_state.hnsw_enabled_phase2: st.sidebar.success(f"✅ Phase 2: HNSW Enabled (m={DENSE_M_PROD})")
    elif st.session_state.indexing_complete_phase1b: st.sidebar.warning("⚪ Phase 2: HNSW Not Enabled")
    else: st.sidebar.caption("⚪ Phase 2: (Requires Ingest)")

    qdrant_operational = (st.session_state.collection_created_phase1a and
                         st.session_state.indexing_complete_phase1b and
                         st.session_state.hnsw_enabled_phase2 and
                         st.session_state.collection_has_data)
    if qdrant_operational: st.sidebar.caption("✅ Qdrant Index Ready for Search")
    else: st.sidebar.caption("⚠️ Qdrant Index Not Fully Ready")


    # --- Chat Section ---
    st.sidebar.divider()
    st.sidebar.subheader("💬 Document Chat")
    # Document selection for chat context
    if st.session_state.processed_documents:
        # Filter for successfully processed documents for chat selection
        valid_doc_options = []
        for doc_result in st.session_state.processed_documents:
             if isinstance(doc_result, dict) and "raw_extracted_content" in doc_result:
                  raw_content = doc_result["raw_extracted_content"]
                  filename = raw_content.get("filename")
                  # Only include docs without processing errors
                  if filename and not raw_content.get("error"):
                       valid_doc_options.append(filename)

        if valid_doc_options:
            st.session_state.selected_docs = st.sidebar.multiselect(
                "Select documents for chat context (max 3 recommended):",
                options=sorted(list(set(valid_doc_options))),
                default=st.session_state.selected_docs,
                key="doc_context_selector_sidebar",
                max_selections=3 # Limit context size for performance/cost
            )
        else:
            st.sidebar.caption("No successfully processed documents available to reference.")
    else:
        st.sidebar.caption("No documents processed yet.")

    # Chat History & Input Container
    with st.sidebar.container(border=True, height=400): # Adjusted height
        st.markdown("**Chat History**")
        # Display chat history (as before)
        if not st.session_state.get("messages", []):
            st.info("Ask a question below!")
        else:
            for message in st.session_state.messages[-10:]: # Show last 10 messages
                role = message.get("role", "")
                content = message.get("content", "")
                with st.chat_message(role):
                    st.markdown(content)

    # Chat input outside the history container
    if prompt := st.sidebar.chat_input("Ask about selected documents...", key="sidebar_chat_input_main", disabled=st.session_state.processing_active):
        current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        st.session_state.messages.append({"role": "user", "content": prompt, "timestamp": current_time})
        # Call async chat processing
        run_async(process_chat_message, prompt)

async def process_chat_message(user_prompt: str):
    """
    Handles the user's chat message, gets context, calls LLM, and updates history.
    """
    logging.info(f"Processing chat prompt: {user_prompt[:50]}...")

    # 1. Get selected document context (Example: using stored summaries)
    context_text = ""
    selected_doc_names = st.session_state.get("selected_docs", [])
    if selected_doc_names:
        context_parts = [f"--- Context from {name} ---"]
        for doc_result in st.session_state.get("processed_documents", []):
             if isinstance(doc_result, dict):
                 raw_content = doc_result.get("raw_extracted_content", {})
                 filename = raw_content.get("filename")
                 if filename in selected_doc_names:
                     summary_data = raw_content.get("summary", {})
                     if isinstance(summary_data, dict):
                          doc_summary = summary_data.get("summary", "")
                          if doc_summary:
                              context_parts.append(f"Summary for {filename}:\n{doc_summary}\n")
                          # Optionally add key concepts/themes here too
        # Limit context length
        context_text = "\n".join(context_parts)
        MAX_CONTEXT_LEN = 5000 # Adjust as needed
        if len(context_text) > MAX_CONTEXT_LEN:
            context_text = context_text[:MAX_CONTEXT_LEN] + "\n... [Context Truncated]"
        logging.debug(f"Using context from {len(selected_doc_names)} documents.")
    else:
        context_text = "No specific document context selected."
        logging.debug("No document context selected for chat.")

    # 2. Construct LLM Prompt
    full_prompt = f"""You are a helpful assistant answering questions based on the provided document context.

    Context:
    {context_text}

    User Question: {user_prompt}

    Answer based only on the provided context. If the answer isn't in the context, say so.
    Answer:
    """

    # 3. Call LLM
    try:
        # Use a model suitable for chat/Q&A
        # Note: Use the actual client object, not the GenerativeModel instance if using the client directly
        # Adjust based on how get_gemini_client is implemented. Assuming it returns a client object usable like this:
        response = await retry_api_call(
            client.aio.models.generate_content,
            model="gemini-2.5-flash-preview-04-17", # Model suitable for text extraction
            contents=[
                types.Content(parts=[types.Part.from_text(text=full_prompt)]),
            ],
        )

        ai_response_text = "Sorry, I couldn't generate a response."
        if response and hasattr(response, 'text'): # Check response structure
             ai_response_text = response.text
        elif response and hasattr(response, 'candidates') and response.candidates:
             # Handle potential list structure
             first_candidate = response.candidates[0]
             if hasattr(first_candidate, 'content') and hasattr(first_candidate.content, 'parts') and first_candidate.content.parts:
                  ai_response_text = first_candidate.content.parts[0].text


    except Exception as e:
        logging.error(f"Error calling LLM for chat: {e}", exc_info=True)
        ai_response_text = f"Sorry, an error occurred: {e}"

    # 4. Update Chat History
    current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.session_state.messages.append({"role": "assistant", "content": ai_response_text, "timestamp": current_time})

    # 5. Trigger UI Update (Rerun needed after async operation completes)
    st.rerun()


def main():
    """Main function to run the Streamlit application."""
    st.set_page_config(page_title="Multimodal Document Processor", page_icon="📊", layout="wide")

    # Initialize session state keys robustly (now includes Qdrant phases)
    default_session_state = {
        "processed_documents": [], "processed_file_names": set(), "processing_active": False,
        "messages": [], "selected_docs": [], "show_results": False,
        "project_ontology": None, "show_ontology": False,
        "qdrant_client": None, "dense_model": None, "sparse_model": None,
        "collection_created_phase1a": False, "indexing_complete_phase1b": False,
        "hnsw_enabled_phase2": False, "collection_has_data": False,
    }
    for key, default_value in default_session_state.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

    # Pre-load client/models (no change needed here)
    if st.session_state.qdrant_client is None: st.session_state.qdrant_client = get_qdrant_client()
    if st.session_state.dense_model is None: st.session_state.dense_model, st.session_state.sparse_model = load_models()

    # --- SIDEBAR ---
    display_sidebar_chat() # Sidebar now includes 3-phase Qdrant controls

    # --- MAIN CONTENT AREA ---
    st.title("📊 Advanced Document Processor & Qdrant KB")
    st.markdown("Upload docs, generate ontology, index with 2-phase Qdrant method, explore.")
    st.divider()

    # Display Logic (Ontology/Reports view unchanged)
    if st.session_state.get("show_ontology") and st.session_state.project_ontology:
        display_project_ontology(st.session_state.project_ontology)

    elif st.session_state.get("show_results") and st.session_state.processed_documents:
        st.header("Document Reports")
        doc_options = {f"{i+1}: {d['raw_extracted_content'].get('filename','?')}":i for i,d in enumerate(st.session_state.processed_documents) if isinstance(d,dict)}
        selected_idx = 0
        if len(st.session_state.processed_documents) > 1 and doc_options:
             selected_name = st.selectbox("Choose report:",options=list(doc_options.keys()),index=0,key="doc_report_selector")
             selected_idx = doc_options.get(selected_name, 0)
        elif len(st.session_state.processed_documents) == 1 and doc_options: selected_idx = 0
        elif not doc_options and st.session_state.processed_documents: st.warning("No valid docs."); selected_idx = -1
        else: selected_idx = -1

        if selected_idx != -1:
            try: render_unified_document_report(st.session_state.processed_documents[selected_idx])
            except Exception as e: st.error(f"Display error: {e}"); logging.error("Display error", exc_info=True)

    else: # Welcome Screen
        if not st.session_state.processing_active:
             st.info("Upload documents, then manage Qdrant indexing via the sidebar.")
             # (Keep detailed welcome message, maybe update Qdrant description)
             st.markdown("""
             ## 📊 Advanced Multimodal Document Analyzer & Qdrant Indexer
             Leverage Google's Gemini models and Qdrant vector database with optimized 2-phase indexing.

             - **Features**: Broad File Support, Structure Extraction, Chapter Generation, Graphing, Summarization.
             - **Qdrant**: Utilizes a **2-phase indexing strategy** (create m=0 -> ingest -> enable HNSW) with **Binary Quantization** for potentially faster ingest and lower memory usage during search.
             - **Workflow**: Process docs -> Generate Ontology -> Use Sidebar buttons for Qdrant (Phase 1a -> 1b -> 2).

             ### Getting Started:
             1.  **Upload & Process:** Use sidebar controls.
             2.  **Generate Ontology:** Use sidebar button.
             3.  **Manage Qdrant Index:** Use the three "Phase" buttons in the sidebar **in order** (1a -> 1b -> 2) after generating the ontology.
             4.  **Explore:** View reports/ontology, use Qdrant node explorer.
             """)

# --- Entry Point ---
if __name__ == "__main__":
    main()
